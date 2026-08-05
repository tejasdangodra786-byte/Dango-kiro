#!/usr/bin/env python3
"""
MET Comprehensive Presentation v4.1
FIXES: No visual overlapping - strict zone layout
NEW: Therapy comparisons, family/caregiver, comorbid disorders sections
Based on: MET Manual (NIAAA Project MATCH) & NIMHANS Substance Use Disorders
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# LAYOUT RULES (strict - NO overlaps)
# Title bar: y=0 to y=1.0 (height=1.0)
# Content zone: y=1.15 to y=6.45 (5.3 inches available)
# Reference bar: y=6.55 to y=7.1 (height=0.55)
# Side margins: 0.4 left/right, gap between columns: 0.3
# ============================================================

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Strict Y positions
TITLE_Y = Inches(0)
TITLE_H = Inches(1.0)
CONTENT_Y = Inches(1.15)
CONTENT_END = Inches(6.45)
REF_Y = Inches(6.6)
REF_H = Inches(0.5)

# X positions
LEFT_M = Inches(0.4)
RIGHT_M = Inches(0.4)
COL_GAP = Inches(0.3)
FULL_W = Inches(12.5)  # full width content
HALF_W = Inches(6.1)   # half width column
COL2_X = Inches(6.8)   # second column start

COLORS = {
    'deep_blue': RGBColor(0x1B, 0x3A, 0x5C),
    'teal': RGBColor(0x00, 0x7B, 0x83),
    'green': RGBColor(0x2E, 0x7D, 0x32),
    'orange': RGBColor(0xE6, 0x5C, 0x00),
    'purple': RGBColor(0x6A, 0x1B, 0x9A),
    'red': RGBColor(0xC6, 0x28, 0x28),
    'gold': RGBColor(0xF9, 0xA8, 0x25),
    'light_blue': RGBColor(0xBB, 0xDE, 0xFB),
    'light_green': RGBColor(0xC8, 0xE6, 0xC9),
    'light_purple': RGBColor(0xE1, 0xBE, 0xE7),
    'light_orange': RGBColor(0xFF, 0xE0, 0xB2),
    'light_teal': RGBColor(0xB2, 0xDF, 0xDB),
    'light_red': RGBColor(0xFF, 0xCD, 0xD2),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'cream': RGBColor(0xFF, 0xF8, 0xE1),
    'navy': RGBColor(0x0D, 0x47, 0xA1),
    'dark_green': RGBColor(0x1B, 0x5E, 0x20),
    'maroon': RGBColor(0x88, 0x00, 0x38),
}

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]



# ============================================================
# HELPER FUNCTIONS
# ============================================================

def add_bg(slide, c1, c2=None):
    """Add solid or gradient background."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    if c2:
        bg.fill.gradient()
        bg.fill.gradient_stops[0].color.rgb = COLORS[c1]
        bg.fill.gradient_stops[1].color.rgb = COLORS[c2]
    else:
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS[c1]
    bg.line.fill.background()
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def title_bar(slide, text, color='deep_blue'):
    """Title bar: y=0, h=1.0"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), TITLE_Y, SLIDE_W, TITLE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[color]
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Times New Roman'

def ref_bar(slide, text):
    """Reference bar: y=6.6, h=0.5"""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), REF_Y, SLIDE_W, REF_H)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p.font.name = 'Times New Roman'

def cbox(slide, x, y, w, h, lines, bg='white', tc='dark_gray', fs=13, title=None, ttc='deep_blue', border=None):
    """Content box with safe text."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS[bg]
    if border:
        box.line.color.rgb = COLORS[border]
        box.line.width = Pt(1.5)
    else:
        box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.05)
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS[ttc]
        p.font.name = 'Times New Roman'
        p.space_after = Pt(4)
    for i, line in enumerate(lines):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fs)
        p.font.color.rgb = COLORS[tc]
        p.font.name = 'Times New Roman'
        p.space_after = Pt(3)

def divider(title, subtitle, color='deep_blue'):
    """Section divider slide."""
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, color, 'white')
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Times New Roman'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLORS['cream']
    p2.font.name = 'Times New Roman'
    p2.alignment = PP_ALIGN.CENTER
    return slide

def tbl_slide(title, headers, rows, color='deep_blue', ref=''):
    """Table slide with proper spacing."""
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, 'white', 'light_blue')
    title_bar(slide, title, color)
    nr = len(rows) + 1
    nc = len(headers)
    t_top = Inches(1.15)
    t_h = Inches(5.2) if not ref else Inches(5.0)
    table = slide.shapes.add_table(nr, nc, LEFT_M, t_top, FULL_W, t_h).table
    cw = int(FULL_W / nc)
    for i in range(nc):
        table.columns[i].width = cw
    for i, h in enumerate(headers):
        c = table.cell(0, i)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = COLORS[color]
        p = c.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = 'Times New Roman'
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = table.cell(ri + 1, ci)
            c.text = val
            if ri % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['dark_gray']
            p.font.name = 'Times New Roman'
    if ref:
        ref_bar(slide, ref)
    return slide

def flow_boxes(slide, y, items, colors):
    """Horizontal process flow."""
    n = len(items)
    bw = (FULL_W - Inches(0.25) * (n - 1)) / n
    for i, (item, ck) in enumerate(zip(items, colors)):
        x = LEFT_M + (bw + Inches(0.25)) * i
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), y, int(bw), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS[ck]
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = 'Times New Roman'
        p.alignment = PP_ALIGN.CENTER
        if i < n - 1:
            ax = int(x + bw)
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, y + Inches(0.3), Inches(0.25), Inches(0.35))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLORS['gold']
            arr.line.fill.background()



# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'deep_blue', 'navy')
tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10), Inches(4.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "MOTIVATIONAL ENHANCEMENT THERAPY (MET)"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.font.name = 'Times New Roman'
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "A Comprehensive Clinical Guide"
p2.font.size = Pt(22)
p2.font.color.rgb = COLORS['gold']
p2.font.name = 'Times New Roman'
p2.alignment = PP_ALIGN.CENTER
p2.space_after = Pt(30)
p3 = tf.add_paragraph()
p3.text = "Based on the NIAAA Project MATCH MET Manual (Miller et al., 1992)"
p3.font.size = Pt(14)
p3.font.color.rgb = COLORS['light_blue']
p3.font.name = 'Times New Roman'
p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph()
p4.text = "& NIMHANS Substance Use Disorders Manual (2016)"
p4.font.size = Pt(14)
p4.font.color.rgb = COLORS['light_blue']
p4.font.name = 'Times New Roman'
p4.alignment = PP_ALIGN.CENTER
p5 = tf.add_paragraph()
p5.text = ""
p5.space_after = Pt(20)
p6 = tf.add_paragraph()
p6.text = "4-Session Brief Intervention  |  Evidence-Based  |  Client-Centered"
p6.font.size = Pt(13)
p6.font.color.rgb = COLORS['gold']
p6.font.name = 'Times New Roman'
p6.alignment = PP_ALIGN.CENTER
ref_bar(slide, "Miller, W.R., Zweben, A., DiClemente, C.C., & Rychtarik, R.G. (1992). MET Manual. NIAAA Project MATCH Series Vol. 2.")

# ============================================================
# SLIDE 2: TABLE OF CONTENTS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "PRESENTATION OUTLINE", 'deep_blue')
sections = [
    ("1. Introduction to MET", "What is MET, origins, how it differs from other approaches"),
    ("2. Theoretical Foundations", "Transtheoretical Model, Self-Efficacy, Cognitive Dissonance - in detail"),
    ("3. Principles & Spirit of MET", "FRAMES, OARS, DARES, MI Spirit, Handling Resistance"),
    ("4. Session-by-Session Guide", "How to DO the therapy - 4-session protocol with dialogues"),
    ("5. Clinical Techniques", "Reflective listening, decisional balance, rulers, therapist traps"),
    ("6. Case Conceptualization", "Applied case study with session-wise MET intervention plan"),
    ("7. Worksheets & Tools", "6 printable clinical worksheets for MET sessions"),
    ("8. Comparison with Other Therapies", "MET vs CBT, 12-Step, BSFT; family & caregiver approaches"),
    ("9. Comorbid Psychiatric Disorders", "Personality disorders, depression, anxiety - what to do"),
    ("10. Research & Effectiveness", "Project MATCH, UKATT, meta-analyses, Indian research"),
]
for i, (sec_t, sec_d) in enumerate(sections):
    y = Inches(1.2) + Inches(0.52) * i
    colors_list = ['light_blue', 'light_teal', 'light_green', 'light_purple', 'light_orange', 'light_red', 'light_blue', 'light_green', 'light_purple', 'light_orange']
    tcolors = ['deep_blue', 'teal', 'green', 'purple', 'orange', 'maroon', 'deep_blue', 'green', 'purple', 'orange']
    cbox(slide, LEFT_M, y, FULL_W, Inches(0.48), [sec_d], bg=colors_list[i], fs=11, title=sec_t, ttc=tcolors[i], border=tcolors[i])
ref_bar(slide, "Structure: Miller et al. (1992). MET Manual; NIMHANS (2016). Substance Use Disorders.")



# ============================================================
# SECTION 1: INTRODUCTION
# ============================================================
divider("SECTION 1", "Introduction to Motivational Enhancement Therapy", 'deep_blue')

# What is MET
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'cream')
title_bar(slide, "What is Motivational Enhancement Therapy (MET)?", 'deep_blue')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.5),
    ["MET is a brief, systematic intervention designed",
     "to produce rapid, internally motivated change.",
     "It does NOT guide the client step-by-step through",
     "recovery. Instead, it uses motivational strategies",
     "to mobilize the client's OWN change resources.",
     "",
     "Key: 4 planned sessions over 12 weeks."],
    bg='light_blue', title="Definition", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.5),
    ["- Developed for NIAAA Project MATCH (1992)",
     "- Based on Motivational Interviewing principles",
     "- Uses personal feedback + MI techniques",
     "- Designed for alcohol/substance use disorders",
     "- Adaptable to various clinical settings",
     "- Therapist as collaborator, NOT expert"],
    bg='light_green', title="Key Facts", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["\"MET is based on principles of motivational psychology and is designed to produce rapid,",
     "internally motivated change. This treatment strategy does not attempt to guide and train the",
     "client, step by step, through recovery, but instead employs motivational strategies to mobilize",
     "the client's own change resources.\" (MET Manual, p. 1)",
     "",
     "Core idea: The motivation and resources for change ALREADY EXIST within the client.",
     "The therapist's job is to create conditions that help these emerge naturally."],
    bg='cream', title="From the Manual:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Miller, W.R. et al. (1992). MET Manual, p. 1. NIAAA Project MATCH Monograph Series, Vol. 2.")

# How MET Differs
tbl_slide("How MET Differs from Other Approaches",
    ["Feature", "Traditional Approaches", "MET Approach"],
    [["Therapist Role", "Expert / Teacher / Director", "Collaborative partner; elicits client's own motivation"],
     ["Client Role", "Passive recipient of treatment", "Active agent of own change"],
     ["View of Resistance", "Denial to be confronted", "Signal to change therapeutic strategy"],
     ["Goals", "Set by therapist or program", "Negotiated with client based on their values"],
     ["Core Techniques", "Skills training, education, confrontation", "Reflective listening, feedback, exploring ambivalence"],
     ["Duration", "Usually 12+ sessions", "4 structured sessions over 12 weeks"],
     ["Change Mechanism", "Teaching new skills/information", "Mobilizing client's own internal resources"]],
    color='green',
    ref="Ref: Miller et al. (1992). MET Manual, pp. 1-5; NIMHANS (2016). Substance Use Disorders, Ch. 8.")

# Origins
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_purple')
title_bar(slide, "Origins and Development of MET", 'purple')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.5),
    ["1983: Miller publishes foundational MI paper",
     "1991: Miller & Rollnick publish first MI book",
     "1992: MET Manual for Project MATCH",
     "  - Largest alcohol treatment trial ever",
     "  - 1,726 participants across 9 US sites",
     "  - Compared MET vs CBT vs 12-Step (TSF)"],
    bg='light_purple', title="Timeline", ttc='purple', border='purple', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.5),
    ["MET = MI + Personalized Assessment Feedback",
     "",
     "1. MOTIVATIONAL INTERVIEWING (MI):",
     "   Client-centered counseling style",
     "   Explores and resolves ambivalence",
     "2. PERSONALIZED FEEDBACK:",
     "   Objective data presented to client",
     "   Creates discrepancy with goals/values"],
    bg='cream', title="Core Components", ttc='orange', border='orange', fs=12)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["Project MATCH Finding: MET achieved comparable outcomes to 12-session CBT and 12-session TSF",
     "in just 4 sessions. This was revolutionary - a brief motivational intervention matching efficacy",
     "of longer treatments. Particularly effective for clients HIGH in anger (non-confrontational approach)",
     "and those LOW in readiness to change (meets them where they are)."],
    bg='light_green', title="Key Finding:", ttc='green', border='green', fs=12)
ref_bar(slide, "Ref: Project MATCH Research Group (1997). J Studies on Alcohol, 58, 7-29; Miller (1983). Behav Psychother, 11, 147-172.")



# ============================================================
# SECTION 2: THEORETICAL FOUNDATIONS
# ============================================================
divider("SECTION 2", "Theoretical Foundations of MET", 'teal')

# TTM Overview
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_teal')
title_bar(slide, "The Transtheoretical Model (TTM) of Change", 'teal')
cbox(slide, LEFT_M, Inches(1.2), FULL_W, Inches(1.8),
    ["The Transtheoretical Model (Prochaska & DiClemente, 1982) is the PRIMARY theoretical foundation",
     "of MET. It describes behavior change as a PROCESS through predictable stages. People do not change",
     "abruptly. The therapist's task is to match interventions to the client's current stage.",
     "\"The responsibility and capability for change lie within the client.\" (MET Manual, p. 2)"],
    bg='cream', title="Foundation of MET", ttc='teal', border='teal', fs=12)
flow_boxes(slide, Inches(3.2),
    ["Pre-\ncontemplation", "Contemplation", "Preparation", "Action", "Maintenance", "Relapse"],
    ['red', 'orange', 'gold', 'green', 'teal', 'purple'])
cbox(slide, LEFT_M, Inches(4.5), FULL_W, Inches(1.8),
    ["Key Principle: Change is CYCLICAL, not linear. Most people cycle through stages 3-7 times",
     "before achieving lasting change. Relapse is a NORMAL part of the process, not failure.",
     "Average smoker cycles through stages 3 times before permanent quit (Prochaska et al., 1992).",
     "MET is designed to move people forward through these stages using motivational techniques."],
    bg='light_red', title="Important:", ttc='red', border='red', fs=12)
ref_bar(slide, "Ref: Prochaska & DiClemente (1984). The Transtheoretical Approach; MET Manual (1992), pp. 6-12.")

# Stage 1: Precontemplation
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_red')
title_bar(slide, "Stage 1: PRECONTEMPLATION - Not Yet Considering Change", 'red')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.5),
    ["Person does NOT see substance use as a problem.",
     "Not thinking about change at all.",
     "",
     "Common statements:",
     "\"I don't have a problem\"",
     "\"Everyone drinks like me\"",
     "\"My family is overreacting\"",
     "\"I'm here because the court sent me\""],
    bg='light_red', title="What It Looks Like:", ttc='red', border='red', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.5),
    ["DO: Raise doubt gently, increase awareness",
     "DO: Provide objective information/feedback",
     "DO: Explore events that brought them in",
     "DO: Plant seeds without pushing",
     "",
     "DON'T: Argue or confront",
     "DON'T: Push for immediate change",
     "DON'T: Label them as 'alcoholic'/'addict'"],
    bg='cream', title="Therapist Tasks:", ttc='orange', border='orange', fs=12)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["MET Technique: Present Personal Feedback Report (PFR) showing objective data.",
     "Example: \"Your liver enzymes are at ___ level. Here's where that falls compared to norms...\"",
     "The therapist presents FACTS without arguing, letting the DATA create discrepancy.",
     "",
     "\"The precontemplation stage is characterized by a lack of awareness that a problem exists.",
     "The individual has no intention to change behavior in the foreseeable future.\" (Manual, p. 7)"],
    bg='white', title="How MET Addresses This Stage:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 6-8; Prochaska & DiClemente (1982). Psychotherapy, 19, 276-288.")

# Stage 2: Contemplation
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_orange')
title_bar(slide, "Stage 2: CONTEMPLATION - Thinking About Change", 'orange')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.5),
    ["Person ACKNOWLEDGES a problem but is",
     "AMBIVALENT. Sees both pros and cons.",
     "",
     "Common statements:",
     "\"I know I drink too much, but...\"",
     "\"I want to quit, but I don't know how\"",
     "\"Sometimes I think I should cut down\"",
     "\"Part of me wants to change, part doesn't\""],
    bg='light_orange', title="What It Looks Like:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.5),
    ["DO: Explore ambivalence (don't resolve it for them)",
     "DO: Use Decisional Balance worksheet",
     "DO: Elicit self-motivational statements",
     "DO: Tip the balance toward change",
     "DO: Develop discrepancy with values",
     "",
     "DON'T: Rush to action planning",
     "DON'T: Tell them what to do"],
    bg='light_green', title="Therapist Tasks:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["MET Key Strategy: EXPLORING AMBIVALENCE with Decisional Balance",
     "Help client articulate BOTH sides - reasons to change AND reasons to stay the same.",
     "\"Ambivalence is the contemplator's defining feature. The person is simultaneously considering",
     "and rejecting change... The therapist's task is to tip the decisional balance.\" (Manual, p. 8)",
     "",
     "Research: DiClemente et al. (1991) showed contemplators who explored ambivalence progressed faster."],
    bg='cream', title="MET Strategy:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 8-10; DiClemente et al. (1991). J Consult Clin Psychol, 59, 295-304.")

# Stage 3: Preparation
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'cream')
title_bar(slide, "Stage 3: PREPARATION - Ready to Plan", 'gold')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.5),
    ["Person has DECIDED to change and is",
     "planning how. Balance has tipped.",
     "",
     "Common statements:",
     "\"I need to do something about this\"",
     "\"What are my options?\"",
     "\"I'm going to quit next Monday\"",
     "\"I've already started cutting down\""],
    bg='light_orange', title="What It Looks Like:", ttc='gold', border='gold', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.5),
    ["DO: Help develop a Change Plan",
     "DO: Offer a MENU of options (not prescribe)",
     "DO: Support self-efficacy strongly",
     "DO: Help set realistic, achievable goals",
     "DO: Negotiate (not impose) a plan",
     "",
     "DON'T: Prescribe a single path",
     "DON'T: Miss the window of readiness"],
    bg='light_green', title="Therapist Tasks:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["MET Key Strategy: CHANGE PLAN WORKSHEET",
     "Help client complete: (1) Changes desired (2) Reasons for change (3) Steps planned",
     "(4) How others can help (5) How they'll know it's working (6) Things that could interfere",
     "",
     "\"The window of determination is open for a period of time. If action is not taken, the person",
     "may slip back into contemplation or precontemplation.\" (MET Manual, p. 10)"],
    bg='light_blue', title="MET Strategy:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 10-11; DiClemente et al. (1991). J Consult Clin Psychol, 59, 295-304.")

# Stage 4-6: Action, Maintenance, Relapse
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_green')
title_bar(slide, "Stages 4-6: Action, Maintenance, and Relapse", 'green')
cbox(slide, LEFT_M, Inches(1.2), Inches(4.0), Inches(2.5),
    ["ACTION: Person is actively",
     "making changes to behavior.",
     "",
     "Therapist tasks:",
     "- Affirm efforts",
     "- Troubleshoot obstacles",
     "- Review Change Plan",
     "- Build self-efficacy"],
    bg='light_green', title="Action", ttc='green', border='green', fs=11)
cbox(slide, Inches(4.7), Inches(1.2), Inches(4.0), Inches(2.5),
    ["MAINTENANCE: Sustaining",
     "gains over time.",
     "",
     "Therapist tasks:",
     "- Prevent complacency",
     "- Identify high-risk situations",
     "- Plan coping strategies",
     "- Celebrate milestones"],
    bg='light_teal', title="Maintenance", ttc='teal', border='teal', fs=11)
cbox(slide, Inches(9.0), Inches(1.2), Inches(4.0), Inches(2.5),
    ["RELAPSE: Return to earlier",
     "stage - NOT a failure.",
     "",
     "Therapist tasks:",
     "- Normalize (3-7 cycles avg)",
     "- Explore without blame",
     "- Revise the Change Plan",
     "- Re-engage motivation"],
    bg='light_red', title="Relapse", ttc='red', border='red', fs=11)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["\"The Transtheoretical Model views relapse not as catastrophic failure but as a normal part of the",
     "cycle of change. With each cycle, the person learns something new and may achieve more stable",
     "maintenance in the next attempt.\" (MET Manual, p. 12)",
     "",
     "MET Sessions 3-4 focus on these stages: reviewing progress, affirming successes, problem-solving",
     "obstacles, and normalizing any slips as information rather than defeat."],
    bg='cream', title="From the Manual:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 11-14; Prochaska, DiClemente & Norcross (1992). Am Psychologist, 47, 1102-1114.")



# Self-Efficacy
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_purple')
title_bar(slide, "Theoretical Foundation: Self-Efficacy (Bandura, 1977)", 'purple')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.5),
    ["Self-efficacy = person's BELIEF in ability to",
     "successfully perform a behavior.",
     "",
     "In MET: belief that one CAN change.",
     "",
     "Sources (Bandura):",
     "1. Past performance accomplishments",
     "2. Vicarious experience (others' success)",
     "3. Verbal persuasion (therapist affirmation)",
     "4. Physiological/emotional states"],
    bg='light_purple', title="What is Self-Efficacy?", ttc='purple', border='purple', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.5),
    ["How MET Builds Self-Efficacy:",
     "",
     "1. AFFIRMING client's strengths/past successes",
     "   \"You stayed sober 3 months - real strength\"",
     "2. SUPPORTING ability to change",
     "   \"You have what it takes to do this\"",
     "3. OFFERING menu of options (sense of control)",
     "4. CELEBRATING small wins",
     "5. Asking \"Why not a lower number?\" on rulers"],
    bg='light_green', title="MET Application:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["\"A person who lacks confidence that he or she can change is unlikely to try. The therapist's",
     "task is to enhance the client's belief in the possibility of change.\" (MET Manual, p. 17)",
     "",
     "Research: Ilgen et al. (2005) found that self-efficacy measured at treatment entry predicted",
     "drinking outcomes at 1-year follow-up. Clients with higher self-efficacy had better outcomes.",
     "This is the 'S' in FRAMES: Support self-efficacy."],
    bg='cream', title="Manual Quote & Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Bandura (1977). Psychol Review, 84; MET Manual (1992), pp. 16-17; Ilgen et al. (2005). Addictive Behaviors, 30.")

# Cognitive Dissonance
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_orange')
title_bar(slide, "Theoretical Foundation: Cognitive Dissonance (Festinger, 1957)", 'orange')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.5),
    ["Cognitive Dissonance = uncomfortable tension",
     "when behavior contradicts values/beliefs.",
     "",
     "Substance use examples:",
     "- \"I am a good parent\" vs \"I drink every",
     "   night and my children suffer\"",
     "- \"I value my health\" vs \"I am damaging",
     "   my liver with alcohol\"",
     "- \"I am independent\" vs \"I can't function",
     "   without my substance\""],
    bg='light_orange', title="The Theory:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.5),
    ["How MET Uses Cognitive Dissonance:",
     "",
     "1. DEVELOPING DISCREPANCY (key principle)",
     "   Show gap between IS and WANT TO BE",
     "2. Personal Feedback creates dissonance:",
     "   \"You said health is your top value...\"",
     "   \"Your liver enzymes show damage...\"",
     "3. Client must resolve the tension",
     "4. Resolution happens through CHANGE",
     "   (not through therapist confrontation)"],
    bg='light_blue', title="MET Application:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["\"A discrepancy between present behavior and important personal goals will motivate change.",
     "The therapist's task is to develop and amplify such discrepancy.\" (MET Manual, p. 14)",
     "",
     "Research: Draycott & Dabbs (1998) confirmed that discrepancy awareness predicts behavior change.",
     "McNally et al. (2005) showed that developing discrepancy is the most potent MI technique for",
     "moving clients from contemplation to preparation stage."],
    bg='cream', title="Manual Quote & Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Festinger (1957). Cognitive Dissonance; MET Manual (1992), pp. 13-15; McNally et al. (2005). Drug Alcohol Rev.")



# ============================================================
# SECTION 3: PRINCIPLES & SPIRIT
# ============================================================
divider("SECTION 3", "Principles and Spirit of MET", 'green')

# Five Principles
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_green')
title_bar(slide, "Five Principles of Motivational Interviewing in MET", 'green')
principles = [
    ("1. Express Empathy", "Accept the client where they are. Ambivalence is NORMAL. Use reflective listening.", 'light_blue', 'deep_blue'),
    ("2. Develop Discrepancy", "Help client see gap between current behavior and important goals/values.", 'light_green', 'green'),
    ("3. Avoid Argumentation", "Arguments are counterproductive. Defending breeds defensiveness.", 'light_orange', 'orange'),
    ("4. Roll with Resistance", "Don't fight resistance - use it. Offer new perspectives without imposing.", 'light_purple', 'purple'),
    ("5. Support Self-Efficacy", "Client's belief in possibility of change is a key motivator.", 'light_teal', 'teal'),
]
for i, (t, d, bg_c, t_c) in enumerate(principles):
    y = Inches(1.15) + Inches(1.05) * i
    cbox(slide, LEFT_M, y, FULL_W, Inches(0.95), [d], bg=bg_c, title=t, ttc=t_c, border=t_c, fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 13-17; Miller & Rollnick (1991). Motivational Interviewing, Ch. 3.")

# FRAMES
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "FRAMES: Elements of Effective Brief Interventions", 'deep_blue')
frames = [
    ("F - Feedback", "Personal feedback about risk/impairment based on objective assessment data", 'light_red', 'red'),
    ("R - Responsibility", "Emphasize responsibility for change lies with the CLIENT", 'light_orange', 'orange'),
    ("A - Advice", "Clear advice to change given in a non-prescriptive manner", 'cream', 'gold'),
    ("M - Menu", "Offer a MENU of strategies/options - client chooses their own path", 'light_green', 'green'),
    ("E - Empathy", "Warm, reflective, empathic counseling style without judgment", 'light_blue', 'deep_blue'),
    ("S - Self-Efficacy", "Reinforce hope and belief that change IS possible for this person", 'light_purple', 'purple'),
]
for i, (t, d, bg_c, t_c) in enumerate(frames):
    y = Inches(1.15) + Inches(0.88) * i
    cbox(slide, LEFT_M, y, FULL_W, Inches(0.8), [d], bg=bg_c, title=t, ttc=t_c, border=t_c, fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 16-17; Miller & Sanchez (1994). Motivating young adults for change.")

# OARS
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_teal')
title_bar(slide, "OARS: Core Microskills of MET", 'teal')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["OPEN-ENDED QUESTIONS:",
     "- Cannot be answered with yes/no",
     "- \"What concerns you about your drinking?\"",
     "- \"How has substance use affected your life?\"",
     "",
     "AFFIRMATIONS:",
     "- Recognize strengths and efforts",
     "- \"It took courage to come here today\"",
     "- \"You've shown real resilience\""],
    bg='light_blue', title="O - Open Questions & A - Affirmations", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["REFLECTIVE LISTENING (Primary skill):",
     "- Repeat back MEANING (not just words)",
     "- Simple: \"So you feel frustrated\"",
     "- Complex: \"You want to quit but worry",
     "  about losing your social circle\"",
     "",
     "SUMMARIES:",
     "- Collect and link what client has said",
     "- \"Let me see if I have this right...\""],
    bg='light_green', title="R - Reflections & S - Summaries", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["PRACTICE RATIO: Aim for 2-3 reflections for every question asked.",
     "Most common mistake: asking too many questions (feels like interrogation).",
     "",
     "\"The principal technique for expressing empathy is reflective listening... seeking through",
     "your responses to understand the client's meaning and feelings.\" (MET Manual, p. 21)",
     "",
     "Research: Moyers et al. (2005) showed therapist MI-consistent behaviors (especially reflections)",
     "predicted better client outcomes. The quality of MI delivery matters."],
    bg='cream', title="Clinical Tip:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 20-35; Moyers et al. (2005). J Subst Abuse Treat, 28, 19-26.")

# DARES
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_orange')
title_bar(slide, "DARES: Self-Motivational Statements to Elicit", 'orange')
cbox(slide, LEFT_M, Inches(1.2), FULL_W, Inches(1.3),
    ["Goal: ELICIT (not provide) self-motivational statements. The more the CLIENT voices arguments",
     "FOR change, the more likely actual change becomes. Listen for and actively evoke these:"],
    bg='cream', title="Core Principle:", ttc='orange', border='gold', fs=12)
dares = [
    ("D - Desire", "\"I want to change\" / \"I wish I could stop\"", 'light_blue'),
    ("A - Ability", "\"I think I could do it\" / \"I was able to quit before\"", 'light_green'),
    ("R - Reasons", "\"My health is suffering\" / \"My family needs me sober\"", 'light_purple'),
    ("E - Emotional", "\"I'm tired of feeling this way\" / \"I hate what I've become\"", 'light_orange'),
    ("S - Steps", "\"I could try meetings\" / \"Maybe I'll call my doctor\"", 'light_teal'),
]
for i, (t, ex, bg_c) in enumerate(dares):
    y = Inches(2.7) + Inches(0.72) * i
    cbox(slide, LEFT_M, y, FULL_W, Inches(0.62), [f"Example: {ex}"], bg=bg_c, title=t, ttc='dark_gray', fs=11)
ref_bar(slide, "Ref: MET Manual (1992), pp. 25-30; Amrhein et al. (2003). J Consult Clin Psychol, 71, 862-878.")

# Spirit of MI
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_green')
title_bar(slide, "The Spirit of Motivational Interviewing", 'green')
cbox(slide, LEFT_M, Inches(1.2), Inches(4.0), Inches(2.5),
    ["Working WITH the client as",
     "equal partners. NOT",
     "expert-to-patient.",
     "",
     "\"The therapeutic relationship",
     "is more like a partnership",
     "than an expert-recipient",
     "one.\" (Manual, p. 13)"],
    bg='light_blue', title="Collaboration", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, Inches(4.7), Inches(1.2), Inches(4.0), Inches(2.5),
    ["Drawing out client's OWN",
     "motivations, strengths,",
     "and reasons for change.",
     "",
     "Motivation is EVOKED",
     "from within, not installed",
     "from outside. The client",
     "has the answers."],
    bg='light_green', title="Evocation", ttc='green', border='green', fs=12)
cbox(slide, Inches(9.0), Inches(1.2), Inches(4.0), Inches(2.5),
    ["Respecting client's right",
     "and capacity to direct",
     "their own life.",
     "",
     "Client decides whether,",
     "when, and how to change.",
     "Even if they choose NOT",
     "to change right now."],
    bg='light_purple', title="Autonomy", ttc='purple', border='purple', fs=12)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["In Practice: Therapist asks more than tells; listens more than instructs.",
     "Client does most of the talking (aim for 70:30 client:therapist talk ratio).",
     "No labeling, no shaming, no arguing for change on behalf of the client.",
     "Resistance is met with curiosity, not confrontation.",
     "",
     "Research: Miller et al. (1993) found therapist directiveness predicted worse outcomes.",
     "Client resistance increased when therapists were confrontational."],
    bg='cream', title="What It Looks Like:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 13-18; Miller et al. (1993). J Stud Alcohol, 54, 455-461.")

# Handling Resistance
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_red')
title_bar(slide, "Handling Client Resistance: Strategies from the Manual", 'red')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.5),
    ["TYPES OF RESISTANCE:",
     "1. Arguing - challenging, hostility",
     "2. Interrupting - cutting off, talking over",
     "3. Denying - blaming, excusing, minimizing",
     "4. Ignoring - inattention, sidetracking",
     "",
     "KEY: Resistance is a signal to CHANGE",
     "your approach, not fight harder."],
    bg='light_red', title="Recognizing Resistance:", ttc='red', border='red', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.5),
    ["STRATEGIES TO ROLL WITH IT:",
     "1. Simple Reflection - acknowledge",
     "2. Amplified Reflection - stronger form",
     "3. Double-Sided Reflection - both sides",
     "4. Shifting Focus - redirect topic",
     "5. Agreement with Twist - agree then reframe",
     "6. Reframing - new interpretation",
     "7. Emphasize Personal Choice - \"It's up to you\""],
    bg='light_green', title="Rolling with Resistance:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.9), FULL_W, Inches(2.4),
    ["Client: \"I don't think I drink more than my friends.\"",
     "Simple Reflection: \"You see yourself as a pretty normal drinker.\"",
     "Amplified: \"So there's really nothing at all to be concerned about.\"",
     "Double-Sided: \"Your drinking feels normal to you, AND something prompted you to come in.\"",
     "Reframe: \"Your friends are important and you want to fit in with them.\"",
     "",
     "Research: Resistance in session predicts WORSE outcomes (Miller et al., 1993)."],
    bg='cream', title="Clinical Examples:", ttc='orange', border='gold', fs=11)
ref_bar(slide, "Ref: MET Manual (1992), pp. 32-38; Miller et al. (1993). J Stud Alcohol, 54, 455-461.")



# ============================================================
# SECTION 4: SESSION-BY-SESSION GUIDE
# ============================================================
divider("SECTION 4", "Session-by-Session Therapy Guide\nHow to Actually DO MET", 'purple')

# Overview
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_purple')
title_bar(slide, "MET: The 4-Session Structure (Overview)", 'purple')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["SESSION 1 (Week 1):",
     "- Build rapport and therapeutic alliance",
     "- Present Personal Feedback Report",
     "- Explore client's reaction to feedback",
     "- Elicit self-motivational statements",
     "- Gauge readiness to change",
     "",
     "SESSION 2 (Week 2):",
     "- Strengthen commitment to change",
     "- Develop Change Plan (collaborative)",
     "- Set specific goals"],
    bg='light_blue', title="Sessions 1 & 2", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["SESSION 3 (Week 6):",
     "- Review progress since Session 2",
     "- Renew motivation and commitment",
     "- Problem-solve obstacles",
     "- Revise Change Plan if needed",
     "",
     "SESSION 4 (Week 12):",
     "- Review overall progress",
     "- Consolidate gains",
     "- Plan for maintenance/relapse prevention",
     "- Termination and future planning"],
    bg='light_orange', title="Sessions 3 & 4", ttc='orange', border='orange', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["KEY PRINCIPLE: Each session uses MI techniques throughout. MET is not a checklist but a clinical",
     "style applied within a structured framework. The therapist continuously uses OARS skills, evokes",
     "change talk, rolls with resistance, and supports self-efficacy in every interaction.",
     "",
     "Session spacing is intentional: Session 1-2 close together (build momentum), then space out to",
     "allow client time to implement changes and experience results before reviewing progress."],
    bg='cream', title="Important Note:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 40-80; Chapters III-V: Session protocols.")

# Session 1 - Part A
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "SESSION 1: Building Motivation for Change", 'deep_blue')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["STEP 1: OPENING (10-15 min)",
     "- Welcome warmly and genuinely",
     "- Explain the process briefly",
     "- Set non-judgmental tone",
     "- Ask open question to start:",
     "  \"What brought you here today?\"",
     "  \"How are things going for you?\"",
     "",
     "Then LISTEN. Use reflections.",
     "Let them tell their story."],
    bg='light_blue', title="Opening:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["STEP 2: PERSONAL FEEDBACK (30-40 min)",
     "Present assessment data objectively:",
     "1. Drinking/drug use patterns",
     "2. Comparison with population norms",
     "3. Blood test results (liver, etc.)",
     "4. Neuropsych test results",
     "5. Consequences reported",
     "6. Risk factors identified",
     "",
     "After each: \"What do you make of this?\""],
    bg='light_green', title="Feedback:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["STEP 3: ELICITING CHANGE TALK (15-20 min)",
     "- Ask evocative questions: \"What worries you about your use?\"",
     "- Use Importance Ruler: \"How important is it to change? (0-10) Why not a lower number?\"",
     "- Explore pros/cons: \"What do you like about using? What concerns you?\"",
     "- Look forward: \"Where do you see yourself in 5 years if nothing changes?\"",
     "",
     "STEP 4: CLOSING (5-10 min) - Summarize key themes, assess readiness, assign reflection task."],
    bg='cream', title="Eliciting & Closing:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 40-58; Chapter III: Session 1 - Building Motivation for Change.")

# Session 1 - Clinical Dialogue
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "SESSION 1: Clinical Dialogue Example (from Manual)", 'deep_blue')
cbox(slide, LEFT_M, Inches(1.2), FULL_W, Inches(5.1),
    ["FEEDBACK PRESENTATION DIALOGUE:",
     "",
     "Therapist: \"I'd like to share some results from your assessment. Here's where your drinking",
     "           falls compared to the general adult population in India...\"  [shows graph]",
     "Client:    \"Hmm. I didn't think it was that high.\"",
     "Therapist: \"It's more than you expected.\" (Simple reflection)",
     "Client:    \"Yeah. But a lot of people drink.\"",
     "Therapist: \"You're not alone in drinking, AND what surprises you is how your amount compares.\"",
     "           (Double-sided reflection)",
     "Client:    \"I mean... maybe I do drink more than I thought.\" (CHANGE TALK - problem recognition)",
     "Therapist: \"You're starting to see that your drinking may be on the higher end.\" (Amplify change talk)",
     "",
     "ELICITING SELF-MOTIVATIONAL STATEMENTS:",
     "",
     "Therapist: \"You mentioned your wife has been concerned. What concerns her the most?\"",
     "Client:    \"She says I'm different when I drink. More angry. I guess she has a point.\"",
     "Therapist: \"You've noticed that too - alcohol changes you in ways you don't like.\" (Reflection)",
     "Client:    \"Yeah... I don't want to be that kind of person. That's not who I am.\"",
     "Therapist: \"Being a good person and partner is really important to you.\" (Affirming values)",
     "",
     "NOTE: The therapist NEVER argues. Every response is a reflection, affirmation, or open question.",
     "The CLIENT is the one voicing reasons for change."],
    bg='white', title="", ttc='deep_blue', border='deep_blue', fs=11)
ref_bar(slide, "Ref: MET Manual (1992), pp. 48-55; Adapted dialogue based on manual principles and examples.")

# Session 2
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_green')
title_bar(slide, "SESSION 2: Strengthening Commitment & Change Plan", 'green')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["STEP 1: OPENING (10 min)",
     "\"What's happened since last time?\"",
     "\"Have you thought more about this?\"",
     "Reflect on any changes already made.",
     "",
     "STEP 2: RECAPITULATION (5 min)",
     "Brief summary of Session 1 themes.",
     "Check: \"Did I capture that correctly?\"",
     "",
     "STEP 3: DEEPENING (15 min)",
     "Continue exploring ambivalence.",
     "Use Decisional Balance if needed."],
    bg='light_green', title="First Half:", ttc='green', border='green', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["STEP 4: CHANGE PLAN (25-30 min)",
     "",
     "\"Would you like to put together a plan?\"",
     "",
     "Change Plan Worksheet:",
     "1. Changes I want to make...",
     "2. Most important reasons...",
     "3. Steps I plan to take...",
     "4. How others can help...",
     "5. I'll know it's working if...",
     "6. Things that could interfere..."],
    bg='light_blue', title="Second Half - The Change Plan:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["ADAPTING TO READINESS:",
     "- NOT READY (precontemplating): Continue building rapport, present more feedback, plant seeds",
     "- UNSURE (contemplating): Use Decisional Balance, explore values vs behavior, ask evocative Qs",
     "- READY (determined): Negotiate Change Plan, offer menu of options, help set specific goals",
     "",
     "\"The Change Plan is the CLIENT's plan - not the therapist's prescription. It is negotiated",
     "collaboratively, and the client's preferences and choices are honored.\" (MET Manual, p. 63)"],
    bg='cream', title="Adapting to Client:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 59-68; Chapter IV: Session 2 - Strengthening Commitment.")

# Sessions 3-4
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_orange')
title_bar(slide, "SESSIONS 3 & 4: Reviewing Progress & Maintenance", 'orange')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["SESSION 3 (Week 6):",
     "1. \"How have things been going?\"",
     "2. Review Change Plan progress",
     "   - What's working? What isn't?",
     "   - Any modifications needed?",
     "3. Renew motivation",
     "   - Re-explore importance/confidence",
     "   - Affirm progress (even small steps)",
     "4. Address slips WITHOUT judgment",
     "5. Problem-solve new obstacles"],
    bg='light_orange', title="Session 3:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["SESSION 4 (Week 12):",
     "1. Review overall progress",
     "   - Compare current to baseline",
     "   - Celebrate achievements",
     "2. Maintenance planning",
     "   - Identify high-risk situations",
     "   - Plan coping strategies",
     "3. Discuss relapse as normal",
     "4. Plan ongoing support",
     "5. Termination with hope"],
    bg='light_purple', title="Session 4:", ttc='purple', border='purple', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["IF CLIENT HAS RELAPSED (Session 3/4 scenario):",
     "1. Normalize: \"Slips are a common part of change. Many people experience this.\"",
     "2. Explore without blame: \"What happened? What was going on?\"",
     "3. Reframe as learning: \"Now you know that [situation] is a trigger. That's useful information.\"",
     "4. Re-engage: \"Where do you want to go from here?\"",
     "",
     "NEVER: Express disappointment, say \"you failed\", lecture, or label them as hopeless."],
    bg='light_red', title="Handling Relapse:", ttc='red', border='red', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 69-80; Chapter V: Sessions 3-4; Marlatt & Gordon (1985).")



# ============================================================
# SECTION 5: CLINICAL TECHNIQUES
# ============================================================
divider("SECTION 5", "Clinical Techniques in Detail", 'teal')

# Reflective Listening
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_teal')
title_bar(slide, "Reflective Listening: The Foundation Skill", 'teal')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["LEVELS OF REFLECTION:",
     "",
     "Level 1 - REPEAT/REPHRASE:",
     "Client: \"I drink every night\"",
     "Therapist: \"You drink every evening\"",
     "",
     "Level 2 - PARAPHRASE (meaning):",
     "Therapist: \"It's become a daily routine\"",
     "",
     "Level 3 - REFLECTION OF FEELING:",
     "Therapist: \"You sound concerned about that\""],
    bg='light_teal', title="Types:", ttc='teal', border='teal', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["STRATEGIC USE IN MET:",
     "",
     "- Reflect CHANGE TALK more than sustain talk",
     "- AMPLIFY motivation:",
     "  C: \"I guess I drink a bit much\"",
     "  T: \"You've realized your drinking has become",
     "      a real problem\" (amplified)",
     "- DOUBLE-SIDED for ambivalence:",
     "  \"Drinking helps you relax AND it's costing",
     "   you your marriage.\""],
    bg='light_green', title="Strategic Use:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["PRACTICE GUIDELINES:",
     "- Aim for 2-3 reflections for every question asked",
     "- Reflections are STATEMENTS (drop voice at end), not questions",
     "- If your reflection is wrong, the client will correct you - that's fine",
     "- Common mistake: Asking too many questions and not enough reflections",
     "",
     "Research: Apodaca & Longabaugh (2009) meta-analysis found that MI-consistent therapist",
     "behaviors (especially reflections) were significantly linked to reduced substance use."],
    bg='cream', title="Tips & Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 20-24; Apodaca & Longabaugh (2009). J Subst Abuse Treat, 37, 68-86.")

# Decisional Balance & Rulers
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_orange')
title_bar(slide, "Technique: Decisional Balance & Rulers", 'orange')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["DECISIONAL BALANCE (4-quadrant grid):",
     "",
     "Good things about using | Costs of using",
     "________________________|________________",
     "Benefits of change     | Costs of change",
     "________________________|________________",
     "",
     "Start with GOOD things about using first!",
     "(Shows understanding, builds trust)",
     "Then: \"Where does that leave you?\""],
    bg='light_orange', title="Decisional Balance:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["IMPORTANCE & CONFIDENCE RULERS:",
     "",
     "\"How important is it to change? (0-10)\"",
     "\"How confident are you? (0-10)\"",
     "",
     "KEY QUESTION: \"Why not a lower number?\"",
     "(Invites client to argue FOR change)",
     "",
     "If importance low: develop discrepancy",
     "If confidence low: build self-efficacy",
     "If both high: move to Change Plan"],
    bg='light_blue', title="Rulers:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["WHY \"WHY NOT A LOWER NUMBER?\" WORKS:",
     "Asking \"Why not higher?\" = client defends why they HAVEN'T changed (sustain talk).",
     "Asking \"Why not lower?\" = client articulates reasons they DO want to change (change talk).",
     "",
     "Research: LaBrie et al. (2006) found that Importance and Confidence rulers used in brief MI",
     "sessions predicted actual behavior change at follow-up. The articulation of one's own reasons",
     "for change is more persuasive than hearing reasons from others (Self-Perception Theory)."],
    bg='cream', title="Clinical Wisdom:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 28-37; LaBrie et al. (2006). Addictive Behaviors, 31(8), 1428-1435.")

# Therapist Traps
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_red')
title_bar(slide, "Therapist Traps to AVOID (from the Manual)", 'red')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["1. QUESTION-ANSWER TRAP",
     "   Too many Qs = feels like interrogation",
     "   Fix: More reflections, fewer questions",
     "",
     "2. CONFRONTATION-DENIAL TRAP",
     "   Arguing = client defends drinking more",
     "   Fix: Roll with resistance, reflect",
     "",
     "3. EXPERT TRAP",
     "   \"You should...\" / \"Research shows...\"",
     "   Fix: Elicit client's own solutions first"],
    bg='light_red', title="Traps 1-3:", ttc='red', border='red', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["4. LABELING TRAP",
     "   Insisting on \"alcoholic\" label",
     "   Fix: Focus on behavior, not labels",
     "",
     "5. PREMATURE FOCUS TRAP",
     "   Jumping to solutions too early",
     "   Fix: Follow client's pace",
     "",
     "6. BLAMING TRAP",
     "   Client feels blamed for problem",
     "   Fix: \"Blame is irrelevant - what matters",
     "   is what you want to do NOW\""],
    bg='light_orange', title="Traps 4-6:", ttc='orange', border='orange', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["\"Direct argumentation and an aggressive confrontational approach tend to increase client",
     "resistance and are associated with poorer outcomes.\" (MET Manual, pp. 18-19)",
     "",
     "REMEMBER: Every time YOU argue FOR change, the client argues AGAINST it.",
     "This is the 'righting reflex' - therapist's natural desire to fix things actually backfires.",
     "",
     "Research: Miller et al. (1993) showed that confrontational therapist style predicted client",
     "drinking at 1-year follow-up: MORE confrontation = MORE drinking. The evidence is clear."],
    bg='cream', title="Key Principle:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 18-20; Miller et al. (1993). J Stud Alcohol, 54, 455-461.")



# ============================================================
# SECTION 6: CASE CONCEPTUALIZATION
# ============================================================
divider("SECTION 6", "Case Conceptualization\nApplying MET to a Clinical Case", 'maroon')

# Case Intro
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_red')
title_bar(slide, "Case Study: Applying MET to a Substance Use Case", 'maroon')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["PRESENTING INFORMATION:",
     "- Age: 25 years, Male",
     "- Multiple substance use (polysubstance)",
     "- History of marital dissolution (divorced)",
     "- Referred for assessment and intervention",
     "",
     "ASSESSMENT FINDINGS:",
     "- Externalizing personality organization",
     "- Adequate psychological resources",
     "- Poorly modulated affect (impulsive)",
     "- Damaged, negative self-concept"],
    bg='light_red', title="Client Background:", ttc='maroon', border='maroon', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["KEY PSYCHOLOGICAL FEATURES:",
     "- Emotion-driven decision making",
     "- Hasty, inefficient information processing",
     "- Raw emotional discharge under stress",
     "- Socially engaged but difficulty with",
     "  deep intimate relationships",
     "- Pervasive pessimism/damage-oriented self",
     "- Elevated risk indicators (monitor closely)",
     "",
     "SUBSTANCE USE FUNCTION:",
     "Substances regulate overwhelming affect"],
    bg='light_orange', title="Psychological Profile:", ttc='orange', border='orange', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["WHY MET IS APPROPRIATE FOR THIS CLIENT:",
     "1. He has ADEQUATE resources but uses them poorly (emotion-first style) - MET helps engage",
     "   his existing capacities without confrontation that would trigger defensiveness",
     "2. Confrontation would activate his externalizing defenses and increase resistance",
     "3. His damaged self-image NEEDS affirmation (the 'S' in FRAMES), not more criticism",
     "4. His autonomy needs to be respected (his impulsive style responds poorly to directives)",
     "5. Substances serve an affect-regulation function - MET explores this non-judgmentally"],
    bg='cream', title="MET Rationale:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: Case data from Rorschach assessment (Exner Comprehensive System); MET Manual (1992).")

# Case - Session Plans
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "Case: Session-by-Session MET Plan", 'deep_blue')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["SESSION 1 FOCUS:",
     "- Build rapport carefully (guarded intimacy)",
     "- Present feedback GENTLY (damaged self)",
     "- Focus discrepancy: VALUES (being good",
     "  partner) vs SITUATION (divorced, dependent)",
     "- AFFIRM strengths heavily (resources exist)",
     "- Monitor emotional escalation closely",
     "",
     "Session 1 caution: His affect can escalate",
     "quickly. Keep pace slow, reflect often."],
    bg='light_blue', title="Sessions 1-2:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["SESSION 2 FOCUS:",
     "- Change Plan addresses FUNCTION of use",
     "  (What does using DO for you? What need?)",
     "- Small, achievable steps (failures feel",
     "  catastrophic given damaged self-image)",
     "- Affect regulation alternatives:",
     "  \"What else helps when feelings get intense?\"",
     "- Reflective delay practice:",
     "  \"What if you took 10 min before deciding?\"",
     "",
     "Key: Plan must match his pace/style."],
    bg='light_green', title="", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["SESSIONS 3-4 FOCUS:",
     "- Session 3: If progress - affirm heavily (counters damaged self-image). If slip - normalize,",
     "  explore emotional triggers, revise plan to add more emotion regulation strategies.",
     "- Session 4: Maintenance plan focused on high-risk situations (intense emotions, conflict,",
     "  rejection, self-criticism). Referral for continued therapy (affect regulation, self-image repair).",
     "",
     "BEYOND MET (this client needs): Affect regulation training, self-image repair work,",
     "interpersonal skills, slower processing practice. MET provides MOTIVATIONAL FOUNDATION."],
    bg='cream', title="Sessions 3-4 & Beyond:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), Chs. III-V; Exner (2003). Rorschach, Vol. 1; NIMHANS (2016).")



# ============================================================
# SECTION 7: WORKSHEETS
# ============================================================
divider("SECTION 7", "Clinical Worksheets & Tools\nPrintable Resources for MET Sessions", 'deep_blue')

# Change Plan Worksheet
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "WORKSHEET 1: Change Plan (from MET Manual)", 'deep_blue')
cbox(slide, LEFT_M, Inches(1.2), FULL_W, Inches(5.1),
    ["1. THE CHANGES I WANT TO MAKE ARE:",
     "   _________________________________________________________________________",
     "",
     "2. THE MOST IMPORTANT REASONS WHY I WANT TO MAKE THESE CHANGES ARE:",
     "   _________________________________________________________________________",
     "",
     "3. THE STEPS I PLAN TO TAKE IN CHANGING ARE:",
     "   _________________________________________________________________________",
     "",
     "4. THE WAYS OTHER PEOPLE CAN HELP ME ARE:",
     "   Person: _________________ How: ___________________________________________",
     "",
     "5. I WILL KNOW THAT MY PLAN IS WORKING IF:",
     "   _________________________________________________________________________",
     "",
     "6. SOME THINGS THAT COULD INTERFERE WITH MY PLAN ARE:",
     "   _________________________________________________________________________"],
    bg='white', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), Appendix B: Change Plan Worksheet.")

# Decisional Balance Worksheet
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_green')
title_bar(slide, "WORKSHEET 2: Decisional Balance Grid", 'green')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["GOOD THINGS about my current behavior",
     "(What I LIKE about using substances):",
     "",
     "1. ____________________________________",
     "2. ____________________________________",
     "3. ____________________________________",
     "4. ____________________________________",
     "5. ____________________________________"],
    bg='light_orange', title="Benefits of Status Quo:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["COSTS of my current behavior",
     "(What CONCERNS me about using):",
     "",
     "1. ____________________________________",
     "2. ____________________________________",
     "3. ____________________________________",
     "4. ____________________________________",
     "5. ____________________________________"],
    bg='light_red', title="Costs of Status Quo:", ttc='red', border='red', fs=12)
cbox(slide, LEFT_M, Inches(3.8), HALF_W, Inches(2.5),
    ["BENEFITS of making a change",
     "(What would be GOOD about changing):",
     "",
     "1. ____________________________________",
     "2. ____________________________________",
     "3. ____________________________________",
     "4. ____________________________________",
     "5. ____________________________________"],
    bg='light_green', title="Benefits of Change:", ttc='green', border='green', fs=12)
cbox(slide, COL2_X, Inches(3.8), HALF_W, Inches(2.5),
    ["COSTS of making a change",
     "(What would be HARD about changing):",
     "",
     "1. ____________________________________",
     "2. ____________________________________",
     "3. ____________________________________",
     "4. ____________________________________",
     "5. ____________________________________"],
    bg='light_purple', title="Costs of Change:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: Janis & Mann (1977). Decision Making; MET Manual (1992), p. 29.")

# Rulers Worksheet
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_purple')
title_bar(slide, "WORKSHEET 3: Readiness, Importance & Confidence Rulers", 'purple')
cbox(slide, LEFT_M, Inches(1.2), FULL_W, Inches(1.5),
    ["IMPORTANCE: How important is it to you to make this change?",
     "Not important  0 --- 1 --- 2 --- 3 --- 4 --- 5 --- 6 --- 7 --- 8 --- 9 --- 10  Extremely important",
     "Why this number and not lower? _______________________________________________________"],
    bg='light_blue', border='deep_blue', fs=12)
cbox(slide, LEFT_M, Inches(2.9), FULL_W, Inches(1.5),
    ["CONFIDENCE: How confident are you that you COULD make this change?",
     "Not confident  0 --- 1 --- 2 --- 3 --- 4 --- 5 --- 6 --- 7 --- 8 --- 9 --- 10  Extremely confident",
     "Why this number and not lower? _______________________________________________________"],
    bg='light_green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(4.6), FULL_W, Inches(1.5),
    ["READINESS: How ready are you to make this change RIGHT NOW?",
     "Not ready  0 --- 1 --- 2 --- 3 --- 4 --- 5 --- 6 --- 7 --- 8 --- 9 --- 10  Completely ready",
     "What would help you move up one point? _______________________________________________"],
    bg='light_orange', border='orange', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 35-37; Rollnick, Mason & Butler (1999). Health Behavior Change.")

# Diary and RP Plan
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_teal')
title_bar(slide, "WORKSHEET 4: Daily Self-Monitoring Diary", 'teal')
tbl_slide("WORKSHEET 4: Daily Self-Monitoring Diary",
    ["Day", "Situation/Trigger", "Feelings (0-10)", "Urge (0-10)", "What I Did Instead", "Result"],
    [["Monday", "____________", "____", "____", "____________", "______"],
     ["Tuesday", "____________", "____", "____", "____________", "______"],
     ["Wednesday", "____________", "____", "____", "____________", "______"],
     ["Thursday", "____________", "____", "____", "____________", "______"],
     ["Friday", "____________", "____", "____", "____________", "______"],
     ["Saturday", "____________", "____", "____", "____________", "______"],
     ["Sunday", "____________", "____", "____", "____________", "______"]],
    color='teal', ref="Ref: Adapted from MET Manual (1992); Marlatt & Gordon (1985).")

# RP Plan
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_orange')
title_bar(slide, "WORKSHEET 5: My Relapse Prevention Plan", 'orange')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["MY HIGH-RISK SITUATIONS:",
     "1. ____________________________________",
     "2. ____________________________________",
     "3. ____________________________________",
     "",
     "MY EARLY WARNING SIGNS:",
     "1. ____________________________________",
     "2. ____________________________________",
     "3. ____________________________________"],
    bg='light_red', title="Identifying Risks:", ttc='red', border='red', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["MY COPING STRATEGIES:",
     "1. ____________________________________",
     "2. ____________________________________",
     "3. ____________________________________",
     "",
     "SUPPORT CONTACTS:",
     "Name: _____________ Phone: _____________",
     "Name: _____________ Phone: _____________",
     "Name: _____________ Phone: _____________"],
    bg='light_green', title="My Resources:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["IF I HAVE A SLIP, I WILL:",
     "1. Remember: a slip is NOT a failure - it's information about a trigger to address",
     "2. Call: _________________________________ (my support person)",
     "3. Do instead: _______________________________________________________________",
     "4. Review: What triggered it? What was I feeling? What can I learn?",
     "",
     "MY TOP REASONS TO STAY ON TRACK (from Decisional Balance):",
     "_______________________________________________________________________________"],
    bg='cream', title="Emergency Plan:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), Session 4; Marlatt & Donovan (2005). Relapse Prevention, 2nd ed.")

# Values Worksheet
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_green')
title_bar(slide, "WORKSHEET 6: Personal Values Exploration", 'green')
cbox(slide, LEFT_M, Inches(1.2), FULL_W, Inches(3.3),
    ["Rate each value: IMPORTANCE (1-5) and how much CURRENT BEHAVIOR aligns (1-5)",
     "",
     "VALUE                    IMPORTANCE    ALIGNMENT    GAP?",
     "Being a good parent      _________     _________    ___",
     "Physical health          _________     _________    ___",
     "Financial security       _________     _________    ___",
     "Being honest             _________     _________    ___",
     "Close relationships      _________     _________    ___",
     "Self-respect             _________     _________    ___",
     "Career success           _________     _________    ___",
     "Independence             _________     _________    ___",
     "Spirituality/faith       _________     _________    ___"],
    bg='white', title="Values-Behavior Alignment:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(4.7), FULL_W, Inches(1.6),
    ["Which values have the BIGGEST gaps? How does substance use affect these values?",
     "What would change look like for the values that matter most?",
     "",
     "Therapist Note: Creates cognitive dissonance by making the gap explicit. Use to develop discrepancy."],
    bg='cream', title="Reflection:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Miller et al. (2001). Personal Values Card Sort; MET Manual (1992); Miller & Rollnick (2002).")



# ============================================================
# SECTION 8: COMPARISON WITH OTHER THERAPIES (NEW)
# ============================================================
divider("SECTION 8", "MET Compared with Other Therapies\nFor Patient, Family, and Caregivers", 'navy')

# MET vs CBT
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "MET vs Cognitive Behavioral Therapy (CBT)", 'navy')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["MET APPROACH:",
     "- Focus: WHY to change (motivation)",
     "- Client provides the answers",
     "- 4 sessions over 12 weeks",
     "- Non-directive, exploratory style",
     "- Works best EARLY in treatment",
     "- Addresses ambivalence and readiness",
     "- No homework assignments given",
     "- Therapist evokes, doesn't teach"],
    bg='light_blue', title="MET:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["CBT APPROACH:",
     "- Focus: HOW to change (skills)",
     "- Therapist teaches coping skills",
     "- 12-16 sessions typically",
     "- Structured, directive style",
     "- Works best DURING action stage",
     "- Addresses triggers and behaviors",
     "- Regular homework/practice assigned",
     "- Therapist educates and trains"],
    bg='light_green', title="CBT:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["EVIDENCE FOR COMBINING MET + CBT:",
     "- Project MATCH (1997): MET alone = CBT alone in outcomes (4 sessions vs 12 sessions)",
     "- Marijuana Treatment Project (Stephens et al., 2004): MET+CBT combined was optimal",
     "- Cannabis Youth Treatment Study (Dennis et al., 2004): MET/CBT combination most effective",
     "- COMBINE Study (Anton et al., 2006, JAMA): MET as prelude to CBT improved engagement",
     "",
     "CLINICAL RECOMMENDATION: Use MET FIRST (Sessions 1-4) to build motivation, then transition",
     "to CBT for skills training once client is in Action stage. Sequential model is most effective."],
    bg='cream', title="Research Evidence:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Project MATCH (1997); Stephens et al. (2004). J Consult Clin Psychol; COMBINE (2006). JAMA, 295.")

# MET vs 12-Step
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_purple')
title_bar(slide, "MET vs 12-Step Facilitation (TSF) / AA Approach", 'purple')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["MET APPROACH:",
     "- No labels required (\"alcoholic\")",
     "- Client sets own goals (abstinence or",
     "  moderation - client's choice)",
     "- Brief (4 sessions)",
     "- Individual therapy",
     "- Therapist is collaborative partner",
     "- No spiritual component",
     "- \"You decide what's best for you\""],
    bg='light_purple', title="MET:", ttc='purple', border='purple', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["12-STEP / TSF APPROACH:",
     "- Acceptance of \"alcoholic\" identity",
     "- Goal: total abstinence only",
     "- Ongoing/lifetime involvement",
     "- Group-based (AA/NA meetings)",
     "- Sponsor system for support",
     "- Spiritual foundation (Higher Power)",
     "- \"You are powerless over alcohol\"",
     "- Structured step-work program"],
    bg='light_orange', title="12-Step:", ttc='orange', border='orange', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["EVIDENCE COMPARISON:",
     "- Project MATCH (1997): Both EQUALLY effective at 1 and 3-year follow-up",
     "- TSF produced slightly higher continuous abstinence rates (but not significantly)",
     "- MET more cost-effective (4 sessions vs 12 sessions for same outcomes)",
     "- MET better for clients HIGH in anger (Karno & Longabaugh, 2005)",
     "- TSF better for clients with high social networks supportive of drinking",
     "",
     "INTEGRATION: Many clinicians use MET to BUILD motivation, then refer to AA/NA for ongoing",
     "support. MET addresses the \"why change\" while 12-Step provides community and structure."],
    bg='cream', title="Research:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: Project MATCH (1997); Karno & Longabaugh (2005). J Stud Alcohol, 66, 488-495.")

# MET vs Contingency Management
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_teal')
title_bar(slide, "MET vs Other Evidence-Based Therapies", 'teal')
tbl_slide("Comparison of Evidence-Based Therapies for Substance Use",
    ["Therapy", "Mechanism", "Duration", "Best For", "Evidence Level"],
    [["MET", "Internal motivation\n(discrepancy + MI)", "4 sessions\n(12 weeks)", "Low-readiness clients\nAngry clients\nEarly treatment", "Level 1\n(Project MATCH, UKATT)"],
     ["CBT", "Skills training\n(coping, triggers)", "12-16 sessions", "Action-stage clients\nAnxiety comorbidity", "Level 1\n(Carroll et al., 1994)"],
     ["12-Step (TSF)", "Peer support\nSpiritual growth", "Ongoing/\nlifetime", "High social support\nAbstinence-oriented", "Level 1\n(Project MATCH)"],
     ["Contingency Mgmt", "External rewards\nfor abstinence", "12-24 weeks", "Stimulant use\nOpioid use\nImmediate reinforcement", "Level 1\n(Higgins et al., 2004)"],
     ["Community\nReinforcement", "Restructure\nenvironment", "12-24 sessions", "Social isolation\nUnemployment", "Level 1\n(Meyers & Smith, 1995)"],
     ["BSFT (Family)", "Family system\nrestructuring", "12-16 sessions", "Adolescents\nFamily conflict", "Level 1\n(Szapocznik et al., 2003)"]],
    color='teal',
    ref="Ref: NICE Guidelines (2011); SAMHSA TIP 35 (1999); APA Practice Guidelines (2006); NIMHANS (2016).")

# Why MET is More Effective (comparative advantages)
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_green')
title_bar(slide, "Why MET is Often MORE Effective Than Traditional Approaches", 'green')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["ADVANTAGES OF MET:",
     "",
     "1. COST-EFFECTIVE: 4 sessions = 12 sessions",
     "   (70% less therapist time for same results)",
     "",
     "2. ENGAGES RESISTANT CLIENTS:",
     "   Works for people who don't want treatment",
     "   (court-mandated, family-pressured)",
     "",
     "3. REDUCES DROPOUT:",
     "   Non-confrontational = clients stay longer",
     "   (Carroll et al., 2006)"],
    bg='light_green', title="Clinical Advantages:", ttc='green', border='green', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["4. MATCHES CLIENT'S STAGE:",
     "   Doesn't force action on unready clients",
     "",
     "5. SCALABLE:",
     "   Can be delivered by trained lay counselors",
     "   (Nadkarni et al., 2017 - Lancet)",
     "",
     "6. UNIVERSAL ENHANCER:",
     "   Improves outcomes when added to ANY",
     "   other treatment (CBT, pharmacotherapy,",
     "   residential treatment)",
     "",
     "7. WORKS ACROSS CULTURES (WHO, 2002)"],
    bg='light_blue', title="Systemic Advantages:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["RESEARCH SUPPORT FOR SUPERIORITY IN SPECIFIC CONTEXTS:",
     "- Angry clients: MET significantly BETTER than CBT/TSF (Project MATCH matching hypothesis)",
     "- Low-readiness clients: MET better at engaging and retaining them in treatment",
     "- Brief settings (primary care): MET achievable where 12-session therapy is not feasible",
     "- Stepped care model: MET as first step, add intensive treatment only if MET insufficient",
     "",
     "Lundahl et al. (2010) meta-analysis: MI/MET effect size d=0.22 overall; d=0.79 for substance use",
     "when compared to no treatment. This is a MEDIUM-LARGE effect for a 4-session intervention."],
    bg='cream', title="Research Evidence:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Lundahl et al. (2010). Patient Ed & Counsel, 80; Project MATCH (1997); Nadkarni et al. (2017). Lancet.")



# Family & Caregiver Approaches
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_orange')
title_bar(slide, "Family & Caregiver Involvement in MET", 'orange')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["ROLE OF FAMILY IN MET:",
     "",
     "The MET Manual acknowledges that a",
     "\"concerned significant other\" can be",
     "involved in Session 2 for Change Plan:",
     "",
     "- Provides support for client's goals",
     "- Helps identify triggers at home",
     "- Can reinforce positive changes",
     "- Learns non-confrontational approach",
     "",
     "KEY: Family supports, doesn't direct."],
    bg='light_orange', title="Family in MET:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["COMPLEMENTARY FAMILY THERAPIES:",
     "",
     "1. CRAFT (Community Reinforcement",
     "   and Family Training):",
     "   - For families of unmotivated users",
     "   - Teaches positive communication",
     "   - Research: 64% engagement rate",
     "   (Meyers et al., 2002)",
     "",
     "2. Behavioral Couples Therapy:",
     "   - Recovery contracts",
     "   - Reduces both substance use AND",
     "     relationship distress simultaneously"],
    bg='light_green', title="Adjunct Therapies:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["CAREGIVER PSYCHOEDUCATION (use alongside MET):",
     "- Teach CRAFT principles: reward sobriety, allow natural consequences, self-care",
     "- Avoid enabling behaviors (covering up, making excuses, giving money)",
     "- Understand stages of change (don't expect overnight transformation)",
     "- Learn MI techniques for home conversations (reflective listening, not lecturing)",
     "",
     "Research: Smith & Meyers (2004) showed CRAFT training for families achieved 64% treatment",
     "engagement vs 30% for Al-Anon and 13% for Johnson Intervention (confrontational approach)."],
    bg='cream', title="Caregiver Education:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), p. 65; Meyers et al. (2002). J Consult Clin Psychol; Smith & Meyers (2004).")

# Family-inclusive MET adaptations
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_green')
title_bar(slide, "Therapies Used ALONGSIDE MET for Comprehensive Care", 'green')
cbox(slide, LEFT_M, Inches(1.2), FULL_W, Inches(1.8),
    ["THE STEPPED CARE MODEL (recommended by NIMHANS and international guidelines):",
     "Step 1: MET/Brief MI (4 sessions) - for ALL clients as starting point",
     "Step 2: If insufficient - add CBT/Skills Training (12 sessions)",
     "Step 3: If still insufficient - add Pharmacotherapy + Intensive outpatient",
     "Step 4: If still insufficient - Residential treatment / Therapeutic community"],
    bg='light_green', title="Sequential Treatment Model:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.2), HALF_W, Inches(3.0),
    ["FOR THE PATIENT (adjunct therapies):",
     "",
     "1. MET + Pharmacotherapy:",
     "   - Naltrexone/Acamprosate (alcohol)",
     "   - Buprenorphine (opioids)",
     "   - Research: COMBINE (2006) showed MI +",
     "     naltrexone as effective combination",
     "",
     "2. MET + CBT (sequential):",
     "   - MET first for motivation (Sessions 1-4)",
     "   - CBT next for skills (Sessions 5-16)",
     "   - Best evidence: Cannabis Youth Study"],
    bg='light_blue', title="For Patient:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2_X, Inches(3.2), HALF_W, Inches(3.0),
    ["FOR FAMILY & CAREGIVERS:",
     "",
     "1. CRAFT Training (6-12 sessions)",
     "   Teaches: positive communication,",
     "   allowing consequences, self-care",
     "",
     "2. Behavioral Family Therapy:",
     "   Communication skills, problem-solving",
     "",
     "3. Family Psychoeducation:",
     "   Understanding addiction as illness",
     "   Reducing expressed emotion (EE)",
     "",
     "4. Caregiver support groups"],
    bg='light_purple', title="For Family:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: NIMHANS (2016), Ch. 8-10; COMBINE (2006). JAMA; Meyers et al. (2002); NICE Guidelines (2011).")



# ============================================================
# SECTION 9: COMORBID PSYCHIATRIC DISORDERS (NEW)
# ============================================================
divider("SECTION 9", "MET with Comorbid Psychiatric Disorders\nWhat to Do When There Are Co-Occurring Conditions", 'purple')

# Overview of Comorbidity
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_purple')
title_bar(slide, "Dual Diagnosis: Substance Use + Psychiatric Disorders", 'purple')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["PREVALENCE OF COMORBIDITY:",
     "",
     "- 50-70% of substance users have at least",
     "  one comorbid psychiatric disorder",
     "  (Regier et al., 1990; Kessler, 2004)",
     "- Common comorbidities:",
     "  * Depression (30-50%)",
     "  * Anxiety disorders (25-40%)",
     "  * Personality disorders (40-70%)",
     "  * PTSD (25-40%)",
     "  * Bipolar disorder (20-30%)"],
    bg='light_purple', title="The Problem:", ttc='purple', border='purple', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["CAN MET BE USED WITH COMORBIDITY?",
     "",
     "YES - with adaptations:",
     "- MET principles are applicable regardless",
     "  of comorbid diagnosis",
     "- MI/MET has been tested in dual diagnosis",
     "  populations with positive results",
     "- Key: integrate, don't separate treatment",
     "",
     "\"Integrated treatment is more effective",
     "than parallel or sequential treatment\"",
     "(Drake et al., 2004, Schizophr Bull)"],
    bg='light_green', title="MET Applicability:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["GENERAL PRINCIPLES FOR MET WITH COMORBIDITY (from research):",
     "1. TREAT BOTH SIMULTANEOUSLY (integrated approach superior to sequential)",
     "2. Acknowledge that substance use may be SELF-MEDICATION for psychiatric symptoms",
     "3. Explore the FUNCTION of use: \"What does drinking do for your anxiety/depression?\"",
     "4. Develop discrepancy: \"You use to feel better, but does it actually make things better long-term?\"",
     "5. Adapt pace - comorbid clients may need more sessions or slower progression",
     "6. Coordinate with prescribing psychiatrist for pharmacotherapy decisions"],
    bg='cream', title="General Approach:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Drake et al. (2004). Schizophr Bull, 30; Kessler (2004). Biol Psych; NIMHANS (2016), Ch. 12.")

# Personality Disorders
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_red')
title_bar(slide, "MET with Comorbid Personality Disorders", 'red')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["CHALLENGES:",
     "- Interpersonal difficulties affect rapport",
     "- Emotional dysregulation (BPD)",
     "- Impulsivity increases relapse risk",
     "- Identity disturbance complicates values work",
     "- Therapist-client dynamic may become complex",
     "- Higher dropout rates in treatment",
     "",
     "PREVALENCE: 40-70% of substance users",
     "meet criteria for at least one PD",
     "(Verheul, 2001; NIMHANS data)"],
    bg='light_red', title="Personality Disorders & SUD:", ttc='red', border='red', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["MET ADAPTATIONS FOR PD:",
     "",
     "1. LONGER ENGAGEMENT PHASE",
     "   - May need extra sessions before feedback",
     "   - Build alliance more carefully",
     "",
     "2. EMPHASIS ON ROLLING WITH RESISTANCE",
     "   - PD clients show more resistance",
     "   - Never argue or confront",
     "",
     "3. VALIDATE EMOTIONS before exploring change",
     "   (especially for Borderline PD)",
     "",
     "4. SMALLER CHANGE GOALS (reduce overwhelm)"],
    bg='light_green', title="Adaptations:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["RESEARCH EVIDENCE:",
     "- Ball et al. (2007): Dual Focus Schema Therapy + MI for PD+SUD: improved outcomes vs",
     "  standard drug counseling (J Nerv Ment Dis, 195, 24-31)",
     "- Bornovalova & Daughters (2007): MI effective for Borderline PD + substance use when",
     "  distress tolerance component added",
     "- Gregory et al. (2008): MI feasible with Antisocial PD (often seen in substance users)",
     "- NIMHANS recommendation: Use MI/MET as engagement strategy for PD patients who are",
     "  difficult to retain in longer psychotherapy; then transition to DBT or Schema Therapy"],
    bg='cream', title="Research:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: Ball et al. (2007). J Nerv Ment Dis; Bornovalova & Daughters (2007); NIMHANS (2016), Ch. 12.")

# Depression Comorbidity
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "MET with Comorbid Depression", 'deep_blue')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["DEPRESSION + SUBSTANCE USE:",
     "",
     "- 30-50% of substance users have depression",
     "- Alcohol is a CNS depressant - worsens mood",
     "- Substance-induced vs independent depression",
     "- Both need addressing simultaneously",
     "",
     "CHALLENGES FOR MET:",
     "- Low energy/motivation (core of depression)",
     "- Hopelessness may affect self-efficacy",
     "- Suicidal ideation must be monitored",
     "- Cognitive distortions complicate reasoning"],
    bg='light_blue', title="The Challenge:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["MET ADAPTATIONS FOR DEPRESSION:",
     "",
     "1. EXTRA EMPHASIS on self-efficacy building",
     "   (counters hopelessness of depression)",
     "2. Explore MOOD-SUBSTANCE link:",
     "   \"When do you drink most? When low?\"",
     "3. Develop discrepancy with MOOD:",
     "   \"You drink to feel better, but next day",
     "    your mood is even worse.\"",
     "4. Coordinate with antidepressant medication",
     "5. Safety planning if suicidal ideation present",
     "6. Smaller goals (energy is limited)"],
    bg='light_green', title="Adaptations:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["RESEARCH EVIDENCE:",
     "- Baker et al. (2010): MI + CBT for comorbid depression and substance use - significantly",
     "  reduced both depressive symptoms AND substance use (Addiction, 105, 1560-1568)",
     "- Satre et al. (2016): MI effective for older adults with depression + alcohol use",
     "- Westra & Dozois (2006): MI as prelude to CBT for depression enhanced treatment engagement",
     "- Rao et al. (2015, NIMHANS): MI-based brief intervention effective for depression + alcohol",
     "  in Indian primary care settings (Asian J Psychiatry, 16, 45-51)"],
    bg='cream', title="Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Baker et al. (2010). Addiction; Westra & Dozois (2006). J Consult Clin Psychol; Rao et al. (2015).")

# Anxiety & PTSD
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_teal')
title_bar(slide, "MET with Comorbid Anxiety Disorders & PTSD", 'teal')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["ANXIETY + SUBSTANCE USE:",
     "",
     "- 25-40% comorbidity rate",
     "- Substances used to \"self-medicate\" anxiety",
     "- Alcohol reduces anxiety short-term,",
     "  increases it long-term (rebound effect)",
     "- Withdrawal mimics anxiety symptoms",
     "",
     "PTSD + SUBSTANCE USE:",
     "- 25-40% of PTSD patients have SUD",
     "- Substances numb trauma memories",
     "- Integrated treatment essential"],
    bg='light_teal', title="The Challenge:", ttc='teal', border='teal', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["MET ADAPTATIONS:",
     "",
     "For ANXIETY:",
     "1. Validate anxiety as real (not just excuse)",
     "2. Explore substance-anxiety cycle",
     "3. Include anxiety management in Change Plan",
     "4. Confidence building is crucial",
     "",
     "For PTSD:",
     "1. Safety first - stabilization before change",
     "2. Respect avoidance (don't push trauma work)",
     "3. MET for engagement, then refer to EMDR/PE",
     "4. Trauma-informed MI adaptation"],
    bg='light_green', title="Adaptations:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["RESEARCH EVIDENCE:",
     "- Hien et al. (2009): Integrated MI + Seeking Safety for PTSD+SUD - effective in reducing",
     "  both trauma symptoms and substance use (J Consult Clin Psychol, 77, 607-619)",
     "- Westra (2012): MI enhanced CBT outcomes for anxiety disorders (Addictive Behaviors, 37, 159)",
     "- Najavits (2002): Seeking Safety + MI as integrated approach for trauma + substance use",
     "- Bolton et al. (2016, Lancet Psychiatry): MI effective in low-resource settings for comorbidity",
     "- Sannibale et al. (2013): MI + exposure therapy for alcohol+PTSD - feasible and effective"],
    bg='cream', title="Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Hien et al. (2009). JCCP; Westra (2012). Addictive Behaviors; Sannibale et al. (2013). Addiction.")

# Bipolar & Psychosis
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_orange')
title_bar(slide, "MET with Bipolar Disorder & Psychotic Disorders", 'orange')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["BIPOLAR + SUBSTANCE USE:",
     "- 40-60% lifetime SUD comorbidity",
     "- Mania: impulsivity increases use",
     "- Depression: self-medication pattern",
     "- Non-adherence to mood stabilizers",
     "",
     "PSYCHOSIS + SUBSTANCE USE:",
     "- 40-50% have comorbid SUD",
     "- Cannabis, alcohol most common",
     "- Complicates symptom management",
     "- Reduced insight may limit engagement"],
    bg='light_orange', title="The Challenge:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["MET ADAPTATIONS:",
     "",
     "For BIPOLAR:",
     "1. Address medication adherence with MI",
     "2. Develop discrepancy: substance use triggers",
     "   mood episodes (show the pattern)",
     "3. Change Plan includes mood monitoring",
     "",
     "For PSYCHOSIS:",
     "1. Simplify language and concepts",
     "2. Focus on immediate, concrete goals",
     "3. Shorter sessions if attention limited",
     "4. Affirm heavily (low self-esteem common)"],
    bg='light_green', title="Adaptations:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["RESEARCH EVIDENCE:",
     "- Barrowclough et al. (2001): MI + CBT for schizophrenia + SUD - significant reduction in",
     "  substance use over 18 months (BMJ, 323, 1-5). Landmark dual diagnosis trial.",
     "- Graeber et al. (2003): MI for psychosis + SUD - improved engagement and treatment retention",
     "- Kemp et al. (1996): Compliance Therapy (MI-based) for bipolar medication adherence",
     "- Haddock et al. (2003): Integrated MI+CBT for psychosis+SUD in community mental health",
     "- Baker et al. (2005): MI effective for comorbid psychosis + cannabis use (Addiction, 100, 1614)"],
    bg='cream', title="Research:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: Barrowclough et al. (2001). BMJ; Baker et al. (2005). Addiction; Haddock et al. (2003).")

# Comprehensive Comorbidity Table
tbl_slide("Summary: MET Adaptations for Comorbid Disorders",
    ["Comorbid Disorder", "Key MET Adaptation", "Adjunct Therapy", "Key Research"],
    [["Personality Disorders\n(Borderline, Antisocial)", "Extra alliance building\nValidation before change\nSmaller goals", "DBT (Borderline)\nSchema Therapy\nMentalisation Therapy", "Ball et al. (2007)\nBornovalova (2007)"],
     ["Major Depression", "Extra self-efficacy\nExplore mood-substance link\nSafety planning", "Antidepressants\nCBT for Depression\nBehavioral Activation", "Baker et al. (2010)\nRao et al. (2015)"],
     ["Anxiety Disorders", "Validate anxiety\nExplore substance-anxiety cycle\nInclude anxiety management", "SSRIs/SNRIs\nCBT for Anxiety\nRelaxation Training", "Westra (2012)\nWestra & Dozois (2006)"],
     ["PTSD", "Safety first\nTrauma-informed MI\nRespect avoidance", "EMDR / PE\nSeeking Safety\nCPT", "Hien et al. (2009)\nSannibale et al. (2013)"],
     ["Bipolar Disorder", "Address med adherence\nLink use to mood episodes\nMood monitoring", "Mood Stabilizers\nIPSRT\nPsychoeducation", "Kemp et al. (1996)\nGraeber et al. (2003)"],
     ["Psychotic Disorders", "Simplify language\nConcrete goals\nShorter sessions", "Antipsychotics\nIntegrated dual Dx\nFamily intervention", "Barrowclough (2001)\nBaker et al. (2005)"]],
    color='purple',
    ref="Ref: All sources cited in individual slides. Review: Drake et al. (2004). Schizophr Bull, 30, 795-808.")



# ============================================================
# SECTION 10: RESEARCH & EFFECTIVENESS
# ============================================================
divider("SECTION 10", "Research Evidence & Effectiveness\nHow Well Does MET Work?", 'navy')

# Project MATCH
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "Project MATCH: The Landmark Trial", 'navy')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["STUDY DESIGN:",
     "- Largest alcohol treatment trial ever",
     "- 1,726 participants across 9 US sites",
     "- Randomized to 3 conditions:",
     "  1. MET (4 sessions / 12 weeks)",
     "  2. CBT (12 sessions)",
     "  3. TSF/12-Step (12 sessions)",
     "- Follow-up: 1 year and 3 years",
     "- Cost: $27 million (NIAAA funded)"],
    bg='light_blue', title="The Study:", ttc='navy', border='navy', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["KEY FINDINGS:",
     "- All three treatments significantly improved",
     "- MET achieved COMPARABLE outcomes to",
     "  CBT and TSF in just 4 sessions vs 12",
     "- Percent Days Abstinent improved in ALL",
     "- Drinks Per Drinking Day decreased in ALL",
     "- At 3-year follow-up: gains maintained",
     "",
     "IMPLICATION: MET is as effective as longer",
     "treatments = highly COST-EFFECTIVE"],
    bg='light_green', title="Results:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["MATCHING FINDING: Clients HIGH in anger did significantly BETTER in MET than other treatments.",
     "Non-confrontational approach works better for angry, reactive clients.",
     "",
     "\"The finding that a 4-session motivational intervention could produce outcomes comparable to",
     "12-session treatments had profound implications for cost-effectiveness.\" (MATCH Group, 1997)",
     "",
     "India relevance: With limited therapist availability, a 4-session effective treatment is ideal",
     "for resource-constrained settings (NIMHANS recommends MET as first-line, 2016)."],
    bg='cream', title="Key Implication:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Project MATCH (1997). J Stud Alcohol, 58, 7-29; (1998). Addiction, 93, 1434-1447.")

# Meta-analyses
tbl_slide("Meta-Analyses: MI/MET Effectiveness Across Studies",
    ["Study", "N Studies", "Key Finding", "Effect Size"],
    [["Burke et al. (2003)\nJ Consult Clin Psychol", "30 RCTs", "MI/MET effective for alcohol,\ndrugs, diet, adherence", "d = 0.25-0.57"],
     ["Hettema et al. (2005)\nAnnual Rev Clin Psychol", "72 studies", "MI effective across substances\nvs no treatment/advice", "d = 0.77 at follow-up"],
     ["Lundahl et al. (2010)\nPatient Ed & Counsel", "119 studies", "MI produces moderate effect\nstronger for substance use", "d = 0.22 (combined)\nd = 0.79 (substance)"],
     ["Vasilaki et al. (2006)\nAddiction", "15 RCTs", "Brief MI effective for reducing\nalcohol consumption", "Significant reduction"],
     ["Smedslund et al. (2011)\nCochrane Review", "59 RCTs", "MI reduces substance use\nmore than no treatment", "SMD = -0.79"]],
    color='navy',
    ref="Ref: All meta-analyses cited in table. All are peer-reviewed systematic reviews/meta-analyses.")

# Indian Research
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_orange')
title_bar(slide, "Research Evidence: Indian Context", 'orange')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["NIMHANS EVIDENCE:",
     "- Brief MI effective in Indian primary care",
     "  (Murthy et al., 2009)",
     "- NIMHANS protocols integrate MET as",
     "  first-line psychosocial intervention",
     "- Community MI reduced alcohol use in",
     "  rural Karnataka (Nadkarni et al., 2017)",
     "",
     "AIIMS, DELHI:",
     "- Brief MI for alcohol dependence positive",
     "  outcomes (Pal et al., 2007)"],
    bg='light_orange', title="Indian Studies:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["PREMIUM TRIAL (Lancet, 2017):",
     "- Lay counselors delivered MI in Goa",
     "- Significant reduction in harmful drinking",
     "- Published in The Lancet (top journal)",
     "- Proved task-shifting feasible in India",
     "",
     "CHAND et al. (2018, NIMHANS):",
     "- Technology-assisted brief intervention",
     "- Showed feasibility in Indian clinics",
     "",
     "OUTCOME: MI/MET effective in India"],
    bg='light_green', title="Key Indian Studies:", ttc='green', border='green', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["CULTURAL CONSIDERATIONS FOR INDIA:",
     "- Family involvement crucial in Indian context (can integrate into MET Change Plan)",
     "- High stigma around substance use - MI's non-judgmental approach especially valuable",
     "- Task-shifting to lay counselors feasible (PREMIUM trial demonstrated this conclusively)",
     "- Adaptations needed: family values, community roles, spiritual beliefs integrated into values work",
     "- NIMHANS (2016) recommends MI/MET as evidence-based first-line psychosocial intervention",
     "- Cost-effective for India's resource-constrained mental health system (4 sessions adequate)"],
    bg='cream', title="Indian Adaptations:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: Nadkarni et al. (2017). Lancet, 389, 186-195; NIMHANS (2016); Chand et al. (2018). Indian J Psych.")

# UKATT
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_green')
title_bar(slide, "UKATT & International Evidence", 'green')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(2.4),
    ["UKATT (UK Alcohol Treatment Trial):",
     "- Largest UK alcohol trial (742 clients)",
     "- MET (3 sessions) vs SBNT (8 sessions)",
     "- Result: BOTH equally effective",
     "- MET significantly more cost-effective",
     "- 5x less therapist time for same results",
     "",
     "CONCLUSION: MET achieves equivalent",
     "outcomes at a fraction of the cost",
     "(UKATT Research Team, 2005, BMJ)"],
    bg='light_green', title="UKATT Trial:", ttc='green', border='green', fs=12)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(2.4),
    ["OTHER INTERNATIONAL EVIDENCE:",
     "- WHO Brief Intervention (2002):",
     "  MI effective across 10 countries",
     "- COMBINE Study (2006, JAMA):",
     "  MI + naltrexone superior combination",
     "- Cannabis Youth Treatment (2004):",
     "  MET effective for adolescent cannabis",
     "- Stephens et al. (2004):",
     "  2 sessions of MI = 6 sessions of CBT",
     "",
     "MI/MET works across cultures, substances,",
     "age groups, and clinical settings."],
    bg='light_blue', title="Global Evidence:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LEFT_M, Inches(3.8), FULL_W, Inches(2.5),
    ["SUMMARY OF RESEARCH:",
     "- EMPIRICALLY SUPPORTED for alcohol use disorders (Level 1 evidence, highest level)",
     "- Effective for cannabis, cocaine, opioids, and polysubstance use",
     "- Works with adolescents, adults, and older adults across cultures",
     "- Effective when delivered by trained lay counselors (not just specialists)",
     "- Cost-effective: same results as longer treatments in fewer sessions",
     "- Compatible with pharmacotherapy and enhances medication adherence"],
    bg='cream', title="Bottom Line:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: UKATT (2005). BMJ, 331, 544; COMBINE (2006). JAMA, 295; WHO (2002).")



# ============================================================
# CLOSING SLIDES
# ============================================================

# Key Takeaways
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'cream')
title_bar(slide, "Key Takeaways: MET in Clinical Practice", 'deep_blue')
takeaways = [
    ("MET is BRIEF but POWERFUL", "4 sessions = outcomes of 12-session treatments", 'light_blue', 'deep_blue'),
    ("Based on SOLID THEORY", "Transtheoretical Model + Self-Efficacy + Cognitive Dissonance", 'light_green', 'green'),
    ("The CLIENT does the work", "Therapist evokes motivation, doesn't install it", 'light_purple', 'purple'),
    ("RESISTANCE is information", "Not to fight - a signal to change approach", 'light_orange', 'orange'),
    ("Works with COMORBIDITY", "Adaptable for PD, depression, anxiety, PTSD, psychosis", 'light_teal', 'teal'),
    ("EVIDENCE-BASED globally", "Project MATCH, UKATT, Lancet PREMIUM trial", 'light_red', 'maroon'),
]
for i, (t, d, bg_c, t_c) in enumerate(takeaways):
    y = Inches(1.15) + Inches(0.88) * i
    cbox(slide, LEFT_M, y, FULL_W, Inches(0.8), [d], bg=bg_c, title=t, ttc=t_c, border=t_c, fs=12)
ref_bar(slide, "Ref: Miller et al. (1992); Project MATCH (1997); UKATT (2005); NIMHANS (2016); Nadkarni et al. (2017).")

# References
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'white', 'light_blue')
title_bar(slide, "Comprehensive References", 'deep_blue')
cbox(slide, LEFT_M, Inches(1.2), HALF_W, Inches(5.1),
    ["PRIMARY SOURCES:",
     "Miller, W.R. et al. (1992). MET Manual. NIAAA",
     "  Project MATCH Monograph Series, Vol. 2.",
     "NIMHANS (2016). Substance Use Disorders Manual.",
     "Miller & Rollnick (2013). MI, 3rd ed. Guilford.",
     "Prochaska & DiClemente (1984). TTM.",
     "Project MATCH (1997). J Stud Alcohol, 58.",
     "UKATT (2005). BMJ, 331, 544.",
     "",
     "COMORBIDITY:",
     "Barrowclough et al. (2001). BMJ, 323, 1-5.",
     "Baker et al. (2010). Addiction, 105, 1560.",
     "Ball et al. (2007). J Nerv Ment Dis, 195.",
     "Hien et al. (2009). JCCP, 77, 607-619.",
     "Drake et al. (2004). Schizophr Bull, 30."],
    bg='white', border='deep_blue', fs=10)
cbox(slide, COL2_X, Inches(1.2), HALF_W, Inches(5.1),
    ["INDIAN RESEARCH:",
     "Nadkarni et al. (2017). Lancet, 389, 186-195.",
     "Chand et al. (2018). Indian J Psychiatry.",
     "Rao et al. (2015). Asian J Psychiatry, 16.",
     "Murthy et al. (2009). NIMHANS Publication.",
     "",
     "META-ANALYSES:",
     "Lundahl et al. (2010). Pat Ed Counsel, 80.",
     "Hettema et al. (2005). Ann Rev Clin Psychol, 1.",
     "Smedslund et al. (2011). Cochrane Review.",
     "Burke et al. (2003). JCCP, 71, 843-861.",
     "",
     "FAMILY:",
     "Smith & Meyers (2004). CRAFT.",
     "Meyers et al. (2002). JCCP, 70, 711."],
    bg='white', border='green', fs=10)
ref_bar(slide, "All references from peer-reviewed journals, published manuals, and institutional publications.")

# Thank You
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, 'deep_blue', 'navy')
tb = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(9), Inches(4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.font.name = 'Times New Roman'
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = ""
p2.space_after = Pt(30)
p3 = tf.add_paragraph()
p3.text = "\"People are generally better persuaded by the reasons"
p3.font.size = Pt(16)
p3.font.italic = True
p3.font.color.rgb = COLORS['gold']
p3.font.name = 'Times New Roman'
p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph()
p4.text = "which they have themselves discovered"
p4.font.size = Pt(16)
p4.font.italic = True
p4.font.color.rgb = COLORS['gold']
p4.font.name = 'Times New Roman'
p4.alignment = PP_ALIGN.CENTER
p5 = tf.add_paragraph()
p5.text = "than by those which have come into the minds of others.\""
p5.font.size = Pt(16)
p5.font.italic = True
p5.font.color.rgb = COLORS['gold']
p5.font.name = 'Times New Roman'
p5.alignment = PP_ALIGN.CENTER
p6 = tf.add_paragraph()
p6.text = "— Blaise Pascal"
p6.font.size = Pt(13)
p6.font.color.rgb = COLORS['light_blue']
p6.font.name = 'Times New Roman'
p6.alignment = PP_ALIGN.CENTER
p6.space_after = Pt(30)
p7 = tf.add_paragraph()
p7.text = "Based on: MET Manual (NIAAA, 1992) & NIMHANS Substance Use Disorders (2016)"
p7.font.size = Pt(11)
p7.font.color.rgb = COLORS['light_blue']
p7.font.name = 'Times New Roman'
p7.alignment = PP_ALIGN.CENTER

# ============================================================
# SAVE
# ============================================================
output_path = '/projects/sandbox/Dango-kiro/MET_Comprehensive_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
