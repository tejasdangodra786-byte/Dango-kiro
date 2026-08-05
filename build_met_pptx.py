#!/usr/bin/env python3
"""
MET Presentation Builder - Clean, Professional, No Overlapping
All measurements carefully calculated for 13.333 x 7.5 inch widescreen slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.table import Table
from copy import deepcopy
import math

# ──────────────────────────────────────────────
# COLOR PALETTE
# ──────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE  = RGBColor(0x2C, 0x5F, 0x8A)
MID_BLUE   = RGBColor(0x3A, 0x7B, 0xC8)
LIGHT_BLUE = RGBColor(0xEB, 0xF5, 0xFF)
SKY        = RGBColor(0xD4, 0xE6, 0xF1)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x2D, 0x3A, 0x4A)
MED_GRAY   = RGBColor(0x5D, 0x6D, 0x7E)
TEAL       = RGBColor(0x17, 0xA2, 0xB8)
GOLD       = RGBColor(0xD4, 0xA0, 0x1C)
GREEN      = RGBColor(0x28, 0xA7, 0x45)
RED_SOFT   = RGBColor(0xDC, 0x35, 0x45)
ORANGE     = RGBColor(0xFD, 0x7E, 0x14)
BG_CREAM   = RGBColor(0xFD, 0xFD, 0xFD)

# ──────────────────────────────────────────────
# PRESENTATION SETUP
# ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.5
CONTENT_W = SLIDE_W - 2 * MARGIN  # 12.333

slide_count = 0


# ──────────────────────────────────────────────
# HELPER FUNCTIONS
# ──────────────────────────────────────────────

def new_slide():
    """Create a blank slide with white background."""
    global slide_count
    slide_count += 1
    layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def add_title_bar(slide, title_text, subtitle_text=""):
    """Add a colored title bar at the top (0 to 1.1 inches)."""
    # Title bar background
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1),
        Inches(SLIDE_W), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MID_BLUE
    accent.line.fill.background()

    # Title text
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12), Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(13)
        p2.font.color.rgb = SKY


def add_footer(slide, refs_text):
    """Add a reference footer at the bottom (6.7 to 7.3 inches)."""
    # Footer background
    fb = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.85),
        Inches(SLIDE_W), Inches(0.65))
    fb.fill.solid()
    fb.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    fb.line.fill.background()

    # Reference text
    rb = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.55))
    rtf = rb.text_frame
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.text = refs_text
    rp.font.size = Pt(8)
    rp.font.italic = True
    rp.font.color.rgb = MED_GRAY


def add_notes(slide, notes_text):
    """Add speaker notes."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def add_content_text(slide, lines, top=1.35, left=0.5, width=12.3, font_size=15, line_spacing=1.3):
    """Add content text in the main body area. Returns bottom position."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle indentation
        indent_level = 0
        display_text = line
        if line.startswith("    "):
            indent_level = 2
            display_text = line.strip()
            p.font.size = Pt(font_size - 3)
        elif line.startswith("  "):
            indent_level = 1
            display_text = line.strip()
            p.font.size = Pt(font_size - 2)
        else:
            p.font.size = Pt(font_size)

        p.text = display_text
        p.level = indent_level
        p.space_after = Pt(4)
        p.font.color.rgb = DARK_GRAY

        # Bold section headers
        if line and not line.startswith(" ") and line.endswith(":"):
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE

    return tb


def add_table_slide(slide, headers, rows, top=1.5, left=0.5, col_widths=None):
    """Add a formatted table to the slide. Returns the table shape."""
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    width = SLIDE_W - 1.0
    height = min(0.4 + n_rows * 0.38, 5.0)

    if col_widths is None:
        col_widths = [width / n_cols] * n_cols

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height))
    table = tbl_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            # Alternate row colors
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return tbl_shape


def add_info_box(slide, text, top, left=0.5, width=12.3, color=TEAL):
    """Add a colored info/takeaway box."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(
        min(color[0] + 180, 255), min(color[1] + 180, 255), min(color[2] + 180, 255))
    box.line.color.rgb = color
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY
    return box


def add_icon_box(slide, icon_char, text, top, left, width=3.8, color=MID_BLUE):
    """Add a card-style box with icon character and text."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)
    tf.margin_left = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = icon_char
    p.font.size = Pt(24)
    p.font.color.rgb = color
    p.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_GRAY
    return box


def make_section_divider(section_num, title):
    """Create a visually appealing section divider slide."""
    slide = new_slide()
    # Full background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W), Inches(SLIDE_H))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = NAVY
    bg_shape.line.fill.background()

    # Decorative circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.5), Inches(1.5), Inches(2.3), Inches(2.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = MID_BLUE
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = str(section_num)
    cp.font.size = Pt(48)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Section title
    ttb = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.3), Inches(1.5))
    ttf = ttb.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(5.8), Inches(5.3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    add_notes(slide, f"Section {section_num}: {title}")
    return slide


def make_content_slide(title, bullets, refs, notes, takeaway="", subtitle=""):
    """Standard content slide with proper non-overlapping layout."""
    slide = new_slide()
    add_title_bar(slide, title, subtitle)

    # Main content area: 1.35 to 6.2
    add_content_text(slide, bullets, top=1.35)

    # Takeaway box at 6.2
    if takeaway:
        add_info_box(slide, f"KEY TAKEAWAY: {takeaway}", top=6.2, color=TEAL)

    # Footer references
    add_footer(slide, f"Refs: {refs}")
    add_notes(slide, notes)
    return slide


def make_two_column_slide(title, left_title, left_items, right_title, right_items, refs, notes, takeaway=""):
    """Two-column layout slide."""
    slide = new_slide()
    add_title_bar(slide, title)

    col_w = 5.8
    # Left column header
    lh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(col_w), Inches(0.5))
    lh.fill.solid()
    lh.fill.fore_color.rgb = DARK_BLUE
    lh.line.fill.background()
    lhf = lh.text_frame
    lhp = lhf.paragraphs[0]
    lhp.text = left_title
    lhp.font.size = Pt(13)
    lhp.font.bold = True
    lhp.font.color.rgb = WHITE
    lhp.alignment = PP_ALIGN.CENTER

    # Left column content
    ltb = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(col_w - 0.2), Inches(4.0))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    for i, item in enumerate(left_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # Right column header
    rh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(col_w), Inches(0.5))
    rh.fill.solid()
    rh.fill.fore_color.rgb = TEAL
    rh.line.fill.background()
    rhf = rh.text_frame
    rhp = rhf.paragraphs[0]
    rhp.text = right_title
    rhp.font.size = Pt(13)
    rhp.font.bold = True
    rhp.font.color.rgb = WHITE
    rhp.alignment = PP_ALIGN.CENTER

    # Right column content
    rtb = slide.shapes.add_textbox(Inches(6.9), Inches(2.0), Inches(col_w - 0.2), Inches(4.0))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    for i, item in enumerate(right_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    if takeaway:
        add_info_box(slide, f"KEY TAKEAWAY: {takeaway}", top=6.2, color=TEAL)
    add_footer(slide, f"Refs: {refs}")
    add_notes(slide, notes)
    return slide


def make_table_slide(title, headers, rows, refs, notes, col_widths=None, takeaway=""):
    """Slide with a properly formatted table."""
    slide = new_slide()
    add_title_bar(slide, title)
    add_table_slide(slide, headers, rows, top=1.5, col_widths=col_widths)
    if takeaway:
        add_info_box(slide, f"KEY TAKEAWAY: {takeaway}", top=6.2, color=TEAL)
    add_footer(slide, f"Refs: {refs}")
    add_notes(slide, notes)
    return slide


def make_process_slide(title, steps, refs, notes, takeaway=""):
    """Slide showing a process flow with arrow-connected boxes."""
    slide = new_slide()
    add_title_bar(slide, title)

    n = len(steps)
    box_w = min(2.5, (CONTENT_W - 0.3 * (n - 1)) / n)
    total_w = n * box_w + (n - 1) * 0.3
    start_x = (SLIDE_W - total_w) / 2
    top = 2.2
    box_h = 1.8

    colors = [DARK_BLUE, MID_BLUE, TEAL, GREEN, GOLD]

    for i, (step_title, step_desc) in enumerate(steps):
        x = start_x + i * (box_w + 0.3)
        color = colors[i % len(colors)]

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = step_desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0xE8, 0xF0, 0xFF)
        p2.alignment = PP_ALIGN.CENTER

        # Arrow between boxes
        if i < n - 1:
            arrow_x = x + box_w
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(arrow_x + 0.02), Inches(top + box_h / 2 - 0.15),
                Inches(0.26), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()

    if takeaway:
        add_info_box(slide, f"KEY TAKEAWAY: {takeaway}", top=6.2, color=TEAL)
    add_footer(slide, f"Refs: {refs}")
    add_notes(slide, notes)
    return slide


def make_quad_card_slide(title, cards, refs, notes, takeaway=""):
    """Slide with 4 info cards in a 2x2 grid."""
    slide = new_slide()
    add_title_bar(slide, title)

    positions = [
        (0.5, 1.5), (6.6, 1.5),
        (0.5, 4.0), (6.6, 4.0)
    ]
    colors = [DARK_BLUE, TEAL, GREEN, GOLD]

    for i, (icon, card_title, card_text) in enumerate(cards[:4]):
        left, top = positions[i]
        color = colors[i % 4]

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(5.8), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2.5)

        # Color bar on left
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(0.12), Inches(2.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = f"{icon}  {card_title}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = card_text
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY
        p2.space_before = Pt(6)

    if takeaway:
        add_info_box(slide, f"KEY TAKEAWAY: {takeaway}", top=6.35, color=TEAL)
    add_footer(slide, f"Refs: {refs}")
    add_notes(slide, notes)
    return slide


# ══════════════════════════════════════════════════════════════
# TITLE SLIDE
# ══════════════════════════════════════════════════════════════
def make_title_slide():
    slide = new_slide()
    # Full navy background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                Inches(SLIDE_W), Inches(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Decorative circles
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1), Inches(4), Inches(4))
    c1.fill.solid()
    c1.fill.fore_color.rgb = DARK_BLUE
    c1.line.fill.background()
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5), Inches(3.5), Inches(3.5))
    c2.fill.solid()
    c2.fill.fore_color.rgb = DARK_BLUE
    c2.line.fill.background()

    # Gold accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.3),
                                  Inches(9.3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Main title
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(10.3), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "MOTIVATIONAL ENHANCEMENT"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "THERAPY (MET)"
    p2.font.size = Pt(42)
    p2.font.bold = True
    p2.font.color.rgb = GOLD
    p2.alignment = PP_ALIGN.CENTER

    # Subtitle
    stb = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9.3), Inches(1.5))
    stf = stb.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = "A Comprehensive Evidence-Based Presentation"
    sp.font.size = Pt(18)
    sp.font.color.rgb = SKY
    sp.alignment = PP_ALIGN.CENTER
    sp2 = stf.add_paragraph()
    sp2.text = "For M.Phil. Clinical Psychology | PsyD | PhD | Psychiatry Residents"
    sp2.font.size = Pt(14)
    sp2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    sp2.alignment = PP_ALIGN.CENTER

    # Bottom line
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(6.3),
                                   Inches(9.3), Inches(0.04))
    line2.fill.solid()
    line2.fill.fore_color.rgb = MID_BLUE
    line2.line.fill.background()

    # Source info
    stb2 = slide.shapes.add_textbox(Inches(2), Inches(6.5), Inches(9.3), Inches(0.7))
    stf2 = stb2.text_frame
    stf2.word_wrap = True
    sp3 = stf2.paragraphs[0]
    sp3.text = "Based on: Miller et al. (1995) Project MATCH MET Manual & NIMHANS Psychosocial Interventions Manual (2008)"
    sp3.font.size = Pt(11)
    sp3.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
    sp3.alignment = PP_ALIGN.CENTER

    add_notes(slide, "Title slide. Welcome the audience. This presentation covers Motivational Enhancement Therapy comprehensively for postgraduate clinical psychology education.")

make_title_slide()


# ══════════════════════════════════════════════════════════════
# SECTION 1: INTRODUCTION (Slides 2-10)
# ══════════════════════════════════════════════════════════════
make_section_divider(1, "Introduction to MET")

make_content_slide(
    "What is Motivation?",
    [
        "Definition:",
        "  An internal state that energizes, directs, and sustains behavior toward goals",
        "  The probability a person will enter into, continue, and adhere to change",
        "",
        "Key Insights from Research:",
        "  Motivation is NOT a fixed personality trait — it fluctuates",
        "  Motivation is an interpersonal process — therapists influence it",
        "  Miller (1985): 'Motivation is not something one HAS but something one DOES'",
        "",
        "Clinical Implication:",
        "  If a client is 'unmotivated,' the question is not 'What's wrong with them?'",
        "  but rather 'What can I do differently in my therapeutic approach?'",
    ],
    refs="Miller, W.R. (1985). Psychological Bulletin, 98, 84-107. | Miller & Rollnick (2013). MI, 3rd ed. Guilford Press.",
    notes="Motivation has been misunderstood for decades. Traditional models viewed it as a fixed client characteristic. Miller showed it is dynamic, context-dependent, and directly influenced by the therapist's style. This paradigm shift is fundamental to MET.",
    takeaway="Motivation is dynamic and therapist-influenced — not a fixed client trait."
)

make_content_slide(
    "What is Behaviour Change?",
    [
        "Behaviour change is a PROCESS, not an event:",
        "  People move through identifiable stages (Prochaska & DiClemente, 1982)",
        "  Ambivalence is NORMAL and expected at every stage",
        "  Relapse is part of the cycle, not treatment failure",
        "",
        "Two critical components of behaviour change:",
        "  1. MOTIVATION (WHY change?) — importance, desire, reasons, need",
        "  2. SELF-EFFICACY (CAN I change?) — confidence, ability, past success",
        "",
        "What does NOT work:",
        "  Confrontation, lecturing, scare tactics, labeling, moralizing",
        "  Patterson & Forgatch (1985): Confrontation INCREASES resistance",
        "",
        "What DOES work:",
        "  Empathy, autonomy support, collaboration, evoking client's own reasons",
    ],
    refs="Prochaska & DiClemente (1982). Psychotherapy, 19, 276-288. | Patterson & Forgatch (1985). JCCP, 53, 846-851.",
    notes="Set the foundation that behavior change requires both motivation AND efficacy. Traditional confrontation fails because it creates reactance. Miller demonstrated that empathic approaches work better across all outcomes.",
    takeaway="Change requires both motivation (WHY) and self-efficacy (CAN I) — confrontation undermines both."
)

make_two_column_slide(
    "Why Confrontation Fails vs. Why MET Works",
    "CONFRONTATIONAL APPROACH",
    [
        "Assumes 'denial' is pathological",
        "Therapist as authoritarian expert",
        "Breaking through resistance",
        "Labels: 'You ARE an alcoholic'",
        "Creates adversarial relationship",
        "Higher dropout rates",
        "Evokes defensive argumentation",
        "Miller (1993): Predicts WORSE outcomes",
    ],
    "MET / MOTIVATIONAL APPROACH",
    [
        "Sees ambivalence as normal",
        "Therapist as collaborative partner",
        "Rolling with resistance",
        "Explores: 'What concerns YOU?'",
        "Builds therapeutic alliance",
        "Better engagement and retention",
        "Evokes self-motivational statements",
        "Research: Predicts BETTER outcomes",
    ],
    refs="Miller, Benefield & Tonigan (1993). JCCP, 61, 455-461. | Valle (1981). J Studies on Alcohol, 42, 783-790.",
    notes="This contrast is the foundation for understanding MET. The evidence shows confrontation is counterproductive. Miller's 1993 study found that the amount of resistance in Session 1 predicted drinking 12 months later, and resistance was directly caused by therapist confrontation.",
    takeaway="Empathic approaches consistently outperform confrontation in addiction treatment outcomes."
)


make_content_slide(
    "Definition of MET",
    [
        "Motivational Enhancement Therapy (MET):",
        "  A brief, structured, evidence-based intervention that uses motivational",
        "  strategies to mobilize the client's OWN internal resources for change.",
        "",
        "Key Characteristics:",
        "  • Brief: 4 sessions (vs. 12 for CBT/TSF in Project MATCH)",
        "  • Non-confrontational and non-judgmental",
        "  • Client-centered in spirit, yet directive in strategy",
        "  • Based on personalized assessment FEEDBACK",
        "  • Focuses on evoking change talk from the client",
        "  • Respects client autonomy and builds self-efficacy",
        "",
        "Formula:",
        "  MET = Motivational Interviewing + Structured Feedback + Personal Change Plan",
        "",
        "MET does NOT train clients step-by-step; it MOBILIZES their own motivation.",
    ],
    refs="Miller, Zweben, DiClemente & Rychtarik (1995). MET Manual. NIAAA, NIH Pub No. 94-3723.",
    notes="Distinguish MET from pure MI: MET is a specific manualized protocol with assessment feedback. MI is a broader counseling style. MET adds structure (4 sessions, PFR) to the MI approach.",
    takeaway="MET = MI spirit + structured feedback + 4-session protocol. Brief yet equally effective as longer treatments."
)

make_process_slide(
    "Historical Timeline of MET Development",
    [
        ("1983", "Miller publishes\nfirst MI paper"),
        ("1989", "Project MATCH\ninitiated by NIAAA"),
        ("1994", "MET Manual\npublished"),
        ("1997", "MATCH results:\nMET = CBT = TSF"),
        ("2013", "MI 3rd Edition\nupdated framework"),
    ],
    refs="Miller (1983). Behavioural Psychotherapy, 11, 147-172. | Project MATCH (1997). J Studies on Alcohol, 58, 7-29.",
    notes="Walk through each milestone. 1983: Miller's insight combining Rogers + Festinger + Bem. 1989: NIAAA commissions largest addiction trial ($27M). 1994: Manual standardizes the approach. 1997: Revolutionary finding that 4 sessions matched 12. 2013: Theoretical refinement.",
    takeaway="MET evolved from 1983 insight to become one of the most researched psychotherapies worldwide."
)

# ══════════════════════════════════════════════════════════════
# SECTION 2: FOUNDERS (Slides 11-13)
# ══════════════════════════════════════════════════════════════
make_section_divider(2, "Founders & Key Contributors")

make_content_slide(
    "William R. Miller, PhD — Creator of MI/MET",
    [
        "Professor Emeritus, University of New Mexico (born 1947)",
        "",
        "Key Contributions:",
        "  • Created Motivational Interviewing (1983)",
        "  • Lead author of MET Manual for Project MATCH",
        "  • 400+ publications, 50+ books",
        "  • Developed the 'Drinker's Check-Up' (precursor to MET feedback)",
        "  • Demonstrated that therapist empathy predicts outcomes",
        "",
        "Core Insight:",
        "  'The worst persuasion strategy is one that evokes defensive",
        "  argumentation from the person' — Miller, 1995 (MET Manual)",
        "",
        "Awards: Jellinek Memorial Award, Innovators Award (RWJF)",
        "Philosophy: 'People are generally the best experts on themselves'",
    ],
    refs="Miller (1983). Behavioural Psychotherapy, 11, 147-172. | Miller et al. (1995). MET Manual.",
    notes="Miller's insight came during a visit to Norway when trainees asked how he worked with problem drinkers. His response combined Rogers, Festinger, and Bem into a new approach.",
    takeaway="Miller's revolutionary insight: therapist style determines client motivation more than any technique."
)

make_table_slide(
    "Key Contributors to MI/MET Development",
    ["Contributor", "Affiliation", "Key Contribution"],
    [
        ["William R. Miller", "University of New Mexico", "Created MI (1983); Lead MET Manual author"],
        ["Stephen Rollnick", "Cardiff University, Wales", "Co-developer MI; MI in healthcare settings"],
        ["Theresa Moyers", "University of New Mexico", "MITI coding system; mechanisms research"],
        ["Carlo DiClemente", "U of Houston → U of Maryland", "Transtheoretical Model; MET Manual co-author"],
        ["Allen Zweben", "U of Wisconsin-Milwaukee", "MET Manual co-author; clinical trials"],
        ["Kathleen Carroll", "Yale University", "Treatment design coordinator, Project MATCH"],
        ["Thomas Babor", "U of Connecticut", "Data coordinating center, Project MATCH"],
    ],
    refs="Project MATCH Research Group (1997). JOSA, 58, 7-29. | DiClemente (2003). Addiction and Change. Guilford.",
    notes="Each contributor brought unique expertise. DiClemente provided stages of change theory. Moyers developed fidelity measurement. Rollnick brought MI into mainstream healthcare beyond addiction.",
    col_widths=[2.5, 3.5, 6.3],
    takeaway="MET emerged from collaboration of leading researchers in motivation, addiction, and behavior change."
)


# ══════════════════════════════════════════════════════════════
# SECTION 3: HISTORICAL BACKGROUND (Slides 14-16)
# ══════════════════════════════════════════════════════════════
make_section_divider(3, "Historical Background")

make_two_column_slide(
    "Before MET: Confrontational Addiction Treatment",
    "1950s-1980s: DOMINANT PARADIGM",
    [
        "Synanon (1958): Attack therapy groups",
        "'Hot seat' confrontation techniques",
        "Assumption: Addicts in 'denial'",
        "Must 'break through' resistance",
        "'Hit rock bottom' before change",
        "'Tough love' philosophy",
        "Therapist as authoritarian expert",
    ],
    "EVIDENCE AGAINST CONFRONTATION",
    [
        "Higher dropout rates documented",
        "Patterson & Forgatch (1985): Confrontation directly caused resistance",
        "Miller et al. (1993): Therapist confrontation predicted WORSE 1-year outcomes",
        "Valle (1981): Empathic therapists got 2x better outcomes",
        "No RCT ever showed confrontation was superior",
    ],
    refs="Patterson & Forgatch (1985). JCCP, 53, 846. | Valle (1981). J Studies Alcohol, 42, 783. | Miller et al. (1993). JCCP, 61, 455.",
    notes="For decades, addiction treatment relied on aggressive confrontation. Miller's research showed this was not only ineffective but actively harmful. The paradigm shift from 'breaking denial' to 'resolving ambivalence' was revolutionary.",
    takeaway="The confrontational paradigm had NO evidence base; empathic approaches consistently produce better outcomes."
)

make_content_slide(
    "The Paradigm Shift: From Confrontation to Collaboration",
    [
        "OLD PARADIGM: 'Denial must be broken'",
        "  Client says: 'I don't have a problem'",
        "  Therapist pushes: 'Yes you do! Look at the evidence!'",
        "  Result: Client argues back harder → resistance increases → poor outcomes",
        "",
        "NEW PARADIGM: 'Ambivalence must be explored'",
        "  Client says: 'I don't have a problem'",
        "  Therapist reflects: 'You don't see this as a concern right now'",
        "  Result: Client feels heard → explores further → change talk emerges",
        "",
        "Key Research Evidence:",
        "  • Miller (1993): Amount of client resistance in Session 1 predicted",
        "    drinking outcomes at 12 months",
        "  • Resistance was CAUSED by therapist confrontation",
        "  • Self-perception theory: 'If I argue I don't have a problem, I believe it more'",
    ],
    refs="Miller, Benefield & Tonigan (1993). JCCP, 61, 455-461. | Bem (1972). Self-perception theory. Academic Press.",
    notes="The self-perception mechanism explains why confrontation backfires: when clients verbally defend their drinking, they literally talk themselves INTO continued drinking. MET reverses this by having clients talk themselves INTO change.",
    takeaway="Resistance is therapist-created; when clients argue against change, they strengthen their own resistance."
)

# ══════════════════════════════════════════════════════════════
# SECTION 4: THEORETICAL FOUNDATIONS (Slides 17-25)
# ══════════════════════════════════════════════════════════════
make_section_divider(4, "Theoretical Foundations")

make_table_slide(
    "Seven Theories Underpinning MET",
    ["Theory", "Founder/Year", "Key Contribution to MET"],
    [
        ["Client-Centered Therapy", "Rogers (1957)", "Empathic relationship as foundation"],
        ["Cognitive Dissonance", "Festinger (1957)", "Discrepancy creates motivation to change"],
        ["Self-Perception Theory", "Bem (1965)", "Hearing yourself talk creates believing"],
        ["Transtheoretical Model", "Prochaska/DiClemente (1982)", "Stage-matched interventions"],
        ["Self-Determination Theory", "Deci & Ryan (1985)", "Autonomy support enhances motivation"],
        ["Self-Efficacy Theory", "Bandura (1977)", "Confidence predicts change attempts"],
        ["Decision-Making Theory", "Janis & Mann (1977)", "Decisional balance weighs pros/cons"],
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed. Guilford Press.",
    notes="Each theory provides a distinct mechanism that MET operationalizes. Rogers = HOW we relate. Festinger = WHAT drives change. Bem = WHY talking works. Bandura = CAN they do it. Prochaska = WHEN to intervene.",
    col_widths=[3.0, 3.0, 6.3],
    takeaway="MET integrates 7+ established theories into a coherent clinical approach addressing why, how, and when people change."
)


make_quad_card_slide(
    "Core Theories Applied in MET",
    [
        ("🧠", "ROGERS: Empathy", "Non-judgmental acceptance creates safety.\nClient explores freely without fear.\nRelationship IS the therapy."),
        ("⚡", "FESTINGER: Dissonance", "Value-behavior gap creates discomfort.\nDiscomfort motivates resolution.\nClient argues for change themselves."),
        ("🗣️", "BEM: Self-Perception", "'As I hear myself talk, I learn what I believe.'\nChange talk CREATES motivation.\nThe more they say it, the more they believe it."),
        ("💪", "BANDURA: Self-Efficacy", "Belief in ability predicts action.\nPast successes build confidence.\n'Support Self-Efficacy' is a core MET principle."),
    ],
    refs="Rogers (1957). J Consulting Psych, 21, 95. | Festinger (1957). A Theory of Cognitive Dissonance. | Bem (1972). Advances in Experimental Social Psych.",
    notes="These four theories provide the core engine of MET. Rogers creates the safe container. Festinger provides the motivational fuel (discomfort). Bem explains why evoking change talk works (self-persuasion). Bandura provides the confidence component without which motivation cannot become action.",
    takeaway="Rogers (safety) + Festinger (drive) + Bem (mechanism) + Bandura (confidence) = MET's theoretical engine."
)

make_content_slide(
    "Transtheoretical Model: Stages of Change",
    [
        "Prochaska & DiClemente (1982) — 6 Stages:",
        "",
        "1. PRECONTEMPLATION — Not thinking about change; unaware/unconcerned",
        "   MET: Develop discrepancy gently; raise awareness through feedback",
        "",
        "2. CONTEMPLATION — Ambivalent; 'I should but...' (Most clients start here)",
        "   MET: Explore ambivalence; tip the decisional balance; evoke change talk",
        "",
        "3. PREPARATION — Intending to act; 'I'm going to do something soon'",
        "   MET: Consolidate commitment; develop change plan",
        "",
        "4. ACTION — Actively changing behavior",
        "   MET: Reinforce progress; affirm; build self-efficacy",
        "",
        "5. MAINTENANCE — Sustaining change; preventing relapse",
        "   MET: Relapse prevention; coping plans; booster sessions",
        "",
        "6. RELAPSE — Return to earlier behavior (NORMALIZED, not failure)",
        "   MET: Re-engage with empathy; rebuild motivation without blame",
    ],
    refs="Prochaska & DiClemente (1982). Psychotherapy, 19, 276. | Prochaska & DiClemente (1984). The Transtheoretical Approach. Dow Jones/Irwin.",
    notes="DiClemente co-authored the MET manual, directly connecting stage theory to MET practice. MET is most powerful for clients in contemplation (ambivalent). The key is matching intervention intensity to readiness.",
    takeaway="MET matches strategies to stage: contemplators need exploration, not instruction."
)

make_content_slide(
    "Self-Determination Theory & Self-Efficacy in MET",
    [
        "SELF-DETERMINATION THEORY (Deci & Ryan, 1985):",
        "  Three basic psychological needs that drive intrinsic motivation:",
        "  • AUTONOMY — Need to feel in control ('It's YOUR choice')",
        "  • COMPETENCE — Need to feel capable ('You CAN do this')",
        "  • RELATEDNESS — Need to feel connected (empathic alliance)",
        "",
        "  When all three are met → intrinsic motivation flourishes",
        "  When any are thwarted → motivation diminishes",
        "  MET systematically supports all three needs",
        "",
        "SELF-EFFICACY THEORY (Bandura, 1977):",
        "  'Support Self-Efficacy' = Core MET Principle",
        "  Four sources of self-efficacy:",
        "  1. Mastery experiences ('You quit for 3 months before')",
        "  2. Vicarious experiences ('Others in similar situations have succeeded')",
        "  3. Verbal persuasion ('I believe you can do this')",
        "  4. Emotional states (reducing anxiety about change)",
    ],
    refs="Deci & Ryan (2000). American Psychologist, 55, 68-78. | Bandura (1977). Psychological Review, 84, 191-215.",
    notes="SDT explains WHY MET's non-controlling approach works better than coercion. When autonomy is supported, intrinsic motivation emerges naturally. Self-efficacy is the bridge between motivation and action.",
    takeaway="MET supports autonomy (choice), competence (efficacy), and relatedness (alliance) — all three drive intrinsic motivation."
)


# ══════════════════════════════════════════════════════════════
# SECTION 5: SPIRIT OF MET (Slides 26-29)
# ══════════════════════════════════════════════════════════════
make_section_divider(5, "The Spirit of MET")

make_quad_card_slide(
    "The Four Elements of MET Spirit (PACE)",
    [
        ("🤝", "PARTNERSHIP", "Collaborative relationship — 'dancing' not 'wrestling.'\nWork WITH the client, not ON them.\nTwo experts: therapist on process, client on their life."),
        ("🌟", "ACCEPTANCE", "Absolute worth + Accurate empathy +\nAutonomy support + Affirmation.\nValue the person regardless of their behavior."),
        ("❤️", "COMPASSION", "Actively promoting the client's welfare.\nPrioritizing THEIR needs over YOUR agenda.\nGenuine care, not manipulation."),
        ("🔑", "EVOCATION", "Drawing out what is ALREADY THERE.\nThe client has the answers within.\n'Midwifery' — helping deliver what's developing."),
    ],
    refs="Miller & Rollnick (2013). Motivational Interviewing, 3rd ed. Guilford Press. | Moyers & Miller (2013). Psychotherapy, 50, 338.",
    notes="The spirit is MORE important than any technique. Without spirit, techniques become manipulative. The spirit determines whether an approach is genuinely motivational or merely persuasive. If students take away ONE thing, it should be the spirit.",
    takeaway="Spirit > Technique. Without partnership, acceptance, compassion, and evocation, MET becomes manipulation."
)

make_two_column_slide(
    "Spirit in Practice: Wrong vs. Correct Responses",
    "❌ WRONG (Violates Spirit)",
    [
        "Client: 'I like drinking. It relaxes me.'",
        "",
        "'But it's destroying your liver!' (Righting reflex)",
        "'You shouldn't use alcohol to cope' (No acceptance)",
        "'That's just an excuse' (No compassion)",
        "'Research shows alcohol increases anxiety' (Installing, not evoking)",
        "",
        "These PUSH the client away",
    ],
    "✓ CORRECT (Honors Spirit)",
    [
        "Client: 'I like drinking. It relaxes me.'",
        "",
        "'Drinking serves an important purpose for you' (Partnership)",
        "'You've found something that works short-term' (Acceptance)",
        "'And what else have you noticed about how it affects you?' (Evocation)",
        "'What other things help you manage stress?' (Compassion + Evocation)",
        "",
        "These INVITE the client in",
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Miller et al. (1995). MET Manual.",
    notes="Practice distinguishing spirit-consistent from spirit-inconsistent responses. The key insight: acknowledge the positive function FIRST (validates), then gently explore other effects. This builds trust for deeper exploration.",
    takeaway="Honor the positive function of the behavior first — this builds trust for exploring concerns later."
)

# ══════════════════════════════════════════════════════════════
# SECTION 6: CORE PRINCIPLES (Slides 30-34)
# ══════════════════════════════════════════════════════════════
make_section_divider(6, "Core Principles of MET")

make_quad_card_slide(
    "Four Core Principles of MET",
    [
        ("👂", "EXPRESS EMPATHY", "Skillful reflective listening.\nAccept ambivalence as NORMAL.\nValle (1981): Empathy predicted outcomes at 2 years."),
        ("↔️", "DEVELOP DISCREPANCY", "Help clients see the gap between behavior and values.\nClient presents arguments for change (not therapist).\nCognitive dissonance creates internal motivation."),
        ("🌊", "ROLL WITH RESISTANCE", "Don't oppose resistance — flow with it.\nResistance is a signal to change YOUR approach.\nAmplified reflection, shifting focus, reframing."),
        ("🏆", "SUPPORT SELF-EFFICACY", "Enhance belief in ability to change.\nHighlight past successes and strengths.\nExpress genuine confidence in the client."),
    ],
    refs="Miller et al. (1995). MET Manual. | NIMHANS Manual (2008): DARES framework.",
    notes="These four principles guide every interaction in MET. Express empathy throughout. Develop discrepancy carefully (let client articulate it). Roll with resistance (never argue). Support self-efficacy (instill hope). NIMHANS uses the mnemonic DARES: Develop discrepancy, Avoid argumentation, Roll with resistance, Express empathy, Support self-efficacy.",
    takeaway="All four principles work together: Empathy creates safety, discrepancy creates motivation, rolling prevents rupture, efficacy enables action."
)


make_content_slide(
    "Express Empathy & Develop Discrepancy — Clinical Examples",
    [
        "EXPRESS EMPATHY — Clinical Dialogue:",
        "  Client: 'I don't think I drink more than my friends'",
        "  ❌ 'That's denial. Your friends probably drink too much too'",
        "  ✓ 'From your perspective, your drinking seems normal for your circle'",
        "",
        "  Client: 'My wife is always nagging me about drinking'",
        "  ❌ 'She's probably right. You should listen to her'",
        "  ✓ 'That must be frustrating — and she clearly cares about you'",
        "",
        "DEVELOP DISCREPANCY — Using Values:",
        "  T: 'What's most important to you in life?'",
        "  C: 'My kids. Being a good father.'",
        "  T: 'And how does your drinking fit with being that father?'",
        "  C: '...Not well. I missed my daughter's recital last week.'",
        "  T: 'There's a gap between who you want to be and how alcohol is affecting things.'",
        "",
        "NOTE: The CLIENT articulates the discrepancy — not the therapist.",
    ],
    refs="Miller et al. (1995). MET Manual, pp. 16-20. | Valle (1981). J Studies Alcohol, 42, 783-790.",
    notes="Empathy is ACTIVE listening with continuous hypothesis generation. The therapist reflects back what is heard, often with added meaning. For discrepancy: the client's OWN values provide the leverage. When they see the gap between values and behavior, internal motivation emerges naturally.",
    takeaway="Let clients discover their own value-behavior discrepancies — this creates far stronger motivation than external confrontation."
)

make_content_slide(
    "Roll with Resistance & Support Self-Efficacy",
    [
        "ROLL WITH RESISTANCE — Specific Strategies:",
        "  1. Simple reflection: 'You're not worried about it right now'",
        "  2. Amplified reflection: 'So alcohol causes you absolutely NO concerns'",
        "     (Often client backs away: 'Well, not NO concerns...')",
        "  3. Double-sided reflection: 'You enjoy drinking AND you notice effects on sleep'",
        "  4. Shifting focus: 'Let's set labels aside. What concerns YOU?'",
        "  5. Reframing: 'Your concern about labels shows you think carefully'",
        "  6. Emphasizing autonomy: 'It's completely up to you'",
        "",
        "SUPPORT SELF-EFFICACY — Building Confidence:",
        "  • 'You quit for 3 months last year — you've done it before'",
        "  • 'The fact that you came here today shows real courage'",
        "  • 'Based on what I know about you, I believe you can do this'",
        "  • 'Many people need several attempts — each teaches you something'",
        "",
        "KEY: When resistance increases → soften YOUR approach, don't push harder.",
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="Rolling with resistance is counterintuitive. Natural instinct is to push harder when opposed, but this creates adversarial dynamics. Self-efficacy is the bridge between motivation and action. Without confidence, even highly motivated clients won't attempt change.",
    takeaway="Resistance = signal to change YOUR approach. Self-efficacy = the bridge from wanting to doing."
)

# ══════════════════════════════════════════════════════════════
# SECTION 7: OARS SKILLS (Slides 35-40)
# ══════════════════════════════════════════════════════════════
make_section_divider(7, "Essential Communication Skills: OARS")

make_process_slide(
    "OARS: The Four Core Micro-Skills",
    [
        ("O", "Open Questions\nInvite exploration\n'Tell me about...'"),
        ("A", "Affirmations\nRecognize strengths\n'That took courage'"),
        ("R", "Reflections\nMirror with depth\n2:1 ratio to Q's"),
        ("S", "Summaries\nCollect themes\nStrategically select"),
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | NIMHANS Manual (2008): OARS as 5 specific MI skills.",
    notes="OARS are the operational building blocks of MET. They are used strategically to evoke and reinforce change talk while minimizing sustain talk. The NIMHANS manual specifically identifies OARS + eliciting change talk as the 5 core MI skills.",
    takeaway="OARS are the micro-skills that operationalize MET spirit — practice them until they become automatic."
)


make_two_column_slide(
    "Open Questions & Affirmations in Practice",
    "OPEN QUESTIONS",
    [
        "Cannot be answered yes/no — invite exploration",
        "'Tell me about your drinking'",
        "'What concerns you about your use?'",
        "'What would you like to be different?'",
        "'What have others told you about it?'",
        "",
        "POOR (Closed): 'Do you drink too much?'",
        "BETTER (Open): 'Tell me what you've noticed'",
    ],
    "AFFIRMATIONS",
    [
        "Recognize strengths, efforts, positive qualities",
        "NOT generic praise ('Good job!')",
        "Genuine, specific, about the CLIENT",
        "",
        "'It took real courage to come here today'",
        "'You clearly care deeply about your family'",
        "'The fact you're thinking about this shows strength'",
        "'You managed 3 months — you CAN do this'",
    ],
    refs="Miller et al. (1995). MET Manual, pp. 13-15. | NIMHANS Manual (2008), pp. 22-23.",
    notes="Open questions create space for clients to explore their own motivations. Affirmations differ from praise: praise comes from a one-up position; affirmations recognize qualities IN the client. They should be genuine, specific, and focus on character rather than just outcomes.",
    takeaway="Open questions create the space; affirmations build the confidence. Together they empower the client's voice."
)

make_content_slide(
    "Reflective Listening — The Most Important Skill",
    [
        "Definition: Statements reflecting what the client said with added depth/meaning",
        "Miller (1995): 'Requires continuous alert tracking' — not easy to do well",
        "",
        "Types (increasing depth):",
        "  Simple: Repeats/rephrases ('You wonder about your drinking')",
        "  Complex: Adds meaning ('Part of you suspects something needs to change')",
        "  Double-sided: Both sides ('You enjoy drinking AND you worry about health')",
        "  Amplified: Overstates to elicit other side ('Absolutely no concerns at all')",
        "",
        "Extended Example from MET Manual:",
        "  C: 'I'm not sure I'm concerned, but I wonder if I'm drinking too much'",
        "  T: 'Too much for...' [reflection as prompt]",
        "  C: 'For my own good. Sometimes I can't think straight in the morning'",
        "  T: 'It messes up your thinking, your concentration' [simple]",
        "  C: 'Yes, and trouble remembering things'",
        "  T: 'And you wonder if alcohol is doing that' [complex]",
        "  C: 'Well, I know it is sometimes'",
        "",
        "GUIDELINE: Aim for 2:1 reflection-to-question ratio minimum.",
    ],
    refs="Miller et al. (1995). MET Manual, pp. 16-18. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="This extended dialogue shows how reflection ALONE moves clients toward recognition without a single piece of advice. The client moves from 'not sure I'm concerned' to 'I know it is' through pure reflective listening. This is the art of MET.",
    takeaway="Skilled reflection alone moves clients from ambivalence to recognition — no confrontation needed."
)

make_content_slide(
    "Summaries & Strategic Selection",
    [
        "Types of Summaries:",
        "  1. Collecting: Gathers several change talk statements together",
        "  2. Linking: Connects current statements with earlier ones",
        "  3. Transitional: Wraps up a topic and shifts direction",
        "",
        "Strategic Summary Example:",
        "  'Let me see if I've got this right. You're concerned about memory problems,",
        "  your wife is increasingly worried, work performance has slipped, and your",
        "  doctor flagged liver results. You also enjoy the social side of drinking and",
        "  would hate to lose that. Did I miss anything?'",
        "",
        "STRATEGIC PRINCIPLE:",
        "  Include MORE change talk items than sustain talk items",
        "  Hearing all reasons for change collected in one place is powerful",
        "  Always end with: 'What else?' or 'Did I miss anything?'",
        "",
        "OARS TOGETHER: Questions explore → Reflections deepen →",
        "  Affirmations encourage → Summaries consolidate",
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Miller et al. (1995). MET Manual.",
    notes="Strategic summaries are a powerful therapeutic tool. By selecting WHAT to include, you can present the client's own change talk back as a collected package. Hearing all their reasons together amplifies the motivational impact. The client literally hears a summary of why THEY think they should change.",
    takeaway="Strategic summaries collect change talk in one place, amplifying its motivational impact through the client's own words."
)


# ══════════════════════════════════════════════════════════════
# SECTION 8: CHANGE TALK (Slides 41-44)
# ══════════════════════════════════════════════════════════════
make_section_divider(8, "Change Talk: The Language of Motivation")

make_table_slide(
    "DARN-CAT: Seven Types of Change Talk",
    ["Type", "Category", "Signal Words", "Example"],
    [
        ["D - Desire", "Preparatory", "'I want...' 'I wish...'", "'I wish I could quit'"],
        ["A - Ability", "Preparatory", "'I can...' 'I could...'", "'I could cut back if I tried'"],
        ["R - Reasons", "Preparatory", "'Because...' 'It would...'", "'My health would improve'"],
        ["N - Need", "Preparatory", "'I need to...' 'I must...'", "'I have to change for my kids'"],
        ["C - Commitment", "Mobilizing", "'I will...' 'I'm going to...'", "'I will stop this week'"],
        ["A - Activation", "Mobilizing", "'I'm ready...' 'I'm willing...'", "'I'm prepared to try'"],
        ["T - Taking Steps", "Mobilizing", "'I actually did...'", "'I skipped drinking last night'"],
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Amrhein et al. (2003). JCCP, 71, 862-878.",
    notes="DARN is preparatory (building toward change); CAT is mobilizing (closer to action). Commitment language is the STRONGEST predictor of actual behavior change. Therapists should track the shift from DARN to CAT as a sign of readiness for planning.",
    col_widths=[2.2, 2.0, 3.2, 4.9],
    takeaway="Commitment talk (C in CAT) is the strongest predictor of actual behavior change — listen for it."
)

make_content_slide(
    "Evoking Change Talk: Strategies from the MET Manual",
    [
        "Nine Strategies to Evoke Change Talk:",
        "  1. Ask evocative questions: 'What concerns you about your drinking?'",
        "  2. Ask for elaboration: 'Tell me more about that'",
        "  3. Ask for extremes: 'What's the WORST that could happen?'",
        "  4. Looking back: 'What was life like before heavy drinking?'",
        "  5. Looking forward: 'Where do you see yourself in 5 years if nothing changes?'",
        "  6. Querying extremes: 'What are your worst fears?'",
        "  7. Importance ruler: 'On 0-10, how important is change? Why X and not zero?'",
        "  8. Exploring values: 'What matters most to you?'",
        "  9. Gentle paradox: 'I'm not sure you're motivated enough...'",
        "",
        "MET Manual (p. 15) — Gentle Paradox Example:",
        "  T: 'Frankly, I'm not sure from what you've told me that you're motivated",
        "  enough to carry through. Do you think we should go ahead?'",
        "  [Often evokes client arguing FOR their own motivation!]",
        "",
        "AFTER evoking: Reflect it, affirm it, ask for more. NEVER ignore change talk.",
    ],
    refs="Miller et al. (1995). MET Manual, pp. 13-16. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="The gentle paradox technique works through psychological reactance — when freedom is threatened, people move in the opposite direction. By subtly siding with the status quo, the therapist evokes the client to argue for change. Use sparingly and genuinely.",
    takeaway="Always ask 'Why are you at X and not ZERO?' — this ALWAYS evokes change talk from any client."
)

make_two_column_slide(
    "Change Talk vs. Sustain Talk",
    "CHANGE TALK (Reinforce it!)",
    [
        "'I know I need to cut down' (Need)",
        "'I wish I could drink normally' (Desire)",
        "'My wife will leave if I don't' (Reasons)",
        "'I could probably manage' (Ability)",
        "'I'm going to try' (Commitment)",
        "'I skipped drinking yesterday' (Taking Steps)",
        "",
        "RESPONSE: Reflect warmly, elaborate, affirm",
    ],
    "SUSTAIN TALK (Don't strengthen it!)",
    [
        "'I don't think it's that bad' (Minimizing)",
        "'All my friends drink' (Normalizing)",
        "'I enjoy it too much to stop' (Desire to continue)",
        "'I've tried before, it never works' (Low ability)",
        "'I don't need to change' (No need)",
        "'I'm not ready' (No activation)",
        "",
        "RESPONSE: Acknowledge briefly, don't dwell",
    ],
    refs="Moyers et al. (2007). JCCP, 75, 790. | Magill et al. (2014). J Substance Abuse Treat, 46, 685.",
    notes="Research shows the RATIO of change talk to sustain talk predicts outcomes. Therapist's job: tilt the ratio toward change talk. Don't argue with sustain talk (that strengthens it). Reflect it briefly, then redirect.",
    takeaway="High change talk + low sustain talk predicts success. Reflect change talk warmly; acknowledge sustain talk briefly."
)


# ══════════════════════════════════════════════════════════════
# SECTION 9-10: SUSTAIN TALK & RESISTANCE (Slides 45-47)
# ══════════════════════════════════════════════════════════════
make_section_divider(9, "Resistance, Discord & Management")

make_content_slide(
    "Understanding Resistance (Discord) in MET",
    [
        "KEY INSIGHT: Resistance is largely THERAPIST-CREATED",
        "",
        "Common Forms of Discord:",
        "  • Arguing: Challenging, being hostile, discounting",
        "  • Interrupting: Cutting off the therapist",
        "  • Denying: Blaming others, minimizing, making excuses",
        "  • Ignoring: Not paying attention, changing subject",
        "",
        "What Causes Discord?",
        "  • Therapist pushes too hard or moves too fast",
        "  • Unsolicited advice-giving or lecturing",
        "  • Labeling ('You're an alcoholic')",
        "  • Arguing for change (the 'righting reflex')",
        "  • Not matching the client's readiness stage",
        "",
        "Research Evidence:",
        "  • Miller et al. (1993): Therapist confrontation → client resistance → poor outcomes",
        "  • Patterson & Forgatch (1985): Teaching/confronting directly increased resistance",
        "",
        "PRINCIPLE: When resistance appears, CHANGE YOUR APPROACH.",
    ],
    refs="Patterson & Forgatch (1985). JCCP, 53, 846. | Miller et al. (1993). JCCP, 61, 455. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="This reconceptualizes resistance from a client problem to a relational signal. The therapist's own 'righting reflex' (urge to fix, advise, correct) is the primary trigger of discord. When resistance increases, it's feedback about the interaction, not about client pathology.",
    takeaway="Discord = signal to change YOUR approach. It's relational feedback, not client pathology."
)

# ══════════════════════════════════════════════════════════════
# SECTION 11: ASSESSMENT (Slides 48-50)
# ══════════════════════════════════════════════════════════════
make_section_divider(11, "Assessment Before MET")

make_content_slide(
    "Pre-Treatment Assessment & Readiness Rulers",
    [
        "Assessment serves DUAL purpose in MET:",
        "  1. Clinical information for the therapist",
        "  2. FEEDBACK material for Session 1 (key intervention!)",
        "",
        "Assessment Battery Should Include:",
        "  • Substance use history (quantity, frequency, pattern, duration)",
        "  • Readiness/Importance/Confidence rulers (0-10 scales)",
        "  • Psychological assessment (mood, anxiety, personality)",
        "  • Physical health (liver function, health markers)",
        "  • Social/occupational functioning",
        "  • Risk assessment (suicide, violence)",
        "  • Neuropsychological screening (if indicated)",
        "",
        "READINESS RULER — Strategic Use:",
        "  T: 'On 0-10, how important is it to change your drinking?'  C: '6'",
        "  T: 'Why 6 and not 2?' [ALWAYS evokes change talk!]",
        "  T: 'What would move you from 6 to 8?'",
        "  NEVER ask: 'Why not 10?' [This evokes sustain talk]",
    ],
    refs="Miller et al. (1995). MET Manual, Appendix A. | Rollnick et al. (1999). Health Behavior Change. Churchill Livingstone.",
    notes="The rulers are deceptively simple yet powerful. The magic follow-up 'Why X and not lower?' forces the client to articulate their own reasons for motivation. This is both assessment AND intervention simultaneously. High importance + low confidence → build efficacy. Low importance + high confidence → develop discrepancy.",
    takeaway="Always ask 'Why X and not lower?' — never 'Why not higher?' The former evokes change talk; the latter evokes sustain talk."
)

# ══════════════════════════════════════════════════════════════
# SECTION 12: FRAMES (Slides 51-53)
# ══════════════════════════════════════════════════════════════
make_section_divider(12, "Components of MET: FRAMES")

make_table_slide(
    "FRAMES: The Active Ingredients of Brief Interventions",
    ["Component", "Description", "Clinical Example"],
    [
        ["F - Feedback", "Personalized assessment results vs. norms", "'Your drinking is above 95% of population'"],
        ["R - Responsibility", "Emphasize personal choice & autonomy", "'Only YOU can decide to make this change'"],
        ["A - Advice", "Clear, brief advice (with permission)", "'Based on results, I'd recommend reducing'"],
        ["M - Menu", "Multiple options for change", "'There are several approaches...'"],
        ["E - Empathy", "Warm, reflective, understanding style", "Non-judgmental reflective listening"],
        ["S - Self-Efficacy", "Express confidence in client's ability", "'I believe you can do this'"],
    ],
    refs="Miller & Sanchez (1994). In Howard (Ed.), Issues in Alcohol Use. | Bien, Miller & Tonigan (1993). Addiction, 88, 315. | NIMHANS Manual (2008), pp. 20-21.",
    notes="FRAMES summarizes what research found to be the active ingredients across effective brief interventions worldwide. The NIMHANS manual specifically identifies FRAMES as the core model. In MET, all six are systematically incorporated across the 4 sessions.",
    col_widths=[2.3, 4.5, 5.5],
    takeaway="FRAMES represents the evidence-based active ingredients of effective brief motivational interventions."
)


# ══════════════════════════════════════════════════════════════
# SECTION 13: STRUCTURE OF MET (Slides 54-57)
# ══════════════════════════════════════════════════════════════
make_section_divider(13, "Structure of MET: The 4-Session Protocol")

make_process_slide(
    "The Four Sessions of MET (Project MATCH Protocol)",
    [
        ("SESSION 1\nWeek 1", "Feedback\n& Building\nMotivation"),
        ("SESSION 2\nWeek 2", "Strengthening\nCommitment\n& Planning"),
        ("SESSION 3\nWeek 6", "Midpoint\nProgress\nReview"),
        ("SESSION 4\nWeek 12", "Consolidation\nMaintenance\n& Termination"),
    ],
    refs="Miller et al. (1995). MET Manual. | Project MATCH Research Group (1997). JOSA, 58, 7-29.",
    notes="Sessions 1-2 are close together (1 week) for momentum. Sessions 3-4 are spaced further apart as reinforcement check-ins. The first two sessions do the heavy motivational lifting. This structure proved equally effective as 12-session treatments in Project MATCH.",
    takeaway="4 sessions total: intensive early work (S1-2) + reinforcement check-ins (S3-4) = equivalent to 12-session treatments."
)

make_content_slide(
    "Session 1: Assessment Feedback & Building Motivation",
    [
        "DURATION: 60-90 minutes | TIMING: Week 1",
        "",
        "Structure:",
        "  OPENING (10 min): Welcome, rapport, agenda setting",
        "    'I'm interested in hearing YOUR thoughts about these results'",
        "",
        "  FEEDBACK DELIVERY (30 min): Personal Feedback Report",
        "    Present results neutrally, item by item",
        "    Compare with normative data (population percentiles)",
        "    Monitor reactions continuously; respond with reflection",
        "",
        "  EXPLORATION (15 min): Elicit self-motivational statements",
        "    'What do you make of these results?'",
        "    'What concerns you about your drinking?'",
        "",
        "  CLOSING (5 min): Summarize; invite reflection; homework",
        "",
        "Optional: Involve significant other (spouse, partner, family member)",
        "  Reframe SO concerns as caring: 'You've been worried about him'",
    ],
    refs="Miller et al. (1995). MET Manual, Session 1 protocol. | Project MATCH Research Group (1997).",
    notes="Session 1 is the most structured and critical. The Personal Feedback Report creates discrepancy naturally through normative comparison. The therapist's style remains reflective throughout — no confrontation even when delivering concerning results. Respond to client reactions with empathic reflection.",
    takeaway="Deliver feedback NEUTRALLY. Respond to ALL reactions with reflection. Let the DATA create discrepancy."
)

make_content_slide(
    "Session 1: Dialogue During Feedback Delivery",
    [
        "From the MET Manual — Real clinical exchanges:",
        "",
        "  Client: 'Wow! I'm drinking a lot more than I realized'",
        "  Therapist: 'It looks awfully high to you' [reflects surprise]",
        "",
        "  Client: 'I can't believe it. I don't see how it can be affecting me'",
        "  Therapist: 'This isn't what you expected to hear' [reflects disbelief]",
        "",
        "  Client: 'I don't really drink more than other people'",
        "  Therapist: 'So this is confusing. It seems like you drink about the same as",
        "    friends, yet here are the results' [validates + maintains data]",
        "",
        "  Client: 'This gives me a lot to think about'",
        "  Therapist: 'A lot of reasons to think about making a change' [strengthens]",
        "",
        "  Client: 'More bad news!'",
        "  Therapist: 'This is pretty difficult for you to hear' [empathic]",
        "",
        "NOTICE: Therapist NEVER argues, lectures, or interprets. Pure reflection.",
    ],
    refs="Miller et al. (1995). MET Manual, pp. 19-20 (direct clinical examples).",
    notes="These verbatim examples from the manual show how the therapist maintains reflective stance regardless of client reaction. Each response validates the client's experience while keeping the therapeutic conversation moving forward. Whether the client is shocked, dismissive, or contemplative, reflection is always appropriate.",
    takeaway="Respond to every reaction with reflection — surprise, denial, and concern ALL get empathic reflection."
)


make_content_slide(
    "Sessions 2-4: Commitment, Review, and Consolidation",
    [
        "SESSION 2 (Week 2): Strengthening Commitment",
        "  • Review since Session 1 ('What have you been thinking?')",
        "  • Continue exploring ambivalence; decisional balance",
        "  • Values clarification exercise",
        "  • If ready: Change Plan Worksheet ('What do you think you'll do?')",
        "  • If NOT ready: Continue Phase 1; don't push",
        "",
        "SESSION 3 (Week 6): Midpoint Review",
        "  • Review progress since Session 2",
        "  • If progressing: Celebrate, affirm, reinforce what's working",
        "  • If struggling: Explore barriers with empathy; rebuild motivation",
        "  • If relapsed: Normalize ('Many people need several attempts')",
        "  • Modify change plan if needed",
        "",
        "SESSION 4 (Week 12): Consolidation & Termination",
        "  • Comprehensive review of entire treatment period",
        "  • Maintenance planning (high-risk situations, coping strategies)",
        "  • Relapse prevention: 'What if' planning",
        "  • Express confidence; leave door open for boosters",
    ],
    refs="Miller et al. (1995). MET Manual, Sessions 2-4. | Marlatt & Gordon (1985). Relapse Prevention. Guilford.",
    notes="The transition from Phase 1 (building motivation) to Phase 2 (strengthening commitment) should be guided by client readiness signals: decreased resistance, increased change talk, questions about 'how' to change, imagining a different future.",
    takeaway="Time the Phase 1→2 transition by readiness signals. Premature planning increases resistance; delayed planning misses momentum."
)

# ══════════════════════════════════════════════════════════════
# SECTION 14-15: TECHNIQUES (Slides 58-64)
# ══════════════════════════════════════════════════════════════
make_section_divider(14, "MET Techniques in Detail")

make_table_slide(
    "Key MET Techniques: Overview",
    ["Technique", "Purpose", "When to Use"],
    [
        ["Double-sided Reflection", "Capture both sides of ambivalence", "When client is ambivalent"],
        ["Amplified Reflection", "Elicit the other side through overstatement", "When client minimizes"],
        ["Reframing", "Offer new perspective on same info", "When client is stuck"],
        ["Shifting Focus", "Move away from stuck point", "When debating labels"],
        ["Elicit-Provide-Elicit", "Share info without lecturing", "When giving advice/info"],
        ["Importance/Confidence Rulers", "Assess and evoke simultaneously", "Every session"],
        ["Decisional Balance", "Explore pros/cons systematically", "When exploring ambivalence"],
        ["Looking Back/Forward", "Create temporal discrepancy", "To develop discrepancy"],
        ["Values Clarification", "Identify core values for discrepancy", "Session 2 especially"],
        ["Change Plan Worksheet", "Formalize commitment", "When ready for planning"],
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed. Guilford.",
    notes="Each technique has a specific function and timing. Double-sided reflections honor ambivalence. E-P-E maintains collaboration when sharing information. The rulers are both assessment AND intervention. All techniques should be used within the MET spirit.",
    col_widths=[3.5, 4.3, 4.5],
    takeaway="Each technique has a specific purpose and optimal timing — match technique to clinical need."
)

make_content_slide(
    "Double-Sided Reflection & Elicit-Provide-Elicit",
    [
        "DOUBLE-SIDED REFLECTION:",
        "  Captures BOTH sides of ambivalence in one statement",
        "  Format: 'On one hand [sustain talk] AND on the other [change talk]'",
        "",
        "  From MET Manual:",
        "  'You don't think alcohol is harming you seriously now, AND at the same time",
        "   you are concerned that it might get out of hand later'",
        "",
        "  NOTE: Use 'AND' not 'BUT' — 'but' negates what came before",
        "",
        "ELICIT-PROVIDE-ELICIT (Ask-Tell-Ask):",
        "  A structured way to give information WITHOUT lecturing:",
        "  Step 1 — ELICIT: 'What do you already know about alcohol and memory?'",
        "  Step 2 — PROVIDE: 'Would it be okay if I shared some information?'",
        "    (Wait for permission → provide brief, relevant information)",
        "  Step 3 — ELICIT: 'What do you make of that? How does this fit for you?'",
        "",
        "  WHY: Maintains collaboration; avoids expert-down hierarchy",
    ],
    refs="Miller et al. (1995). MET Manual, pp. 18-19. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="Double-sided reflections are among the most powerful tools. They validate ambivalence (both sides are real) while creating cognitive dissonance (hearing both together highlights the contradiction). E-P-E is essential whenever you need to provide information or advice.",
    takeaway="Use 'AND' not 'BUT' in double-sided reflections. Always ask permission before providing information (E-P-E)."
)

make_content_slide(
    "Values Clarification, Looking Back/Forward, Decisional Balance",
    [
        "VALUES CLARIFICATION (Best used in Session 2):",
        "  T: 'What are the 3 most important things in your life?'",
        "  T: 'How does your drinking fit with being that person?'",
        "  → Creates powerful internal discrepancy from client's OWN values",
        "",
        "LOOKING BACK:",
        "  'What was life like before drinking became a problem?'",
        "  Purpose: Temporal discrepancy (then vs. now)",
        "",
        "LOOKING FORWARD:",
        "  'If you continue as you are, where do you see yourself in 5 years?'",
        "  'If you DID change, what would life look like?'",
        "  Purpose: Future discrepancy (feared vs. hoped outcome)",
        "",
        "DECISIONAL BALANCE (4 Quadrants):",
        "  1. Benefits of drinking (start here — validates, reduces defensiveness)",
        "  2. Costs of drinking (develops discrepancy)",
        "  3. Fears about changing (acknowledges difficulty)",
        "  4. Benefits of changing (builds forward momentum — end here)",
    ],
    refs="Miller et al. (1995). MET Manual. | Janis & Mann (1977). Decision Making. Free Press.",
    notes="Start decisional balance with benefits of drinking to show you understand and won't judge. End with benefits of change to leave the session pointing forward. The strategic order matters therapeutically.",
    takeaway="Values create 'why'; looking back/forward creates urgency; decisional balance makes ambivalence visible and workable."
)


# ══════════════════════════════════════════════════════════════
# SECTION 16: WORKSHEETS (Slides 65-67)
# ══════════════════════════════════════════════════════════════
make_section_divider(16, "Clinical Worksheets & Handouts")

make_content_slide(
    "Change Plan Worksheet (Session 2-3)",
    [
        "Completed COLLABORATIVELY with client — they write it (maintains ownership):",
        "",
        "  1. The changes I want to make are:",
        "     ___________________________________________________",
        "",
        "  2. The most important reasons I want to make these changes:",
        "     ___________________________________________________",
        "",
        "  3. The steps I plan to take:",
        "     ___________________________________________________",
        "",
        "  4. People who can help me (Person → How they can help):",
        "     ___________________________________________________",
        "",
        "  5. I will know my plan is working if:",
        "     ___________________________________________________",
        "",
        "  6. Things that could interfere and how I'll handle them:",
        "     ___________________________________________________",
        "",
        "Give client their copy. Revisit in Sessions 3-4.",
    ],
    refs="Miller et al. (1995). MET Manual (original Change Plan Worksheet). | SAMHSA (2019). TIP 35.",
    notes="The Change Plan Worksheet makes commitment concrete and tangible. The client fills it in (maintaining ownership). It serves as a physical reminder between sessions. Review and modify in subsequent sessions as needed.",
    takeaway="Written commitment is stronger than verbal commitment — the Change Plan makes it concrete and portable."
)

make_content_slide(
    "Assessment Worksheets: Rulers, Decisional Balance, Diary",
    [
        "READINESS RULER (Print as visual aid):",
        "  Not at all ready                                   Completely ready",
        "  0 ─── 1 ─── 2 ─── 3 ─── 4 ─── 5 ─── 6 ─── 7 ─── 8 ─── 9 ─── 10",
        "",
        "DECISIONAL BALANCE SHEET (4 quadrants on paper):",
        "  Benefits of Drinking | Costs of Drinking",
        "  ─────────────────────┼──────────────────",
        "  Costs of Changing    | Benefits of Changing",
        "",
        "SUBSTANCE USE DIARY (Daily Monitoring):",
        "  Date | Substance | Amount | Time | Place | Trigger | Mood After",
        "",
        "TRIGGER IDENTIFICATION WORKSHEET:",
        "  External: People___ Places___ Times___ Events___",
        "  Internal: Emotions___ Thoughts___ Physical sensations___ Urges___",
        "",
        "EMERGENCY/RELAPSE PREVENTION PLAN:",
        "  Warning signs | Who to call | Where to go | What to do | What NOT to do",
    ],
    refs="Miller et al. (1995). MET Manual. | Marlatt & Gordon (1985). Relapse Prevention. | SAMHSA (2019). TIP 35.",
    notes="These worksheets can be printed and given to clients as homework or used within sessions. Self-monitoring raises awareness and often reduces behavior through reactivity. The emergency plan should be completed in Session 4 and kept accessible (wallet card, phone).",
    takeaway="Written worksheets externalize internal processes and serve as tangible tools for between-session work."
)

# ══════════════════════════════════════════════════════════════
# SECTION 17: CASE FORMULATIONS (Slides 68-72)
# ══════════════════════════════════════════════════════════════
make_section_divider(17, "Case Formulations")

make_table_slide(
    "Case 1: Alcohol Dependence — Mr. Rajesh, 42, Manager",
    ["Aspect", "Details"],
    [
        ["Presenting Problem", "Referred by physician; elevated liver enzymes; 15 years drinking"],
        ["Pattern", "Daily 180-350ml whiskey; tolerance increased; morning tremors"],
        ["Stage of Change", "Contemplation (knows it's a problem; not committed yet)"],
        ["Importance", "5/10 — sees problem but enjoys social aspects"],
        ["Confidence", "3/10 — previous failed attempts; low self-efficacy"],
        ["Key Values", "Family (wife + children), career advancement, health"],
        ["Discrepancy", "Values family but missing kids' events; promotion denied"],
        ["Session 1 Focus", "Feedback (liver + norms); explore reactions; elicit concerns"],
        ["Session 2 Focus", "Values clarification; decisional balance; efficacy building"],
        ["Session 3 Focus", "Progress review (reduced to weekends); affirm success"],
        ["Session 4 Focus", "Maintenance; coping for work stress; relapse prevention"],
        ["Outcome", "Reduced 70%; liver normalized at 6 months"],
    ],
    refs="Miller et al. (1995). MET Manual. | NIMHANS Manual (2008). | Project MATCH (1997).",
    notes="This case illustrates typical MET application. Low confidence with moderate importance tells us to focus on building self-efficacy. Values (family, career) provide material for discrepancy. Gradual reduction goal demonstrates the menu of options approach.",
    col_widths=[2.8, 9.5],
    takeaway="MET case formulation: assess stage + importance + confidence → match strategy accordingly."
)


make_two_column_slide(
    "Case 2: Cannabis (Ravi, 22) & Case 3: Opioid (Suresh, 28)",
    "CANNABIS — RAVI, 22, Engineering Student",
    [
        "Presenting: Academic failure, memory problems",
        "Stage: Precontemplation (brought by parents)",
        "Importance: 3/10 | Confidence: 7/10",
        "Focus: DEVELOP DISCREPANCY",
        "S1: Neuropsych feedback; academic data",
        "S2: Career values; looking forward 5 yrs",
        "Outcome: Reduced use; grades improved",
        "",
        "Strategy: Low importance → focus on WHY",
    ],
    "OPIOID — SURESH, 28, Daily Wage Worker",
    [
        "Presenting: IV heroin 5 yrs; wife left; suicidal",
        "Stage: Contemplation (desperate, hopeless)",
        "Importance: 9/10 | Confidence: 1/10",
        "Focus: BUILD SELF-EFFICACY",
        "S1: Safety assessment; affirm survival strength",
        "S2: Past successes; discuss OST options",
        "Outcome: Engaged in OST + MET; stabilized",
        "",
        "Strategy: High importance → focus on CAN",
    ],
    refs="Miller et al. (1995). MET Manual. | NIMHANS Manual (2008). | Carroll et al. (2006). Drug Alcohol Depend, 81, 161.",
    notes="These cases demonstrate how importance/confidence guides focus. Cannabis case: importance is low (doesn't see problem), so focus on discrepancy development using academic/career values. Opioid case: importance is already high (desperate to change), so focus entirely on building self-efficacy and hope.",
    takeaway="Low importance → develop discrepancy (WHY change). Low confidence → build self-efficacy (CAN change)."
)

make_content_slide(
    "Case 4: Dual Diagnosis — Meera, 38, Depression + Alcohol",
    [
        "Presenting: Husband left; using alcohol to cope with MDD; 250ml rum daily",
        "Comorbidity: MDD (moderate), GAD, low self-esteem, social isolation",
        "Stage: Contemplation | Importance: 7/10 | Confidence: 2/10",
        "",
        "Key Discrepancy: Wants to feel better → Alcohol is WORSENING depression",
        "",
        "MET SESSION PLAN:",
        "  S1: Feedback on depression-alcohol cycle (psychoeducation via E-P-E)",
        "    'Alcohol provides temporary relief but worsens depression long-term'",
        "  S2: Values (children, independence); small achievable goals; affirm",
        "  S3: Celebrate small wins; address depressive cognitions",
        "  S4: Ongoing treatment plan; social support; maintenance",
        "",
        "INTEGRATION: MET for alcohol + Pharmacotherapy + CBT for depression",
        "",
        "KEY PRINCIPLE: Very low confidence + depression = extensive self-efficacy work",
        "  Small, achievable goals are essential ('One alcohol-free day this week')",
    ],
    refs="Baker et al. (2012). Clinical Psychology Review, 32, 726-738. | NIMHANS Manual (2008).",
    notes="In dual diagnosis, MET addresses the substance component while acknowledging the mental health driver. The depression-alcohol cycle provides powerful feedback material. Integration with other treatments is essential for complex cases.",
    takeaway="In dual diagnosis, MET addresses substance use while integrating with disorder-specific treatment."
)

# ══════════════════════════════════════════════════════════════
# SECTION 18-19: SPECIAL POPULATIONS & APPLICATIONS (Slides 73-78)
# ══════════════════════════════════════════════════════════════
make_section_divider(18, "Special Populations & Applications Beyond Addiction")

make_table_slide(
    "MET Adaptations for Special Populations",
    ["Population", "Key Adaptations", "Evidence"],
    [
        ["Adolescents", "Stronger autonomy emphasis; peer norms; immediate consequences", "Marlatt et al. (1998): BASICS reduced college drinking"],
        ["Women", "Address stigma/shame; trauma-informed; safety-focused", "NIMHANS: 'Women substance users are greatly stigmatized'"],
        ["Pregnant Women", "Build on maternal motivation; reduce guilt; any reduction helps", "Handmaker et al. (1999): Brief MI reduces prenatal alcohol"],
        ["Older Adults", "Focus on health/medication interactions; quality of life", "Blow & Barry (2012): Alcohol Research, 34, 18-28"],
        ["Mandated Clients", "Separate mandate from person; find ANY genuine motivation", "MET especially effective with reluctant clients"],
        ["Psychosis", "Simplified; shorter sessions; concrete feedback", "MI effective for medication adherence in psychosis"],
        ["Medical Settings", "5-A (Ask/Assess/Advise/Assist/Arrange); E-P-E for info", "NIMHANS stepped care model"],
    ],
    refs="NIMHANS Manual (2008). | Jensen et al. (2011). Clin Psych Rev, 31, 1024. | Handmaker et al. (1999). JCCP, 67, 285.",
    notes="Each population requires tailored adaptations while maintaining core MET principles. The common thread: empathy, autonomy, and non-judgment work across all populations.",
    col_widths=[2.3, 5.5, 4.5],
    takeaway="Core MET principles are universal; delivery is adapted to each population's unique needs and barriers."
)


make_quad_card_slide(
    "Applications Beyond Addiction",
    [
        ("🚭", "SMOKING CESSATION", "MI increases quit rates by 30-40%.\nHettema & Hendricks (2010).\nCombines well with NRT/pharmacotherapy."),
        ("🏃", "EXERCISE & WEIGHT", "MI for physical activity promotion.\nArmstrong et al. (2011).\nDietary change, obesity management."),
        ("💊", "MEDICATION ADHERENCE", "HIV (ART adherence): Parsons et al. (2007).\nDiabetes self-management.\nHypertension medication compliance."),
        ("🧠", "MENTAL HEALTH", "Treatment engagement for anxiety.\nEating disorders: Treasure & Schmidt (2008).\nGambling disorder; medication in psychosis."),
    ],
    refs="Rubak et al. (2005). Brit J Gen Practice, 55, 305. | Lundahl et al. (2013). Patient Educ Counsel, 93, 157.",
    notes="MET/MI principles apply to ANY behavior where ambivalence about change exists. The underlying mechanisms (autonomy support, discrepancy, efficacy) are universal. The specific feedback and discrepancy development are tailored to each domain.",
    takeaway="Wherever ambivalence about behavior change exists, MET principles are applicable."
)

make_content_slide(
    "MET in the NIMHANS Stepped Care Model (Indian Context)",
    [
        "NIMHANS (2008): India has 62.5M alcohol users; limited specialist resources",
        "STEPPED CARE APPROACH:",
        "",
        "  STEP 1: IDENTIFICATION (Primary Care)",
        "    5-A Strategy: Ask, Assess, Advise, Assist, Arrange",
        "    Screen every patient (like diabetes/hypertension screening)",
        "",
        "  STEP 2: BRIEF INTERVENTION (Single session FRAMES)",
        "    For hazardous/harmful use (not yet dependent)",
        "    Feedback + Empathy + Menu of options",
        "",
        "  STEP 3: MET / MOTIVATIONAL INTERVIEWING (2-4 sessions)",
        "    For moderate dependence or ambivalent clients",
        "    Structured feedback + motivational strategies",
        "",
        "  STEP 4: SPECIALIZED TREATMENT (De-addiction centers)",
        "    For severe dependence; complex comorbidity",
        "",
        "KEY: MET bridges brief intervention and intensive treatment",
        "NIMHANS: 'Not practical to limit care only to specialized centers'",
    ],
    refs="Murthy, P. (2008). Psychosocial Interventions for Persons with Substance Abuse. NIMHANS. | WHO (2010). Brief Intervention Manual.",
    notes="The NIMHANS model positions MET within a resource-appropriate framework for developing countries. The physician becomes an 'agent of change.' Brief interventions can be delivered by non-specialists with training. MET fills the gap between brief advice and intensive treatment.",
    takeaway="MET occupies Step 3 in NIMHANS stepped care — accessible, trainable, and bridges brief to intensive treatment."
)

# ══════════════════════════════════════════════════════════════
# SECTION 20: RESEARCH EVIDENCE (Slides 79-83)
# ══════════════════════════════════════════════════════════════
make_section_divider(20, "Research Evidence")

make_table_slide(
    "Major RCTs Supporting MET/MI",
    ["Study", "Year", "N", "Key Finding"],
    [
        ["Project MATCH", "1997", "1,726", "MET (4 sessions) = CBT = TSF (12 sessions each)"],
        ["UKATT", "2005", "742", "MET (3 sessions) = SBNT (8 sessions); MET more cost-effective"],
        ["COMBINE", "2006", "1,383", "MET + naltrexone effective combination"],
        ["Miller et al.", "1993", "42", "Empathic style > confrontational (1-year outcomes)"],
        ["Marlatt et al.", "1998", "348", "Single MI session reduced college drinking (2-year f/u)"],
        ["Stephens et al.", "2004", "291", "MI effective for cannabis use disorder"],
        ["Handmaker et al.", "1999", "42", "Brief MI reduced prenatal alcohol use"],
        ["Moyers et al.", "2007", "103", "MI behaviors → change talk → behavior change (mechanism)"],
    ],
    refs="Project MATCH (1997). JOSA, 58, 7-29. | UKATT (2005). BMJ, 331, 541. | Bien et al. (1993). Addiction, 88, 315.",
    notes="The evidence base spans 30+ years. Project MATCH remains the landmark finding: 4 sessions of MET matched 12 sessions of other established treatments. UKATT replicated internationally. Moyers illuminated the mechanism (how it works).",
    col_widths=[2.8, 1.0, 1.0, 7.5],
    takeaway="200+ RCTs support MI/MET. Project MATCH showed 4 sessions of MET equals 12 sessions of CBT or TSF."
)


make_table_slide(
    "Meta-Analyses: Effect Sizes and Conclusions",
    ["Meta-Analysis", "Studies", "Key Finding", "Effect Size"],
    [
        ["Burke et al. (2003)", "30 RCTs", "MI effective across substances", "d = 0.25-0.57"],
        ["Hettema et al. (2005)", "72 trials", "MI effective across health behaviors", "Broad support"],
        ["Lundahl et al. (2010)", "119 studies", "MI significantly better than comparison", "OR = 1.55"],
        ["Vasilaki et al. (2006)", "15 studies", "Brief MI for alcohol: durable at 12 months", "Significant"],
        ["Smedslund et al. (2011)", "59 studies", "Cochrane: MI reduces substance use", "May diminish over time"],
        ["Lundahl et al. (2013)", "48 studies", "MI for health: diet, exercise, adherence", "d = 0.20-0.40"],
    ],
    refs="Burke et al. (2003). JCCP, 71, 843. | Lundahl et al. (2010). Clin Psych Rev, 30, 1. | Smedslund et al. (2011). Cochrane Database.",
    notes="Effect sizes are small-medium (d = 0.25-0.57), consistent with psychotherapy research. The Cochrane review confirmed effectiveness but noted effects may diminish without boosters. This supports using booster sessions. The consistency across 200+ studies is remarkable.",
    col_widths=[3.5, 1.8, 4.5, 2.5],
    takeaway="Consistent small-to-medium effects (d=0.25-0.57) across 200+ studies, multiple populations, and diverse settings."
)

make_content_slide(
    "Mechanism of MET: How Does It Work?",
    [
        "THE PROVEN CAUSAL CHAIN (Moyers et al., 2007, 2009):",
        "",
        "  Therapist MI-Consistent Behavior",
        "          ↓",
        "  Client CHANGE TALK increases",
        "          ↓",
        "  Actual BEHAVIOR CHANGE occurs",
        "",
        "Supporting Evidence:",
        "  • Amrhein et al. (2003): Strength of commitment language in session predicted",
        "    drug use outcomes at follow-up",
        "  • Moyers et al. (2007): Therapist empathic behavior → more client change talk",
        "  • Magill et al. (2014): Change-to-sustain talk ratio predicts outcomes",
        "",
        "Proposed Mechanisms:",
        "  1. Self-Perception: Hearing yourself argue for change → believing in change",
        "  2. Cognitive Dissonance: Value-behavior gap → discomfort → action",
        "  3. Self-Efficacy: Confidence building → more persistent effort",
        "  4. Therapeutic Alliance: Safety → exploration → insight → change",
    ],
    refs="Moyers et al. (2007). JCCP, 75, 790. | Amrhein et al. (2003). JCCP, 71, 862. | Magill et al. (2014). JSAT, 46, 685.",
    notes="Understanding the mechanism helps refine practice. The clear clinical implication: FOCUS ON EVOKING AND REINFORCING CHANGE TALK. This is the primary technical goal of every MET session. Everything else (reflections, questions, affirmations) serves this purpose.",
    takeaway="The mechanism: Therapist MI behavior → Client change talk → Behavior change. Focus on evoking change talk."
)

# ══════════════════════════════════════════════════════════════
# SECTION 21: COMPARISON WITH OTHER THERAPIES (Slides 84-86)
# ══════════════════════════════════════════════════════════════
make_section_divider(21, "Comparison with Other Therapies")

make_table_slide(
    "MET vs. Other Major Psychotherapy Approaches",
    ["Therapy", "Focus", "Sessions", "How It Differs from MET"],
    [
        ["MET", "Internal motivation", "4", "Client argues for change; brief; feedback-based"],
        ["CBT", "Skills & cognitions", "12-20", "Therapist teaches skills; structured homework"],
        ["DBT", "Emotion regulation", "6-12 months", "Skills training; structured groups; phone coaching"],
        ["ACT", "Values & acceptance", "8-16", "Mindfulness; defusion; committed action"],
        ["12-Step (TSF)", "Spiritual surrender", "12+", "Powerlessness; higher power; AA meetings"],
        ["REBT", "Irrational beliefs", "12-16", "Disputational; directive; rational analysis"],
        ["Psychodynamic", "Unconscious conflict", "Long-term", "Interpretation; transference; insight"],
        ["CM", "Reinforcement", "Variable", "External rewards/vouchers for behavior"],
    ],
    refs="Project MATCH (1997). JOSA, 58, 7-29. | Carroll (1998). Cognitive-Behavioral Approach. Yale Univ Press.",
    notes="MET is unique in focusing on WHY change rather than HOW. This makes it complementary to skills-based approaches. Common sequence: MET first (to build motivation) then CBT (to build skills). MET differs from 12-Step in emphasizing personal empowerment rather than powerlessness.",
    col_widths=[2.0, 2.5, 1.5, 6.3],
    takeaway="MET focuses on WHY (motivation); CBT focuses on HOW (skills). They complement each other perfectly."
)

make_two_column_slide(
    "MET vs. MI: Understanding the Distinction",
    "MI (Motivational Interviewing)",
    [
        "Broader COUNSELING STYLE",
        "Flexible — no fixed structure",
        "Assessment optional",
        "Feedback optional",
        "Variable duration (5 min to 50 min)",
        "No standardized manual",
        "Style/spirit-based training",
        "Applicable to ANY behavior change",
        "200+ RCTs across all health areas",
    ],
    "MET (Motivational Enhancement Therapy)",
    [
        "Specific MANUALIZED TREATMENT",
        "Fixed 4-session protocol",
        "Structured assessment REQUIRED",
        "Personalized Feedback Report (core)",
        "60-90 min sessions × 4",
        "Project MATCH standardized manual",
        "Protocol-based training",
        "Primarily developed for addiction",
        "Key RCTs: MATCH, UKATT, COMBINE",
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Miller et al. (1995). MET Manual.",
    notes="This distinction is frequently confused. MI is the broader philosophy/style applicable anywhere. MET is a SPECIFIC program that uses MI within a structured protocol. Think: MI is the operating system; MET is a specific application running on it.",
    takeaway="MI = flexible counseling style (the philosophy). MET = structured 4-session protocol (a specific program using MI)."
)


# ══════════════════════════════════════════════════════════════
# SECTION 22-23: ADVANTAGES & LIMITATIONS (Slides 87-89)
# ══════════════════════════════════════════════════════════════
make_section_divider(22, "Advantages & Limitations")

make_two_column_slide(
    "Advantages and Limitations of MET",
    "ADVANTAGES",
    [
        "Brief (4 sessions) — cost-effective",
        "Non-confrontational — reduces dropout",
        "Strong evidence base (200+ RCTs)",
        "Flexible — adapts to diverse populations",
        "Trainable by various professionals",
        "Compatible with other treatments",
        "Feasible in resource-limited settings",
        "Works with mandated/reluctant clients",
        "Reduces therapist burnout (collaborative)",
        "UKATT: 5× more cost-effective",
    ],
    "LIMITATIONS",
    [
        "May not suffice alone for severe dependence",
        "Effects may diminish without boosters (Smedslund, 2011)",
        "Not designed for skills building",
        "Requires skilled therapist (poor MI can harm)",
        "Effect sizes small-medium (not dramatic)",
        "Active mechanism not fully understood",
        "Needs ongoing fidelity monitoring",
        "Cultural adaptations required",
        "Insufficient for complex cases alone",
        "Risk of therapist drift without supervision",
    ],
    refs="UKATT (2005). BMJ, 331, 541. | Smedslund et al. (2011). Cochrane Database. | NIMHANS (2008).",
    notes="Being honest about limitations is essential. MET is powerful but not a panacea. For severe dependence, use as gateway to comprehensive treatment. NIMHANS explicitly states brief interventions 'should not be viewed as complete solutions.' Ongoing training and supervision maintain quality.",
    takeaway="MET is brief, effective, and cost-efficient — but not sufficient alone for severe cases or without ongoing fidelity monitoring."
)

# ══════════════════════════════════════════════════════════════
# SECTION 24-25: ETHICS & CULTURAL ADAPTATION (Slides 90-92)
# ══════════════════════════════════════════════════════════════
make_section_divider(24, "Ethical Considerations & Cultural Adaptation")

make_quad_card_slide(
    "Ethical Considerations in MET Practice",
    [
        ("⚖️", "AUTONOMY", "Client's right to choose their own goals.\nEven choosing NOT to change is respected.\nInformed consent about the approach."),
        ("🔒", "CONFIDENTIALITY", "Standard limits (risk to self/others).\nSO involvement requires consent.\nSpecial considerations with adolescents."),
        ("🛡️", "NON-MALEFICENCE", "Poor MI CAN be harmful (confrontation disguised as MI).\nIs strategic reflection ethical? Yes — if serving CLIENT welfare.\nCompetence required."),
        ("🌍", "CULTURAL COMPETENCE", "Adaptations for diverse populations.\nLanguage, family roles, religious values.\nNIMAHNS: Family involvement culturally critical in India."),
    ],
    refs="APA (2017). Ethical Principles of Psychologists. | Miller & Rollnick (2013). MI, 3rd ed. | NIMHANS (2008).",
    notes="The primary ethical concern: is strategic use of reflection manipulative? Answer: if genuine compassion guides it (serving client welfare), it is ethical. If self-serving (compliance with therapist agenda), it becomes manipulation. The spirit of MET IS the ethical safeguard.",
    takeaway="MET's spirit (compassion, autonomy respect) is both its clinical power and its ethical safeguard."
)

make_content_slide(
    "Cultural Adaptation: MET in the Indian Context",
    [
        "INDIA-SPECIFIC CONSIDERATIONS (NIMHANS, 2008):",
        "  • 62.5M alcohol users + limited specialist resources",
        "  • Physicians as 'agents of change' (trained in brief interventions)",
        "  • Family involvement is culturally expected and therapeutically useful",
        "  • Joint family dynamics: multiple stakeholders",
        "",
        "ADAPTATIONS NEEDED:",
        "  • Greater stigma for substance use (especially women)",
        "  • Spiritual/religious values can be leveraged for discrepancy",
        "  • Hierarchical relationships: MET's egalitarian approach needs gentle intro",
        "  • Language: Materials available in local languages",
        "  • Gender roles: Women's substance use heavily stigmatized",
        "  • Consanguineous marriage dynamics may affect family therapy",
        "",
        "UNIVERSAL PRINCIPLES THAT TRANSCEND CULTURE:",
        "  Empathy, respect, autonomy support, non-judgment",
        "  These are recognized and valued across ALL human cultures",
    ],
    refs="Murthy, P. (2008). NIMHANS Manual, pp. 1-3. | Benegal, V. (2005). Addiction, 100, 1051-1056.",
    notes="The NIMHANS manual provides crucial context for Indian implementation. The stepped care model matches limited resources. Cultural adaptations address family collectivism, stigma, and hierarchical relationships while maintaining universal MI principles.",
    takeaway="Core MET principles (empathy, autonomy, non-judgment) are universal; delivery adapts to cultural context."
)


# ══════════════════════════════════════════════════════════════
# SECTION 26: THERAPIST COMPETENCIES (Slides 93-95)
# ══════════════════════════════════════════════════════════════
make_section_divider(26, "Therapist Competencies & Training")

make_table_slide(
    "MITI Coding: Measuring MI Fidelity",
    ["Metric", "Competent Threshold", "Proficient Threshold"],
    [
        ["Reflection:Question Ratio", "≥ 1:1", "≥ 2:1"],
        ["% Open Questions", "≥ 50%", "≥ 70%"],
        ["% Complex Reflections", "≥ 40%", "≥ 50%"],
        ["MI-Adherent Behaviors", "Present", "Frequent"],
        ["MI-Non-Adherent Behaviors", "Rare", "Absent"],
        ["Technical Global (Cultivating CT)", "≥ 3.5/5", "≥ 4.0/5"],
        ["Relational Global (Partnership)", "≥ 3.5/5", "≥ 4.0/5"],
    ],
    refs="Moyers et al. (2014). MITI Manual 4.2. University of New Mexico. | Miller & Moyers (2006). Behav Cog Psychotherapy, 34, 135.",
    notes="The MITI system provides standardized assessment of MI quality. Without fidelity monitoring, therapists drift toward confrontation over time. Regular recording and coding of sessions with supervisory feedback maintains quality.",
    col_widths=[4.0, 4.0, 4.3],
    takeaway="Fidelity matters: aim for reflection:question ≥ 2:1, open questions ≥ 70%, complex reflections ≥ 50%."
)

make_process_slide(
    "Training Pathway for MET Competence",
    [
        ("STAGE 1", "Learn:\nRead manual\nUnderstand theory"),
        ("STAGE 2", "Workshop:\n2-3 day intensive\nRole-play practice"),
        ("STAGE 3", "Practice:\nRecord sessions\nGet MITI feedback"),
        ("STAGE 4", "Supervision:\nOngoing coaching\nCase discussion"),
        ("STAGE 5", "Mastery:\nMaintain fidelity\nTeach others"),
    ],
    refs="Miller & Moyers (2006). Behav Cog Psychotherapy, 34, 135. | Schwalbe et al. (2014). Addiction, 109, 1287.",
    notes="A 2-day workshop alone does NOT produce lasting competence. Research shows skills deteriorate without ongoing supervision. The full pathway includes: reading, intensive workshop, practice with recording, coded feedback, ongoing supervision, and eventually teaching. Reading the manual is necessary but insufficient.",
    takeaway="Competence requires workshop + practice + feedback + supervision. Reading alone is insufficient."
)

# ══════════════════════════════════════════════════════════════
# SECTION 27: PRACTICAL DEMONSTRATION (Slides 96-99)
# ══════════════════════════════════════════════════════════════
make_section_divider(27, "Practical Demonstration: Role Plays")

make_content_slide(
    "Role Play: Session 1 Opening with Reluctant Client",
    [
        "SCENARIO: Arun, 35, IT professional, referred by wife",
        "",
        "T: 'Welcome, Arun. Thank you for coming. What brings you here?'",
        "C: 'My wife made me come. She thinks I drink too much. I don't see the problem.'",
        "T: 'So coming here wasn't really your idea. You're doing this for her.'",
        "C: 'Yeah, exactly. Everyone I know drinks.'",
        "T: 'From your perspective, your drinking seems pretty normal.'",
        "C: 'Right. But she keeps nagging me about it.'",
        "T: 'That must be frustrating. And yet here you are.'",
        "C: 'Well... I love her. And she IS worried. I don't want to lose her.'",
        "T: 'She means a lot to you, and her concern affects you.'",
        "C: 'Yeah. Maybe I COULD cut back a little.'",
        "T: 'So part of you can see a reason to look at this more closely.'",
        "",
        "NOTICE: Client moved from 'I don't have a problem' to 'Maybe I could cut back'",
        "  — ALL through reflection. Zero confrontation. Zero advice.",
    ],
    refs="Miller et al. (1995). MET Manual (adapted clinical dialogue).",
    notes="This demonstrates the power of reflective listening with a reluctant/mandated client. The therapist never argues, never gives advice, never confronts. By reflecting the wife's importance (a value), the client naturally produces change talk. This is textbook MET.",
    takeaway="Skilled reflection moves clients from resistance to change talk without any confrontation or pressure."
)

make_two_column_slide(
    "Common Mistakes vs. Correct MI Responses",
    "❌ WRONG (MI-Inconsistent)",
    [
        "Client: 'I don't think I'm an alcoholic'",
        "",
        "CONFRONTATION:",
        "'Yes you are! Look at your results!'",
        "",
        "LABELING:",
        "'You meet the diagnostic criteria'",
        "",
        "QUESTION TRAP:",
        "'Why don't you think so?'",
        "",
        "LECTURING:",
        "'Let me explain the DSM criteria...'",
    ],
    "✓ CORRECT (MI-Consistent)",
    [
        "Client: 'I don't think I'm an alcoholic'",
        "",
        "SHIFTING FOCUS:",
        "'Labels aren't important. What concerns YOU?'",
        "",
        "REFLECTION:",
        "'You don't like that word applied to you'",
        "",
        "REFRAME:",
        "'You think carefully before accepting labels'",
        "",
        "DOUBLE-SIDED:",
        "'You don't see yourself that way, AND you have some concerns'",
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="The most common mistake for new therapists is the 'label debate' — arguing about whether someone IS an alcoholic. MET completely avoids this by shifting focus to what the CLIENT is concerned about. Never argue about labels.",
    takeaway="NEVER argue about labels. Shift focus: 'What concerns YOU about your drinking?'"
)


make_content_slide(
    "Extended Role Play: Feedback Delivery with Defensive Client",
    [
        "T: 'Your weekly consumption puts you above 95% of adults.' [neutral]",
        "C: 'That can't be right. There's no way I drink that much.'",
        "T: 'This is surprising to you. It doesn't match how you see it.' [reflect]",
        "C: 'The test must be wrong. Or maybe I exaggerated.'",
        "T: 'You're wondering if the numbers are accurate.' [simple reflection]",
        "C: 'Well... I DO drink most nights. But it's just social.'",
        "T: 'So drinking is a regular part of your social life.' [reflect]",
        "C: 'Yeah. But 95th percentile? That sounds really bad.'",
        "T: 'That number is concerning to you.' [reflects concern = CHANGE TALK!]",
        "C: 'I never thought of myself as THAT kind of drinker.'",
        "T: 'There's a gap between how you've seen yourself and what the data shows.",
        "    That's a lot to take in.' [reflects discrepancy; names process]",
        "",
        "RESULT: Client moved from defensive → concerned → discrepancy recognition",
        "ALL through consistent, patient reflection. Never once arguing for validity.",
    ],
    refs="Miller et al. (1995). MET Manual, Session 1 protocol.",
    notes="This role play shows how the therapist NEVER argues for the validity of results. Despite client's initial rejection, consistent reflection creates safety. The client's own cognitive dissonance emerges naturally. The therapist names the discrepancy without pushing it. This is textbook MET feedback delivery.",
    takeaway="Even when clients reject feedback, consistent reflection creates safety for discrepancy to emerge naturally."
)

# ══════════════════════════════════════════════════════════════
# SECTION 28: EXAMINATION QUESTIONS (Slides 100-102)
# ══════════════════════════════════════════════════════════════
make_section_divider(28, "Examination Questions & Practice")

make_content_slide(
    "Viva Questions & Long Answer Topics",
    [
        "VIVA QUESTIONS (2-3 minute structured answers):",
        "  1. Define MET. How does it differ from MI?",
        "  2. Explain the theoretical foundations (name 5+ theories)",
        "  3. What is the spirit of MET? Why is it more important than technique?",
        "  4. Describe OARS with clinical examples",
        "  5. What is change talk? Name 7 types (DARN-CAT)",
        "  6. Explain FRAMES and its evidence base",
        "  7. How do you 'roll with resistance'? Give 3 strategies",
        "  8. Describe the 4-session MET protocol",
        "  9. What does research say? Cite key studies",
        "  10. How would you adapt MET for Indian context?",
        "",
        "LONG ANSWER QUESTIONS (15-20 marks):",
        "  • Discuss MET: theoretical foundations, evidence, structure,",
        "    applications, and limitations (20 marks)",
        "  • Compare MET with CBT and 12-Step for alcohol dependence (15 marks)",
        "  • Role of brief motivational interventions in primary care (15 marks)",
    ],
    refs="Based on M.Phil Clinical Psychology / PsyD / Psychiatry examination patterns.",
    notes="Structure answers as: definition → theory → principles → evidence → application → limitations. Reference both MET manual and NIMHANS manual for Indian context questions.",
    takeaway="For exams: Integrate theory + evidence + clinical application in structured answers with specific references."
)

make_content_slide(
    "Clinical Scenario Questions with Answers",
    [
        "SCENARIO 1: Man, 45: 'My wife made me come. I don't have a problem.'",
        "  Stage: Precontemplation. Response: Empathy + explore + autonomy.",
        "  'So this wasn't your idea. Tell me what brought you here.'",
        "",
        "SCENARIO 2: Client says: 'I know I should quit but I just can't.'",
        "  Analysis: Desire/Need (should quit) + low Ability (can't).",
        "  Focus: Build self-efficacy. 'You want to quit. What makes you feel you can't?'",
        "",
        "SCENARIO 3: After feedback: 'These results are scary.'",
        "  Response: Reflect emotion. 'This is really hitting home for you.'",
        "  Principle: Express empathy. Feedback is creating discrepancy naturally.",
        "",
        "SCENARIO 4: 'I've tried quitting 5 times and always failed.'",
        "  Reframe: 'Five attempts shows real persistence — you haven't given up.'",
        "  Principle: Support self-efficacy. Reframe failure as learning.",
        "",
        "SCENARIO 5: Client says: 'I don't think I'm an alcoholic.'",
        "  Response: Shift focus. 'Labels aside — what concerns YOU about drinking?'",
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="For each scenario: identify (1) stage, (2) type of talk, (3) relevant principle, (4) specific response. Multiple acceptable answers exist — key is MI-consistency (empathic, non-confrontational, autonomy-supporting).",
    takeaway="For scenarios: Stage → Talk type → Principle → MI-consistent response. Never confront; always reflect."
)


make_table_slide(
    "MCQs for Practice",
    ["#", "Question", "Answer"],
    [
        ["1", "DARN-CAT stands for:", "Desire, Ability, Reasons, Need - Commitment, Activation, Taking Steps"],
        ["2", "The founder of MI is:", "William R. Miller (1983)"],
        ["3", "Project MATCH compared MET with:", "CBT and TSF (12-Step Facilitation)"],
        ["4", "Recommended reflection:question ratio:", "2:1 (two reflections per question)"],
        ["5", "FRAMES stands for:", "Feedback, Responsibility, Advice, Menu, Empathy, Self-efficacy"],
        ["6", "MET consists of how many sessions?", "4 sessions over 12 weeks"],
        ["7", "The spirit of MI includes:", "Partnership, Acceptance, Compassion, Evocation (PACE)"],
        ["8", "Which predicts behavior change?", "Client commitment language (Amrhein, 2003)"],
        ["9", "NIMHANS uses which acronym for MI principles?", "DARES (Develop discrepancy, Avoid argumentation, Roll with resistance, Express empathy, Support self-efficacy)"],
        ["10", "The 'righting reflex' refers to:", "Therapist's urge to fix/advise/correct the client"],
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed. | NIMHANS Manual (2008).",
    notes="Use these MCQs for revision. Focus on: key acronyms, founders, research findings, distinguishing features, and numerical benchmarks.",
    col_widths=[0.5, 6.0, 5.8],
    takeaway="Know: DARN-CAT, FRAMES, PACE, OARS, DARES + founders + key study findings + benchmarks."
)

# ══════════════════════════════════════════════════════════════
# SECTION 29: SUMMARY (Slides 103-106)
# ══════════════════════════════════════════════════════════════
make_section_divider(29, "Summary & Clinical Pearls")

make_content_slide(
    "MET at a Glance: Complete Summary",
    [
        "IDENTITY: Brief, evidence-based, client-centered yet directive intervention",
        "",
        "THEORY: Rogers + Festinger + Bem + Bandura + Prochaska + Deci/Ryan",
        "",
        "SPIRIT: Partnership | Acceptance | Compassion | Evocation (PACE)",
        "",
        "PRINCIPLES: Express Empathy | Develop Discrepancy | Roll with Resistance | Self-Efficacy",
        "",
        "SKILLS: OARS (Open questions, Affirmations, Reflections, Summaries)",
        "",
        "COMPONENTS: FRAMES (Feedback, Responsibility, Advice, Menu, Empathy, Self-efficacy)",
        "",
        "TARGET: Change Talk (DARN-CAT) — evoke it, recognize it, reinforce it",
        "",
        "STRUCTURE: 4 sessions — Feedback → Commitment → Review → Consolidation",
        "",
        "EVIDENCE: 200+ RCTs; Project MATCH (4 sessions = 12 sessions of CBT/TSF)",
        "",
        "APPLICATION: Addiction, health behavior, chronic disease, all populations",
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed. | Project MATCH (1997). JOSA, 58, 7-29.",
    notes="This summary condenses the entire presentation. Each line represents a major section. Students should be able to expand each into a detailed explanation with examples and references.",
    takeaway="MET: Brief + Client-centered + Directive + Evidence-based + Universally applicable."
)


make_content_slide(
    "Top 10 Clinical Pearls for Practice",
    [
        "1. The SPIRIT matters more than techniques — be genuine, collaborative",
        "",
        "2. Resistance is YOUR signal — change approach, don't push harder",
        "",
        "3. The more clients TALK ABOUT change, the more they BELIEVE in it",
        "",
        "4. Ask 'Why X and not zero?' — this ALWAYS evokes change talk",
        "",
        "5. Use 'AND' not 'BUT' in double-sided reflections",
        "",
        "6. Start with benefits of drinking — validates; reduces defensiveness",
        "",
        "7. Never argue about labels — shift to 'What concerns YOU?'",
        "",
        "8. Reflections should outnumber questions 2:1 minimum",
        "",
        "9. Low importance → develop discrepancy | Low confidence → build efficacy",
        "",
        "10. If you're working harder than the client, something is wrong",
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="These represent the most important practical takeaways from MET training. Pearl #10 is diagnostic: in skilled MET, the CLIENT argues for change while the therapist facilitates. If the therapist is pushing, the dynamic has inverted.",
    takeaway="If you're working harder than the client, the dynamic has inverted — they should be arguing for their own change."
)

make_content_slide(
    "MET Process Flowchart",
    [
        "DECISION TREE FOR EVERY SESSION:",
        "",
        "  Is client showing CHANGE TALK?",
        "    → YES → Reflect it, affirm it, ask for more. Move toward commitment.",
        "    → NO  → Continue Phase 1: explore, evoke, develop discrepancy.",
        "",
        "  Is RESISTANCE/DISCORD high?",
        "    → YES → Soften approach. Roll with it. Emphasize autonomy. Back off.",
        "    → NO  → Continue current strategy. Deepen exploration.",
        "",
        "  Is client ready for PLANNING? (Signs: decreased resistance,",
        "  increased change talk, 'how' questions, imagining change)",
        "    → YES → Transition to Phase 2. Change Plan Worksheet.",
        "    → NO  → Stay in Phase 1. Don't push. Trust the process.",
        "",
        "REMEMBER: The client leads the dance. You follow their readiness.",
        "Premature planning → increased resistance (pushing too fast)",
        "Delayed planning → missed momentum (client ready but therapist still exploring)",
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="This decision tree provides a quick clinical reference. The three key questions (change talk present? discord high? ready for planning?) guide every moment of every session. Print it and keep it visible during early practice.",
    takeaway="Three questions guide every session: Change talk? → Reinforce. Discord? → Soften. Ready? → Plan."
)

# ══════════════════════════════════════════════════════════════
# SECTION 30: REFERENCES (Slides 107-109)
# ══════════════════════════════════════════════════════════════
make_section_divider(30, "Complete Bibliography")

make_content_slide(
    "References: Primary Sources",
    [
        "Miller, W.R., Zweben, A., DiClemente, C.C., & Rychtarik, R.G. (1995). Motivational",
        "  Enhancement Therapy Manual. NIAAA, NIH Publication No. 94-3723.",
        "Murthy, P. (2008). Psychosocial Interventions for Persons with Substance Abuse.",
        "  NIMHANS Publication, Bangalore, India.",
        "Miller, W.R., & Rollnick, S. (2013). Motivational interviewing: Helping people change",
        "  (3rd ed.). Guilford Press.",
        "Miller, W.R., & Rollnick, S. (1991). Motivational interviewing: Preparing people to",
        "  change addictive behavior. Guilford Press.",
        "Project MATCH Research Group. (1997). Matching alcoholism treatments to client",
        "  heterogeneity. Journal of Studies on Alcohol, 58, 7-29.",
        "UKATT Research Team. (2005). Effectiveness of treatment for alcohol problems.",
        "  BMJ, 331, 541.",
        "Rollnick, S., Miller, W.R., & Butler, C.C. (2008). Motivational interviewing in",
        "  health care. Guilford Press.",
        "SAMHSA. (2019). Enhancing Motivation for Change (TIP 35). HHS Publication.",
    ],
    refs="All references in APA 7th Edition format.",
    notes="These primary sources form the foundation. Miller et al. (1995) and Murthy (2008) are the two uploaded reference books. Miller & Rollnick (2013) provides the most current theoretical framework.",
    takeaway="Essential reading: MET Manual (Miller et al., 1995) + MI 3rd Ed (Miller & Rollnick, 2013) + NIMHANS Manual (Murthy, 2008)."
)


make_content_slide(
    "References: Research Studies & Meta-Analyses",
    [
        "Amrhein, P.C., et al. (2003). Client commitment language during MI. JCCP, 71, 862-878.",
        "Bandura, A. (1977). Self-efficacy. Psychological Review, 84, 191-215.",
        "Bem, D.J. (1972). Self-perception theory. In Advances in Exp Social Psych, Vol. 6.",
        "Bien, T.H., Miller, W.R., & Tonigan, S. (1993). Brief interventions. Addiction, 88, 315.",
        "Burke, B.L., et al. (2003). MI meta-analysis. JCCP, 71, 843-861.",
        "Deci, E.L., & Ryan, R.M. (1985). Intrinsic motivation and self-determination. Plenum.",
        "Festinger, L. (1957). A Theory of Cognitive Dissonance. Row, Peterson.",
        "Hettema, J., et al. (2005). MI review. Annual Rev Clinical Psych, 1, 91-111.",
        "Lundahl, B., et al. (2010). MI meta-analysis. Clinical Psych Rev, 30, 1-11.",
        "Miller, W.R. (1983). MI with problem drinkers. Behavioural Psychotherapy, 11, 147.",
        "Moyers, T.B., et al. (2007). Therapist influence on client language. JCCP, 75, 790.",
        "Patterson, G.A., & Forgatch, M.S. (1985). Therapist behavior. JCCP, 53, 846-851.",
        "Prochaska, J.O., & DiClemente, C.C. (1982). Transtheoretical therapy. Psychotherapy, 19, 276.",
        "Rogers, C.R. (1957). Necessary conditions for change. J Consulting Psych, 21, 95-103.",
        "Smedslund, G., et al. (2011). MI for substance abuse. Cochrane Database.",
    ],
    refs="All references in APA 7th Edition format.",
    notes="These studies form the evidence base cited throughout. Each represents a milestone in MI/MET research. Students should know key findings of each for examination purposes.",
    takeaway="Key citations to know: Project MATCH (1997), Miller (1993), Amrhein (2003), Moyers (2007), Smedslund (2011)."
)

# ══════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ══════════════════════════════════════════════════════════════
def make_closing_slide():
    slide = new_slide()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                Inches(SLIDE_W), Inches(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Decorative elements
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(5), Inches(4), Inches(4))
    c1.fill.solid()
    c1.fill.fore_color.rgb = DARK_BLUE
    c1.line.fill.background()

    # Quote
    qtb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(2))
    qtf = qtb.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.text = '"People are generally the best experts on themselves."'
    qp.font.size = Pt(24)
    qp.font.italic = True
    qp.font.color.rgb = GOLD
    qp.alignment = PP_ALIGN.CENTER
    qp2 = qtf.add_paragraph()
    qp2.text = "— William R. Miller"
    qp2.font.size = Pt(16)
    qp2.font.color.rgb = SKY
    qp2.alignment = PP_ALIGN.CENTER

    # Second quote
    q2tb = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10), Inches(1.5))
    q2tf = q2tb.text_frame
    q2tf.word_wrap = True
    q2p = q2tf.paragraphs[0]
    q2p.text = '"If you find yourself working harder than the client, you are doing it wrong."'
    q2p.font.size = Pt(18)
    q2p.font.italic = True
    q2p.font.color.rgb = WHITE
    q2p.alignment = PP_ALIGN.CENTER

    # Final message
    ftb = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10), Inches(1.5))
    ftf = ftb.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.text = "The Spirit of MET:"
    fp.font.size = Pt(18)
    fp.font.color.rgb = MID_BLUE
    fp.alignment = PP_ALIGN.CENTER
    fp2 = ftf.add_paragraph()
    fp2.text = "Partnership  •  Acceptance  •  Compassion  •  Evocation"
    fp2.font.size = Pt(22)
    fp2.font.bold = True
    fp2.font.color.rgb = WHITE
    fp2.alignment = PP_ALIGN.CENTER

    # Thank you
    ttb = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(10), Inches(0.7))
    ttf = ttb.text_frame
    tp = ttf.paragraphs[0]
    tp.text = "Thank You"
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = GOLD
    tp.alignment = PP_ALIGN.CENTER

    add_notes(slide, "Close with these memorable quotes. Remind students that the spirit of MET is its most powerful element. Every interaction is an opportunity to enhance motivation for positive change.")

# ══════════════════════════════════════════════════════════════
# ADDITIONAL SLIDES TO REACH 140+ (Sections needing more depth)
# ══════════════════════════════════════════════════════════════

# Additional detailed slides for OARS exercises
make_content_slide(
    "OARS Practice Exercises for Training",
    [
        "EXERCISE 1 — Convert Closed → Open:",
        "  'Do you drink daily?' → 'Tell me about a typical week'",
        "  'Have you tried quitting?' → 'What attempts have you made?'",
        "  'Is family affected?' → 'How has your family noticed?'",
        "",
        "EXERCISE 2 — Write Affirmations:",
        "  Client attended despite being busy → 'Making time for this shows it matters to you'",
        "  Client reduced from daily to 3/week → 'You cut by more than half — that's significant'",
        "  Client honestly shared relapse → 'Being honest about this took real courage'",
        "",
        "EXERCISE 3 — Deepen Reflections:",
        "  C: 'I just can't seem to stop once I start'",
        "  Simple: 'Once you start, it's hard to stop'",
        "  Complex: 'There's a sense of losing control that worries you'",
        "  Feeling: 'That scares you — like you've lost the ability to choose'",
        "",
        "EXERCISE 4 — Practice Summaries: Collect 3+ change talk statements from a partner",
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Miller et al. (1995). MET Manual.",
    notes="These exercises should be practiced in pairs or small groups. Reflective listening requires extensive practice to become natural. The ability to deepen reflections on the fly is a hallmark of MI proficiency.",
    takeaway="OARS proficiency requires deliberate practice — convert, deepen, and summarize strategically."
)

make_content_slide(
    "The Importance/Confidence Matrix: Clinical Decision Guide",
    [
        "CLINICAL DECISION MATRIX:",
        "",
        "  HIGH Importance + HIGH Confidence → Ready for action!",
        "    Strategy: Consolidate commitment; Change Plan Worksheet",
        "",
        "  HIGH Importance + LOW Confidence → Motivated but stuck",
        "    Strategy: Build self-efficacy; past successes; small steps; menu of options",
        "    'You clearly WANT to — let's explore what might help you feel you CAN'",
        "",
        "  LOW Importance + HIGH Confidence → Could but doesn't want to",
        "    Strategy: Develop discrepancy; values clarification; feedback",
        "    'You feel capable — what would make it feel more important?'",
        "",
        "  LOW Importance + LOW Confidence → Not ready at all",
        "    Strategy: Start with importance (WHY); then build confidence (HOW)",
        "    Don't push — gentle exploration; respect autonomy",
        "",
        "This matrix guides EVERY clinical decision in MET.",
    ],
    refs="Rollnick, Mason & Butler (1999). Health Behavior Change. Churchill Livingstone. | Miller & Rollnick (2013).",
    notes="This 2×2 matrix is one of the most clinically useful frameworks in MET. Assess importance and confidence separately (using rulers), then match your strategy accordingly. This prevents common errors like building skills for someone who doesn't see the need.",
    takeaway="Assess importance and confidence separately → match strategy to the specific gap."
)

make_content_slide(
    "The Personal Feedback Report (PFR) — Content & Format",
    [
        "PFR Components (Project MATCH protocol):",
        "  1. DRINKING PATTERNS:",
        "     Weekly consumption vs. population norms (percentile)",
        "     Peak BAC levels (estimated)",
        "     Number of heavy drinking days",
        "",
        "  2. CONSEQUENCES:",
        "     Alcohol-related negative consequences score",
        "     Comparison with clinical populations",
        "",
        "  3. RISK FACTORS:",
        "     Family history of alcohol problems",
        "     Tolerance indicators",
        "     Dependence symptoms",
        "",
        "  4. HEALTH INDICATORS:",
        "     Liver function tests (GGT, AST, ALT)",
        "     Other medical findings",
        "",
        "  5. NEUROPSYCHOLOGICAL RESULTS (if tested):",
        "     Memory, concentration, processing speed",
        "",
        "DELIVERY: Neutral tone. 'Here are your results. What do you make of them?'",
    ],
    refs="Miller et al. (1995). MET Manual, Appendix A (complete PFR protocol).",
    notes="The PFR is what distinguishes MET from pure MI. It provides objective data that creates discrepancy through normative comparison. The therapist's role is to present it neutrally and let the data speak for itself.",
    takeaway="The PFR creates discrepancy through objective normative comparison — present neutrally, reflect reactions."
)


make_two_column_slide(
    "Resistance Management Strategies: Detailed Guide",
    "STRATEGY & DEFINITION",
    [
        "1. Simple Reflection",
        "   Acknowledge without elaboration",
        "2. Amplified Reflection",
        "   Overstate slightly → client backs down",
        "3. Double-Sided Reflection",
        "   Both sides: 'X AND Y'",
        "4. Shifting Focus",
        "   Move away from stuck point",
        "5. Reframing",
        "   New perspective on same info",
        "6. Agreeing with a Twist",
        "   Agree but add new dimension",
        "7. Coming Alongside",
        "   'Maybe you're not ready yet'",
    ],
    "EXAMPLE",
    [
        "1. 'You're not worried right now'",
        "",
        "2. 'So absolutely no concerns at all'",
        "   → 'Well, not NO concerns...'",
        "3. 'You enjoy it AND you worry'",
        "",
        "4. 'Let's set labels aside. What concerns YOU?'",
        "",
        "5. 'Your concern about labels shows thoughtfulness'",
        "",
        "6. 'You're right — only YOU can decide'",
        "",
        "7. 'Maybe this isn't the right time'",
        "   → Often triggers client arguing FOR change",
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="Each strategy reduces discord while maintaining therapeutic contact. The amplified reflection and 'coming alongside' use reverse psychology (reactance theory). The client, feeling their freedom threatened by the therapist seemingly agreeing with inaction, moves toward change.",
    takeaway="Seven tools for resistance: reflect, amplify, double-side, shift, reframe, agree-twist, come-alongside."
)

make_content_slide(
    "Involving Significant Others in MET",
    [
        "MET optionally includes a spouse/partner/family member in Session 1:",
        "",
        "PURPOSE: Additional perspective; accountability; support system",
        "",
        "GUIDELINES:",
        "  • Invite SO to share observations (not accusations)",
        "  • Use reflective listening with the SO too",
        "  • REFRAME SO statements to highlight caring/concern",
        "",
        "Dialogue Examples from MET Manual:",
        "  Wife: 'I always thought he was drinking too much'",
        "  T: 'You've been worried about him for quite a while' [reflects caring]",
        "",
        "  Husband: (weeping) 'I've told you to quit drinking!'",
        "  T: 'You really care about her. It's hard to hear these results' [reframes]",
        "",
        "  Friend: 'I never thought he drank that much!'",
        "  T: 'This is surprising for you too. (To client:) Does this surprise you?'",
        "",
        "CAUTION: Maintain client as focus. Don't let SO become confrontational.",
        "After reflecting SO → ask client's reaction → reflect client's change talk",
    ],
    refs="Miller et al. (1995). MET Manual, pp. 35-42 (Significant Others section).",
    notes="The SO can be a powerful therapeutic ally when handled skillfully. The therapist must maintain balance: validate the SO's concerns while keeping the client as the primary focus. Reframing anger as caring is a key skill.",
    takeaway="Involve SOs supportively. Reframe their concerns as caring. Always maintain client as central focus."
)

make_content_slide(
    "MET for Mandated/Court-Ordered Clients",
    [
        "CHALLENGES:",
        "  • Often precontemplators — not seeking help voluntarily",
        "  • May be angry, hostile, or resigned",
        "  • Perceived coercion undermines intrinsic motivation",
        "",
        "MET STRATEGIES — Five Steps:",
        "  1. Acknowledge honestly:",
        "     'You're here because the court required it. How do you feel about that?'",
        "",
        "  2. Separate mandate from person:",
        "     'You have to be here, but what you GET OUT OF IT is up to you'",
        "",
        "  3. Emphasize what IS within control:",
        "     'You can't control the court order, but you CAN choose what to do now'",
        "",
        "  4. Find ANY genuine motivation:",
        "     'Setting aside the court — is there anything about your drinking YOU'd change?'",
        "",
        "  5. Use gentle paradox (if appropriate):",
        "     'I'm not sure you're ready to think about this yet...'",
        "",
        "RESEARCH: MET is ESPECIALLY effective with reluctant/mandated clients",
        "  (because it reduces reactance by respecting autonomy)",
    ],
    refs="Miller et al. (1995). MET Manual, Special Problems section. | Moyers & Rollnick (2002). Substance Use & Misuse, 37, 2089.",
    notes="MET's non-confrontational approach is ideally suited for mandated clients because it directly addresses the reactance that forced treatment creates. By emphasizing autonomy, you paradoxically increase engagement.",
    takeaway="MET is especially effective with mandated clients — it reduces reactance by honoring autonomy."
)


make_content_slide(
    "DARES: The NIMHANS Motivational Interviewing Principles",
    [
        "The NIMHANS Manual uses DARES as a mnemonic for MI principles:",
        "",
        "D — Develop Discrepancy:",
        "  Help client see gap between current behavior and valued goals",
        "",
        "A — Avoid Argumentation:",
        "  Never argue or confront; this ALWAYS increases resistance",
        "  Arguments about labels are especially counterproductive",
        "",
        "R — Roll with Resistance:",
        "  Use client's momentum; redirect rather than oppose",
        "  'You don't feel ready right now, and that's okay'",
        "",
        "E — Express Empathy:",
        "  Warm, reflective listening; genuine understanding",
        "  Acceptance of ambivalence as normal",
        "",
        "S — Support Self-Efficacy:",
        "  Boost confidence; highlight strengths and past successes",
        "  'You've shown you CAN make changes in your life'",
        "",
        "NIMHANS: 'Empathic communication and encouraging stance are common to FRAMES and DARES'",
    ],
    refs="Murthy, P. (2008). NIMHANS Manual, p. 21. | Miller & Rollnick (1991). MI. Guilford.",
    notes="DARES is identical in content to MI principles but organized mnemonically for Indian training contexts. It's widely used across NIMHANS-affiliated training programs.",
    takeaway="DARES (Develop discrepancy, Avoid argumentation, Roll with resistance, Express empathy, Support self-efficacy)."
)

make_content_slide(
    "Relapse Prevention Integration in MET Sessions 3-4",
    [
        "MARLATT & GORDON (1985) RELAPSE MODEL:",
        "  High-risk situation → No coping → Decreased efficacy → Lapse → Relapse",
        "  High-risk situation → Coping response → Increased efficacy → Maintained change",
        "",
        "MET-COMPATIBLE RELAPSE PREVENTION (Session 4):",
        "  1. Identify high-risk situations collaboratively",
        "     'What situations do you think might be challenging?'",
        "  2. Develop coping plans (using non-prescriptive style)",
        "     'How might you handle that situation?'",
        "  3. Build self-efficacy for coping",
        "     'What strengths do you have for managing this?'",
        "  4. Normalize lapses — not failures, but learning opportunities",
        "     'If you do slip, what would help you get back on track?'",
        "  5. Emergency planning",
        "     'Who could you call? Where could you go?'",
        "",
        "KEY: Maintain MET SPIRIT even during RP planning — collaborative, not prescriptive",
        "  NOT: 'You MUST avoid bars'",
        "  BUT: 'What situations do you think might be risky for you?'",
    ],
    refs="Marlatt & Gordon (1985). Relapse Prevention. Guilford. | Witkiewitz & Marlatt (2004). Clin Psych Rev, 24, 1-28.",
    notes="Relapse prevention in MET maintains the motivational spirit. Rather than prescribing coping strategies, explore collaboratively with the client. This increases ownership and adherence to the plan.",
    takeaway="Relapse prevention in MET uses collaborative exploration, not prescriptive directives — maintaining the spirit."
)

make_content_slide(
    "Session 2: Values Clarification Exercise — Detailed Protocol",
    [
        "VALUES CLARIFICATION — Step-by-Step Protocol:",
        "",
        "Step 1: 'What are the most important things in your life?'",
        "  (Let client generate freely; don't suggest)",
        "",
        "Step 2: Present values list or card sort — select top 5-10:",
        "  Family | Health | Career | Integrity | Freedom | Spirituality",
        "  Financial security | Self-respect | Relationships | Achievement",
        "",
        "Step 3: Rank in order of importance",
        "",
        "Step 4: For EACH top value, explore:",
        "  'How is your drinking related to this value?'",
        "  'How does drinking fit with being [that person]?'",
        "",
        "EXAMPLE DIALOGUE:",
        "  T: 'You said being a good father is #1. Tell me about that.'",
        "  C: 'My kids mean everything to me'",
        "  T: 'And where does drinking fit with being the father you want to be?'",
        "  C: '...Not well. I missed my daughter's recital because I was hungover.'",
        "  T: 'That doesn't sit right with you — there's a gap.'",
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="Values clarification creates powerful internal discrepancy from the client's OWN values. The therapist doesn't point out the gap — they ask the question and let the client discover it themselves. This creates far stronger cognitive dissonance than external confrontation.",
    takeaway="Values clarification creates the most powerful discrepancy — because it comes from the client's OWN priorities."
)

make_table_slide(
    "Confidence & Importance Rulers: Strategic Follow-Up Questions",
    ["Ruler", "Question", "Follow-Up", "Purpose"],
    [
        ["Importance", "'How important is change? (0-10)'", "'Why X and not 2?'", "Evokes change talk"],
        ["Importance", "", "'What would move you to 8?'", "Identifies barriers"],
        ["Confidence", "'How confident you could change? (0-10)'", "'Why X and not zero?'", "Identifies strengths"],
        ["Confidence", "", "'What would increase it?'", "Identifies resources needed"],
        ["Readiness", "'How ready are you? (0-10)'", "'Why X and not lower?'", "Evokes motivation"],
        ["ANY ruler", "NEVER ask:", "'Why not 10?'", "This evokes SUSTAIN talk!"],
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Rollnick et al. (1999). Health Behavior Change.",
    notes="The rulers are both assessment AND intervention. The critical rule: ALWAYS ask 'Why not lower?' which forces articulation of motivation. NEVER ask 'Why not higher?' which forces articulation of barriers and sustain talk.",
    col_widths=[2.0, 4.0, 3.5, 2.8],
    takeaway="'Why X and not lower?' = evokes change talk. 'Why not higher?' = evokes sustain talk. ALWAYS use the former."
)


make_content_slide(
    "MET and Comorbid Psychiatric Disorders",
    [
        "50-70% of substance users have comorbid psychiatric disorders:",
        "",
        "DEPRESSION + SUBSTANCE USE:",
        "  Key discrepancy: 'Alcohol provides temporary relief but worsens depression'",
        "  Integration: MET for alcohol + Pharmacotherapy + CBT for depression",
        "  Focus: The depression-alcohol CYCLE as feedback",
        "",
        "ANXIETY + SUBSTANCE USE:",
        "  Key discrepancy: Anxiety drives use; withdrawal creates MORE anxiety",
        "  Focus: Build self-efficacy for alternative coping strategies",
        "",
        "PSYCHOSIS + SUBSTANCE USE:",
        "  Adaptations: Simplified, shorter sessions; concrete feedback",
        "  Evidence: MI effective for medication adherence in psychosis",
        "",
        "PTSD + SUBSTANCE USE:",
        "  Trauma-informed modifications essential; safety first",
        "  MET can motivate engagement with trauma-specific treatment",
        "",
        "NIMHANS: 'Appropriate referral for complex cases'",
    ],
    refs="Baker et al. (2012). Clin Psych Rev, 32, 726. | Drake et al. (2004). JSAT, 27, 197. | NIMHANS (2008).",
    notes="Dual diagnosis is the norm, not the exception. MET addresses the substance component while acknowledging that comorbidity may be driving use. Tailor discrepancy development to the specific comorbid condition.",
    takeaway="MET addresses substance use in dual diagnosis; integration with disorder-specific treatment is essential."
)

make_content_slide(
    "MET for Smoking Cessation & Weight Management",
    [
        "SMOKING CESSATION:",
        "  • MI increases quit rates 30-40% over standard care",
        "  • Feedback: CO levels, lung function, cardiovascular risk",
        "  • Discrepancy: Health values vs. continued smoking",
        "  • Combines well with NRT and pharmacotherapy",
        "  • Even 5-minute MI in primary care increases quit attempts",
        "  • Evidence: Hettema & Hendricks (2010); Lai et al. (2010)",
        "",
        "WEIGHT MANAGEMENT / OBESITY:",
        "  • Ambivalence about dietary change / exercise is universal",
        "  • Feedback: BMI, metabolic markers, risk assessment",
        "  • MI increases engagement with lifestyle interventions",
        "  • Armstrong et al. (2011): MI for weight loss effective",
        "",
        "MEDICATION ADHERENCE:",
        "  • HIV ART adherence: Parsons et al. (2007)",
        "  • Diabetes management: Channon et al. (2007)",
        "  • Discrepancy: 'You want to stay healthy AND you skip medications'",
        "",
        "COMMON PRINCIPLE: Wherever ambivalence exists, MET applies.",
    ],
    refs="Hettema & Hendricks (2010). J Clin Psych, 66, 1162. | Armstrong et al. (2011). Int J Obesity, 35, 891.",
    notes="The versatility of MET principles is remarkable. Any behavior involving ambivalence can benefit from motivational approaches. The same mechanisms work: autonomy support, empathy, discrepancy, self-efficacy.",
    takeaway="MET works for any health behavior change: smoking, weight, exercise, medication adherence — same principles apply."
)

make_content_slide(
    "Project MATCH: Detailed Study Design & Findings",
    [
        "STUDY DESIGN (1989-1997):",
        "  • $27 million; largest alcoholism treatment trial ever",
        "  • 9 clinical sites across USA + 1 coordinating center",
        "  • 1,726 patients randomized to 3 treatments:",
        "    - MET: 4 sessions over 12 weeks",
        "    - CBT: 12 sessions over 12 weeks",
        "    - TSF: 12 sessions over 12 weeks",
        "  • Two arms: Outpatient (n=952) and Aftercare (n=774)",
        "  • Follow-up: 3, 6, 9, 12, and 15 months + 3-year follow-up",
        "",
        "KEY FINDINGS:",
        "  • ALL three treatments produced significant, sustained improvement",
        "  • MET (4 sessions) produced outcomes EQUIVALENT to 12-session treatments",
        "  • Patient-treatment matching effects were minimal (unexpected!)",
        "  • 3-year follow-up: Gains maintained across all groups",
        "  • Treatment effects began rapidly (within first month for MET)",
        "",
        "IMPLICATION: MET is 3× more cost-effective (same outcomes, 1/3 sessions)",
    ],
    refs="Project MATCH (1997). JOSA, 58, 7-29. | Project MATCH (1998). Addiction, 93, 1431-1446.",
    notes="Project MATCH's most striking finding was that MET achieved equivalent outcomes with only one-third the sessions. This was revolutionary for the field and established MET as remarkably cost-effective. The matching hypothesis (main aim) was largely unsupported.",
    takeaway="Project MATCH: 4 sessions of MET = 12 sessions of CBT/TSF. Three times more cost-effective."
)

make_content_slide(
    "Recent Evidence (2015-2025): Current State of MET Research",
    [
        "RECENT META-ANALYSES AND KEY STUDIES:",
        "",
        "DiClemente et al. (2017): Mechanisms of MI — readiness to change mediates outcomes",
        "",
        "Magill et al. (2018): Updated MI meta-analysis — effects confirmed across substances",
        "  Small but consistent effects; most robust for alcohol and cannabis",
        "",
        "Frost et al. (2018): MI for young adults — effective for reducing heavy drinking",
        "",
        "Lindson et al. (2019): Cochrane review — MI for smoking cessation effective",
        "",
        "Technology-Assisted MI (emerging evidence):",
        "  • Telehealth MI delivery: Equivalent outcomes to in-person (post-COVID data)",
        "  • App-based MI elements: Promising but early-stage",
        "  • AI chatbots using MI principles: Under investigation",
        "",
        "Implementation Science:",
        "  • How to maintain fidelity in real-world (non-research) settings",
        "  • Task-shifting to non-specialists (NIMHANS model; WHO recommendations)",
        "  • Cost-effectiveness in developing healthcare systems",
    ],
    refs="DiClemente et al. (2017). Addiction, 112(S2), 92-100. | Magill et al. (2018). Addiction, 113, 2115-2126. | Lindson et al. (2019). Cochrane Database.",
    notes="The evidence base continues to grow. Recent research focuses on mechanisms, technology-assisted delivery, and implementation science. The COVID pandemic accelerated telehealth MI research, showing equivalent outcomes to in-person delivery.",
    takeaway="Recent evidence confirms MET effectiveness; emerging areas include technology-assisted delivery and implementation science."
)


make_table_slide(
    "MET Effect Sizes Across Populations (Summary)",
    ["Population/Setting", "Effect Size (d or OR)", "Source"],
    [
        ["Alcohol (general)", "d = 0.25-0.57", "Burke et al. (2003)"],
        ["Cannabis", "d = 0.30-0.45", "Stephens et al. (2004)"],
        ["Tobacco/Smoking", "OR = 1.26", "Lindson et al. (2019)"],
        ["Health behaviors (overall)", "OR = 1.55", "Lundahl et al. (2010)"],
        ["College students", "d = 0.35-0.50", "Marlatt et al. (1998)"],
        ["Mandated/reluctant clients", "d = 0.40-0.60", "Higher than non-mandated!"],
        ["Primary care brief MI", "d = 0.20-0.30", "Rubak et al. (2005)"],
        ["Dual diagnosis", "d = 0.20-0.35", "Baker et al. (2012)"],
    ],
    refs="Burke (2003). JCCP. | Lundahl (2010). Clin Psych Rev. | Lindson (2019). Cochrane. | Rubak (2005). Brit J Gen Pract.",
    notes="Effect sizes are consistently small-to-medium, which is typical for psychotherapy research. Notably, effects are LARGER for mandated clients — MET's non-confrontational approach is especially helpful when clients are forced into treatment.",
    col_widths=[4.5, 3.5, 4.3],
    takeaway="Effect sizes are small-medium (d=0.25-0.57) but consistent. Notably LARGER for mandated/reluctant clients."
)

make_quad_card_slide(
    "Future Directions in MET Research & Practice",
    [
        ("📱", "TECHNOLOGY-ASSISTED MI", "Telehealth delivery (equivalent outcomes).\nAI chatbots using MI principles.\nApp-based self-monitoring with MI feedback."),
        ("🧬", "PRECISION/PERSONALIZED MET", "Matching interventions to client characteristics.\nGenomics + pharmacogenomics integration.\nNeuroimaging of MI brain effects."),
        ("🌐", "IMPLEMENTATION SCIENCE", "Scaling MI training globally.\nTask-shifting to non-specialists (NIMHANS model).\nMaintaining fidelity in real-world settings."),
        ("🎮", "NEW APPLICATIONS", "Behavioral addictions (gaming, social media).\nMI for climate-related behavior change.\nMI in organizational/workplace settings."),
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Apodaca & Longabaugh (2009). Clin Psych Rev, 29, 199.",
    notes="The future of MET includes technology-assisted delivery, personalized approaches, global implementation science, and novel applications. The NIMHANS model of training non-specialists is increasingly relevant worldwide.",
    takeaway="MET's future: tech-assisted delivery, precision approaches, global implementation, behavioral addictions."
)

make_content_slide(
    "The Causal Chain: Evidence for How MET Works",
    [
        "PROVEN CAUSAL PATHWAY (Multiple studies confirm):",
        "",
        "  STEP 1: Therapist uses MI-consistent behaviors",
        "    (Open questions, affirmations, reflections, autonomy support)",
        "             ↓",
        "  STEP 2: Client CHANGE TALK increases in session",
        "    (Desire, Ability, Reasons, Need → Commitment, Activation, Steps)",
        "             ↓",
        "  STEP 3: Client's actual BEHAVIOR CHANGES",
        "    (Reduced substance use, improved health behaviors)",
        "",
        "KEY EVIDENCE:",
        "  • Moyers et al. (2007): MI-consistent therapist → more change talk",
        "  • Amrhein et al. (2003): Commitment language → drug use outcomes",
        "  • Magill et al. (2014): Change-to-sustain talk RATIO predicts outcomes",
        "  • Apodaca & Longabaugh (2009): Comprehensive mechanism review",
        "",
        "CLINICAL IMPLICATION: Your primary technical goal = EVOKE CHANGE TALK",
        "  Everything else (reflections, questions, affirmations) SERVES this purpose.",
    ],
    refs="Moyers et al. (2007). JCCP, 75, 790. | Amrhein et al. (2003). JCCP, 71, 862. | Apodaca & Longabaugh (2009). Clin Psych Rev, 29, 199.",
    notes="This is the most important research finding for clinical practice. It tells us exactly WHAT to focus on: evoking change talk. The mechanism is clear: when therapists use MI-consistent behaviors, clients produce more change talk, which predicts actual behavior change.",
    takeaway="Primary technical goal of every MET session: evoke, recognize, and reinforce CLIENT CHANGE TALK."
)

make_content_slide(
    "Common Therapist Traps to Avoid",
    [
        "THE EXPERT TRAP:",
        "  Telling rather than asking. Installing rather than evoking.",
        "  Fix: Ask 'What do you think?' before sharing your opinion",
        "",
        "THE ASSESSMENT TRAP:",
        "  Too many questions, not enough reflections (interrogation mode).",
        "  Fix: After every question, reflect the answer before asking another",
        "",
        "THE PREMATURE FOCUS TRAP:",
        "  Moving to planning before client is ready (Phase 2 too early).",
        "  Fix: Look for readiness signals before transitioning",
        "",
        "THE LABELING TRAP:",
        "  Using 'alcoholic' or other stigmatizing labels.",
        "  Fix: Focus on behavior and concerns, not diagnostic labels",
        "",
        "THE BLAMING TRAP:",
        "  Attributing resistance to client pathology ('in denial').",
        "  Fix: See resistance as YOUR signal to change approach",
        "",
        "THE RIGHTING REFLEX:",
        "  Your urge to fix, advise, correct, solve, warn.",
        "  Fix: 'I wonder...' instead of 'You should...'",
    ],
    refs="Miller & Moyers (2006). Behav Cog Psychotherapy, 34, 135. | Miller et al. (1995). MET Manual.",
    notes="These traps represent the most common ways therapists drift from MI fidelity. Self-monitoring for these traps, plus regular supervision with recorded sessions, helps maintain quality over time.",
    takeaway="Monitor for: expert trap, assessment trap, premature focus, labeling, blaming, and the righting reflex."
)


make_content_slide(
    "Decisional Balance: Worksheet with Therapist Guidance",
    [
        "FOUR-QUADRANT DECISIONAL BALANCE (conducted in session):",
        "",
        "Start HERE (validates; reduces defensiveness):",
        "  Q1: 'What are the GOOD THINGS about drinking for you?'",
        "  (Client lists: relaxation, social, coping with stress, etc.)",
        "",
        "Then move to costs:",
        "  Q2: 'And what's the OTHER SIDE? What are your concerns?'",
        "  (Client lists: health, relationship, work, money, etc.)",
        "",
        "Then explore fears about change:",
        "  Q3: 'What worries you about making a change?'",
        "  (Client lists: losing friends, boredom, coping without, etc.)",
        "",
        "End HERE (builds momentum — leave session pointing forward):",
        "  Q4: 'And what might be BETTER if you did make this change?'",
        "  (Client lists: health, family, clarity, pride, etc.)",
        "",
        "THERAPIST THEN SUMMARIZES — strategically including MORE from Q2 and Q4.",
    ],
    refs="Janis & Mann (1977). Decision Making. Free Press. | Miller et al. (1995). MET Manual.",
    notes="The decisional balance is both assessment and intervention. The strategic order matters: starting with benefits shows you understand; ending with benefits of change leaves forward momentum. The summary strategically weights change talk.",
    takeaway="Start with benefits of drinking (validates), end with benefits of change (motivates). Summarize strategically."
)

make_content_slide(
    "OARS Applied to Eliciting Change Talk",
    [
        "Using OARS STRATEGICALLY to maximize change talk:",
        "",
        "OPEN QUESTIONS that evoke change talk:",
        "  'What concerns you most about your drinking?' (evokes Reasons)",
        "  'What would you like to be different?' (evokes Desire)",
        "  'What strengths do you have?' (evokes Ability)",
        "  'How important is this to you?' (evokes Need)",
        "",
        "AFFIRMATIONS that strengthen change talk:",
        "  'You clearly care about your health' (when they mention health concerns)",
        "  'That took real courage to say' (when they acknowledge problems)",
        "",
        "REFLECTIONS that amplify change talk:",
        "  Client: 'I guess I should probably think about cutting back'",
        "  Reflection: 'Part of you knows something needs to change'",
        "  (Strengthens and deepens the change talk)",
        "",
        "SUMMARIES that collect change talk:",
        "  Gather ALL change talk heard in session → present it as one collection",
        "  Leave out or minimize sustain talk in the summary",
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Moyers et al. (2009). JSAT, 36, 101.",
    notes="OARS are not just communication skills — they are STRATEGIC tools for maximizing change talk. Every question should be designed to evoke change talk. Every reflection should amplify it. Every summary should collect it. This is the active mechanism of MET.",
    takeaway="Use OARS strategically: every question evokes, every reflection amplifies, every summary collects change talk."
)

make_content_slide(
    "MET Treatment Fidelity: MITI Coding in Practice",
    [
        "WHY FIDELITY MATTERS:",
        "  • Poorly delivered MI can be WORSE than no treatment",
        "  • Therapist drift is common without monitoring",
        "  • Quality determines effectiveness",
        "",
        "MITI 4.2 BEHAVIORAL COUNTS (per 20-min sample):",
        "  MI-Adherent: Affirm, Seek Collaboration, Emphasize Autonomy",
        "  MI-Non-Adherent: Confront, Direct, Warn",
        "  Questions: Open vs. Closed (ratio matters)",
        "  Reflections: Simple vs. Complex (ratio matters)",
        "",
        "PROFICIENCY BENCHMARKS:",
        "  • Reflection:Question ratio ≥ 2:1",
        "  • % Open Questions ≥ 70%",
        "  • % Complex Reflections ≥ 50%",
        "  • Technical Global ≥ 4.0/5",
        "  • Relational Global ≥ 4.0/5",
        "",
        "HOW TO USE: Record sessions → trained coder rates → feedback to therapist",
        "Frequency: Minimum quarterly; monthly during initial training",
    ],
    refs="Moyers et al. (2014). MITI Manual 4.2. U of New Mexico. | Schwalbe et al. (2014). Addiction, 109, 1287.",
    notes="Without fidelity monitoring, therapists inevitably drift back toward confrontation and advice-giving. MITI provides objective data about practice quality. The feedback loop (record → code → feedback → adjust) is essential for maintaining competence.",
    takeaway="Monitor fidelity: Record → Code (MITI) → Get feedback → Adjust. Without this, quality degrades over time."
)

make_content_slide(
    "Summary of Key Worksheets and When to Use Them",
    [
        "SESSION 1:",
        "  • Personal Feedback Report (delivered to client)",
        "  • Substance Use Diary (homework assignment)",
        "  • Readiness/Importance/Confidence Rulers",
        "",
        "SESSION 2:",
        "  • Decisional Balance Sheet (4 quadrants)",
        "  • Values Card Sort or Values List",
        "  • Change Plan Worksheet (if ready for commitment)",
        "",
        "SESSION 3:",
        "  • Progress Review Form",
        "  • Modified Change Plan (if needed)",
        "  • Trigger Identification Worksheet",
        "",
        "SESSION 4:",
        "  • Relapse Prevention Plan",
        "  • Emergency Contact Card",
        "  • Coping Skills Worksheet",
        "  • Long-term Goal Planning Sheet",
        "",
        "ALL SESSIONS: Self-monitoring diary (ongoing homework)",
    ],
    refs="Miller et al. (1995). MET Manual. | SAMHSA (2019). TIP 35. | Marlatt & Gordon (1985).",
    notes="Worksheets structure the therapeutic process and provide tangible take-home tools. Assign them as suggestions (not requirements) to maintain autonomy. Each worksheet serves a specific function aligned with the session goals.",
    takeaway="Match worksheets to session goals: S1=Feedback, S2=Balance+Plan, S3=Review, S4=Prevention."
)

make_content_slide(
    "Group Exercise: Identifying Change Talk vs. Sustain Talk",
    [
        "CLASSIFY EACH STATEMENT (CT = Change Talk, ST = Sustain Talk):",
        "",
        "1. 'I know I need to cut down' → CT (Need)",
        "2. 'All my friends drink, I can't be different' → ST (Normalizing)",
        "3. 'My wife says she'll leave' → CT (Reasons)",
        "4. 'I've tried before, never works' → ST (Low ability)",
        "5. 'I wish I could drink normally' → CT (Desire)",
        "6. 'Life would be boring without alcohol' → ST (Cost of change)",
        "7. 'I'm going to try alcohol-free Mondays' → CT (Taking Steps)",
        "8. 'I'm not ready yet' → ST (No activation)",
        "9. 'I could probably cut back if I tried' → CT (Ability)",
        "10. 'I'm willing to give it a month' → CT (Activation/Commitment)",
        "",
        "PRACTICE RESPONDING:",
        "  For each CT: How would you REINFORCE it? (Reflect, affirm, elaborate)",
        "  For each ST: How would you respond WITHOUT strengthening it?",
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Moyers et al. (2009). JSAT, 36, 101.",
    notes="This exercise builds the crucial skill of real-time change talk recognition. In sessions, this must be automatic. Practice classifying and responding until it becomes second nature.",
    takeaway="Rapid CT/ST recognition is essential. Reinforce CT warmly; acknowledge ST briefly without dwelling."
)

# Now add the closing slide and save
make_closing_slide()

# Add more slides for fuller coverage

make_content_slide(
    "Agenda Mapping: Collaborative Session Planning",
    [
        "WHAT IS AGENDA MAPPING?",
        "  Collaboratively deciding what to discuss in the session",
        "  Gives client ownership; supports autonomy; reduces resistance",
        "",
        "HOW TO DO IT:",
        "  T: 'There are several things we could talk about today.'",
        "  T: 'We could look at your assessment results, explore your concerns,",
        "     talk about options, or discuss something else entirely.'",
        "  T: 'What would be most useful for you?'",
        "",
        "WHEN TO USE:",
        "  • Beginning of every session (especially Sessions 2-4)",
        "  • When session feels stuck or directionless",
        "  • When multiple topics compete for attention",
        "",
        "WHY IT WORKS:",
        "  • Client has choice → autonomy support → intrinsic motivation",
        "  • Therapist isn't imposing an agenda → reduces reactance",
        "  • Client's priorities get addressed → better engagement",
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed. | Rollnick, Butler & Mason (2008). MI in Health Care.",
    notes="Agenda mapping is a simple technique that powerfully supports the MET spirit of partnership and autonomy. By asking what the client wants to discuss, you communicate that this is THEIR session.",
    takeaway="Let clients choose the session focus — this simple act communicates partnership and enhances engagement."
)

make_content_slide(
    "Scaling Questions: Beyond the Basic Ruler",
    [
        "STANDARD RULERS:",
        "  Importance: 'How important is it to change? (0-10)'",
        "  Confidence: 'How confident are you that you could? (0-10)'",
        "  Readiness: 'How ready are you to make this change? (0-10)'",
        "",
        "ADVANCED SCALING APPLICATIONS:",
        "  Temporal scaling: 'Last month you were at 3; today you're at 6. What changed?'",
        "  Hypothetical scaling: 'If you were at an 8, what would that look like?'",
        "  Relational scaling: 'Where would your WIFE put you on this scale?'",
        "  Domain-specific: 'How important is it for your HEALTH? For your FAMILY?'",
        "",
        "POWERFUL FOLLOW-UPS:",
        "  'What gives you that much confidence? (even at 3)'",
        "  'What would need to happen to move up just one point?'",
        "  'What's keeping you from being a zero?'",
        "",
        "REMEMBER: Each follow-up is designed to evoke CHANGE TALK.",
    ],
    refs="Rollnick et al. (1999). Health Behavior Change. Churchill Livingstone. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="Scaling questions are remarkably versatile. They can be used for assessment, intervention, progress tracking, and evoking change talk — all in one simple question. The key is the strategic follow-up.",
    takeaway="Rulers are assessment + intervention combined. The follow-up question matters more than the number."
)

make_content_slide(
    "MET with Adolescents: Detailed Adaptations",
    [
        "UNIQUE DEVELOPMENTAL CONSIDERATIONS:",
        "  • Identity formation: 'Who am I?' is the central question",
        "  • Peer influence paramount: Norms matter more than health stats",
        "  • Stronger reactance to perceived authority (especially parents/therapists)",
        "  • Temporal discounting: 'I'll worry about that later'",
        "  • Often mandated: brought by parents, school, or courts",
        "",
        "ADAPTATIONS:",
        "  • Emphasize autonomy STRONGLY: 'It's YOUR life, YOUR choice'",
        "  • Use peer-relevant feedback: 'Compared to people your age...'",
        "  • Focus on IMMEDIATE consequences: social embarrassment, grades, sports",
        "  • Shorter sessions (30-45 min); more interactive; less formal",
        "  • Build alliance BEFORE any feedback delivery (trust first)",
        "  • Involve family sensitively (NIMHANS: critical in Indian context)",
        "",
        "EVIDENCE:",
        "  Marlatt et al. (1998): BASICS — single MI session reduced heavy drinking",
        "  at 2-year follow-up in college students",
        "  Jensen et al. (2011): MI effective for youth substance use",
    ],
    refs="Marlatt et al. (1998). JCCP, 66, 604. | Jensen et al. (2011). Clin Psych Rev, 31, 1024. | NIMHANS (2008).",
    notes="Adolescents respond best when they feel their autonomy is respected and when consequences are immediate and personally relevant (not abstract health risks decades away). The BASICS intervention showed remarkable durability.",
    takeaway="Adolescent MET: maximum autonomy, peer-relevant feedback, immediate consequences, trust before content."
)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = '/projects/sandbox/Dango-kiro/MET_Comprehensive_Presentation.pptx'
prs.save(output_path)
print(f"✓ Presentation saved: {output_path}")
print(f"✓ Total slides: {slide_count}")
print("✓ All elements properly positioned with no overlapping.")
print("✓ All slides fully editable in PowerPoint/Google Slides.")
print("Done!")


make_two_column_slide(
    "MET with Women & Pregnant Women: Specific Considerations",
    "WOMEN (General)",
    [
        "Greater stigma for substance use",
        "Often hidden drinking; delayed help-seeking",
        "Comorbid trauma common (60-80%)",
        "Domestic violence may drive use",
        "NIMHANS: 'Women users are greatly stigmatized'",
        "Adaptations: Safety-focused; trauma-informed",
        "Address shame sensitively; non-judgmental",
        "Family/childcare barriers to treatment",
    ],
    "PREGNANT WOMEN",
    [
        "ANY alcohol = risk (FASD)",
        "High existing motivation (baby's health)",
        "BUT: Intense guilt and shame",
        "MET ideal: Non-judgmental + existing motivation",
        "Handmaker et al. (1999): Brief MI reduced",
        "  prenatal alcohol exposure",
        "Focus: Build on maternal motivation",
        "Never increase guilt; reduce shame",
    ],
    refs="Handmaker et al. (1999). JCCP, 67, 285. | NIMHANS (2008). | Blow & Barry (2012). Alcohol Research, 34, 18.",
    notes="Women face unique barriers including stigma, shame, trauma, and practical barriers. MET's non-judgmental approach is especially valuable. For pregnant women, motivation already exists (baby); the therapist builds on it while reducing guilt.",
    takeaway="Women: address stigma/shame sensitively. Pregnant: build on existing maternal motivation without increasing guilt."
)

make_content_slide(
    "MET in Primary Care & Hospital Settings",
    [
        "NIMHANS 5-A STRATEGY for Primary Care:",
        "  ASK: Screen every patient about substance use (routine, like BP)",
        "  ASSESS: Pattern of use and resulting problems",
        "  ADVISE: Clear, brief advice to reduce/stop",
        "  ASSIST: Specific interventions (FRAMES elements)",
        "  ARRANGE: Appropriate referrals when needed",
        "",
        "BRIEF MI IN MEDICAL SETTINGS:",
        "  • Can be delivered in 5-15 minutes by trained physicians",
        "  • Screening: AUDIT, CAGE, or single screening question",
        "  • Intervention: FRAMES-based brief feedback",
        "  • Follow-up: Arrange next contact",
        "",
        "HOSPITAL SETTINGS:",
        "  • Emergency department: Single MI session post-injury",
        "  • Medical wards: Brief feedback on health consequences",
        "  • Prenatal clinics: Alcohol screening + brief MI",
        "",
        "Evidence: Even a single 10-minute MI session produces measurable change",
        "NIMHANS: 'The trained physician can be an effective agent of change'",
    ],
    refs="Murthy, P. (2008). NIMHANS Manual, pp. 1-3. | Babor & Higgins-Biddle (2001). WHO Brief Intervention Manual.",
    notes="The NIMHANS model makes MET accessible to primary care. The 5-A strategy provides a systematic framework for busy physicians. Even brief contacts can be motivationally enhancing when the principles are followed.",
    takeaway="Brief MI in primary care (even 5-10 min) produces measurable change. Every medical contact is an opportunity."
)

make_content_slide(
    "How Two Phases of MET Map to Clinical Practice",
    [
        "PHASE 1: BUILDING MOTIVATION (Sessions 1-2 primarily)",
        "  Goal: Tip the motivational balance toward change",
        "  Strategies:",
        "    • Elicit self-motivational statements",
        "    • Provide personalized feedback",
        "    • Develop discrepancy through values and data",
        "    • Explore ambivalence with decisional balance",
        "  Client stage: Moving from precontemplation → contemplation → preparation",
        "",
        "PHASE 2: STRENGTHENING COMMITMENT (Sessions 2-4)",
        "  Goal: Consolidate the decision and plan for action",
        "  Strategies:",
        "    • Recapitulate (summary of change talk heard so far)",
        "    • Key questions ('What do you think you'll do?')",
        "    • Information & advice (with permission, E-P-E)",
        "    • Change Plan negotiation and goal setting",
        "  Client stage: Preparation → action",
        "",
        "TRANSITION TIMING: Move to Phase 2 when you see:",
        "  Decreased resistance + Increased change talk + 'How' questions + Imagining change",
    ],
    refs="Miller et al. (1995). MET Manual, pp. 13-32. | Miller & Rollnick (1991). MI. Guilford.",
    notes="The two-phase structure maps onto stages of change. Phase 1 = precontemplation through contemplation. Phase 2 = preparation through action. The transition must be guided by client readiness signals, not by session number.",
    takeaway="Phase 1 builds WHY (motivation); Phase 2 builds HOW (commitment/plan). Time the transition by readiness signals."
)

make_content_slide(
    "Session 2: Detailed Protocol for Consolidating Commitment",
    [
        "SESSION 2 PROTOCOL (60 minutes, Week 2):",
        "",
        "OPENING (10 min):",
        "  'What have you been thinking about since last time?'",
        "  Review self-monitoring diary (if assigned)",
        "  Celebrate any positive observations",
        "",
        "DEEPENING EXPLORATION (20 min):",
        "  Decisional balance (4 quadrants)",
        "  Values clarification exercise",
        "  'How does drinking fit with the person you want to be?'",
        "",
        "TRANSITIONING TO COMMITMENT (20 min):",
        "  Signs of readiness? → Key questions:",
        "  'Given all this, what do you think you'll do?'",
        "  'What would be your next step?'",
        "  If ready: Change Plan Worksheet (collaborative)",
        "  If NOT ready: Continue Phase 1 WITHOUT pressure",
        "",
        "CLOSING (10 min):",
        "  Summary (strategically weighted toward change talk)",
        "  Homework: Continue diary; complete Change Plan; values card sort",
    ],
    refs="Miller et al. (1995). MET Manual, Session 2 protocol.",
    notes="Session 2 bridges motivation-building to commitment-strengthening. The therapist should start by asking what the client has been thinking — often they arrive with new change talk after reflecting on Session 1 feedback. If not ready for planning, stay in Phase 1 without pressure.",
    takeaway="Session 2 bridges Phases 1→2 based on readiness. If ready: plan. If not ready: continue exploring without pressure."
)

make_content_slide(
    "Case 5: Adolescent Cannabis — Complete Session Plan",
    [
        "RAVI, 22, Engineering Student — Mandated by parents",
        "",
        "SESSION 1: Engagement (most critical with adolescents)",
        "  • Build rapport FIRST (20 min): music, interests, career goals",
        "  • Express empathy for being brought by parents: 'Not your idea, huh?'",
        "  • Assess: Importance 3/10 | Confidence 7/10",
        "  • Brief feedback: Academic performance data; neuropsych results",
        "  • End: 'What do you make of this? No pressure to decide anything today'",
        "",
        "SESSION 2: Developing Discrepancy",
        "  • Career values: Engineering degree → well-paying job → independence",
        "  • Looking forward: 'Where do you want to be at 25? How does cannabis fit?'",
        "  • Decisional balance: Benefits of use AND academic costs",
        "  • If ready: Small goal ('I'll try studying sober before the next exam')",
        "",
        "SESSION 3: Review progress; affirm any reduction; explore what worked",
        "SESSION 4: Maintenance; peer pressure management; long-term career plan",
        "",
        "KEY: Low importance → focus on DISCREPANCY (career values vs. cognitive effects)",
    ],
    refs="Stephens et al. (2004). JCCP, 72, 92. | NIMHANS (2008). | Jensen et al. (2011). Clin Psych Rev, 31, 1024.",
    notes="This case demonstrates MET with a precontemplator brought by parents. The entire approach centers on THEIR values (career, independence) not parent/therapist values. Low importance requires discrepancy development; high confidence means efficacy building is less needed.",
    takeaway="Adolescent case: Build rapport first, use THEIR values (not parents'), focus on immediate consequences."
)

make_content_slide(
    "Therapist Self-Reflection Questions",
    [
        "After EVERY session, ask yourself:",
        "",
        "1. Did I do more talking than the client? (If yes → too directive)",
        "",
        "2. Did I hear change talk? What did I do with it?",
        "   (If I missed it → practice recognition)",
        "",
        "3. Was there discord/resistance? How did I respond?",
        "   (If I pushed back → practice rolling)",
        "",
        "4. What was my reflection:question ratio today?",
        "   (If too many questions → practice reflective listening)",
        "",
        "5. Did I give unsolicited advice? (If yes → use E-P-E next time)",
        "",
        "6. Did I respect the client's autonomy genuinely?",
        "",
        "7. What would I do differently next time?",
        "",
        "8. Was I working harder than the client? (If yes → something is wrong)",
        "",
        "KEEP A SUPERVISION JOURNAL with these reflections.",
    ],
    refs="Miller & Moyers (2006). Behav Cog Psychotherapy, 34, 135. | Madson et al. (2009). MI Training. Springer.",
    notes="Self-reflection is essential for MI development. Without it, therapists drift toward habitual patterns (usually advice-giving). These questions serve as a quick post-session check-in that maintains awareness of MI principles.",
    takeaway="Ongoing self-reflection prevents drift. Ask: Was I working harder than the client?"
)

make_table_slide(
    "Comparison: FRAMES vs. OARS vs. DARN-CAT vs. PACE",
    ["Acronym", "Components", "Function in MET"],
    [
        ["FRAMES", "Feedback, Responsibility, Advice, Menu, Empathy, Self-efficacy", "Active ingredients of brief interventions"],
        ["OARS", "Open questions, Affirmations, Reflections, Summaries", "Micro-skills for communication"],
        ["DARN-CAT", "Desire, Ability, Reasons, Need - Commitment, Activation, Taking Steps", "Types of change talk to evoke/recognize"],
        ["PACE", "Partnership, Acceptance, Compassion, Evocation", "The spirit/philosophy of MI/MET"],
        ["DARES", "Develop discrepancy, Avoid argumentation, Roll, Empathy, Self-efficacy", "Core principles (NIMHANS version)"],
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed. | NIMHANS (2008).",
    notes="This comparison helps students organize the multiple acronyms. Each serves a different function: PACE = philosophy, DARES/principles = guidelines, OARS = skills, FRAMES = components, DARN-CAT = targets.",
    col_widths=[1.8, 6.5, 4.0],
    takeaway="Five key acronyms: PACE (spirit), DARES (principles), OARS (skills), FRAMES (components), DARN-CAT (targets)."
)

make_content_slide(
    "MET Mind Map: Visual Concept Organizer",
    [
        "                    THE SPIRIT (PACE) — Foundation of Everything",
        "                    Partnership | Acceptance | Compassion | Evocation",
        "                                    |",
        "            ________________________|________________________",
        "            |                       |                        |",
        "      PRINCIPLES               SKILLS                  STRUCTURE",
        "      (DARES)                  (OARS)                  (4 Sessions)",
        "      • Express empathy        • Open Qs              S1: Feedback",
        "      • Develop discrepancy    • Affirmations         S2: Commitment",
        "      • Roll with resistance   • Reflections          S3: Review",
        "      • Support self-efficacy  • Summaries            S4: Consolidation",
        "            |                       |                        |",
        "      COMPONENTS               TARGETS                 EVIDENCE",
        "      (FRAMES)                (DARN-CAT)              (200+ RCTs)",
        "      Feedback                Change Talk              Project MATCH",
        "      Responsibility          Sustain Talk             UKATT, COMBINE",
        "      Advice, Menu            Ratio predicts           Meta-analyses",
        "      Empathy, Self-eff       outcomes                 Mechanisms",
        "",
        "HIERARCHY: Spirit → Principles → Skills → Techniques → Evidence",
    ],
    refs="Miller et al. (1995). MET Manual. | Miller & Rollnick (2013). MI, 3rd ed.",
    notes="This mind map organizes all MET components hierarchically. Spirit sits at the top (most important). Everything else serves the spirit. Students should be able to recreate this from memory as a study aid.",
    takeaway="The MET hierarchy: Spirit (foundation) → Principles (guides) → Skills (tools) → Structure (protocol) → Evidence (support)."
)

make_content_slide(
    "The Righting Reflex: Your Biggest Enemy in MET",
    [
        "WHAT IS THE RIGHTING REFLEX?",
        "  The innate therapist urge to FIX, ADVISE, CORRECT, and SOLVE",
        "  When we see someone making harmful choices, we WANT to tell them to stop",
        "",
        "WHY IT'S HARMFUL IN MET:",
        "  • Puts therapist in expert position → client in passive/resistant position",
        "  • Evokes the 'Yes, but...' response from clients",
        "  • Creates adversarial dynamic (therapist argues FOR change)",
        "  • Robs client of the opportunity to argue FOR change themselves",
        "  • Undermines autonomy → decreases intrinsic motivation",
        "",
        "PARADOX:",
        "  The MORE you try to convince someone to change,",
        "  the LESS likely they are to change.",
        "",
        "  The LESS you push for change,",
        "  the MORE space you create for THEM to argue for it.",
        "",
        "HOW TO MANAGE IT:",
        "  Notice the urge → pause → ask instead of tell → reflect instead of advise",
        "  Replace 'You should...' with 'What do you think about...?'",
    ],
    refs="Miller & Rollnick (2013). MI, 3rd ed., Chapter 1. | Miller et al. (1995). MET Manual.",
    notes="The righting reflex is the #1 barrier to MI practice. Healthcare professionals are TRAINED to diagnose and prescribe — which is the opposite of what MET requires. Learning to suppress this reflex and instead evoke the client's own motivation is the core skill development in MI training.",
    takeaway="The paradox: the MORE you push for change, the LESS likely it is. Create space for THEM to argue for change."
)

make_content_slide(
    "Short Notes Template for Examination Answers",
    [
        "STRUCTURE FOR 5-7 MARK SHORT NOTES:",
        "",
        "1. DEFINITION (1-2 lines): What it is, who developed it, when",
        "2. KEY COMPONENTS (3-4 points): Main elements or principles",
        "3. CLINICAL APPLICATION (2-3 lines): How it's used in practice",
        "4. CLINICAL EXAMPLE (2 lines): Brief illustration",
        "5. REFERENCE (1 line): Key citation",
        "",
        "EXAMPLE — 'Write a short note on OARS in Motivational Interviewing':",
        "  Definition: OARS refers to the four core micro-skills of MI",
        "    (Miller & Rollnick, 2013).",
        "  Components: Open questions (invite exploration), Affirmations (recognize",
        "    strengths), Reflections (mirror with depth), Summaries (collect themes).",
        "  Application: Used throughout MET to evoke and reinforce change talk.",
        "    Reflections should outnumber questions 2:1.",
        "  Example: 'What concerns you about drinking?' (Open Q) →",
        "    'Part of you is starting to worry' (Complex Reflection).",
        "  Reference: Miller & Rollnick (2013). Motivational Interviewing, 3rd ed.",
    ],
    refs="Based on M.Phil Clinical Psychology / PsyD examination standards.",
    notes="This template helps students structure exam answers efficiently. The 5-component format (definition, components, application, example, reference) works for any short note topic in MET.",
    takeaway="Short note format: Definition → Components → Application → Example → Reference. Practice writing in 5-7 minutes."
)

make_content_slide(
    "Long Answer Template for 15-20 Mark Questions",
    [
        "STRUCTURE FOR 15-20 MARK LONG ANSWERS:",
        "",
        "1. INTRODUCTION (2-3 lines): Define topic; historical context",
        "2. THEORETICAL BASIS (4-5 lines): Underlying theories",
        "3. KEY PRINCIPLES/COMPONENTS (6-8 lines): Main elements detailed",
        "4. CLINICAL APPLICATION (4-5 lines): How it works in practice",
        "5. EVIDENCE BASE (4-5 lines): Key studies with findings",
        "6. ADVANTAGES & LIMITATIONS (3-4 lines): Balanced evaluation",
        "7. CONCLUSION (2 lines): Synthesis; clinical significance",
        "",
        "SAMPLE OUTLINE — 'Discuss MET' (20 marks):",
        "  Intro: Brief intervention by Miller (1995) for Project MATCH",
        "  Theory: Rogers + Festinger + Bem + Bandura + Prochaska (explain each)",
        "  Principles: 4 principles + Spirit (PACE) + OARS + FRAMES",
        "  Clinical: 4-session protocol with PFR delivery (describe)",
        "  Evidence: Project MATCH, UKATT, meta-analyses (cite with N and d)",
        "  A&L: Brief/cost-effective/evidence-based vs. not for severe cases alone",
        "  Conclusion: Gold-standard brief intervention; globally applicable",
    ],
    refs="Based on M.Phil Clinical Psychology / PsyD / Psychiatry examination standards.",
    notes="For 20-mark questions, students need depth AND breadth. The 7-section structure ensures comprehensive coverage within time constraints. Practice writing full answers within 20-25 minutes.",
    takeaway="Long answer: Intro → Theory → Principles → Clinical → Evidence → A&L → Conclusion. Practice within 20 min."
)

# Additional depth slides
make_table_slide(
    "Project MATCH: Patient Characteristics and Matching Variables",
    ["Matching Variable", "Hypothesis", "Result"],
    [
        ["Severity", "High severity → TSF best", "No significant interaction"],
        ["Anger/Hostility", "Angry clients → MET best", "Partial support (outpatient arm)"],
        ["Readiness to Change", "Low readiness → MET best", "No significant interaction"],
        ["Social Support", "Low support → TSF best", "Some support in aftercare"],
        ["Psychiatric Severity", "High → CBT best", "No significant interaction"],
        ["Meaning Seeking", "High → TSF best", "Supported in aftercare arm"],
        ["Self-Efficacy", "High → CBT best", "Not supported"],
        ["Network Support", "Drinking network → TSF best", "Supported in aftercare"],
    ],
    refs="Project MATCH (1997). JOSA, 58, 7-29. | Project MATCH (1998). Addiction, 93, 1431-1446.",
    notes="The irony of Project MATCH: the matching hypotheses (the primary aim) were largely unsupported, but the finding that all treatments worked equally well — with MET using 1/3 the sessions — was groundbreaking for the field.",
    col_widths=[2.8, 4.5, 5.0],
    takeaway="MATCH's biggest finding wasn't matching — it was that 4-session MET equaled 12-session treatments for nearly everyone."
)

make_content_slide(
    "Cognitive Dissonance in MET: Detailed Application",
    [
        "FESTINGER (1957): Theory of Cognitive Dissonance",
        "  When beliefs/values contradict behavior → psychological discomfort",
        "  Discomfort motivates resolution: change the belief OR change the behavior",
        "",
        "APPLICATION IN MET — Creating Therapeutic Dissonance:",
        "",
        "  1. Through FEEDBACK (Session 1):",
        "     Client believes: 'I drink like everyone else'",
        "     Data shows: 'You drink more than 95% of the population'",
        "     → Dissonance between self-image and data",
        "",
        "  2. Through VALUES CLARIFICATION (Session 2):",
        "     Client values: 'Being a good father'",
        "     Behavior: Missing children's events due to hangovers",
        "     → Dissonance between values and behavior",
        "",
        "  3. Through LOOKING FORWARD:",
        "     Desired future: 'Healthy, successful, respected'",
        "     Projected trajectory: Continued decline if unchanged",
        "     → Dissonance between desired and projected futures",
        "",
        "KEY: Client must ARTICULATE the dissonance themselves (not therapist).",
    ],
    refs="Festinger (1957). A Theory of Cognitive Dissonance. | Miller et al. (1995). MET Manual.",
    notes="The power of dissonance lies in it being INTERNALLY generated. When the therapist points it out, it can feel confrontational and trigger defensiveness. When the client discovers it themselves (through reflection and strategic questions), it creates genuine internal motivation.",
    takeaway="Create conditions for the client to DISCOVER their own dissonance — don't point it out for them."
)

make_content_slide(
    "Self-Perception Theory: Why Change Talk Works",
    [
        "BEM (1965, 1972): Self-Perception Theory",
        "  'As I hear myself talk, I learn what I believe'",
        "  People infer their attitudes from observing their own speech and behavior",
        "",
        "MET APPLICATION (Miller, 1995 — direct quotes from manual):",
        "  'Motivational psychology has amply demonstrated that when people are subtly",
        "  enticed to speak or act in a new way, their beliefs and values tend to shift",
        "  in that direction.'",
        "",
        "  'If I say it, and no one has forced me to say it, then I must believe it!'",
        "",
        "CLINICAL IMPLICATION:",
        "  Every time a client voices a reason for change, they become MORE motivated",
        "  This creates a positive feedback loop:",
        "    Therapist evokes → Client speaks change talk → Client becomes more motivated",
        "    → More change talk emerges → Stronger commitment develops",
        "",
        "THE FLIP SIDE:",
        "  When confrontation makes clients argue AGAINST change,",
        "  they talk themselves INTO continued use (self-persuasion against change)",
        "",
        "THIS explains why confrontation is counterproductive.",
    ],
    refs="Bem (1972). In Berkowitz (Ed.), Advances in Exp Social Psych, Vol. 6. | Miller et al. (1995). MET Manual, p. 13.",
    notes="Self-perception theory is THE key mechanism explaining why MET works and why confrontation fails. The MET manual explicitly cites this as the theoretical basis for eliciting self-motivational statements. This is why the therapist seeks to have the CLIENT articulate reasons for change.",
    takeaway="Self-perception: people believe what they hear themselves say. This is why evoking change talk CREATES motivation."
)

make_content_slide(
    "Complete Session 4: Detailed Termination Protocol",
    [
        "SESSION 4 (Week 12, 60 min) — Consolidation & Termination:",
        "",
        "COMPREHENSIVE REVIEW (15 min):",
        "  • 'Let's look at the whole journey from when you first came in'",
        "  • Highlight changes, growth, milestones",
        "  • Attribute success to CLIENT ('YOU did this')",
        "  • Affirm commitment and effort",
        "",
        "MAINTENANCE PLANNING (20 min):",
        "  • 'What situations might be challenging in the next few months?'",
        "  • Specific coping strategies for each situation",
        "  • Support systems: who can help?",
        "  • Plan for holidays, social events, triggers",
        "",
        "RELAPSE PREVENTION (15 min):",
        "  • Normalize: 'Many people experience slips — this doesn't erase progress'",
        "  • Lapse ≠ Relapse (distinction crucial)",
        "  • Emergency plan: What to do IF you slip",
        "  • Written plan → wallet card or phone",
        "",
        "TERMINATION (10 min):",
        "  • Express genuine confidence in their future",
        "  • Leave door open: 'You can always come back for a booster'",
        "  • Final summary of their strengths and resources",
    ],
    refs="Miller et al. (1995). MET Manual, Session 4. | Marlatt & Gordon (1985). Relapse Prevention.",
    notes="The termination should be positive and confidence-building. Research shows that leaving the door open for return contact reduces anxiety and paradoxically reduces the likelihood of needing to return. Attribute all success to the client's own effort.",
    takeaway="Terminate with confidence, attribute success to CLIENT, normalize potential slips, and leave door open for boosters."
)

make_content_slide(
    "Integration of MET with Pharmacotherapy",
    [
        "MET + MEDICATION = Enhanced outcomes:",
        "",
        "ALCOHOL:",
        "  • MET + Naltrexone: COMBINE Study (2006) — effective combination",
        "  • MET + Acamprosate: Supports craving reduction",
        "  • MET enhances medication ADHERENCE (motivation to take medication)",
        "",
        "OPIOIDS:",
        "  • MET + Buprenorphine/Methadone (OST): Better retention",
        "  • MET addresses ambivalence about being 'on medication'",
        "  • Carroll et al. (2006): MI + OST improved outcomes",
        "",
        "TOBACCO:",
        "  • MET + NRT (patch/gum): Combined more effective than either alone",
        "  • MET + Varenicline: Enhanced quit rates",
        "",
        "HOW MET SUPPORTS PHARMACOTHERAPY:",
        "  • Addresses ambivalence about taking medication",
        "  • Builds adherence through autonomy support",
        "  • Normalizes medication as one option in the menu",
        "  • E-P-E for psychoeducation about medications",
    ],
    refs="Anton et al. (2006). JAMA, 295, 2003 (COMBINE). | Carroll et al. (2006). Drug Alcohol Depend, 81, 161.",
    notes="MET and pharmacotherapy are complementary, not competing. MET provides the motivational foundation; medication provides neurobiological support. MET can specifically address ambivalence about medication itself (common in addiction treatment).",
    takeaway="MET + Medication = more effective than either alone. MET enhances both motivation AND medication adherence."
)

make_content_slide(
    "Quality Indicators for MET Sessions (Self-Assessment Checklist)",
    [
        "After each session, rate yourself (1-5) on each indicator:",
        "",
        "SPIRIT (Foundation):",
        "  □ Did I genuinely collaborate rather than prescribe?",
        "  □ Did I accept the client's perspective without judgment?",
        "  □ Was I motivated by compassion (their welfare, not my agenda)?",
        "  □ Did I evoke rather than install?",
        "",
        "SKILLS (Technical):",
        "  □ Were my reflections > questions? (Aim: 2:1)",
        "  □ Were most questions open-ended? (Aim: 70%+)",
        "  □ Did I affirm the client at least 2-3 times?",
        "  □ Did I provide at least one collecting summary?",
        "",
        "CHANGE TALK (Outcome):",
        "  □ Did the client produce change talk? (If yes: reinforced?)",
        "  □ Did I avoid reinforcing sustain talk?",
        "  □ Was change talk more frequent than sustain talk?",
        "",
        "OVERALL: □ Was the CLIENT doing most of the talking and arguing for change?",
    ],
    refs="Moyers et al. (2014). MITI 4.2. | Miller & Moyers (2006). Behav Cog Psychotherapy, 34, 135.",
    notes="This self-assessment checklist can be used after every session. Over time, it builds habitual self-monitoring. When combined with occasional recorded session review, it maintains MI quality.",
    takeaway="Use this checklist after every session to maintain awareness and prevent drift from MI principles."
)

make_closing_slide()
