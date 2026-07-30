#!/usr/bin/env python3
"""
Comprehensive Motivational Enhancement Therapy (MET) Presentation Generator
130-160 slides for M.Phil/PsyD/PhD Clinical Psychology teaching
Based on: Miller et al. (1995) MET Manual (Project MATCH) and NIMHANS Substance Abuse Manual
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# Color palette - Professional medical/psychology blue-white theme
NAVY = RGBColor(0x0F, 0x27, 0x44)
DARK_BLUE = RGBColor(0x1D, 0x5F, 0xA6)
MEDIUM_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
ACCENT_TEAL = RGBColor(0x0E, 0x7C, 0x6B)
ACCENT_GOLD = RGBColor(0xB7, 0x95, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_slide(title, bullets, notes="", refs="", learning_obj="", takeaway=""):
    """Add a formatted slide with all required elements."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Clear default placeholders
    for shape in list(slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Title bar background
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Inches(1.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE


    # Learning Objective (if provided)
    y_pos = 1.2
    if learning_obj:
        obj_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.3), Inches(0.5))
        otf = obj_box.text_frame
        otf.word_wrap = True
        op = otf.paragraphs[0]
        op.text = f"Learning Objective: {learning_obj}"
        op.font.size = Pt(12)
        op.font.italic = True
        op.font.color.rgb = ACCENT_TEAL
        y_pos += 0.5
    
    # Content bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.3), Inches(4.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = ctf.paragraphs[0]
        else:
            p = ctf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
        if bullet.startswith("  "):
            p.level = 1
            p.font.size = Pt(14)


    # Key Take-Home Message box
    if takeaway:
        ta_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(9), Inches(0.6))
        ta_shape.fill.solid()
        ta_shape.fill.fore_color.rgb = LIGHT_BLUE
        ta_shape.line.color.rgb = MEDIUM_BLUE
        ta_tf = ta_shape.text_frame
        ta_tf.word_wrap = True
        ta_p = ta_tf.paragraphs[0]
        ta_p.text = f"Key Takeaway: {takeaway}"
        ta_p.font.size = Pt(11)
        ta_p.font.bold = True
        ta_p.font.color.rgb = NAVY
    
    # References box at bottom
    if refs:
        ref_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.9))
        rtf = ref_box.text_frame
        rtf.word_wrap = True
        rp = rtf.paragraphs[0]
        rp.text = f"References: {refs}"
        rp.font.size = Pt(9)
        rp.font.italic = True
        rp.font.color.rgb = RGBColor(0x5D, 0x6D, 0x7E)
    
    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_title_slide(title, subtitle):
    """Add a section title slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1.5))
    stf = sub_box.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = LIGHT_BLUE
    sp.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_divider(section_num, section_title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE


    # Section number
    num_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    ntf = num_box.text_frame
    np = ntf.paragraphs[0]
    np.text = f"SECTION {section_num}"
    np.font.size = Pt(18)
    np.font.color.rgb = ACCENT_GOLD
    np.alignment = PP_ALIGN.CENTER
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(2))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = section_title
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Section {section_num}: {section_title}. Transition to this new section of the presentation."
    
    return slide

# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
add_title_slide(
    "MOTIVATIONAL ENHANCEMENT THERAPY (MET)",
    "A Comprehensive Evidence-Based Presentation for Postgraduate Clinical Psychology Education\n"
    "Based on Project MATCH Manual (Miller et al., 1995) & NIMHANS Psychosocial Interventions Manual"
)


# ============================================================
# SECTION 1: INTRODUCTION (Slides 2-10)
# ============================================================
add_section_divider(1, "Introduction to Motivational Enhancement Therapy")

add_slide(
    "Definition of Motivation",
    [
        "Motivation: An internal state that energizes, directs, and sustains behavior toward goals",
        "In clinical context: The probability that a person will enter into, continue, and adhere to a specific change strategy",
        "Motivation is NOT a personality trait or static characteristic",
        "It is a state of readiness or eagerness to change, which may fluctuate over time and across situations",
        "Miller (1985): 'Motivation is not something that one has but rather something one does'",
        "Motivation involves the interaction between the person and their environment",
    ],
    notes="Explain that motivation has been misunderstood in addiction treatment for decades. Traditional view saw motivation as something clients either had or didn't have. Miller revolutionized this by showing motivation is a dynamic, interpersonal process that therapists can influence. Discuss how this shifts responsibility partially to the therapeutic interaction rather than solely to the client.",
    refs="Miller, W.R. (1985). Motivation for treatment: A review with special emphasis on alcoholism. Psychological Bulletin, 98, 84-107. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing: Helping people change (3rd ed.). Guilford Press.",
    learning_obj="Define motivation in clinical context and understand its dynamic nature",
    takeaway="Motivation is not a fixed trait but a dynamic state that therapists can actively influence through their interaction style."
)

add_slide(
    "What is Behaviour Change?",
    [
        "Behaviour change: The process of modifying habitual patterns of action",
        "Involves cognitive, emotional, and behavioral components",
        "Key principles of behaviour change:",
        "  1. Change is a process, not an event",
        "  2. People move through identifiable stages",
        "  3. Ambivalence is normal and expected",
        "  4. The therapeutic relationship strongly influences change",
        "  5. Intrinsic motivation is more durable than extrinsic pressure",
        "In substance use: Moving from harmful patterns to healthier alternatives",
        "Behaviour change requires both motivation (why) and self-efficacy (how)",
    ],
    notes="Discuss how behaviour change in substance use is complex and multifaceted. Explain that traditional models assumed change was simply a matter of willpower. Modern understanding recognizes that change involves neurobiological, psychological, social, and environmental factors. The therapeutic relationship is a key ingredient. Reference the NIMHANS manual's emphasis on stepped care approaches.",
    refs="Prochaska, J.O., & DiClemente, C.C. (1984). The transtheoretical approach. Dow Jones/Irwin. | Murthy, P. (2008). Psychosocial interventions for persons with substance abuse. NIMHANS Publication.",
    learning_obj="Understand the multidimensional nature of behaviour change",
    takeaway="Change is a process with identifiable stages; ambivalence is normal and the therapeutic relationship is a key driver of change."
)


add_slide(
    "Why People Resist Change",
    [
        "Psychological Reactance: Perceived threats to personal freedom trigger opposition",
        "Ambivalence: Simultaneous desire to change AND to maintain status quo",
        "Fear of the unknown: Familiar patterns feel safe even when harmful",
        "Loss of perceived benefits: Substances serve functions (coping, social, pleasure)",
        "Cognitive dissonance avoidance: Minimization, denial, rationalization",
        "Low self-efficacy: 'I can't do it' beliefs prevent attempts",
        "Environmental reinforcement: Social networks that support continued use",
        "Neurobiological factors: Brain reward circuits maintain addictive behaviors",
        "Secondary gains: Attention, sick role, avoidance of responsibilities",
    ],
    notes="This slide sets the stage for understanding WHY MET was developed. Traditional approaches interpreted resistance as pathological denial. Miller argued that resistance is largely an interpersonal phenomenon - it is significantly influenced by the therapist's style. Aggressive confrontation increases resistance. Empathic, reflective approaches decrease it. Reference Patterson & Forgatch (1985) who showed therapist behavior directly determined client resistance.",
    refs="Miller, W.R., et al. (1995). Motivational Enhancement Therapy Manual. NIAAA. | Brehm, S.S., & Brehm, J.W. (1981). Psychological reactance. Academic Press. | Patterson, G.A., & Forgatch, M.S. (1985). Journal of Consulting and Clinical Psychology, 53, 846-851.",
    learning_obj="Identify psychological and biological factors underlying resistance to change",
    takeaway="Resistance is a normal interpersonal phenomenon influenced by therapist style, not solely a client deficit."
)

add_slide(
    "Why Traditional Confrontation Fails",
    [
        "Confrontational approaches assume denial is a defense that must be 'broken through'",
        "Research evidence: Aggressive confrontation INCREASES resistance (Miller et al., 1993)",
        "Confrontation evokes defensive argumentation from clients",
        "Self-perception theory: 'If I argue I don't have a problem, I believe it more'",
        "Miller (1995): 'The worst persuasion strategy evokes defensive argumentation'",
        "Confrontation damages therapeutic alliance - the strongest predictor of outcome",
        "High therapist confrontation correlates with poor outcomes and higher dropout",
        "Valle (1981): Empathic therapist style predicted 2-year outcomes better than any other factor",
        "The Dodo Bird Effect: Relationship factors outweigh technique across therapies",
    ],
    notes="Emphasize the historical context. For decades, addiction treatment relied on aggressive confrontation (e.g., Synanon, 'hot seat' approaches). Miller's research showed this was counterproductive. When therapists argued that clients were alcoholics, clients predictably argued back that they weren't. This verbal behavior actually reinforced their resistance. The manual explicitly states: 'An aggressive argument that you're an alcoholic will usually evoke: No I'm not, and no I don't.' This insight was revolutionary.",
    refs="Miller, W.R., Benefield, R.G., & Tonigan, J.S. (1993). Journal of Consulting and Clinical Psychology, 61, 455-461. | Valle, S.K. (1981). Journal of Studies on Alcohol, 42, 783-790. | Miller, W.R., et al. (1995). MET Manual, Project MATCH.",
    learning_obj="Understand evidence against confrontational approaches in addiction treatment",
    takeaway="Confrontation increases resistance; empathic approaches produce better outcomes in addiction treatment."
)


add_slide(
    "Introduction to Motivational Enhancement Therapy (MET)",
    [
        "MET: A systematic intervention approach based on principles of motivational psychology",
        "Designed to produce rapid, internally motivated change",
        "Does NOT guide client step-by-step through recovery",
        "Instead employs motivational strategies to mobilize client's OWN resources",
        "Originally developed for Project MATCH (NIAAA, 1989-1997)",
        "A brief, 4-session structured intervention",
        "Combines motivational interviewing with personalized assessment feedback",
        "Evidence-based: One of three treatments in Project MATCH alongside CBT and TSF",
        "Cost-effective and time-efficient compared to longer treatments",
    ],
    notes="MET was developed specifically for Project MATCH, the largest clinical trial of alcoholism treatments ever conducted. It was designed by William R. Miller and colleagues as a brief, 4-session treatment. Unlike CBT (12 sessions) and TSF (12 sessions), MET achieved comparable outcomes in just 4 sessions, making it remarkably cost-effective. The therapy is distinguished by its reliance on the client's own motivation rather than external direction.",
    refs="Miller, W.R., Zweben, A., DiClemente, C.C., & Rychtarik, R.G. (1995). Motivational Enhancement Therapy Manual. NIAAA, NIH Publication No. 94-3723. | Project MATCH Research Group. (1997). Matching alcoholism treatments to client heterogeneity. Journal of Studies on Alcohol, 58, 7-29.",
    learning_obj="Define MET and understand its position among evidence-based addiction treatments",
    takeaway="MET is a brief, 4-session evidence-based intervention that mobilizes clients' own internal motivation for change."
)

add_slide(
    "Definition and Key Features of MET",
    [
        "DEFINITION: MET is a client-centered, directive therapeutic style that enhances",
        "  intrinsic motivation to change by exploring and resolving ambivalence",
        "KEY FEATURES:",
        "  • Brief (4 sessions in Project MATCH protocol)",
        "  • Structured but flexible and individualized",
        "  • Based on personalized assessment feedback",
        "  • Non-confrontational and non-judgmental",
        "  • Client-centered in spirit, directive in strategy",
        "  • Focuses on eliciting change talk from the client",
        "  • Supports client autonomy and self-efficacy",
        "  • Uses specific motivational strategies (FRAMES, OARS)",
        "MET = Motivational Interviewing + Structured Feedback + Personal Plan",
    ],
    notes="Clarify the distinction between MI and MET. Motivational Interviewing is a broader counseling style/spirit. MET is a specific manualized treatment that incorporates MI principles within a structured 4-session format that includes assessment feedback. MET adds the structured feedback component (Personal Feedback Report) to the MI framework. This makes it more directive and protocol-driven while maintaining the MI spirit.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (1991). Motivational interviewing: Preparing people to change addictive behavior. Guilford Press. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Distinguish MET from Motivational Interviewing and identify key features",
    takeaway="MET combines the spirit of MI with structured personalized feedback in a brief 4-session protocol."
)


add_slide(
    "Origin and Historical Development of MET",
    [
        "1983: Miller publishes 'Motivational Interviewing with Problem Drinkers' - seminal paper",
        "1989: Project MATCH initiated by NIAAA Treatment Research Branch",
        "1991: Miller & Rollnick publish first MI textbook",
        "1992: MET manual drafted for Project MATCH (9 clinical sites)",
        "1994: MET Manual published (NIH Publication No. 94-3723)",
        "1997: Project MATCH results published - MET equivalent to 12-session treatments",
        "2002: Miller & Rollnick publish 2nd edition incorporating MI advances",
        "2008: NIMHANS Manual integrates MET for Indian clinical settings",
        "2013: Miller & Rollnick 3rd edition - updated spirit and processes",
        "2020s: MET adapted globally for multiple health behaviors beyond addiction",
    ],
    notes="Walk through the historical timeline emphasizing key milestones. Miller's 1983 paper was influenced by his experience that empathic therapists got better outcomes and his observation that confrontational approaches increased resistance. Project MATCH was the largest ($27 million) multisite clinical trial of alcoholism treatment. The finding that 4 sessions of MET matched 12 sessions of CBT or TSF was groundbreaking for the field.",
    refs="Miller, W.R. (1983). Motivational interviewing with problem drinkers. Behavioural Psychotherapy, 11, 147-172. | Project MATCH Research Group. (1997). Journal of Studies on Alcohol, 58, 7-29. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Trace the historical development of MET from its origins to current applications",
    takeaway="MET evolved from Miller's 1983 paper through Project MATCH, proving that brief motivational therapy matches longer treatments."
)

add_slide(
    "Difference Between Motivation and Compliance",
    [
        "MOTIVATION vs. COMPLIANCE - A Critical Distinction:",
        "  Motivation: Internal desire to change; autonomous, self-directed",
        "  Compliance: External conformity to demands; controlled, other-directed",
        "MOTIVATION characteristics:",
        "  • Internally driven, self-determined",
        "  • Sustainable long-term",
        "  • Associated with genuine behavior change",
        "  • Resilient to setbacks",
        "COMPLIANCE characteristics:",
        "  • Externally driven by rewards/punishments",
        "  • Temporary, depends on external monitoring",
        "  • Associated with surface-level change only",
        "  • Collapses when external pressure removed",
        "Self-Determination Theory: Autonomy, Competence, Relatedness drive true motivation",
    ],
    notes="This distinction is crucial for understanding why MET works differently from coercive approaches. Many mandated clients appear compliant but are not internally motivated. True behavior change requires internalization of the motivation. SDT research shows that when people feel autonomous (not controlled), competent (self-efficacious), and connected (therapeutic relationship), they develop intrinsic motivation. MET systematically fosters all three of these needs.",
    refs="Deci, E.L., & Ryan, R.M. (2000). The 'what' and 'why' of goal pursuits. Psychological Inquiry, 11, 227-268. | Miller, W.R., et al. (1995). MET Manual. | Ryan, R.M., & Deci, E.L. (2000). American Psychologist, 55, 68-78.",
    learning_obj="Differentiate intrinsic motivation from compliance in therapeutic contexts",
    takeaway="MET targets genuine internal motivation rather than mere compliance, leading to more durable behavior change."
)


add_slide(
    "Why MET Was Developed",
    [
        "Clinical need: Many clients dropped out of confrontational programs",
        "Research gap: No brief, empirically-supported motivational treatment existed",
        "Project MATCH required a brief comparison therapy (4 sessions vs. 12 sessions)",
        "Cost-effectiveness: Healthcare systems needed shorter, effective interventions",
        "Therapist variability: Need for standardized motivational approach",
        "Growing evidence that therapist empathy predicted outcomes (Valle, 1981; Miller et al., 1980)",
        "Recognition that 'denial' was largely therapist-created resistance",
        "Need to match treatment intensity to client readiness",
        "NIMHANS (2008): Need for culturally appropriate, feasible interventions in developing countries",
        "Primary care integration: Brief, trainable approach for non-specialists",
    ],
    notes="Explain the convergence of factors that led to MET's development. The field needed: (1) evidence that brief interventions could work, (2) a standardized protocol that could be taught and replicated, (3) a humane alternative to confrontation, and (4) something cost-effective. The NIMHANS manual emphasizes the Indian context where resources are limited and training non-specialists is essential. The combination of these needs made MET an ideal solution.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Murthy, P. (2008). NIMHANS Manual. | Holder, H.D., et al. (1991). Journal of Studies on Alcohol, 52, 517-540.",
    learning_obj="Understand the clinical and research contexts that necessitated MET's development",
    takeaway="MET was developed to provide a brief, empirically-supported, non-confrontational alternative to traditional addiction treatment."
)

# ============================================================
# SECTION 2: FOUNDERS AND CONTRIBUTORS (Slides 11-14)
# ============================================================
add_section_divider(2, "Founders and Key Contributors")

add_slide(
    "William R. Miller, PhD - Father of MI/MET",
    [
        "Professor Emeritus, University of New Mexico",
        "Born 1947 - Over 50 years of contributions to addiction science",
        "Key contributions:",
        "  • Created Motivational Interviewing (1983)",
        "  • Lead author of MET Manual for Project MATCH",
        "  • Published over 400 articles and 50 books",
        "  • Developed the Drinker's Check-Up (precursor to MET feedback)",
        "  • Pioneered research on therapist effects in treatment",
        "  • Showed that therapist empathy predicts client outcomes",
        "Theoretical influences: Carl Rogers, cognitive dissonance, self-perception theory",
        "Awards: Jellinek Memorial Award, Innovators Award (Robert Wood Johnson Foundation)",
        "Philosophy: 'People are generally the best experts on themselves'",
    ],
    notes="Miller is widely regarded as one of the most influential figures in addiction treatment. His 1983 paper was written after visiting Norway and being asked by trainees how he worked with problem drinkers. His response - drawing on Rogers, Festinger, and Bem - became motivational interviewing. His research consistently showed that the way therapists interact with clients matters more than the specific technique used. He demonstrated that a single session of empathic counseling could produce lasting change.",
    refs="Miller, W.R. (1983). Behavioural Psychotherapy, 11, 147-172. | Miller, W.R., & Rollnick, S. (1991, 2002, 2013). Motivational Interviewing. Guilford Press. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Appreciate William R. Miller's foundational contributions to MI and MET",
    takeaway="Miller's insight that therapist style determines client motivation revolutionized addiction treatment worldwide."
)


add_slide(
    "Stephen Rollnick, PhD & Other Key Contributors",
    [
        "STEPHEN ROLLNICK (Cardiff University, Wales):",
        "  • Co-developer of Motivational Interviewing with Miller",
        "  • Focused on MI applications in healthcare settings",
        "  • Developed brief MI adaptations for medical consultations",
        "THERESA MOYERS (University of New Mexico):",
        "  • Leading MI process researcher",
        "  • Developer of MITI (Motivational Interviewing Treatment Integrity) coding",
        "  • Research on mechanisms of MI effectiveness",
        "PROJECT MATCH INVESTIGATORS:",
        "  • Allen Zweben, DSW - MET manual co-author (University of Wisconsin-Milwaukee)",
        "  • Carlo DiClemente, PhD - Transtheoretical Model co-creator",
        "  • Roberto Rychtarik, PhD - Research Institute on Alcoholism",
        "  • Thomas Babor, PhD - Coordinating Center (University of Connecticut)",
        "  • Kathleen Carroll, PhD - Yale University, treatment design coordinator",
    ],
    notes="Rollnick brought MI into mainstream healthcare beyond addiction. His books on MI in health settings made the approach accessible to physicians, nurses, and other healthcare professionals. Moyers' work on MI process research has been crucial in understanding HOW MI works - she showed that therapist MI-consistent behaviors predict client change talk, which in turn predicts behavior change outcomes. DiClemente's Transtheoretical Model provided the theoretical staging framework that MET uses.",
    refs="Rollnick, S., Miller, W.R., & Butler, C.C. (2008). Motivational interviewing in health care. Guilford Press. | Moyers, T.B., et al. (2009). Journal of Substance Abuse Treatment, 36, 101-109. | DiClemente, C.C. (2003). Addiction and change. Guilford Press.",
    learning_obj="Identify key contributors to MI/MET development and their specific contributions",
    takeaway="MET was developed through collaboration of leading researchers in motivation, addiction, and behavior change."
)

add_slide(
    "Project MATCH Research Group",
    [
        "Largest clinical trial of alcoholism treatment ($27 million, 1989-1997)",
        "9 clinical research sites across the United States",
        "DESIGN: 1,726 patients randomized to 3 treatments:",
        "  • MET: 4 sessions over 12 weeks",
        "  • CBT (Cognitive-Behavioral): 12 sessions over 12 weeks",
        "  • TSF (Twelve-Step Facilitation): 12 sessions over 12 weeks",
        "Two parallel arms: Outpatient (5 sites) and Aftercare (4 sites)",
        "Follow-up: 3, 6, 9, 12, and 15 months post-treatment",
        "KEY FINDING: All three treatments produced significant improvements",
        "  MET with only 4 sessions matched 12-session treatments on most outcomes",
        "Principal sites: UNM, Buffalo, Charleston, Connecticut, Houston, Milwaukee, Brown, Seattle",
    ],
    notes="Project MATCH was groundbreaking in its scope and rigor. The primary hypothesis was that patient-treatment matching would improve outcomes, but the main finding was that matching effects were minimal. The most striking result was that MET, with only one-third the sessions, produced essentially equivalent outcomes to the longer treatments. This established MET as a remarkably cost-effective intervention. The study used standardized manuals, trained therapists, and rigorous fidelity monitoring.",
    refs="Project MATCH Research Group. (1997). Journal of Studies on Alcohol, 58, 7-29. | Project MATCH Research Group. (1998). Addiction, 93, 1431-1446. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Understand Project MATCH design and its significance for establishing MET's evidence base",
    takeaway="Project MATCH demonstrated that 4 sessions of MET achieved outcomes comparable to 12-session treatments."
)


# ============================================================
# SECTION 3: HISTORICAL BACKGROUND (Slides 15-17)
# ============================================================
add_section_divider(3, "Historical Background")

add_slide(
    "Origins of Motivational Interviewing",
    [
        "1950s-1970s: Addiction treatment dominated by confrontational models",
        "  • Synanon (1958), 'hot seat' therapy, attack therapy",
        "  • Assumption: Addicts are in 'denial' that must be broken",
        "  • Therapist as authoritarian expert, client as resistant patient",
        "1950s: Carl Rogers develops client-centered therapy (unconditional positive regard)",
        "1957: Festinger publishes A Theory of Cognitive Dissonance",
        "1965-1972: Daryl Bem develops Self-Perception Theory",
        "1980: Miller observes that empathic therapists get better outcomes",
        "1983: Miller's seminal paper combines Rogers + Festinger + Bem",
        "Key insight: Help people TALK THEMSELVES INTO change rather than being told to change",
    ],
    notes="The historical context is essential for understanding why MI/MET was revolutionary. Before 1983, the dominant paradigm in addiction treatment (especially in the US) was confrontational. Programs like Synanon used attack therapy, humiliation, and breaking down defenses. Miller noticed in his own research that when therapists were more empathic, clients did better. He combined Rogerian principles (creating safe space), Festinger's cognitive dissonance (awareness of discrepancy), and Bem's self-perception theory (we believe what we hear ourselves say) into a new approach.",
    refs="Miller, W.R. (1983). Behavioural Psychotherapy, 11, 147-172. | Rogers, C.R. (1957). Journal of Consulting Psychology, 21, 95-103. | Festinger, L. (1957). A Theory of Cognitive Dissonance. Row, Peterson.",
    learning_obj="Trace the intellectual origins of motivational approaches in addiction treatment",
    takeaway="MI emerged from combining Rogerian empathy, cognitive dissonance theory, and self-perception theory."
)

add_slide(
    "Addiction Treatment Before MET: The Confrontational Model",
    [
        "Traditional Confrontational Assumptions:",
        "  • Addiction = moral weakness or character defect",
        "  • 'Denial' is a defense mechanism requiring forceful confrontation",
        "  • Therapist must 'break through' resistance",
        "  • Client must 'hit rock bottom' before change is possible",
        "  • 'Tough love' is necessary and therapeutic",
        "Evidence AGAINST confrontation:",
        "  • Higher dropout rates with confrontational counselors",
        "  • Patterson & Forgatch (1985): Therapist confrontation directly increased resistance",
        "  • Miller et al. (1993): Confrontational style predicted WORSE drinking outcomes at 1 year",
        "  • No evidence that confrontation improves outcomes over any alternative approach",
        "The Paradigm Shift: From 'breaking denial' to 'resolving ambivalence'",
    ],
    notes="This slide contrasts the old and new paradigms. Historically, if a client said 'I don't have a problem,' the therapist would argue harder. This created an adversarial dynamic. Miller's research showed that the amount of client resistance in session 1 predicted drinking outcomes at 1 year - and this resistance was largely determined by therapist style. The new paradigm reframes 'denial' as ambivalence and 'resistance' as a signal that the therapist needs to change approach.",
    refs="Miller, W.R., Benefield, R.G., & Tonigan, J.S. (1993). JCCP, 61, 455-461. | Patterson, G.A., & Forgatch, M.S. (1985). JCCP, 53, 846-851. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Compare confrontational and motivational paradigms in addiction treatment",
    takeaway="The shift from confrontation to motivational approaches was driven by clear evidence that confrontation worsens outcomes."
)


# ============================================================
# SECTION 4: THEORETICAL FOUNDATIONS (Slides 18-26)
# ============================================================
add_section_divider(4, "Theoretical Foundations of MET")

add_slide(
    "Humanistic Psychology & Carl Rogers",
    [
        "HUMANISTIC PSYCHOLOGY (Maslow, Rogers, May):",
        "  • People have inherent capacity for growth and self-actualization",
        "  • Focus on human potential rather than pathology",
        "  • Subjective experience is central to understanding behavior",
        "CARL ROGERS (1902-1987) - Client-Centered Therapy:",
        "  • Core conditions: Empathy, Unconditional Positive Regard, Congruence",
        "  • The relationship IS the therapy",
        "  • Non-directive: Trust the client's own wisdom",
        "  • Rogers (1957): These conditions are necessary and sufficient for change",
        "RELATION TO MET:",
        "  • MET adopts Rogerian empathy as fundamental therapeutic stance",
        "  • Client is viewed as expert on their own life",
        "  • Therapist provides conditions for self-exploration",
    ],
    notes="Rogers' influence on MET cannot be overstated. The therapeutic relationship in MET is fundamentally Rogerian - the therapist creates a safe, accepting, empathic environment where the client can explore ambivalence without judgment. However, MET adds a DIRECTIVE component that pure Rogerian therapy lacks - the therapist strategically reinforces change talk and uses structured feedback. This makes MET 'client-centered yet directive' - a unique combination.",
    refs="Rogers, C.R. (1957). Journal of Consulting Psychology, 21, 95-103. | Rogers, C.R. (1959). In Koch, S. (Ed.), Psychology: A Study of a Science, Vol. 3. McGraw-Hill. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Understand Rogers' client-centered principles and their integration into MET",
    takeaway="MET adopts Rogers' empathic stance while adding strategic direction - creating a 'client-centered directive' approach."
)

add_slide(
    "Self-Determination Theory (Deci & Ryan)",
    [
        "FOUNDERS: Edward Deci & Richard Ryan (University of Rochester, 1985-present)",
        "CORE ASSUMPTIONS:",
        "  • Humans have innate psychological needs that drive behavior",
        "  • Three basic needs: Autonomy, Competence, Relatedness",
        "  • Motivation exists on a continuum from extrinsic to intrinsic",
        "THREE BASIC NEEDS:",
        "  • Autonomy: Need to feel in control of one's own behavior",
        "  • Competence: Need to feel effective and capable (= self-efficacy)",
        "  • Relatedness: Need to feel connected to others (= therapeutic alliance)",
        "RELATION TO MET:",
        "  • MET supports autonomy (client decides, not therapist)",
        "  • MET builds competence (supporting self-efficacy)",
        "  • MET provides relatedness (empathic therapeutic relationship)",
        "  • Extrinsic controls undermine intrinsic motivation",
    ],
    notes="SDT provides the theoretical backbone for WHY MET's non-controlling approach works better than authoritarian approaches. When we tell people what to do, we undermine their autonomy, which reduces intrinsic motivation. When we support choice, we enhance autonomy and thereby intrinsic motivation. MET systematically supports all three basic needs. The clinical implication is clear: offering a menu of options (autonomy), affirming client strengths (competence), and providing empathic listening (relatedness) naturally enhance motivation.",
    refs="Deci, E.L., & Ryan, R.M. (1985). Intrinsic motivation and self-determination in human behavior. Plenum. | Ryan, R.M., & Deci, E.L. (2000). American Psychologist, 55, 68-78. | Markland, D., et al. (2005). Addictive Behaviors, 30, 1859-1869.",
    learning_obj="Apply Self-Determination Theory principles to understand MET's effectiveness",
    takeaway="MET supports autonomy, competence, and relatedness - the three basic needs that drive intrinsic motivation."
)


add_slide(
    "Cognitive Dissonance Theory (Festinger, 1957)",
    [
        "FOUNDER: Leon Festinger (1957)",
        "CORE PRINCIPLE: Psychological discomfort when behavior conflicts with beliefs/values",
        "ASSUMPTIONS:",
        "  • People seek internal consistency between attitudes, beliefs, and behaviors",
        "  • Discrepancy between cognitions creates uncomfortable 'dissonance'",
        "  • People are motivated to reduce dissonance (change belief OR behavior)",
        "CLINICAL IMPLICATIONS:",
        "  • When clients see discrepancy between values and behavior → motivation increases",
        "  • MET's 'Develop Discrepancy' principle directly applies this theory",
        "  • Personalized feedback highlights the gap between goals and current behavior",
        "ADVANTAGES: Powerful motivator; changes come from within",
        "LIMITATIONS: Too much dissonance → defensiveness; requires careful calibration",
        "RELATION TO MET: The therapist helps clients see their OWN discrepancies, gently",
    ],
    notes="Festinger showed that when people hold contradictory beliefs simultaneously, they experience psychological discomfort and seek to resolve it. MET exploits this by helping clients see the gap between their stated values (e.g., 'I want to be a good father') and their current behavior (e.g., drinking heavily). The key is that the client must articulate this discrepancy themselves - if the therapist points it out, it may feel confrontational. The feedback session in MET is designed to create optimal discrepancy.",
    refs="Festinger, L. (1957). A Theory of Cognitive Dissonance. Row, Peterson. | Miller, W.R., et al. (1995). MET Manual. | Miller, W.R. (1983). Behavioural Psychotherapy, 11, 147-172.",
    learning_obj="Explain how cognitive dissonance theory underpins MET's develop discrepancy principle",
    takeaway="MET creates therapeutic cognitive dissonance by helping clients recognize gaps between their values and behaviors."
)

add_slide(
    "Transtheoretical Model (Prochaska & DiClemente)",
    [
        "FOUNDERS: James Prochaska & Carlo DiClemente (1982)",
        "THE STAGES OF CHANGE:",
        "  1. Precontemplation: Not considering change; unaware or in denial",
        "  2. Contemplation: Ambivalent; aware of problem but not committed to action",
        "  3. Preparation: Intending to take action; making small steps",
        "  4. Action: Actively modifying behavior, environment, experiences",
        "  5. Maintenance: Sustaining change; preventing relapse",
        "  6. Relapse: Return to earlier stage (normalized, not failure)",
        "KEY PRINCIPLES:",
        "  • Change is cyclical, not linear",
        "  • Different stages require different interventions",
        "  • MET is especially effective for Precontemplation → Contemplation → Preparation",
        "  • Matching intervention to stage increases effectiveness",
    ],
    notes="The Transtheoretical Model is foundational to MET. Carlo DiClemente was a co-author of the MET manual, directly connecting stage theory to MET practice. The model tells us that most clients entering treatment are in the contemplation stage - they're ambivalent. MET is specifically designed for this ambivalent population. For precontemplators, MET uses gentle discrepancy development. For contemplators, it explores ambivalence. For those in preparation, it consolidates commitment.",
    refs="Prochaska, J.O., & DiClemente, C.C. (1982). Psychotherapy: Theory, Research and Practice, 19, 276-288. | Prochaska, J.O., & DiClemente, C.C. (1984). The Transtheoretical Approach. Dow Jones/Irwin. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Apply the Stages of Change model to understand MET's stage-matched approach",
    takeaway="MET matches motivational strategies to the client's stage of change, especially targeting ambivalence in contemplation."
)


add_slide(
    "Social Learning Theory & Self-Efficacy (Bandura)",
    [
        "FOUNDER: Albert Bandura (1977, 1982)",
        "SOCIAL LEARNING THEORY:",
        "  • Behavior is learned through observation, modeling, and reinforcement",
        "  • Reciprocal determinism: behavior, cognition, and environment interact",
        "  • Drinking behavior is functionally related to psychosocial problems",
        "SELF-EFFICACY THEORY:",
        "  • Self-efficacy: Belief in one's capability to execute behaviors required for outcomes",
        "  • Four sources: Mastery experiences, Vicarious experiences, Verbal persuasion, Emotional states",
        "  • Higher self-efficacy → more persistent effort → better outcomes",
        "RELATION TO MET:",
        "  • 'Support Self-Efficacy' is a core MET principle",
        "  • Affirmations build client's belief in ability to change",
        "  • Past successes are highlighted to enhance efficacy",
        "  • Client autonomy preserves sense of personal control",
    ],
    notes="Bandura's self-efficacy theory is directly operationalized in MET through the principle 'Support Self-Efficacy.' The MET manual explicitly states that the therapist must convey belief in the client's ability to change. This is done through affirmations, highlighting past successes, normalizing the change process, and offering a range of options (which communicates that change is possible in multiple ways). Low self-efficacy is one of the strongest predictors of relapse; building it is therefore therapeutic.",
    refs="Bandura, A. (1977). Self-efficacy: Toward a unifying theory of behavioral change. Psychological Review, 84, 191-215. | Bandura, A. (1982). American Psychologist, 37, 122-147. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Understand Bandura's self-efficacy theory and its operationalization in MET",
    takeaway="MET systematically builds self-efficacy through affirmations, past successes, and supporting client autonomy."
)

add_slide(
    "Self-Perception Theory (Bem) & Decision-Making Theory",
    [
        "SELF-PERCEPTION THEORY (Daryl Bem, 1965, 1972):",
        "  • 'As I hear myself talk, I learn what I believe'",
        "  • People infer their attitudes from their own behavior and speech",
        "  • If I say it and no one forced me → I must believe it",
        "  • IMPLICATION: Eliciting change talk literally CREATES motivation",
        "  • Therapist's role: Create conditions where client argues FOR change",
        "DECISION-MAKING THEORY (Janis & Mann, 1977):",
        "  • Decisions involve weighing pros and cons (decisional balance)",
        "  • People use heuristics and are influenced by framing effects",
        "  • MET uses decisional balance to explore both sides of ambivalence",
        "COMBINED APPLICATION IN MET:",
        "  • Evoke change talk (self-perception) while exploring the balance (decision theory)",
        "  • The MORE clients talk about change, the MORE they believe in it",
        "  • Therapist strategically reinforces change talk through reflection",
    ],
    notes="Bem's self-perception theory is perhaps the most direct theoretical underpinning of MET technique. The MET manual states explicitly: 'motivational psychology has amply demonstrated that when people are subtly enticed to speak or act in a new way, their beliefs and values tend to shift in that direction.' This is why the therapist seeks to elicit SELF-MOTIVATIONAL STATEMENTS from the client. Every time the client voices a reason for change, they become slightly more motivated. This creates a positive feedback loop.",
    refs="Bem, D.J. (1972). In Berkowitz, L. (Ed.), Advances in Experimental Social Psychology, Vol. 6. Academic Press. | Janis, I.L., & Mann, L. (1977). Decision Making. Free Press. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Apply self-perception and decision-making theories to MET practice",
    takeaway="People believe what they hear themselves say - MET leverages this by strategically eliciting change talk."
)


add_slide(
    "Theoretical Integration: How All Theories Connect in MET",
    [
        "INTEGRATED FRAMEWORK - Each theory contributes a key element:",
        "  Rogers → Empathic therapeutic relationship (HOW we relate)",
        "  Festinger → Develop discrepancy (WHAT drives change)",
        "  Bem → Elicit change talk (WHY talking creates believing)",
        "  Bandura → Support self-efficacy (CAN I do it?)",
        "  Prochaska/DiClemente → Stage matching (WHEN to intervene)",
        "  Deci/Ryan → Autonomy support (WHO controls the process)",
        "  Janis/Mann → Decisional balance (HOW to weigh options)",
        "THE MET FORMULA:",
        "  Empathic relationship + Discrepancy awareness + Change talk",
        "  + Self-efficacy + Stage-appropriate + Autonomy-supportive",
        "  = Internally motivated, sustainable behavior change",
        "[FLOWCHART SUGGESTION: Visual showing theories feeding into MET]",
    ],
    notes="This integration slide shows students how MET is not an atheoretical technique but is deeply grounded in multiple established psychological theories. Each theory addresses a different component of the motivation and change process. Rogers provides the relational foundation, Festinger provides the motivational engine (discomfort drives change), Bem explains the mechanism of change talk, Bandura provides the confidence component, Prochaska provides the timing framework, and Deci/Ryan explain why autonomy-support works better than coercion.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Miller, W.R., et al. (1995). MET Manual. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Synthesize multiple theoretical frameworks into a unified understanding of MET",
    takeaway="MET integrates 7+ theoretical traditions into a coherent clinical approach that addresses why, how, and when people change."
)

# ============================================================
# SECTION 5: SPIRIT OF MET (Slides 27-30)
# ============================================================
add_section_divider(5, "The Spirit of MET")

add_slide(
    "The Four Elements of MET Spirit",
    [
        "The SPIRIT of MET is more important than any specific technique",
        "Miller & Rollnick (2013): Spirit = Partnership + Acceptance + Compassion + Evocation",
        "",
        "1. PARTNERSHIP: Collaborative relationship; 'dancing' not 'wrestling'",
        "2. ACCEPTANCE: Absolute worth, accurate empathy, autonomy support, affirmation",
        "3. COMPASSION: Actively promoting the client's welfare, prioritizing their needs",
        "4. EVOCATION: Drawing out what is already there; the client has the answers",
        "",
        "Without the spirit, techniques become manipulative",
        "The spirit determines whether the approach is truly motivational",
        "Assessment: 'Am I trying to make them change, or help them explore their own reasons?'",
        "Common error: Using MI techniques without MI spirit = 'righting reflex'",
    ],
    notes="The spirit of MET/MI is the MOST important concept to convey. Techniques without spirit become manipulative tools. The spirit distinguishes genuine motivational enhancement from persuasion or manipulation. Partnership means working WITH the client, not doing TO them. Acceptance means valuing the person regardless of their behavior. Compassion means genuinely caring about their wellbeing. Evocation means drawing out their own motivation rather than installing yours. If students take away only one thing, it should be the spirit.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Miller, W.R., et al. (1995). MET Manual. | Moyers, T.B., & Miller, W.R. (2013). Psychotherapy, 50, 338-345.",
    learning_obj="Internalize the four elements of MET spirit as the foundation of the approach",
    takeaway="The spirit (partnership, acceptance, compassion, evocation) is more important than any technique in MET."
)


add_slide(
    "Partnership & Acceptance in Practice",
    [
        "PARTNERSHIP - Clinical Examples:",
        "  WRONG: 'I think you need to quit drinking. Here's what you should do...'",
        "  CORRECT: 'What are your thoughts about your drinking? What concerns you?'",
        "  WRONG: 'Your test results show you're an alcoholic'",
        "  CORRECT: 'Here are your results. What do you make of these?'",
        "",
        "ACCEPTANCE - Four Components:",
        "  1. Absolute Worth: Every person deserves respect regardless of behavior",
        "  2. Accurate Empathy: Deeply understanding the client's perspective",
        "  3. Autonomy Support: 'It's your life and your choice'",
        "  4. Affirmation: Recognizing strengths and efforts",
        "",
        "WRONG: 'You're in denial and you need to accept you're an alcoholic'",
        "CORRECT: 'I can see you're genuinely struggling with this decision, and I respect that'",
    ],
    notes="Use role-play demonstrations to illustrate these principles. The key distinction is WHO is the expert. In traditional approaches, the therapist is the expert telling the client what to do. In MET, the client is the expert on their own life; the therapist provides information and creates conditions for self-exploration. Acceptance does NOT mean agreeing with harmful behavior - it means accepting the PERSON while recognizing that behavior may need to change. This is identical to Rogers' unconditional positive regard.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Miller, W.R., et al. (1995). MET Manual. | Rogers, C.R. (1957). Journal of Consulting Psychology, 21, 95-103.",
    learning_obj="Demonstrate partnership and acceptance through clinical dialogue examples",
    takeaway="Partnership treats clients as equals; acceptance values the person unconditionally while the behavior may need to change."
)

add_slide(
    "Compassion & Evocation in Practice",
    [
        "COMPASSION:",
        "  • Actively promoting client welfare as primary concern",
        "  • NOT feeling sorry for the client (sympathy ≠ compassion)",
        "  • Keeping the client's best interests at the center",
        "  • Avoiding self-serving motivations (proving you're right, getting compliance)",
        "  EXAMPLE: Prioritizing what the client needs to hear over what you want to say",
        "",
        "EVOCATION - Drawing Out What Is Already There:",
        "  • The client already has the arguments for change inside them",
        "  • Therapist's role: Create conditions for these to emerge",
        "  • WRONG: Installing motivation from outside ('You should...')",
        "  • CORRECT: Evoking motivation from within ('What would you like to be different?')",
        "  • Like midwifery: helping deliver what is already developing",
        "  • Miller (1995): 'The words which come out of a person's mouth are quite persuasive to that person'",
    ],
    notes="Evocation is perhaps the most distinctive element of MET spirit. Traditional therapy installs new beliefs or behaviors from outside. MET assumes the client already has the seeds of motivation within them - perhaps buried under ambivalence, fear, or habit. The therapist's skill lies in creating conditions where these seeds can germinate and grow. Every time the client voices a reason for change, that motivation grows stronger (self-perception theory). The therapist's job is to ask the right questions and reflect strategically.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Apply compassion and evocation principles in clinical dialogue",
    takeaway="Evocation means drawing out the client's own motivation - the answers are already within them."
)


# ============================================================
# SECTION 6: CORE PRINCIPLES (Slides 31-35)
# ============================================================
add_section_divider(6, "Core Principles of MET")

add_slide(
    "Express Empathy",
    [
        "DEFINITION: Skillful reflective listening that seeks to understand the client's perspective",
        "NOT sympathy ('I feel sorry for you') but UNDERSTANDING ('I see what you mean')",
        "Rogers' empathy: Sensing the client's private world as if it were your own",
        "KEY COMPONENTS:",
        "  • Active listening with full attention",
        "  • Reflecting back what is heard (verbal and emotional content)",
        "  • Acceptance of ambivalence as NORMAL, not pathological",
        "  • Non-judgmental stance throughout",
        "CLINICAL DIALOGUE:",
        "  Client: 'I don't think I drink more than my friends'",
        "  WRONG: 'That's denial. Your friends probably drink too much too'",
        "  CORRECT: 'From your perspective, your drinking seems pretty normal compared to people around you'",
        "RESEARCH: Valle (1981) - therapist empathy predicted 2-year outcomes",
    ],
    notes="Empathy in MET is not passive listening. It is an ACTIVE skill where the therapist continuously generates hypotheses about what the client means and reflects these back for verification. The MET manual emphasizes that empathy serves multiple functions: (1) reduces resistance, (2) encourages continued exploration, (3) builds therapeutic alliance, (4) clarifies meaning, and (5) can be used to SELECTIVELY reinforce certain ideas. This last point is strategic - you can reflect change talk more than sustain talk.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Valle, S.K. (1981). Journal of Studies on Alcohol, 42, 783-790. | Rogers, C.R. (1957). Journal of Consulting Psychology, 21, 95-103. | Truax, C.B., & Carkhuff, R.R. (1967). Toward Effective Counseling and Psychotherapy. Aldine.",
    learning_obj="Master the principle of expressing empathy through accurate reflective listening",
    takeaway="Empathy is an active skill of reflecting understanding; it reduces resistance and builds alliance."
)

add_slide(
    "Develop Discrepancy",
    [
        "DEFINITION: Helping clients see the gap between present behavior and important goals/values",
        "MECHANISM: Cognitive dissonance creates motivation to resolve the discrepancy",
        "KEY PRINCIPLE: The CLIENT should present the arguments for change, not the therapist",
        "HOW TO DEVELOP DISCREPANCY:",
        "  • Explore personal values: 'What's most important to you in life?'",
        "  • Compare behavior with values: 'How does your drinking fit with being a good father?'",
        "  • Use assessment feedback: Show normative comparisons",
        "  • Looking back: 'What was life like before heavy drinking?'",
        "  • Looking forward: 'Where do you see yourself in 5 years if nothing changes?'",
        "CAUTION: Too much discrepancy too fast → defensiveness",
        "CLINICAL EXAMPLE:",
        "  Client values: family, career, health",
        "  Current behavior: Missing children's events, poor performance reviews, elevated liver enzymes",
        "  Therapist: 'On one hand you deeply value your family, and on the other hand...'",
    ],
    notes="This principle directly applies Festinger's cognitive dissonance theory. The crucial point is that discrepancy must be developed from the CLIENT's own values and goals - not from the therapist's agenda. The Personal Feedback Report in MET Session 1 is designed to create discrepancy by comparing the client's drinking with normative data. When clients see they drink more than 95% of the population, this creates powerful discrepancy. But it must be presented neutrally - 'What do you make of this?' rather than 'See how bad you are?'",
    refs="Miller, W.R., et al. (1995). MET Manual. | Festinger, L. (1957). A Theory of Cognitive Dissonance. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Apply the develop discrepancy principle using client values and assessment feedback",
    takeaway="Discrepancy between values and behavior creates internal motivation; let the client articulate it."
)


add_slide(
    "Roll with Resistance",
    [
        "DEFINITION: Avoid arguing; use the client's momentum rather than opposing it directly",
        "Resistance is a SIGNAL that the therapist should change approach, not push harder",
        "Miller (1995): 'Resistance is an interpersonal phenomenon - therapist behavior influences it'",
        "STRATEGIES FOR ROLLING WITH RESISTANCE:",
        "  • Simple reflection: Acknowledge what client is saying",
        "  • Amplified reflection: Slightly overstate to elicit the other side",
        "  • Double-sided reflection: 'On one hand... and on the other...'",
        "  • Shifting focus: Move away from the stuck point",
        "  • Reframing: Offer a new interpretation of the information",
        "  • Agreeing with a twist: Agree but add a new dimension",
        "EXAMPLE OF AMPLIFIED REFLECTION:",
        "  Client: 'I can handle my drinking just fine'",
        "  Therapist: 'So there's absolutely nothing about your drinking that concerns you at all'",
        "  Client: 'Well... I wouldn't say NOTHING...'",
    ],
    notes="Rolling with resistance is counterintuitive for most therapists. The natural instinct when a client resists is to push harder - but this creates an adversarial dynamic. Instead, MET uses the client's own energy. Like a martial arts practitioner who uses the opponent's momentum, the MET therapist flows with resistance rather than opposing it. Amplified reflection works by overstating the client's position slightly, which often causes them to back away from the extreme and acknowledge the other side. Double-sided reflection captures both sides of ambivalence simultaneously.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Murthy, P. (2008). NIMHANS Manual: DARES approach.",
    learning_obj="Demonstrate strategies for rolling with resistance rather than opposing it",
    takeaway="Resistance signals a need to change approach; rolling with it reduces defensiveness and opens exploration."
)

add_slide(
    "Support Self-Efficacy",
    [
        "DEFINITION: Enhancing the client's belief in their ability to change",
        "Self-efficacy is the CONFIDENCE component: 'I CAN do this'",
        "Without self-efficacy, even highly motivated clients won't attempt change",
        "STRATEGIES TO SUPPORT SELF-EFFICACY:",
        "  • Affirm past successes: 'You quit for 3 months last year - you've done it before'",
        "  • Highlight strengths: 'You've shown real courage in coming here today'",
        "  • Normalize struggle: 'Many people need several attempts - this is normal'",
        "  • Offer menu of options: Multiple paths = greater sense of possibility",
        "  • Attribution to client: 'YOU did that, not the program'",
        "  • Express confidence: 'Based on what I know about you, I believe you can do this'",
        "RESEARCH: Self-efficacy at intake predicts outcomes across addiction treatments",
        "NOTE: Hope is therapeutic - genuine belief in the client's potential communicates powerfully",
    ],
    notes="Self-efficacy is Bandura's concept directly applied in MET. The manual explicitly states that the therapist must convey confidence in the client's ability to change. This is done both explicitly (verbal encouragement) and implicitly (offering choices implies change is possible). Research shows that self-efficacy at intake is one of the strongest predictors of treatment outcome. Therefore, anything that builds efficacy is therapeutic. The therapist's own belief in the client matters enormously - clients can sense genuine vs. fake optimism.",
    refs="Bandura, A. (1982). American Psychologist, 37, 122-147. | Miller, W.R., et al. (1995). MET Manual. | DiClemente, C.C., et al. (1995). Alcoholism: Clinical and Experimental Research, 19, 1062-1071.",
    learning_obj="Implement strategies to build client self-efficacy within MET sessions",
    takeaway="Self-efficacy is the bridge between motivation and action; therapists must actively nurture it."
)


# ============================================================
# SECTION 7: ESSENTIAL COMMUNICATION SKILLS - OARS (Slides 36-42)
# ============================================================
add_section_divider(7, "Essential Communication Skills: OARS")

add_slide(
    "OARS: The Core Skills of MET",
    [
        "OARS = Open Questions + Affirmations + Reflective Listening + Summaries",
        "These are the MICRO-SKILLS that operationalize the MET spirit",
        "Every MET session relies heavily on OARS throughout",
        "",
        "O - Open Questions: Invite exploration, cannot be answered with yes/no",
        "A - Affirmations: Acknowledge client strengths and efforts",
        "R - Reflective Listening: Mirror back what client says with added depth",
        "S - Summaries: Collect and present key themes",
        "",
        "NIMHANS Manual: 'OARS and eliciting change talk are the 5 specific skills in MI'",
        "Ratio guideline: Reflections should outnumber questions 2:1",
        "These skills work together synergistically",
    ],
    notes="OARS are the building blocks of MET communication. They are not unique to MET but are used with specific intention and strategic direction in MET. The therapist uses open questions to explore, affirmations to build confidence, reflections to deepen understanding and selectively reinforce, and summaries to organize and transition. The NIMHANS manual specifically identifies OARS plus eliciting change talk as the five core skills. Mastery of these skills requires extensive practice.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Murthy, P. (2008). NIMHANS Manual, pp. 22-23.",
    learning_obj="Identify and define the four OARS communication skills used in MET",
    takeaway="OARS (Open questions, Affirmations, Reflections, Summaries) are the operational micro-skills of MET."
)

add_slide(
    "Open Questions: Inviting Exploration",
    [
        "DEFINITION: Questions that cannot be answered with a simple yes/no",
        "PURPOSE: Invite client to do the talking; explore their perspective",
        "TYPES OF OPEN QUESTIONS IN MET:",
        "  • Exploring: 'Tell me about your drinking' / 'What concerns you?'",
        "  • Evoking: 'What would you like to be different?' / 'Why might you make this change?'",
        "  • Planning: 'What do you think you might do?' / 'How could you go about it?'",
        "EXAMPLES FROM MET MANUAL:",
        "  'What are your worries about drinking?'",
        "  'Tell me what you've noticed about your drinking. How has it changed over time?'",
        "  'What makes you think that perhaps you need to make a change?'",
        "  'What have other people told you about your drinking?'",
        "POOR EXAMPLES (Closed questions):",
        "  'Do you think you drink too much?' / 'Are you going to stop?'",
        "PRACTICE: Convert closed → open questions",
    ],
    notes="Open questions are the engine that keeps the client talking. In MET, the therapist's goal is to get the CLIENT to articulate reasons for change. Open questions create space for this. The MET manual provides specific examples: 'Tell me what concerns you about your drinking. What are your worries?' Notice these are genuinely exploratory - they don't presume the answer. A key practice skill is converting closed questions to open ones. Instead of 'Do you drink every day?' try 'Tell me about a typical week with your drinking.'",
    refs="Miller, W.R., et al. (1995). MET Manual, pp. 13-15. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Formulate effective open questions that invite exploration of ambivalence",
    takeaway="Open questions create space for clients to explore their own motivations; they are the engine of change talk."
)


add_slide(
    "Affirmations: Recognizing Strengths",
    [
        "DEFINITION: Statements that recognize client strengths, efforts, and positive qualities",
        "PURPOSE: Build self-efficacy, strengthen therapeutic alliance, reduce defensiveness",
        "CHARACTERISTICS OF GOOD AFFIRMATIONS:",
        "  • Genuine and specific (not generic praise)",
        "  • Focus on efforts and character, not just outcomes",
        "  • Acknowledge difficulty of the process",
        "  • Attribute positive qualities to the client",
        "EXAMPLES:",
        "  'It took real courage to come here today'",
        "  'You clearly care deeply about your family'",
        "  'The fact that you're thinking about this shows real strength'",
        "  'You managed 3 months without drinking - that shows you CAN do this'",
        "POOR EXAMPLES:",
        "  'Good job!' (too vague) / 'I'm proud of you' (patronizing, about therapist)",
        "KEY: Affirmations are about the CLIENT, not about the therapist's approval",
    ],
    notes="Affirmations differ from praise. Praise often comes from a one-up position ('I'm proud of you' implies the therapist is the judge). Affirmations recognize qualities IN the client ('You showed real determination'). They should be genuine - clients can detect insincerity. They should be specific - 'You managed to stay sober for your daughter's birthday despite pressure from your friends' is more powerful than 'Good job staying sober.' Affirmations build self-efficacy by helping clients see their own strengths.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Miller, W.R., et al. (1995). MET Manual. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Distinguish effective affirmations from generic praise and practice delivery",
    takeaway="Effective affirmations are genuine, specific, and attribute positive qualities to the client themselves."
)

add_slide(
    "Reflective Listening: The Core Skill",
    [
        "DEFINITION: Statements that reflect back what the client has said with added meaning",
        "Miller (1995): 'Neither easy nor to be done poorly - requires continuous alert tracking'",
        "TYPES OF REFLECTIONS:",
        "  1. Simple Reflection: Repeats or slightly rephrases what was said",
        "  2. Complex Reflection: Adds meaning, feeling, or implication beyond what was stated",
        "  3. Double-Sided Reflection: Captures both sides of ambivalence",
        "  4. Amplified Reflection: Overstates slightly to elicit the other side",
        "EXAMPLE FROM MET MANUAL:",
        "  Client: 'I'm not sure I'm concerned, but I do wonder if I'm drinking too much'",
        "  Simple: 'You wonder about your drinking'",
        "  Complex: 'Part of you suspects something needs to change'",
        "  Double-sided: 'You're not very worried, AND you wonder if it's too much'",
        "STRATEGIC USE: Selectively reflect change talk more than sustain talk",
        "Guideline: 2 reflections for every 1 question (2:1 ratio)",
    ],
    notes="Reflective listening is THE most important skill in MET. The manual devotes significant attention to it. The therapist listens carefully, formulates a guess about what the client means, and offers it back as a statement (not a question). This serves multiple purposes: (1) shows understanding, (2) reduces resistance, (3) keeps the client talking, (4) can be used strategically to reinforce change talk. The strategic aspect is crucial - by choosing WHICH elements to reflect, the therapist can guide the conversation toward change without being directive.",
    refs="Miller, W.R., et al. (1995). MET Manual, pp. 16-18. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Egan, G. (1982). The Skilled Helper. Brooks/Cole.",
    learning_obj="Differentiate types of reflections and practice strategic reflective listening",
    takeaway="Reflective listening is MET's most important skill - it shows understanding and strategically reinforces change talk."
)


add_slide(
    "Reflective Listening: Clinical Dialogue Examples",
    [
        "FROM THE MET MANUAL - Extended Example:",
        "  T: 'What else concerns you about your drinking?'",
        "  C: 'I'm not sure I'm concerned, but I wonder if I'm drinking too much'",
        "  T: 'Too much for...' [reflection as prompt]",
        "  C: 'For my own good. Sometimes I wake up feeling awful, can't think straight'",
        "  T: 'It messes up your thinking, your concentration' [simple reflection]",
        "  C: 'Yes, and sometimes I have trouble remembering things'",
        "  T: 'And you wonder if that might be because you're drinking too much' [complex]",
        "  C: 'Well, I know it is sometimes'",
        "  T: 'You're pretty sure about that. But maybe there's more' [reflection + prompt]",
        "  C: 'Even when I'm not drinking, I mix things up, and I wonder about that'",
        "  T: 'Wonder if...' [continuing the thought]",
        "  C: 'If alcohol's pickling my brain, I guess'",
        "Notice: Therapist uses ONLY reflections to move conversation forward!",
    ],
    notes="This extended dialogue from the MET manual demonstrates how reflective listening alone can build momentum toward change talk. Notice that the therapist asks only ONE question at the beginning, then uses nothing but reflections. Each reflection gently moves the client toward deeper acknowledgment of the problem. The client moves from 'I'm not sure I'm concerned' to 'alcohol is pickling my brain' in just a few exchanges. This is the power of skilled reflective listening - it creates a safe space for the client to explore without feeling pushed.",
    refs="Miller, W.R., et al. (1995). MET Manual, pp. 16-17 (direct clinical example from the manual).",
    learning_obj="Observe skilled reflective listening in an extended clinical dialogue",
    takeaway="Skilled reflection alone can move clients from ambivalence to recognition without a single piece of advice."
)

add_slide(
    "Summaries: Collecting and Transitioning",
    [
        "DEFINITION: Longer reflections that collect multiple things the client has said",
        "TYPES OF SUMMARIES:",
        "  1. Collecting Summary: Gathers several change talk statements together",
        "  2. Linking Summary: Connects current statements with earlier ones",
        "  3. Transitional Summary: Wraps up a topic and shifts direction",
        "PURPOSE:",
        "  • Shows you've been listening carefully",
        "  • Reinforces change talk by collecting it in one place",
        "  • Organizes the session",
        "  • Creates natural transition points",
        "EXAMPLE (Collecting Summary):",
        "  'Let me see if I've got this right. You're concerned about your memory problems,",
        "  you've noticed your wife is increasingly worried, your work performance has slipped,",
        "  and your doctor has flagged your liver results. Did I miss anything?'",
        "STRATEGIC NOTE: Include more change talk items than sustain talk items in summaries",
    ],
    notes="Summaries serve both relational and strategic functions. Relationally, they show the client you've been tracking carefully. Strategically, you can choose WHAT to include. A well-crafted summary will collect the client's change talk statements and present them back as a bouquet - hearing all their reasons for change in one place can be very powerful. The transitional summary is particularly useful at session transitions and when moving from exploring to planning. Always end summaries with 'What else?' or 'Did I miss anything?' to invite correction and further exploration.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Miller, W.R., et al. (1995). MET Manual. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Construct different types of summaries for different clinical purposes",
    takeaway="Strategic summaries collect change talk in one place, amplifying its impact on the client's motivation."
)


# ============================================================
# SECTION 8: CHANGE TALK (Slides 43-46)
# ============================================================
add_section_divider(8, "Change Talk: The Language of Motivation")

add_slide(
    "Change Talk: Definition and Importance",
    [
        "DEFINITION: Any client speech that favors movement toward change",
        "Miller (1995): 'Self-motivational statements' - speech that argues FOR change",
        "WHY IT MATTERS:",
        "  • Self-perception theory: Saying it makes you believe it",
        "  • Client change talk in session PREDICTS actual behavior change",
        "  • Research: Amount of change talk predicts outcomes (Amrhein et al., 2003)",
        "  • Conversely: Sustain talk predicts continued use",
        "THE THERAPIST'S ROLE:",
        "  • Evoke change talk through strategic questions",
        "  • Recognize change talk when it occurs",
        "  • Reinforce/strengthen change talk through reflection",
        "  • Summarize change talk to amplify its effect",
        "MET MANUAL: 'The ME therapist seeks to elicit self-motivating statements'",
        "  These include: acknowledging problems, expressing concern, and willingness to change",
    ],
    notes="Change talk is the therapeutic target in MET. The therapist's primary job is to elicit, recognize, and reinforce change talk. Research by Amrhein et al. (2003) showed that client commitment language during MI sessions predicted actual substance use outcomes. Moyers et al. (2007) demonstrated the causal chain: therapist MI-consistent behavior → client change talk → behavior change. This means the therapist has a direct path to influencing outcomes through the technical skill of evoking change talk.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Amrhein, P.C., et al. (2003). Journal of Consulting and Clinical Psychology, 71, 862-878. | Moyers, T.B., et al. (2007). Journal of Consulting and Clinical Psychology, 75, 790-798.",
    learning_obj="Define change talk and understand its predictive relationship to outcomes",
    takeaway="Change talk predicts behavior change; the therapist's primary task is to evoke, recognize, and reinforce it."
)

add_slide(
    "Types of Change Talk: DARN-CAT",
    [
        "PREPARATORY CHANGE TALK (DARN) - Building momentum:",
        "  D - Desire: 'I want to...' / 'I wish I could...' / 'I'd like to...'",
        "  A - Ability: 'I can...' / 'I could...' / 'I'm able to...'",
        "  R - Reasons: 'I would feel better if...' / 'My health would improve...'",
        "  N - Need: 'I need to...' / 'I have to...' / 'I must...'",
        "",
        "MOBILIZING CHANGE TALK (CAT) - Stronger commitment signals:",
        "  C - Commitment: 'I will...' / 'I'm going to...' / 'I promise...'",
        "  A - Activation: 'I'm ready to...' / 'I'm prepared to...' / 'I'm willing to...'",
        "  T - Taking Steps: 'I actually did...' / 'This week I...' / 'I tried...'",
        "",
        "CLINICAL SIGNIFICANCE:",
        "  • DARN = Preparatory (earlier in process)",
        "  • CAT = Mobilizing (closer to action)",
        "  • Commitment language is STRONGEST predictor of behavior change (Amrhein et al., 2003)",
    ],
    notes="The DARN-CAT framework helps therapists recognize different strengths of change talk. Preparatory change talk (DARN) signals that the client is moving toward change but hasn't committed yet. Mobilizing change talk (CAT) signals readiness for action. Commitment language ('I will') is the strongest predictor. Therapists should track the shift from DARN to CAT as an indicator that the client is ready for planning. Premature focus on planning when the client is still in DARN territory can backfire.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Amrhein, P.C., et al. (2003). JCCP, 71, 862-878. | Moyers, T.B., et al. (2009). Journal of Substance Abuse Treatment, 36, 101-109.",
    learning_obj="Identify and classify the seven types of change talk using the DARN-CAT framework",
    takeaway="DARN-CAT: Desire, Ability, Reasons, Need (preparatory) → Commitment, Activation, Taking Steps (mobilizing)."
)


add_slide(
    "Evoking Change Talk: Strategies and Exercises",
    [
        "STRATEGIES TO EVOKE CHANGE TALK (from MET Manual):",
        "  1. Ask evocative questions: 'What concerns you about your drinking?'",
        "  2. Ask for elaboration: 'Tell me more about that' / 'Give me an example'",
        "  3. Ask for extremes: 'What's the worst that could happen?'",
        "  4. Looking back: 'What was life like before the drinking got heavy?'",
        "  5. Looking forward: 'If you made this change, how would life be different?'",
        "  6. Querying extremes: 'What are your worst fears if you DON'T change?'",
        "  7. Using importance ruler: 'On a scale of 0-10, how important is this change?'",
        "  8. Exploring goals and values: 'What matters most to you?'",
        "  9. Gentle paradox: 'I'm not sure you're motivated enough...' (evokes opposition)",
        "FROM MET MANUAL - Gentle Paradox Example:",
        "  Therapist: 'I'm not sure from what you've told me that you're motivated enough",
        "  to carry through. Do you think we should go ahead?'",
        "  [This often evokes client arguing FOR their own motivation]",
    ],
    notes="The MET manual provides specific techniques for evoking self-motivational statements. The gentle paradox technique is particularly interesting - by subtly taking the side of status quo, the therapist evokes the client to argue for change. This must be used carefully and genuinely, not manipulatively. The manual provides the example: 'You haven't convinced me yet that you are seriously concerned... Is that all you're concerned about?' This table-turning approach works because of psychological reactance - people resist having their freedom limited.",
    refs="Miller, W.R., et al. (1995). MET Manual, pp. 13-16. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Practice specific strategies for evoking change talk in clinical sessions",
    takeaway="Strategic questions, looking back/forward, rulers, and gentle paradox all evoke change talk from clients."
)

# ============================================================
# SECTION 9: SUSTAIN TALK (Slide 47)
# ============================================================
add_section_divider(9, "Sustain Talk")

add_slide(
    "Sustain Talk: Recognition and Response",
    [
        "DEFINITION: Any client speech that favors maintaining status quo (not changing)",
        "EXAMPLES OF SUSTAIN TALK:",
        "  • 'I don't think my drinking is that bad' (minimizing)",
        "  • 'I enjoy drinking and don't want to stop' (desire for status quo)",
        "  • 'I've tried before and it never works' (inability)",
        "  • 'All my friends drink' (reasons to continue)",
        "  • 'I don't need to change' (no need)",
        "THERAPIST RESPONSE TO SUSTAIN TALK:",
        "  • Do NOT argue against it (this creates more sustain talk)",
        "  • Reflect it simply (without dwelling): 'You enjoy the social aspects'",
        "  • Use double-sided reflection: 'You enjoy drinking, AND you're worried about your health'",
        "  • Shift focus: Move to a different topic",
        "  • DO NOT reflect sustain talk in detail or elaborate on it",
        "RESEARCH: High sustain talk + low change talk predicts poor outcomes",
        "RATIO: Aim for change talk to outweigh sustain talk over the session",
    ],
    notes="Sustain talk is the opposite of change talk - it's the voice of status quo. It is NORMAL, especially early in therapy. The key principle is: don't feed it. When you argue against sustain talk, you paradoxically strengthen it (the client hears themselves defending their behavior). Instead, acknowledge it briefly and move on. Or use a double-sided reflection that pairs it with change talk. Research shows the ratio of change talk to sustain talk is what matters for outcomes. The therapist's job is to tilt this ratio toward change talk.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Moyers, T.B., et al. (2007). JCCP, 75, 790-798. | Magill, M., et al. (2014). Journal of Substance Abuse Treatment, 46, 685-697.",
    learning_obj="Recognize sustain talk and respond in ways that do not strengthen resistance",
    takeaway="Sustain talk is normal; acknowledge it briefly without elaboration, then steer toward change talk."
)


# ============================================================
# SECTION 10: RESISTANCE / DISCORD (Slide 48-49)
# ============================================================
add_section_divider(10, "Resistance and Discord")

add_slide(
    "Understanding Resistance (Discord) in MET",
    [
        "REFRAMING RESISTANCE: Miller & Rollnick (2013) replaced 'resistance' with 'discord'",
        "Discord = signals of disharmony in the therapeutic relationship",
        "COMMON FORMS OF DISCORD:",
        "  • Arguing: Challenging, discounting, being hostile",
        "  • Interrupting: Cutting off the therapist",
        "  • Denying: Blaming others, minimizing, disagreeing, excusing",
        "  • Ignoring: Not paying attention, not responding, changing subject",
        "KEY INSIGHT: Resistance is largely THERAPIST-CREATED",
        "  • Patterson & Forgatch (1985): Teaching/confronting → increased resistance",
        "  • Miller et al. (1993): Amount of client resistance in Session 1 predicted",
        "    drinking outcomes at 12 months",
        "MANAGEMENT: Discord is a signal to CHANGE YOUR APPROACH, not push harder",
        "The 'righting reflex' (desire to fix/advise) triggers discord",
    ],
    notes="This slide reframes how we think about resistance. In MET, resistance is NOT a client characteristic but an interpersonal phenomenon. When resistance increases, it signals that the therapist is pushing too hard, moving too fast, or not aligning with the client's readiness. The therapist should view resistance as useful feedback about the therapeutic interaction, not as evidence of client pathology. Miller's landmark research showed that therapist behavior (confrontation) directly caused client resistance, which then predicted poor outcomes a year later.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Patterson, G.A., & Forgatch, M.S. (1985). JCCP, 53, 846-851. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Reconceptualize resistance as an interpersonal signal rather than a client deficit",
    takeaway="Resistance is a therapist-influenced interpersonal phenomenon; when it increases, change YOUR approach."
)

add_slide(
    "Strategies for Managing Discord",
    [
        "WHEN DISCORD APPEARS, USE THESE STRATEGIES:",
        "  1. Simple Reflection: 'You don't see this as a problem right now'",
        "  2. Amplified Reflection: 'So alcohol causes you absolutely no issues whatsoever'",
        "  3. Double-Sided Reflection: 'You enjoy drinking AND you notice it's affecting sleep'",
        "  4. Shifting Focus: 'Let's set that aside. What brought you here today?'",
        "  5. Reframing: 'Your concern about being labeled shows you think carefully'",
        "  6. Agreeing with a Twist: 'You're right, only YOU can decide...'",
        "  7. Emphasizing Autonomy: 'It's completely up to you what happens here'",
        "  8. Coming Alongside: 'Maybe this isn't the right time for you to change'",
        "WHAT NOT TO DO:",
        "  ✗ Argue / confront / lecture / warn / threaten / label",
        "  ✗ Push harder when you feel resistance",
        "  ✗ Take it personally",
        "  ✗ Interpret it as denial or pathology",
    ],
    notes="Practice each of these responses through role-play. The amplified reflection is counterintuitive - by slightly overstating the client's position, you often get them to back away from the extreme. 'Coming alongside' is also counterintuitive - by agreeing that maybe they don't need to change, you often trigger reactance where the client argues FOR change. Each of these techniques respects the client's autonomy while creating space for the other side of ambivalence to emerge. The key is to stay calm, curious, and non-defensive yourself.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Murthy, P. (2008). NIMHANS Manual: 'Roll with resistance.'",
    learning_obj="Apply specific strategies for managing discord without damaging the therapeutic relationship",
    takeaway="When discord appears, soften your approach rather than pushing harder; use reflections and autonomy emphasis."
)


# ============================================================
# SECTION 11: ASSESSMENT BEFORE MET (Slides 50-53)
# ============================================================
add_section_divider(11, "Assessment Before MET")

add_slide(
    "Pre-Treatment Assessment for MET",
    [
        "MET requires comprehensive assessment BEFORE therapy begins",
        "Assessment serves DUAL purpose: (1) Clinical information, (2) Feedback material",
        "ASSESSMENT BATTERY SHOULD INCLUDE:",
        "  • Substance use history (quantity, frequency, duration, pattern)",
        "  • Readiness to change assessment",
        "  • Psychological assessment (mood, anxiety, personality)",
        "  • Mental Status Examination",
        "  • Physical health assessment (liver function, etc.)",
        "  • Social and occupational functioning",
        "  • Risk assessment (suicide, violence, self-harm)",
        "  • Family history and social support",
        "  • Previous treatment history",
        "  • Neuropsychological screening (if indicated)",
        "Project MATCH: 'Personal Feedback Report' was generated from this assessment",
    ],
    notes="Assessment in MET serves a unique dual purpose. Unlike other therapies where assessment simply informs the therapist, in MET the assessment results are given BACK to the client as personalized feedback. This is a key therapeutic intervention in Session 1. The assessment should be comprehensive enough to generate meaningful normative comparisons. In Project MATCH, the assessment included drinking patterns, consequences, blood alcohol level estimation, liver function, psychological problems, and readiness to change.",
    refs="Miller, W.R., et al. (1995). MET Manual, Appendix A. | Murthy, P. (2008). NIMHANS Manual. | SAMHSA. (2019). Enhancing Motivation for Change in Substance Use Disorder Treatment. TIP 35.",
    learning_obj="Design a comprehensive pre-treatment assessment battery suitable for MET",
    takeaway="MET assessment serves dual purposes: informing the therapist AND providing personalized feedback to the client."
)

add_slide(
    "Readiness, Importance, and Confidence Rulers",
    [
        "READINESS RULER: 'On a scale of 0-10, how ready are you to make a change?'",
        "  Follow-up: 'Why are you at a ___ and not a zero?' (evokes change talk)",
        "  Follow-up: 'What would it take to move from ___ to ___?' (identifies barriers)",
        "",
        "IMPORTANCE RULER: 'How important is it to you to change your drinking? (0-10)'",
        "  Assesses: Value placed on change (motivational component)",
        "  Low score → Need to develop discrepancy",
        "",
        "CONFIDENCE RULER: 'How confident are you that you could change? (0-10)'",
        "  Assesses: Self-efficacy (capability belief)",
        "  Low score → Need to build self-efficacy",
        "",
        "CLINICAL SIGNIFICANCE:",
        "  High importance + Low confidence → Focus on building efficacy",
        "  Low importance + High confidence → Focus on developing discrepancy",
        "  Low both → Start with importance (why), then build confidence (how)",
    ],
    notes="The rulers are elegantly simple yet powerful assessment and intervention tools. The key follow-up question 'Why are you at X and not zero?' is brilliant because it evokes change talk - the client must articulate reasons they ARE somewhat motivated. Asking 'Why aren't you at 10?' would evoke sustain talk. The combination of importance and confidence tells you where to focus: if someone wants to change but doesn't feel able, build efficacy. If they feel able but don't see it as important, develop discrepancy.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Rollnick, S., Mason, P., & Butler, C. (1999). Health Behavior Change. Churchill Livingstone. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Administer and interpret readiness, importance, and confidence rulers",
    takeaway="Rulers assess motivation (importance) and capability (confidence) separately, guiding therapeutic focus."
)


add_slide(
    "Decisional Balance Assessment",
    [
        "DECISIONAL BALANCE: Systematic exploration of pros and cons",
        "Based on Janis & Mann (1977) decision-making theory",
        "",
        "FOUR-QUADRANT MODEL:",
        "  ┌─────────────────────┬─────────────────────┐",
        "  │ PROS of Drinking    │ CONS of Drinking    │",
        "  │ (Benefits of        │ (Costs of           │",
        "  │  status quo)        │  status quo)        │",
        "  ├─────────────────────┼─────────────────────┤",
        "  │ CONS of Changing    │ PROS of Changing    │",
        "  │ (Costs of change)   │ (Benefits of        │",
        "  │                     │  change)            │",
        "  └─────────────────────┴─────────────────────┘",
        "",
        "THERAPEUTIC USE: Explore ALL four quadrants non-judgmentally",
        "  Start with pros of drinking (validates, reduces defensiveness)",
        "  Then explore cons of drinking (develops discrepancy)",
        "  Explore fears about change (acknowledges difficulty)",
        "  End with benefits of change (builds forward momentum)",
    ],
    notes="The decisional balance is both an assessment tool and an intervention. By exploring all four quadrants, the therapist validates the client's ambivalence (both sides are real and understandable). Starting with the benefits of drinking shows you're not going to lecture, which reduces defensiveness. The strategic order matters: ending with benefits of change leaves the client thinking about positive possibilities. The balance can be done on paper as a worksheet or verbally in session.",
    refs="Janis, I.L., & Mann, L. (1977). Decision Making. Free Press. | Miller, W.R., et al. (1995). MET Manual. | Prochaska, J.O., et al. (1994). Cancer, 73, 3191-3197.",
    learning_obj="Conduct a four-quadrant decisional balance assessment in clinical practice",
    takeaway="The decisional balance validates ambivalence while systematically developing discrepancy favoring change."
)

# ============================================================
# SECTION 12: COMPONENTS OF MET - FRAMES (Slides 54-57)
# ============================================================
add_section_divider(12, "Components of MET: The FRAMES Model")

add_slide(
    "FRAMES: The Active Ingredients of Brief Interventions",
    [
        "FRAMES (Miller & Sanchez, 1994) - Core components of effective brief interventions:",
        "",
        "F - Feedback: Personal risk/impairment feedback based on assessment",
        "R - Responsibility: Emphasis on personal responsibility for change",
        "A - Advice: Clear advice to change (when appropriate)",
        "M - Menu: A menu of alternative change options",
        "E - Empathy: Warm, reflective, understanding counseling style",
        "S - Self-efficacy: Optimism about the client's ability to change",
        "",
        "RESEARCH SUPPORT: Bien, Miller & Tonigan (1993) meta-analysis found",
        "  these elements consistently present in effective brief interventions",
        "NIMHANS Manual: FRAMES identified as the essential elements of brief intervention",
        "Each element can be present to varying degrees",
        "Not all elements required in every interaction",
    ],
    notes="FRAMES summarizes what research has found to be the active ingredients across effective brief interventions worldwide. It emerged from a systematic review of what the successful interventions had in common. The NIMHANS manual specifically identifies FRAMES as the core model for brief intervention in Indian settings. In MET, all six elements are systematically incorporated: feedback comes in Session 1 (Personal Feedback Report), responsibility is emphasized throughout, advice is offered with permission, menu provides options, empathy is the communication style, and self-efficacy is actively supported.",
    refs="Miller, W.R., & Sanchez, V.C. (1994). In Howard, G. (Ed.), Issues in Alcohol Use and Misuse. Notre Dame Press. | Bien, T.H., Miller, W.R., & Tonigan, S. (1993). Addiction, 88, 315-336. | Murthy, P. (2008). NIMHANS Manual, pp. 20-21.",
    learning_obj="Identify and explain the six FRAMES components of effective brief interventions",
    takeaway="FRAMES (Feedback, Responsibility, Advice, Menu, Empathy, Self-efficacy) are the proven active ingredients of MET."
)


add_slide(
    "FRAMES: Detailed Explanation (Feedback, Responsibility, Advice)",
    [
        "F - FEEDBACK:",
        "  • Personalized assessment results compared to norms",
        "  • 'Your drinking is higher than 95% of the population'",
        "  • Objective, non-judgmental presentation",
        "  • Creates discrepancy between self-perception and data",
        "",
        "R - RESPONSIBILITY:",
        "  • 'No one can make you change. It's your choice'",
        "  • Emphasizes personal control and autonomy",
        "  • Avoids coercion or external pressure",
        "  • NIMHANS: 'Reinforces sense of personal control and responsibility'",
        "",
        "A - ADVICE (offered with permission):",
        "  • Clear, specific, direct but non-coercive",
        "  • 'Based on these results, I would recommend reducing your drinking'",
        "  • Elicit-Provide-Elicit: Ask permission → Give info → Ask reaction",
        "  • Brief and impactful; not lecturing",
    ],
    notes="Feedback is the distinguishing feature of MET vs. pure MI. The Personal Feedback Report provides objective data that the client cannot easily dismiss. Presenting it neutrally ('Here are your results. What do you make of them?') is more effective than interpreting it for them. Responsibility is crucial - by emphasizing that the choice is theirs, you paradoxically make change more likely (autonomy support). Advice should be brief, clear, and offered with permission ('Would you like to hear my professional opinion?').",
    refs="Miller, W.R., et al. (1995). MET Manual, Appendix A. | Murthy, P. (2008). NIMHANS Manual. | Miller, W.R., & Sanchez, V.C. (1994). In Howard, G. (Ed.), Issues in Alcohol Use and Misuse.",
    learning_obj="Deliver personalized feedback, emphasize responsibility, and offer advice appropriately",
    takeaway="Feedback creates discrepancy, responsibility empowers autonomy, advice is offered briefly with permission."
)

add_slide(
    "FRAMES: Detailed Explanation (Menu, Empathy, Self-Efficacy)",
    [
        "M - MENU OF OPTIONS:",
        "  • Provide multiple strategies/choices for change",
        "  • 'There are several approaches people have found helpful...'",
        "  • Options might include: reduction, abstinence, timing changes, support groups",
        "  • NIMHANS: 'Providing choices reinforces sense of personal control'",
        "  • Examples: diary keeping, guidelines, relapse prevention, counseling, self-help",
        "",
        "E - EMPATHY:",
        "  • Warm, reflective, understanding communication style",
        "  • NOT confrontation, lecturing, or moralizing",
        "  • Communicates respect and genuine concern",
        "  • 'The therapist communicates respect, encourages exploration' (NIMHANS)",
        "",
        "S - SELF-EFFICACY:",
        "  • Communicate confidence in the client's ability to change",
        "  • 'I believe you can do this'",
        "  • Elicit self-efficacy statements from clients",
        "  • Highlight past successes and strengths",
    ],
    notes="The Menu component directly supports autonomy (SDT theory) - by offering choices, you communicate that the client has control AND that multiple paths to change exist (building hope). Empathy is the relational glue that holds everything together - without it, feedback can feel like attack, advice like lecturing, and menu like homework assignment. Self-efficacy is about instilling hope and confidence - not false optimism, but genuine belief based on the client's demonstrated strengths. Together, FRAMES provides a complete brief intervention framework.",
    refs="Murthy, P. (2008). NIMHANS Manual, pp. 20-21. | Miller, W.R., et al. (1995). MET Manual. | Bien, T.H., Miller, W.R., & Tonigan, S. (1993). Addiction, 88, 315-336.",
    learning_obj="Provide a menu of options, demonstrate empathy, and build self-efficacy in clinical practice",
    takeaway="Menu offers choices (autonomy), empathy provides safety, self-efficacy instills hope - together completing FRAMES."
)


# ============================================================
# SECTION 13: STRUCTURE OF MET (Slides 58-60)
# ============================================================
add_section_divider(13, "Structure of MET")

add_slide(
    "Overall Structure: The Four-Session MET Protocol",
    [
        "SESSION 1 (Week 1): Assessment Feedback & Building Motivation",
        "  • Personal Feedback Report delivery",
        "  • Explore reactions to feedback",
        "  • Elicit self-motivational statements",
        "  • Optional: Involve significant other",
        "",
        "SESSION 2 (Week 2): Strengthening Commitment",
        "  • Review progress since Session 1",
        "  • Continue developing discrepancy and change talk",
        "  • Begin consolidating commitment to change",
        "  • Develop change plan (if ready)",
        "",
        "SESSION 3 (Week 6 - Midpoint): Progress Review",
        "  • Reinforce progress / address setbacks",
        "  • Reassess motivation and commitment",
        "",
        "SESSION 4 (Week 12 - Final): Consolidation & Termination",
        "  • Review overall progress",
        "  • Strengthen long-term commitment",
        "  • Relapse prevention planning",
    ],
    notes="The structure of MET in Project MATCH was specifically designed to be brief (4 sessions) compared to CBT and TSF (12 sessions each). Sessions 1 and 2 are close together (1 week apart) to build momentum. Sessions 3 and 4 are spaced further apart (weeks 6 and 12) serving as check-ins and reinforcement. The first two sessions do the heavy motivational lifting; the later sessions consolidate and maintain gains. This structure means the most intensive work happens early when motivation is being built.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Project MATCH Research Group. (1997). Journal of Studies on Alcohol, 58, 7-29.",
    learning_obj="Describe the overall structure and timing of the four MET sessions",
    takeaway="MET uses 4 strategically timed sessions: intensive early work (Sessions 1-2) followed by reinforcement check-ins (Sessions 3-4)."
)

add_slide(
    "Two Phases of MET: Building Motivation & Strengthening Commitment",
    [
        "PHASE 1: BUILDING MOTIVATION FOR CHANGE (Sessions 1-2 primarily)",
        "  Goal: Tip the motivational balance in favor of change",
        "  Strategies: Elicit self-motivational statements, empathic listening,",
        "    personal feedback, develop discrepancy, explore ambivalence",
        "  Client varies from precontemplation → contemplation → preparation",
        "",
        "PHASE 2: STRENGTHENING COMMITMENT TO CHANGE (Sessions 2-4)",
        "  Goal: Consolidate decision and plan for action",
        "  Strategies: Recapitulate, key questions, information & advice (with permission),",
        "    change plan negotiation, goal setting",
        "  Signs of readiness: Decreased resistance, increased change talk,",
        "    questions about change, imagining change, experimenting",
        "",
        "TIMING THE TRANSITION: Move to Phase 2 when client shows readiness signs",
        "  Premature transition → increased resistance (pushing too fast)",
        "  Delayed transition → missed opportunity (client ready but therapist still exploring)",
    ],
    notes="The two-phase structure maps onto the Stages of Change model. Phase 1 corresponds to moving clients from precontemplation through contemplation. Phase 2 corresponds to preparation and action. The MET manual describes the 'motivational seesaw' - one side favoring status quo (benefits of drinking, fears of change) and the other favoring change (benefits of change, costs of drinking). The therapist's job in Phase 1 is to shift this balance. Signs of readiness for Phase 2 include: decreased resistance, more change talk, asking 'how' questions, and imagining a changed future.",
    refs="Miller, W.R., et al. (1995). MET Manual, pp. 13-32. | Miller, W.R., & Rollnick, S. (1991). Motivational interviewing. Guilford Press.",
    learning_obj="Distinguish the two phases of MET and identify signs of readiness to transition",
    takeaway="Phase 1 builds motivation (why change); Phase 2 strengthens commitment (how to change) - timing the transition is key."
)


# ============================================================
# SECTION 14: SESSION-BY-SESSION MANUAL (Slides 61-72)
# ============================================================
add_section_divider(14, "Complete Session-by-Session Manual")

add_slide(
    "Session 1: Opening, Rapport, and Agenda Setting",
    [
        "OBJECTIVES: Establish rapport, deliver feedback, elicit motivation",
        "DURATION: 60-90 minutes",
        "PREPARATION: Personal Feedback Report ready, quiet room, assessment data reviewed",
        "",
        "OPENING (First 10 minutes):",
        "  • Warm welcome; introduce yourself and your role",
        "  • Explain session structure: 'We'll spend time looking at your assessment results'",
        "  • Set collaborative tone: 'I'm interested in hearing YOUR thoughts about this'",
        "  • Address any concerns about the process",
        "",
        "RAPPORT BUILDING:",
        "  • Use open questions about general life (not just drinking)",
        "  • Express genuine interest in the person",
        "  • Affirm their decision to attend",
        "  • 'It took real effort to come here today'",
    ],
    notes="Session 1 is the most structured and critical session in MET. The therapist must balance delivering substantial assessment feedback while maintaining a client-centered, empathic style. The opening sets the tone - if the client feels judged or lectured in the first minutes, the entire therapy may be compromised. Start with genuine warmth and curiosity. The agenda should be set collaboratively, not imposed. Emphasize that this is a conversation, not a lecture.",
    refs="Miller, W.R., et al. (1995). MET Manual, Session 1 protocol. | Project MATCH Research Group. (1997). JOSA, 58, 7-29.",
    learning_obj="Structure the opening of MET Session 1 to establish rapport and set agenda",
    takeaway="Session 1's opening must establish safety, collaboration, and genuine interest before feedback delivery."
)

add_slide(
    "Session 1: Delivering the Personal Feedback Report",
    [
        "THE PERSONAL FEEDBACK REPORT (PFR) - Core of Session 1:",
        "  Give client their written PFR copy; retain yours for reference",
        "  Go through step by step, explaining each item",
        "  Compare client's scores with normative data",
        "",
        "PFR CONTENT (from MET Manual):",
        "  • Drinking patterns (quantity/frequency vs. population norms)",
        "  • Peak BAC levels (estimated)",
        "  • Alcohol-related negative consequences",
        "  • Risk factors for alcohol problems",
        "  • Neuropsychological test results (if applicable)",
        "  • Physiological indicators (liver function, health markers)",
        "  • Comparison with general population percentiles",
        "",
        "DELIVERY STYLE: Neutral, non-judgmental, curious",
        "  'Here are your results. What do you make of this?'",
        "  NOT: 'These results show you clearly have a serious problem'",
        "MONITOR client reactions continuously; use reflective listening",
    ],
    notes="The PFR delivery is the heart of Session 1. The manual is very specific about HOW to deliver it: neutrally, item by item, allowing the client to react at each step. The therapist monitors reactions and uses reflective listening. When clients express concern ('Wow, I'm drinking more than I realized'), reflect and amplify ('It looks quite high to you'). When clients minimize ('I don't see how it's affecting me'), reflect without arguing ('This isn't what you expected'). The feedback creates discrepancy naturally - the data speaks for itself when presented in normative context.",
    refs="Miller, W.R., et al. (1995). MET Manual, Appendix A (complete PFR protocol). | Miller, W.R., & Sovereign, R.G. (1989). In Loberg, T., et al. (Eds.), Addictive Behaviors. Swets & Zeitlinger.",
    learning_obj="Deliver personalized assessment feedback in a motivationally enhancing style",
    takeaway="Deliver feedback neutrally, compare with norms, monitor reactions, and reflect - let the data create discrepancy."
)


add_slide(
    "Session 1: Clinical Dialogue During Feedback",
    [
        "DIALOGUE EXAMPLES FROM MET MANUAL:",
        "",
        "Client: 'Wow! I'm drinking a lot more than I realized'",
        "Therapist: 'It looks awfully high to you' [reflects surprise]",
        "",
        "Client: 'I can't believe it. I don't see how my drinking can be affecting me that much'",
        "Therapist: 'This isn't what you expected to hear' [reflects disbelief without arguing]",
        "",
        "Client: 'No, I don't really drink that much more than other people'",
        "Therapist: 'So this is confusing to you. It seems like you drink about the same",
        "  as your friends, yet here are the results' [validates AND maintains data]",
        "",
        "Client: 'This gives me a lot to think about'",
        "Therapist: 'A lot of reasons to think about making a change' [strengthens change talk]",
        "",
        "Client: 'More bad news!'",
        "Therapist: 'This is pretty difficult for you to hear' [empathic reflection]",
    ],
    notes="These dialogue examples come directly from the MET manual. They illustrate how the therapist responds to EVERY reaction with reflective listening, never with confrontation or interpretation. Notice the range of client responses - surprise, disbelief, minimization, contemplation - and how the therapist uses the same basic skill (reflection) differently for each. When the client is surprised, reflect the surprise. When they minimize, reflect the confusion without arguing. When they show concern, reflect and strengthen it. This is the art of selective reflection.",
    refs="Miller, W.R., et al. (1995). MET Manual, pp. 19-20 (direct quotes from clinical examples).",
    learning_obj="Respond to various client reactions during feedback delivery using reflective listening",
    takeaway="Respond to every client reaction with reflection - surprise, denial, concern all get empathic reflection."
)

add_slide(
    "Session 1: Involving a Significant Other",
    [
        "MET OPTIONALLY involves a spouse/partner/family member in Session 1",
        "PURPOSE: Provides additional perspective; increases accountability; builds support",
        "",
        "SIGNIFICANT OTHER (SO) GUIDELINES:",
        "  • Invite the SO to share their observations",
        "  • Use reflective listening with the SO too",
        "  • Reframe SO statements to highlight caring/concern",
        "",
        "DIALOGUE EXAMPLES FROM MANUAL:",
        "  Wife: 'I always thought he was drinking too much'",
        "  T: 'You've been worried about him for quite a while' [reflects caring]",
        "",
        "  Husband: (weeping) 'I've told you to quit drinking!'",
        "  T: 'You really care about her a lot. It's hard to hear these results' [reframes]",
        "",
        "  Friend: 'I never really thought he drank that much!'",
        "  T: 'This is taking you by surprise. (To client:) Does this surprise you too?'",
        "",
        "CAUTION: Maintain balance; don't let SO become confrontational",
    ],
    notes="The involvement of a significant other is optional but recommended when feasible. The MET manual provides specific guidance on how to manage the three-way interaction. The therapist must: (1) maintain the client as the focus, (2) use SO information supportively, (3) reframe SO expressions of anger or frustration into caring statements, and (4) prevent the session from becoming confrontational. After reflecting an SO statement, the manual recommends turning to the client and asking for their perceptions, then reflecting self-motivational elements.",
    refs="Miller, W.R., et al. (1995). MET Manual, pp. 19-20, 35-42 (Involving Significant Others section).",
    learning_obj="Skillfully involve a significant other in MET feedback sessions",
    takeaway="Involve significant others supportively; reframe their concerns as caring; maintain the client as central focus."
)


add_slide(
    "Session 2: Exploring Ambivalence and Strengthening Commitment",
    [
        "OBJECTIVES: Deepen exploration, strengthen commitment, develop change plan",
        "TIMING: 1 week after Session 1",
        "DURATION: 60 minutes",
        "",
        "SESSION 2 STRUCTURE:",
        "  1. OPENING (10 min): Review since last session, what client has been thinking",
        "  2. EXPLORATION (20 min): Continue exploring ambivalence",
        "     • Decisional balance: Pros/cons of drinking AND changing",
        "     • Values clarification: What matters most?",
        "     • Discrepancy development: Values vs. current behavior",
        "  3. COMMITMENT (20 min): If ready, consolidate decision",
        "     • Key question: 'What do you think you'll do?'",
        "     • Change plan worksheet",
        "  4. CLOSING (10 min): Summarize, homework, next appointment",
        "",
        "IF CLIENT NOT READY: Continue Phase 1 strategies; don't push",
        "HOMEWORK: Self-monitoring, change plan development, values card sort",
    ],
    notes="Session 2 is the bridge between building motivation (Phase 1) and strengthening commitment (Phase 2). The therapist should start by asking what the client has been thinking since the feedback session. Often clients have been reflecting and may arrive with new change talk. If the client shows readiness signs (decreased resistance, questions about change, imagining a different future), the therapist can begin transitioning to commitment and planning. If not ready, continue Phase 1 strategies without pressure.",
    refs="Miller, W.R., et al. (1995). MET Manual, Session 2 protocol. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Structure and conduct MET Session 2 with appropriate phase transitions",
    takeaway="Session 2 bridges motivation-building to commitment-strengthening based on client readiness signals."
)

add_slide(
    "Session 3: Midpoint Review and Reinforcement",
    [
        "OBJECTIVES: Reinforce progress, address setbacks, reassess commitment",
        "TIMING: Week 6 (midpoint of 12-week treatment period)",
        "DURATION: 60 minutes",
        "",
        "SESSION 3 STRUCTURE:",
        "  1. CHECK-IN (15 min): How has it been going?",
        "     • Review drinking behavior since Session 2",
        "     • Celebrate any positive changes (affirm!)",
        "  2. PROGRESS REVIEW (20 min):",
        "     • If making progress: Reinforce, affirm, explore what's working",
        "     • If struggling: Express empathy, explore barriers, problem-solve",
        "     • If relapsed: Normalize, explore triggers, rebuild motivation",
        "  3. REASSESSMENT (15 min):",
        "     • Re-evaluate change plan; modify if needed",
        "     • Address new barriers or challenges",
        "  4. FORWARD PLANNING (10 min):",
        "     • Set goals for next 6 weeks",
        "     • Plan for high-risk situations",
        "THERAPIST STANCE: Maintain all MET principles; do not become directive/prescriptive",
    ],
    notes="Session 3 serves as a midpoint check-in. By week 6, the client has had time to implement changes (or not). The therapist's role is to reinforce any progress (however small), address setbacks with empathy rather than disappointment, and reassess the change plan. If the client has been successful, the session celebrates this and plans for maintenance. If struggling, it returns to Phase 1 principles. The key is maintaining the MET spirit throughout - even when discussing setbacks, the therapist uses empathy and supports efficacy.",
    refs="Miller, W.R., et al. (1995). MET Manual, Session 3 protocol. | Marlatt, G.A., & Gordon, J.R. (1985). Relapse Prevention. Guilford Press.",
    learning_obj="Conduct a midpoint MET session that addresses both progress and setbacks",
    takeaway="Session 3 reinforces progress, normalizes setbacks, and adjusts the change plan based on real-world experience."
)


add_slide(
    "Session 4: Consolidation, Maintenance, and Termination",
    [
        "OBJECTIVES: Consolidate gains, plan for maintenance, prepare for termination",
        "TIMING: Week 12 (final session)",
        "DURATION: 60 minutes",
        "",
        "SESSION 4 STRUCTURE:",
        "  1. COMPREHENSIVE REVIEW (15 min):",
        "     • Review entire treatment period",
        "     • Highlight changes and growth",
        "     • Affirm commitment and effort",
        "  2. MAINTENANCE PLANNING (20 min):",
        "     • Identify ongoing high-risk situations",
        "     • Develop specific coping strategies",
        "     • Establish support systems",
        "     • Plan for holidays, social events, triggers",
        "  3. RELAPSE PREVENTION (15 min):",
        "     • Normalize possibility of slips",
        "     • Distinguish lapse from relapse",
        "     • Emergency plan: What to do if you slip",
        "  4. TERMINATION (10 min):",
        "     • Express confidence in client's future",
        "     • Leave door open for booster sessions",
        "     • Final summary of strengths and resources",
        "     • Provide written relapse prevention plan",
    ],
    notes="Session 4 consolidates the work of the entire therapy. The therapist reviews progress, highlights the client's own agency in making changes, and plans for the future. Relapse prevention is incorporated but not in a lecturing way - rather, the therapist explores with the client what situations might be challenging and how they could handle them. The termination should be positive and confidence-building. Research shows that leaving the door open for return contact reduces relapse anxiety. Booster sessions can be offered if appropriate.",
    refs="Miller, W.R., et al. (1995). MET Manual, Session 4 protocol. | Marlatt, G.A., & Gordon, J.R. (1985). Relapse Prevention. Guilford Press. | Witkiewitz, K., & Marlatt, G.A. (2004). Clinical Psychology Review, 24, 1-28.",
    learning_obj="Conduct the final MET session with appropriate consolidation and termination strategies",
    takeaway="Session 4 consolidates gains, builds relapse prevention skills, and terminates with confidence and hope."
)

add_slide(
    "Session 1: Complete Worksheet and Homework",
    [
        "SESSION 1 HOMEWORK ASSIGNMENTS:",
        "  1. Review Personal Feedback Report at home",
        "  2. Self-Monitoring: Keep a drinking diary for 1 week",
        "     (Record: when, where, how much, with whom, mood before/after)",
        "  3. Reflect on: 'What do these results mean to me?'",
        "  4. Optional: Share feedback with significant other",
        "",
        "SELF-MONITORING DIARY FORMAT:",
        "  ┌────────┬────────┬────────┬────────┬────────┬────────┐",
        "  │ Day    │ Time   │ Amount │ Where  │ Mood   │ Trigger│",
        "  ├────────┼────────┼────────┼────────┼────────┼────────┤",
        "  │Mon     │        │        │        │        │        │",
        "  │Tue     │        │        │        │        │        │",
        "  └────────┴────────┴────────┴────────┴────────┴────────┘",
        "",
        "SESSION 1 THERAPIST NOTES:",
        "  Document: Change talk heard, stage of change, readiness/confidence scores,",
        "  key themes, plan for Session 2",
    ],
    notes="Homework in MET is offered as a suggestion, not a requirement - maintaining autonomy. Self-monitoring serves multiple purposes: (1) it raises awareness of actual drinking patterns, (2) it provides data for future sessions, (3) the act of monitoring often reduces behavior (reactivity effect), and (4) it demonstrates engagement with the change process. The therapist should explain the rationale and ask if the client is willing, rather than prescribing it.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Sobell, L.C., & Sobell, M.B. (1992). Timeline Follow-Back method. In Litten, R.Z., & Allen, J.P. (Eds.), Measuring Alcohol Consumption. Humana Press.",
    learning_obj="Assign appropriate homework that reinforces MET Session 1 themes",
    takeaway="Homework is offered (not prescribed), focusing on self-monitoring and reflection to maintain momentum."
)


# ============================================================
# SECTION 15: TECHNIQUES USED IN MET (Slides 73-80)
# ============================================================
add_section_divider(15, "Techniques Used in MET")

add_slide(
    "Double-Sided Reflection & Complex Reflection",
    [
        "DOUBLE-SIDED REFLECTION:",
        "  Captures BOTH sides of ambivalence in one statement",
        "  Format: 'On one hand [sustain talk], AND on the other [change talk]'",
        "  Examples from MET Manual:",
        "    'You don't think alcohol is harming you seriously now, AND at the same time",
        "     you are concerned that it might get out of hand later'",
        "    'You really enjoy drinking and would hate to give it up, AND you can see",
        "     that it is causing serious problems for your family and your job'",
        "  NOTE: Use 'AND' rather than 'BUT' (but negates the first half)",
        "",
        "COMPLEX REFLECTION:",
        "  Goes BEYOND what the client said; adds meaning, feeling, or implication",
        "  Example: Client says 'I guess I drink a lot'",
        "    Simple: 'You drink quite a bit'",
        "    Complex: 'Part of you is starting to worry about where this is heading'",
        "  Complex reflections deepen exploration and show deeper understanding",
    ],
    notes="Double-sided reflections are one of the most powerful tools in MET. They validate the client's ambivalence (both sides are real and understandable) while creating cognitive dissonance (hearing both sides together highlights the contradiction). The MET manual specifically recommends using 'and' rather than 'but' because 'but' negates what came before it. Complex reflections demonstrate deeper understanding and can gently guide the client toward greater self-awareness. Both require practice to use naturally.",
    refs="Miller, W.R., et al. (1995). MET Manual, pp. 18-19. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Formulate double-sided and complex reflections from clinical material",
    takeaway="Double-sided reflections capture both sides of ambivalence; use 'AND' not 'BUT' to honor both sides."
)

add_slide(
    "Amplified Reflection, Reframing, and Shifting Focus",
    [
        "AMPLIFIED REFLECTION:",
        "  Overstates the client's position slightly to elicit the other side",
        "  Client: 'I can handle my drinking just fine'",
        "  Amplified: 'So your drinking causes you absolutely no concerns whatsoever'",
        "  Effect: Client often backs away from the extreme → 'Well, not NO concerns...'",
        "  CAUTION: Must not be sarcastic; delivered sincerely",
        "",
        "REFRAMING:",
        "  Offers a new perspective on the same information",
        "  Client: 'My wife keeps nagging me about my drinking'",
        "  Reframe: 'She cares about you so much that she keeps trying to reach you'",
        "  MET Manual: 'Reframing is offering a different interpretation of information'",
        "",
        "SHIFTING FOCUS:",
        "  Moving attention away from a stuck point",
        "  'Let's not worry about labels right now. Tell me what concerns YOU'",
        "  Useful when client is fixated on whether they're 'an alcoholic'",
        "  Avoids power struggles over definitions",
    ],
    notes="These three techniques are all used to manage resistance/discord. Amplified reflection uses reverse psychology - by overstating, you create reactance where the client argues AGAINST the extreme (and thereby toward change). It must be done sincerely and warmly, never sarcastically. Reframing takes the same facts but puts them in a different light, often highlighting caring or positive intent. Shifting focus avoids unproductive debates (especially about labels) by redirecting to what matters clinically.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Apply amplified reflection, reframing, and shifting focus to manage resistance",
    takeaway="Amplified reflection elicits the other side; reframing offers new perspective; shifting focus avoids power struggles."
)


add_slide(
    "Elicit-Provide-Elicit (E-P-E) and Agenda Mapping",
    [
        "ELICIT-PROVIDE-ELICIT (Ask-Tell-Ask):",
        "  A structured way to give information/advice without lecturing",
        "  Step 1 - ELICIT: 'What do you already know about alcohol and memory?'",
        "  Step 2 - PROVIDE: 'Would it be okay if I shared some information about this?'",
        "    (Wait for permission, then provide brief, relevant information)",
        "  Step 3 - ELICIT: 'What do you make of that?' / 'How does this fit for you?'",
        "  WHY: Maintains collaboration; avoids expert-patient hierarchy",
        "",
        "AGENDA MAPPING:",
        "  Collaboratively deciding what to discuss",
        "  'There are several things we could talk about today. Would you like to look at",
        "   your assessment results, or explore your concerns, or talk about options?'",
        "  PURPOSE: Client has ownership of the session; reduces resistance",
        "  Supports autonomy and collaboration",
        "",
        "SCALING QUESTIONS:",
        "  'On a scale of 0-10, how important/confident/ready are you?'",
        "  Follow-up is key: 'Why X and not lower?' (evokes change talk)",
    ],
    notes="E-P-E is essential for maintaining the spirit of MET while still providing expert information. The traditional medical model has the expert telling the patient what to do. E-P-E democratizes this by first checking what the client knows, then asking permission to add information, then checking how they process it. Agenda mapping at the start of sessions gives clients ownership and choice. Scaling questions are versatile tools that can assess and simultaneously evoke change talk through the strategic follow-up.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Rollnick, S., et al. (2008). Motivational interviewing in health care. Guilford Press. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Use E-P-E for information exchange, agenda mapping for collaboration, and scaling for assessment",
    takeaway="E-P-E maintains collaboration when sharing information; always ask permission before providing advice."
)

add_slide(
    "Values Clarification and Goal Setting",
    [
        "VALUES CLARIFICATION:",
        "  Purpose: Identify what matters most to the client (foundation for discrepancy)",
        "  Methods: Values card sort, open-ended exploration, life priorities list",
        "  Questions: 'What are the 3 most important things in your life?'",
        "    'What kind of person do you want to be?' 'What do you want to be remembered for?'",
        "  Connect values to behavior: 'How does drinking fit with being that person?'",
        "",
        "GOAL SETTING:",
        "  Client-generated goals (not therapist-imposed)",
        "  SMART framework: Specific, Measurable, Achievable, Relevant, Time-bound",
        "  Examples: 'I will not drink on weeknights for the next month'",
        "    'I will limit to 3 drinks maximum on weekends'",
        "    'I will have 4 alcohol-free days per week'",
        "",
        "ACTION PLANNING - Change Plan Worksheet:",
        "  1. Changes I want to make: ___",
        "  2. Most important reasons: ___",
        "  3. Steps I plan to take: ___",
        "  4. How others can help: ___",
        "  5. What might get in the way: ___",
    ],
    notes="Values clarification is a powerful tool for developing discrepancy. When clients articulate what truly matters to them (family, health, career, spirituality) and then look at how their behavior conflicts with these values, genuine internal motivation emerges. Goal setting should always be client-driven - the therapist guides the process but the client decides the goals. The Change Plan Worksheet formalizes the commitment and can be revisited in later sessions. Goals should be realistic and incremental.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Locke, E.A., & Latham, G.P. (2002). American Psychologist, 57, 705-717.",
    learning_obj="Conduct values clarification and collaborative goal setting in MET sessions",
    takeaway="Values create the 'why' for change; goals create the 'what'; action plans create the 'how'."
)


add_slide(
    "Looking Back, Looking Forward, and Decisional Balance",
    [
        "LOOKING BACK:",
        "  'What was life like before drinking became a problem?'",
        "  'What were you like before? What has changed?'",
        "  Purpose: Creates temporal discrepancy (then vs. now)",
        "  Evokes desire to return to better functioning",
        "",
        "LOOKING FORWARD:",
        "  'If you continue drinking as you are, where do you see yourself in 5 years?'",
        "  'If you DID make this change, what would life look like?'",
        "  Purpose: Creates future discrepancy (feared future vs. hoped future)",
        "  Makes consequences tangible and personal",
        "",
        "DECISIONAL BALANCE IN SESSION:",
        "  'What are the good things about drinking for you?' (Start here - validates)",
        "  'And what's the other side? What concerns you?' (Then explore costs)",
        "  'What might be better if you changed?' (Then benefits of change)",
        "  'What worries you about changing?' (Finally, fears about change)",
        "  Therapist SUMMARIZES: Include all sides; weight toward change talk",
    ],
    notes="These three techniques are among the most frequently used in MET. Looking back helps clients reconnect with their pre-addiction self and the losses they've experienced. Looking forward makes abstract consequences concrete and personal. Decisional balance explores ambivalence thoroughly. The order matters: starting with benefits of drinking shows you understand and won't judge, which reduces defensiveness. Ending with benefits of change leaves the session pointing forward. The summary strategically includes more change-favoring elements.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Janis, I.L., & Mann, L. (1977). Decision Making. Free Press.",
    learning_obj="Use temporal perspective techniques and decisional balance to develop discrepancy",
    takeaway="Looking back/forward creates temporal discrepancy; decisional balance explores ambivalence completely."
)

# ============================================================
# SECTION 16: WORKSHEETS (Slides 81-84)
# ============================================================
add_section_divider(16, "Clinical Worksheets and Handouts")

add_slide(
    "Change Plan Worksheet",
    [
        "CHANGE PLAN WORKSHEET (To be completed with client in Session 2-3):",
        "",
        "1. The changes I want to make are:",
        "   _____________________________________________",
        "",
        "2. The most important reasons I want to make these changes are:",
        "   _____________________________________________",
        "",
        "3. The steps I plan to take in changing are:",
        "   _____________________________________________",
        "",
        "4. The ways other people can help me are:",
        "   Person: _________ How they can help: _________",
        "",
        "5. I will know that my plan is working if:",
        "   _____________________________________________",
        "",
        "6. Some things that could interfere with my plan are:",
        "   _____________________________________________",
        "",
        "7. What I will do if the plan isn't working:",
        "   _____________________________________________",
    ],
    notes="The Change Plan Worksheet is a key tool in MET Phase 2. It should be completed collaboratively - the therapist helps the client think through each section but the client fills it in (maintains ownership). This worksheet makes the commitment concrete and tangible. It can be revisited in Sessions 3 and 4 to assess progress and modify the plan if needed. The worksheet should be given to the client to keep as a physical reminder of their commitment.",
    refs="Miller, W.R., et al. (1995). MET Manual (original Change Plan Worksheet). | SAMHSA. (2019). TIP 35: Enhancing Motivation for Change.",
    learning_obj="Complete a Change Plan Worksheet collaboratively with clients",
    takeaway="The Change Plan Worksheet makes commitment tangible and serves as a concrete roadmap for change."
)


add_slide(
    "Decisional Balance Worksheet & Relapse Prevention Plan",
    [
        "DECISIONAL BALANCE WORKSHEET:",
        "  ┌──────────────────────────┬──────────────────────────┐",
        "  │ BENEFITS of Drinking     │ COSTS of Drinking        │",
        "  │ (What I like about it)   │ (What concerns me)       │",
        "  │ 1.                       │ 1.                       │",
        "  │ 2.                       │ 2.                       │",
        "  │ 3.                       │ 3.                       │",
        "  ├──────────────────────────┼──────────────────────────┤",
        "  │ COSTS of Changing        │ BENEFITS of Changing     │",
        "  │ (What I fear/will miss)  │ (What I'd gain)          │",
        "  │ 1.                       │ 1.                       │",
        "  │ 2.                       │ 2.                       │",
        "  │ 3.                       │ 3.                       │",
        "  └──────────────────────────┴──────────────────────────┘",
        "",
        "RELAPSE PREVENTION PLAN:",
        "  My triggers: ___ | Warning signs: ___ | Coping strategies: ___",
        "  Emergency contacts: ___ | If I slip, I will: ___",
    ],
    notes="These worksheets should be printed as handouts for clients. The decisional balance worksheet makes ambivalence visible and concrete. Writing it down externalizes the internal conflict. The relapse prevention plan is introduced in Session 3-4 and helps the client anticipate and prepare for challenges. Both worksheets can be taken home and referred to when motivation wavers. In the Indian context (NIMHANS), these can be adapted linguistically and culturally.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Marlatt, G.A., & Gordon, J.R. (1985). Relapse Prevention. Guilford Press. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Use worksheets as clinical tools for externalizing ambivalence and planning",
    takeaway="Written worksheets make internal processes visible and serve as take-home reminders of commitment."
)

add_slide(
    "Additional Clinical Worksheets",
    [
        "READINESS/IMPORTANCE/CONFIDENCE VISUAL RULER:",
        "  Not at all                                    Extremely",
        "  0----1----2----3----4----5----6----7----8----9----10",
        "",
        "SUBSTANCE USE DIARY (Daily Monitoring Sheet):",
        "  Date | Substance | Amount | Time | Situation | Thoughts | Feelings | Consequences",
        "",
        "TRIGGER IDENTIFICATION WORKSHEET:",
        "  External triggers: People___ Places___ Times___ Events___",
        "  Internal triggers: Emotions___ Thoughts___ Physical___ Urges___",
        "",
        "COPING SKILLS WORKSHEET:",
        "  Situation → Usual response → Alternative response → Outcome",
        "",
        "VALUES CARD SORT (Top 10 Values):",
        "  Family | Health | Career | Spirituality | Freedom | Integrity",
        "  Love | Adventure | Security | Achievement | Friendship | Peace",
        "",
        "EMERGENCY PLAN: Who to call | Where to go | What to do | What NOT to do",
    ],
    notes="These worksheets complement the core MET sessions. They can be assigned as between-session activities or used within sessions as structured exercises. The values card sort is particularly powerful for developing discrepancy - once clients identify their top values, the therapist can explore how substance use aligns (or conflicts) with these values. The emergency plan is crucial for relapse prevention - having a concrete plan reduces the likelihood of a lapse becoming a full relapse.",
    refs="Miller, W.R., et al. (1995). MET Manual. | SAMHSA. (2019). TIP 35. | Murthy, P. (2008). NIMHANS Manual. | Marlatt, G.A., & Gordon, J.R. (1985). Relapse Prevention.",
    learning_obj="Select and use appropriate worksheets to support different phases of MET",
    takeaway="Worksheets structure the therapeutic process and give clients tangible tools for between-session work."
)


# ============================================================
# SECTION 17: CASE FORMULATIONS (Slides 85-89)
# ============================================================
add_section_divider(17, "Case Formulations")

add_slide(
    "Case 1: Alcohol Dependence - Mr. Rajesh, 42 years",
    [
        "PRESENTING PROBLEM: Referred by physician for elevated liver enzymes; 15 years drinking",
        "HISTORY: Daily drinking 180-350ml whiskey; tolerance increased; morning tremors",
        "STAGE OF CHANGE: Contemplation (aware of problem, not committed to change)",
        "",
        "MET FORMULATION:",
        "  Importance: 5/10 (knows it's a problem but enjoys social aspects)",
        "  Confidence: 3/10 (previous failed attempts; low self-efficacy)",
        "  Values: Family (wife and children), career advancement, health",
        "  Discrepancy: Values family but missing children's school events; promotion denied",
        "",
        "MET IMPLEMENTATION:",
        "  Session 1: Feedback (liver results, normative comparison); explore concerns",
        "  Session 2: Values clarification; decisional balance; build efficacy",
        "  Session 3: Review progress (reduced to weekends only); affirm success",
        "  Session 4: Maintenance plan; coping strategies for work stress; relapse prevention",
        "",
        "OUTCOME: Reduced drinking by 70%; liver enzymes normalized at 6 months",
    ],
    notes="This case illustrates a typical MET application for alcohol dependence in the Indian context. Key points: (1) The referral from a physician provides a natural opportunity for feedback. (2) His low confidence despite moderate importance tells us to focus on building self-efficacy. (3) His values (family, career) provide rich material for discrepancy development. (4) The gradual reduction goal (rather than abstinence) demonstrates the menu of options approach. In Indian settings, family involvement and cultural norms around drinking are important considerations.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Murthy, P. (2008). NIMHANS Manual. | Project MATCH Research Group. (1997). JOSA, 58, 7-29.",
    learning_obj="Formulate and implement a MET treatment plan for alcohol dependence",
    takeaway="MET case formulation integrates readiness assessment, values, discrepancy, and stage-matched intervention."
)

add_slide(
    "Case 2: Cannabis Use Disorder & Case 3: Nicotine Dependence",
    [
        "CASE 2 - CANNABIS: Ravi, 22, Engineering Student",
        "  Presenting: Academic failure, memory problems, social withdrawal",
        "  Stage: Precontemplation → Contemplation (came due to parents)",
        "  MET Focus: Develop discrepancy (career goals vs. cognitive effects)",
        "    Session 1: Feedback on neuropsychological testing; academic performance",
        "    Session 2: Values (engineering career); looking forward (5 years)",
        "    Outcome: Reduced use; improved grades over 3 months",
        "",
        "CASE 3 - NICOTINE: Priya, 35, Software Professional",
        "  Presenting: 15 cigarettes/day x 12 years; chest pain; wants to quit",
        "  Stage: Preparation (motivated but multiple failed attempts)",
        "  MET Focus: Build self-efficacy (past successes); address confidence",
        "    Session 1: Feedback (CO levels, lung function); health data",
        "    Session 2: Previous quit attempts - what worked? Build on strengths",
        "    Session 3: Using NRT; coping with triggers; celebrating 4 weeks smoke-free",
        "    Outcome: Quit at 6 months with NRT support",
    ],
    notes="These cases demonstrate MET adaptation for different substances. Cannabis case shows MET with a precontemplator (mandated by parents) - gentle approach needed. The neuropsychological feedback creates discrepancy between his engineering career goals and cognitive impairment. Nicotine case shows MET with someone already motivated but lacking confidence - focus shifts to self-efficacy building. Both cases show how MET is individualized based on the client's specific profile.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Stephens, R.S., et al. (2004). JCCP, 72, 92-103. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Apply MET to cannabis and nicotine use disorders with appropriate adaptations",
    takeaway="MET adapts to different substances by tailoring feedback and discrepancy to substance-specific concerns."
)


add_slide(
    "Case 4: Dual Diagnosis - Depression with Alcohol Use",
    [
        "CASE: Meera, 38, Homemaker with Major Depression and Alcohol Dependence",
        "  Presenting: Husband left; using alcohol to cope with depression; 250ml rum daily",
        "  Comorbidity: MDD (moderate), GAD, low self-esteem, social isolation",
        "  Stage: Contemplation (sees connection between drinking and depression worsening)",
        "",
        "MET FORMULATION:",
        "  Importance: 7/10 (can see alcohol worsening depression)",
        "  Confidence: 2/10 (feels hopeless; 'I can't do anything right')",
        "  Key discrepancy: Wants to feel better ↔ Alcohol is making depression worse",
        "",
        "MET SESSION PLAN:",
        "  Session 1: Feedback on depression-alcohol cycle; psychoeducation (E-P-E)",
        "  Session 2: Values (children, independence); self-efficacy building; small goals",
        "  Session 3: Progress review; address depressive cognitions; celebrate small wins",
        "  Session 4: Ongoing depression treatment plan; social support; maintenance",
        "",
        "INTEGRATION: MET for substance use + pharmacotherapy + CBT for depression",
        "OUTCOME: Reduced drinking; engaged in depression treatment; improved functioning",
    ],
    notes="Dual diagnosis cases are common and challenging. The NIMHANS manual emphasizes that substance use and mental health disorders are frequently comorbid in Indian settings. In this case, MET addresses the alcohol use while acknowledging that depression is a driver. The feedback focuses on the alcohol-depression cycle (alcohol provides temporary relief but worsens depression long-term). The very low confidence requires significant self-efficacy work. Small, achievable goals are essential. MET can be integrated with pharmacotherapy and CBT.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Murthy, P. (2008). NIMHANS Manual. | Baker, A.L., et al. (2012). Clinical Psychology Review, 32, 726-738.",
    learning_obj="Formulate MET for dual diagnosis cases with comorbid depression",
    takeaway="In dual diagnosis, MET addresses the substance use while integrating with treatment for the comorbid condition."
)

# ============================================================
# SECTION 18: SPECIAL POPULATIONS (Slides 90-93)
# ============================================================
add_section_divider(18, "Special Populations")

add_slide(
    "MET with Adolescents and Young Adults",
    [
        "UNIQUE CONSIDERATIONS:",
        "  • Developmental stage: Identity formation, peer influence paramount",
        "  • Often mandated (parents, school, courts) - may be precontemplators",
        "  • Less insight into consequences (temporal discounting)",
        "  • Stronger reactance to perceived authority",
        "ADAPTATIONS:",
        "  • Emphasize autonomy even more strongly: 'It's YOUR life'",
        "  • Use peer-relevant feedback (comparison with age norms)",
        "  • Shorter sessions; more interactive; less formal",
        "  • Focus on immediate consequences (social, academic) over long-term health",
        "  • Build alliance before any feedback delivery",
        "  • Involve family sensitively (NIMHANS: family context critical in Indian settings)",
        "EVIDENCE: MET effective with college student heavy drinkers (Marlatt et al., 1998)",
        "  BASICS intervention: Single MI session reduced heavy drinking at 2-year follow-up",
        "NIMHANS: Adolescent substance use requires family involvement and cultural sensitivity",
    ],
    notes="Adolescents present unique challenges for MET. Their developmental stage means peer influence is stronger than for adults, they discount future consequences more heavily, and they react more strongly to perceived control. MET adaptations include: emphasizing autonomy, using peer-relevant normative feedback, keeping sessions shorter and more interactive, and focusing on consequences that matter to THEM (social embarrassment, academic problems) rather than long-term health. In Indian settings, family context is critical.",
    refs="Marlatt, G.A., et al. (1998). Journal of Consulting and Clinical Psychology, 66, 604-615. | Jensen, C.D., et al. (2011). Clinical Psychology Review, 31, 1024-1036. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Adapt MET principles for adolescent and young adult populations",
    takeaway="Adolescent MET emphasizes autonomy strongly, uses peer-relevant feedback, and addresses immediate consequences."
)


add_slide(
    "MET with Women, Pregnant Women, and Older Adults",
    [
        "WOMEN:",
        "  • Greater stigma associated with substance use (especially in India)",
        "  • Often hidden drinking; delayed help-seeking",
        "  • Comorbid trauma, domestic violence, depression common",
        "  • NIMHANS: 'Women substance users are greatly stigmatized'",
        "  • Adaptations: Safety-focused; trauma-informed; address shame sensitively",
        "",
        "PREGNANT WOMEN:",
        "  • Any alcohol use carries risk (Fetal Alcohol Spectrum Disorders)",
        "  • High motivation (baby's health) but high guilt and shame",
        "  • MET well-suited: Non-judgmental, builds on existing maternal motivation",
        "  • Evidence: Brief MI reduces alcohol use in pregnant women (Handmaker et al., 1999)",
        "",
        "OLDER ADULTS:",
        "  • Increased sensitivity to alcohol (pharmacokinetic changes)",
        "  • Loss, loneliness, retirement may drive use",
        "  • Fewer social consequences may reduce discrepancy",
        "  • Adaptations: Focus on health consequences, medication interactions, quality of life",
    ],
    notes="These populations require specific adaptations. Women in Indian society face tremendous stigma around substance use, leading to hidden drinking and delayed help-seeking. MET's non-judgmental approach is especially valuable here. Pregnant women are often already motivated (baby's health) but experience intense shame - MET can build on existing motivation while reducing shame. Older adults may have fewer social consequences of drinking but face serious health risks and medication interactions. Each population requires tailored feedback and sensitivity.",
    refs="Murthy, P. (2008). NIMHANS Manual. | Handmaker, N.S., et al. (1999). Journal of Consulting and Clinical Psychology, 67, 285-290. | Blow, F.C., & Barry, K.L. (2012). Alcohol Research, 34, 18-28.",
    learning_obj="Identify specific adaptations of MET for women, pregnant women, and older adults",
    takeaway="Each population requires tailored MET adaptations addressing unique barriers, stigma, and motivational factors."
)

add_slide(
    "MET in Medical Settings: HIV, Diabetes, Cancer",
    [
        "MEDICAL SETTINGS - MET Applications Beyond Addiction:",
        "",
        "HIV/AIDS:",
        "  • Medication adherence (antiretroviral therapy)",
        "  • Risk reduction (safe sex, needle exchange)",
        "  • Substance use complicating HIV treatment",
        "  • Evidence: MI improves ART adherence (Parsons et al., 2007)",
        "",
        "DIABETES:",
        "  • Lifestyle change (diet, exercise, medication adherence)",
        "  • Blood sugar monitoring compliance",
        "  • MET for ambivalence about lifestyle restrictions",
        "",
        "CANCER:",
        "  • Smoking cessation post-diagnosis",
        "  • Treatment adherence (chemotherapy)",
        "  • Pain management without substance misuse",
        "",
        "PRIMARY CARE INTEGRATION (NIMHANS Model):",
        "  • Brief interventions by non-specialists",
        "  • 5-A strategy: Ask, Assess, Advise, Assist, Arrange",
        "  • Physicians as agents of change",
    ],
    notes="The NIMHANS manual strongly advocates for integrating brief motivational interventions into primary care and medical settings. The 5-A strategy parallels FRAMES in a format accessible to busy physicians. MET principles are increasingly applied to any health behavior where ambivalence exists - medication adherence, lifestyle change, self-management of chronic illness. The key insight is that the same motivational principles (autonomy, empathy, discrepancy) apply to any behavior change context.",
    refs="Murthy, P. (2008). NIMHANS Manual, pp. 1-3. | Parsons, J.T., et al. (2007). AIDS and Behavior, 11, 725-734. | Rubak, S., et al. (2005). British Journal of General Practice, 55, 305-312.",
    learning_obj="Apply MET principles in medical settings for health behavior change",
    takeaway="MET principles apply to any health behavior change: medication adherence, lifestyle modification, risk reduction."
)


# ============================================================
# SECTION 19: APPLICATIONS BEYOND ADDICTION (Slides 94-96)
# ============================================================
add_section_divider(19, "Applications Beyond Addiction")

add_slide(
    "MET for Health Behavior Change",
    [
        "MET/MI principles have been successfully applied to:",
        "",
        "LIFESTYLE BEHAVIORS:",
        "  • Smoking cessation (Hettema & Hendricks, 2010)",
        "  • Weight management and obesity (Armstrong et al., 2011)",
        "  • Physical activity/exercise promotion",
        "  • Dietary change (fruit/vegetable intake, salt reduction)",
        "",
        "CHRONIC DISEASE MANAGEMENT:",
        "  • Diabetes self-management (Channon et al., 2007)",
        "  • Hypertension medication adherence",
        "  • Asthma management",
        "  • Pain management (Alperstein & Sharpe, 2016)",
        "",
        "MENTAL HEALTH:",
        "  • Treatment engagement for anxiety disorders",
        "  • Medication adherence in psychosis",
        "  • Eating disorders (Treasure & Schmidt, 2008)",
        "  • Gambling disorder",
        "",
        "COMMON PRINCIPLE: Wherever ambivalence about change exists, MET is applicable",
    ],
    notes="The versatility of MET principles is remarkable. Any situation where a person is ambivalent about changing a behavior can benefit from motivational approaches. The underlying mechanisms are the same: autonomy support, empathy, discrepancy, self-efficacy. The specific feedback and discrepancy development are tailored to the domain (e.g., HbA1c results for diabetes, BMI for weight). Meta-analyses confirm MI's effectiveness across these domains with small-to-medium effect sizes.",
    refs="Rubak, S., et al. (2005). British Journal of General Practice, 55, 305-312. | Lundahl, B., et al. (2013). Patient Education and Counseling, 93, 157-168. | Hettema, J., et al. (2005). Annual Review of Clinical Psychology, 1, 91-111.",
    learning_obj="Identify diverse health behavior change applications of MET beyond addiction",
    takeaway="MET principles apply wherever ambivalence exists - from substance use to exercise to medication adherence."
)

# ============================================================
# SECTION 20: RESEARCH EVIDENCE (Slides 97-101)
# ============================================================
add_section_divider(20, "Research Evidence for MET")

add_slide(
    "Major Randomized Controlled Trials",
    [
        "PROJECT MATCH (1997): N=1,726",
        "  • 3-arm RCT: MET (4 sessions) vs. CBT (12) vs. TSF (12)",
        "  • Result: All effective; MET equivalent with 1/3 sessions",
        "  • 3-year follow-up: Gains maintained across all groups",
        "",
        "UKATT (2005): N=742 (UK Alcohol Treatment Trial)",
        "  • MET (3 sessions) vs. Social Behaviour Network Therapy (8 sessions)",
        "  • Result: Equivalent outcomes; MET more cost-effective",
        "",
        "COMBINE Study (2006): N=1,383",
        "  • MET + CBT + Naltrexone combinations",
        "  • Result: Combined medical management + MI = good outcomes",
        "",
        "Brief Interventions Meta-analysis (Bien et al., 1993):",
        "  • FRAMES-based interventions reduce drinking by 20-30%",
        "  • Effect comparable to more intensive treatments",
        "",
        "Overall: >200 RCTs support MI/MET across substances and populations",
    ],
    notes="The evidence base for MET is among the strongest in psychotherapy research. Project MATCH remains the landmark study - its finding that 4 sessions matched 12 sessions was revolutionary. UKATT replicated this in the UK. The COMBINE study showed MET integrates well with pharmacotherapy. Meta-analyses consistently find MI/MET effective with small-to-medium effect sizes (d = 0.25-0.57). Importantly, MET is cost-effective - achieving similar results with fewer sessions means lower cost per unit of improvement.",
    refs="Project MATCH Research Group. (1997). JOSA, 58, 7-29. | UKATT Research Team. (2005). BMJ, 331, 541. | Anton, R.F., et al. (2006). JAMA, 295, 2003-2017. | Bien, T.H., Miller, W.R., & Tonigan, S. (1993). Addiction, 88, 315-336.",
    learning_obj="Summarize major RCTs supporting MET effectiveness",
    takeaway="200+ RCTs support MI/MET; Project MATCH showed 4 sessions of MET matches 12 sessions of other treatments."
)


add_slide(
    "Meta-Analyses and Systematic Reviews",
    [
        "KEY META-ANALYSES:",
        "",
        "Burke et al. (2003): 30 studies; MI effect size d=0.25-0.57",
        "  Strongest effects for alcohol, drugs, diet, exercise",
        "",
        "Hettema et al. (2005): 72 clinical trials reviewed",
        "  MI effective across multiple problem behaviors",
        "  Effects enhanced when combined with other active treatment",
        "",
        "Lundahl et al. (2010): 119 studies; MI significantly better than comparison",
        "  Small-medium effect across health behaviors (OR = 1.55)",
        "",
        "Vasilaki et al. (2006): 15 studies; Brief MI for alcohol",
        "  Significant reduction in alcohol consumption; effects durable at 12 months",
        "",
        "Smedslund et al. (2011): Cochrane Review - 59 studies",
        "  MI reduces substance use compared to no treatment",
        "  Effects may diminish over time without booster sessions",
        "",
        "MECHANISMS: Therapist empathy, change talk, and working alliance mediate outcomes",
    ],
    notes="The meta-analytic evidence provides strong support for MI/MET. Effect sizes are typically in the small-to-medium range, which is consistent with other psychotherapy research. The Cochrane Review by Smedslund et al. confirmed effectiveness but noted effects may diminish over time, supporting the use of booster sessions. Mechanisms research suggests the causal chain: therapist MI-consistent behavior → client change talk → behavior change. This understanding helps us refine training.",
    refs="Burke, B.L., et al. (2003). JCCP, 71, 843-861. | Hettema, J., et al. (2005). Annual Review of Clinical Psychology, 1, 91-111. | Lundahl, B., et al. (2010). Clinical Psychology Review, 30, 1-11. | Smedslund, G., et al. (2011). Cochrane Database of Systematic Reviews, 11.",
    learning_obj="Evaluate the meta-analytic evidence supporting MET effectiveness",
    takeaway="Meta-analyses confirm MI/MET effectiveness (d=0.25-0.57) across substances and health behaviors."
)

add_slide(
    "Evidence Table: Key Studies Summary",
    [
        "┌─────────────────────┬──────┬────────┬────────────────────────────────┐",
        "│ Study               │ Year │ N      │ Key Finding                    │",
        "├─────────────────────┼──────┼────────┼────────────────────────────────┤",
        "│ Project MATCH       │ 1997 │ 1,726  │ MET = CBT = TSF (fewer sess.)  │",
        "│ UKATT               │ 2005 │ 742    │ MET = SBNT (more cost-eff.)    │",
        "│ COMBINE             │ 2006 │ 1,383  │ MET + meds effective           │",
        "│ Marlatt et al.      │ 1998 │ 348    │ Brief MI reduces college drink │",
        "│ Stephens et al.     │ 2004 │ 291    │ MI for cannabis effective      │",
        "│ Handmaker et al.    │ 1999 │ 42     │ MI reduces prenatal alcohol    │",
        "│ Miller et al.       │ 1993 │ 42     │ Empathic > confrontational     │",
        "│ Bien et al.         │ 1993 │ meta   │ Brief interventions effective  │",
        "│ Moyers et al.       │ 2007 │ 103    │ MI behaviors → change talk     │",
        "│ Amrhein et al.      │ 2003 │ 84     │ Commitment talk predicts Δ     │",
        "└─────────────────────┴──────┴────────┴────────────────────────────────┘",
        "Overall evidence: STRONG support for MET across populations and substances",
    ],
    notes="This evidence table provides a quick reference for students. Each study represents a milestone in the MET evidence base. Project MATCH established equivalence with longer treatments. UKATT replicated internationally. Miller's 1993 study directly demonstrated that therapist style matters. Moyers and Amrhein illuminated the mechanism (change talk). The overall picture is one of consistent, replicated support across multiple research groups, populations, and settings.",
    refs="All citations as listed in the table. See bibliography for full APA references.",
    learning_obj="Reference key studies that form the evidence base for MET",
    takeaway="The evidence base for MET spans large RCTs, meta-analyses, and mechanism studies across 30+ years."
)


# ============================================================
# SECTION 21: COMPARISON WITH OTHER THERAPIES (Slides 102-104)
# ============================================================
add_section_divider(21, "Comparison with Other Therapies")

add_slide(
    "MET vs. CBT, DBT, ACT, and Other Approaches",
    [
        "┌────────────┬────────────────────────┬─────────────────────────────────┐",
        "│ Therapy    │ Focus                  │ How it differs from MET         │",
        "├────────────┼────────────────────────┼─────────────────────────────────┤",
        "│ MET        │ Internal motivation    │ Client argues for change        │",
        "│ CBT        │ Skills & cognitions    │ Therapist teaches skills        │",
        "│ DBT        │ Emotion regulation     │ Structured skills training      │",
        "│ ACT        │ Acceptance & values    │ Mindfulness + committed action  │",
        "│ 12-Step    │ Spiritual surrender    │ Powerlessness; higher power     │",
        "│ REBT       │ Irrational beliefs     │ Disputational; directive        │",
        "│ Psychodyn. │ Unconscious conflict   │ Interpretation; long-term       │",
        "│ CM         │ Reinforcement          │ External rewards for behavior   │",
        "│ Support.   │ Emotional support      │ Non-directive; no structure     │",
        "└────────────┴────────────────────────┴─────────────────────────────────┘",
        "",
        "KEY DISTINCTION: MET focuses on WHY change (motivation) not HOW (skills)",
        "MET can be COMBINED with other approaches (MET → CBT is common sequence)",
        "MET often used as 'front door' to engage clients in further treatment",
    ],
    notes="This comparison helps students understand where MET fits in the therapeutic landscape. MET is unique in focusing primarily on motivation (the 'why') rather than skills (the 'how'). This makes it complementary to skills-based approaches like CBT. A common clinical sequence is MET first (to build motivation) followed by CBT (to build skills). MET differs from 12-Step in emphasizing personal empowerment rather than powerlessness. It differs from CBT in being less prescriptive and more client-led.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Project MATCH Research Group. (1997). JOSA, 58, 7-29. | Carroll, K.M. (1998). A Cognitive-Behavioral Approach. Yale University Press.",
    learning_obj="Compare and contrast MET with other major psychotherapy approaches",
    takeaway="MET focuses on WHY change (motivation) rather than HOW (skills), making it complementary to CBT and other approaches."
)

add_slide(
    "MET vs. MI: Key Differences",
    [
        "┌──────────────────────┬────────────────────────────────────────────┐",
        "│ Feature              │ MI                    MET                  │",
        "├──────────────────────┼────────────────────────────────────────────┤",
        "│ Definition           │ Counseling style      Manualized treatment │",
        "│ Structure            │ Flexible              4-session protocol   │",
        "│ Assessment           │ Optional              Required (PFR)       │",
        "│ Feedback             │ May/may not include   Core component       │",
        "│ Duration             │ Variable              4 sessions fixed     │",
        "│ Manual               │ No fixed manual       Project MATCH manual │",
        "│ Training             │ Style/spirit-based    Protocol-based       │",
        "│ Applicability        │ Any behavior change   Primarily addiction  │",
        "│ Evidence base        │ Broader (200+ RCTs)   Specific (MATCH etc) │",
        "└──────────────────────┴────────────────────────────────────────────┘",
        "",
        "MI = The broader approach/spirit (applicable anywhere)",
        "MET = MI + Structured Assessment Feedback + Fixed Protocol",
        "MET is a SPECIFIC APPLICATION of MI principles in a structured format",
    ],
    notes="This distinction is frequently confused by students. MI is a general counseling style that can be applied in any setting for any duration. MET is a SPECIFIC manualized treatment that USES MI principles within a structured 4-session format that includes personalized assessment feedback. Think of MI as the philosophy and MET as a specific program built on that philosophy. MI can be used in a 5-minute primary care encounter or a 50-minute therapy session. MET always follows the 4-session Project MATCH protocol.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Clearly distinguish between MI (counseling style) and MET (manualized treatment)",
    takeaway="MI is a flexible counseling style; MET is MI embedded in a structured 4-session protocol with assessment feedback."
)


# ============================================================
# SECTION 22-23: ADVANTAGES & LIMITATIONS (Slides 105-107)
# ============================================================
add_section_divider(22, "Advantages and Limitations of MET")

add_slide(
    "Advantages of MET",
    [
        "CLINICAL ADVANTAGES:",
        "  • Brief (4 sessions) - accessible and cost-effective",
        "  • Non-confrontational - reduces dropout and resistance",
        "  • Client-centered - respects autonomy and dignity",
        "  • Evidence-based - strong research support (200+ trials)",
        "  • Flexible - adaptable to diverse populations and settings",
        "  • Trainable - can be learned by various professionals",
        "  • Compatible - can be combined with other treatments (CBT, pharmacotherapy)",
        "",
        "PRACTICAL ADVANTAGES:",
        "  • Feasible in resource-limited settings (NIMHANS: suitable for India)",
        "  • Can be delivered by non-specialists with training",
        "  • Applicable to primary care (brief formats)",
        "  • Effective across cultures when adapted appropriately",
        "  • Works with mandated/reluctant clients",
        "  • Reduces therapist burnout (collaborative vs. adversarial)",
        "",
        "COST-EFFECTIVENESS: UKATT showed MET was 5x more cost-effective than comparison",
    ],
    notes="MET's advantages make it uniquely suitable for widespread implementation. Its brevity means more clients can be served with limited resources. Its non-confrontational nature reduces dropout (a major problem in addiction treatment). Its evidence base satisfies funders and policymakers. Its trainability means it can be disseminated broadly. The NIMHANS manual specifically identifies brief motivational approaches as suitable for the Indian healthcare system where specialist resources are limited.",
    refs="Miller, W.R., et al. (1995). MET Manual. | UKATT Research Team. (2005). BMJ, 331, 541. | Murthy, P. (2008). NIMHANS Manual. | Lundahl, B., et al. (2010). Clinical Psychology Review, 30, 1-11.",
    learning_obj="Articulate the clinical, practical, and economic advantages of MET",
    takeaway="MET is brief, non-confrontational, evidence-based, cost-effective, and adaptable to diverse settings."
)

add_slide(
    "Limitations of MET",
    [
        "CLINICAL LIMITATIONS:",
        "  • May not be sufficient alone for severe dependence (needs integrated care)",
        "  • Effects may diminish over time without boosters (Smedslund et al., 2011)",
        "  • Not designed for skills building (may need CBT follow-up)",
        "  • Limited evidence for some populations (psychosis, severe cognitive impairment)",
        "  • Requires skilled therapist (poorly done MI can be harmful)",
        "",
        "METHODOLOGICAL CONCERNS:",
        "  • Effect sizes are small-medium (not dramatic improvement)",
        "  • Active mechanism not fully understood",
        "  • Difficult to separate effects of relationship from specific techniques",
        "  • Treatment fidelity varies widely across studies",
        "",
        "PRACTICAL LIMITATIONS:",
        "  • Requires quality training and supervision (not just reading the manual)",
        "  • NIMHANS: 'Should not be viewed as complete solutions to substance use problems'",
        "  • May be insufficient for complex cases (severe trauma, psychosis)",
        "  • Cultural adaptations needed (not 'one size fits all')",
        "  • Risk of therapist drift without ongoing fidelity monitoring",
    ],
    notes="Being honest about limitations is essential for academic integrity. MET is powerful but not a panacea. The NIMHANS manual explicitly states that brief interventions 'should not be viewed as complete solutions.' For severe dependence, MET works best as a motivational gateway leading to comprehensive treatment. The small-medium effect sizes are typical for psychotherapy but remind us that many clients will not respond. Ongoing training, supervision, and fidelity monitoring are needed to maintain quality.",
    refs="Smedslund, G., et al. (2011). Cochrane Database. | Murthy, P. (2008). NIMHANS Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Critically evaluate the limitations and appropriate boundaries of MET",
    takeaway="MET has limitations: small-medium effects, may need boosters, insufficient alone for severe cases, requires skilled delivery."
)


# ============================================================
# SECTION 24-25: ETHICS & CULTURAL ADAPTATION (Slides 108-111)
# ============================================================
add_section_divider(24, "Ethical Considerations and Cultural Adaptation")

add_slide(
    "Ethical Considerations in MET",
    [
        "AUTONOMY:",
        "  • Client's right to choose their own goals (including continued use)",
        "  • Informed consent about the therapeutic approach",
        "  • Mandated clients: 'You have to be here, but what you do is still your choice'",
        "",
        "CONFIDENTIALITY:",
        "  • Standard limits (risk to self/others, child abuse)",
        "  • Significant other involvement requires client consent",
        "  • Special considerations with adolescents and mandated clients",
        "",
        "NON-MALEFICENCE:",
        "  • Poor MI can be harmful (confrontation disguised as MI)",
        "  • Manipulation concern: Is strategic reflection ethical? (Yes, if serving client welfare)",
        "  • Maintaining boundaries between influence and coercion",
        "",
        "COMPETENCE:",
        "  • Therapist must be trained, supervised, and fidelity-monitored",
        "  • 'Reading the manual is not sufficient' - requires practice and feedback",
        "  • Cultural competence essential for diverse populations",
    ],
    notes="Ethical issues in MET warrant careful consideration. The primary ethical concern is whether strategic use of reflective listening constitutes manipulation. Miller & Rollnick address this directly: if the therapist's intent is genuinely compassionate (serving the client's welfare), strategic techniques are ethical. If intent is self-serving (making the client comply with OUR goals), it becomes manipulation. The spirit of MET (compassion, partnership) is the ethical safeguard. Autonomy must be genuinely honored - even if the client decides not to change.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | APA. (2017). Ethical principles of psychologists. | NIMHANS ethical guidelines.",
    learning_obj="Identify ethical considerations specific to MET practice",
    takeaway="MET's ethical foundation rests on genuine compassion and respect for autonomy; strategic techniques serve the client."
)

add_slide(
    "Cultural Adaptation of MET: Indian Context",
    [
        "NIMHANS PERSPECTIVE ON CULTURAL ADAPTATION:",
        "  • India: 62.5 million alcohol users; limited specialist resources",
        "  • Brief motivational interventions particularly relevant for primary care",
        "  • Need for trained physicians as 'agents of change' (NIMHANS, 2008)",
        "",
        "CULTURAL CONSIDERATIONS FOR INDIA:",
        "  • Family involvement is culturally expected and therapeutically useful",
        "  • Stigma around substance use (especially for women) - heightened sensitivity needed",
        "  • Spiritual/religious values can be leveraged for discrepancy development",
        "  • Hierarchical relationships: MET's egalitarian approach may need gentle introduction",
        "  • Language adaptation: Materials must be available in local languages",
        "  • Joint family dynamics: Consider multiple stakeholders",
        "  • Gender roles: Women's substance use is heavily stigmatized",
        "",
        "UNIVERSAL PRINCIPLES THAT TRANSCEND CULTURE:",
        "  • Empathy, respect, autonomy support, non-judgment",
        "  • These values are recognized across all cultures",
    ],
    notes="The NIMHANS manual provides the Indian context for MET implementation. India's large population of substance users combined with limited specialist resources makes brief, trainable interventions essential. Cultural adaptations include: greater family involvement (consistent with collectivist culture), attention to stigma (especially for women), use of spiritual/religious values as motivational levers, and respect for hierarchical relationships while still promoting client autonomy. The core principles of empathy and respect are universal.",
    refs="Murthy, P. (2008). NIMHANS Manual, pp. 1-3. | D'Amico, E.J., et al. (2015). Journal of Substance Abuse Treatment, 59, 91-99. | Benegal, V. (2005). Addiction, 100, 1051-1056.",
    learning_obj="Adapt MET for cultural contexts, particularly Indian settings",
    takeaway="MET's core principles are universal; cultural adaptations address family involvement, stigma, and local values."
)


# ============================================================
# SECTION 26: THERAPIST COMPETENCIES (Slides 112-114)
# ============================================================
add_section_divider(26, "Therapist Competencies and Training")

add_slide(
    "Skills Required and MITI Coding",
    [
        "CORE COMPETENCIES FOR MET THERAPISTS:",
        "  • Reflective listening (the foundational skill)",
        "  • Asking open questions",
        "  • Affirming client strengths",
        "  • Summarizing strategically",
        "  • Recognizing and reinforcing change talk",
        "  • Rolling with resistance",
        "  • Delivering feedback non-judgmentally",
        "  • Managing own 'righting reflex'",
        "",
        "MITI (Motivational Interviewing Treatment Integrity) CODING:",
        "  Developed by Moyers et al. for assessing MI fidelity",
        "  Measures: Technical (reflection:question ratio, % open questions, % complex reflections)",
        "  Measures: Relational (partnership, empathy)",
        "  THRESHOLDS: Competent = reflection:question ≥ 1:1; Open questions ≥ 50%",
        "             Proficient = reflection:question ≥ 2:1; Open questions ≥ 70%",
        "",
        "TRAINING: Workshop + practice + feedback + supervision (not just reading)",
    ],
    notes="Therapist competency is critical for MET effectiveness. Research shows that poorly delivered MI can be no better than - or worse than - no treatment. The MITI coding system provides objective measures of fidelity. Key metrics include the reflection-to-question ratio (should be at least 2:1 for proficiency), percentage of open questions (should be >70% for proficiency), and percentage of complex reflections. Training typically involves a 2-day workshop followed by ongoing practice, recorded sessions, feedback, and supervision.",
    refs="Moyers, T.B., et al. (2005). MITI Manual (Version 3.0). University of New Mexico. | Miller, W.R., & Moyers, T.B. (2006). Behavioural and Cognitive Psychotherapy, 34, 135-143. | Madson, M.B., et al. (2009). Motivational Interviewing Training. Springer.",
    learning_obj="Identify competencies required for MET practice and methods for assessing fidelity",
    takeaway="MI fidelity requires specific measurable skills; MITI coding provides objective assessment of competence."
)

add_slide(
    "Training, Supervision, and Maintaining Fidelity",
    [
        "TRAINING PATHWAY:",
        "  1. Foundational knowledge: Read manual; understand theories",
        "  2. Workshop training: 2-3 day intensive with practice",
        "  3. Practice with feedback: Record sessions; get MITI coding",
        "  4. Supervision: Regular case discussion with trained supervisor",
        "  5. Ongoing development: Coaching, peer groups, booster training",
        "",
        "SUPERVISION FOCUS:",
        "  • Review recorded sessions (audio/video)",
        "  • Identify moments of drift (confrontation, advice-giving)",
        "  • Celebrate skillful moments",
        "  • Address therapist 'righting reflex'",
        "  • Role-play challenging scenarios",
        "",
        "COMMON THERAPIST MISTAKES:",
        "  • The 'expert trap': Telling rather than asking",
        "  • The 'assessment trap': Too many questions, not enough reflections",
        "  • The 'premature focus trap': Moving to planning before client is ready",
        "  • The 'labeling trap': Using 'alcoholic' or other stigmatizing labels",
        "  • The 'blaming trap': Attributing resistance to client pathology",
    ],
    notes="Training in MI/MET is a developmental process, not a one-time event. Research shows that a 2-day workshop alone does not produce lasting competence - ongoing supervision and feedback are essential. Common mistakes include falling back into expert mode (telling), asking too many questions (interrogation mode), moving to planning too quickly (premature focus), and using labels that trigger resistance. The therapist's own 'righting reflex' (desire to fix problems) is the primary barrier to maintaining fidelity.",
    refs="Miller, W.R., & Moyers, T.B. (2006). Behavioural and Cognitive Psychotherapy, 34, 135-143. | Miller, W.R., et al. (2004). Journal of Substance Abuse Treatment, 26, 3-12. | Schwalbe, C.S., et al. (2014). Addiction, 109, 1287-1294.",
    learning_obj="Design a training and supervision plan for developing MET competence",
    takeaway="Competence requires workshop + practice + feedback + supervision; reading alone is insufficient."
)


# ============================================================
# SECTION 27: PRACTICAL DEMONSTRATION (Slides 115-118)
# ============================================================
add_section_divider(27, "Practical Demonstration: Role Plays")

add_slide(
    "Role Play Script: Therapist-Client Dialogue (Session 1 Opening)",
    [
        "SCENARIO: First MET session with Arun, 35, IT professional, referred by wife",
        "",
        "T: 'Welcome, Arun. Thank you for coming in today. I'd like to hear from you -",
        "   what brings you here?'",
        "C: 'My wife made me come. She thinks I drink too much. I don't see the problem.'",
        "T: 'So coming here wasn't really your idea. You're doing this for her.' [reflection]",
        "C: 'Yeah, exactly. I mean, everyone I know drinks.'",
        "T: 'From your perspective, your drinking is pretty normal for your circle.' [reflection]",
        "C: 'Right. But she keeps nagging me about it.'",
        "T: 'That must be frustrating. And yet here you are.' [reflection + observation]",
        "C: 'Well... I love her. And she IS worried. I don't want to lose her.'",
        "T: 'She means a lot to you, and her concern affects you.' [reflects caring value]",
        "C: 'Yeah. I mean, maybe I COULD cut back a little if it would make her happy.'",
        "T: 'So part of you can see a reason to look at this more closely.' [change talk!]",
    ],
    notes="This role play demonstrates a typical Session 1 opening with a reluctant/mandated client. Notice: (1) The therapist NEVER argues with the client's position, (2) Every response is a reflection, (3) The therapist picks up on the value (love for wife) and reflects it, (4) By the end of this brief exchange, the client has moved from 'I don't see the problem' to 'maybe I could cut back' - without any confrontation. This is the power of reflective listening and rolling with resistance. The change talk emerges naturally when the client feels safe.",
    refs="Miller, W.R., et al. (1995). MET Manual (clinical dialogue examples adapted from the manual).",
    learning_obj="Observe and practice MET dialogue techniques in a realistic clinical scenario",
    takeaway="Skilled reflection moves clients from resistance to change talk without any confrontation or pressure."
)

add_slide(
    "Role Play: Common Mistakes vs. Correct Responses",
    [
        "SCENARIO: Client says 'I don't think I'm an alcoholic'",
        "",
        "❌ WRONG RESPONSES:",
        "  Confrontation: 'Yes you are! Look at your test results!'",
        "  Labeling: 'Whether you call it that or not, you meet the criteria'",
        "  Question trap: 'Why do you think you're NOT an alcoholic?'",
        "  Lecturing: 'Let me tell you about the diagnostic criteria...'",
        "",
        "✓ CORRECT RESPONSES:",
        "  Shifting focus: 'Labels aren't important here. What concerns YOU about your drinking?'",
        "  Reflection: 'You don't like that word being applied to you'",
        "  Reframe: 'You're someone who thinks carefully before accepting conclusions'",
        "  Double-sided: 'You don't think you're alcoholic, AND you do have some concerns'",
        "  Amplified: 'So your drinking is absolutely fine in every way' [elicits: 'Well..not EVERY way']",
        "",
        "KEY PRINCIPLE: NEVER argue about labels; explore what the person IS concerned about",
    ],
    notes="This slide demonstrates the contrast between MI-inconsistent and MI-consistent responses. Students should practice each response type and notice how the wrong responses create an adversarial dynamic while the correct responses maintain collaboration. The most common mistake for new therapists is the 'label debate' - arguing about whether someone is an alcoholic. MET completely avoids this by shifting focus to the client's own concerns. This is one of the most important clinical skills to develop.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Identify common therapist mistakes and practice MI-consistent alternatives",
    takeaway="Never argue about labels; shift focus to the client's own concerns and reflect their perspective."
)


add_slide(
    "Extended Role Play: Feedback Delivery with Resistance",
    [
        "SCENARIO: Delivering high drinking results to a defensive client",
        "",
        "T: 'Here are your results. Your weekly consumption is about 35 standard drinks.",
        "   That puts you above 95% of the adult population.' [neutral feedback]",
        "C: 'That can't be right. There's no way I drink that much.'",
        "T: 'This is surprising to you. It doesn't match how you see your drinking.' [reflect]",
        "C: 'The test must be wrong. Or maybe I exaggerated.'",
        "T: 'You're wondering if the numbers are accurate.' [simple reflection]",
        "C: 'Well... I guess I DO drink most nights. But it's just social.'",
        "T: 'So drinking is a regular part of your social life.' [reflect]",
        "C: 'Yeah. But 95th percentile? That sounds really bad.'",
        "T: 'That number is concerning to you.' [reflects concern = change talk!]",
        "C: 'I mean... I never thought of myself as THAT kind of drinker.'",
        "T: 'There's a gap between how you've seen yourself and what these data suggest.",
        "   That's a lot to take in.' [reflects discrepancy; names the process]",
        "RESULT: Client moves from denial → concern → discrepancy - all via reflection",
    ],
    notes="This extended role play shows the feedback delivery process from the MET manual. Notice the therapist's discipline - despite the client's initial rejection of the data, the therapist NEVER argues for the validity of the results. Instead, they reflect the client's experience. This creates safety. Gradually, the client's own cognitive dissonance emerges - 'I never thought of myself as THAT kind of drinker.' The therapist names the discrepancy without pushing it. This is textbook MET.",
    refs="Miller, W.R., et al. (1995). MET Manual, Session 1 protocol.",
    learning_obj="Practice feedback delivery with resistant clients using pure reflective responding",
    takeaway="Even when clients reject feedback, consistent reflection creates safety for the discrepancy to emerge naturally."
)

# ============================================================
# SECTION 28: EXAMINATION QUESTIONS (Slides 119-122)
# ============================================================
add_section_divider(28, "Examination Questions and Practice")

add_slide(
    "Viva Questions and Long Answer Questions",
    [
        "POSSIBLE VIVA QUESTIONS:",
        "  1. Define MET. How does it differ from MI?",
        "  2. Explain the theoretical foundations of MET (name 5 theories).",
        "  3. What is the spirit of MI/MET? Why is it more important than technique?",
        "  4. Describe the OARS framework with clinical examples.",
        "  5. What is change talk? Name the 7 types (DARN-CAT).",
        "  6. Explain the FRAMES model and its evidence base.",
        "  7. How do you 'roll with resistance'? Give 3 specific strategies.",
        "  8. Describe the structure of the 4-session MET protocol.",
        "  9. What does the research evidence say about MET? Cite key studies.",
        "  10. How would you adapt MET for the Indian context?",
        "",
        "LONG ANSWER QUESTIONS (10-15 marks):",
        "  • Discuss MET with reference to its theoretical foundations, evidence base,",
        "    session structure, and clinical applications. (20 marks)",
        "  • Compare MET with CBT and 12-Step approaches for alcohol dependence. (15 marks)",
        "  • Discuss the role of brief motivational interventions in primary care. (15 marks)",
    ],
    notes="These questions reflect common examination patterns in M.Phil Clinical Psychology and PsyD programs. Students should practice answering each within the typical time constraints. For viva questions, aim for 2-3 minute structured answers. For long answers, use a structured format: definition → theoretical basis → key principles → evidence → clinical application → limitations. Reference both the MET manual and NIMHANS manual for Indian context questions.",
    refs="Based on examination patterns from M.Phil Clinical Psychology programs (NIMHANS, RCI-recognized institutions).",
    learning_obj="Prepare for examination questions on MET at postgraduate level",
    takeaway="Prepare structured answers that integrate theory, evidence, and clinical application for each topic."
)


add_slide(
    "MCQs and Short Notes Topics",
    [
        "MULTIPLE CHOICE QUESTIONS (Sample):",
        "  1. DARN-CAT stands for: (a) Desire, Ability, Reasons, Need - Commitment,",
        "     Activation, Taking Steps [CORRECT] (b) Denial, Ambivalence... etc.",
        "  2. The founder of MI is: (a) Aaron Beck (b) William Miller [CORRECT]",
        "     (c) Albert Ellis (d) Carl Rogers",
        "  3. Project MATCH compared MET with: (a) CBT and DBT (b) CBT and TSF [CORRECT]",
        "  4. The recommended reflection:question ratio in MI is: (a) 1:2 (b) 2:1 [CORRECT]",
        "  5. FRAMES includes all EXCEPT: (a) Feedback (b) Responsibility",
        "     (c) Confrontation [CORRECT] (d) Menu",
        "",
        "SHORT NOTES TOPICS (5-7 marks each):",
        "  • OARS in motivational interviewing",
        "  • Stages of change model",
        "  • Self-efficacy in MET",
        "  • FRAMES model of brief intervention",
        "  • Double-sided reflection with examples",
        "  • Decisional balance",
        "  • Change talk vs. sustain talk",
    ],
    notes="MCQs test factual knowledge while short notes require concise explanations with examples. For MCQs, focus on: key terminology (DARN-CAT, OARS, FRAMES), founders/developers, research findings (Project MATCH results), and distinguishing features. For short notes, structure as: definition (1 line), explanation (3-4 lines), clinical example (2-3 lines), and one reference. Practice writing short notes within 5-7 minutes for exam preparation.",
    refs="Based on examination patterns from M.Phil Clinical Psychology, PsyD, and Psychiatry residency programs.",
    learning_obj="Practice examination-style questions on MET",
    takeaway="Know key acronyms (DARN-CAT, OARS, FRAMES), founders, research findings, and clinical applications."
)

add_slide(
    "Clinical Scenario Questions",
    [
        "SCENARIO 1: A 45-year-old man says 'My wife made me come. I don't have a problem.'",
        "  Q: What stage of change is he likely in? How would you respond using MET principles?",
        "  A: Precontemplation. Use: Express empathy, explore his perspective, avoid confrontation,",
        "     emphasize autonomy. Response: 'So this wasn't your idea. Tell me what brought you here.'",
        "",
        "SCENARIO 2: A client says 'I know I should quit but I just can't.'",
        "  Q: Identify the type of change talk. What is the therapeutic focus?",
        "  A: Desire/Need ('should quit') + low Ability ('can't'). Focus: Build self-efficacy.",
        "  Response: 'You clearly want to quit. What makes you feel you can't?'",
        "",
        "SCENARIO 3: After receiving feedback, client says 'These results are scary.'",
        "  Q: How do you respond? What principle are you using?",
        "  A: Reflect the emotion + explore. 'This is really hitting home for you.'",
        "  Principle: Express empathy; the feedback is creating discrepancy naturally.",
        "",
        "SCENARIO 4: Client says 'I've tried quitting 5 times and always failed.'",
        "  Q: Reframe this using MET principles.",
        "  A: 'You've shown real persistence - 5 attempts shows you haven't given up.'",
    ],
    notes="Clinical scenarios test application of principles. Students should identify: (1) stage of change, (2) type of talk (change/sustain), (3) appropriate MET principle, (4) specific response. The key is responding in a way that maintains alliance while moving toward change. Each scenario has multiple acceptable responses - the important thing is that the response is MI-consistent (empathic, non-confrontational, autonomy-supporting).",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Apply MET principles to clinical scenario questions",
    takeaway="For clinical scenarios: identify the stage, type of talk, relevant principle, and formulate an MI-consistent response."
)


# ============================================================
# SECTION 29: SUMMARY (Slides 123-126)
# ============================================================
add_section_divider(29, "Summary and Clinical Pearls")

add_slide(
    "Summary: Key Concepts in MET",
    [
        "CORE IDENTITY: MET is a brief, evidence-based, client-centered yet directive approach",
        "",
        "THEORETICAL BASE: Rogers + Festinger + Bem + Bandura + Prochaska + Deci/Ryan",
        "",
        "SPIRIT: Partnership | Acceptance | Compassion | Evocation",
        "",
        "PRINCIPLES: Express Empathy | Develop Discrepancy | Roll with Resistance | Support Self-Efficacy",
        "",
        "SKILLS: OARS (Open questions, Affirmations, Reflections, Summaries)",
        "",
        "COMPONENTS: FRAMES (Feedback, Responsibility, Advice, Menu, Empathy, Self-efficacy)",
        "",
        "TARGET: Change Talk (DARN-CAT) - evoke it, recognize it, reinforce it",
        "",
        "STRUCTURE: 4 sessions - Feedback → Commitment → Review → Consolidation",
        "",
        "EVIDENCE: 200+ RCTs; Project MATCH (4 sessions = 12 sessions of CBT/TSF)",
        "",
        "APPLICATION: Addiction, health behavior, chronic disease, diverse populations",
    ],
    notes="This summary slide condenses the entire presentation into key bullet points. Use it as a revision aid and as a framework for answering examination questions. Each line represents a major section of the presentation. Students should be able to expand each line into a detailed explanation with clinical examples and references.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). MI (3rd ed.). | Murthy, P. (2008). NIMHANS Manual. | Project MATCH Research Group. (1997). JOSA, 58, 7-29.",
    learning_obj="Integrate all key MET concepts into a coherent summary framework",
    takeaway="MET: Brief + Client-centered + Directive + Evidence-based + Applicable across populations and settings."
)

add_slide(
    "Clinical Pearls: Top 10 Takeaways for Practice",
    [
        "1. The SPIRIT matters more than techniques - be genuine, collaborative, compassionate",
        "2. Resistance is YOUR signal to change approach, not the client's pathology",
        "3. The more clients TALK ABOUT change, the more they BELIEVE in change",
        "4. Ask 'Why are you at X and not zero?' - this ALWAYS evokes change talk",
        "5. Use 'AND' not 'BUT' in double-sided reflections",
        "6. Start with benefits of drinking - this validates and reduces defensiveness",
        "7. Never argue about labels ('alcoholic') - shift to 'What concerns YOU?'",
        "8. Reflections should outnumber questions 2:1 minimum",
        "9. Low importance → develop discrepancy; Low confidence → build self-efficacy",
        "10. The client is the expert on their own life; you facilitate, not prescribe",
        "",
        "BONUS PEARL: If you find yourself working harder than the client, something is wrong.",
        "  The client should be doing most of the talking and most of the arguing for change.",
    ],
    notes="These clinical pearls represent the most important practical lessons from MET training. They should be memorized and regularly reviewed. Each pearl addresses a common mistake or crucial insight. Pearl #12 (bonus) is perhaps the most diagnostic - if the therapist is working harder than the client (arguing for change, giving advice, pushing), the dynamic has inverted. In skilled MET, the CLIENT argues for change while the therapist facilitates the process.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Internalize the most important clinical practice principles of MET",
    takeaway="If you're working harder than the client, something is wrong - they should be arguing for their own change."
)


add_slide(
    "MET Process Flowchart",
    [
        "MET THERAPY PROCESS FLOW:",
        "",
        "PRE-TREATMENT → SESSION 1 → SESSION 2 → SESSION 3 → SESSION 4",
        "Assessment      Feedback    Commitment   Review      Consolidation",
        "                ↓           ↓            ↓           ↓",
        "                Build       Strengthen   Reinforce   Maintain",
        "                Motivation  Commitment   Progress    Change",
        "",
        "WITHIN EACH SESSION:",
        "  Opening → Agenda → Explore → Develop Discrepancy → Elicit Change Talk",
        "  → Reflect → Summarize → Close",
        "",
        "DECISION TREE:",
        "  Client shows change talk? → YES → Move toward commitment",
        "                            → NO  → Continue Phase 1 (build motivation)",
        "",
        "  High resistance? → YES → Roll with it; soften approach; emphasize autonomy",
        "                   → NO  → Continue current strategy; deepen exploration",
    ],
    notes="This flowchart provides a visual map of the MET process. Use it as a quick reference during clinical work. The key decision points are: (1) Is the client showing change talk? If yes, move toward commitment. If no, continue building motivation. (2) Is resistance high? If yes, soften approach. The flowchart also shows the within-session flow: opening → exploration → development → commitment → closing. Each session follows this general pattern while adapting to the client's current stage.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Visualize the MET process as a decision-guided flowchart",
    takeaway="MET follows a clear process flow with key decision points based on client readiness and resistance signals."
)

add_slide(
    "Mind Map: MET at a Glance",
    [
        "                          ┌── Spirit (PACE)────────────────┐",
        "                          │   Partnership, Acceptance,     │",
        "                          │   Compassion, Evocation        │",
        "                          └────────────────────────────────┘",
        "                                        │",
        "        ┌──────────────────────── MET ──────────────────────┐",
        "        │                           │                        │",
        "   Principles                   Skills                  Structure",
        "  (Express Empathy            (OARS)                  (4 Sessions)",
        "   Develop Discrepancy     Open Q's                 S1: Feedback",
        "   Roll w/ Resistance      Affirmations            S2: Commitment",
        "   Support Self-Efficacy)  Reflections             S3: Review",
        "                           Summaries)              S4: Consolidate",
        "        │                           │                        │",
        "   Components                 Targets                 Evidence",
        "   (FRAMES)              (Change Talk               (200+ RCTs",
        "                          DARN-CAT)                Project MATCH)",
    ],
    notes="This mind map organizes all MET components in a visual hierarchy. At the top is the Spirit (most important), branching into Principles, Skills, and Structure. Under each branch are the specific elements. This visual can serve as a study aid and quick reference. Students should be able to fill in this map from memory. The hierarchical structure emphasizes that spirit > principles > techniques.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Organize all MET concepts into a coherent visual mind map",
    takeaway="The MET hierarchy: Spirit (foundation) → Principles (guides) → Skills (tools) → Structure (protocol)."
)


# ============================================================
# SECTION 30: REFERENCES (Slides 127-130)
# ============================================================
add_section_divider(30, "Complete Bibliography")

add_slide(
    "References (1/4): Primary Sources & Foundational Texts",
    [
        "Miller, W.R., Zweben, A., DiClemente, C.C., & Rychtarik, R.G. (1995). Motivational Enhancement",
        "  Therapy Manual. NIAAA, NIH Publication No. 94-3723.",
        "Murthy, P. (2008). Psychosocial Interventions for Persons with Substance Abuse. NIMHANS.",
        "Miller, W.R., & Rollnick, S. (1991). Motivational interviewing: Preparing people to change",
        "  addictive behavior. Guilford Press.",
        "Miller, W.R., & Rollnick, S. (2002). Motivational interviewing (2nd ed.). Guilford Press.",
        "Miller, W.R., & Rollnick, S. (2013). Motivational interviewing: Helping people change (3rd ed.).",
        "  Guilford Press.",
        "Miller, W.R. (1983). Motivational interviewing with problem drinkers. Behavioural Psychotherapy,",
        "  11, 147-172.",
        "Project MATCH Research Group. (1997). Matching alcoholism treatments to client heterogeneity.",
        "  Journal of Studies on Alcohol, 58, 7-29.",
        "Project MATCH Research Group. (1998). Therapist effects in three treatments for alcohol problems.",
        "  Psychotherapy Research, 8, 455-474.",
        "Rollnick, S., Miller, W.R., & Butler, C.C. (2008). Motivational interviewing in health care.",
        "  Guilford Press.",
    ],
    notes="These are the primary sources upon which this presentation is based. The two uploaded reference books (Miller et al., 1995 MET Manual and Murthy, 2008 NIMHANS Manual) form the foundation. Miller & Rollnick's three editions of the MI textbook provide the theoretical evolution. Project MATCH publications provide the evidence base.",
    refs="All references in APA 7th Edition format.",
    learning_obj="Access primary sources for further study of MET",
    takeaway="The MET Manual (Miller et al., 1995) and MI textbook (Miller & Rollnick, 2013) are essential reading."
)

add_slide(
    "References (2/4): Research Studies",
    [
        "Amrhein, P.C., Miller, W.R., Yahne, C.E., et al. (2003). Client commitment language during MI",
        "  predicts drug use outcomes. JCCP, 71, 862-878.",
        "Bien, T.H., Miller, W.R., & Tonigan, S. (1993). Brief interventions for alcohol problems.",
        "  Addiction, 88, 315-336.",
        "Burke, B.L., Arkowitz, H., & Menchola, M. (2003). MI effectiveness meta-analysis. JCCP, 71, 843.",
        "Hettema, J., Steele, J., & Miller, W.R. (2005). MI: Annual Review of Clin. Psychology, 1, 91-111.",
        "Lundahl, B., et al. (2010). A meta-analysis of MI outcomes. Clinical Psychology Review, 30, 1-11.",
        "Miller, W.R., Benefield, R.G., & Tonigan, J.S. (1993). Enhancing motivation for change. JCCP, 61.",
        "Moyers, T.B., Martin, T., et al. (2007). Therapist influence on client language. JCCP, 75, 790-798.",
        "Patterson, G.A., & Forgatch, M.S. (1985). Therapist behavior as determinant. JCCP, 53, 846-851.",
        "Smedslund, G., et al. (2011). MI for substance abuse. Cochrane Database of Systematic Reviews.",
        "UKATT Research Team. (2005). Effectiveness of treatment for alcohol problems. BMJ, 331, 541.",
        "Valle, S.K. (1981). Interpersonal functioning of counselors. J. Studies on Alcohol, 42, 783-790.",
    ],
    notes="These research studies provide the evidence base for MET. Each study is cited multiple times throughout the presentation. Students should be familiar with the key findings of each. For examination purposes, knowing the authors, year, sample size, and main finding of each study is essential.",
    refs="All references in APA 7th Edition format.",
    learning_obj="Access key research studies supporting MET",
    takeaway="Key studies: Project MATCH (1997), Miller et al. (1993), Moyers et al. (2007), Amrhein et al. (2003)."
)


add_slide(
    "References (3/4): Theoretical Foundations",
    [
        "Bandura, A. (1977). Self-efficacy: Toward a unifying theory of behavioral change.",
        "  Psychological Review, 84, 191-215.",
        "Bandura, A. (1982). Self-efficacy mechanism in human agency. American Psychologist, 37, 122-147.",
        "Bem, D.J. (1967). Self-perception: An alternative interpretation of cognitive dissonance.",
        "  Psychological Review, 74, 183-200.",
        "Bem, D.J. (1972). Self-perception theory. In Berkowitz, L. (Ed.), Advances in Experimental",
        "  Social Psychology, Vol. 6. Academic Press.",
        "Deci, E.L., & Ryan, R.M. (1985). Intrinsic motivation and self-determination. Plenum.",
        "Festinger, L. (1957). A Theory of Cognitive Dissonance. Row, Peterson.",
        "Janis, I.L., & Mann, L. (1977). Decision Making. Free Press.",
        "Prochaska, J.O., & DiClemente, C.C. (1982). Transtheoretical therapy. Psychotherapy, 19, 276.",
        "Prochaska, J.O., & DiClemente, C.C. (1984). The Transtheoretical Approach. Dow Jones/Irwin.",
        "Rogers, C.R. (1957). Necessary and sufficient conditions for therapeutic personality change.",
        "  Journal of Consulting Psychology, 21, 95-103.",
        "Ryan, R.M., & Deci, E.L. (2000). Self-determination theory. American Psychologist, 55, 68-78.",
    ],
    notes="These theoretical references provide the foundations upon which MET is built. Each theory contributes a key element: Rogers (empathic relationship), Festinger (cognitive dissonance/discrepancy), Bem (self-perception/change talk), Bandura (self-efficacy), Prochaska & DiClemente (stages of change), Deci & Ryan (autonomy/self-determination). Students should understand how each theory connects to specific MET principles and techniques.",
    refs="All references in APA 7th Edition format.",
    learning_obj="Access foundational theoretical texts underpinning MET",
    takeaway="MET is grounded in established psychological theories: Rogers, Festinger, Bem, Bandura, Prochaska, Deci/Ryan."
)

add_slide(
    "References (4/4): Guidelines, Manuals, and Additional Resources",
    [
        "APA. (2017). Ethical Principles of Psychologists and Code of Conduct. Washington, DC.",
        "DiClemente, C.C. (2003). Addiction and Change. Guilford Press.",
        "Egan, G. (1982). The Skilled Helper. Brooks/Cole.",
        "Marlatt, G.A., & Gordon, J.R. (1985). Relapse Prevention. Guilford Press.",
        "Miller, W.R. (1985). Motivation for treatment: A review. Psychological Bulletin, 98, 84-107.",
        "Miller, W.R., & Moyers, T.B. (2006). Eight stages in learning MI. Behavioural and Cognitive",
        "  Psychotherapy, 34, 135-143.",
        "Moyers, T.B., et al. (2005). MITI Manual (Version 3.0). University of New Mexico.",
        "NICE. (2011). Alcohol-use disorders: Diagnosis, assessment and management. NICE guideline CG115.",
        "SAMHSA. (2019). Enhancing Motivation for Change in Substance Use Disorder Treatment. TIP 35.",
        "Truax, C.B., & Carkhuff, R.R. (1967). Toward Effective Counseling and Psychotherapy. Aldine.",
        "WHO. (2010). Brief Intervention for Substance Use. WHO Press.",
        "Witkiewitz, K., & Marlatt, G.A. (2004). Relapse prevention for alcohol and drug problems.",
        "  Clinical Psychology Review, 24, 1-28.",
    ],
    notes="These additional resources include clinical guidelines (NICE, SAMHSA, WHO), training manuals (MITI), and supplementary texts that support MET practice. SAMHSA TIP 35 is a comprehensive government resource freely available online. The WHO brief intervention manual supports global application. NICE guidelines provide UK-specific clinical recommendations. These resources are essential for continued learning beyond this presentation.",
    refs="All references in APA 7th Edition format.",
    learning_obj="Identify additional resources for continued MET learning and practice",
    takeaway="Key resources: SAMHSA TIP 35, NICE guidelines, WHO brief intervention manual, MITI coding manual."
)


# ============================================================
# ADDITIONAL SLIDES: FILLING TO 140+ SLIDES
# ============================================================

# Additional Section 5 slides
add_slide(
    "Spirit of MET: Wrong vs. Correct Therapist Responses",
    [
        "SCENARIO: Client says 'I like drinking. It relaxes me after a long day.'",
        "",
        "❌ WRONG (Violates Spirit):",
        "  • 'But it's destroying your liver!' (Righting reflex; no partnership)",
        "  • 'You shouldn't use alcohol to cope' (Advice without permission; no acceptance)",
        "  • 'That's just an excuse' (Judgmental; no compassion)",
        "  • 'Research shows that alcohol actually increases anxiety' (Installing; no evocation)",
        "",
        "✓ CORRECT (Honors Spirit):",
        "  • 'Drinking serves an important purpose for you - stress relief.' (Partnership)",
        "  • 'You've found something that works in the short term.' (Acceptance)",
        "  • 'And I wonder what other things you've noticed about how it affects you.' (Evocation)",
        "  • 'What else have you tried for managing stress?' (Evocation + Compassion)",
        "",
        "PRINCIPLE: Honor the positive function before exploring costs",
        "The spirit guides HOW you respond, not just WHAT you say",
    ],
    notes="This practice slide helps students distinguish responses that honor the MET spirit from those that violate it. Each wrong response demonstrates a violation of one or more spirit elements. Each correct response maintains collaboration while still moving toward exploration. The key insight is that acknowledging the positive function of drinking (stress relief) is NOT endorsing the behavior - it's building rapport and trust that will allow deeper exploration later.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Practice responding to clinical material while honoring the four elements of MET spirit",
    takeaway="Honor the positive function of the behavior first; this builds trust for exploring concerns later."
)

# Additional Section 7 slide
add_slide(
    "OARS Practice Exercises",
    [
        "EXERCISE 1: Convert Closed → Open Questions",
        "  Closed: 'Do you drink every day?' → Open: 'Tell me about a typical week with drinking'",
        "  Closed: 'Have you tried quitting?' → Open: 'What attempts have you made to change?'",
        "  Closed: 'Is your family affected?' → Open: 'How has your family noticed your drinking?'",
        "",
        "EXERCISE 2: Write Affirmations for:",
        "  • Client who attended session despite being busy",
        "  • Client who reduced from daily to 3 days/week",
        "  • Client who honestly shared about a relapse",
        "",
        "EXERCISE 3: Practice Reflections (at increasing depth):",
        "  Client: 'I just can't seem to stop once I start'",
        "  Simple: 'Once you start drinking, it's hard to stop'",
        "  Complex: 'There's a sense of being out of control that worries you'",
        "  Feeling: 'That scares you - like you've lost the ability to choose'",
        "",
        "EXERCISE 4: Summarize the session themes you've heard so far",
    ],
    notes="These exercises should be done in pairs or small groups. For Exercise 1, practice converting at least 10 closed questions to open ones. For Exercise 2, write affirmations that are specific, genuine, and focused on the client's qualities. For Exercise 3, practice making reflections at progressively deeper levels - this is the hardest skill and requires extensive practice. The ability to shift between simple and complex reflections based on clinical need is a hallmark of MI proficiency.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Miller, W.R., et al. (1995). MET Manual.",
    learning_obj="Practice OARS skills through structured exercises",
    takeaway="OARS proficiency requires deliberate practice; convert closed to open, deepen reflections, and summarize strategically."
)


# Additional Section 14 slides
add_slide(
    "Session 2: Values Clarification Exercise",
    [
        "VALUES CLARIFICATION - Conducted in Session 2:",
        "",
        "Step 1: 'What are the most important things in your life?' (Open exploration)",
        "Step 2: Present values card sort (or list) - client selects top 5-10 values",
        "Step 3: Rank in order of importance",
        "Step 4: For each top value, explore: 'How is your drinking related to this value?'",
        "",
        "COMMON VALUES IDENTIFIED:",
        "  Family | Health | Career | Integrity | Freedom | Spirituality",
        "  Financial security | Self-respect | Relationships | Achievement",
        "",
        "DEVELOPING DISCREPANCY FROM VALUES:",
        "  T: 'You said being a good father is your #1 value. Tell me about that.'",
        "  C: 'My kids mean everything to me'",
        "  T: 'And where does drinking fit with being the father you want to be?'",
        "  C: '... Not very well. Last week I missed my daughter's recital because I was hungover.'",
        "  T: 'That really doesn't sit right with you - there's a gap between who you want",
        "     to be and how alcohol has been affecting things.'",
    ],
    notes="Values clarification is one of the most powerful tools for developing discrepancy. When clients explicitly name their highest values and then examine how their substance use conflicts with those values, genuine internal motivation emerges. The therapist does NOT point out the discrepancy - instead asks the question ('How does drinking fit with this value?') and lets the client discover it themselves. This creates far more powerful cognitive dissonance than any external confrontation could achieve.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Conduct a values clarification exercise that develops discrepancy naturally",
    takeaway="Values clarification creates powerful internal discrepancy - let clients discover their own value-behavior gaps."
)

# Additional Section 15 slides  
add_slide(
    "Confidence Ruler and Importance Ruler: Detailed Application",
    [
        "IMPORTANCE RULER PROTOCOL:",
        "  T: 'On a scale of 0-10, how important is it for you to change your drinking?'",
        "  C: '6'",
        "  T: 'A 6. And why are you at a 6 and not a 2?' [Evokes change talk!]",
        "  C: 'Because my health is suffering and my wife is threatening to leave'",
        "  T: 'So health and your marriage are pushing you toward thinking change is important'",
        "  T: 'What would it take to move from a 6 to an 8?'",
        "  C: 'If the doctor said my liver is getting worse, probably'",
        "",
        "CONFIDENCE RULER PROTOCOL:",
        "  T: 'And how confident are you that you COULD change? (0-10)'",
        "  C: '3'",
        "  T: 'A 3. What gives you that much confidence - why not a zero?'",
        "  C: 'Because I managed to quit for 3 months last year before my friend's party'",
        "  T: 'So you HAVE done it before. You know you're capable of this.'",
        "",
        "CLINICAL DECISION:",
        "  Importance > Confidence → Focus on building self-efficacy (HOW to change)",
        "  Confidence > Importance → Focus on developing discrepancy (WHY to change)",
    ],
    notes="The rulers are deceptively simple tools with enormous clinical utility. The magic is in the follow-up questions. 'Why X and not lower?' ALWAYS evokes change talk because the client must articulate reasons they have SOME motivation or confidence. Never ask 'Why not higher?' because that evokes sustain talk (reasons they're NOT motivated). The clinical decision matrix (importance vs. confidence) provides clear guidance on therapeutic focus. This can be used in every session as a check-in.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Rollnick, S., et al. (1999). Health Behavior Change. Churchill Livingstone.",
    learning_obj="Master the ruler technique including strategic follow-up questions",
    takeaway="Always ask 'Why not lower?' (evokes change talk); Never ask 'Why not higher?' (evokes sustain talk)."
)


# Additional slides to reach target count

add_slide(
    "MET in the NIMHANS Stepped Care Model",
    [
        "NIMHANS STEPPED CARE APPROACH (Murthy, 2008):",
        "",
        "STEP 1: IDENTIFICATION (Screening by primary care physician)",
        "  • Ask about substance use as routine (like diabetes/hypertension)",
        "  • 5-A Strategy: Ask, Assess, Advise, Assist, Arrange",
        "",
        "STEP 2: BRIEF INTERVENTION (FRAMES-based)",
        "  • Single session: Feedback + Empathy + Menu of options",
        "  • Suitable for hazardous/harmful use (not yet dependent)",
        "",
        "STEP 3: MET/MI (2-4 sessions)",
        "  • For moderate dependence or ambivalent clients",
        "  • Structured feedback + motivational strategies",
        "",
        "STEP 4: SPECIALIZED TREATMENT",
        "  • For severe dependence; complex comorbidity",
        "  • Refer to de-addiction centers, psychiatrists",
        "",
        "KEY POINT: MET bridges brief intervention and intensive treatment",
        "NIMHANS: 'Not practical to limit care only to specialized centers'",
    ],
    notes="The NIMHANS manual situates MET within a broader stepped care framework appropriate for the Indian healthcare system. Step 1 involves routine screening in primary care - the 5-A strategy makes this systematic. Step 2 is a single brief intervention session using FRAMES. Step 3 is where MET fits - for those needing more than brief intervention but not requiring intensive treatment. Step 4 is specialized care. This model maximizes resource efficiency by matching treatment intensity to need.",
    refs="Murthy, P. (2008). NIMHANS Manual, pp. 1-3. | WHO. (2010). Brief Intervention for Substance Use. | Babor, T.F., & Higgins-Biddle, J.C. (2001). WHO Manual.",
    learning_obj="Position MET within the NIMHANS stepped care model for substance use",
    takeaway="MET occupies Step 3 in NIMHANS stepped care - bridging brief intervention and intensive treatment."
)

add_slide(
    "DARES: The NIMHANS Motivational Interviewing Principles",
    [
        "NIMHANS Manual uses the DARES acronym for MI principles:",
        "",
        "D - Develop Discrepancy",
        "  Help client see gap between current behavior and valued goals",
        "  'Where do you want to be? Where are you now?'",
        "",
        "A - Avoid Argumentation",
        "  Never argue or confront; this increases resistance",
        "  Arguments about labels are especially counterproductive",
        "",
        "R - Roll with Resistance",
        "  Use client's momentum; redirect rather than oppose",
        "  'You don't feel ready right now, and that's okay'",
        "",
        "E - Express Empathy",
        "  Warm, reflective listening; genuine understanding",
        "  Acceptance of ambivalence as normal",
        "",
        "S - Support Self-Efficacy",
        "  Boost confidence; highlight strengths and past successes",
        "  'You've shown you CAN make changes in your life'",
    ],
    notes="The NIMHANS manual uses DARES (slightly different order from Miller & Rollnick's presentation) as a memorable acronym for MI principles. This is identical in content to the core MET principles but organized mnemonically for training purposes. Note that 'Avoid Argumentation' replaces the original 'Roll with Resistance' was later updated in Miller & Rollnick's 3rd edition where they replaced these four principles with the four processes of MI. However, DARES remains widely used in training.",
    refs="Murthy, P. (2008). NIMHANS Manual, p. 21. | Miller, W.R., & Rollnick, S. (1991). Motivational interviewing. Guilford Press.",
    learning_obj="Apply the DARES framework from the NIMHANS manual",
    takeaway="DARES (Develop discrepancy, Avoid argumentation, Roll with resistance, Express empathy, Support self-efficacy)."
)


add_slide(
    "Eliciting Change Talk: The NIMHANS Approach",
    [
        "NIMHANS MANUAL - Eliciting Change Talk involves:",
        "",
        "1. Recognizing DISADVANTAGES of staying the same:",
        "  'What worries you about your current drinking pattern?'",
        "  'What might happen if things continue as they are?'",
        "",
        "2. Recognizing ADVANTAGES of change:",
        "  'If you did make this change, what would be different?'",
        "  'What would be the best things about cutting down?'",
        "",
        "3. Expressing OPTIMISM about change:",
        "  'What gives you hope that you could change?'",
        "  'What strengths do you have that might help?'",
        "",
        "4. Expressing INTENT to change:",
        "  'What do you think you might do?'",
        "  'What's your next step?'",
        "",
        "NIMHANS: 'Clients are encouraged for clarifications and elaborations'",
        "KEY: 'Goals and values are explored to find discrepancies'",
    ],
    notes="The NIMHANS manual outlines a four-step process for eliciting change talk that parallels the DARN model. Step 1 (disadvantages of status quo) evokes Reasons and Need. Step 2 (advantages of change) evokes Desire. Step 3 (optimism) evokes Ability. Step 4 (intent) evokes Commitment. This provides a practical sequence for therapists to follow during sessions. The emphasis on encouraging elaboration and exploring goals/values to find discrepancies aligns with core MET practice.",
    refs="Murthy, P. (2008). NIMHANS Manual, p. 23. | Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press.",
    learning_obj="Follow the NIMHANS four-step process for eliciting change talk",
    takeaway="Elicit change talk by exploring: disadvantages of status quo, advantages of change, optimism, and intent."
)

add_slide(
    "MET and Relapse Prevention Integration",
    [
        "RELAPSE PREVENTION IN MET (Sessions 3-4):",
        "",
        "MARLATT & GORDON (1985) MODEL:",
        "  High-risk situation → Coping response → Outcome",
        "  If NO coping response → Decreased self-efficacy → Lapse → Relapse",
        "  If coping response → Increased self-efficacy → Continued abstinence/reduction",
        "",
        "MET-COMPATIBLE RELAPSE PREVENTION STRATEGIES:",
        "  1. Identify high-risk situations (triggers, cravings)",
        "  2. Develop coping plans (using collaborative, non-prescriptive style)",
        "  3. Build self-efficacy for coping (affirm past successes)",
        "  4. Normalize lapses (not failures, just opportunities to learn)",
        "  5. Emergency planning (who to call, what to do)",
        "",
        "KEY PRINCIPLE: Use MET SPIRIT even during RP planning",
        "  Not: 'You MUST avoid bars' (directive)",
        "  But: 'What situations do you think might be challenging, and how might you handle them?'",
        "",
        "NIMHANS: 'Strategies to anticipate and prevent relapses can be taught'",
    ],
    notes="Relapse prevention is incorporated into MET's later sessions but must maintain the MET spirit. Rather than prescribing coping strategies, the therapist explores with the client what situations might be risky and how they might handle them. This collaborative approach to RP planning increases ownership and adherence. The Marlatt & Gordon model shows that self-efficacy is the key mediator - MET's emphasis on building self-efficacy naturally supports relapse prevention.",
    refs="Marlatt, G.A., & Gordon, J.R. (1985). Relapse Prevention. Guilford Press. | Witkiewitz, K., & Marlatt, G.A. (2004). Clinical Psychology Review, 24, 1-28. | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Integrate relapse prevention strategies within the MET framework and spirit",
    takeaway="Relapse prevention in MET uses collaborative exploration of high-risk situations, not prescriptive directives."
)


add_slide(
    "MET for Opioid Dependence: Case Application",
    [
        "CASE: Suresh, 28, Daily wage worker, Opioid (heroin) dependence x 5 years",
        "  Presenting: Injecting heroin; wife left; job loss; considering suicide",
        "  Stage: Contemplation (wants to change but feels hopeless)",
        "  Importance: 9/10 | Confidence: 1/10",
        "",
        "MET APPROACH:",
        "  Focus: BUILD SELF-EFFICACY (importance already high)",
        "  Session 1: Acknowledge desperation; safety assessment; feedback (medical)",
        "    T: 'You've been through so much and you're still fighting. That shows strength.'",
        "  Session 2: Explore past attempts; what worked even briefly?",
        "    Build on any success; discuss OST (Opioid Substitution Therapy) as option",
        "  Session 3: Integrated approach - OST + MET; rebuilding life structure",
        "  Session 4: Maintenance; social reintegration; family reconnection plan",
        "",
        "INTEGRATION WITH PHARMACOTHERAPY:",
        "  MET + Buprenorphine/Methadone = better retention and outcomes",
        "  MET enhances medication adherence and engagement",
        "NIMHANS: Stepped care with appropriate referral for severe cases",
    ],
    notes="This case demonstrates MET with severe opioid dependence. Key insight: importance is already very high (9/10) so developing discrepancy is unnecessary and potentially harmful (could increase hopelessness). The focus must be entirely on building self-efficacy (confidence = 1/10). MET is combined with pharmacotherapy (OST) which the NIMHANS manual supports. The therapist's role is to instill hope, identify even tiny past successes, and support engagement with medication.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Murthy, P. (2008). NIMHANS Manual. | Carroll, K.M., et al. (2006). Drug and Alcohol Dependence, 81, 161-167.",
    learning_obj="Apply MET to opioid dependence with integrated pharmacotherapy",
    takeaway="When importance is high but confidence is low, focus entirely on building self-efficacy and hope."
)

add_slide(
    "MET with Mandated/Court-Ordered Clients",
    [
        "CHALLENGES WITH MANDATED CLIENTS:",
        "  • Often precontemplators - not seeking help voluntarily",
        "  • May be angry, hostile, or resigned",
        "  • Perceived coercion can undermine intrinsic motivation",
        "  • Therapist may be seen as extension of the punishing system",
        "",
        "MET STRATEGIES FOR MANDATED CLIENTS:",
        "  1. Acknowledge the situation honestly:",
        "     'You're here because the court required it. How do you feel about that?'",
        "  2. Separate the mandate from the person:",
        "     'You have to be here, but what you get out of it is completely up to you'",
        "  3. Emphasize what IS within their control:",
        "     'You can't control the court order, but you CAN choose what to do with this time'",
        "  4. Find ANY genuine motivation:",
        "     'Setting aside the court, is there anything about your drinking YOU'd like to change?'",
        "  5. Use gentle paradox if appropriate:",
        "     'I'm not sure you're ready to think about this yet...'",
        "",
        "RESEARCH: MET is ESPECIALLY effective with mandated/reluctant clients",
        "  (because it reduces reactance and respects autonomy)",
    ],
    notes="Mandated clients are often considered the most difficult, but MET is actually well-suited for them because it addresses the reactance that mandated treatment creates. By emphasizing autonomy ('What you do here is your choice'), the therapist paradoxically increases engagement. The key is to find ANY genuine motivation beyond the mandate - even a small personal reason for change provides a foundation. Research shows that MI/MET is particularly effective with angry, resistant clients compared to other approaches.",
    refs="Miller, W.R., et al. (1995). MET Manual, Special Problems section. | Moyers, T.B., & Rollnick, S. (2002). Substance Use and Misuse, 37, 2089-2102.",
    learning_obj="Adapt MET strategies for mandated/involuntary treatment clients",
    takeaway="MET is especially effective with mandated clients because it reduces reactance by emphasizing autonomy."
)


add_slide(
    "Mechanisms of MET: How Does It Work?",
    [
        "THE CAUSAL CHAIN (Research-Supported):",
        "  Therapist MI-Consistent Behavior → Client Change Talk → Behavior Change",
        "  (Moyers et al., 2007; 2009)",
        "",
        "PROPOSED MECHANISMS:",
        "  1. Self-Perception: Hearing oneself argue for change creates belief in change",
        "  2. Cognitive Dissonance: Awareness of value-behavior gap creates discomfort → action",
        "  3. Self-Efficacy: Therapist confidence + past success reflection → belief in ability",
        "  4. Therapeutic Alliance: Empathic relationship → safety → exploration → insight",
        "  5. Autonomy Support: Feeling in control → intrinsic motivation (SDT)",
        "  6. Commitment Strengthening: Verbal commitment → behavioral follow-through",
        "",
        "MEDIATORS IDENTIFIED IN RESEARCH:",
        "  • Client change talk strength (Amrhein et al., 2003)",
        "  • Working alliance (Boardman et al., 2006)",
        "  • Readiness to change (DiClemente et al., 2017)",
        "",
        "WHAT WE DON'T KNOW YET:",
        "  • Exact dose-response relationship",
        "  • Which specific techniques are most active",
        "  • Why effects diminish for some clients over time",
    ],
    notes="Understanding mechanisms helps refine practice. The best-supported mechanism is the Moyers/Amrhein causal chain: therapist behavior influences client in-session language, which predicts out-of-session behavior change. This has been replicated across multiple studies. The clinical implication is clear: focus on evoking and reinforcing change talk. Other proposed mechanisms (self-perception, dissonance, self-efficacy, alliance) are theoretically supported but have less direct empirical evidence of mediation.",
    refs="Moyers, T.B., et al. (2007). JCCP, 75, 790-798. | Amrhein, P.C., et al. (2003). JCCP, 71, 862-878. | Magill, M., et al. (2014). Journal of Substance Abuse Treatment, 46, 685-697.",
    learning_obj="Understand the research-supported mechanisms through which MET produces change",
    takeaway="The primary mechanism: Therapist MI-consistent behavior → Client change talk → Behavior change."
)

add_slide(
    "MET Treatment Fidelity: Ensuring Quality",
    [
        "WHY FIDELITY MATTERS:",
        "  • Poorly delivered MI can be equivalent to or worse than no treatment",
        "  • Therapist drift is common without monitoring",
        "  • Quality determines effectiveness",
        "",
        "FIDELITY INDICATORS (MITI 4.2 Coding):",
        "  TECHNICAL COMPONENT:",
        "  • Cultivating change talk (evoking reasons, desires, needs)",
        "  • Softening sustain talk (not reinforcing reasons against change)",
        "  RELATIONAL COMPONENT:",
        "  • Partnership (collaborative vs. expert)",
        "  • Empathy (understanding client perspective)",
        "",
        "BEHAVIORAL COUNTS:",
        "  • Giving information / Persuade / Persuade with permission",
        "  • Questions (open vs. closed)",
        "  • Reflections (simple vs. complex)",
        "  • Affirm / Seek collaboration / Emphasize autonomy / Confront",
        "",
        "THRESHOLDS: Reflection:Question ≥ 2:1 | Complex reflections ≥ 40% | Open Q's ≥ 70%",
    ],
    notes="Treatment fidelity is the bridge between research efficacy and clinical effectiveness. The MITI (Motivational Interviewing Treatment Integrity) coding system provides standardized assessment. Sessions are recorded and coded by trained raters. Key metrics include reflection-to-question ratio, percentage of complex reflections, and presence of MI-consistent vs. MI-inconsistent behaviors. Regular fidelity monitoring with feedback is essential for maintaining quality in clinical practice.",
    refs="Moyers, T.B., et al. (2014). MITI Manual 4.2. University of New Mexico. | Miller, W.R., & Moyers, T.B. (2006). Behavioural and Cognitive Psychotherapy, 34, 135-143.",
    learning_obj="Assess MET treatment fidelity using the MITI coding system",
    takeaway="Fidelity monitoring ensures quality: aim for reflection:question ≥ 2:1, complex reflections ≥ 40%, open Q's ≥ 70%."
)


add_slide(
    "Group Exercise: Identifying Change Talk vs. Sustain Talk",
    [
        "IDENTIFY EACH STATEMENT AS CHANGE TALK (CT) OR SUSTAIN TALK (ST):",
        "",
        "1. 'I know I need to cut down' → CT (Need)",
        "2. 'All my friends drink, I can't be the odd one out' → ST (Reasons to continue)",
        "3. 'My wife says she'll leave if I don't change' → CT (Reasons)",
        "4. 'I've tried before and it never works' → ST (Low ability)",
        "5. 'I wish I could drink like a normal person' → CT (Desire)",
        "6. 'Life would be so boring without alcohol' → ST (Cost of change)",
        "7. 'I'm going to start by having alcohol-free Mondays' → CT (Taking Steps)",
        "8. 'I'm not ready to think about this yet' → ST (No activation)",
        "9. 'I could probably cut down if I really tried' → CT (Ability)",
        "10. 'I'm willing to give it a go for a month' → CT (Activation/Commitment)",
        "",
        "THERAPIST RESPONSE PRACTICE:",
        "  For each CT: How would you REINFORCE it? (Reflection, affirmation, elaboration)",
        "  For each ST: How would you RESPOND without strengthening it?",
    ],
    notes="This group exercise builds the crucial skill of recognizing change talk and sustain talk in real-time. In actual sessions, this recognition must be automatic. Each statement should be classified by type (DARN-CAT for change talk; reasons/desire/ability/need to maintain status quo for sustain talk). Then students practice responding: for change talk, reflect it warmly and ask for elaboration ('Tell me more about that'). For sustain talk, acknowledge briefly without dwelling and redirect.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Moyers, T.B., et al. (2009). Journal of Substance Abuse Treatment, 36, 101-109.",
    learning_obj="Practice real-time identification of change talk and sustain talk",
    takeaway="Rapid recognition of change talk vs. sustain talk is essential; respond differently to each."
)

add_slide(
    "MET and Comorbid Psychiatric Disorders",
    [
        "DUAL DIAGNOSIS CONSIDERATIONS:",
        "  • 50-70% of substance users have comorbid psychiatric disorders",
        "  • Depression, anxiety, PTSD, psychosis, personality disorders",
        "  • MET addresses substance use; may need integration with other treatments",
        "",
        "DEPRESSION + SUBSTANCE USE:",
        "  • Alcohol worsens depression; depression drives self-medication",
        "  • MET focus: Develop discrepancy (depression-alcohol cycle)",
        "  • Integration: MET + pharmacotherapy + CBT/IPT for depression",
        "",
        "ANXIETY + SUBSTANCE USE:",
        "  • Anxiety drives use (self-medication); withdrawal creates more anxiety",
        "  • MET focus: Build self-efficacy for alternative coping",
        "",
        "PSYCHOSIS + SUBSTANCE USE:",
        "  • Requires simplified, shorter sessions; concrete feedback",
        "  • Evidence: MI effective for medication adherence in psychosis",
        "",
        "PTSD + SUBSTANCE USE:",
        "  • Trauma-informed approach essential",
        "  • MET can motivate engagement with trauma treatment",
        "  • NIMHANS: Appropriate referral for complex cases",
    ],
    notes="Dual diagnosis is the norm rather than the exception in substance use treatment. MET can be applied to the substance use component while integrated with disorder-specific treatments. The key adaptation is recognizing which comorbidity is driving the use and tailoring the discrepancy development accordingly. For depression: the cycle of alcohol worsening depression is a powerful discrepancy. For anxiety: building confidence in alternative coping methods is crucial. For psychosis: simplify and shorten sessions. For PTSD: trauma-informed modifications are essential.",
    refs="Baker, A.L., et al. (2012). Clinical Psychology Review, 32, 726-738. | Murthy, P. (2008). NIMHANS Manual. | Drake, R.E., et al. (2004). Journal of Substance Abuse Treatment, 27, 197-213.",
    learning_obj="Adapt MET for clients with comorbid psychiatric disorders",
    takeaway="MET addresses the substance component in dual diagnosis; integration with disorder-specific treatment is essential."
)


add_slide(
    "Future Directions in MET Research and Practice",
    [
        "EMERGING RESEARCH AREAS:",
        "  • Technology-assisted MI (telehealth, AI chatbots, apps)",
        "  • Neuroimaging studies of MI effects on brain reward circuits",
        "  • Personalized/precision MET based on client characteristics",
        "  • Integration with genomic/pharmacogenomic data",
        "  • MET for behavioral addictions (gaming, social media, gambling)",
        "",
        "IMPLEMENTATION SCIENCE:",
        "  • How to scale MET training effectively across healthcare systems",
        "  • Maintaining fidelity in real-world (non-research) settings",
        "  • Cost-effectiveness analyses in different healthcare systems",
        "  • Task-shifting to non-specialists in low-resource settings (NIMHANS model)",
        "",
        "CLINICAL INNOVATIONS:",
        "  • Group-based MET formats",
        "  • Single-session MI (e.g., emergency department interventions)",
        "  • MI combined with contingency management",
        "  • MI for medication-assisted treatment engagement",
        "  • Culturally-adapted MI protocols for diverse populations",
        "",
        "THE FUTURE: MET principles integrated into standard clinical training worldwide",
    ],
    notes="The future of MET/MI is dynamic. Technology offers new delivery methods (apps, telehealth, even AI-assisted interventions), while neuroscience provides deeper understanding of mechanisms. Implementation science addresses the gap between research efficacy and real-world effectiveness. The NIMHANS model of training non-specialists is increasingly relevant globally. Behavioral addictions (gaming, social media) represent new application areas. The vision is MI principles becoming a foundational clinical skill for ALL healthcare professionals.",
    refs="Miller, W.R., & Rollnick, S. (2013). Motivational interviewing (3rd ed.). Guilford Press. | Lundahl, B., et al. (2013). Patient Education and Counseling, 93, 157-168. | Apodaca, T.R., & Longabaugh, R. (2009). Clinical Psychology Review, 29, 199-218.",
    learning_obj="Identify future directions in MET research and clinical application",
    takeaway="MET's future includes technology-assisted delivery, precision approaches, and integration into standard healthcare training."
)

# Final closing slide
add_slide(
    "Thank You and Final Reflections",
    [
        "MOTIVATIONAL ENHANCEMENT THERAPY",
        "",
        "'People are generally the best experts on themselves.'",
        "  — William R. Miller",
        "",
        "'If you find yourself in a struggle with the patient, you are doing it wrong.'",
        "  — Adapted from Miller & Rollnick, 2013",
        "",
        "'The therapist's task is to tip the motivational balance in favor of change.'",
        "  — Miller et al., 1995 (MET Manual)",
        "",
        "'Brief interventions provide a method of health promotion and disease prevention.'",
        "  — Murthy, 2008 (NIMHANS Manual)",
        "",
        "REMEMBER: The spirit of MET is its most powerful element.",
        "  Partnership. Acceptance. Compassion. Evocation.",
        "",
        "Every interaction is an opportunity to enhance motivation for positive change.",
    ],
    notes="Close the presentation with these memorable quotes that encapsulate the MET philosophy. Remind students that they are now equipped with knowledge that can genuinely help people change. The spirit of MET - genuine care, respect, and belief in people's capacity for change - is both a clinical tool and a way of being. Encourage students to practice these skills in their everyday interactions, not just in formal therapy sessions.",
    refs="Miller, W.R., et al. (1995). MET Manual. | Miller, W.R., & Rollnick, S. (2013). MI (3rd ed.). | Murthy, P. (2008). NIMHANS Manual.",
    learning_obj="Integrate the spirit of MET as a professional value and clinical competency",
    takeaway="The spirit of MET - partnership, acceptance, compassion, evocation - transforms both therapy and the therapist."
)


# ============================================================
# SAVE THE PRESENTATION
# ============================================================
output_path = '/projects/sandbox/Dango-kiro/MET_Comprehensive_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
