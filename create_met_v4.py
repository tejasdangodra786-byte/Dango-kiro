#!/usr/bin/env python3
"""
MET Comprehensive Presentation v4
Based primarily on: Motivational Enhancement Therapy Manual (NIAAA Project MATCH)
                    Substance Use Disorders (NIMHANS Publication)
With case conceptualization and detailed worksheets.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
import os

# Slide dimensions - widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color palette - vibrant and professional
COLORS = {
    'deep_blue': RGBColor(0x1B, 0x3A, 0x5C),
    'teal': RGBColor(0x00, 0x7B, 0x83),
    'green': RGBColor(0x2E, 0x7D, 0x32),
    'orange': RGBColor(0xE6, 0x5C, 0x00),
    'purple': RGBColor(0x6A, 0x1B, 0x9A),
    'red': RGBColor(0xC6, 0x28, 0x28),
    'gold': RGBColor(0xF9, 0xA8, 0x25),
    'light_blue': RGBColor(0xBB, 0xDE, 0xFB),
    'light_green': RGBColor(0xC8, 0xE6, 0xC9),
    'light_purple': RGBColor(0xE1, 0xBE, 0xE7),
    'light_orange': RGBColor(0xFF, 0xE0, 0xB2),
    'light_teal': RGBColor(0xB2, 0xDF, 0xDB),
    'light_red': RGBColor(0xFF, 0xCD, 0xD2),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'cream': RGBColor(0xFF, 0xF8, 0xE1),
    'navy': RGBColor(0x0D, 0x47, 0xA1),
    'dark_green': RGBColor(0x1B, 0x5E, 0x20),
    'maroon': RGBColor(0x88, 0x00, 0x38),
}

# Section colors for variety
SECTION_COLORS = [
    ('deep_blue', 'light_blue'),
    ('teal', 'light_teal'),
    ('green', 'light_green'),
    ('purple', 'light_purple'),
    ('orange', 'light_orange'),
    ('maroon', 'light_red'),
]


prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]

def add_background(slide, color_key):
    """Add a colored background rectangle to slide."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS[color_key]
    bg.line.fill.background()
    # Send to back
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def add_gradient_bg(slide, color1_key, color2_key):
    """Add gradient background."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.gradient()
    bg.fill.gradient_stops[0].color.rgb = COLORS[color1_key]
    bg.fill.gradient_stops[1].color.rgb = COLORS[color2_key]
    bg.line.fill.background()
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def set_text_props(tf, text, font_size=14, bold=False, color_key='dark_gray', alignment=PP_ALIGN.LEFT):
    """Set text frame properties."""
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color_key]
    p.font.name = 'Times New Roman'
    p.alignment = alignment

def add_title_bar(slide, title_text, color_key='deep_blue'):
    """Add a colored title bar at top of slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[color_key]
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Times New Roman'

def add_content_box(slide, left, top, width, height, text_lines, bg_color='white',
                    text_color='dark_gray', font_size=14, title=None, title_color='deep_blue',
                    border_color=None):
    """Add a colored content box with text."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS[bg_color]
    if border_color:
        box.line.color.rgb = COLORS[border_color]
        box.line.width = Pt(2)
    else:
        box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS[title_color]
        p.font.name = 'Times New Roman'
        p.space_after = Pt(6)
    start_idx = 1 if title else 0
    for i, line in enumerate(text_lines):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS[text_color]
        p.font.name = 'Times New Roman'
        p.space_after = Pt(4)


def add_reference_bar(slide, ref_text):
    """Add reference bar at bottom of slide."""
    ref_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), SLIDE_WIDTH, Inches(0.6)
    )
    ref_box.fill.solid()
    ref_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    ref_box.line.fill.background()
    tf = ref_box.text_frame
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = ref_text
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.font.name = 'Times New Roman'

def add_section_divider(title, subtitle, color_key='deep_blue'):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(blank_layout)
    add_gradient_bg(slide, color_key, 'white')
    # Large centered title
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Times New Roman'
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLORS['cream']
    p2.font.name = 'Times New Roman'
    p2.alignment = PP_ALIGN.CENTER
    return slide

def add_bullet_content(slide, left, top, width, height, bullets, font_size=14,
                       color_key='dark_gray', bullet_char='\u2022'):
    """Add bulleted text content."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS[color_key]
        p.font.name = 'Times New Roman'
        p.space_after = Pt(6)


def add_table_slide(title, headers, rows, color_key='deep_blue', ref_text=''):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(blank_layout)
    add_gradient_bg(slide, 'white', 'light_blue')
    add_title_bar(slide, title, color_key)
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl_left = Inches(0.5)
    tbl_top = Inches(1.3)
    tbl_width = Inches(12.3)
    tbl_height = Inches(5.2)
    
    table = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height).table
    
    # Set column widths evenly
    col_width = int(Inches(12.3) / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width
    
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS[color_key]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = 'Times New Roman'
    
    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['dark_gray']
            p.font.name = 'Times New Roman'
    
    if ref_text:
        add_reference_bar(slide, ref_text)
    return slide

def add_process_flow(slide, left, top, width, items, color_keys):
    """Add a horizontal process flow with arrows."""
    n = len(items)
    box_w = (width - Inches(0.3) * (n - 1)) / n
    for i, (item, ck) in enumerate(zip(items, color_keys)):
        x = left + (box_w + Inches(0.3)) * i
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(x), top, int(box_w), Inches(1.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS[ck]
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = 'Times New Roman'
        p.alignment = PP_ALIGN.CENTER
        # Arrow between boxes
        if i < n - 1:
            arr_x = int(x + box_w)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arr_x, top + Inches(0.4), Inches(0.3), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['gold']
            arrow.line.fill.background()


# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'deep_blue', 'navy')

# Title
tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(2))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "MOTIVATIONAL ENHANCEMENT THERAPY (MET)"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.font.name = 'Times New Roman'
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "A Comprehensive Clinical Guide"
p2.font.size = Pt(24)
p2.font.color.rgb = COLORS['gold']
p2.font.name = 'Times New Roman'
p2.alignment = PP_ALIGN.CENTER

p3 = tf.add_paragraph()
p3.text = ""
p3.space_after = Pt(20)

p4 = tf.add_paragraph()
p4.text = "Based on the NIAAA Project MATCH MET Manual"
p4.font.size = Pt(16)
p4.font.color.rgb = COLORS['light_blue']
p4.font.name = 'Times New Roman'
p4.alignment = PP_ALIGN.CENTER

p5 = tf.add_paragraph()
p5.text = "& NIMHANS Substance Use Disorders Publication"
p5.font.size = Pt(16)
p5.font.color.rgb = COLORS['light_blue']
p5.font.name = 'Times New Roman'
p5.alignment = PP_ALIGN.CENTER

# Decorative box
deco = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(5.5), Inches(7), Inches(1.2)
)
deco.fill.solid()
deco.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
deco.line.color.rgb = COLORS['gold']
deco.line.width = Pt(2)
tf2 = deco.text_frame
tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf2.paragraphs[0]
p.text = "4-Session Brief Intervention | Evidence-Based | Client-Centered"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['gold']
p.font.name = 'Times New Roman'
p.alignment = PP_ALIGN.CENTER

add_reference_bar(slide, "References: Miller, W.R. et al. (1992). MET Manual, NIAAA Project MATCH Series Vol. 2; NIMHANS (2016). Substance Use Disorders.")


# ============================================================
# SLIDE 2: TABLE OF CONTENTS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "PRESENTATION OUTLINE", 'deep_blue')

sections = [
    ("1. Introduction to MET", "What is MET, its origins, and core philosophy"),
    ("2. Theoretical Foundations", "Transtheoretical Model, Self-Efficacy, Cognitive Dissonance"),
    ("3. Principles & Spirit of MET", "FRAMES, OARS, DARES, and MI Spirit"),
    ("4. Session-by-Session Guide", "Detailed 4-session protocol with techniques"),
    ("5. Clinical Techniques", "Reflective listening, eliciting change talk, handling resistance"),
    ("6. Case Conceptualization", "Applied case study with session-wise intervention"),
    ("7. Worksheets & Tools", "Printable clinical worksheets for sessions"),
    ("8. Research & Effectiveness", "Project MATCH, meta-analyses, Indian research"),
]

for i, (sec_title, sec_desc) in enumerate(sections):
    y = Inches(1.4) + Inches(0.68) * i
    color_idx = i % len(SECTION_COLORS)
    box_color = SECTION_COLORS[color_idx][1]
    text_color = SECTION_COLORS[color_idx][0]
    add_content_box(slide, Inches(0.8), y, Inches(11.5), Inches(0.6),
                    [sec_desc], bg_color=box_color, text_color=text_color,
                    font_size=12, title=sec_title, title_color=text_color)

add_reference_bar(slide, "Structure adapted from: Miller et al. (1992). MET Manual, NIAAA Project MATCH Monograph Series, Vol. 2.")


# ============================================================
# SECTION 1: INTRODUCTION TO MET
# ============================================================
add_section_divider("SECTION 1", "Introduction to Motivational Enhancement Therapy", 'deep_blue')

# SLIDE: What is MET?
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'cream')
add_title_bar(slide, "What is Motivational Enhancement Therapy (MET)?", 'deep_blue')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "MET is a brief, systematic intervention",
        "designed to produce rapid, internally",
        "motivated change in addictive behaviors.",
        "",
        "It does NOT guide the client step-by-step",
        "through recovery. Instead, it uses motivational",
        "strategies to mobilize the client's OWN",
        "change resources and motivation.",
    ],
    bg_color='light_blue', title="Definition", title_color='deep_blue',
    border_color='deep_blue', font_size=14)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "- Developed for NIAAA Project MATCH",
        "- Based on principles of Motivational Interviewing",
        "- 4 planned sessions over 12 weeks",
        "- Therapist uses personal feedback + MI techniques",
        "- Designed for alcohol/substance use disorders",
        "- Adaptable to various clinical settings",
    ],
    bg_color='light_green', title="Key Facts", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.3),
    [
        "\"MET is based on principles of motivational psychology and is designed to produce rapid,",
        "internally motivated change. This treatment strategy does not attempt to guide and train the",
        "client, step by step, through recovery, but instead employs motivational strategies to mobilize",
        "the client's own change resources.\" (MET Manual, p. 1)",
    ],
    bg_color='cream', title="From the Manual:", title_color='orange',
    border_color='gold', font_size=13)

add_reference_bar(slide, "Reference: Miller, W.R., Zweben, A., DiClemente, C.C., & Rychtarik, R.G. (1992). MET Manual, p. 1. NIAAA Project MATCH.")


# SLIDE: How MET Differs from Other Approaches
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "How MET Differs from Other Approaches", 'green')

add_table_slide("How MET Differs from Other Approaches",
    ["Feature", "Traditional Approaches", "MET Approach"],
    [
        ["Therapist Role", "Expert/Teacher/Director", "Collaborative partner, elicits client's own motivation"],
        ["Client Role", "Passive recipient of treatment", "Active agent of own change"],
        ["Resistance", "Seen as denial to be confronted", "Signal to change therapeutic strategy"],
        ["Goals", "Set by therapist/program", "Negotiated with client based on their values"],
        ["Techniques", "Skills training, education", "Reflective listening, feedback, exploring ambivalence"],
        ["Duration", "Usually 12+ sessions", "4 structured sessions over 12 weeks"],
        ["Mechanism", "Learning new skills", "Mobilizing client's own internal resources"],
    ],
    color_key='green',
    ref_text="Reference: Miller et al. (1992). MET Manual, pp. 1-5; NIMHANS (2016). Substance Use Disorders, Ch. 8.")

# SLIDE: Origins and Development
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_purple')
add_title_bar(slide, "Origins and Development of MET", 'purple')

add_content_box(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(3),
    [
        "1983: William R. Miller publishes foundational",
        "  paper on Motivational Interviewing (MI)",
        "",
        "1991: Miller & Rollnick publish first MI book",
        "",
        "1992: MET Manual developed for Project MATCH",
        "  - Largest alcohol treatment trial ever conducted",
        "  - 1,726 participants across 9 sites",
        "  - Compared MET vs CBT vs 12-Step Facilitation",
    ],
    bg_color='light_purple', title="Timeline", title_color='purple',
    border_color='purple', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(3),
    [
        "MET integrates two core components:",
        "",
        "1. MOTIVATIONAL INTERVIEWING (MI)",
        "   - Client-centered counseling style",
        "   - Explores and resolves ambivalence",
        "",
        "2. PERSONALIZED ASSESSMENT FEEDBACK",
        "   - Objective data presented to client",
        "   - Creates discrepancy between current",
        "     behavior and personal goals/values",
    ],
    bg_color='cream', title="Core Components", title_color='orange',
    border_color='orange', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.8),
    [
        "Project MATCH Finding: MET achieved comparable outcomes to 12-session CBT and 12-session",
        "TSF in just 4 sessions, making it highly cost-effective. This was revolutionary in demonstrating",
        "that a brief motivational intervention could match the efficacy of longer treatments.",
    ],
    bg_color='light_green', title="Key Finding:", title_color='green',
    border_color='green', font_size=13)

add_reference_bar(slide, "Reference: Project MATCH Research Group (1997). Matching treatments to patient heterogeneity. J Stud Alcohol, 58, 7-29.")


# ============================================================
# SECTION 2: THEORETICAL FOUNDATIONS
# ============================================================
add_section_divider("SECTION 2", "Theoretical Foundations of MET", 'teal')

# SLIDE: Transtheoretical Model - Overview
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_teal')
add_title_bar(slide, "The Transtheoretical Model (TTM) of Change", 'teal')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.2),
    [
        "The Transtheoretical Model (Prochaska & DiClemente, 1982, 1984) is the primary theoretical",
        "foundation of MET. It describes behavior change as a PROCESS that unfolds over time through",
        "a series of stages. People do not change abruptly - they move through predictable stages,",
        "and the therapist's task is to match their interventions to the client's current stage.",
        "",
        "\"MET is based on the assumption that the responsibility and capability for change lie",
        "within the client. The therapist's task is to create a set of conditions that will enhance",
        "the client's own motivation for and commitment to change.\" (MET Manual, p. 2)",
    ],
    bg_color='cream', title="Foundation of MET", title_color='teal',
    border_color='teal', font_size=13)

# Process flow of stages
add_process_flow(slide, Inches(0.5), Inches(3.8), Inches(12.3),
    ["Pre-\ncontemplation", "Contemplation", "Preparation", "Action", "Maintenance", "Relapse\n(Recycling)"],
    ['red', 'orange', 'gold', 'green', 'teal', 'purple'])

add_content_box(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.3),
    [
        "Key Principle: Change is cyclical, not linear. Most people cycle through stages multiple times",
        "before achieving lasting change. Relapse is a NORMAL part of the change process, not a failure.",
    ],
    bg_color='light_red', title="Important:", title_color='red',
    border_color='red', font_size=13)

add_reference_bar(slide, "Reference: Prochaska & DiClemente (1984). The Transtheoretical Approach; MET Manual (1992), pp. 6-12; NIMHANS (2016), Ch. 8.")


# SLIDE: Stage 1 - Precontemplation
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_red')
add_title_bar(slide, "Stage 1: PRECONTEMPLATION", 'red')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "The person does NOT see their substance",
        "use as a problem. They are not thinking",
        "about change and may be unaware of risks.",
        "",
        "Common statements:",
        "\"I don't have a problem\"",
        "\"Everyone drinks like me\"",
        "\"My family is overreacting\"",
    ],
    bg_color='light_red', title="What It Looks Like", title_color='red',
    border_color='red', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "DO: Raise doubt, increase awareness",
        "DO: Provide objective information/feedback",
        "DO: Explore events that brought them in",
        "DO: Help them see risks they haven't noticed",
        "",
        "DON'T: Argue or confront",
        "DON'T: Push for immediate change",
        "DON'T: Label them as 'alcoholic' or 'addict'",
    ],
    bg_color='cream', title="Therapist Tasks", title_color='orange',
    border_color='orange', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.3),
    [
        "MET Technique: Present Personal Feedback Report (PFR) showing objective assessment data.",
        "Example: \"Your blood alcohol level on the assessment day was ___ . Here's where that falls",
        "compared to the general population...\" The therapist presents facts WITHOUT arguing,",
        "letting the DATA create discrepancy between client's belief and reality.",
        "",
        "According to the manual: \"The precontemplation stage is characterized by a lack of awareness",
        "that a problem exists. The individual has no intention to change behavior in the foreseeable future.\"",
    ],
    bg_color='white', title="How MET Addresses This Stage:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 6-8; Prochaska & DiClemente (1982). Transtheoretical therapy. Psychotherapy, 19, 276-288.")


# SLIDE: Stage 2 - Contemplation
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_orange')
add_title_bar(slide, "Stage 2: CONTEMPLATION", 'orange')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "The person ACKNOWLEDGES a problem and",
        "is thinking about change, but is AMBIVALENT.",
        "They see both pros and cons of their behavior.",
        "",
        "Common statements:",
        "\"I know I drink too much, but...\"",
        "\"I want to quit, but I don't know how\"",
        "\"Sometimes I think I should cut down\"",
    ],
    bg_color='light_orange', title="What It Looks Like", title_color='orange',
    border_color='orange', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "DO: Explore ambivalence (don't resolve it)",
        "DO: Use Decisional Balance worksheet",
        "DO: Elicit self-motivational statements",
        "DO: Tip the balance toward change",
        "DO: Highlight discrepancy between values",
        "     and current behavior",
        "",
        "DON'T: Rush to action planning",
        "DON'T: Tell them what to do",
    ],
    bg_color='light_green', title="Therapist Tasks", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.3),
    [
        "MET Key Strategy: EXPLORING AMBIVALENCE",
        "The therapist helps the client articulate BOTH sides - reasons to change AND reasons to stay",
        "the same. This is done through reflective listening and the Decisional Balance exercise.",
        "",
        "\"Ambivalence is the contemplator's defining feature. The person is simultaneously considering",
        "and rejecting the prospect of change... The therapist's task is to tip the balance.\" (Manual, p. 8)",
    ],
    bg_color='cream', title="How MET Addresses This Stage:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 8-10; Miller & Rollnick (1991). Motivational Interviewing, Ch. 4.")

# SLIDE: Stage 3 - Preparation/Determination
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'cream')
add_title_bar(slide, "Stage 3: PREPARATION (Determination)", 'gold')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "The person has DECIDED to change and is",
        "planning how to do it. They have tipped the",
        "balance and are ready to take action.",
        "",
        "Common statements:",
        "\"I need to do something about this\"",
        "\"What are my options?\"",
        "\"I'm going to quit next Monday\"",
    ],
    bg_color='cream', title="What It Looks Like", title_color='gold',
    border_color='gold', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "DO: Help develop a Change Plan",
        "DO: Offer a menu of options",
        "DO: Support self-efficacy",
        "DO: Help set realistic, achievable goals",
        "DO: Negotiate a plan (not impose one)",
        "",
        "DON'T: Prescribe a single path",
        "DON'T: Ignore client's preferences",
        "DON'T: Miss the window of readiness",
    ],
    bg_color='light_green', title="Therapist Tasks", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.3),
    [
        "MET Key Strategy: CHANGE PLAN WORKSHEET",
        "The therapist helps the client complete a structured Change Plan that includes: specific changes",
        "desired, main reasons for change, steps planned, how others can help, how they'll know it's working,",
        "and things that could interfere. This plan is the client's OWN plan, negotiated collaboratively.",
        "",
        "\"The window of determination is open for a period of time. If action is not taken, the person",
        "may slip back into contemplation or precontemplation.\" (MET Manual, p. 10)",
    ],
    bg_color='light_blue', title="How MET Addresses This Stage:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 10-11; DiClemente et al. (1991). J Consult Clin Psychol, 59, 295-304.")


# SLIDE: Stage 4 - Action
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "Stage 4: ACTION", 'green')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "The person is ACTIVELY making changes.",
        "They are modifying their behavior, experiences,",
        "or environment to overcome their problem.",
        "",
        "Common statements:",
        "\"I stopped drinking 2 weeks ago\"",
        "\"I'm attending AA meetings now\"",
        "\"I threw away my bottles\"",
    ],
    bg_color='light_green', title="What It Looks Like", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "DO: Affirm and support their efforts",
        "DO: Help troubleshoot obstacles",
        "DO: Review and refine the Change Plan",
        "DO: Explore what's working and what isn't",
        "DO: Build self-efficacy (\"You can do this\")",
        "",
        "DON'T: Assume the work is done",
        "DON'T: Become complacent",
        "DON'T: Ignore emerging difficulties",
    ],
    bg_color='light_purple', title="Therapist Tasks", title_color='purple',
    border_color='purple', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.3),
    [
        "MET Key Strategy: REVIEWING PROGRESS & RENEWING COMMITMENT",
        "Sessions 3 and 4 of MET focus on reviewing the Change Plan, celebrating successes,",
        "problem-solving obstacles, and renewing motivation when it flags. The therapist continues",
        "to use MI techniques to maintain momentum and prevent discouragement.",
        "",
        "\"Action without preparation is a recipe for failure. The Action stage is where the therapist",
        "supports implementation of the plan developed in previous sessions.\" (MET Manual, p. 11)",
    ],
    bg_color='cream', title="How MET Addresses This Stage:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 11-12; Prochaska, DiClemente & Norcross (1992). Am Psychologist, 47, 1102-1114.")

# SLIDE: Stages 5 & 6 - Maintenance and Relapse
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_teal')
add_title_bar(slide, "Stages 5 & 6: MAINTENANCE and RELAPSE", 'teal')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.8),
    [
        "MAINTENANCE:",
        "- Sustaining gains made during action",
        "- Preventing return to old patterns",
        "- Building new lifestyle and identity",
        "- Developing coping for high-risk situations",
        "",
        "Duration: 6 months to lifetime",
        "Challenge: Remaining vigilant without",
        "  becoming rigid or fearful",
    ],
    bg_color='light_teal', title="Maintenance Stage", title_color='teal',
    border_color='teal', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.8),
    [
        "RELAPSE (Recycling):",
        "- Return to earlier stage behavior",
        "- NOT a failure - a learning opportunity",
        "- Average person cycles 3-7 times before",
        "  achieving stable change",
        "- MET normalizes this as part of the process",
        "",
        "Key MET response to relapse: Explore what",
        "happened WITHOUT blame, renew motivation,",
        "revise the Change Plan accordingly",
    ],
    bg_color='light_red', title="Relapse Stage", title_color='red',
    border_color='red', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2),
    [
        "\"The Transtheoretical Model views relapse not as a catastrophic failure but as a normal part of",
        "the cycle of change. The revolving-door model emphasizes that with each cycle, the person learns",
        "something new and may achieve more stable maintenance in the next attempt.\" (MET Manual, p. 12)",
        "",
        "MET implication: Even in Session 4, the therapist prepares the client for possible slips and reframes",
        "them as information rather than defeat, thereby reducing shame and maintaining engagement.",
    ],
    bg_color='cream', title="From the Manual:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 11-14; NIMHANS (2016). Substance Use Disorders, pp. 89-94; Marlatt & Gordon (1985).")


# SLIDE: Self-Efficacy Theory
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_purple')
add_title_bar(slide, "Theoretical Foundation: Self-Efficacy (Bandura, 1977)", 'purple')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3),
    [
        "Self-efficacy = a person's BELIEF in their",
        "ability to successfully perform a behavior.",
        "",
        "In MET context:",
        "- Belief that one CAN change drinking/drug use",
        "- Confidence in ability to cope without substances",
        "- \"I can do this\" vs. \"It's impossible\"",
        "",
        "Sources of Self-Efficacy:",
        "1. Past performance accomplishments",
        "2. Vicarious experience (seeing others succeed)",
        "3. Verbal persuasion (therapist affirmation)",
        "4. Physiological/emotional states",
    ],
    bg_color='light_purple', title="What is Self-Efficacy?", title_color='purple',
    border_color='purple', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3),
    [
        "How MET Builds Self-Efficacy:",
        "",
        "1. AFFIRMING client's strengths and past successes",
        "   \"You managed to stay sober for 3 months last",
        "    year - that shows real strength\"",
        "",
        "2. SUPPORTING the client's own ability to change",
        "   \"No one can make this decision for you. You",
        "    have what it takes to change this.\"",
        "",
        "3. OFFERING a menu of options (not prescribing)",
        "   Gives client sense of control and agency",
        "",
        "4. CELEBRATING small wins and progress",
    ],
    bg_color='light_green', title="MET Application:", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.8),
    [
        "\"A person who lacks confidence that he or she can change is unlikely to try, regardless of how",
        "important the change may seem. The therapist's task is to enhance the client's belief in the",
        "possibility of change.\" (MET Manual, p. 17) - This is the 'S' in FRAMES: Support self-efficacy.",
    ],
    bg_color='cream', title="Manual Quote:", title_color='orange',
    border_color='gold', font_size=13)

add_reference_bar(slide, "Reference: Bandura (1977). Self-efficacy: Toward a unifying theory. Psychol Review, 84; MET Manual (1992), pp. 16-17.")

# SLIDE: Cognitive Dissonance Theory
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_orange')
add_title_bar(slide, "Theoretical Foundation: Cognitive Dissonance (Festinger, 1957)", 'orange')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3),
    [
        "Cognitive Dissonance = the uncomfortable tension",
        "when a person holds two conflicting beliefs or when",
        "their behavior contradicts their values/beliefs.",
        "",
        "In substance use context:",
        "- \"I am a good parent\" vs \"I drink every night",
        "   and my children are suffering\"",
        "- \"I value my health\" vs \"I am damaging my",
        "   liver with alcohol\"",
        "- \"I am independent\" vs \"I can't function",
        "   without my substance\"",
    ],
    bg_color='light_orange', title="The Theory:", title_color='orange',
    border_color='orange', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3),
    [
        "How MET Uses Cognitive Dissonance:",
        "",
        "1. DEVELOPING DISCREPANCY (key principle)",
        "   - Show gap between where client IS and",
        "     where they WANT TO BE",
        "   - Let the discrepancy be the motivator",
        "",
        "2. Personal Feedback Report creates dissonance:",
        "   - \"You said health is your top value...\"",
        "   - \"Your liver enzymes show damage...\"",
        "   - Client must resolve this tension",
        "",
        "3. Resolution happens through CHANGE",
        "   (not through therapist confrontation)",
    ],
    bg_color='light_blue', title="MET Application:", title_color='deep_blue',
    border_color='deep_blue', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.8),
    [
        "\"A discrepancy between present behavior and important personal goals will motivate change.",
        "The therapist's task is to develop and amplify such discrepancy, making the client more aware",
        "of the costs of current behavior relative to their own values and goals.\" (MET Manual, p. 14)",
    ],
    bg_color='cream', title="Manual Quote:", title_color='orange',
    border_color='gold', font_size=13)

add_reference_bar(slide, "Reference: Festinger (1957). Theory of Cognitive Dissonance; MET Manual (1992), pp. 13-15; NIMHANS (2016), Ch. 8.")


# ============================================================
# SECTION 3: PRINCIPLES & SPIRIT OF MET
# ============================================================
add_section_divider("SECTION 3", "Principles and Spirit of MET", 'green')

# SLIDE: Five Principles of MI in MET
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "Five Principles of Motivational Interviewing in MET", 'green')

principles = [
    ("1. Express Empathy", "Accept the client where they are. Ambivalence is NORMAL. Reflective listening is fundamental.", 'light_blue', 'deep_blue'),
    ("2. Develop Discrepancy", "Help client see gap between current behavior and important goals/values. Let THEM argue for change.", 'light_green', 'green'),
    ("3. Avoid Argumentation", "Arguments are counterproductive. Defending breeds defensiveness. Resistance signals to shift approach.", 'light_orange', 'orange'),
    ("4. Roll with Resistance", "Don't fight resistance - use it. Offer new perspectives, don't impose. Client is the primary resource.", 'light_purple', 'purple'),
    ("5. Support Self-Efficacy", "Client's belief in possibility of change is key motivator. Client is responsible for choosing change.", 'light_teal', 'teal'),
]

for i, (title, desc, bg_c, txt_c) in enumerate(principles):
    y = Inches(1.3) + Inches(1.1) * i
    add_content_box(slide, Inches(0.5), y, Inches(12.3), Inches(1.0),
        [desc], bg_color=bg_c, title=title, title_color=txt_c,
        border_color=txt_c, font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 13-17; Miller & Rollnick (1991). Motivational Interviewing, Ch. 3.")

# SLIDE: FRAMES
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "FRAMES: Elements of Effective Brief Interventions", 'deep_blue')

frames = [
    ("F - Feedback", "Provide personal feedback about the client's risk/impairment based on objective assessment data (blood tests, questionnaire scores, comparison with norms)", 'light_red', 'red'),
    ("R - Responsibility", "Emphasize that the responsibility for change lies with the CLIENT. \"Only you can decide what to do with this information. It's your choice.\"", 'light_orange', 'orange'),
    ("A - Advice", "Clear advice to change is given but in a non-prescriptive manner. \"My professional advice would be to consider reducing/stopping...\"", 'cream', 'gold'),
    ("M - Menu", "Offer a MENU of strategies/options. Client chooses their own path. Multiple pathways to change exist. Choice increases commitment.", 'light_green', 'green'),
    ("E - Empathy", "Warm, reflective, empathic counseling style. Understanding without judging. Acceptance of the person.", 'light_blue', 'deep_blue'),
    ("S - Self-Efficacy", "Reinforce client's hope and belief that change IS possible. \"You have the strength to do this. Many people in similar situations have succeeded.\"", 'light_purple', 'purple'),
]

for i, (title, desc, bg_c, txt_c) in enumerate(frames):
    y = Inches(1.3) + Inches(0.93) * i
    add_content_box(slide, Inches(0.5), y, Inches(12.3), Inches(0.85),
        [desc], bg_color=bg_c, title=title, title_color=txt_c,
        border_color=txt_c, font_size=11)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 16-17; Miller & Sanchez (1994). Motivating young adults for change.")


# SLIDE: OARS Skills
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_teal')
add_title_bar(slide, "OARS: Core Microskills of MET", 'teal')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.4),
    [
        "OPEN-ENDED QUESTIONS",
        "- Cannot be answered with yes/no",
        "- Encourage client to explore and elaborate",
        "- \"What concerns you about your drinking?\"",
        "- \"How has substance use affected your life?\"",
        "- \"What would be different if you changed?\"",
    ],
    bg_color='light_blue', title="O - Open-Ended Questions", title_color='deep_blue',
    border_color='deep_blue', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.4),
    [
        "AFFIRMATIONS",
        "- Recognize client's strengths and efforts",
        "- Build confidence and self-efficacy",
        "- \"It took courage to come here today\"",
        "- \"You've shown real resilience in coping\"",
        "- \"That's a significant step you've taken\"",
    ],
    bg_color='light_green', title="A - Affirmations", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.0), Inches(6), Inches(2.4),
    [
        "REFLECTIVE LISTENING",
        "- The PRIMARY skill in MI/MET",
        "- Repeat back meaning (not just words)",
        "- Simple: \"So you feel frustrated\"",
        "- Complex: \"You want to quit but you're",
        "  worried you'll lose your social circle\"",
        "- Shows empathy, builds rapport",
    ],
    bg_color='light_purple', title="R - Reflective Listening", title_color='purple',
    border_color='purple', font_size=13)

add_content_box(slide, Inches(6.8), Inches(4.0), Inches(6), Inches(2.4),
    [
        "SUMMARIES",
        "- Collect and link what client has said",
        "- Show you've been listening carefully",
        "- Prepare transitions between topics",
        "- \"Let me see if I have this right...\"",
        "- Include both sides of ambivalence",
        "- End with open question to continue",
    ],
    bg_color='light_orange', title="S - Summaries", title_color='orange',
    border_color='orange', font_size=13)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 20-35; Miller & Rollnick (2002). MI: Preparing People for Change, 2nd ed.")

# SLIDE: DARES
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_orange')
add_title_bar(slide, "DARES: Self-Motivational Statements to Elicit", 'orange')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.5),
    [
        "The therapist's goal is to ELICIT (not provide) self-motivational statements from the client.",
        "These are statements where the CLIENT argues FOR change. The more the client voices these,",
        "the more likely actual change becomes. DARES represents the categories to listen for and evoke.",
    ],
    bg_color='cream', title="Core Principle:", title_color='orange',
    border_color='gold', font_size=13)

dares = [
    ("D - Desire", "\"I want to change\" / \"I wish I could stop\"", 'light_blue'),
    ("A - Ability", "\"I think I could do it\" / \"I was able to quit before\"", 'light_green'),
    ("R - Reasons", "\"My health is suffering\" / \"My family needs me sober\"", 'light_purple'),
    ("E - Emotional", "\"I'm tired of feeling this way\" / \"I hate what I've become\"", 'light_orange'),
    ("S - Steps", "\"I could try going to meetings\" / \"Maybe I'll call my doctor\"", 'light_teal'),
]

for i, (title, example, bg_c) in enumerate(dares):
    y = Inches(3.1) + Inches(0.72) * i
    add_content_box(slide, Inches(0.5), y, Inches(12.3), Inches(0.65),
        [f"Example: {example}"], bg_color=bg_c, title=title, title_color='dark_gray',
        font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 25-30; Miller & Rollnick (2002). MI, 2nd ed.; Amrhein et al. (2003). J Consult Clin Psychol.")


# SLIDE: Spirit of MI
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "The Spirit of Motivational Interviewing", 'green')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(3),
    [
        "Working WITH the client",
        "as equal partners.",
        "NOT expert-to-patient.",
        "",
        "\"The therapeutic",
        "relationship is more like",
        "a partnership than an",
        "expert-recipient one.\"",
        "(MET Manual, p. 13)",
    ],
    bg_color='light_blue', title="Collaboration", title_color='deep_blue',
    border_color='deep_blue', font_size=13)

add_content_box(slide, Inches(4.6), Inches(1.3), Inches(3.8), Inches(3),
    [
        "Drawing out the client's",
        "own motivations,",
        "strengths, and reasons",
        "for change.",
        "",
        "Motivation is EVOKED",
        "from within, not",
        "installed from outside.",
    ],
    bg_color='light_green', title="Evocation", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(8.7), Inches(1.3), Inches(4.1), Inches(3),
    [
        "Respecting the client's",
        "right and capacity to",
        "direct their own life.",
        "",
        "The client ultimately",
        "decides whether, when,",
        "and how to change.",
        "Therapist respects this.",
    ],
    bg_color='light_purple', title="Autonomy", title_color='purple',
    border_color='purple', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2),
    [
        "What the Spirit Looks Like in Practice:",
        "- Therapist asks more than tells; listens more than instructs",
        "- Client does most of the talking (therapist aims for 70:30 client:therapist talk ratio)",
        "- No labeling, no shaming, no arguing for change on behalf of the client",
        "- Resistance is met with curiosity, not confrontation",
        "- The client's autonomy is honored even when they choose NOT to change immediately",
    ],
    bg_color='cream', title="In Practice:", title_color='orange',
    border_color='gold', font_size=13)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 13-18; Miller & Rollnick (2013). MI, 3rd ed.; NIMHANS (2016), Ch. 8.")


# SLIDE: Handling Resistance - Detailed
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_red')
add_title_bar(slide, "Handling Client Resistance: Strategies from the Manual", 'red')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "TYPES OF RESISTANCE:",
        "",
        "1. Arguing - challenging, discounting, hostility",
        "2. Interrupting - cutting off, talking over",
        "3. Denying - blaming, disagreeing, excusing,",
        "   claiming impunity, minimizing, pessimism",
        "4. Ignoring - inattention, non-answer,",
        "   no response, sidetracking",
    ],
    bg_color='light_red', title="Recognizing Resistance", title_color='red',
    border_color='red', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "STRATEGIES TO ROLL WITH RESISTANCE:",
        "",
        "1. Simple Reflection - acknowledge what they said",
        "2. Amplified Reflection - reflect back in stronger form",
        "3. Double-Sided Reflection - both sides of ambivalence",
        "4. Shifting Focus - redirect to another topic",
        "5. Agreement with a Twist - agree then reframe",
        "6. Reframing - offer new interpretation of facts",
        "7. Emphasizing Personal Choice - \"It's up to you\"",
    ],
    bg_color='light_green', title="Rolling with Resistance", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.4),
    [
        "CLINICAL EXAMPLES from the Manual:",
        "",
        "Client: \"I don't think I drink any more than my friends do.\"",
        "Simple Reflection: \"You see yourself as a pretty normal drinker.\"",
        "Amplified: \"So you really don't think there's anything at all to be concerned about.\"",
        "Double-Sided: \"On one hand you feel your drinking is normal, and on the other hand,",
        "               you're here because something prompted you to come in.\"",
        "Reframe: \"It sounds like your friends are important to you, and you want to fit in with them.\"",
    ],
    bg_color='cream', title="Clinical Dialogue:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 32-38; Miller & Rollnick (1991). MI, Ch. 6: Handling Resistance.")


# ============================================================
# SECTION 4: SESSION-BY-SESSION GUIDE
# ============================================================
add_section_divider("SECTION 4", "Session-by-Session Therapy Guide\nHow to Actually DO MET", 'purple')

# SLIDE: Overview of 4-Session Structure
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_purple')
add_title_bar(slide, "MET: The 4-Session Structure (Overview)", 'purple')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.3),
    [
        "SESSION 1 (Week 1):",
        "- Build rapport and therapeutic alliance",
        "- Present Personal Feedback Report",
        "- Explore client's reaction to feedback",
        "- Elicit self-motivational statements",
        "- Gauge readiness to change",
    ],
    bg_color='light_blue', title="Sessions 1 & 2", title_color='deep_blue',
    border_color='deep_blue', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.3),
    [
        "SESSION 2 (Week 2):",
        "- Strengthen commitment to change",
        "- Develop Change Plan (collaborative)",
        "- Complete Change Plan Worksheet",
        "- Explore ambivalence further if needed",
        "- Set specific goals",
    ],
    bg_color='light_green', title="", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(3.9), Inches(6), Inches(2.3),
    [
        "SESSION 3 (Week 6):",
        "- Review progress since Session 2",
        "- Renew motivation and commitment",
        "- Problem-solve any obstacles",
        "- Revise Change Plan if needed",
        "- Affirm successes and efforts",
    ],
    bg_color='light_orange', title="Sessions 3 & 4", title_color='orange',
    border_color='orange', font_size=13)

add_content_box(slide, Inches(6.8), Inches(3.9), Inches(6), Inches(2.3),
    [
        "SESSION 4 (Week 12):",
        "- Review overall progress",
        "- Consolidate gains",
        "- Plan for maintenance",
        "- Discuss relapse prevention",
        "- Termination and future planning",
    ],
    bg_color='light_purple', title="", title_color='purple',
    border_color='purple', font_size=13)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 40-80; NIMHANS (2016). Substance Use Disorders, Ch. 8.")


# SLIDE: Session 1 - Part A (Opening)
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "SESSION 1: Building Motivation for Change (Part A)", 'deep_blue')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.8),
    [
        "GOALS: (1) Build rapport (2) Present personal feedback (3) Elicit self-motivational statements",
        "(4) Assess stage of readiness (5) Begin to explore ambivalence",
        "",
        "TIMING: 60-90 minutes | WEEK: 1 of treatment",
    ],
    bg_color='cream', title="Session Goals & Structure:", title_color='deep_blue',
    border_color='deep_blue', font_size=13)

add_content_box(slide, Inches(0.5), Inches(3.4), Inches(6), Inches(3.1),
    [
        "STEP 1: OPENING (10-15 min)",
        "",
        "How to do it:",
        "- Welcome client warmly and genuinely",
        "- Explain the process: \"Today we'll look",
        "  at some information from your assessment",
        "  and talk about what it means to you\"",
        "- Set the tone: non-judgmental, curious",
        "- Ask an open question to start:",
        "  \"What brought you here today?\"",
        "  \"How are things going for you?\"",
    ],
    bg_color='light_blue', title="Opening the Session:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_content_box(slide, Inches(6.8), Inches(3.4), Inches(6), Inches(3.1),
    [
        "STEP 2: PERSONAL FEEDBACK (30-40 min)",
        "",
        "Present data from the assessment battery:",
        "1. Drinking/drug use patterns & quantities",
        "2. Comparison with population norms",
        "3. Blood test results (liver function, etc.)",
        "4. Neuropsychological test results",
        "5. Consequences reported in assessment",
        "6. Risk factors identified",
        "",
        "Key: Present information, DON'T argue.",
        "Ask: \"What do you make of this?\"",
    ],
    bg_color='light_green', title="Presenting Feedback:", title_color='green',
    border_color='green', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 40-55; Chapter III: Session 1 - Building Motivation for Change.")

# SLIDE: Session 1 - Part B (Feedback Details)
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "SESSION 1: Personal Feedback Report (Part B)", 'deep_blue')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3.5),
    [
        "PERSONAL FEEDBACK REPORT contains:",
        "",
        "1. DRINKING/DRUG USE PATTERN",
        "   - Quantity/frequency data",
        "   - Peak blood alcohol levels",
        "   - Where client falls on population norms",
        "",
        "2. NEGATIVE CONSEQUENCES",
        "   - Physical, psychological, social",
        "   - Legal, occupational problems",
        "",
        "3. RISK FACTORS/INDICATORS",
        "   - Family history, tolerance, dependence signs",
        "",
        "4. TEST RESULTS",
        "   - Liver function, neuropsych scores",
    ],
    bg_color='white', title="Components of the PFR:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3.5),
    [
        "HOW TO PRESENT FEEDBACK:",
        "",
        "1. Use a structured form/booklet",
        "2. Go through each section systematically",
        "3. After each piece of data, PAUSE and ASK:",
        "   - \"What do you make of this?\"",
        "   - \"Does this surprise you?\"",
        "   - \"How does this compare with what",
        "      you expected?\"",
        "",
        "4. Use reflective listening on their response",
        "5. Do NOT argue if they minimize/dismiss",
        "6. Highlight discrepancies gently:",
        "   \"So on one hand you feel fine, but your",
        "    liver results suggest something else...\"",
    ],
    bg_color='light_green', title="Therapist Technique:", title_color='green',
    border_color='green', font_size=12)

add_content_box(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.4),
    [
        "\"The feedback session is not a lecture. It is a structured clinical interaction in which the",
        "therapist presents information and invites the client's reaction. The style is empathic,",
        "non-confrontational, and curious rather than didactic.\" (MET Manual, p. 44)",
    ],
    bg_color='cream', title="Manual Guidance:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 44-55; Appendix A: Personal Feedback Report format.")


# SLIDE: Session 1 - Part C (Eliciting Change Talk)
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "SESSION 1: Eliciting Self-Motivational Statements (Part C)", 'deep_blue')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "WHAT ARE SELF-MOTIVATIONAL STATEMENTS?",
        "",
        "These are statements from the CLIENT that",
        "argue FOR change. Categories (from Manual):",
        "",
        "1. Problem Recognition:",
        "   \"I guess this is more serious than I thought\"",
        "2. Concern:",
        "   \"I am worried about what this is doing to me\"",
        "3. Intention to Change:",
        "   \"I think I need to do something about this\"",
        "4. Optimism:",
        "   \"I think I can do this if I really try\"",
    ],
    bg_color='light_blue', title="Change Talk Categories:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "HOW TO ELICIT CHANGE TALK:",
        "",
        "1. EVOCATIVE QUESTIONS:",
        "   \"What worries you about your drinking?\"",
        "   \"What would be different if you changed?\"",
        "",
        "2. IMPORTANCE RULER (0-10):",
        "   \"On a scale of 0-10, how important is it to",
        "    you to make a change? Why not a lower number?\"",
        "",
        "3. EXPLORING PROS & CONS:",
        "   \"What do you like about using? What concerns you?\"",
        "",
        "4. LOOKING FORWARD/BACK:",
        "   \"Where do you see yourself in 5 years if nothing changes?\"",
    ],
    bg_color='light_green', title="Techniques to Elicit:", title_color='green',
    border_color='green', font_size=12)

add_content_box(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.4),
    [
        "CLINICAL DIALOGUE EXAMPLE (from MET Manual):",
        "",
        "Therapist: \"You mentioned that your wife has been concerned about your drinking.\"",
        "Client: \"Yeah, she's always on my case about it. But I don't think it's that bad.\"",
        "Therapist: \"So she sees it differently than you do. What concerns her the most?\" (Open Q)",
        "Client: \"She says I'm different when I drink. More angry. I guess she has a point.\"",
        "Therapist: \"You've noticed that too - that alcohol changes your behavior in ways you don't like.\"",
        "           (Reflection - amplifying the change talk)",
        "Client: \"Yeah... I don't want to be that kind of person. That's not who I am.\"",
        "Therapist: \"Being a good person and partner is really important to you.\" (Affirming values)",
    ],
    bg_color='cream', title="Session 1 Dialogue:", title_color='orange',
    border_color='gold', font_size=11)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 48-55; Miller & Rollnick (1991). MI, Ch. 5: Eliciting Self-Motivational Statements.")

# SLIDE: Session 1 - Closing
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "SESSION 1: Closing the Session", 'deep_blue')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.8),
    [
        "SUMMARIZING (5-10 min):",
        "",
        "Provide a comprehensive summary that includes:",
        "1. Key facts from the feedback",
        "2. Client's reactions to the feedback",
        "3. Any self-motivational statements made",
        "4. Both sides of client's ambivalence",
        "5. Any expressions of concern or intention",
        "",
        "End with: \"What else would you add?\"",
        "or \"Does that capture it accurately?\"",
    ],
    bg_color='light_blue', title="Closing Summary:", title_color='deep_blue',
    border_color='deep_blue', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.8),
    [
        "ASSESSING READINESS:",
        "",
        "Use the Readiness Ruler:",
        "\"On a scale of 1-10, how ready are you",
        " to make a change right now?\"",
        "",
        "If 1-3: Precontemplation focus",
        "  - Don't push, plant seeds",
        "If 4-6: Contemplation focus",
        "  - Explore ambivalence more",
        "If 7-10: Preparation focus",
        "  - Move toward Change Plan",
    ],
    bg_color='light_green', title="Readiness Assessment:", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.1),
    [
        "HOMEWORK/BETWEEN-SESSION TASKS:",
        "- Ask client to think about what was discussed: \"Between now and our next meeting, I'd like",
        "  you to think about what we've talked about today.\"",
        "- Give them the Personal Feedback sheet to take home and review",
        "- If ready: \"You might want to start thinking about what changes you'd like to make\"",
        "- Do NOT assign too much - respect their pace and readiness stage",
    ],
    bg_color='light_orange', title="Between Sessions:", title_color='orange',
    border_color='orange', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 55-58; Chapter III: Closing the first session.")


# SLIDE: Session 2 - Strengthening Commitment
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "SESSION 2: Strengthening Commitment to Change", 'green')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.5),
    [
        "GOALS: (1) Review and deepen Session 1 themes (2) Strengthen commitment (3) Develop a Change Plan",
        "(4) Help client identify specific steps (5) Address remaining ambivalence",
        "",
        "TIMING: 60 minutes | WEEK: 2 of treatment | KEY OUTPUT: Completed Change Plan Worksheet",
    ],
    bg_color='cream', title="Session 2 Goals:", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(3.1), Inches(6), Inches(3.4),
    [
        "STEP 1: OPENING (10 min)",
        "- \"What's been happening since last time?\"",
        "- \"Have you thought more about what we discussed?\"",
        "- Reflect on any changes already made",
        "",
        "STEP 2: RECAPITULATION (5 min)",
        "- Provide brief summary of Session 1",
        "- Check: \"Did I capture that correctly?\"",
        "",
        "STEP 3: DEEPENING (15 min)",
        "- Continue exploring ambivalence",
        "- Elicit more self-motivational statements",
        "- Use Decisional Balance if needed",
    ],
    bg_color='light_green', title="Session Flow:", title_color='green',
    border_color='green', font_size=12)

add_content_box(slide, Inches(6.8), Inches(3.1), Inches(6), Inches(3.4),
    [
        "STEP 4: CHANGE PLAN (25-30 min)",
        "",
        "Help client complete the Change Plan:",
        "\"Now that we've talked about this, would you",
        " like to put together a plan?\"",
        "",
        "Change Plan Worksheet sections:",
        "1. The changes I want to make are...",
        "2. The most important reasons are...",
        "3. The steps I plan to take are...",
        "4. How others can help me...",
        "5. I will know my plan is working if...",
        "6. Things that could interfere...",
        "",
        "STEP 5: CLOSING (5 min)",
        "- Summarize plan, affirm commitment",
    ],
    bg_color='light_blue', title="The Change Plan:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 59-68; Chapter IV: Session 2 - Strengthening Commitment.")

# SLIDE: Session 2 - Strategies for Different Readiness
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "SESSION 2: Adapting to Client's Readiness Level", 'green')

add_table_slide("Session 2: Adapting to Client's Readiness Level",
    ["Client Readiness", "What to Do", "What to Avoid"],
    [
        ["NOT READY\n(Still precontemplating)", "Continue building rapport\nPresent more feedback\nExplore pros of substance use\nPlant seeds of doubt gently", "Pushing for Change Plan\nSetting goals prematurely\nExpressing disappointment"],
        ["UNSURE\n(Contemplating)", "Use Decisional Balance\nExplore values vs behavior\nAsk evocative questions\n\"What would need to happen?\"", "Arguing for change\nTelling them what to do\nBeing directive"],
        ["READY\n(Determined to change)", "Negotiate Change Plan\nOffer menu of options\nHelp set specific goals\nBuild self-efficacy", "Being too cautious\nOver-exploring ambivalence\nMissing readiness window"],
    ],
    color_key='green',
    ref_text="Reference: MET Manual (1992), pp. 59-64; Adapting to client readiness, Chapter IV.")


# SLIDE: Sessions 3 & 4 - Review and Consolidation
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_orange')
add_title_bar(slide, "SESSIONS 3 & 4: Reviewing Progress & Renewing Motivation", 'orange')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3.2),
    [
        "SESSION 3 (Week 6):",
        "",
        "1. OPENING:",
        "   \"How have things been going?\"",
        "   \"What's happened since we last met?\"",
        "",
        "2. REVIEW CHANGE PLAN:",
        "   - What's working? What isn't?",
        "   - Any modifications needed?",
        "   - New obstacles encountered?",
        "",
        "3. RENEW MOTIVATION:",
        "   - Re-explore importance/confidence",
        "   - Affirm progress (even small steps)",
        "   - Address any slips without judgment",
    ],
    bg_color='light_orange', title="Session 3 Structure:", title_color='orange',
    border_color='orange', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3.2),
    [
        "SESSION 4 (Week 12):",
        "",
        "1. OPENING:",
        "   \"We're coming to the end of our time",
        "    together. How are things now?\"",
        "",
        "2. REVIEW OVERALL PROGRESS:",
        "   - Compare current status to baseline",
        "   - Celebrate achievements",
        "   - Acknowledge ongoing challenges",
        "",
        "3. MAINTENANCE PLANNING:",
        "   - Identify high-risk situations",
        "   - Plan coping strategies",
        "   - Discuss relapse as normal part of process",
        "   - Plan for ongoing support/follow-up",
    ],
    bg_color='light_purple', title="Session 4 Structure:", title_color='purple',
    border_color='purple', font_size=12)

add_content_box(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.7),
    [
        "KEY PRINCIPLE for Sessions 3 & 4: These sessions are NOT just \"check-ins.\" The therapist continues",
        "to use MI techniques actively - eliciting the client's own evaluation of progress, reflecting back their",
        "successes, exploring any ambivalence that has returned, and maintaining the collaborative spirit.",
        "\"The later sessions are opportunities to consolidate gains and solve problems.\" (MET Manual, p. 70)",
    ],
    bg_color='cream', title="Remember:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 69-80; Chapter V: Sessions 3 and 4 - Review and Renewal.")

# SLIDE: What to do if client has relapsed
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_red')
add_title_bar(slide, "When the Client Has Relapsed (Sessions 3/4)", 'red')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3),
    [
        "IF CLIENT REPORTS A SLIP/RELAPSE:",
        "",
        "1. NORMALIZE it:",
        "   \"Slips are a common part of change.\"",
        "   \"Many people experience this.\"",
        "",
        "2. EXPLORE without blame:",
        "   \"What happened? What was going on?\"",
        "   \"What did you learn from it?\"",
        "",
        "3. REFRAME as learning:",
        "   \"Now you know that [situation] is a",
        "    trigger for you. That's useful information.\"",
        "",
        "4. RE-ENGAGE motivation:",
        "   \"Where do you want to go from here?\"",
    ],
    bg_color='light_red', title="Therapist Response:", title_color='red',
    border_color='red', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3),
    [
        "WHAT NOT TO DO:",
        "",
        "- Do NOT express disappointment",
        "- Do NOT say \"you failed\"",
        "- Do NOT lecture about consequences",
        "- Do NOT use the relapse to prove they",
        "  \"need\" more intensive treatment",
        "- Do NOT label them as hopeless",
        "",
        "REMEMBER:",
        "\"The therapist maintains the same empathic,",
        "non-judgmental stance regardless of whether",
        "the client has been successful or has relapsed.\"",
        "(MET Manual, p. 72)",
    ],
    bg_color='cream', title="Avoid These Traps:", title_color='orange',
    border_color='orange', font_size=12)

add_content_box(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.9),
    [
        "CLINICAL DIALOGUE - Relapse Scenario:",
        "Client: \"I had a drink last weekend. I feel terrible about it. I've failed.\"",
        "Therapist: \"You're feeling discouraged because you had a drink.\" (Simple reflection)",
        "Client: \"Yeah. What's the point of trying?\"",
        "Therapist: \"It feels like all your hard work doesn't count anymore.\" (Amplified reflection)",
        "Client: \"Well... I mean, I DID do well for 4 weeks. That's more than I've managed in years.\"",
        "Therapist: \"That's right - four weeks is significant. What helped you during that time?\" (Affirm + explore)",
    ],
    bg_color='light_green', title="", title_color='green',
    border_color='green', font_size=11)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 70-74; Marlatt & Gordon (1985). Relapse Prevention.")


# ============================================================
# SECTION 5: CLINICAL TECHNIQUES IN DETAIL
# ============================================================
add_section_divider("SECTION 5", "Clinical Techniques in Detail\nMastering the Art of MET", 'teal')

# SLIDE: Reflective Listening - The Core Skill
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_teal')
add_title_bar(slide, "Reflective Listening: The Foundation Skill of MET", 'teal')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "LEVELS OF REFLECTION:",
        "",
        "Level 1 - REPEAT/REPHRASE:",
        "Client: \"I drink every night\"",
        "Therapist: \"You drink every evening\"",
        "",
        "Level 2 - PARAPHRASE (meaning):",
        "Client: \"I drink every night\"",
        "Therapist: \"It's become a daily routine for you\"",
        "",
        "Level 3 - REFLECTION OF FEELING:",
        "Client: \"I drink every night\"",
        "Therapist: \"You sound concerned about that pattern\"",
    ],
    bg_color='light_teal', title="Types of Reflection:", title_color='teal',
    border_color='teal', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "STRATEGIC USE IN MET:",
        "",
        "- Reflect CHANGE TALK more than sustain talk",
        "- Use reflections to AMPLIFY motivation:",
        "  Client: \"I guess I drink a bit much\"",
        "  Therapist: \"You've come to realize that your",
        "  drinking has become a real problem\" (amplified)",
        "",
        "- Use DOUBLE-SIDED reflections for ambivalence:",
        "  \"On one hand, drinking helps you relax,",
        "   and on the other hand, it's costing you",
        "   your marriage.\"",
    ],
    bg_color='light_green', title="Strategic Reflections:", title_color='green',
    border_color='green', font_size=12)

add_content_box(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.4),
    [
        "PRACTICE GUIDELINES from the Manual:",
        "- Aim for 2-3 reflections for every question asked",
        "- Reflections should be STATEMENTS, not questions (drop your voice at the end)",
        "- If your reflection is wrong, the client will correct you - that's fine, you learn more",
        "- \"The principal technique for expressing empathy is reflective listening... seeking through",
        "   your responses to understand the client's meaning and feelings.\" (MET Manual, p. 21)",
        "",
        "Common Mistake: Asking too many questions and not enough reflections.",
        "Good ratio: 2:1 or 3:1 (reflections:questions)",
    ],
    bg_color='cream', title="Practice Tips:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 20-24; Miller & Rollnick (1991). MI, Ch. 4: Building Motivation for Change.")


# SLIDE: Decisional Balance
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_orange')
add_title_bar(slide, "Technique: Decisional Balance (Exploring Ambivalence)", 'orange')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.3),
    [
        "GOOD THINGS about         LESS GOOD THINGS about",
        "my substance use:          my substance use:",
        "________________________   ________________________",
        "- Helps me relax           - Health problems",
        "- Social connection         - Family conflicts",
        "- Coping with stress        - Money problems",
        "- Feels good               - Work issues",
        "________________________   ________________________",
    ],
    bg_color='light_orange', title="The Decisional Balance Grid:", title_color='orange',
    border_color='orange', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.3),
    [
        "GOOD THINGS about         LESS GOOD THINGS about",
        "making a change:           making a change:",
        "________________________   ________________________",
        "- Better health            - Losing drinking friends",
        "- Family happy             - Boredom",
        "- Save money               - Coping without substance",
        "- Self-respect             - Social pressure",
        "________________________   ________________________",
    ],
    bg_color='light_green', title="The Other Side:", title_color='green',
    border_color='green', font_size=12)

add_content_box(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(2.6),
    [
        "HOW TO USE THE DECISIONAL BALANCE IN SESSION:",
        "",
        "1. Draw a 2x2 grid on paper or use the worksheet",
        "2. Start with the GOOD things about substance use (counterintuitive but important!)",
        "   - Shows you understand, builds trust, explores their motivation to continue",
        "3. Then ask about LESS GOOD things about substance use",
        "4. Explore good things and less good things about change",
        "5. Summarize BOTH sides back to the client",
        "6. Ask: \"Where does that leave you?\" or \"What would you like to do with this?\"",
        "",
        "KEY INSIGHT: By exploring the POSITIVES of substance use first, you demonstrate acceptance",
        "and avoid the trap of being the one arguing for change. Let the client discover their own reasons.",
    ],
    bg_color='cream', title="Clinical Application:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 28-30; Janis & Mann (1977). Decision Making; NIMHANS (2016), Ch. 8.")

# SLIDE: Importance and Confidence Rulers
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_purple')
add_title_bar(slide, "Technique: Importance & Confidence Rulers", 'purple')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.8),
    [
        "IMPORTANCE RULER:",
        "",
        "\"On a scale of 0 to 10, where 0 is not",
        "at all important and 10 is extremely",
        "important, how important is it to you",
        "to make this change?\"",
        "",
        "Follow-up questions:",
        "- \"Why did you say ___ and not a 0?\"",
        "  (Elicits reasons for change)",
        "- \"What would it take to move from ___ to",
        "  a higher number?\"",
        "  (Identifies what's needed)",
    ],
    bg_color='light_purple', title="Assessing Importance:", title_color='purple',
    border_color='purple', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.8),
    [
        "CONFIDENCE RULER:",
        "",
        "\"On a scale of 0 to 10, where 0 is not",
        "at all confident and 10 is extremely",
        "confident, how confident are you that",
        "you could make this change?\"",
        "",
        "Follow-up questions:",
        "- \"Why did you say ___ and not a 0?\"",
        "  (Elicits self-efficacy statements)",
        "- \"What has worked for you in the past?\"",
        "  (Builds on past successes)",
        "- \"What would help increase your confidence?\"",
    ],
    bg_color='light_teal', title="Assessing Confidence:", title_color='teal',
    border_color='teal', font_size=12)

add_content_box(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.1),
    [
        "WHY THIS WORKS: The key therapeutic move is asking \"Why not a lower number?\" instead of",
        "\"Why not a higher number?\" This invites the client to argue FOR change (articulate their own",
        "reasons and strengths) rather than defending why they haven't changed yet.",
        "",
        "WHEN TO USE: Session 1 (after feedback), Session 2 (before Change Plan), Sessions 3-4 (tracking progress)",
        "If importance is low: focus on developing discrepancy and exploring values",
        "If confidence is low: focus on building self-efficacy, past successes, small steps",
    ],
    bg_color='cream', title="Clinical Wisdom:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 35-37; Miller & Rollnick (2002). MI 2nd ed., Ch. 4; Rollnick et al. (1999).")


# SLIDE: Things to Avoid (Therapist Traps)
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_red')
add_title_bar(slide, "Therapist Traps to AVOID (from the MET Manual)", 'red')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "1. THE QUESTION-ANSWER TRAP",
        "   Too many questions = client feels interrogated",
        "   Solution: Use more reflections, fewer questions",
        "",
        "2. THE CONFRONTATION-DENIAL TRAP",
        "   Arguing = client defends drinking more",
        "   Solution: Roll with resistance, use reflections",
        "",
        "3. THE EXPERT TRAP",
        "   \"You should...\" \"The research shows...\"",
        "   Solution: Elicit client's own solutions first",
    ],
    bg_color='light_red', title="Common Traps:", title_color='red',
    border_color='red', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "4. THE LABELING TRAP",
        "   Insisting client accept label \"alcoholic\"",
        "   Solution: Focus on behavior, not labels",
        "",
        "5. THE PREMATURE FOCUS TRAP",
        "   Jumping to solutions before exploring problem",
        "   Solution: Follow client's pace",
        "",
        "6. THE BLAMING TRAP",
        "   Client feels blamed for the problem",
        "   Solution: Emphasize that blame is irrelevant;",
        "   what matters is what they want to do now",
    ],
    bg_color='light_orange', title="More Traps:", title_color='orange',
    border_color='orange', font_size=12)

add_content_box(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.4),
    [
        "\"Direct argumentation and an aggressive confrontational approach tend to increase client",
        "resistance and are associated with poorer outcomes... Labeling is unnecessary. The focus of",
        "MET is on eliciting the client's own concern, not on coercing acceptance of a label.\"",
        "(MET Manual, pp. 18-19)",
        "",
        "REMEMBER: Every time you argue FOR change, the client argues AGAINST it.",
        "This is called the 'righting reflex' - the therapist's natural desire to fix things actually backfires.",
        "In MET, it should be the CLIENT who voices the arguments for change, not the therapist.",
    ],
    bg_color='cream', title="Key Manual Quote:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 18-20; Miller & Rollnick (1991). MI, Ch. 7: Traps for the Unwary.")


# ============================================================
# SECTION 6: CASE CONCEPTUALIZATION
# ============================================================
add_section_divider("SECTION 6", "Case Conceptualization\nApplying MET to a Clinical Case", 'maroon')

# SLIDE: Case Introduction
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_red')
add_title_bar(slide, "Case Study: Applying MET to a Substance Use Case", 'maroon')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3),
    [
        "PRESENTING INFORMATION:",
        "",
        "- Age: 25 years, Male",
        "- Multiple substance use (polysubstance)",
        "- History of marital dissolution (divorced)",
        "- Referred for psychological assessment and",
        "  intervention",
        "",
        "ASSESSMENT FINDINGS:",
        "- Externalizing personality organization",
        "- Adequate psychological resources (EA = 12)",
        "- Poorly modulated affect (impulsive expression)",
        "- Damaged, negative self-concept (MOR = 5)",
    ],
    bg_color='light_red', title="Client Background:", title_color='maroon',
    border_color='maroon', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3),
    [
        "PSYCHOLOGICAL PROFILE:",
        "",
        "- Emotion-driven decision making (affect",
        "  dominates over reflection)",
        "- Hasty, inefficient information processing",
        "  (jumps to conclusions without full data)",
        "- Raw emotional discharge (bypasses cognitive",
        "  control when stressed)",
        "- Socially engaged but difficulty with deep",
        "  intimate relationships",
        "- Pervasive pessimism and damage-oriented",
        "  self-view (\"I am fundamentally flawed\")",
        "- Elevated risk indicators requiring monitoring",
    ],
    bg_color='light_orange', title="Key Findings:", title_color='orange',
    border_color='orange', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.9),
    [
        "MET RELEVANCE: This client has ADEQUATE resources but uses them poorly due to an emotion-first,",
        "impulsive style. MET is ideal because: (1) He doesn't lack capacity - he lacks motivation and",
        "reflective delay (2) Confrontation would trigger his externalizing defenses (3) His damaged",
        "self-image needs affirmation, not more criticism (4) His autonomy needs to be respected.",
    ],
    bg_color='cream', title="Why MET is Appropriate:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: Case data from Rorschach Inkblot Assessment (Exner Comprehensive System); MET Manual (1992).")


# SLIDE: Case - Session 1 Plan
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "Case: SESSION 1 - Building Motivation", 'deep_blue')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3),
    [
        "SPECIFIC APPROACH FOR THIS CLIENT:",
        "",
        "1. Build rapport carefully - he is socially",
        "   engaged but emotionally guarded",
        "",
        "2. Present feedback GENTLY - his damaged self-",
        "   image means he may become defensive or",
        "   overwhelmed if data feels like an attack",
        "",
        "3. Focus on discrepancy between his VALUES",
        "   (being a good person, having relationships)",
        "   and his CURRENT situation (divorced,",
        "   substance-dependent)",
        "",
        "4. Use LOTS of affirmation - he needs to hear",
        "   that he has strengths (his resources ARE there)",
    ],
    bg_color='light_blue', title="Session 1 Strategy:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3),
    [
        "SAMPLE DIALOGUE:",
        "",
        "T: \"You mentioned that your marriage ended.",
        "   What role, if any, did substance use play?\"",
        "C: \"She said I was different when I used. Angry.\"",
        "T: \"Being the kind of partner you want to be is",
        "   important to you, and the substances were",
        "   getting in the way of that.\" (Reflection)",
        "C: \"Yeah... I didn't mean to be like that.\"",
        "T: \"There's a gap between who you WANT to",
        "   be and what happens when you use.\"",
        "   (Developing discrepancy)",
        "C: \"I guess so. I just... I don't know how to",
        "   handle things without it.\"",
        "T: \"You're wondering if you can cope differently.\"",
    ],
    bg_color='light_green', title="Clinical Dialogue:", title_color='green',
    border_color='green', font_size=11)

add_content_box(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.9),
    [
        "CAUTIONS FOR THIS CLIENT:",
        "- Avoid confrontation (will trigger impulsive defensive reactions)",
        "- Monitor affect closely (his emotional expression can escalate quickly due to poor modulation)",
        "- Do NOT interpret his externalizing style as \"denial\" - it's his personality structure",
        "- Build self-efficacy continuously (counters damaged self-image)",
        "- Active risk assessment should be ongoing given elevated S-CON indicators",
    ],
    bg_color='light_red', title="Clinical Cautions:", title_color='red',
    border_color='red', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), Ch. III; Exner (2003). The Rorschach: A Comprehensive System, Vol. 1.")

# SLIDE: Case - Session 2 Plan
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "Case: SESSION 2 - Developing the Change Plan", 'green')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3.2),
    [
        "CHANGE PLAN FOR THIS CLIENT:",
        "",
        "Changes I want to make:",
        "- Reduce/stop polysubstance use",
        "- Learn to handle emotions without substances",
        "- Rebuild my self-image and relationships",
        "",
        "Most important reasons:",
        "- I don't want to lose more relationships",
        "- I want to feel good about myself again",
        "- My health is being affected",
        "",
        "Steps I plan to take:",
        "- Identify my triggers (emotional situations)",
        "- Learn to pause before reacting",
        "- Find one healthy coping activity",
    ],
    bg_color='light_green', title="Sample Change Plan:", title_color='green',
    border_color='green', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3.2),
    [
        "THERAPEUTIC FOCUS - SESSION 2:",
        "",
        "Given his personality profile:",
        "",
        "1. AFFECT REGULATION focus:",
        "   - Substances serve as his affect regulator",
        "   - Plan must include alternative strategies",
        "   - \"What else helps when feelings get intense?\"",
        "",
        "2. REFLECTIVE DELAY:",
        "   - His hasty processing means he acts before",
        "     thinking - build in \"pause\" strategies",
        "   - \"What if you took 10 minutes before deciding",
        "     to use? What might happen in those 10 min?\"",
        "",
        "3. SMALL, ACHIEVABLE STEPS:",
        "   - His damaged self-image means failures feel",
        "     catastrophic - ensure early success",
    ],
    bg_color='light_purple', title="Tailored Approach:", title_color='purple',
    border_color='purple', font_size=12)

add_content_box(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.7),
    [
        "KEY PRINCIPLE: For this client, the Change Plan must address the FUNCTION of substance use",
        "(emotional regulation, escape from negative self-view, social lubrication) not just the behavior.",
        "MET helps by exploring: \"What does using DO for you? What need does it meet? How else might",
        "you meet that need?\" This approach respects the client's intelligence and autonomy.",
    ],
    bg_color='cream', title="Integration:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), Ch. IV; NIMHANS (2016). Psychosocial Interventions for Substance Use Disorders.")


# SLIDE: Case - Sessions 3 & 4
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_orange')
add_title_bar(slide, "Case: SESSIONS 3 & 4 - Progress Review & Maintenance", 'orange')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.8),
    [
        "SESSION 3 - POTENTIAL SCENARIOS:",
        "",
        "Scenario A: Client has made progress",
        "- Affirm heavily (counters damaged self-image)",
        "- \"You proved to yourself you can do this\"",
        "- Link success to his own strengths",
        "",
        "Scenario B: Client has slipped",
        "- Normalize without minimizing",
        "- Explore triggers (likely emotional situations)",
        "- \"Your emotions overwhelmed you - that's",
        "   information about what to work on\"",
        "- Revise plan: add more emotion regulation",
    ],
    bg_color='light_orange', title="Session 3 Approach:", title_color='orange',
    border_color='orange', font_size=12)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.8),
    [
        "SESSION 4 - MAINTENANCE PLANNING:",
        "",
        "For this client, high-risk situations include:",
        "- Intense emotional experiences (anger, shame)",
        "- Social situations with substance-using peers",
        "- Relationship conflicts or rejection",
        "- Moments of deep self-criticism",
        "",
        "Maintenance plan should include:",
        "- Ongoing therapy for affect regulation",
        "- Support network identification",
        "- Crisis plan for intense emotional episodes",
        "- Regular self-assessment check-ins",
        "- Referral for continued psychotherapy",
    ],
    bg_color='light_purple', title="Session 4 Approach:", title_color='purple',
    border_color='purple', font_size=12)

add_content_box(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.1),
    [
        "LONG-TERM THERAPEUTIC NEEDS (beyond MET's 4 sessions):",
        "- Affect regulation training (his primary vulnerability is unmodulated emotional discharge)",
        "- Self-image repair work (therapeutic work on the deeply damaged, pessimistic self-view)",
        "- Interpersonal skills development (capacity for genuine intimacy is limited)",
        "- Slower information processing practice (reducing impulsive, hasty decision-making)",
        "",
        "MET provides the MOTIVATIONAL FOUNDATION. Further therapy builds on the motivation MET creates.",
    ],
    bg_color='cream', title="Beyond MET:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), Ch. V; NIMHANS (2016); Marlatt & Donovan (2005). Relapse Prevention, 2nd ed.")


# ============================================================
# SECTION 7: WORKSHEETS & TOOLS
# ============================================================
add_section_divider("SECTION 7", "Clinical Worksheets & Tools\nPrintable Resources for MET Sessions", 'deep_blue')

# SLIDE: Change Plan Worksheet
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "WORKSHEET 1: Change Plan (from MET Manual)", 'deep_blue')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.2),
    [
        "1. THE CHANGES I WANT TO MAKE ARE:",
        "   ___________________________________________________________________________",
        "   ___________________________________________________________________________",
        "",
        "2. THE MOST IMPORTANT REASONS WHY I WANT TO MAKE THESE CHANGES ARE:",
        "   ___________________________________________________________________________",
        "   ___________________________________________________________________________",
        "",
        "3. THE STEPS I PLAN TO TAKE IN CHANGING ARE:",
        "   ___________________________________________________________________________",
        "   ___________________________________________________________________________",
        "",
        "4. THE WAYS OTHER PEOPLE CAN HELP ME ARE:",
        "   Person: _________________ How: _____________________________________________",
        "   Person: _________________ How: _____________________________________________",
        "",
        "5. I WILL KNOW THAT MY PLAN IS WORKING IF:",
        "   ___________________________________________________________________________",
        "",
        "6. SOME THINGS THAT COULD INTERFERE WITH MY PLAN ARE:",
        "   ___________________________________________________________________________",
        "   ___________________________________________________________________________",
    ],
    bg_color='white', text_color='dark_gray',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), Appendix B: Change Plan Worksheet. Adapted for clinical use.")

# SLIDE: Decisional Balance Worksheet
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "WORKSHEET 2: Decisional Balance", 'green')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "BENEFITS OF MY CURRENT BEHAVIOR",
        "(What I LIKE about using substances):",
        "",
        "1. ______________________________________",
        "2. ______________________________________",
        "3. ______________________________________",
        "4. ______________________________________",
        "5. ______________________________________",
    ],
    bg_color='light_orange', title="Side A: Status Quo", title_color='orange',
    border_color='orange', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "COSTS OF MY CURRENT BEHAVIOR",
        "(What CONCERNS me about using):",
        "",
        "1. ______________________________________",
        "2. ______________________________________",
        "3. ______________________________________",
        "4. ______________________________________",
        "5. ______________________________________",
    ],
    bg_color='light_red', title="Side B: Costs", title_color='red',
    border_color='red', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.1), Inches(6), Inches(2.4),
    [
        "BENEFITS OF MAKING A CHANGE",
        "(What would be GOOD about changing):",
        "",
        "1. ______________________________________",
        "2. ______________________________________",
        "3. ______________________________________",
        "4. ______________________________________",
        "5. ______________________________________",
    ],
    bg_color='light_green', title="Side C: Benefits of Change", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(6.8), Inches(4.1), Inches(6), Inches(2.4),
    [
        "COSTS OF MAKING A CHANGE",
        "(What would be HARD about changing):",
        "",
        "1. ______________________________________",
        "2. ______________________________________",
        "3. ______________________________________",
        "4. ______________________________________",
        "5. ______________________________________",
    ],
    bg_color='light_purple', title="Side D: Costs of Change", title_color='purple',
    border_color='purple', font_size=13)

add_reference_bar(slide, "Reference: Janis & Mann (1977). Decision Making; MET Manual (1992), p. 29; Miller & Rollnick (2002), Ch. 7.")


# SLIDE: Readiness Ruler Worksheet
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_purple')
add_title_bar(slide, "WORKSHEET 3: Readiness, Importance & Confidence Rulers", 'purple')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.8),
    [
        "IMPORTANCE: How important is it to you to make this change?",
        "",
        "Not at all important  0 --- 1 --- 2 --- 3 --- 4 --- 5 --- 6 --- 7 --- 8 --- 9 --- 10  Extremely important",
        "",
        "Why did you choose this number and not a lower one? ___________________________________________",
    ],
    bg_color='light_blue', title="", title_color='deep_blue',
    border_color='deep_blue', font_size=13)

add_content_box(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.8),
    [
        "CONFIDENCE: How confident are you that you COULD make this change?",
        "",
        "Not at all confident  0 --- 1 --- 2 --- 3 --- 4 --- 5 --- 6 --- 7 --- 8 --- 9 --- 10  Extremely confident",
        "",
        "Why did you choose this number and not a lower one? ___________________________________________",
    ],
    bg_color='light_green', title="", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.3),
    [
        "READINESS: How ready are you to make this change RIGHT NOW?",
        "",
        "Not at all ready  0 --- 1 --- 2 --- 3 --- 4 --- 5 --- 6 --- 7 --- 8 --- 9 --- 10  Completely ready",
    ],
    bg_color='light_orange', title="", title_color='orange',
    border_color='orange', font_size=13)

add_reference_bar(slide, "Reference: MET Manual (1992), pp. 35-37; Rollnick, Mason & Butler (1999). Health Behavior Change.")

# SLIDE: Daily Self-Monitoring Diary
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_teal')
add_title_bar(slide, "WORKSHEET 4: Daily Self-Monitoring Diary", 'teal')

add_table_slide("WORKSHEET 4: Daily Self-Monitoring Diary",
    ["Day", "Situation/Trigger", "Feelings (0-10)", "Urge to Use (0-10)", "What I Did Instead", "Outcome"],
    [
        ["Monday", "________________", "____", "____", "________________", "________"],
        ["Tuesday", "________________", "____", "____", "________________", "________"],
        ["Wednesday", "________________", "____", "____", "________________", "________"],
        ["Thursday", "________________", "____", "____", "________________", "________"],
        ["Friday", "________________", "____", "____", "________________", "________"],
        ["Saturday", "________________", "____", "____", "________________", "________"],
        ["Sunday", "________________", "____", "____", "________________", "________"],
    ],
    color_key='teal',
    ref_text="Reference: Adapted from MET Manual (1992); Marlatt & Gordon (1985). Self-monitoring in Relapse Prevention.")

# SLIDE: Relapse Prevention Plan
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_orange')
add_title_bar(slide, "WORKSHEET 5: My Relapse Prevention Plan", 'orange')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5),
    [
        "MY HIGH-RISK SITUATIONS:",
        "",
        "1. ______________________________________",
        "2. ______________________________________",
        "3. ______________________________________",
        "",
        "MY EARLY WARNING SIGNS:",
        "",
        "1. ______________________________________",
        "2. ______________________________________",
        "3. ______________________________________",
    ],
    bg_color='light_red', title="Identifying Risks:", title_color='red',
    border_color='red', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
    [
        "MY COPING STRATEGIES:",
        "",
        "1. ______________________________________",
        "2. ______________________________________",
        "3. ______________________________________",
        "",
        "PEOPLE I CAN CALL FOR SUPPORT:",
        "",
        "Name: ______________ Phone: ______________",
        "Name: ______________ Phone: ______________",
        "Name: ______________ Phone: ______________",
    ],
    bg_color='light_green', title="My Resources:", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.4),
    [
        "IF I HAVE A SLIP, I WILL:",
        "1. Remember that a slip is NOT a failure - it is information about a trigger I need to address",
        "2. Call: _________________________________ (my support person)",
        "3. Do this instead: ___________________________________________________________________",
        "4. Review what happened: What was the trigger? What was I feeling? What can I learn?",
        "",
        "MY REASONS TO STAY ON TRACK (from my Decisional Balance):",
        "___________________________________________________________________________________",
        "___________________________________________________________________________________",
    ],
    bg_color='cream', title="My Emergency Plan:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: MET Manual (1992), Session 4; Marlatt & Donovan (2005). Relapse Prevention, 2nd ed.; NIMHANS (2016).")


# SLIDE: Values Card Sort Worksheet
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "WORKSHEET 6: Personal Values Exploration", 'green')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.5),
    [
        "INSTRUCTIONS: Rate each value on how IMPORTANT it is to you (1=Not important, 5=Extremely important)",
        "Then rate how much your CURRENT BEHAVIOR aligns with this value (1=Not at all, 5=Completely aligned)",
        "",
        "VALUE                    IMPORTANCE (1-5)    ALIGNMENT (1-5)    GAP?",
        "Being a good parent      ___________         ___________        ___",
        "Physical health          ___________         ___________        ___",
        "Financial security       ___________         ___________        ___",
        "Being honest             ___________         ___________        ___",
        "Having close relationships ___________       ___________        ___",
        "Self-respect             ___________         ___________        ___",
        "Career success           ___________         ___________        ___",
        "Being independent        ___________         ___________        ___",
    ],
    bg_color='white', title="Values-Behavior Alignment:", title_color='green',
    border_color='green', font_size=12)

add_content_box(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.4),
    [
        "REFLECTION QUESTIONS (for use in session):",
        "",
        "1. Which values have the BIGGEST gaps between importance and alignment?",
        "   ___________________________________________________________________________",
        "2. How does your substance use affect these values?",
        "   ___________________________________________________________________________",
        "3. What would change look like for the values that matter most to you?",
        "   ___________________________________________________________________________",
        "",
        "THERAPIST NOTE: This exercise creates cognitive dissonance by making the gap between values and",
        "behavior explicit. Use this to develop discrepancy in Session 1 feedback or Session 2 planning.",
    ],
    bg_color='cream', title="Processing:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: Miller et al. (2001). Personal Values Card Sort; MET Manual (1992); Miller & Rollnick (2002), Ch. 5.")


# ============================================================
# SECTION 8: RESEARCH & EFFECTIVENESS
# ============================================================
add_section_divider("SECTION 8", "Research Evidence & Effectiveness\nHow Well Does MET Work?", 'navy')

# SLIDE: Project MATCH Results
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "Project MATCH: The Landmark Trial", 'navy')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.8),
    [
        "STUDY DESIGN:",
        "- Largest alcohol treatment trial ever",
        "- 1,726 participants across 9 US sites",
        "- Randomized to 3 conditions:",
        "  1. MET (4 sessions over 12 weeks)",
        "  2. CBT (12 sessions)",
        "  3. TSF/12-Step (12 sessions)",
        "- Follow-up: 1 year and 3 years",
        "",
        "Funded by NIAAA",
        "Published 1997-1998",
    ],
    bg_color='light_blue', title="The Study:", title_color='navy',
    border_color='navy', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.8),
    [
        "KEY FINDINGS:",
        "",
        "- All three treatments produced significant",
        "  and sustained improvements",
        "- MET achieved COMPARABLE outcomes to CBT",
        "  and TSF in just 4 sessions (vs 12)",
        "- Percent Days Abstinent improved in ALL groups",
        "- Drinks Per Drinking Day decreased in ALL groups",
        "- At 3-year follow-up, gains were maintained",
        "",
        "IMPLICATION: MET is as effective as longer",
        "treatments, making it highly COST-EFFECTIVE",
    ],
    bg_color='light_green', title="Results:", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.1),
    [
        "MATCHING HYPOTHESES: The study found limited support for patient-treatment matching. However,",
        "clients HIGH in anger did significantly better in MET than in other treatments (a non-confrontational",
        "approach works better for angry clients). Clients low in readiness to change also responded well to MET.",
        "",
        "\"The finding that a 4-session motivational intervention could produce outcomes comparable to",
        "12-session cognitive-behavioral or twelve-step approaches had profound implications for the",
        "cost-effectiveness of substance abuse treatment.\" (Project MATCH Research Group, 1997)",
    ],
    bg_color='cream', title="Clinical Significance:", title_color='orange',
    border_color='gold', font_size=12)

add_reference_bar(slide, "Reference: Project MATCH Research Group (1997). J Studies on Alcohol, 58, 7-29; (1998). Addiction, 93, 1434-1447.")


# SLIDE: Meta-Analyses
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_green')
add_title_bar(slide, "Meta-Analyses: MI/MET Effectiveness Across Studies", 'green')

add_table_slide("Meta-Analyses: MI/MET Effectiveness Across Studies",
    ["Study", "N Studies", "Key Finding", "Effect Size"],
    [
        ["Burke et al. (2003)\nJ Consult Clin Psychol", "30 RCTs", "MI/MET effective for alcohol, drugs,\ndiet, and treatment adherence", "d = 0.25-0.57\n(small-medium)"],
        ["Hettema et al. (2005)\nAnnual Rev Clin Psychol", "72 studies", "MI effective across substances;\nbetter than no treatment and advice", "d = 0.77 at follow-up\nvs no treatment"],
        ["Lundahl et al. (2010)\nClient Education & Counsel", "119 studies", "MI produces moderate effect;\nstronger for substance use", "d = 0.22 (combined)\nd = 0.79 (substance)"],
        ["Vasilaki et al. (2006)\nAddiction", "15 RCTs", "Brief MI effective for reducing\nalcohol consumption", "Significant reduction\nin drinking quantity"],
        ["Smedslund et al. (2011)\nCochrane Review", "59 RCTs", "MI reduces substance use more\nthan no treatment", "SMD = -0.79\n(post-intervention)"],
    ],
    color_key='green',
    ref_text="References: As cited in table. All are peer-reviewed meta-analyses of MI/MET interventions.")

# SLIDE: Indian Research
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_orange')
add_title_bar(slide, "Research Evidence: Indian Context", 'orange')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.8),
    [
        "NIMHANS STUDIES:",
        "",
        "- Brief interventions using MI principles",
        "  found effective in Indian primary care",
        "  settings (Murthy et al., 2009)",
        "",
        "- NIMHANS addiction treatment protocols",
        "  integrate MET principles as first-line",
        "  psychosocial intervention",
        "",
        "- Community-based MI interventions reduced",
        "  alcohol use in rural Karnataka",
        "  (Nadkarni et al., 2017 - Lancet)",
    ],
    bg_color='light_orange', title="Indian Evidence:", title_color='orange',
    border_color='orange', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.8),
    [
        "KEY INDIAN STUDIES:",
        "",
        "1. Nadkarni et al. (2017) - PREMIUM Trial",
        "   - Lay counselors delivered MI in Goa",
        "   - Significant reduction in harmful drinking",
        "   - Published in The Lancet",
        "",
        "2. Pal et al. (2007) - AIIMS, Delhi",
        "   - Brief MI for alcohol dependence",
        "   - Positive outcomes in Indian setting",
        "",
        "3. Chand et al. (2018) - NIMHANS",
        "   - Technology-assisted brief intervention",
        "   - Showed feasibility in Indian clinics",
    ],
    bg_color='light_green', title="Specific Studies:", title_color='green',
    border_color='green', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.1),
    [
        "CULTURAL CONSIDERATIONS FOR INDIA:",
        "- Family involvement is crucial in Indian context (can be integrated into MET)",
        "- Stigma around substance use is high - MI's non-judgmental approach is especially valuable",
        "- Task-shifting to lay counselors is feasible (PREMIUM trial demonstrated this)",
        "- Adaptations needed for collectivist culture: family values, community roles, spiritual beliefs",
        "- NIMHANS recommends MI/MET as evidence-based first-line psychosocial intervention (2016 manual)",
    ],
    bg_color='cream', title="Indian Adaptations:", title_color='deep_blue',
    border_color='deep_blue', font_size=12)

add_reference_bar(slide, "Reference: NIMHANS (2016); Nadkarni et al. (2017). Lancet, 389, 186-195; Chand et al. (2018). Indian J Psychiatry.")


# SLIDE: UKATT and International Research
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_purple')
add_title_bar(slide, "International Research: UKATT and Global Evidence", 'purple')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.8),
    [
        "UKATT TRIAL (UK Alcohol Treatment Trial):",
        "",
        "- Largest alcohol trial in UK (742 clients)",
        "- Compared MET (3 sessions) vs SBNT (8 sessions)",
        "- Results: BOTH equally effective",
        "- MET was significantly MORE cost-effective",
        "- Savings: 5x less therapist time for same results",
        "",
        "CONCLUSION: MET achieves equivalent outcomes",
        "at a fraction of the cost of longer treatments.",
        "(UKATT Research Team, 2005, BMJ)",
    ],
    bg_color='light_purple', title="UKATT Trial:", title_color='purple',
    border_color='purple', font_size=13)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.8),
    [
        "GLOBAL RESEARCH SUMMARY:",
        "",
        "- WHO Brief Intervention Study (2002):",
        "  MI effective across 10 countries",
        "- COMBINE Study (2006, JAMA):",
        "  MI + naltrexone superior combination",
        "- Cannabis Youth Treatment Study (2004):",
        "  MET effective for adolescent cannabis use",
        "- TOPPS Study (Stephens et al., 2004):",
        "  2 sessions of MI = 6 sessions of CBT",
        "",
        "MI/MET works across cultures, substances,",
        "age groups, and clinical settings.",
    ],
    bg_color='light_teal', title="Other Key Trials:", title_color='teal',
    border_color='teal', font_size=13)

add_content_box(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.1),
    [
        "SUMMARY OF RESEARCH EVIDENCE:",
        "- MET is EMPIRICALLY SUPPORTED for alcohol use disorders (Level 1 evidence)",
        "- Effective for cannabis, cocaine, opioids, and polysubstance use",
        "- Works with adolescents, adults, and older adults",
        "- Effective when delivered by trained lay counselors (not just specialists)",
        "- Cost-effective: achieves same results as longer treatments in fewer sessions",
        "- Compatible with pharmacotherapy (enhances medication adherence)",
    ],
    bg_color='light_green', title="Bottom Line:", title_color='green',
    border_color='green', font_size=12)

add_reference_bar(slide, "Reference: UKATT Research Team (2005). BMJ, 331, 544; COMBINE (2006). JAMA, 295, 2003; WHO (2002).")

# SLIDE: Summary of Evidence Table
add_table_slide("Summary: When to Use MET/MI (Evidence-Based Indications)",
    ["Population", "Evidence Level", "Key Study", "Outcome"],
    [
        ["Alcohol Use Disorder", "Level 1 (Strong)", "Project MATCH (1997)\nUKATT (2005)", "Equal to 12-session CBT/TSF\nin just 4 sessions"],
        ["Cannabis Use", "Level 1 (Strong)", "DRINC/CYT (2004)\nStephens et al.", "Significant reduction in use\nand related problems"],
        ["Polysubstance Use", "Level 2 (Moderate)", "Multiple RCTs\nNIMHANS protocols", "Effective when tailored to\nprimary substance"],
        ["Dual Diagnosis", "Level 2 (Moderate)", "Barrowclough et al. (2001)\nBaker et al. (2005)", "Feasible and effective for\ncomorbid presentations"],
        ["Treatment Engagement", "Level 1 (Strong)", "Carroll et al. (2006)\nMartino et al. (2007)", "Increases retention and\nadherence to treatment"],
        ["Adolescents", "Level 1 (Strong)", "Cannabis Youth Treatment\nJensen et al. (2011)", "Effective and developmentally\nappropriate"],
    ],
    color_key='navy',
    ref_text="References: As cited; NICE Guidelines (2011); SAMHSA TIP 35 (1999); APA Practice Guidelines (2006).")


# ============================================================
# CLOSING SLIDES
# ============================================================

# SLIDE: Key Takeaways
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'cream')
add_title_bar(slide, "Key Takeaways: MET in Clinical Practice", 'deep_blue')

takeaways = [
    ("MET is BRIEF but POWERFUL", "4 sessions achieve outcomes comparable to 12-session treatments", 'light_blue', 'deep_blue'),
    ("It's based on SOLID THEORY", "Transtheoretical Model, Self-Efficacy, Cognitive Dissonance", 'light_green', 'green'),
    ("The CLIENT does the work", "Therapist evokes motivation, doesn't install it", 'light_purple', 'purple'),
    ("RESISTANCE is information", "Not something to fight - a signal to change approach", 'light_orange', 'orange'),
    ("FEEDBACK creates discrepancy", "Personal data + empathy = internal motivation to change", 'light_teal', 'teal'),
    ("It WORKS across settings", "Strong evidence from Project MATCH, UKATT, Indian studies", 'light_red', 'maroon'),
]

for i, (title, desc, bg_c, txt_c) in enumerate(takeaways):
    y = Inches(1.3) + Inches(0.93) * i
    add_content_box(slide, Inches(0.5), y, Inches(12.3), Inches(0.85),
        [desc], bg_color=bg_c, title=title, title_color=txt_c,
        border_color=txt_c, font_size=12)

add_reference_bar(slide, "Reference: Miller et al. (1992). MET Manual; NIMHANS (2016); Project MATCH (1997); UKATT (2005).")

# SLIDE: Comprehensive References
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'white', 'light_blue')
add_title_bar(slide, "References", 'deep_blue')

add_content_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(5.2),
    [
        "PRIMARY SOURCES:",
        "",
        "Miller, W.R., Zweben, A., DiClemente, C.C., &",
        "  Rychtarik, R.G. (1992). Motivational Enhancement",
        "  Therapy Manual. NIAAA Project MATCH, Vol. 2.",
        "",
        "NIMHANS (2016). Substance Use Disorders: A Manual",
        "  for Physicians. National Institute of Mental Health",
        "  and Neuro Sciences, Bangalore.",
        "",
        "Miller, W.R. & Rollnick, S. (2013). Motivational",
        "  Interviewing, 3rd ed. Guilford Press.",
        "",
        "Prochaska, J.O. & DiClemente, C.C. (1984). The",
        "  Transtheoretical Approach. Dow Jones-Irwin.",
        "",
        "Project MATCH Research Group (1997). Matching",
        "  treatments to patient heterogeneity. J Studies on",
        "  Alcohol, 58, 7-29.",
    ],
    bg_color='white', title="", title_color='deep_blue',
    border_color='deep_blue', font_size=11)

add_content_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(5.2),
    [
        "ADDITIONAL REFERENCES:",
        "",
        "UKATT Research Team (2005). Effectiveness of",
        "  treatment for alcohol problems. BMJ, 331, 544.",
        "",
        "Nadkarni, A. et al. (2017). Counselling for alcohol",
        "  problems in India. The Lancet, 389, 186-195.",
        "",
        "Lundahl, B.W. et al. (2010). A meta-analysis of MI.",
        "  Patient Education & Counseling, 80(1), 94-109.",
        "",
        "Hettema, J. et al. (2005). MI. Annual Review of",
        "  Clinical Psychology, 1, 91-111.",
        "",
        "Bandura, A. (1977). Self-efficacy. Psychological",
        "  Review, 84(2), 191-215.",
        "",
        "Marlatt, G.A. & Donovan, D.M. (2005). Relapse",
        "  Prevention, 2nd ed. Guilford Press.",
    ],
    bg_color='white', title="", title_color='green',
    border_color='green', font_size=11)

add_reference_bar(slide, "All references are from peer-reviewed journals, published manuals, and institutional publications.")

# SLIDE: Thank You
slide = prs.slides.add_slide(blank_layout)
add_gradient_bg(slide, 'deep_blue', 'navy')

tb = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.font.name = 'Times New Roman'
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = ""
p2.space_after = Pt(20)

p3 = tf.add_paragraph()
p3.text = "\"People are generally better persuaded by the reasons"
p3.font.size = Pt(18)
p3.font.italic = True
p3.font.color.rgb = COLORS['gold']
p3.font.name = 'Times New Roman'
p3.alignment = PP_ALIGN.CENTER

p4 = tf.add_paragraph()
p4.text = "which they have themselves discovered"
p4.font.size = Pt(18)
p4.font.italic = True
p4.font.color.rgb = COLORS['gold']
p4.font.name = 'Times New Roman'
p4.alignment = PP_ALIGN.CENTER

p5 = tf.add_paragraph()
p5.text = "than by those which have come into the minds of others.\""
p5.font.size = Pt(18)
p5.font.italic = True
p5.font.color.rgb = COLORS['gold']
p5.font.name = 'Times New Roman'
p5.alignment = PP_ALIGN.CENTER

p6 = tf.add_paragraph()
p6.text = "— Blaise Pascal"
p6.font.size = Pt(14)
p6.font.color.rgb = COLORS['light_blue']
p6.font.name = 'Times New Roman'
p6.alignment = PP_ALIGN.CENTER

add_reference_bar(slide, "Based on: MET Manual (NIAAA, 1992) & NIMHANS Substance Use Disorders Manual (2016)")


# ============================================================
# SAVE THE PRESENTATION
# ============================================================
output_path = '/projects/sandbox/Dango-kiro/MET_Comprehensive_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
