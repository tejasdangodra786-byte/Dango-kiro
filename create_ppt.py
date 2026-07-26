#!/usr/bin/env python3
"""
Research Proposal PPT - DETAILED VERSION
Rorschach Cognitive Triad → Subjective Loneliness in Schizophrenia
Each slide has substantial academic content with visual depth.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==================== COLOR PALETTE ====================
PRIMARY = RGBColor(0x0D, 0x1B, 0x2A)      # Deep midnight navy
SECONDARY = RGBColor(0x1B, 0x6B, 0x93)    # Ocean teal
ACCENT = RGBColor(0xC0, 0x39, 0x2B)       # Deep crimson
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)     # Soft grey-blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1C, 0x1C, 0x2E)
GOLD = RGBColor(0xC8, 0x9B, 0x2B)         # Rich gold
GREEN = RGBColor(0x1E, 0x8C, 0x4E)        # Forest green
PURPLE = RGBColor(0x5B, 0x2C, 0x8C)       # Royal purple
GREY = RGBColor(0x5D, 0x6D, 0x7E)         # Neutral grey

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)



# ==================== HELPER FUNCTIONS ====================

def bg(slide, color):
    """Full-slide background"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def bar(slide, left, top, width, height, color):
    """Colored rectangle bar"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def box(slide, left, top, width, height, color, radius=True):
    """Rounded box"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, left, top, width, height):
    """Add textbox and return text_frame"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf

def para(tf, text, size=14, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, space_after=6):
    """Add a paragraph to a text frame"""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(space_after)
    return p



# ==================== SLIDE 1: TITLE ====================
def slide_title():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, PRIMARY)
    # Top gold line
    bar(slide, 0, 0, prs.slide_width, Inches(0.12), GOLD)
    # Bottom gold line
    bar(slide, 0, Inches(7.38), prs.slide_width, Inches(0.12), GOLD)
    # Left vertical accent
    bar(slide, Inches(0.8), Inches(1.0), Inches(0.08), Inches(5.5), SECONDARY)

    tf = txt(slide, Inches(1.2), Inches(1.2), Inches(11), Inches(4.5))
    para(tf, "Rorschach Indices of the Cognitive Triad", size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    para(tf, "(Processing, Mediation, and Ideation)", size=26, color=SECONDARY, align=PP_ALIGN.LEFT)
    para(tf, "", size=12)
    para(tf, "as Predictors of Subjective Loneliness", size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    para(tf, "in Patients with Schizophrenia", size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    para(tf, "", size=20)
    para(tf, "MPhil Clinical Psychology — Research Proposal Presentation", size=16, color=GOLD, align=PP_ALIGN.LEFT)

    tf2 = txt(slide, Inches(1.2), Inches(6.0), Inches(11), Inches(1.2))
    para(tf2, "Candidate: [Your Name]    |    Guide: [Guide Name, Designation]", size=13, color=RGBColor(0xAA,0xBB,0xCC))
    para(tf2, "[Institution Name]    |    Department of Clinical Psychology    |    2026-2027", size=13, color=RGBColor(0xAA,0xBB,0xCC))

slide_title()

# ==================== SLIDE 2: INTRODUCTION PART 1 ====================
def slide_intro1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, SECONDARY)
    bar(slide, 0, Inches(7.38), prs.slide_width, Inches(0.12), SECONDARY)

    tf = txt(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    para(tf, "Introduction: The Problem of Loneliness in Schizophrenia", size=26, bold=True, color=PRIMARY)

    # Accent underline
    bar(slide, Inches(0.6), Inches(1.0), Inches(4), Pt(4), GOLD)

    tf2 = txt(slide, Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.8))
    para(tf2, "Schizophrenia & The Loneliness Paradox:", size=18, bold=True, color=SECONDARY, space_after=8)
    para(tf2, "Schizophrenia affects 0.3-0.7% of the global population (APA, 2013). While antipsychotics effectively manage hallucinations and delusions, social and functional impairments remain largely treatment-resistant. Among these, subjective loneliness — defined as the painful discrepancy between desired and actual social connection — is one of the most pervasive yet clinically neglected experiences. Critically, loneliness is NOT the same as social isolation: a patient can be surrounded by family yet feel profoundly alone.", size=14, color=DARK_TEXT, space_after=12)
    para(tf2, "Clinical Significance — Why Loneliness Matters:", size=18, bold=True, color=SECONDARY, space_after=8)
    para(tf2, "Meta-analysis (Michalska da Rocha et al., 2018): Robust positive relationship between loneliness and psychotic symptoms across multiple studies. Loneliness in schizophrenia is now linked to: (a) increased psychiatric re-hospitalization (Fortuna et al., 2020), (b) elevated suicide risk — loneliness significantly predicts suicidal ideation even after controlling for depression (Yen et al., 2023), (c) accelerated cognitive decline (Wang et al., 2026), (d) increased all-cause mortality (Green et al., 2023). The WHO Commission on Social Connection (2023) has declared loneliness a global public health priority. A 2024 Indian scoping review identified loneliness as an 'emerging public mental health predicament' in the country.", size=14, color=DARK_TEXT, space_after=12)
    para(tf2, "The Question: If loneliness is this destructive, what CAUSES it in schizophrenia? Can we identify its cognitive roots?", size=15, bold=True, color=ACCENT, space_after=4)

slide_intro1()



# ==================== SLIDE 3: INTRODUCTION PART 2 - COGNITIVE MODEL ====================
def slide_intro2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, PURPLE)

    tf = txt(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    para(tf, "Theoretical Foundation: The Cognitive Model of Loneliness", size=26, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(1.0), Inches(4), Pt(4), PURPLE)

    tf2 = txt(slide, Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.8))
    para(tf2, "Cacioppo & Hawkley (2009) — Cognitive Model of Loneliness:", size=17, bold=True, color=PURPLE, space_after=8)
    para(tf2, "Loneliness is NOT simply the absence of social contact. It is a COGNITIVE-PERCEPTUAL phenomenon arising from how the individual processes, perceives, and interprets their social environment. The model proposes a self-reinforcing cycle:", size=14, color=DARK_TEXT, space_after=10)
    para(tf2, "    Step 1: Impaired PROCESSING of social environment (hypervigilant or careless scanning)", size=14, color=DARK_TEXT, space_after=4)
    para(tf2, "    Step 2: MISPERCEPTION of neutral social cues as threatening or rejecting", size=14, color=DARK_TEXT, space_after=4)
    para(tf2, "    Step 3: Distorted REASONING about others' intentions ('they don't want me')", size=14, color=DARK_TEXT, space_after=4)
    para(tf2, "    Step 4: WITHDRAWAL from social situations to avoid perceived threat", size=14, color=DARK_TEXT, space_after=4)
    para(tf2, "    Step 5: Withdrawal CONFIRMS and DEEPENS loneliness (vicious cycle)", size=14, color=DARK_TEXT, space_after=12)
    para(tf2, "Application to Schizophrenia (Michalska da Rocha et al., 2018; Green et al., 2023):", size=17, bold=True, color=PURPLE, space_after=8)
    para(tf2, "In schizophrenia, where cognitive impairment in perception, processing speed, and reasoning is a CARDINAL FEATURE of the disorder, this cognitive pathway to loneliness is substantially magnified. The cognitive deficits characteristic of schizophrenia (perceptual distortion, disorganized thinking, impaired scanning) map directly onto the cognitive mechanisms that the loneliness model identifies as causal. This raises a testable prediction: the MORE impaired a patient's cognitive processing is, the MORE lonely they should feel.", size=14, color=DARK_TEXT, space_after=10)
    para(tf2, "Key Insight: The Rorschach Cognitive Triad (Processing → Mediation → Ideation) measures EXACTLY these three sequential cognitive steps.", size=15, bold=True, color=ACCENT, space_after=4)

slide_intro2()

# ==================== SLIDE 4: THE RORSCHACH COGNITIVE TRIAD EXPLAINED ====================
def slide_triad():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, PRIMARY)
    bar(slide, 0, 0, prs.slide_width, Inches(0.08), GOLD)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.8))
    para(tf, "The Rorschach Cognitive Triad (Exner, 2003) — A Sequential Model", size=24, bold=True, color=GOLD)

    # Three detailed boxes
    # BOX 1: PROCESSING
    box(slide, Inches(0.4), Inches(1.2), Inches(4.0), Inches(5.8), SECONDARY)
    tf1 = txt(slide, Inches(0.6), Inches(1.3), Inches(3.6), Inches(5.6))
    para(tf1, "STEP 1: PROCESSING", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=6)
    para(tf1, "\"How do I SCAN my world?\"", size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER, space_after=10)
    para(tf1, "What it measures:", size=12, bold=True, color=WHITE, space_after=4)
    para(tf1, "How the person organizes and scans the stimulus field — the quality of INPUT entering the cognitive system.", size=11, color=WHITE, space_after=8)
    para(tf1, "Key Variables:", size=12, bold=True, color=WHITE, space_after=4)
    para(tf1, "Zf = Processing effort/motivation", size=11, color=WHITE, space_after=2)
    para(tf1, "Zd = Scanning efficiency (+3=over-incorporative, -3=under-incorporative/hasty)", size=11, color=WHITE, space_after=2)
    para(tf1, "DQ+ = Cognitive complexity (synthesis capacity)", size=11, color=WHITE, space_after=2)
    para(tf1, "DQv = Vague/impressionistic processing", size=11, color=WHITE, space_after=2)
    para(tf1, "W:D:Dd = Attentional distribution", size=11, color=WHITE, space_after=8)
    para(tf1, "In Schizophrenia:", size=12, bold=True, color=GOLD, space_after=4)
    para(tf1, "Patients often show extreme Zd (hasty scanning), low DQ+ (cannot synthesize), elevated DQv (vague thinking).", size=11, color=WHITE, space_after=6)
    para(tf1, "Link to Loneliness:", size=12, bold=True, color=GOLD, space_after=4)
    para(tf1, "If you scan the social world carelessly, you MISS connection opportunities.", size=11, color=WHITE)

    # BOX 2: MEDIATION
    box(slide, Inches(4.6), Inches(1.2), Inches(4.0), Inches(5.8), ACCENT)
    tf2 = txt(slide, Inches(4.8), Inches(1.3), Inches(3.6), Inches(5.6))
    para(tf2, "STEP 2: MEDIATION", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=6)
    para(tf2, "\"How ACCURATELY do I perceive?\"", size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER, space_after=10)
    para(tf2, "What it measures:", size=12, bold=True, color=WHITE, space_after=4)
    para(tf2, "How accurately the person TRANSLATES perceptual input — the fidelity of reality testing.", size=11, color=WHITE, space_after=8)
    para(tf2, "Key Variables:", size=12, bold=True, color=WHITE, space_after=4)
    para(tf2, "XA% = Overall reality testing adequacy", size=11, color=WHITE, space_after=2)
    para(tf2, "WDA% = Reality testing for obvious stimuli", size=11, color=WHITE, space_after=2)
    para(tf2, "X-% = Frequency of outright MISPERCEPTION", size=11, color=WHITE, space_after=2)
    para(tf2, "P = Popular responses (seeing the obvious)", size=11, color=WHITE, space_after=2)
    para(tf2, "X+% = Perceptual conventionality", size=11, color=WHITE, space_after=8)
    para(tf2, "In Schizophrenia:", size=12, bold=True, color=GOLD, space_after=4)
    para(tf2, "Markedly elevated X-% (>29%), low XA% (<70%), low P. Pervasive misperception.", size=11, color=WHITE, space_after=6)
    para(tf2, "Link to Loneliness:", size=12, bold=True, color=GOLD, space_after=4)
    para(tf2, "If you MISREAD people's faces and intentions, you feel rejected when you are not.", size=11, color=WHITE)

    # BOX 3: IDEATION
    box(slide, Inches(8.8), Inches(1.2), Inches(4.0), Inches(5.8), PURPLE)
    tf3 = txt(slide, Inches(9.0), Inches(1.3), Inches(3.6), Inches(5.6))
    para(tf3, "STEP 3: IDEATION", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=6)
    para(tf3, "\"How LOGICALLY do I think?\"", size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER, space_after=10)
    para(tf3, "What it measures:", size=12, bold=True, color=WHITE, space_after=4)
    para(tf3, "The quality of reasoning and conceptualization — whether thinking is organized, logical, and reality-based.", size=11, color=WHITE, space_after=8)
    para(tf3, "Key Variables:", size=12, bold=True, color=WHITE, space_after=4)
    para(tf3, "WSum6 = Formal thought disorder composite", size=11, color=WHITE, space_after=2)
    para(tf3, "Level 2 = Severe/bizarre thought slippage", size=11, color=WHITE, space_after=2)
    para(tf3, "M quality = Reasoning about PEOPLE", size=11, color=WHITE, space_after=2)
    para(tf3, "Ma:Mp = Active vs passive thinking", size=11, color=WHITE, space_after=2)
    para(tf3, "a:p = Cognitive flexibility/rigidity", size=11, color=WHITE, space_after=8)
    para(tf3, "In Schizophrenia:", size=12, bold=True, color=GOLD, space_after=4)
    para(tf3, "Elevated WSum6, Level 2 present, M- (distorted person-reasoning), passive Mp.", size=11, color=WHITE, space_after=6)
    para(tf3, "Link to Loneliness:", size=12, bold=True, color=GOLD, space_after=4)
    para(tf3, "If you REASON about people illogically, your conclusions about relationships are wrong.", size=11, color=WHITE)

    # Arrows between boxes
    for x in [Inches(4.35), Inches(8.55)]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(3.8), Inches(0.3), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

slide_triad()



# ==================== SLIDE 5: REVIEW OF LITERATURE - LONELINESS IN PSYCHOSIS ====================
def slide_rol1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, SECONDARY)
    bar(slide, 0, Inches(7.38), prs.slide_width, Inches(0.12), SECONDARY)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Review of Literature I: Loneliness in Schizophrenia — Prevalence & Impact", size=24, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(5), Pt(4), SECONDARY)

    tf2 = txt(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.1))
    para(tf2, "1. Michalska da Rocha et al. (2018) — Schizophrenia Bulletin, 44(1), 114-125", size=14, bold=True, color=SECONDARY, space_after=4)
    para(tf2, "   Landmark meta-analysis pooling all available studies on loneliness in psychotic disorders. Found a statistically significant positive relationship between loneliness and psychotic symptoms. Patients with psychosis consistently reported higher loneliness than general population. Loneliness was associated with positive symptoms, negative symptoms, depression, reduced QoL, and lower social functioning. CRITICAL GAP IDENTIFIED: the cognitive MECHANISMS driving this relationship remained unknown. Authors explicitly called for research into cognitive-perceptual factors explaining heightened loneliness — making the present study a direct response.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "2. Green, Catalano, Engel & Horan (2023) — Schizophrenia Bulletin, 49(5), 1112-1126", size=14, bold=True, color=SECONDARY, space_after=4)
    para(tf2, "   Proposed the 'social homeostasis' model synthesizing evolutionary loneliness theory with cognitive neuroscience of serious mental illness. Key argument: loneliness in SMI is driven by impairments in the cognitive processing of social information — including perception, interpretation, and reasoning about social stimuli. Chronic loneliness heightens social threat sensitivity and alters brain networks supporting social cognition. The paper explicitly identifies cognitive-perceptual processing deficits as the PRIMARY MECHANISM through which social disconnection becomes chronic and entrenched in schizophrenia.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "3. Yen, Lin et al. (2023) — npj Schizophrenia, 9, Article 40", size=14, bold=True, color=SECONDARY, space_after=4)
    para(tf2, "   Cross-sectional study of 300 participants (267 schizophrenia, 33 schizoaffective). Loneliness was significantly associated with increased SUICIDE RISK and depression severity, even after controlling for demographics. Self-esteem moderated the loneliness-depression pathway; perceived friend support moderated the loneliness-suicide pathway. Establishes loneliness as a variable with LIFE-THREATENING consequences — not merely an uncomfortable feeling but a clinical priority requiring identification of its predictors.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "4. Fortuna et al. (2020) — Social Psychiatry & Psychiatric Epidemiology, 55, 1421-1428", size=14, bold=True, color=SECONDARY, space_after=4)
    para(tf2, "   Established that loneliness in SMI is independently associated with physical health conditions and higher psychiatric hospitalization rates. Lonely patients with schizophrenia had significantly more emergency department visits and inpatient admissions, representing both a human cost and an enormous healthcare burden that justifies urgent research into modifiable cognitive predictors.", size=12, color=DARK_TEXT, space_after=4)

slide_rol1()

# ==================== SLIDE 6: ROL - COGNITIVE CORRELATES ====================
def slide_rol2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, PURPLE)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Review of Literature II: Cognitive-Perceptual Correlates of Loneliness", size=24, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(5), Pt(4), PURPLE)

    tf2 = txt(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.1))
    para(tf2, "5. Thibaudeau et al. (2023) — Schizophrenia Research, 256, 29-37", size=14, bold=True, color=PURPLE, space_after=4)
    para(tf2, "   Examined loneliness and social cognitive abilities in 87 schizophrenia-spectrum patients vs. controls. Used Latent Class Analysis to identify patient subgroups. KEY FINDING: In patients with impaired social cognition (deficits in mentalizing and emotion recognition), loneliness was SIGNIFICANTLY AND NEGATIVELY associated with these cognitive abilities. In patients with intact social cognition, NO association existed. This demonstrates that cognitive-perceptual impairment is a SPECIFIC MECHANISM driving loneliness in a SUBGROUP of patients — and that identifying this subgroup through cognitive assessment (e.g., Rorschach) could enable targeted intervention.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "6. Hajduk, Harvey, Penn & Pinkham (2023) — Schizophrenia Research, 256, 38-46", size=14, bold=True, color=PURPLE, space_after=4)
    para(tf2, "   Used structural equation modeling in schizophrenia patients. SOCIAL THREAT BIAS — the tendency to misperceive neutral social stimuli as threatening — was DIRECTLY linked to increased loneliness AND indirectly through decreased social connection. Negative symptoms were directly linked to loneliness; affective symptoms operated through threat bias. Authors recommended that social threat bias should be a treatment target for loneliness. CRITICAL LINK TO PRESENT STUDY: The Rorschach Mediation cluster (X-%, XA%, FQ-) measures PRECISELY this construct — the tendency to misperceive and distort perceptual input. If computerized threat bias predicts loneliness, Rorschach perceptual distortion should similarly predict it.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "7. Jaya, Hempel & Lincoln (2024) — Schizophrenia Research, 269, 118-125", size=14, bold=True, color=PURPLE, space_after=4)
    para(tf2, "   Tested whether Selective Attention to Threat (ATB) and External Attribution Bias (EAB) mediated the loneliness-psychosis link. Confirmed biased cognition plays a key causal role. Attention to threat partially mediated the relationship. RELEVANCE: The Rorschach Processing cluster (Zd) measures scanning efficiency — whether a person hastily scans or excessively scans their environment. This is the projective equivalent of selective attention bias. The present study tests whether scanning pattern (Zd) predicts loneliness.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "8. Wang, Chen & Liu (2026) — European Archives of Psychiatry & Clinical Neuroscience", size=14, bold=True, color=PURPLE, space_after=4)
    para(tf2, "   Most recent study: loneliness DIRECTLY associated with cognitive impairment in schizophrenia, with sleep quality and anxiety as mediating pathways. Confirms bidirectional relationship: cognitive impairment may cause loneliness, and loneliness may worsen cognition. Present study examines the first direction.", size=12, color=DARK_TEXT, space_after=4)

slide_rol2()



# ==================== SLIDE 7: ROL - RORSCHACH IN SCHIZOPHRENIA ====================
def slide_rol3():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, GREEN)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Review of Literature III: Rorschach Cognitive Assessment in Schizophrenia", size=24, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(5), Pt(4), GREEN)

    tf2 = txt(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.1))
    para(tf2, "9. Ilonen, Leinonen, Wallin & Wahlberg (2012) — Psychological Assessment, 25(1), 253-263", size=14, bold=True, color=GREEN, space_after=4)
    para(tf2, "   SEMINAL STUDY: Examined Rorschach cognitive variables (psychological complexity, thought disorder, interpersonal representations) in 72 outpatients with schizophrenia (mean age 51.2). Psychological complexity — a Rorschach composite reflecting capacity for problem-solving and organizing surroundings — correlated with functional capacity (r=.30) and social skills capacity (r=.34). In multiple regression, Rorschach complexity predicted functioning (B=0.23, p=0.02) and social skills (B=0.35, p<0.01) AFTER controlling for neurocognition and psychopathology. LANDMARK CONCLUSION: The Rorschach captures cognitive dimensions that predict real-world outcomes ABOVE AND BEYOND standard neuropsychological tests. Present study extends this to loneliness as outcome.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "10. Singh, Shukla & Mishra (2021) — Indian J Psychiatric Nursing, 30(1), 49-55", size=14, bold=True, color=GREEN, space_after=4)
    para(tf2, "   Indian study examining Rorschach structural profiles and their correlation with psychopathology (PANSS) in schizophrenia patients. Confirmed that Rorschach provides specific knowledge about coping style, emotional processing, stress management, mediation (reality testing), ideation (thought quality), self-perception, and interpersonal relationships. Rorschach variables significantly correlated with PANSS severity. ESTABLISHES: Rorschach cognitive variables are clinically relevant in Indian schizophrenia populations. GAP: Did NOT examine relationship with subjective experiential outcomes like loneliness.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "11. Kimoto et al. (2016) — Neuropsychiatric Disease & Treatment, 12, 2403-2410", size=14, bold=True, color=GREEN, space_after=4)
    para(tf2, "   Compared Rorschach profiles between schizophrenia and autism spectrum disorder in young adults. Schizophrenia showed significantly higher FQ- (perceptual distortion) and DQo (cognitively simple responses). Confirmed that Rorschach mediation variables (X-%, XA%) are DISTINCTIVELY impaired in schizophrenia — stronger perception distortions and simpler recognition. Present study builds on this by testing whether this distinctive perceptual distortion predicts the SUBJECTIVE EXPERIENCE of loneliness.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "12. Jo, Lee, Lee & Joo (2024) — Science Progress, 107(3)", size=14, bold=True, color=GREEN, space_after=4)
    para(tf2, "   Compared Rorschach patterns between Kraepelinian (severe) vs. DSM-defined schizophrenia. Kraepelinian patients showed significantly increased DV2 scores (severe thought disorder) and decreased D scores. Demonstrates that Rorschach ideation variables (special scores) are SENSITIVE TO SEVERITY GRADIENTS within schizophrenia — supporting their use as continuous predictors of loneliness degree.", size=12, color=DARK_TEXT, space_after=4)

slide_rol3()

# ==================== SLIDE 8: ROL - ADDITIONAL EVIDENCE ====================
def slide_rol4():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, GOLD)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Review of Literature IV: Additional Supporting Evidence", size=24, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(5), Pt(4), GOLD)

    tf2 = txt(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0))
    para(tf2, "13. Hsu, Huang & Lin (2024) — Schizophrenia Research, 266, 171-178", size=14, bold=True, color=GOLD, space_after=4)
    para(tf2, "   Longitudinal study (2019-2021, N=166 inpatients). Social isolation significantly related to poor language abilities (B=-0.17, p=0.013) and delayed recall (B=-0.15, p=0.023). First study establishing longitudinal, directional link between social disconnection and cognitive decline in schizophrenia. Demonstrates that cognitive impairment and social isolation are BIDIRECTIONALLY related — present study examines whether specific cognitive processes (processing/mediation/ideation) are differentially associated with the SUBJECTIVE experience of aloneness.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "14. Stain, Galletly, Clark et al. (2016) — Psychiatry Research, 235, 170-176", size=14, bold=True, color=GOLD, space_after=4)
    para(tf2, "   Assessed loneliness (UCLA-R) in 87 non-depressed schizophrenia/schizoaffective patients vs. 58 controls. Patients reported significantly greater loneliness. Loneliness was associated with impaired social cognition and identified as risk factor for broad-based morbidity even after controlling for demographics. However, SPECIFIC cognitive processes driving loneliness were NOT examined — present study fills this gap using Rorschach Cognitive Triad.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "15. Regev et al. (2024) — Schizophrenia Research, 270, 215-223", size=14, bold=True, color=GOLD, space_after=4)
    para(tf2, "   Patients with HIGH social exclusion were 4.24 times more likely to have cognitive impairment than those with LOW exclusion. Verbal learning was the cognitive function most related to exclusion domains. Establishes dose-response relationship: MORE cognitive impairment = MORE disconnection. Present study extends this using performance-based projective assessment (Rorschach) rather than neuropsychological tests.", size=12, color=DARK_TEXT, space_after=12)

    para(tf2, "16. Lin, Yen et al. (2022) — Int J Environ Res Public Health, 19(14), 8443", size=14, bold=True, color=GOLD, space_after=4)
    para(tf2, "   Evaluated psychometric properties of three UCLA Loneliness Scale versions SPECIFICALLY in schizophrenia/schizoaffective populations. All versions showed acceptable properties; Version 3 performed best overall. Confirmed self-reported loneliness CAN be reliably measured in psychotic populations — addressing the common objection that these patients cannot validly self-report. Supports use of UCLA-LS V3 in present study.", size=12, color=DARK_TEXT, space_after=4)

    para(tf2, "17. Sharma, Narang & Kumar (2023) — Asian J Psychiatry, 88, 103742 [INDIAN]", size=14, bold=True, color=GOLD, space_after=4)
    para(tf2, "   Indian cross-sectional study finding cognitive deficits affect functioning MORE than positive/negative symptoms. Called for research linking cognitive assessment to clinically meaningful outcomes in India. Present study directly responds to this call.", size=12, color=DARK_TEXT, space_after=4)

slide_rol4()



# ==================== SLIDE 9: RESEARCH GAP ====================
def slide_gap():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, PRIMARY)
    bar(slide, 0, 0, prs.slide_width, Inches(0.08), ACCENT)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Research Gap: What Remains Unknown", size=26, bold=True, color=ACCENT)

    gaps = [
        ("GAP 1:", "NO study globally has examined whether Rorschach cognitive triad variables (processing, mediation, ideation) predict subjective loneliness in ANY psychiatric population. The Rorschach-loneliness connection is entirely unexplored territory."),
        ("GAP 2:", "Existing loneliness research in schizophrenia relies exclusively on computerized social cognition tasks (TASIT, ER-40, Hinting Task) that are UNAVAILABLE in most Indian clinical settings. The Rorschach — universally available in Indian MPhil programs — has never been tested as a loneliness predictor."),
        ("GAP 3:", "Hajduk et al. (2023) showed computerized 'social threat bias' predicts loneliness. The Rorschach Mediation cluster (X-%, XA%) is the projective EQUIVALENT of threat bias measurement — but nobody has tested this equivalence."),
        ("GAP 4:", "Ilonen et al. (2012) showed Rorschach cognition predicts FUNCTIONAL CAPACITY — but functional capacity ('can the patient work?') is NOT the same as subjective loneliness ('does the patient feel alone?'). These are distinct outcomes requiring separate investigation."),
        ("GAP 5:", "The question of WHICH cognitive step (processing vs. mediation vs. ideation) most strongly predicts loneliness is completely unanswered — yet this specificity is ESSENTIAL for targeted cognitive remediation planning."),
        ("GAP 6:", "NO Indian study has examined either (a) Rorschach cognitive variables in relation to any outcome beyond symptom severity, OR (b) loneliness and its cognitive correlates in Indian schizophrenia populations."),
    ]

    tf2 = txt(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.2))
    for label, text in gaps:
        para(tf2, label, size=13, bold=True, color=ACCENT, space_after=2)
        para(tf2, "  " + text, size=12, color=WHITE, space_after=10)

slide_gap()

# ==================== SLIDE 10: RATIONALE ====================
def slide_rationale():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, GREEN)
    bar(slide, 0, Inches(7.38), prs.slide_width, Inches(0.12), GREEN)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Rationale of the Study: Why This Research is Necessary", size=26, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(4), Pt(4), GREEN)

    tf2 = txt(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0))
    para(tf2, "FIRST — Clinical Urgency:", size=16, bold=True, color=GREEN, space_after=4)
    para(tf2, "Loneliness in schizophrenia is no longer a 'soft' or peripheral concern. It is an independent clinical variable associated with psychiatric hospitalization (Fortuna, 2020), suicide risk (Yen, 2023), accelerated cognitive decline (Wang, 2026), and all-cause mortality (Green, 2023). Identifying modifiable predictors of loneliness is clinically urgent.", size=13, color=DARK_TEXT, space_after=12)

    para(tf2, "SECOND — Theoretical Precision:", size=16, bold=True, color=GREEN, space_after=4)
    para(tf2, "The Cognitive Model of Loneliness (Cacioppo & Hawkley, 2009), now applied to psychosis (Michalska da Rocha, 2018; Green, 2023), establishes that loneliness arises from cognitive-perceptual disturbances — specifically, impaired scanning, misperception of social stimuli, and distorted reasoning. The Rorschach Cognitive Triad (Processing → Mediation → Ideation) measures EXACTLY these sequential operations. The theoretical fit is not approximate — it is precise.", size=13, color=DARK_TEXT, space_after=12)

    para(tf2, "THIRD — Methodological Advantage of Projective Assessment:", size=16, bold=True, color=GREEN, space_after=4)
    para(tf2, "The Rorschach offers critical advantages over neuropsychological testing in psychotic populations: (a) cannot be faked or manipulated (50-80% of patients lack insight — self-report is unreliable), (b) measures cognitive PROCESS not just product, (c) ecologically valid (ambiguous stimuli simulate real-world social ambiguity), (d) no floor effects, (e) single 30-min administration yields all predictor variables simultaneously.", size=13, color=DARK_TEXT, space_after=12)

    para(tf2, "FOURTH — Practical Clinical Utility:", size=16, bold=True, color=GREEN, space_after=4)
    para(tf2, "If we identify WHICH cognitive step most strongly predicts loneliness, it directly informs individualized rehabilitation: Should we train perceptual accuracy? Scanning efficiency? Logical reasoning? This enables TARGETED cognitive remediation for lonely patients — a priority identified in Indian clinical practice guidelines (IPS, 2024).", size=13, color=DARK_TEXT, space_after=4)

slide_rationale()



# ==================== SLIDE 11: AIM, OBJECTIVES ====================
def slide_aim():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, SECONDARY)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Aim & Objectives", size=26, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(3), Pt(4), SECONDARY)

    # AIM box
    box(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.4), PRIMARY)
    tf_aim = txt(slide, Inches(0.8), Inches(1.3), Inches(11.8), Inches(1.2))
    para(tf_aim, "AIM:", size=16, bold=True, color=GOLD, space_after=4)
    para(tf_aim, "To examine whether Rorschach indices of the Cognitive Triad (Information Processing, Cognitive Mediation, and Ideation) predict subjective loneliness in clinically stable patients with schizophrenia, after controlling for negative symptom severity.", size=14, color=WHITE, space_after=4)

    tf2 = txt(slide, Inches(0.5), Inches(2.9), Inches(12.3), Inches(4.3))
    para(tf2, "OBJECTIVES:", size=16, bold=True, color=SECONDARY, space_after=8)
    objectives = [
        "O1: To assess the Rorschach Cognitive Triad profile (Processing, Mediation, and Ideation cluster variables) of clinically stable patients with schizophrenia.",
        "O2: To assess the level of subjective loneliness in these patients using the UCLA Loneliness Scale Version 3.",
        "O3: To examine correlations between Rorschach PROCESSING variables (Zf, Zd, DQ+, DQv, W:D:Dd) and UCLA Loneliness scores.",
        "O4: To examine correlations between Rorschach MEDIATION variables (XA%, WDA%, X-%, X+%, Xu%, P) and UCLA Loneliness scores.",
        "O5: To examine correlations between Rorschach IDEATION variables (WSum6, Level 2, Ma:Mp, a:p, M quality) and UCLA Loneliness scores.",
        "O6: To determine which cluster of the Cognitive Triad is the STRONGEST predictor of loneliness after controlling for PANSS Negative Subscale score."
    ]
    for obj in objectives:
        para(tf2, obj, size=13, color=DARK_TEXT, space_after=8)

slide_aim()

# ==================== SLIDE 12: HYPOTHESES ====================
def slide_hypotheses():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, GOLD)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Hypotheses", size=26, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(3), Pt(4), GOLD)

    tf2 = txt(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(6.0))
    hyps = [
        ("H1 (Processing):", "There will be a significant negative correlation between Rorschach Processing efficiency (normal-range Zd, higher DQ+, lower DQv) and subjective loneliness scores. Patients who process their environment more efficiently and with greater cognitive complexity will report LESS loneliness — because efficient processing enables detection of social connection opportunities.", SECONDARY),
        ("H2 (Mediation):", "There will be a significant positive correlation between Rorschach Mediation distortion (higher X-%, lower XA%, fewer Populars) and subjective loneliness scores. Patients who more severely MISPERCEIVE reality will report GREATER loneliness — because persistent misperception of social cues as threatening or rejecting drives withdrawal and felt isolation.", ACCENT),
        ("H3 (Ideation):", "There will be a significant positive correlation between Rorschach Ideation disturbance (higher WSum6, more M-, elevated passive Mp) and subjective loneliness scores. Patients with more disordered or passive thinking will report GREATER loneliness — because illogical reasoning about relationships produces faulty conclusions that prevent connection.", PURPLE),
        ("H4 (Incremental Validity):", "The Rorschach Cognitive Triad variables will significantly predict subjective loneliness ABOVE AND BEYOND the contribution of negative symptoms (PANSS-N) — demonstrating that cognitive-perceptual processing explains loneliness beyond what negative symptom overlap accounts for.", GREEN),
        ("H5 (Specificity):", "The MEDIATION cluster (reality testing / perceptual accuracy) will emerge as the strongest predictor of loneliness among the three clusters — because misperception of social reality is the most immediate and direct cognitive barrier to feeling connected.", GOLD),
    ]
    for label, text, color in hyps:
        para(tf2, label, size=14, bold=True, color=color, space_after=2)
        para(tf2, "   " + text, size=12, color=DARK_TEXT, space_after=12)

slide_hypotheses()



# ==================== SLIDE 13: METHODOLOGY - DESIGN, SAMPLE, CRITERIA ====================
def slide_method1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, ACCENT)
    bar(slide, 0, Inches(7.38), prs.slide_width, Inches(0.12), ACCENT)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Methodology: Research Design, Sample & Criteria", size=26, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(4), Pt(4), ACCENT)

    tf2 = txt(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(6.0))
    para(tf2, "RESEARCH DESIGN:", size=16, bold=True, color=ACCENT, space_after=6)
    para(tf2, "Cross-sectional, correlational design.", size=13, color=DARK_TEXT, space_after=4)
    para(tf2, "Single-point assessment of Rorschach cognitive variables (IVs) and UCLA Loneliness score (DV) with PANSS-N as covariate.", size=13, color=DARK_TEXT, space_after=14)

    para(tf2, "SAMPLE:", size=16, bold=True, color=ACCENT, space_after=6)
    para(tf2, "N = 40-45 patients with schizophrenia", size=13, bold=True, color=DARK_TEXT, space_after=4)
    para(tf2, "Sampling: Purposive, from Psychiatry OPD / follow-up clinic", size=13, color=DARK_TEXT, space_after=4)
    para(tf2, "Justification (G*Power 3.1): Based on published effect sizes r=.25-.35 (Thibaudeau, 2023; Ilonen, 2012); f-squared=0.10-0.15 (small-medium); alpha=.05; power=.80; predictors=4 --> Required N = 38-45.", size=13, color=DARK_TEXT, space_after=4)
    para(tf2, "Feasibility: 40 clinically stable schizophrenia patients from a single OPD register within 3-4 months (3-4 patients/week).", size=13, color=DARK_TEXT, space_after=4)

    # Right column - Criteria
    tf3 = txt(slide, Inches(6.8), Inches(1.1), Inches(6.0), Inches(6.0))
    para(tf3, "INCLUSION CRITERIA:", size=15, bold=True, color=GREEN, space_after=6)
    inc = [
        "1. Diagnosis: Schizophrenia (F20.x, ICD-10/11) confirmed by consultant psychiatrist",
        "2. Age: 18-55 years, both sexes",
        "3. Clinical stability: Same antipsychotic regimen for minimum 8 weeks",
        "4. PANSS total < 80 (partial/full remission)",
        "5. Minimum literacy to comprehend UCLA-LS items (read aloud if needed)",
        "6. Valid Rorschach protocol: R >= 14 responses (per Exner guidelines)",
        "7. Informed written consent (+ guardian where applicable)",
    ]
    for item in inc:
        para(tf3, item, size=12, color=DARK_TEXT, space_after=5)

    para(tf3, "", size=8, space_after=4)
    para(tf3, "EXCLUSION CRITERIA:", size=15, bold=True, color=ACCENT, space_after=6)
    exc = [
        "1. Active substance use disorder (past 3 months)",
        "2. Intellectual disability (IQ < 70)",
        "3. Organic brain syndrome / neurological disorder",
        "4. ECT received in past 3 months",
        "5. Acute psychotic exacerbation at time of assessment",
        "6. Comorbid Major Depressive Episode (PANSS G6 > 4)",
    ]
    for item in exc:
        para(tf3, item, size=12, color=DARK_TEXT, space_after=5)

slide_method1()

# ==================== SLIDE 14: TOOLS ====================
def slide_tools():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, SECONDARY)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Tools & Measures", size=26, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(3), Pt(4), SECONDARY)

    # Tool 1
    tf2 = txt(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.2))
    para(tf2, "TOOL 1: Rorschach Inkblot Test — Exner Comprehensive System (Exner, 2003)", size=15, bold=True, color=SECONDARY, space_after=4)
    para(tf2, "10 bilaterally symmetrical inkblot cards administered under standard conditions (Response Phase + Inquiry Phase). Variables extracted from the Cognitive Triad: PROCESSING (Zf, Zd, DQ+, DQv, W:D:Dd), MEDIATION (XA%, WDA%, X-%, X+%, Xu%, P), IDEATION (WSum6, Level 2, Ma:Mp, a:p, M quality). Psychometrics: Inter-rater kappa = .85-.97 for location, .80-.90 for determinants, .85+ for form quality (Meyer et al., 2002). Indian norms available (Dubey, 2011). Administration time: 30-40 minutes. Inter-rater reliability established on 20% of protocols by independent scorer.", size=12, color=DARK_TEXT, space_after=14)

    para(tf2, "TOOL 2: UCLA Loneliness Scale — Version 3 (Russell, 1996)", size=15, bold=True, color=SECONDARY, space_after=4)
    para(tf2, "20-item self-report scale measuring subjective loneliness on 4-point Likert (1=Never to 4=Always). Total range: 20-80 (higher = lonelier). Internal consistency: alpha = .89-.94. Test-retest ICC = .73 over 12 months. Validated SPECIFICALLY in schizophrenia populations (Lin et al., 2022 — all three versions acceptable; V3 best). Indian psychometric properties established (Suri et al., 2020 — strong reliability in Indian samples). Can be administered in Hindi translation or read aloud. Administration: 5-8 minutes.", size=12, color=DARK_TEXT, space_after=14)

    para(tf2, "TOOL 3: PANSS — Negative Subscale (Kay, Fiszbein & Opler, 1987)", size=15, bold=True, color=SECONDARY, space_after=4)
    para(tf2, "7-item clinician-rated scale (blunted affect, emotional withdrawal, poor rapport, passive social withdrawal, difficulty in abstract thinking, lack of spontaneity, stereotyped thinking). Serves as COVARIATE to control for overlap between negative symptoms and loneliness. Also provides PANSS total to confirm inclusion criterion (total < 80). ICC = .83-.87; alpha = .73. Widely validated in Indian research. Administration: 30 minutes (semi-structured interview).", size=12, color=DARK_TEXT, space_after=14)

    para(tf2, "TOOL 4: Sociodemographic & Clinical Data Sheet (researcher-designed)", size=15, bold=True, color=SECONDARY, space_after=4)
    para(tf2, "Age, sex, education, marital status, occupation, residence, duration of illness, number of episodes, current antipsychotic and dose (chlorpromazine equivalent), duration of medication stability.", size=12, color=DARK_TEXT, space_after=4)

slide_tools()



# ==================== SLIDE 15: DATA ANALYSIS ====================
def slide_analysis():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, PURPLE)
    bar(slide, 0, Inches(7.38), prs.slide_width, Inches(0.12), PURPLE)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Data Analysis Plan", size=26, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(3), Pt(4), PURPLE)

    tf2 = txt(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0))
    para(tf2, "STEP 1 — Descriptive Statistics:", size=15, bold=True, color=PURPLE, space_after=4)
    para(tf2, "Mean, SD, frequencies, and range for all sociodemographic, clinical (PANSS), Rorschach cognitive triad variables, and UCLA Loneliness scores. This characterizes the sample and identifies distributional properties.", size=12, color=DARK_TEXT, space_after=10)

    para(tf2, "STEP 2 — Normality & Reliability:", size=15, bold=True, color=PURPLE, space_after=4)
    para(tf2, "Shapiro-Wilk test for normality → determines parametric vs. non-parametric pathway. Inter-rater reliability (Cohen's Kappa / ICC) computed on 20% of Rorschach protocols (n=8-9) rescored by independent trained scorer. Expected kappa > .80 based on published Exner CS reliability.", size=12, color=DARK_TEXT, space_after=10)

    para(tf2, "STEP 3 — Bivariate Correlations:", size=15, bold=True, color=PURPLE, space_after=4)
    para(tf2, "Pearson's r (or Spearman's rho if non-normal) between EACH individual Rorschach cognitive variable and UCLA Loneliness total score. This identifies which specific variables show significant zero-order relationships with loneliness before entering the multivariate model.", size=12, color=DARK_TEXT, space_after=10)

    para(tf2, "STEP 4 — PRIMARY ANALYSIS: Hierarchical Multiple Regression:", size=15, bold=True, color=PURPLE, space_after=4)
    para(tf2, "DV: UCLA Loneliness Score (continuous, 20-80)", size=13, bold=True, color=DARK_TEXT, space_after=4)
    para(tf2, "  Block 1 (Covariate): PANSS Negative Subscale score — controls for overlap between negative symptoms (social withdrawal, blunted affect) and loneliness. Reports R-squared for symptoms alone.", size=12, color=DARK_TEXT, space_after=4)
    para(tf2, "  Block 2 (Processing): Zd (efficiency), DQ+ (complexity), Lambda (simplification). Reports DELTA R-squared — how much ADDITIONAL variance in loneliness does Processing explain beyond symptoms?", size=12, color=DARK_TEXT, space_after=4)
    para(tf2, "  Block 3 (Mediation): XA% (adequacy), X-% (distortion), P (conventionality). Reports DELTA R-squared — does reality testing add further prediction?", size=12, color=DARK_TEXT, space_after=4)
    para(tf2, "  Block 4 (Ideation): WSum6 (thought disorder), M quality composite. Reports DELTA R-squared — does thinking quality add still more?", size=12, color=DARK_TEXT, space_after=8)
    para(tf2, "KEY OUTPUT: Which block produces the LARGEST significant Delta-R-squared? That cluster is the strongest cognitive predictor of loneliness. Software: SPSS 26 / JASP. Alpha: p < .05 (two-tailed).", size=13, bold=True, color=ACCENT, space_after=4)

slide_analysis()

# ==================== SLIDE 16: CLINICAL SIGNIFICANCE ====================
def slide_significance():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, PRIMARY)
    bar(slide, 0, 0, prs.slide_width, Inches(0.08), GOLD)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Clinical Significance & Expected Contribution", size=26, bold=True, color=GOLD)

    tf2 = txt(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.2))
    para(tf2, "1. FIRST STUDY GLOBALLY linking Rorschach cognitive triad to loneliness in any population.", size=14, bold=True, color=SECONDARY, space_after=4)
    para(tf2, "   This is not a replication or extension — it opens an entirely new line of inquiry at the intersection of projective assessment and loneliness research in psychosis.", size=13, color=WHITE, space_after=12)

    para(tf2, "2. IDENTIFIES COGNITIVE TARGETS for loneliness intervention.", size=14, bold=True, color=GREEN, space_after=4)
    para(tf2, "   If Mediation (perceptual distortion) emerges as the strongest predictor, treatment should prioritize social perception training and reality testing. If Processing (scanning) dominates, attention training is indicated. If Ideation (thought disorder), metacognitive therapy is the target. This enables INDIVIDUALIZED rehabilitation.", size=13, color=WHITE, space_after=12)

    para(tf2, "3. PERFORMANCE-BASED assessment that CANNOT be faked.", size=14, bold=True, color=PURPLE, space_after=4)
    para(tf2, "   Unlike self-report cognition measures (unreliable in 50-80% of patients who lack insight), the Rorschach captures implicit cognitive processes the patient cannot deliberately manipulate. Combined with self-reported loneliness, this multi-method approach is the gold standard.", size=13, color=WHITE, space_after=12)

    para(tf2, "4. PRACTICAL for Indian clinical settings.", size=14, bold=True, color=GOLD, space_after=4)
    para(tf2, "   A single 30-minute Rorschach administration replaces a multi-hour neuropsych battery. Available in EVERY clinical psychology department. No computerized equipment needed. If validated as a loneliness predictor, it becomes a feasible screening tool in Indian DMHP settings.", size=13, color=WHITE, space_after=12)

    para(tf2, "5. ADDRESSES a clinical priority endorsed by WHO (2023) and Indian clinical guidelines (IPS, 2024).", size=14, bold=True, color=ACCENT, space_after=4)
    para(tf2, "   Loneliness assessment and intervention are now recommended as part of comprehensive schizophrenia care. This study provides the assessment TOOL and the cognitive MODEL to implement that recommendation.", size=13, color=WHITE, space_after=4)

slide_significance()



# ==================== SLIDE 17: REFERENCES ====================
def slide_references():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    bar(slide, 0, 0, Inches(0.12), prs.slide_height, GREY)

    tf = txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    para(tf, "Key References", size=26, bold=True, color=PRIMARY)
    bar(slide, Inches(0.6), Inches(0.9), Inches(3), Pt(4), GREY)

    tf2 = txt(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.2))
    refs = [
        "Cacioppo, J.T. & Hawkley, L.C. (2009). Perceived social isolation and cognition. Trends in Cognitive Sciences, 13(10), 447-454.",
        "Exner, J.E. (2003). The Rorschach: A Comprehensive System (4th ed.). Wiley.",
        "Fortuna, K.L. et al. (2020). Loneliness and physical health in SMI. Social Psychiatry & Psychiatric Epidemiology, 55, 1421-1428.",
        "Green, M.F. et al. (2023). Social withdrawal, loneliness, and health in schizophrenia. Schizophrenia Bulletin, 49(5), 1112-1126.",
        "Hajduk, M. et al. (2023). Social cognitive bias increases loneliness in schizophrenia. Schizophrenia Research, 256, 38-46.",
        "Hsu, Y.C. et al. (2024). Social isolation and cognitive function in schizophrenia. Schizophrenia Research, 266, 171-178.",
        "Ilonen, T. et al. (2012). Rorschach measures of cognition relate to functioning in schizophrenia. Psychological Assessment, 25(1), 253-263.",
        "Jaya, E.S. et al. (2024). Selective attention to threat mediates loneliness-psychosis link. Schizophrenia Research, 269, 118-125.",
        "Kay, S.R. et al. (1987). The PANSS for schizophrenia. Schizophrenia Bulletin, 13(2), 261-276.",
        "Lin, C.Y. et al. (2022). UCLA Loneliness Scale validation in schizophrenia. Int J Environ Res Public Health, 19(14), 8443.",
        "Michalska da Rocha, B. et al. (2018). Loneliness in psychosis: Meta-analytical review. Schizophrenia Bulletin, 44(1), 114-125.",
        "Russell, D.W. (1996). UCLA Loneliness Scale (Version 3). J Personality Assessment, 66(1), 20-40.",
        "Singh, G. et al. (2021). Rorschach and psychopathology in schizophrenia. Indian J Psychiatric Nursing, 30(1), 49-55.",
        "Thibaudeau, E. et al. (2023). Loneliness linked to mentalizing deficits in schizophrenia. Schizophrenia Research, 256, 29-37.",
        "Wang, X. et al. (2026). Loneliness and cognitive impairment in schizophrenia. European Archives of Psychiatry & Clinical Neuroscience.",
        "Yen, C.F. et al. (2023). Loneliness, suicide risk and depression in schizophrenia. npj Schizophrenia, 9, Article 40.",
    ]
    for ref in refs:
        para(tf2, ref, size=11, color=DARK_TEXT, space_after=5)

slide_references()

# ==================== SLIDE 18: THANK YOU ====================
def slide_thankyou():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, PRIMARY)
    bar(slide, 0, 0, prs.slide_width, Inches(0.12), GOLD)
    bar(slide, 0, Inches(7.38), prs.slide_width, Inches(0.12), GOLD)
    bar(slide, Inches(0.8), Inches(2.0), Inches(0.08), Inches(3.5), SECONDARY)

    tf = txt(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(4))
    para(tf, "Thank You", size=48, bold=True, color=GOLD, align=PP_ALIGN.LEFT, space_after=20)
    para(tf, "\"The deepest suffering of a patient with schizophrenia is not", size=16, color=WHITE, align=PP_ALIGN.LEFT, space_after=2)
    para(tf, "the voice they hear — it is the silence around them.\"", size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT, space_after=30)
    para(tf, "Questions, Suggestions & Discussion Welcome", size=18, color=SECONDARY, align=PP_ALIGN.LEFT, space_after=30)
    para(tf, "[Candidate Name]", size=14, color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.LEFT, space_after=4)
    para(tf, "MPhil Clinical Psychology | [Institution] | 2026-2027", size=13, color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.LEFT)

slide_thankyou()

# ==================== SAVE ====================
output_path = "/projects/sandbox/Dango-kiro/Research_Proposal_Rorschach_Loneliness_PPT.pptx"
prs.save(output_path)
print(f"SUCCESS: Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
