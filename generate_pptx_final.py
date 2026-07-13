#!/usr/bin/env python3
"""
MBRP Research Synopsis PPT - FINAL AESTHETIC VERSION
Professional design with:
- Gradient-style headers
- Charts and graphs
- Visual flow diagrams
- Color-coded boxes
- Accent shapes and decorative elements
- Professional academic formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === PREMIUM COLOR PALETTE ===
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
DEEP_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x7B, 0x83)
LIGHT_TEAL = RGBColor(0x00, 0xA8, 0xA8)
GOLD = RGBColor(0xC9, 0x8A, 0x1A)
SOFT_GOLD = RGBColor(0xE8, 0xB8, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF7, 0xFA)
LIGHT_GRAY = RGBColor(0xEC, 0xEF, 0xF1)
MID_GRAY = RGBColor(0x90, 0xA4, 0xAE)
DARK_TEXT = RGBColor(0x1A, 0x23, 0x2F)
BODY_TEXT = RGBColor(0x2C, 0x3E, 0x50)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
SOFT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
SOFT_TEAL_BG = RGBColor(0xE8, 0xF8, 0xF5)



# === HELPER FUNCTIONS ===

def add_decorative_bar(slide, top=Inches(7.2), color=GOLD):
    """Add a thin gold accent bar at the bottom"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(13.333), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def add_side_accent(slide, color=TEAL):
    """Add left side accent stripe"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def add_slide_number(slide, num, total=42):
    """Add slide number at bottom right"""
    txBox = slide.shapes.add_textbox(Inches(11.8), Inches(7.1), Inches(1.3), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.RIGHT

def create_header(slide, title, slide_num):
    """Create a professional header with gradient-style bar"""
    # Main header background
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DEEP_BLUE
    hdr.line.fill.background()
    # Gold accent line below header
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12.0), Inches(0.85))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    add_slide_number(slide, slide_num)
    add_decorative_bar(slide)

def add_content_slide(prs, title, bullets, slide_num, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = OFF_WHITE
    add_side_accent(slide)
    create_header(slide, title, slide_num)
    top = 1.4
    if subtitle:
        txS = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.45))
        tfS = txS.text_frame
        pS = tfS.paragraphs[0]
        pS.text = subtitle
        pS.font.size = Pt(13)
        pS.font.bold = True
        pS.font.color.rgb = TEAL
        top = 1.85
    txB = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.8), Inches(7.2 - top - 0.4))
    tfB = txB.text_frame
    tfB.word_wrap = True
    for i, b in enumerate(bullets):
        pp = tfB.paragraphs[0] if i == 0 else tfB.add_paragraph()
        pp.text = b
        pp.font.size = Pt(13)
        pp.font.color.rgb = BODY_TEXT
        pp.space_before = Pt(5)
    return slide

def add_paragraph_slide(prs, title, paragraphs, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = OFF_WHITE
    add_side_accent(slide, ACCENT_PURPLE)
    create_header(slide, title, slide_num)
    txB = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12.0), Inches(5.7))
    tfB = txB.text_frame
    tfB.word_wrap = True
    for i, para in enumerate(paragraphs):
        pp = tfB.paragraphs[0] if i == 0 else tfB.add_paragraph()
        pp.text = para
        pp.font.size = Pt(11)
        pp.font.color.rgb = BODY_TEXT
        pp.space_before = Pt(8)
        pp.space_after = Pt(3)
    return slide



def add_box(slide, left, top, width, height, fill_color, text, font_size=11, text_color=BODY_TEXT, bold_first=False):
    """Add a colored rounded box with text"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        if i == 0 and bold_first:
            p.font.bold = True
            p.font.size = Pt(font_size + 1)

def add_arrow(slide, left, top, width=Inches(0.8), height=Inches(0.3)):
    """Add a right-pointing arrow"""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEAL
    arrow.line.fill.background()

def add_down_arrow(slide, left, top, width=Inches(0.35), height=Inches(0.5)):
    """Add a downward arrow"""
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()


# ===================================================================
# SLIDE 1: TITLE SLIDE (Premium Design)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = NAVY

# Top accent bar
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = GOLD
top_bar.line.fill.background()

# Decorative circle (top-right)
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.5), Inches(2.5), Inches(2.5))
circle.fill.solid()
circle.fill.fore_color.rgb = DEEP_BLUE
circle.line.color.rgb = TEAL
circle.line.width = Pt(2)

# Title
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(9.5), Inches(2.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Efficacy of Brief Mindfulness-Based"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p2 = tf.add_paragraph()
p2.text = "Relapse Prevention (MBRP) Intervention"
p2.font.size = Pt(28)
p2.font.bold = True
p2.font.color.rgb = WHITE
p3 = tf.add_paragraph()
p3.text = "on Craving, Impulsivity, and Mindfulness"
p3.font.size = Pt(24)
p3.font.color.rgb = LIGHT_TEAL
p3.space_before = Pt(8)
p4 = tf.add_paragraph()
p4.text = "in Substance Dependent Patients"
p4.font.size = Pt(24)
p4.font.color.rgb = LIGHT_TEAL

# Subtitle box
sub_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.6), Inches(5.5), Inches(0.55))
sub_box.fill.solid()
sub_box.fill.fore_color.rgb = TEAL
sub_box.line.fill.background()
stf = sub_box.text_frame
stf.vertical_anchor = MSO_ANCHOR.MIDDLE
sp = stf.paragraphs[0]
sp.text = "  MPhil Clinical Psychology Research Synopsis"
sp.font.size = Pt(14)
sp.font.bold = True
sp.font.color.rgb = WHITE

# Details
dtx = slide.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(8), Inches(1.8))
dtf = dtx.text_frame
dp = dtf.paragraphs[0]
dp.text = "Researcher: [Name]  |  Guide: [Supervisor Name]"
dp.font.size = Pt(12)
dp.font.color.rgb = MID_GRAY
dp2 = dtf.add_paragraph()
dp2.text = "Department of Clinical Psychology"
dp2.font.size = Pt(12)
dp2.font.color.rgb = MID_GRAY
dp3 = dtf.add_paragraph()
dp3.text = "Man Nasha Mukti Kendra  |  [University Name]  |  2026"
dp3.font.size = Pt(12)
dp3.font.color.rgb = SOFT_GOLD

# Bottom bar
add_decorative_bar(slide, Inches(7.35), GOLD)



# ===================================================================
# SLIDE 2: INTRODUCTION 1/3 (with abbreviation box)
# ===================================================================
slide = add_content_slide(prs, "Introduction & Background (1/3)", [
    "India faces a significant substance use crisis; MAGNITUDE study (2019) estimated ~3.1 crore individuals affected by substance use disorders",
    "Opioid dependence constitutes a major public health burden, particularly in Punjab, Rajasthan, Northeast India, and metropolitan areas",
    "WHO estimates India accounts for ~25% of global opioid-related deaths in South-East Asia",
    "Indian de-addiction centers, including Man Nasha Mukti Kendra, primarily offer pharmacotherapy (OST, naltrexone) with limited structured psychotherapy",
    "Psychosocial interventions remain under-utilized despite evidence of superior combined treatment outcomes"
], 2, subtitle="Substance Dependence: Indian Context")

# Abbreviation box
add_box(slide, Inches(0.8), Inches(5.5), Inches(12.0), Inches(1.4), SOFT_BLUE,
    "KEY ABBREVIATIONS: MBRP = Mindfulness-Based Relapse Prevention | TAU = Treatment As Usual | OCDUS = Obsessive Compulsive Drug Use Scale\nBIS-11 = Barratt Impulsiveness Scale | FFMQ = Five Facet Mindfulness Questionnaire | ASSIST = WHO Substance Screening Test\nOST = Opioid Substitution Therapy | PFC = Prefrontal Cortex | ANCOVA = Analysis of Covariance",
    10, DEEP_BLUE)

# ===================================================================
# SLIDE 3: INTRODUCTION 2/3
# ===================================================================
add_content_slide(prs, "Introduction & Background (2/3)", [
    "Relapse rates in substance dependence range from 40-60% within first year post-treatment (NIDA, 2020)",
    "Indian studies report even higher relapse rates (~70-80%) in opioid dependence (Mattoo et al., 2009)",
    "Triggers for relapse: craving, negative affect, interpersonal conflict, environmental cues",
    "Traditional Relapse Prevention (Marlatt & Gordon, 1985) has moderate efficacy but limited mindfulness integration",
    "Automatic cognitive-behavioral patterns perpetuate relapse cycles",
    "Brief interventions essential for Indian rehab settings (Man Nasha Mukti Kendra) with limited resources"
], 3, subtitle="The Relapse Problem in De-Addiction")

# ===================================================================
# SLIDE 4: INTRODUCTION 3/3
# ===================================================================
add_content_slide(prs, "Introduction & Background (3/3)", [
    "Mindfulness: intentional, non-judgmental, present-moment awareness (Kabat-Zinn, 1990)",
    "MBRP developed by Bowen, Chawla, and Marlatt (2011): integrates MBSR + Cognitive-Behavioral RP",
    "Core mechanisms: decentering from craving, disrupting automatic reactivity, increasing distress tolerance",
    "MBRP targets 'craving -> use' automaticity by cultivating awareness of triggers",
    "Brief adaptations (6 sessions / 3 weeks) show promise for resource-limited rehabilitation settings",
    "Growing global evidence supports MBRP; Indian validation remains critically scarce",
    "Rationale: Brief MBRP can be feasibly implemented at Man Nasha Mukti Kendra"
], 4, subtitle="Mindfulness in Addiction Treatment")



# ===================================================================
# SLIDE 5: VARIABLE 1 - CRAVING (Visual card layout)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, ACCENT_RED)
create_header(slide, "Variable 1: Craving", 5)

# Definition box
add_box(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(1.5), RGBColor(0xFD, 0xED, 0xEC),
    "DEFINITION\nClinical: Intense subjective urge to use a substance triggered by cues\nTheoretical: Motivational state from incentive-sensitization\n(Robinson & Berridge, 1993)",
    11, DARK_TEXT, True)

# Relevance box
add_box(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(1.5), RGBColor(0xFD, 0xED, 0xEC),
    "RELEVANCE IN SUBSTANCE DEPENDENCE\nPrimary predictor of relapse in opioid dependence\nMediates cue exposure -> substance use behavior\nCorrelates with severity and treatment dropout",
    11, DARK_TEXT, True)

# Neuropsychological box
add_box(slide, Inches(0.6), Inches(3.1), Inches(6.0), Inches(1.5), RGBColor(0xFA, 0xF0, 0xDB),
    "NEUROPSYCHOLOGICAL BASIS\nMesolimbic dopamine: VTA -> Nucleus Accumbens\nPFC hypoactivation during craving episodes\nConditioned cue-reward associations",
    11, DARK_TEXT, True)

# Techniques box
add_box(slide, Inches(6.8), Inches(3.1), Inches(6.0), Inches(2.2), RGBColor(0xE8, 0xF8, 0xF5),
    "MBRP TECHNIQUES TARGETING CRAVING\n- Urge Surfing: Observe craving as transient wave\n- Mindfulness Exposure: Non-reactive awareness\n- SOBER Space: Stop-Observe-Breathe-Expand-Respond\n- Cognitive Decentering: 'I am having a craving'",
    11, DARK_TEXT, True)

# Visual icon - wave symbol
wave = slide.shapes.add_shape(MSO_SHAPE.WAVE, Inches(0.6), Inches(5.0), Inches(2.5), Inches(1.2))
wave.fill.solid()
wave.fill.fore_color.rgb = LIGHT_TEAL
wave.line.fill.background()
wtf = wave.text_frame
wtf.vertical_anchor = MSO_ANCHOR.MIDDLE
wp = wtf.paragraphs[0]
wp.text = "URGE\nSURFING"
wp.font.size = Pt(12)
wp.font.bold = True
wp.font.color.rgb = WHITE
wp.alignment = PP_ALIGN.CENTER

# ===================================================================
# SLIDE 6: VARIABLE 2 - IMPULSIVITY (Visual card layout)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, ACCENT_ORANGE)
create_header(slide, "Variable 2: Impulsivity", 6)

add_box(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(1.5), RGBColor(0xFE, 0xF9, 0xE7),
    "DEFINITION\nClinical: Rapid, unplanned actions without considering\nconsequences (Moeller et al., 2001)\nTheoretical: Motor + Attentional + Non-Planning (Patton et al., 1995)",
    11, DARK_TEXT, True)

add_box(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(1.5), RGBColor(0xFE, 0xF9, 0xE7),
    "RELEVANCE IN SUBSTANCE DEPENDENCE\nHigher trait impulsivity predicts relapse\nMediates craving-to-use behavior\nAssociated with treatment non-adherence & dropout",
    11, DARK_TEXT, True)

add_box(slide, Inches(0.6), Inches(3.1), Inches(6.0), Inches(1.5), RGBColor(0xFD, 0xED, 0xEC),
    "NEUROPSYCHOLOGICAL BASIS\nPFC dysfunction -> impaired executive control\nReduced inhibitory control (Go/No-Go paradigms)\nImpaired delay discounting; DLPFC hypoactivation",
    11, DARK_TEXT, True)

add_box(slide, Inches(6.8), Inches(3.1), Inches(6.0), Inches(2.2), RGBColor(0xE8, 0xF8, 0xF5),
    "MBRP TECHNIQUES TARGETING IMPULSIVITY\n- Response Inhibition: Mindful pause before action\n- Awareness Training: Notice impulse-action sequences\n- Mindful Decision-Making: Space between stimulus & response\n- STOP: Stop-Take a breath-Observe-Proceed",
    11, DARK_TEXT, True)

# Stop sign shape
stop = slide.shapes.add_shape(MSO_SHAPE.OCTAGON, Inches(0.8), Inches(5.0), Inches(1.8), Inches(1.8))
stop.fill.solid()
stop.fill.fore_color.rgb = ACCENT_RED
stop.line.fill.background()
stf = stop.text_frame
stf.vertical_anchor = MSO_ANCHOR.MIDDLE
sp = stf.paragraphs[0]
sp.text = "STOP\nTechnique"
sp.font.size = Pt(11)
sp.font.bold = True
sp.font.color.rgb = WHITE
sp.alignment = PP_ALIGN.CENTER



# ===================================================================
# SLIDE 7: VARIABLE 3 - MINDFULNESS (Visual card layout)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, ACCENT_GREEN)
create_header(slide, "Variable 3: Mindfulness", 7)

add_box(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(1.5), RGBColor(0xE8, 0xF8, 0xF5),
    "DEFINITION\nClinical: Present-moment awareness with openness & non-judgment\nTheoretical: 5 facets - Observing, Describing, Acting with\nAwareness, Non-Judging, Non-Reactivity (Baer et al., 2006)",
    11, DARK_TEXT, True)

add_box(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(1.5), RGBColor(0xE8, 0xF8, 0xF5),
    "RELEVANCE IN SUBSTANCE DEPENDENCE\nSubstance users show significantly lower mindfulness\nActs as protective factor against relapse triggers\nImprovements mediate MBRP treatment outcomes",
    11, DARK_TEXT, True)

add_box(slide, Inches(0.6), Inches(3.1), Inches(6.0), Inches(1.5), RGBColor(0xD6, 0xEA, 0xF8),
    "NEUROPSYCHOLOGICAL BASIS\nACC activation -> enhanced self-regulation\nInsula -> interoceptive awareness\nPFC-amygdala connectivity -> emotional regulation\nDMN regulation -> reduced rumination",
    11, DARK_TEXT, True)

add_box(slide, Inches(6.8), Inches(3.1), Inches(6.0), Inches(2.2), RGBColor(0xFA, 0xF0, 0xDB),
    "MBRP TECHNIQUES FOR MINDFULNESS\n- Sitting Meditation: Focused attention on breath\n- Body Scan: Non-judgmental bodily awareness\n- Mindful Movement: Present-moment yoga/walking\n- Non-Judgmental Awareness: Labeling without evaluation",
    11, DARK_TEXT, True)

# Lotus/meditation symbol
lotus = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(5.2), Inches(2.0), Inches(1.5))
lotus.fill.solid()
lotus.fill.fore_color.rgb = ACCENT_GREEN
lotus.line.fill.background()
ltf = lotus.text_frame
ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
lp = ltf.paragraphs[0]
lp.text = "MINDFUL\nAWARENESS"
lp.font.size = Pt(11)
lp.font.bold = True
lp.font.color.rgb = WHITE
lp.alignment = PP_ALIGN.CENTER

# ===================================================================
# SLIDES 8-12: REVIEW OF LITERATURE (5 paragraph slides)
# ===================================================================
add_paragraph_slide(prs, "Review of Literature (1/5) - MBRP & Relapse Prevention", [
    "Bowen et al. (2014) investigate the relative efficacy of Mindfulness-Based Relapse Prevention compared to standard Relapse Prevention and Treatment As Usual for substance use disorders in JAMA Psychiatry. The randomized clinical trial involved 286 participants who had completed initial treatment for substance use disorders. The study demonstrates that at 12-month follow-up, MBRP participants reported significantly fewer days of substance use and heavy drinking compared to both standard RP and TAU groups. The researchers highlight that while all three groups showed initial improvements, MBRP participants maintained superior long-term outcomes, suggesting that the cultivation of mindfulness skills provides a durable protective mechanism against relapse. The study establishes that MBRP's integration of present-moment awareness with traditional cognitive-behavioral relapse prevention strategies creates a synergistic therapeutic effect that addresses both the automatic reactivity underlying craving and the cognitive distortions that precipitate relapse.",
    "",
    "Bowen and Marlatt (2009) examine the effects of brief mindfulness-based intervention on craving among substance users in Psychology of Addictive Behaviors. The study recruited incarcerated individuals with substance use histories and administered a brief urge surfing meditation intervention. Results demonstrated significant reductions in craving intensity and frequency in the mindfulness condition compared to controls. The authors propose that even brief exposure to mindfulness techniques can disrupt the automaticity of craving responses by introducing a meta-cognitive awareness layer between trigger and behavioral response. This study is particularly relevant to the present research as it validates the premise that abbreviated mindfulness interventions can produce meaningful clinical effects on craving, supporting the feasibility of brief MBRP protocols (6 sessions over 3 weeks) in settings like Man Nasha Mukti Kendra where extended programs are impractical."
], 8)



add_paragraph_slide(prs, "Review of Literature (2/5) - Craving & Mindfulness Mechanisms", [
    "Garland et al. (2014) explore how mindfulness training targets neurocognitive mechanisms of addiction at the attention-appraisal-emotion interface in Frontiers in Psychiatry. The study presents a theoretical and empirical framework demonstrating that Mindfulness-Oriented Recovery Enhancement (MORE) reduces opioid craving through three interconnected mechanisms: attentional reorientation away from drug-related cues, positive reappraisal of previously neutral stimuli to generate natural reward, and enhanced savoring of healthy pleasures. The researchers provide neuroimaging evidence showing that mindfulness practice modulates activity in prefrontal and limbic circuits associated with craving and emotional regulation. The study is significant because it elucidates the precise cognitive-neural pathways through which mindfulness reduces craving in opioid users, providing a mechanistic rationale for MBRP techniques such as urge surfing and mindful awareness of craving sensations.",
    "",
    "Witkiewitz et al. (2013) investigate mindfulness-based relapse prevention effects on substance craving in Addictive Behaviors. The study conducted secondary analyses from a randomized controlled trial comparing MBRP to TAU among individuals in aftercare. Over a 4-month follow-up, MBRP participants demonstrated significantly lower craving levels, and importantly, the relationship between negative affect and subsequent craving was significantly attenuated in the MBRP group. The authors conclude that mindfulness practice weakens the affect-craving pathway by cultivating non-reactive awareness of emotional states, thereby preventing negative emotions from automatically triggering craving responses. This decoupling of affect and craving represents a critical therapeutic mechanism that distinguishes MBRP from traditional relapse prevention approaches."
], 9)

add_paragraph_slide(prs, "Review of Literature (3/5) - Impulsivity & Mindfulness", [
    "Garland et al. (2016) examine the efficacy of Mindfulness-Oriented Recovery Enhancement versus CBT for co-occurring substance dependence, traumatic stress, and psychiatric disorders in the Journal of Consulting and Clinical Psychology. The study employed a randomized controlled design with substance-dependent adults exhibiting elevated impulsivity. Results demonstrated that the mindfulness-based intervention produced significant reductions in impulsivity scores as measured by the Barratt Impulsiveness Scale (BIS-11), particularly in the motor and attentional impulsivity subscales. The authors theorize that mindfulness meditation strengthens prefrontal cortical inhibitory mechanisms by repeatedly engaging participants in exercises requiring sustained attention, response monitoring, and deliberate non-reactivity. This enhanced top-down cognitive control translates into improved ability to inhibit prepotent impulsive responses, particularly in high-risk situations where automatic substance-seeking behavior would otherwise occur.",
    "",
    "Murphy and MacKillop (2012) explore the interrelationships between impulsivity, mindfulness, and alcohol misuse in Psychopharmacology. Using a cross-sectional design with 340 participants, the researchers found that trait mindfulness was inversely associated with impulsive decision-making as measured by delay discounting tasks. Critically, mindfulness moderated the relationship between impulsivity and substance use problems, such that individuals with higher mindfulness showed weaker associations between impulsivity and alcohol-related consequences. The authors propose that mindfulness functions as a cognitive resource enabling impulsive individuals to override automatic behavioral tendencies through enhanced metacognitive awareness and response flexibility."
], 10)

add_paragraph_slide(prs, "Review of Literature (4/5) - Meta-Analyses & Brief Models", [
    "Li et al. (2017) conduct a systematic review and meta-analysis of mindfulness treatment for substance misuse in the Journal of Substance Abuse Treatment, encompassing 42 randomized controlled trials. The meta-analytic findings reveal that mindfulness-based interventions produce significant effect sizes for reducing substance misuse (d = 0.33), craving (d = 0.68), and stress (d = 0.44) across diverse populations. The authors highlight that effect sizes for craving reduction are particularly robust, supporting the theoretical premise that mindfulness directly targets craving mechanisms through enhanced interoceptive awareness. Furthermore, the review identifies that interventions of shorter duration (4-8 sessions) demonstrated comparable efficacy to longer protocols when appropriately structured, providing empirical justification for the present study's brief 6-session MBRP adaptation.",
    "",
    "Glasner-Edwards et al. (2017) examine a pilot randomized clinical trial of mindfulness-based relapse prevention for stimulant-dependent adults using an abbreviated 6-session protocol in the journal Mindfulness. The study demonstrates that a condensed mindfulness intervention is both feasible and effective in outpatient settings. Participants showed significant reductions in substance use frequency and craving intensity compared to the health education control. This study directly validates the brief intervention model proposed in the present research, demonstrating that 6-session MBRP protocols can be successfully implemented without substantial loss of efficacy, making them particularly appropriate for Indian rehabilitation settings like Man Nasha Mukti Kendra with typical admission durations of 4-6 weeks."
], 11)

add_paragraph_slide(prs, "Review of Literature (5/5) - Indian Context & Research Gap", [
    "Ghosh et al. (2018) conduct a comprehensive review of relapse in opioid dependence from an Indian perspective in the Indian Journal of Psychiatry. The study reports alarmingly high relapse rates exceeding 70% among opioid-dependent patients treated in North Indian de-addiction centers, with the majority occurring within three months post-discharge. The authors identify craving, peer influence, negative emotional states, and lack of structured psychological aftercare as primary relapse determinants. Significantly, the review highlights that Indian treatment facilities predominantly rely on pharmacological approaches with minimal integration of evidence-based psychological interventions. The authors advocate for development and validation of structured psychosocial protocols tailored to Indian treatment infrastructure and resource constraints. This directly establishes the clinical need for the present research.",
    "",
    "Sarkar and Balhara (2016) highlight the underutilization of structured psychological interventions in Indian de-addiction settings, identifying barriers including limited trained personnel, absence of culturally validated protocols, and institutional emphasis on pharmacological management. Jain et al. (2013) provide preliminary evidence from a mindfulness-based intervention study with alcohol-dependent patients in India showing initial craving reductions. CRITICAL GAP: No published Indian RCT has tested Brief MBRP specifically in opioid-dependent populations, simultaneously assessing craving, impulsivity, and mindfulness as treatment outcomes. The present study addresses this gap as the first brief MBRP trial designed for Indian de-addiction center infrastructure like Man Nasha Mukti Kendra."
], 12)



# ===================================================================
# SLIDE 13: RESEARCH GAP (with visual gap diagram)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, ACCENT_RED)
create_header(slide, "Research Gap", 13)

# Gap items as colored boxes
gaps = [
    ("No Indian RCT has tested MBRP (standard or brief) in opioid-dependent populations", ACCENT_RED),
    ("Most studies examine craving OR mindfulness alone; few assess all three variables simultaneously", ACCENT_ORANGE),
    ("No brief MBRP (6 sessions) validated for Indian rehab settings like Man Nasha Mukti Kendra", ACCENT_ORANGE),
    ("Indian treatment lacks integration of evidence-based psychological interventions alongside OST", GOLD),
    ("Global MBRP studies focus on alcohol/polysubstance; opioid-specific evidence is limited", GOLD),
]
for i, (text, color) in enumerate(gaps):
    y = Inches(1.5 + i * 0.85)
    # Colored left marker
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(0.12), Inches(0.6))
    marker.fill.solid()
    marker.fill.fore_color.rgb = color
    marker.line.fill.background()
    # Text
    txB = slide.shapes.add_textbox(Inches(1.1), y, Inches(11.5), Inches(0.7))
    tf = txB.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = BODY_TEXT

# Present study box
add_box(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.0), RGBColor(0xE8, 0xF8, 0xF5),
    "PRESENT STUDY ADDRESSES: First brief MBRP trial (6 sessions/3 weeks) in Indian opioid-dependent sample at Man Nasha Mukti Kendra\nassessing craving, impulsivity, and mindfulness simultaneously",
    12, TEAL, True)

# ===================================================================
# SLIDE 14: AIM
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, TEAL)
create_header(slide, "Aim of the Study", 14)

# Large aim box centered
add_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5), SOFT_TEAL_BG,
    "To evaluate the efficacy of a Brief Mindfulness-Based Relapse Prevention\n(MBRP) intervention (6 sessions over 3 weeks) in reducing craving and\nimpulsivity, and enhancing mindfulness, among male substance-dependent\npatients at Man Nasha Mukti Kendra\n\nComparing: Brief MBRP + TAU (Experimental) vs. Psychoeducation + TAU (Control)",
    14, DEEP_BLUE, True)

# ===================================================================
# SLIDE 15: OBJECTIVES
# ===================================================================
add_content_slide(prs, "Objectives", [
    "1. To assess and compare craving levels (pre vs. post) in Experimental (Brief MBRP + TAU) and Control (Psychoeducation + TAU) groups",
    "2. To assess and compare impulsivity levels (pre vs. post) in both groups",
    "3. To assess and compare mindfulness levels (pre vs. post) in both groups",
    "4. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in reducing craving",
    "5. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in reducing impulsivity",
    "6. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in enhancing mindfulness"
], 15)



# ===================================================================
# SLIDE 16: HYPOTHESES (with visual layout)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, TEAL)
create_header(slide, "Hypotheses (Research & Null)", 16)

# Research hypotheses box
add_box(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.5), SOFT_TEAL_BG,
    "RESEARCH HYPOTHESES (Directional)\nH1: Brief MBRP + TAU will show significantly GREATER REDUCTION in craving (OCDUS) than Psychoeducation + TAU\nH2: Brief MBRP + TAU will show significantly GREATER REDUCTION in impulsivity (BIS-11) than Psychoeducation + TAU\nH3: Brief MBRP + TAU will show significantly GREATER INCREASE in mindfulness (FFMQ) than Psychoeducation + TAU",
    12, DEEP_BLUE, True)

# Null hypotheses box
add_box(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.5), RGBColor(0xFD, 0xED, 0xEC),
    "NULL HYPOTHESES\nH01: No significant difference in craving (OCDUS) between Experimental and Control groups pre-to-post\nH02: No significant difference in impulsivity (BIS-11) between Experimental and Control groups pre-to-post\nH03: No significant difference in mindfulness (FFMQ) between Experimental and Control groups pre-to-post",
    12, DARK_TEXT, True)

# ===================================================================
# SLIDE 17: RESEARCH DESIGN (with visual diagram)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, TEAL)
create_header(slide, "Research Design: Pre-test Post-test Control Group", 17)

# Design diagram using boxes and arrows
# Randomization box
add_box(slide, Inches(0.5), Inches(2.0), Inches(2.0), Inches(1.0), GOLD,
    "RANDOMIZATION\n(N = 60)", 12, DARK_TEXT, True)

# Arrow right
add_arrow(slide, Inches(2.6), Inches(2.2))

# Pre-test box
add_box(slide, Inches(3.5), Inches(1.7), Inches(2.2), Inches(1.5), SOFT_BLUE,
    "PRE-TEST (O1)\nOCDUS\nBIS-11\nFFMQ\nASSIST", 10, DEEP_BLUE, True)

# Arrow right
add_arrow(slide, Inches(5.8), Inches(2.2))

# Intervention boxes (split)
add_box(slide, Inches(6.7), Inches(1.5), Inches(3.0), Inches(1.0), RGBColor(0xE8, 0xF8, 0xF5),
    "X1: Brief MBRP + TAU\n6 sessions / 3 weeks", 10, TEAL, True)
add_box(slide, Inches(6.7), Inches(2.7), Inches(3.0), Inches(1.0), RGBColor(0xFE, 0xF9, 0xE7),
    "X2: Psychoeducation + TAU\n6 sessions / 3 weeks", 10, ACCENT_ORANGE, True)

# Arrow right
add_arrow(slide, Inches(9.8), Inches(2.2))

# Post-test box
add_box(slide, Inches(10.7), Inches(1.7), Inches(2.2), Inches(1.5), SOFT_BLUE,
    "POST-TEST (O2)\nOCDUS\nBIS-11\nFFMQ", 10, DEEP_BLUE, True)

# Notation below
txN = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(12.0), Inches(2.0))
tfN = txN.text_frame
tfN.word_wrap = True
pN = tfN.paragraphs[0]
pN.text = "R   O1   X1   O2   -->  Experimental Group (Brief MBRP + TAU)"
pN.font.size = Pt(13)
pN.font.color.rgb = TEAL
pN.font.bold = True
pN2 = tfN.add_paragraph()
pN2.text = "R   O1   X2   O2   -->  Control Group (Psychoeducation + TAU)"
pN2.font.size = Pt(13)
pN2.font.color.rgb = ACCENT_ORANGE
pN2.font.bold = True
pN3 = tfN.add_paragraph()
pN3.text = ""
pN4 = tfN.add_paragraph()
pN4.text = "Features: True experimental design | Active control (attention-matched) | TAU continued for all | Setting: Man Nasha Mukti Kendra"
pN4.font.size = Pt(11)
pN4.font.color.rgb = MID_GRAY



# ===================================================================
# SLIDE 18: SAMPLE (with male justification visual)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, TEAL)
create_header(slide, "Sample & Sampling Strategy", 18)

# Two-stage process boxes
add_box(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(1.8), SOFT_BLUE,
    "STAGE 1: PURPOSIVE SELECTION\nEligible participants identified based on\ninclusion/exclusion criteria from patients\nadmitted to Man Nasha Mukti Kendra",
    11, DEEP_BLUE, True)

add_arrow(slide, Inches(6.2), Inches(2.0))

add_box(slide, Inches(7.1), Inches(1.4), Inches(5.5), Inches(1.8), RGBColor(0xE8, 0xF8, 0xF5),
    "STAGE 2: RANDOM ASSIGNMENT\nComputer-generated randomization\nallocating eligible participants to\nExperimental (n=30) or Control (n=30)",
    11, TEAL, True)

# Sample details
add_box(slide, Inches(0.6), Inches(3.6), Inches(6.0), Inches(1.2), LIGHT_GRAY,
    "N = 60 (30 per group) | Recruit 70 (35/group) for attrition\nMale | Age 18-50 | Opioid-dependent | Detoxified",
    11, DARK_TEXT, True)

# Male justification box
add_box(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.8), RGBColor(0xFE, 0xF9, 0xE7),
    "JUSTIFICATION FOR MALE-ONLY SAMPLE\n- Indian de-addiction centers (including Man Nasha Mukti Kendra) admit ~90-95% male patients\n- MAGNITUDE study reports male:female ratio ~10:1 for opioid dependence in India\n- Gender differences in craving, impulsivity, and mindfulness may confound results if mixed sample used\n- Homogeneous sample strengthens internal validity for this initial efficacy trial\n- Female-specific MBRP studies recommended as future direction",
    11, DARK_TEXT, True)

# ===================================================================
# SLIDE 19: SAMPLE SIZE (with formula visual)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, TEAL)
create_header(slide, "Sample Size Estimation", 19)

# Formula box
add_box(slide, Inches(2.5), Inches(1.5), Inches(8.3), Inches(1.0), DEEP_BLUE,
    "n = [(Za/2 + Zb)^2  x  2  x  sigma^2] / d^2",
    16, WHITE, True)

# Parameters
add_box(slide, Inches(0.6), Inches(2.8), Inches(6.0), Inches(1.5), SOFT_BLUE,
    "PARAMETERS\nEffect size (d) = 0.50 (medium)\nPower (1-beta) = 0.80 -> Zb = 0.84\nAlpha = 0.05 (two-tailed) -> Za/2 = 1.96",
    11, DEEP_BLUE, True)

# Calculation
add_box(slide, Inches(6.8), Inches(2.8), Inches(6.0), Inches(1.5), RGBColor(0xE8, 0xF8, 0xF5),
    "CALCULATION\nn = [(1.96 + 0.84)^2 x 2 x 1] / (0.50)^2\nn = [7.84 x 2] / 0.25 = 62.72\nn ~ 63 total (~32 per group)",
    11, TEAL, True)

# Final decision
add_box(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(1.8), RGBColor(0xFE, 0xF9, 0xE7),
    "FINAL DECISION: N = 60 (30 per group)\nJustification: ANCOVA as primary analysis reduces required n | G*Power 3.1 verification: ANCOVA with 1 covariate -> ~34/group\nRecruit 70 total (35/group) to account for ~15% attrition in substance-dependent populations\nConsistent with: Glasner-Edwards et al. (2017), Bowen & Marlatt (2009)",
    11, DARK_TEXT, True)

# ===================================================================
# SLIDE 20-21: INCLUSION & EXCLUSION
# ===================================================================
add_content_slide(prs, "Inclusion Criteria", [
    "1. Diagnosis of Substance Dependence as per ICD-10 (F10-F19) / ICD-11 criteria",
    "2. Primary substance: Opioids (heroin, pharmaceutical opioids); polysubstance with primary opioid dependence included",
    "3. Male participants aged 18-50 years",
    "4. Completed detoxification phase (minimum 7 days post-withdrawal)",
    "5. Currently admitted at Man Nasha Mukti Kendra",
    "6. Minimum education: 5th standard (ability to comprehend psychometric tools)",
    "7. Willingness to provide written informed consent",
    "8. Able to attend all 6 intervention sessions during 3-week period"
], 20)

add_content_slide(prs, "Exclusion Criteria", [
    "1. Severe psychiatric comorbidity: Psychotic disorders, Bipolar I, severe MDE with suicidality",
    "2. Significant cognitive impairment (MMSE < 24) or intellectual disability",
    "3. Active withdrawal symptoms (COWS score > 12)",
    "4. History of traumatic brain injury with LOC > 30 minutes",
    "5. Current participation in another structured psychological intervention study",
    "6. Medical instability requiring acute/intensive care",
    "7. History of prior formal mindfulness/meditation training exceeding 1 month"
], 21)



# ===================================================================
# SLIDE 22: VARIABLES (visual)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, TEAL)
create_header(slide, "Variables", 22)

# IV box
add_box(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.3), SOFT_TEAL_BG,
    "INDEPENDENT VARIABLE (IV): Type of Intervention\nLevel 1: Brief MBRP + TAU (Experimental Group)  |  Level 2: Psychoeducation + TAU (Control Group)",
    12, DEEP_BLUE, True)

# DV boxes
add_box(slide, Inches(0.6), Inches(3.0), Inches(3.8), Inches(1.5), RGBColor(0xFD, 0xED, 0xEC),
    "DV 1: CRAVING\nMeasured by OCDUS\n(Obsessive Compulsive\nDrug Use Scale)",
    11, ACCENT_RED, True)

add_box(slide, Inches(4.7), Inches(3.0), Inches(3.8), Inches(1.5), RGBColor(0xFE, 0xF9, 0xE7),
    "DV 2: IMPULSIVITY\nMeasured by BIS-11\n(Barratt Impulsiveness\nScale)",
    11, ACCENT_ORANGE, True)

add_box(slide, Inches(8.8), Inches(3.0), Inches(3.8), Inches(1.5), RGBColor(0xE8, 0xF8, 0xF5),
    "DV 3: MINDFULNESS\nMeasured by FFMQ\n(Five Facet Mindfulness\nQuestionnaire)",
    11, ACCENT_GREEN, True)

# Controlled variables
add_box(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.0), LIGHT_GRAY,
    "CONTROLLED VARIABLES: Age, education, duration of use, severity (ASSIST baseline), TAU components constant, session duration equalized",
    11, BODY_TEXT, True)

# ===================================================================
# SLIDE 23-25: TOOLS
# ===================================================================
add_content_slide(prs, "Tools: Mindfulness & Impulsivity Measures", [
    "1. FIVE FACET MINDFULNESS QUESTIONNAIRE (FFMQ) - Baer et al. (2006)",
    "   39 items | 5 facets: Observing, Describing, Acting with Awareness, Non-Judging, Non-Reactivity",
    "   5-point Likert (1-5) | Higher = greater mindfulness",
    "   Reliability: alpha = 0.75-0.91 | Validity: Convergent + sensitive to intervention | Hindi adaptation available",
    "",
    "2. BARRATT IMPULSIVENESS SCALE (BIS-11) - Patton et al. (1995)",
    "   30 items | 3 factors: Attentional, Motor, Non-Planning Impulsivity",
    "   4-point scale (1-4) | Higher = greater impulsivity",
    "   Reliability: alpha = 0.79-0.83; test-retest r = 0.83 | Discriminates SUD from controls | Hindi validated"
], 23)

add_content_slide(prs, "Tools: Craving Assessment - OCDUS", [
    "OBSESSIVE COMPULSIVE DRUG USE SCALE (OCDUS) - Franken et al. (2002)",
    "",
    "Description: 12-item self-report measuring obsessive thoughts about drug use and compulsive urges",
    "Scoring: 5-point scale (0-4) | Total range: 0-48 | Higher = greater craving",
    "Subscales: (1) Obsessive thoughts/interference, (2) Desire/control, (3) Resistance to thoughts",
    "Reliability: Internal consistency alpha = 0.86-0.90; test-retest r = 0.78",
    "Validity: Convergent with VAS craving (r = 0.55-0.67); predicts relapse; sensitive to treatment changes",
    "Indian Usability: Applicable to Indian SUD populations; adaptable for Hindi; brief (5 min); suitable for pre-post",
    "Justification: Captures both cognitive (obsessive) and behavioral (compulsive) craving dimensions; applicable across substances including opioids"
], 24)

add_content_slide(prs, "Tools: Severity Assessment - ASSIST (Pre-Test Only)", [
    "ALCOHOL, SMOKING AND SUBSTANCE INVOLVEMENT SCREENING TEST (ASSIST) - WHO (2002)",
    "",
    "Description: 8-item screening for risk level across 10 substance categories",
    "Scoring: Substance-specific risk -> Low (0-3), Moderate (4-26), High (27+)",
    "Reliability: Test-retest r = 0.58-0.90; alpha = 0.77-0.94",
    "Validity: Sensitivity = 0.80, Specificity = 0.71 for substance dependence",
    "Indian: WHO-validated; Hindi version available; used in NDDTC studies; 5-10 min",
    "PURPOSE: Used at PRE-TEST ONLY to establish baseline severity and ensure group equivalence",
    "NOT used as an outcome measure"
], 25)



# ===================================================================
# SLIDE 26: PROCEDURE (visual flow diagram)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, TEAL)
create_header(slide, "Procedure: Research Flow at Man Nasha Mukti Kendra", 26)

# Flow diagram - vertical steps
steps = [
    ("STEP 1: SCREENING", "Identify eligible male patients\n(ICD-10, detoxified, criteria met)", SOFT_BLUE),
    ("STEP 2: CONSENT", "Informed consent in Hindi\n(purpose, procedures, rights)", SOFT_BLUE),
    ("STEP 3: PRE-TEST", "Administer: ASSIST + OCDUS\n+ BIS-11 + FFMQ", SOFT_TEAL_BG),
    ("STEP 4: RANDOMIZE", "Computer-generated allocation\n-> Experimental vs. Control", RGBColor(0xFE, 0xF9, 0xE7)),
    ("STEP 5: INTERVENTION", "MBRP (6x60min, 3wks) OR\nPsychoeducation (6x60min, 3wks)", RGBColor(0xE8, 0xF8, 0xF5)),
    ("STEP 6: POST-TEST", "Administer: OCDUS + BIS-11\n+ FFMQ (within 1 week)", SOFT_BLUE),
    ("STEP 7: ANALYSIS", "Data compilation &\nstatistical analysis", LIGHT_GRAY),
]

x_positions = [Inches(0.3), Inches(2.1), Inches(3.9), Inches(5.7), Inches(7.5), Inches(9.3), Inches(11.1)]
for i, (title, desc, color) in enumerate(steps):
    x = x_positions[i]
    # Box
    add_box(slide, x, Inches(1.5), Inches(1.7), Inches(2.2), color,
        f"{title}\n{desc}", 9, DARK_TEXT, True)
    # Arrow between boxes
    if i < 6:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.75), Inches(2.4), Inches(0.3), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background()

# Note
txN = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(12.0), Inches(0.8))
tfN = txN.text_frame
tfN.word_wrap = True
pN = tfN.paragraphs[0]
pN.text = "Both groups receive TAU (pharmacotherapy, routine counseling) throughout | Assessments by blinded RA | Intervention by MPhil Clinical Psychologist"
pN.font.size = Pt(11)
pN.font.color.rgb = MID_GRAY

# ===================================================================
# SLIDE 27-28: INTERVENTION PLANS (visual)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, ACCENT_GREEN)
create_header(slide, "Intervention: Brief MBRP (6 Sessions / 3 Weeks) + TAU", 27)

sessions = [
    ("S1", "Introduction & Autopilot", "Raisin exercise; automatic patterns"),
    ("S2", "Triggers & Body Scan", "Body scan; personal trigger mapping"),
    ("S3", "Daily Mindfulness & SOBER", "Breath meditation; SOBER space"),
    ("S4", "Urge Surfing & High-Risk", "Urge surfing; decentering"),
    ("S5", "Acceptance & Non-Judgment", "Non-reactive awareness; skillful action"),
    ("S6", "Integration & Maintenance", "Loving-kindness; relapse prevention plan"),
]

for i, (num, topic, activity) in enumerate(sessions):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.4 + row * 2.5)
    color = RGBColor(0xE8, 0xF8, 0xF5) if i % 2 == 0 else SOFT_TEAL_BG
    add_box(slide, x, y, Inches(4.0), Inches(2.0), color,
        f"{num}: {topic}\n\n{activity}", 11, DARK_TEXT, True)

# Footer note
txF = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.5))
tfF = txF.text_frame
pF = tfF.paragraphs[0]
pF.text = "Delivery: Twice weekly | 60 min/session | Group (6-8 participants) | Hindi audio meditations | Practice logs"
pF.font.size = Pt(11)
pF.font.color.rgb = MID_GRAY

# CONTROL GROUP
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, ACCENT_ORANGE)
create_header(slide, "Intervention: Psychoeducation (6 Sessions / 3 Weeks) + TAU", 28)

ctrl_sessions = [
    ("S1", "Understanding Addiction", "Disease model; brain changes"),
    ("S2", "Effects of Opioids", "Physical & psychological consequences"),
    ("S3", "Understanding Relapse", "Warning signs; high-risk situations"),
    ("S4", "Health & Nutrition", "Physical recovery; sleep hygiene"),
    ("S5", "Social Consequences", "Family; legal; stigma; rehabilitation"),
    ("S6", "Motivation & Goals", "Stages of change; recovery planning"),
]

for i, (num, topic, activity) in enumerate(ctrl_sessions):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.4 + row * 2.5)
    color = RGBColor(0xFE, 0xF9, 0xE7) if i % 2 == 0 else RGBColor(0xFA, 0xF0, 0xDB)
    add_box(slide, x, y, Inches(4.0), Inches(2.0), color,
        f"{num}: {topic}\n\n{activity}", 11, DARK_TEXT, True)

txF = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.5))
tfF = txF.text_frame
pF = tfF.paragraphs[0]
pF.text = "Active control: Matched for time, format, attention | NO mindfulness component | Handouts + visual aids + discussion"
pF.font.size = Pt(11)
pF.font.color.rgb = MID_GRAY



# ===================================================================
# SLIDE 29: DATA ANALYSIS (with ANCOVA justification)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, TEAL)
create_header(slide, "Data Analysis", 29)

# Statistics overview
add_box(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.0), SOFT_BLUE,
    "STATISTICAL TESTS\nDescriptive: Mean, SD, frequencies\nNormality: Shapiro-Wilk test\nWithin-group: Paired t-test / Wilcoxon\nBetween-group: Independent t-test / Mann-Whitney\nPrimary: ANCOVA | Effect size: Partial eta^2, Cohen's d\nAlpha = 0.05 (two-tailed) | Software: SPSS 26.0",
    11, DEEP_BLUE, True)

# ANCOVA justification box
add_box(slide, Inches(6.6), Inches(1.4), Inches(6.1), Inches(3.5), RGBColor(0xE8, 0xF8, 0xF5),
    "DETAILED ANCOVA JUSTIFICATION\n\n(a) Controls for pre-existing baseline differences\n     between groups even after randomization\n(b) Increases statistical power by reducing\n     within-group error variance\n(c) Provides more precise treatment effect estimate\n     by adjusting post-test means for baseline\n(d) Reduces required sample size compared to\n     independent t-test (critical for clinical populations)\n(e) Recommended for pre-post designs by\n     Tabachnick & Fidell (2013) and Field (2018)",
    11, DARK_TEXT, True)

# ANCOVA model box
add_box(slide, Inches(0.6), Inches(3.7), Inches(5.8), Inches(1.5), RGBColor(0xFE, 0xF9, 0xE7),
    "ANCOVA MODEL\nDV: Post-test scores (OCDUS / BIS-11 / FFMQ)\nIV: Group (Experimental vs. Control)\nCovariate: Pre-test scores (same measure)\nITT: Last Observation Carried Forward (LOCF)",
    11, DARK_TEXT, True)

# ===================================================================
# SLIDE 30: ETHICAL CONSIDERATIONS
# ===================================================================
add_content_slide(prs, "Ethical Considerations", [
    "Informed Consent: Written consent in Hindi; participants fully informed of purpose, procedures, duration, risks, benefits",
    "Voluntary Participation: Right to withdraw anytime without penalty or impact on treatment at Man Nasha Mukti Kendra",
    "Confidentiality: Data coded with participant IDs; no identifying information in publications; secure storage",
    "Non-Maleficence: Control receives active psychoeducation (not waitlist); TAU continued for all",
    "Institutional Approval: Ethical clearance from IEC prior to data collection",
    "Debriefing: Control group offered brief MBRP orientation post-study",
    "Compliance: ICMR (2017) National Ethical Guidelines for Biomedical and Health Research"
], 30)



# ===================================================================
# SLIDE 31-32: EXPECTED RESULTS (with bar chart)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, ACCENT_GREEN)
create_header(slide, "Expected Results (1/2)", 31)

# Add a bar chart showing expected pattern
chart_data = CategoryChartData()
chart_data.categories = ['Craving\n(OCDUS)', 'Impulsivity\n(BIS-11)', 'Mindfulness\n(FFMQ)']
chart_data.add_series('Pre-Test (Both Groups)', (35, 72, 95))
chart_data.add_series('Post-Test: MBRP Group', (18, 55, 130))
chart_data.add_series('Post-Test: Control Group', (28, 65, 105))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.6), Inches(1.4), Inches(7.0), Inches(5.5),
    chart_data
).chart

chart.has_legend = True
chart.legend.include_in_layout = False

# Expected results text
txR = slide.shapes.add_textbox(Inches(7.8), Inches(1.5), Inches(5.0), Inches(5.5))
tfR = txR.text_frame
tfR.word_wrap = True
bullets = [
    "CRAVING (OCDUS):",
    "Significant REDUCTION expected in",
    "MBRP group vs. Control",
    "Effect: d = 0.50-0.80",
    "Mechanism: Urge surfing disrupts",
    "automatic craving-use cycle",
    "",
    "IMPULSIVITY (BIS-11):",
    "Significant REDUCTION expected",
    "(Motor + Attentional subscales)",
    "Effect: d = 0.40-0.60",
    "Mechanism: Mindfulness strengthens",
    "prefrontal inhibitory control",
]
for i, b in enumerate(bullets):
    p = tfR.paragraphs[0] if i == 0 else tfR.add_paragraph()
    p.text = b
    p.font.size = Pt(11)
    p.font.color.rgb = BODY_TEXT
    if "CRAVING" in b or "IMPULSIVITY" in b:
        p.font.bold = True
        p.font.color.rgb = TEAL

# Expected results 2/2
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, ACCENT_GREEN)
create_header(slide, "Expected Results (2/2)", 32)

# Mindfulness chart - line showing increase
chart_data2 = CategoryChartData()
chart_data2.categories = ['Pre-Test', 'Post-Test']
chart_data2.add_series('MBRP Group - Mindfulness', (95, 130))
chart_data2.add_series('Control Group - Mindfulness', (95, 105))
chart_data2.add_series('MBRP Group - Craving', (35, 18))
chart_data2.add_series('Control Group - Craving', (35, 28))

chart2 = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(0.6), Inches(1.4), Inches(6.5), Inches(4.5),
    chart_data2
).chart
chart2.has_legend = True

# Text
txR = slide.shapes.add_textbox(Inches(7.4), Inches(1.5), Inches(5.5), Inches(5.5))
tfR = txR.text_frame
tfR.word_wrap = True
bullets2 = [
    "MINDFULNESS (FFMQ):",
    "Significant INCREASE expected in",
    "MBRP group (Acting with Awareness",
    "& Non-Reactivity facets)",
    "Effect: d = 0.50-0.70",
    "Mechanism: Structured meditation",
    "cultivates dispositional mindfulness",
    "",
    "OVERALL PATTERN:",
    "Brief MBRP + TAU demonstrates",
    "superiority across all 3 DVs",
    "Null hypotheses (H01, H02, H03)",
    "expected to be REJECTED",
    "",
    "Supports feasibility of 6-session",
    "MBRP at Man Nasha Mukti Kendra"
]
for i, b in enumerate(bullets2):
    p = tfR.paragraphs[0] if i == 0 else tfR.add_paragraph()
    p.text = b
    p.font.size = Pt(11)
    p.font.color.rgb = BODY_TEXT
    if "MINDFULNESS" in b or "OVERALL" in b:
        p.font.bold = True
        p.font.color.rgb = TEAL



# ===================================================================
# SLIDE 33: CLINICAL IMPLICATIONS
# ===================================================================
add_content_slide(prs, "Clinical Implications", [
    "Validates a brief MBRP model (6 sessions/3 weeks) feasible for Indian de-addiction settings",
    "Provides evidence-based psychological intervention to complement pharmacotherapy (OST, naltrexone)",
    "Demonstrates mindfulness-based approaches are culturally compatible with Indian populations",
    "Addresses multiple relapse risk factors simultaneously through single integrated protocol",
    "Supports task-shifting: Brief MBRP deliverable by MPhil-trained Clinical Psychologists",
    "Informs national treatment policy (NIMHANS, NDDTC, State Mental Health Authorities)",
    "Contributes to RCI-recognized intervention repertoire for clinical psychology practice in India"
], 33)

# ===================================================================
# SLIDE 34: LIMITATIONS
# ===================================================================
add_content_slide(prs, "Limitations", [
    "Male-only sample from single center (Man Nasha Mukti Kendra) limits generalizability",
    "Short-term assessment: Post-test immediately after 3-week intervention; no follow-up",
    "Self-report measures (OCDUS, BIS-11, FFMQ) susceptible to social desirability bias",
    "No biological markers: Craving measured subjectively without physiological corroboration",
    "Therapist effects: Single therapist delivery may introduce confounds",
    "Attention-matched control does not fully account for specific mindfulness mechanisms",
    "Potential attrition despite over-recruitment planning"
], 34)

# ===================================================================
# SLIDE 35: FUTURE DIRECTIONS
# ===================================================================
add_content_slide(prs, "Future Directions", [
    "3-month and 6-month follow-up assessments for sustained MBRP effects",
    "Multi-site RCTs across multiple Indian de-addiction centers",
    "Include female participants to examine gender-specific effects",
    "Neuroimaging (fMRI/EEG) to examine neural mechanisms of MBRP",
    "Dose-response analysis: Compare 4 vs. 6 vs. 8 sessions for optimal dosage",
    "Mediator analysis: Test whether mindfulness mediates craving/impulsivity reduction",
    "Technology-assisted: App-based/digital MBRP for post-discharge and rural access",
    "Comparative effectiveness: MBRP vs. CBT vs. ACT in Indian samples"
], 35)

# ===================================================================
# SLIDE 36: SUMMARY
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = OFF_WHITE
add_side_accent(slide, DEEP_BLUE)
create_header(slide, "Summary", 36)

add_box(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5), SOFT_BLUE,
    "STUDY SUMMARY\n\n"
    "Design: Pre-test Post-test Control Group Experimental Design\n"
    "Setting: Man Nasha Mukti Kendra | Population: Male opioid-dependent patients (N=60)\n"
    "Experimental: Brief MBRP (6 sessions / 3 weeks) + TAU | Control: Psychoeducation + TAU\n"
    "DVs: Craving (OCDUS) + Impulsivity (BIS-11) + Mindfulness (FFMQ)\n"
    "Sampling: Two-stage (Purposive selection -> Computer-generated random assignment)\n"
    "Primary Analysis: ANCOVA controlling for pre-test scores as covariates\n"
    "Hypothesis: Brief MBRP expected to show significantly greater reduction in craving/impulsivity\n"
    "and greater enhancement of mindfulness compared to psychoeducation\n\n"
    "SIGNIFICANCE: First brief MBRP trial in Indian opioid-dependent sample addressing\n"
    "critical research gap in evidence-based psychological interventions for de-addiction",
    12, DEEP_BLUE, True)

# ===================================================================
# SLIDE 37: CONCLUSION
# ===================================================================
add_content_slide(prs, "Conclusion", [
    "Substance dependence (particularly opioid) represents a significant public health crisis in India with relapse rates >70%",
    "Current treatment at Indian de-addiction centers relies heavily on pharmacotherapy with limited evidence-based psychological interventions",
    "MBRP offers a theoretically grounded approach targeting craving (urge surfing), impulsivity (mindful pause), and mindfulness enhancement simultaneously",
    "A brief 6-session MBRP protocol is clinically practical, culturally appropriate, and feasible at Man Nasha Mukti Kendra",
    "If supported, Brief MBRP can be integrated into standard de-addiction care as a cost-effective, replicable intervention",
    "This study contributes to evidence for mindfulness-based interventions in addiction while addressing Indian clinical psychology practice needs",
    "Results will inform RCI-recognized training, national treatment policy, and relapse burden reduction"
], 37)



# ===================================================================
# SLIDES 38-41: REFERENCES (APA 7)
# ===================================================================
add_content_slide(prs, "References (1/4)", [
    "Baer, R. A., Smith, G. T., Hopkins, J., Krietemeyer, J., & Toney, L. (2006). Using self-report assessment methods to explore facets of mindfulness. Assessment, 13(1), 27-45.",
    "Bowen, S., Chawla, N., & Marlatt, G. A. (2011). Mindfulness-based relapse prevention for addictive behaviors: A clinician's guide. Guilford Press.",
    "Bowen, S., & Marlatt, G. A. (2009). Surfing the urge: Brief mindfulness-based intervention. Psychology of Addictive Behaviors, 23(4), 666-671.",
    "Bowen, S., Witkiewitz, K., Clifasefi, S. L., et al. (2014). Relative efficacy of MBRP, standard RP, and TAU. JAMA Psychiatry, 71(5), 547-556.",
    "Brewer, J. A., Mallik, S., Babuscio, T. A., et al. (2011). Mindfulness training for smoking cessation. Drug and Alcohol Dependence, 119(1-2), 72-80.",
    "Chiesa, A., & Serretti, A. (2014). Are MBIs effective for substance use disorders? Substance Use & Misuse, 49(5), 492-512."
], 38)

add_content_slide(prs, "References (2/4)", [
    "Franken, I. H. A., Hendriks, V. M., & van den Brink, W. (2002). New instruments to measure impulsivity and compulsivity. Psychiatrie en Neurobiologie, 1, 10-14.",
    "Garland, E. L., Froeliger, B., & Howard, M. O. (2014). Mindfulness targets neurocognitive mechanisms of addiction. Frontiers in Psychiatry, 4, 173.",
    "Garland, E. L., Manusov, E. G., Froeliger, B., et al. (2014). MORE for chronic pain and opioid misuse. J. Consulting and Clinical Psychology, 82(3), 448-459.",
    "Garland, E. L., Roberts-Lewis, A., Tronnier, C. D., et al. (2016). MORE vs. CBT for co-occurring SUDs. J. Consulting and Clinical Psychology, 84(4), 281-293.",
    "Ghosh, A., Basu, D., & Avasthi, A. (2018). Relapse in opioid dependence: Indian perspective. Indian Journal of Psychiatry, 60(Suppl 4), S469-S476.",
    "Glasner-Edwards, S., Mooney, L. J., Ang, A., et al. (2017). MBRP for stimulant dependent adults: Pilot RCT. Mindfulness, 8(1), 126-135."
], 39)

add_content_slide(prs, "References (3/4)", [
    "Grant, S., Colaiaco, B., Motala, A., et al. (2017). MBRP for SUDs: Meta-analysis. J. Addiction Medicine, 11(5), 386-396.",
    "Humeniuk, R., Ali, R., Babor, T. F., et al. (2008). Validation of ASSIST. Addiction, 103(6), 1039-1047.",
    "Jain, R., Majumder, P., & Gupta, T. (2013). Indian Journal of Psychiatry, 55(Suppl 1), S86-S92.",
    "Kabat-Zinn, J. (1990). Full catastrophe living. Delacorte Press.",
    "Karyadi, K. A., VanderVeen, J. D., & Cyders, M. A. (2014). Trait mindfulness and substance use: Meta-analysis. Drug and Alcohol Dependence, 143, 1-10.",
    "Li, W., Howard, M. O., Garland, E. L., et al. (2017). Mindfulness for substance misuse: Meta-analysis. J. Substance Abuse Treatment, 75, 62-96.",
    "Marlatt, G. A., & Gordon, J. R. (1985). Relapse prevention. Guilford Press."
], 40)

add_content_slide(prs, "References (4/4)", [
    "Mattoo, S. K., Chakrabarti, S., & Anjaiah, M. (2009). Psychosocial factors in relapse. Indian J. Medical Research, 130(6), 702-708.",
    "Ministry of Social Justice & Empowerment. (2019). Magnitude of substance use in India. Government of India.",
    "Moeller, F. G., Barratt, E. S., Dougherty, D. M., et al. (2001). Psychiatric aspects of impulsivity. Am. J. Psychiatry, 158(11), 1783-1793.",
    "Murphy, C., & MacKillop, J. (2012). Impulsivity, mindfulness, and alcohol misuse. Psychopharmacology, 219(2), 527-536.",
    "NIDA. (2020). Drugs, brains, and behavior: The science of addiction.",
    "Patton, J. H., Stanford, M. S., & Barratt, E. S. (1995). Factor structure of BIS. J. Clinical Psychology, 51(6), 768-774.",
    "Robinson, T. E., & Berridge, K. C. (1993). Neural basis of drug craving. Brain Research Reviews, 18(3), 247-291.",
    "Sarkar, S., & Balhara, Y. P. S. (2016). Indian Journal of Psychiatry, 58(3), 290-295.",
    "Tabachnick, B. G., & Fidell, L. S. (2013). Using multivariate statistics (6th ed.). Pearson.",
    "WHO ASSIST Working Group. (2002). ASSIST development. Addiction, 97(9), 1183-1194.",
    "Witkiewitz, K., Bowen, S., Douglas, H., & Hsu, S. H. (2013). MBRP for craving. Addictive Behaviors, 38(2), 1563-1571."
], 41)

# ===================================================================
# SLIDE 42: THANK YOU (Premium)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = NAVY

# Decorative elements
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(4.5), Inches(3.5), Inches(3.5))
circle1.fill.solid()
circle1.fill.fore_color.rgb = DEEP_BLUE
circle1.line.color.rgb = TEAL
circle1.line.width = Pt(2)

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(-0.5), Inches(2.5), Inches(2.5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = DEEP_BLUE
circle2.line.color.rgb = GOLD
circle2.line.width = Pt(1.5)

# Thank you text
txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Questions & Discussion"
p2.font.size = Pt(18)
p2.font.color.rgb = LIGHT_TEAL
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)
p3 = tf.add_paragraph()
p3.text = ""
p4 = tf.add_paragraph()
p4.text = "Man Nasha Mukti Kendra  |  [University Name]"
p4.font.size = Pt(13)
p4.font.color.rgb = MID_GRAY
p4.alignment = PP_ALIGN.CENTER

# Gold bar
add_decorative_bar(slide, Inches(7.2), GOLD)

# ===================================================================
# SAVE
# ===================================================================
output = "/projects/sandbox/Dango-kiro/MBRP_Research_Synopsis_PPT.pptx"
prs.save(output)
print(f"SUCCESS: {output}")
print(f"Total slides: {len(prs.slides)}")
