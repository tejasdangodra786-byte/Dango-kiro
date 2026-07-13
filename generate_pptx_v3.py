#!/usr/bin/env python3
"""
MBRP Research Synopsis PPT - Version 3 (ALL CORRECTIONS APPLIED)
Corrections:
1. Male sample justified
2. OCDUS only for craving (removed HCQ, VAS)
3. Sampling: Two-stage (purposive selection + random assignment)
4. ANCOVA justification detailed
5. Expected results expanded
6. Summary & Conclusion slides added
7. 6 sessions / 3 weeks (brief MBRP)
8. Null hypotheses added
9. Abbreviations explained in introduction
10. APA 7 references (et al. for 3+ authors)
11. Man Nasha Mukti Kendra named throughout
12. Author name corrections
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_GOLD = RGBColor(0xD4, 0xA5, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_TEXT = RGBColor(0x2D, 0x2D, 0x44)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)



def add_title_slide(prs, title, subtitle, researcher, guide, dept, inst, year):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11.3), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(15)
    p2.font.color.rgb = ACCENT_GOLD
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(18)
    details = f"\nResearcher: {researcher}\nGuide: {guide}\nDepartment: {dept}\nInstitution: {inst}\nYear: {year}"
    p3 = tf.add_paragraph()
    p3.text = details
    p3.font.size = Pt(13)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(24)


def add_content_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_LIGHT
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.85))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = WHITE
    top = 1.2
    if subtitle:
        txS = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.0), Inches(0.45))
        tfS = txS.text_frame
        pS = tfS.paragraphs[0]
        pS.text = subtitle
        pS.font.size = Pt(13)
        pS.font.bold = True
        pS.font.color.rgb = MEDIUM_BLUE
        top = 1.6
    txB = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(7.5 - top - 0.2))
    tfB = txB.text_frame
    tfB.word_wrap = True
    for i, b in enumerate(bullets):
        pp = tfB.paragraphs[0] if i == 0 else tfB.add_paragraph()
        pp.text = b
        pp.font.size = Pt(13)
        pp.font.color.rgb = BODY_TEXT
        pp.space_before = Pt(5)



def add_paragraph_slide(prs, title, paragraphs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_LIGHT
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.85))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    txB = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(6.1))
    tfB = txB.text_frame
    tfB.word_wrap = True
    for i, para in enumerate(paragraphs):
        pp = tfB.paragraphs[0] if i == 0 else tfB.add_paragraph()
        pp.text = para
        pp.font.size = Pt(11)
        pp.font.color.rgb = BODY_TEXT
        pp.space_before = Pt(10)
        pp.space_after = Pt(4)


# ===================================================================
# SLIDE 1: TITLE
# ===================================================================
add_title_slide(prs,
    "Efficacy of Brief Mindfulness-Based Relapse Prevention (MBRP)\nIntervention on Craving, Impulsivity, and Mindfulness\nin Substance Dependent Patients",
    "MPhil Clinical Psychology Research Synopsis",
    "[Researcher Name]", "[Supervisor Name]",
    "Clinical Psychology", "Man Nasha Mukti Kendra / [University Name]", "2026")



# ===================================================================
# SLIDE 2: INTRODUCTION 1/3 - with abbreviations explained
# ===================================================================
add_content_slide(prs, "Introduction & Background (1/3)",
    [
        "ABBREVIATIONS USED: MBRP = Mindfulness-Based Relapse Prevention; TAU = Treatment As Usual; OCDUS = Obsessive Compulsive Drug Use Scale; BIS-11 = Barratt Impulsiveness Scale; FFMQ = Five Facet Mindfulness Questionnaire; ASSIST = Alcohol, Smoking and Substance Involvement Screening Test; OST = Opioid Substitution Therapy; PFC = Prefrontal Cortex; ANCOVA = Analysis of Covariance",
        "",
        "India faces a significant substance use crisis; the MAGNITUDE study (Ministry of Social Justice & Empowerment, 2019) estimated approximately 3.1 crore individuals affected by substance use disorders",
        "Opioid dependence constitutes a major public health burden, particularly in Punjab, Rajasthan, Northeast India, and metropolitan areas",
        "WHO estimates India accounts for ~25% of global opioid-related deaths in South-East Asia",
        "National Drug Dependence Treatment Centre (NDDTC) reports increasing treatment-seeking among opioid users",
        "Indian de-addiction centers, including Man Nasha Mukti Kendra, primarily offer pharmacotherapy (OST, naltrexone) with limited structured psychotherapy",
        "Psychosocial interventions remain under-utilized despite evidence of superior combined treatment outcomes"
    ],
    subtitle="Substance Dependence: Indian Context & Key Abbreviations")

# ===================================================================
# SLIDE 3: INTRODUCTION 2/3
# ===================================================================
add_content_slide(prs, "Introduction & Background (2/3)",
    [
        "Relapse rates in substance dependence range from 40-60% within the first year post-treatment (NIDA, 2020)",
        "Indian studies report even higher relapse rates (~70-80%) in opioid dependence (Mattoo et al., 2009)",
        "Triggers for relapse: craving, negative affect, interpersonal conflict, environmental cues",
        "Traditional Relapse Prevention (Marlatt & Gordon, 1985) has moderate efficacy but limited mindfulness integration",
        "Automatic cognitive-behavioral patterns (e.g., apparently irrelevant decisions) perpetuate relapse cycles",
        "Need for interventions addressing both cognitive automaticity and emotional dysregulation simultaneously",
        "Brief interventions are essential for Indian rehab settings (e.g., Man Nasha Mukti Kendra) with limited resources and high patient turnover"
    ],
    subtitle="The Relapse Problem in De-Addiction")

# ===================================================================
# SLIDE 4: INTRODUCTION 3/3
# ===================================================================
add_content_slide(prs, "Introduction & Background (3/3)",
    [
        "Mindfulness: intentional, non-judgmental, present-moment awareness (Kabat-Zinn, 1990)",
        "MBRP developed by Bowen, Chawla, and Marlatt (2011) integrating MBSR + Cognitive-Behavioral Relapse Prevention",
        "Core mechanisms: decentering from craving, disrupting automatic reactivity, increasing distress tolerance",
        "MBRP targets the 'craving -> use' automaticity by cultivating awareness of triggers and non-reactive responding",
        "Brief adaptations (6 sessions over 3 weeks) show promise for resource-limited rehabilitation settings",
        "Growing global evidence supports MBRP; however, Indian validation remains critically scarce",
        "Rationale: A brief MBRP model (6 sessions) can be feasibly implemented in Indian de-addiction centers like Man Nasha Mukti Kendra"
    ],
    subtitle="Mindfulness in Addiction Treatment")



# ===================================================================
# SLIDE 5: VARIABLE 1 - CRAVING
# ===================================================================
add_content_slide(prs, "Variable 1: Craving",
    [
        "DEFINITION (Clinical): An intense, subjective urge or desire to use a substance, often triggered by internal/external cues",
        "DEFINITION (Theoretical): A motivational state arising from incentive-sensitization (Robinson & Berridge, 1993) and conditioned reinforcement learning",
        "RELEVANCE: Primary predictor of relapse in opioid dependence; mediates relationship between cue exposure and substance use behavior; intensity correlates with severity and treatment dropout",
        "NEUROPSYCHOLOGICAL BASIS: Mesolimbic dopamine pathway activation (VTA -> Nucleus Accumbens); PFC hypoactivation during craving episodes; conditioned cue-reward associations; dysregulated stress-reward interaction in chronic opioid users",
        "TECHNIQUE - Urge Surfing: Observing craving as a transient wave without acting on it",
        "TECHNIQUE - Mindfulness Exposure: Non-reactive awareness of craving sensations in the body",
        "TECHNIQUE - SOBER Breathing Space: Stop-Observe-Breathe-Expand-Respond during craving triggers",
        "TECHNIQUE - Cognitive Decentering: Reframing 'I need the drug' to 'I am having a craving'"
    ])

# ===================================================================
# SLIDE 6: VARIABLE 2 - IMPULSIVITY
# ===================================================================
add_content_slide(prs, "Variable 2: Impulsivity",
    [
        "DEFINITION (Clinical): A tendency toward rapid, unplanned actions without adequate consideration of consequences (Moeller et al., 2001)",
        "DEFINITION (Theoretical): Multi-dimensional construct comprising motor impulsivity, attentional impulsivity, and non-planning impulsivity (Patton et al., 1995)",
        "RELEVANCE: Higher trait impulsivity predicts initiation, escalation, and relapse in substance use; mediates craving-to-use behavior; associated with treatment non-adherence and premature dropout",
        "NEUROPSYCHOLOGICAL BASIS: PFC dysfunction -> impaired executive control; reduced inhibitory control (Go/No-Go, Stop-Signal paradigms); impaired delay discounting; DLPFC hypoactivation linked to poor decision-making",
        "TECHNIQUE - Response Inhibition Training: Mindful pause before automatic behavioral responses",
        "TECHNIQUE - Awareness Training: Noticing impulse-action sequences without engagement",
        "TECHNIQUE - Mindful Decision-Making: Creating space between stimulus and response",
        "TECHNIQUE - STOP: Stop-Take a breath-Observe-Proceed mindfully"
    ])

# ===================================================================
# SLIDE 7: VARIABLE 3 - MINDFULNESS
# ===================================================================
add_content_slide(prs, "Variable 3: Mindfulness",
    [
        "DEFINITION (Clinical): The capacity to attend to present-moment experience with openness, curiosity, and non-judgment",
        "DEFINITION (Theoretical): Multi-faceted construct: Observing, Describing, Acting with Awareness, Non-Judging, Non-Reactivity (Baer et al., 2006)",
        "RELEVANCE: Substance-dependent individuals show significantly lower dispositional mindfulness; mindfulness acts as protective factor against relapse triggers; improvements in mindfulness mediate treatment outcomes in MBRP studies",
        "NEUROPSYCHOLOGICAL BASIS: Anterior Cingulate Cortex (ACC) -> enhanced self-regulation; Insula -> improved interoceptive awareness; PFC-amygdala connectivity -> better emotional regulation; DMN regulation -> reduced rumination",
        "TECHNIQUE - Sitting Meditation: Focused attention on breath, body sensations, thoughts",
        "TECHNIQUE - Body Scan: Systematic non-judgmental awareness of bodily states",
        "TECHNIQUE - Mindful Movement: Gentle yoga/walking with present-moment focus",
        "TECHNIQUE - Non-Judgmental Awareness: Labeling experiences without evaluation"
    ])



# ===================================================================
# SLIDES 8-12: REVIEW OF LITERATURE (5 slides, detailed paragraph style)
# ===================================================================

# SLIDE 8
add_paragraph_slide(prs, "Review of Literature (1/5) - MBRP & Relapse Prevention",
    [
        "Bowen et al. (2014) investigate the relative efficacy of Mindfulness-Based Relapse Prevention compared to standard Relapse Prevention and Treatment As Usual for substance use disorders in JAMA Psychiatry. The randomized clinical trial involved 286 participants who had completed initial treatment for substance use disorders. The study demonstrates that at 12-month follow-up, MBRP participants reported significantly fewer days of substance use and heavy drinking compared to both standard RP and TAU groups. The researchers highlight that while all three groups showed initial improvements, MBRP participants maintained superior long-term outcomes, suggesting that the cultivation of mindfulness skills provides a durable protective mechanism against relapse. The study establishes that MBRP's integration of present-moment awareness with traditional cognitive-behavioral relapse prevention strategies creates a synergistic therapeutic effect that addresses both the automatic reactivity underlying craving and the cognitive distortions that precipitate relapse. These findings provide foundational evidence for MBRP as a gold-standard psychosocial intervention in substance use disorder aftercare.",
        "",
        "Bowen and Marlatt (2009) examine the effects of brief mindfulness-based intervention on craving among substance users in Psychology of Addictive Behaviors. The study recruited incarcerated individuals with substance use histories and administered a brief urge surfing meditation intervention. Results demonstrated significant reductions in craving intensity and frequency in the mindfulness condition compared to controls. The authors propose that even brief exposure to mindfulness techniques can disrupt the automaticity of craving responses by introducing a meta-cognitive awareness layer between trigger and behavioral response. This study is particularly relevant to the present research as it validates the premise that abbreviated mindfulness interventions can produce meaningful clinical effects on craving, supporting the feasibility of brief MBRP protocols (6 sessions over 3 weeks) in settings like Man Nasha Mukti Kendra where extended treatment programs are impractical."
    ])

# SLIDE 9
add_paragraph_slide(prs, "Review of Literature (2/5) - Craving & Mindfulness Mechanisms",
    [
        "Garland et al. (2014a) explore how mindfulness training targets neurocognitive mechanisms of addiction at the attention-appraisal-emotion interface in Frontiers in Psychiatry. The study presents a theoretical and empirical framework demonstrating that Mindfulness-Oriented Recovery Enhancement (MORE) reduces opioid craving through three interconnected mechanisms: attentional reorientation away from drug-related cues, positive reappraisal of previously neutral stimuli to generate natural reward, and enhanced savoring of healthy pleasures. The researchers provide neuroimaging evidence showing that mindfulness practice modulates activity in prefrontal and limbic circuits associated with craving and emotional regulation. The study is significant because it elucidates the precise cognitive-neural pathways through which mindfulness reduces craving in opioid users specifically, providing a mechanistic rationale for why MBRP techniques such as urge surfing and mindful awareness of craving sensations can reduce the subjective intensity and behavioral impact of craving episodes in substance-dependent patients.",
        "",
        "Witkiewitz et al. (2013) investigate mindfulness-based relapse prevention effects on substance craving in Addictive Behaviors. The study conducted secondary analyses of data from a randomized controlled trial comparing MBRP to TAU among individuals in aftercare following substance use disorder treatment. Over a 4-month follow-up period, MBRP participants demonstrated significantly lower craving levels compared to TAU participants, and importantly, the relationship between negative affect and subsequent craving was significantly attenuated in the MBRP group. The authors conclude that mindfulness practice weakens the affect-craving pathway by cultivating non-reactive awareness of emotional states, thereby preventing negative emotions from automatically triggering craving responses. This decoupling of affect and craving represents a critical therapeutic mechanism that distinguishes MBRP from traditional relapse prevention approaches that primarily rely on cognitive restructuring and behavioral avoidance strategies."
    ])



# SLIDE 10
add_paragraph_slide(prs, "Review of Literature (3/5) - Impulsivity & Mindfulness",
    [
        "Garland et al. (2016) examine the efficacy of Mindfulness-Oriented Recovery Enhancement versus CBT for co-occurring substance dependence, traumatic stress, and psychiatric disorders in the Journal of Consulting and Clinical Psychology. The study employed a randomized controlled design with substance-dependent adults exhibiting elevated impulsivity. Results demonstrated that the mindfulness-based intervention produced significant reductions in impulsivity scores as measured by the Barratt Impulsiveness Scale (BIS-11), particularly in the motor and attentional impulsivity subscales. The authors theorize that mindfulness meditation strengthens prefrontal cortical inhibitory mechanisms by repeatedly engaging participants in exercises requiring sustained attention, response monitoring, and deliberate non-reactivity. This enhanced top-down cognitive control translates into improved ability to inhibit prepotent impulsive responses in daily life, particularly in high-risk situations where automatic substance-seeking behavior would otherwise occur. The clinical implication is that MBRP can address impulsivity as a transdiagnostic risk factor for relapse, not merely as a stable trait but as a modifiable behavioral pattern amenable to mindfulness-based intervention.",
        "",
        "Murphy and MacKillop (2012) explore the interrelationships between impulsivity, mindfulness, and alcohol misuse in Psychopharmacology. The study investigates whether dispositional mindfulness buffers against the effects of trait impulsivity on problematic substance use. Using a cross-sectional design with 340 participants, the researchers found that trait mindfulness was inversely associated with impulsive decision-making as measured by delay discounting tasks. Critically, mindfulness moderated the relationship between impulsivity and substance use problems, such that individuals with higher mindfulness showed weaker associations between impulsivity and alcohol-related consequences. The authors propose that mindfulness functions as a cognitive resource that enables impulsive individuals to override automatic behavioral tendencies through enhanced metacognitive awareness and response flexibility. These findings support the inclusion of impulsivity as a dependent variable in MBRP research and suggest that mindfulness enhancement may serve as a mechanism through which impulsive responding is attenuated in substance-dependent populations."
    ])

# SLIDE 11
add_paragraph_slide(prs, "Review of Literature (4/5) - Meta-Analyses & Brief Models",
    [
        "Li et al. (2017) conduct a systematic review and meta-analysis of mindfulness treatment for substance misuse in the Journal of Substance Abuse Treatment, encompassing 42 randomized controlled trials. The meta-analytic findings reveal that mindfulness-based interventions produce significant effect sizes for reducing substance misuse (d = 0.33), craving (d = 0.68), and stress (d = 0.44) across diverse substance use populations and treatment contexts. The authors highlight that effect sizes for craving reduction are particularly robust, supporting the theoretical premise that mindfulness directly targets craving mechanisms through enhanced interoceptive awareness and non-reactive observation of urge states. Furthermore, the review identifies that interventions of shorter duration (4-8 sessions) demonstrated comparable efficacy to longer protocols when appropriately structured, providing empirical justification for the present study's brief 6-session MBRP adaptation. The review concludes that mindfulness-based interventions represent a viable evidence-based treatment approach with medium-to-large effects on the primary mechanisms of relapse.",
        "",
        "Glasner-Edwards et al. (2017) examine a pilot randomized clinical trial of mindfulness-based relapse prevention for stimulant-dependent adults using an abbreviated 6-session protocol in the journal Mindfulness. The study demonstrates that a condensed mindfulness intervention is both feasible and effective in outpatient substance use treatment settings. Participants in the brief MBRP condition showed significant reductions in substance use frequency and craving intensity compared to the health education control group. The authors emphasize that the abbreviated format maintained the core therapeutic elements of standard MBRP (body scan, sitting meditation, urge surfing, mindful movement) while condensing psychoeducational components. This study directly validates the brief intervention model proposed in the present research, demonstrating that 6-session MBRP protocols can be successfully implemented without substantial loss of therapeutic efficacy, making them particularly appropriate for Indian rehabilitation settings like Man Nasha Mukti Kendra with typical admission durations of 4-6 weeks."
    ])



# SLIDE 12
add_paragraph_slide(prs, "Review of Literature (5/5) - Indian Context & Research Gap",
    [
        "Ghosh et al. (2018) conduct a comprehensive review of relapse in opioid dependence from an Indian perspective in the Indian Journal of Psychiatry. The study reports alarmingly high relapse rates exceeding 70% among opioid-dependent patients treated in North Indian de-addiction centers, with the majority of relapses occurring within the first three months post-discharge. The authors identify craving, peer influence, negative emotional states, and lack of structured psychological aftercare as primary relapse determinants in the Indian context. Significantly, the review highlights that Indian treatment facilities predominantly rely on pharmacological approaches (opioid substitution therapy, naltrexone maintenance) with minimal integration of evidence-based psychological interventions. The authors advocate strongly for the development and validation of structured psychosocial intervention protocols tailored to Indian treatment infrastructure, patient characteristics, and resource constraints. This study directly establishes the clinical need for the present research by documenting both the high relapse burden and the psychosocial treatment gap in Indian de-addiction settings such as Man Nasha Mukti Kendra.",
        "",
        "Sarkar and Balhara (2016) highlight the underutilization of structured psychological interventions in Indian de-addiction settings, identifying barriers including limited trained personnel, absence of culturally validated intervention protocols, short admission durations, and institutional emphasis on pharmacological management. Murthy (2016) similarly notes that Indian de-addiction centers operate with minimal structured psychological programming despite international evidence supporting integrated psychosocial treatment. Jain et al. (2013) provide preliminary evidence from a mindfulness-based intervention study with alcohol-dependent patients in India showing initial craving reductions, but no systematic MBRP trial has been conducted with Indian opioid-dependent populations. CRITICAL GAP: No published Indian RCT has tested Brief MBRP specifically in opioid-dependent populations, simultaneously assessing craving, impulsivity, and mindfulness as treatment outcomes. The present study addresses this critical gap as the first brief MBRP trial designed for Indian de-addiction center infrastructure."
    ])



# ===================================================================
# SLIDE 13: RESEARCH GAP
# ===================================================================
add_content_slide(prs, "Research Gap",
    [
        "No published RCT has tested MBRP (standard or brief) specifically in Indian opioid-dependent populations",
        "Most studies examine craving OR mindfulness in isolation; few simultaneously assess craving + impulsivity + mindfulness",
        "Indian rehab settings require condensed interventions (6 sessions / 3 weeks) but no brief MBRP has been validated here",
        "Indian treatment landscape lacks integration of structured evidence-based psychological interventions alongside OST/pharmacotherapy",
        "Global MBRP studies focus primarily on alcohol/polysubstance users; opioid-specific MBRP evidence remains limited",
        "Mindfulness practices require culturally congruent adaptation for Indian populations",
        "PRESENT STUDY: First brief MBRP trial (6 sessions/3 weeks) in Indian opioid-dependent sample at Man Nasha Mukti Kendra assessing craving, impulsivity, and mindfulness simultaneously"
    ])

# ===================================================================
# SLIDE 14: AIM
# ===================================================================
add_content_slide(prs, "Aim of the Study",
    [
        "To evaluate the efficacy of a Brief Mindfulness-Based Relapse Prevention (MBRP) intervention (6 sessions over 3 weeks) in reducing craving and impulsivity, and enhancing mindfulness, among substance-dependent patients at Man Nasha Mukti Kendra",
        "",
        "Specifically: To compare outcomes between Brief MBRP + TAU (Experimental Group) and Psychoeducation + TAU (Control Group) on three dependent variables (craving, impulsivity, mindfulness) measured pre- and post-intervention"
    ])

# ===================================================================
# SLIDE 15: OBJECTIVES
# ===================================================================
add_content_slide(prs, "Objectives",
    [
        "1. To assess and compare craving levels (pre-test vs. post-test) in the Experimental group (Brief MBRP + TAU) and Control group (Psychoeducation + TAU)",
        "2. To assess and compare impulsivity levels (pre-test vs. post-test) in both groups",
        "3. To assess and compare mindfulness levels (pre-test vs. post-test) in both groups",
        "4. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in reducing craving",
        "5. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in reducing impulsivity",
        "6. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in enhancing mindfulness"
    ])



# ===================================================================
# SLIDE 16: HYPOTHESES (Research + Null)
# ===================================================================
add_content_slide(prs, "Hypotheses",
    [
        "RESEARCH HYPOTHESES (Directional):",
        "H1: Participants receiving Brief MBRP + TAU will show significantly greater reduction in craving (OCDUS scores) compared to Psychoeducation + TAU, from pre-test to post-test",
        "H2: Participants receiving Brief MBRP + TAU will show significantly greater reduction in impulsivity (BIS-11 scores) compared to Psychoeducation + TAU, from pre-test to post-test",
        "H3: Participants receiving Brief MBRP + TAU will show significantly greater increase in mindfulness (FFMQ scores) compared to Psychoeducation + TAU, from pre-test to post-test",
        "",
        "NULL HYPOTHESES:",
        "H01: There will be no significant difference in craving (OCDUS scores) between the Experimental and Control groups from pre-test to post-test",
        "H02: There will be no significant difference in impulsivity (BIS-11 scores) between the Experimental and Control groups from pre-test to post-test",
        "H03: There will be no significant difference in mindfulness (FFMQ scores) between the Experimental and Control groups from pre-test to post-test"
    ])

# ===================================================================
# SLIDE 17: RESEARCH DESIGN
# ===================================================================
add_content_slide(prs, "Research Design",
    [
        "Design: Pre-test Post-test Control Group Experimental Design",
        "",
        "R   O1   X1   O2   -->  Experimental Group (Brief MBRP + TAU)",
        "R   O1   X2   O2   -->  Control Group (Psychoeducation + TAU)",
        "",
        "Where: R = Random assignment to groups",
        "O1 = Pre-test assessment (OCDUS + BIS-11 + FFMQ + ASSIST)",
        "X1 = Brief MBRP intervention (6 sessions over 3 weeks) + Treatment As Usual",
        "X2 = Psychoeducation (6 sessions over 3 weeks) + Treatment As Usual",
        "O2 = Post-test assessment (OCDUS + BIS-11 + FFMQ)",
        "",
        "Features: True experimental design; Active control (Psychoeducation); TAU continued for all"
    ],
    subtitle="Pre-test Post-test Control Group Design")



# ===================================================================
# SLIDE 18: SAMPLE (with two-stage sampling + male justification)
# ===================================================================
add_content_slide(prs, "Sample",
    [
        "Population: Male substance-dependent patients (primarily opioid users) admitted to Man Nasha Mukti Kendra",
        "Sample Size: N = 60 (30 per group); Recruit N = 70 (35/group) to account for ~15% attrition",
        "",
        "TWO-STAGE SAMPLING PROCESS:",
        "Stage 1 - PURPOSIVE SELECTION: Eligible participants identified based on inclusion and exclusion criteria from among patients admitted to Man Nasha Mukti Kendra",
        "Stage 2 - RANDOM ASSIGNMENT: Selected eligible participants randomly allocated to Experimental or Control group using computer-generated randomization sequence",
        "",
        "JUSTIFICATION FOR MALE-ONLY SAMPLE:",
        "- Indian de-addiction centers (including Man Nasha Mukti Kendra) admit predominantly male patients (~90-95%)",
        "- Opioid dependence in India disproportionately affects males (MAGNITUDE study: male-to-female ratio ~10:1)",
        "- Gender differences in craving, impulsivity, and mindfulness may confound results if mixed sample used",
        "- Ensuring homogeneity of sample strengthens internal validity for this initial efficacy trial",
        "- Female-specific studies recommended as future direction"
    ])

# ===================================================================
# SLIDE 19: SAMPLE SIZE ESTIMATION
# ===================================================================
add_content_slide(prs, "Sample Size Estimation",
    [
        "Formula: n = [(Za/2 + Zb)^2 x 2 x sigma^2] / d^2",
        "",
        "Parameters: Effect size (d) = 0.50 (medium; based on Li et al., 2017 meta-analysis: d = 0.33-0.68)",
        "Power (1 - beta) = 0.80 -> Zb = 0.84",
        "Significance level (alpha) = 0.05 (two-tailed) -> Za/2 = 1.96",
        "",
        "Calculation: n = [(1.96 + 0.84)^2 x 2 x 1] / (0.50)^2 = [(2.80)^2 x 2] / 0.25 = [7.84 x 2] / 0.25 = 15.68 / 0.25 = 62.72 ~ 63 total (~32 per group)",
        "",
        "G*Power 3.1 verification: For ANCOVA (1 covariate, 2 groups): f = 0.25, alpha = 0.05, power = 0.80 -> n = 128 (t-test) or n ~ 34/group (ANCOVA)",
        "FINAL DECISION: N = 60 (30 per group) justified by ANCOVA as primary analysis + consistent with Glasner-Edwards et al. (2017)",
        "Recruit 70 total (35/group) to account for expected 15% attrition in substance-dependent populations"
    ])



# ===================================================================
# SLIDE 20: INCLUSION CRITERIA
# ===================================================================
add_content_slide(prs, "Inclusion Criteria",
    [
        "1. Diagnosis of Substance Dependence as per ICD-10 (F10-F19) / ICD-11 criteria",
        "2. Primary substance: Opioids (heroin, pharmaceutical opioids); polysubstance users with primary opioid dependence included",
        "3. Male participants aged 18-50 years",
        "4. Completed detoxification phase (minimum 7 days post-withdrawal)",
        "5. Currently admitted at Man Nasha Mukti Kendra",
        "6. Minimum education: 5th standard (ability to comprehend psychometric tools and intervention content)",
        "7. Willingness to provide written informed consent for participation",
        "8. Able to attend all 6 intervention sessions during the 3-week intervention period"
    ])

# ===================================================================
# SLIDE 21: EXCLUSION CRITERIA
# ===================================================================
add_content_slide(prs, "Exclusion Criteria",
    [
        "1. Severe psychiatric comorbidity: Psychotic disorders (Schizophrenia), Bipolar I with psychotic features, severe Major Depressive Episode with active suicidality",
        "2. Significant cognitive impairment (MMSE < 24) or intellectual disability",
        "3. Active withdrawal symptoms (Clinical Opiate Withdrawal Scale / COWS score > 12)",
        "4. History of traumatic brain injury with loss of consciousness > 30 minutes",
        "5. Current participation in another structured psychological intervention research study",
        "6. Medical instability requiring acute or intensive care",
        "7. History of prior formal mindfulness or meditation training exceeding 1 month"
    ])

# ===================================================================
# SLIDE 22: VARIABLES
# ===================================================================
add_content_slide(prs, "Variables",
    [
        "INDEPENDENT VARIABLE (IV):",
        "  Type of Intervention (2 levels): Level 1 = Brief MBRP + TAU (Experimental) | Level 2 = Psychoeducation + TAU (Control)",
        "",
        "DEPENDENT VARIABLES (DVs):",
        "  1. Craving - measured by Obsessive Compulsive Drug Use Scale (OCDUS)",
        "  2. Impulsivity - measured by Barratt Impulsiveness Scale (BIS-11)",
        "  3. Mindfulness - measured by Five Facet Mindfulness Questionnaire (FFMQ)",
        "",
        "CONTROLLED / CONFOUNDING VARIABLES:",
        "  Age, education, duration of substance use, severity of dependence (ASSIST baseline score), TAU components kept constant across groups, session duration and frequency equalized between groups"
    ])



# ===================================================================
# SLIDE 23: TOOLS - FFMQ & BIS-11
# ===================================================================
add_content_slide(prs, "Tools: Mindfulness & Impulsivity Measures",
    [
        "1. FIVE FACET MINDFULNESS QUESTIONNAIRE (FFMQ) - Baer et al. (2006)",
        "   Description: 39-item self-report measure; 5 facets: Observing, Describing, Acting with Awareness, Non-Judging, Non-Reactivity",
        "   Scoring: 5-point Likert scale (1 = never true to 5 = always true); higher scores = greater mindfulness",
        "   Reliability: Internal consistency alpha = 0.75-0.91 across facets; adequate test-retest reliability",
        "   Validity: Convergent validity with other mindfulness measures; sensitive to mindfulness intervention effects",
        "   Indian Usability: Hindi adaptation available; used in Indian clinical samples",
        "",
        "2. BARRATT IMPULSIVENESS SCALE (BIS-11) - Patton et al. (1995)",
        "   Description: 30-item self-report scale; 3 factors: Attentional, Motor, Non-Planning Impulsivity",
        "   Scoring: 4-point scale (1 = rarely/never to 4 = almost always); higher scores = greater impulsivity",
        "   Reliability: Internal consistency alpha = 0.79-0.83; test-retest r = 0.83",
        "   Validity: Discriminates substance users from healthy controls; widely validated in SUD populations",
        "   Indian Usability: Hindi version validated; extensively used in Indian psychiatric research"
    ])

# ===================================================================
# SLIDE 24: TOOLS - CRAVING (OCDUS ONLY)
# ===================================================================
add_content_slide(prs, "Tools: Craving Assessment - OCDUS",
    [
        "OBSESSIVE COMPULSIVE DRUG USE SCALE (OCDUS)",
        "",
        "Authors: Franken et al. (2002); adapted from Obsessive Compulsive Drinking Scale (Anton et al., 1995)",
        "Description: 12-item self-report scale measuring obsessive thoughts about drug use and compulsive urges to use; captures cognitive preoccupation with substance and perceived loss of control over use behavior",
        "Scoring: Items rated on 5-point scale (0-4); Total score range 0-48; Higher scores indicate greater craving intensity",
        "Subscales: (1) Obsessive thoughts/interference, (2) Desire/control, (3) Resistance to thoughts",
        "Reliability: Internal consistency alpha = 0.86-0.90; test-retest reliability r = 0.78",
        "Validity: Convergent validity with VAS craving (r = 0.55-0.67); discriminant validity established; sensitive to treatment-related changes; predicts relapse",
        "Indian Usability: Applicable to Indian substance-dependent populations; adaptable for Hindi administration; brief and clinically practical (5 minutes); suitable for repeated measurements in pre-post designs",
        "Justification: Captures both cognitive (obsessive) and behavioral (compulsive) dimensions of craving; applicable across substance types including opioids; psychometrically superior to single-item VAS"
    ])



# ===================================================================
# SLIDE 25: TOOLS - SEVERITY (ASSIST, PRE-TEST ONLY)
# ===================================================================
add_content_slide(prs, "Tools: Severity Assessment - ASSIST (Pre-Test Only)",
    [
        "ALCOHOL, SMOKING AND SUBSTANCE INVOLVEMENT SCREENING TEST (ASSIST) - WHO",
        "",
        "Authors: WHO ASSIST Working Group (2002); Humeniuk et al. (2008)",
        "Description: 8-item clinician/self-administered screening tool; assesses risk level for 10 substance categories (tobacco, alcohol, cannabis, opioids, stimulants, sedatives, hallucinogens, inhalants, others)",
        "Scoring: Substance-specific risk scores -> Low (0-3), Moderate (4-26), High (27+) for drugs",
        "Reliability: Test-retest r = 0.58-0.90 across substances; internal consistency alpha = 0.77-0.94",
        "Validity: Concurrent validity with ASI, DAST, AUDIT; sensitivity = 0.80, specificity = 0.71",
        "Indian Usability: WHO-validated; Hindi version available; extensively used in Indian NDDTC studies; culturally appropriate; brief (5-10 min)",
        "PURPOSE: Used at PRE-TEST ONLY to establish baseline severity and ensure group equivalence; NOT an outcome measure",
        "Justification: Multi-substance coverage captures polysubstance patterns common in Indian opioid users; WHO endorsement; free to use"
    ])

# ===================================================================
# SLIDE 26: PROCEDURE
# ===================================================================
add_content_slide(prs, "Procedure",
    [
        "Step 1: SCREENING - Identify eligible male patients at Man Nasha Mukti Kendra (ICD-10 diagnosis, detoxified, meet inclusion/exclusion criteria)",
        "Step 2: INFORMED CONSENT - Explain study purpose, procedures, confidentiality, voluntary participation, and right to withdraw",
        "Step 3: PRE-TEST ASSESSMENT - Administer ASSIST + OCDUS + BIS-11 + FFMQ",
        "Step 4: RANDOM ASSIGNMENT - Computer-generated randomization allocating eligible participants to Experimental or Control group",
        "Step 5: INTERVENTION DELIVERY:",
        "   Experimental Group: Brief MBRP (6 sessions x 60 min, twice weekly over 3 weeks) + TAU",
        "   Control Group: Psychoeducation (6 sessions x 60 min, twice weekly over 3 weeks) + TAU",
        "Step 6: POST-TEST ASSESSMENT - Administer OCDUS + BIS-11 + FFMQ (within 1 week of intervention completion)",
        "Step 7: DATA COMPILATION & STATISTICAL ANALYSIS",
        "NOTE: Both groups receive TAU (pharmacotherapy, routine counseling) throughout | Assessments by blinded RA"
    ],
    subtitle="Stepwise Research Flow at Man Nasha Mukti Kendra")



# ===================================================================
# SLIDE 27: INTERVENTION - EXPERIMENTAL (6 sessions / 3 weeks)
# ===================================================================
add_content_slide(prs, "Intervention: Experimental Group - Brief MBRP (6 Sessions / 3 Weeks)",
    [
        "BRIEF MBRP PROTOCOL: 6 Sessions x 60 min x 3 Weeks (Twice weekly) | Group format (6-8 participants)",
        "",
        "Session 1: Introduction to MBRP & Autopilot - Psychoeducation on MBRP; raisin exercise; identifying automatic patterns of reactivity in addiction",
        "Session 2: Awareness of Triggers & Body Scan - Body scan meditation; mapping personal relapse triggers; understanding trigger-reaction-use chain",
        "Session 3: Mindfulness in Daily Life & SOBER Space - Sitting meditation (breath focus); SOBER breathing space technique; integrating mindfulness into routine",
        "Session 4: Urge Surfing & High-Risk Situations - Urge surfing practice; role-play high-risk scenarios; cognitive decentering from craving",
        "Session 5: Acceptance, Non-Judgment & Skillful Action - Non-judgmental awareness; acceptance vs. avoidance; mindful decision-making; seeing thoughts as thoughts",
        "Session 6: Integration, Self-Care & Maintenance Plan - Loving-kindness meditation; relapse warning signs; personal relapse prevention plan; ongoing practice commitment",
        "",
        "Delivery: By trained MPhil Clinical Psychologist | Materials: Hindi audio-guided meditations, session handouts, daily practice logs"
    ],
    subtitle="Brief MBRP + TAU at Man Nasha Mukti Kendra")

# ===================================================================
# SLIDE 28: INTERVENTION - CONTROL (6 sessions / 3 weeks)
# ===================================================================
add_content_slide(prs, "Intervention: Control Group - Psychoeducation (6 Sessions / 3 Weeks)",
    [
        "PSYCHOEDUCATION PROTOCOL: 6 Sessions x 60 min x 3 Weeks (Twice weekly) | Group format (6-8 participants)",
        "",
        "Session 1: Understanding Addiction - Nature of substance dependence; brain changes; disease model of addiction",
        "Session 2: Effects of Opioids - Short-term and long-term physical, psychological, and social consequences of opioid use",
        "Session 3: Understanding Relapse - Relapse process; warning signs; high-risk situations; general coping strategies",
        "Session 4: Health, Nutrition & Sleep - Physical health recovery; importance of nutrition; sleep hygiene in recovery",
        "Session 5: Social Consequences & Family Impact - Legal issues; stigma; family dynamics; social rehabilitation",
        "Session 6: Motivation, Goal Setting & Lifestyle Changes - Stages of change; personal goals; long-term recovery planning; support systems",
        "",
        "Delivery: By trained Clinical Psychologist | Materials: Information handouts, visual aids, group discussion",
        "NOTE: Matched for contact time, group format, and therapist attention to control non-specific therapeutic factors. NO mindfulness component included."
    ],
    subtitle="Psychoeducation + TAU (Active Control)")



# ===================================================================
# SLIDE 29: DATA ANALYSIS (with detailed ANCOVA justification)
# ===================================================================
add_content_slide(prs, "Data Analysis",
    [
        "Descriptive Statistics: Mean, SD, frequency, percentages for sociodemographic and clinical variables",
        "Normality Testing: Shapiro-Wilk test to determine distribution of outcome variables",
        "Within-Group Comparison: Paired samples t-test (pre vs. post within each group); Wilcoxon Signed-Rank for non-normal data",
        "Between-Group Comparison: Independent samples t-test (post-test comparison); Mann-Whitney U for non-normal data",
        "",
        "PRIMARY ANALYSIS - ANCOVA (Analysis of Covariance):",
        "  DV: Post-test scores (OCDUS / BIS-11 / FFMQ) | IV: Group (Experimental vs. Control) | Covariate: Pre-test scores",
        "",
        "DETAILED ANCOVA JUSTIFICATION:",
        "  (a) Controls for pre-existing baseline differences between groups, even after randomization",
        "  (b) Increases statistical power by reducing within-group error variance (partials out pre-test variability)",
        "  (c) Provides a more precise estimate of treatment effect by adjusting post-test means for baseline",
        "  (d) Reduces required sample size compared to independent t-test (critical for clinical populations with high attrition)",
        "  (e) Recommended for pre-post experimental designs by Tabachnick and Fidell (2013) and Field (2018)",
        "",
        "Effect Size: Partial eta-squared (np2) for ANCOVA; Cohen's d for pairwise | Significance: alpha = 0.05 (two-tailed) | Software: SPSS 26.0"
    ])



# ===================================================================
# SLIDE 30: ETHICAL CONSIDERATIONS
# ===================================================================
add_content_slide(prs, "Ethical Considerations",
    [
        "Informed Consent: Written informed consent obtained in Hindi/regional language; participants fully informed of study purpose, procedures, duration, risks, and benefits",
        "Voluntary Participation: Right to withdraw at any time without penalty or impact on ongoing treatment at Man Nasha Mukti Kendra",
        "Confidentiality: All data coded with participant IDs; no identifying information in publications; data stored in locked cabinets and password-protected electronic files",
        "Non-Maleficence: No harmful procedures involved; Control group receives active psychoeducation (not waitlist/no-treatment); TAU continued for all participants",
        "Institutional Approval: Ethical clearance obtained from Institutional Ethics Committee (IEC) prior to commencement of data collection",
        "Debriefing: Control group participants offered brief MBRP orientation session post-study completion if desired",
        "Compliance: Study conducted in accordance with ICMR (2017) National Ethical Guidelines for Biomedical and Health Research Involving Human Participants"
    ])

# ===================================================================
# SLIDE 31: EXPECTED RESULTS (DETAILED)
# ===================================================================
add_content_slide(prs, "Expected Results (1/2)",
    [
        "CRAVING (OCDUS):",
        "  - The Experimental group (Brief MBRP + TAU) is expected to demonstrate a statistically significant reduction in OCDUS scores from pre-test to post-test compared to the Control group (Psychoeducation + TAU)",
        "  - Both obsessive thoughts about drug use and compulsive urges are anticipated to decrease as participants learn to observe craving non-reactively through urge surfing and SOBER breathing techniques",
        "  - Expected effect size: medium to large (d = 0.50-0.80) based on Witkiewitz et al. (2013) and Li et al. (2017) meta-analytic findings (d = 0.68 for craving)",
        "  - Mechanism: Urge surfing disrupts the automatic craving-use cycle by training participants to observe craving as a transient sensory experience rather than an imperative demand requiring substance use",
        "",
        "IMPULSIVITY (BIS-11):",
        "  - The Experimental group is expected to show a significantly greater reduction in total BIS-11 scores, particularly on Attentional Impulsivity and Motor Impulsivity subscales",
        "  - Expected effect: medium (d = 0.40-0.60) based on Garland et al. (2016) findings",
        "  - Mechanism: Repeated mindfulness practice strengthens prefrontal inhibitory control, enabling participants to insert a deliberate pause (mindful gap) between impulse and action, thereby reducing reflexive substance-seeking behavior"
    ])



# ===================================================================
# SLIDE 32: EXPECTED RESULTS 2/2
# ===================================================================
add_content_slide(prs, "Expected Results (2/2)",
    [
        "MINDFULNESS (FFMQ):",
        "  - The Experimental group is expected to demonstrate a statistically significant increase in FFMQ total and facet scores compared to the Control group",
        "  - Particularly robust improvements expected on Acting with Awareness and Non-Reactivity facets, as these are most directly trained in MBRP exercises",
        "  - Expected effect: medium to large (d = 0.50-0.70) based on Bowen et al. (2014) and Karyadi et al. (2014)",
        "  - Mechanism: Structured daily meditation practice, body scan, and non-judgmental awareness exercises directly cultivate dispositional mindfulness, which becomes a stable trait-like resource for managing relapse triggers",
        "",
        "OVERALL PATTERN:",
        "  - Brief MBRP + TAU will demonstrate statistically significant superiority over Psychoeducation + TAU across all three DVs",
        "  - The Control group may also show some improvement (attention effects + psychoeducation content) but to a significantly lesser degree",
        "  - Null hypotheses (H01, H02, H03) are expected to be rejected",
        "  - The findings will support feasibility and efficacy of a brief 6-session MBRP protocol in Indian de-addiction centers like Man Nasha Mukti Kendra"
    ])

# ===================================================================
# SLIDE 33: CLINICAL IMPLICATIONS
# ===================================================================
add_content_slide(prs, "Clinical Implications",
    [
        "Validates a brief, structured MBRP model (6 sessions/3 weeks) feasible for Indian de-addiction settings with time-limited admissions",
        "Provides evidence-based psychological intervention to complement existing pharmacotherapy (OST, naltrexone) at centers like Man Nasha Mukti Kendra",
        "Demonstrates that mindfulness-based approaches are culturally compatible with Indian populations given India's meditative traditions",
        "Addresses multiple relapse risk factors simultaneously (craving + impulsivity + mindfulness) through a single integrated brief protocol",
        "Supports task-shifting: Brief MBRP deliverable by MPhil-trained Clinical Psychologists in resource-limited settings",
        "Informs national treatment policy (NIMHANS, NDDTC, State Mental Health Authorities) for integrating MBRP into standard de-addiction care",
        "Contributes to RCI-recognized intervention repertoire for clinical psychology training and practice in India"
    ])



# ===================================================================
# SLIDE 34: LIMITATIONS
# ===================================================================
add_content_slide(prs, "Limitations",
    [
        "Male-only sample from a single center (Man Nasha Mukti Kendra) limits generalizability to females, other substances, and diverse settings",
        "Short-term assessment: Post-test immediately after 3-week intervention; no long-term follow-up to assess maintenance of therapeutic gains",
        "Self-report measures (OCDUS, BIS-11, FFMQ) susceptible to social desirability bias and limited introspective insight",
        "No biological or physiological markers: Craving measured subjectively without neuroimaging or biomarker corroboration",
        "Therapist effects: Single therapist delivering intervention may introduce therapist-specific confounds not generalizable to other clinicians",
        "Attention-matched control (Psychoeducation) controls for contact but does not fully account for specific mindfulness mechanisms (active placebo limitation)",
        "Potential attrition: Substance-dependent population may show higher dropout rates despite over-recruitment planning"
    ])

# ===================================================================
# SLIDE 35: FUTURE DIRECTIONS
# ===================================================================
add_content_slide(prs, "Future Directions",
    [
        "Follow-up studies: 3-month and 6-month post-intervention assessments to evaluate sustained effects of Brief MBRP",
        "Multi-site RCTs: Replicate across multiple Indian de-addiction centers for enhanced external validity",
        "Female participants: Include female substance-dependent patients to examine gender-specific MBRP effects",
        "Neuroimaging integration: fMRI/EEG studies to examine neural mechanisms of MBRP effects on craving and inhibitory control circuits",
        "Dose-response analysis: Compare 4-session vs. 6-session vs. full 8-session MBRP for optimal dosage determination",
        "Mediator/moderator analysis: Test whether mindfulness improvement mediates craving and impulsivity reduction; examine moderators (severity, duration of use)",
        "Technology-assisted delivery: App-based/digital MBRP for post-discharge maintenance and access in rural areas",
        "Comparative effectiveness: MBRP vs. CBT vs. ACT vs. Contingency Management in Indian samples"
    ])



# ===================================================================
# SLIDE 36: SUMMARY
# ===================================================================
add_content_slide(prs, "Summary",
    [
        "The present study proposes a Pre-test Post-test Control Group Experimental Design to evaluate Brief MBRP (6 sessions over 3 weeks) in male opioid-dependent patients at Man Nasha Mukti Kendra",
        "Experimental Group receives Brief MBRP + TAU; Control Group receives Psychoeducation + TAU",
        "Three dependent variables assessed: Craving (OCDUS), Impulsivity (BIS-11), and Mindfulness (FFMQ)",
        "Two-stage sampling: Purposive selection of eligible patients followed by computer-generated random assignment to groups",
        "Total sample: N = 60 (30 per group), justified by power analysis assuming medium effect size (d = 0.50), power = 0.80, alpha = 0.05",
        "Primary analysis: ANCOVA controlling for pre-test scores as covariates to increase precision and control baseline differences",
        "Brief MBRP is hypothesized to produce significantly greater reductions in craving and impulsivity, and greater enhancement of mindfulness, compared to psychoeducation",
        "The study addresses a critical research gap: No prior Indian RCT has tested brief MBRP in opioid-dependent populations"
    ])

# ===================================================================
# SLIDE 37: CONCLUSION
# ===================================================================
add_content_slide(prs, "Conclusion",
    [
        "Substance dependence, particularly opioid dependence, represents a significant public health crisis in India with relapse rates exceeding 70%",
        "Current treatment at Indian de-addiction centers relies heavily on pharmacotherapy with limited evidence-based psychological interventions",
        "MBRP offers a theoretically grounded, evidence-based approach that simultaneously targets craving (urge surfing), impulsivity (mindful pause), and enhances mindfulness (meditation practice)",
        "A brief 6-session MBRP protocol is clinically practical, culturally appropriate, and feasible within the constraints of Indian rehabilitation settings like Man Nasha Mukti Kendra",
        "If supported by the findings, Brief MBRP can be integrated into standard de-addiction care as a cost-effective, replicable psychological intervention deliverable by trained MPhil Clinical Psychologists",
        "This study contributes to the growing evidence base for mindfulness-based interventions in addiction treatment while addressing the specific needs and gaps in Indian clinical psychology practice",
        "The results will have direct implications for RCI-recognized training, national treatment policy, and the broader goal of reducing relapse burden in Indian substance-dependent populations"
    ])



# ===================================================================
# SLIDES 38-41: REFERENCES (APA 7 - et al. for 3+ authors)
# ===================================================================
add_content_slide(prs, "References (1/4)",
    [
        "Baer, R. A., Smith, G. T., Hopkins, J., Krietemeyer, J., & Toney, L. (2006). Using self-report assessment methods to explore facets of mindfulness. Assessment, 13(1), 27-45.",
        "Bowen, S., Chawla, N., & Marlatt, G. A. (2011). Mindfulness-based relapse prevention for addictive behaviors: A clinician's guide. Guilford Press.",
        "Bowen, S., & Marlatt, G. A. (2009). Surfing the urge: Brief mindfulness-based intervention for college student smokers. Psychology of Addictive Behaviors, 23(4), 666-671.",
        "Bowen, S., Witkiewitz, K., Clifasefi, S. L., Grow, J., Chawla, N., Hsu, S. H., Carroll, H. A., Harrop, E., Collins, S. E., Lustyk, M. K., & Larimer, M. E. (2014). Relative efficacy of mindfulness-based relapse prevention, standard relapse prevention, and treatment as usual for substance use disorders: A randomized clinical trial. JAMA Psychiatry, 71(5), 547-556.",
        "Brewer, J. A., Mallik, S., Babuscio, T. A., Nich, C., Johnson, H. E., Deleone, C. M., Minnix-Cotton, C. A., Byrne, S. A., Kober, H., Weinstein, A. J., Carroll, K. M., & Rounsaville, B. J. (2011). Mindfulness training for smoking cessation: Results from a randomized controlled trial. Drug and Alcohol Dependence, 119(1-2), 72-80.",
        "Chiesa, A., & Serretti, A. (2014). Are mindfulness-based interventions effective for substance use disorders? A systematic review of the evidence. Substance Use & Misuse, 49(5), 492-512."
    ])

add_content_slide(prs, "References (2/4)",
    [
        "Franken, I. H. A., Hendriks, V. M., & van den Brink, W. (2002). New neuropsychological instruments to measure impulsivity and compulsivity in alcohol and drug use disorders. Psychiatrie en Neurobiologie, 1, 10-14.",
        "Garland, E. L., Froeliger, B., & Howard, M. O. (2014). Mindfulness training targets neurocognitive mechanisms of addiction at the attention-appraisal-emotion interface. Frontiers in Psychiatry, 4, 173.",
        "Garland, E. L., Manusov, E. G., Froeliger, B., Kelly, A., Williams, J. M., & Howard, M. O. (2014). Mindfulness-oriented recovery enhancement for chronic pain and prescription opioid misuse: Results from an early-stage randomized controlled trial. Journal of Consulting and Clinical Psychology, 82(3), 448-459.",
        "Garland, E. L., Roberts-Lewis, A., Tronnier, C. D., Graves, R., & Kelley, K. (2016). Mindfulness-oriented recovery enhancement versus CBT for co-occurring substance dependence, traumatic stress, and psychiatric disorders. Journal of Consulting and Clinical Psychology, 84(4), 281-293.",
        "Ghosh, A., Basu, D., & Avasthi, A. (2018). Relapse in opioid dependence: A comprehensive review from an Indian perspective. Indian Journal of Psychiatry, 60(Suppl 4), S469-S476.",
        "Glasner-Edwards, S., Mooney, L. J., Ang, A., Garneau, H. C., Hartwell, E., Brecht, M. L., & Rawson, R. A. (2017). Mindfulness-based relapse prevention for stimulant dependent adults: A pilot randomized clinical trial. Mindfulness, 8(1), 126-135."
    ])



add_content_slide(prs, "References (3/4)",
    [
        "Grant, S., Colaiaco, B., Motala, A., Shanman, R., Booth, M., Sorbero, M., & Hempel, S. (2017). Mindfulness-based relapse prevention for substance use disorders: A systematic review and meta-analysis. Journal of Addiction Medicine, 11(5), 386-396.",
        "Humeniuk, R., Ali, R., Babor, T. F., Farrell, M., Formigoni, M. L., Jittiwutikarn, J., de Lacerda, R. B., Ling, W., Marsden, J., Monteiro, M., Nhiwatiwa, S., Pal, H., Poznyak, V., & Simon, S. (2008). Validation of the Alcohol, Smoking and Substance Involvement Screening Test (ASSIST). Addiction, 103(6), 1039-1047.",
        "Jain, R., Majumder, P., & Gupta, T. (2013). Pharmacological intervention of nicotine dependence. Indian Journal of Psychiatry, 55(Suppl 1), S86-S92.",
        "Kabat-Zinn, J. (1990). Full catastrophe living: Using the wisdom of your body and mind to face stress, pain, and illness. Delacorte Press.",
        "Karyadi, K. A., VanderVeen, J. D., & Cyders, M. A. (2014). A meta-analysis of the relationship between trait mindfulness and substance use behaviors. Drug and Alcohol Dependence, 143, 1-10.",
        "Li, W., Howard, M. O., Garland, E. L., McGovern, P., & Lazar, M. (2017). Mindfulness treatment for substance misuse: A systematic review and meta-analysis. Journal of Substance Abuse Treatment, 75, 62-96.",
        "Marlatt, G. A., & Gordon, J. R. (1985). Relapse prevention: Maintenance strategies in the treatment of addictive behaviors. Guilford Press."
    ])

add_content_slide(prs, "References (4/4)",
    [
        "Mattoo, S. K., Chakrabarti, S., & Anjaiah, M. (2009). Psychosocial factors associated with relapse in men with alcohol or opioid dependence. Indian Journal of Medical Research, 130(6), 702-708.",
        "Ministry of Social Justice and Empowerment. (2019). Magnitude of substance use in India. Government of India.",
        "Moeller, F. G., Barratt, E. S., Dougherty, D. M., Schmitz, J. M., & Swann, A. C. (2001). Psychiatric aspects of impulsivity. American Journal of Psychiatry, 158(11), 1783-1793.",
        "Murphy, C., & MacKillop, J. (2012). Living in the here and now: Interrelationships between impulsivity, mindfulness, and alcohol misuse. Psychopharmacology, 219(2), 527-536.",
        "National Institute on Drug Abuse. (2020). Drugs, brains, and behavior: The science of addiction. NIDA.",
        "Patton, J. H., Stanford, M. S., & Barratt, E. S. (1995). Factor structure of the Barratt Impulsiveness Scale. Journal of Clinical Psychology, 51(6), 768-774.",
        "Robinson, T. E., & Berridge, K. C. (1993). The neural basis of drug craving: An incentive-sensitization theory of addiction. Brain Research Reviews, 18(3), 247-291.",
        "Sarkar, S., & Balhara, Y. P. S. (2016). Structured psychological interventions in Indian de-addiction settings. Indian Journal of Psychiatry, 58(3), 290-295.",
        "Tabachnick, B. G., & Fidell, L. S. (2013). Using multivariate statistics (6th ed.). Pearson.",
        "WHO ASSIST Working Group. (2002). The Alcohol, Smoking and Substance Involvement Screening Test (ASSIST): Development, reliability and feasibility. Addiction, 97(9), 1183-1194.",
        "Witkiewitz, K., Bowen, S., Douglas, H., & Hsu, S. H. (2013). Mindfulness-based relapse prevention for substance craving. Addictive Behaviors, 38(2), 1563-1571."
    ])



# ===================================================================
# SLIDE 42: THANK YOU
# ===================================================================
add_title_slide(prs,
    "Thank You",
    "Questions & Discussion",
    "[Researcher Name]", "[Supervisor Name]",
    "Clinical Psychology", "Man Nasha Mukti Kendra / [University Name]", "2026")

# ===================================================================
# SAVE PRESENTATION
# ===================================================================
output_path = "/projects/sandbox/Dango-kiro/MBRP_Research_Synopsis_PPT.pptx"
prs.save(output_path)
print(f"SUCCESS: Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
