#!/usr/bin/env python3
"""
MET Comprehensive Presentation v4.2
- Keeps ALL original v4 content (62 slides)
- FIXES overlapping with strict zone layout
- ADDS ~18 new slides: therapy comparisons, family/caregiver, comorbid disorders
- Target: ~80 slides total
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# LAYOUT RULES (strict zones - NO overlaps)
# Title bar:     y=0    to y=1.0    (height=1.0)
# Content zone:  y=1.15 to y=6.4    (5.25 inches)
# Reference bar: y=6.55 to y=7.1    (height=0.55)
# ============================================================

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_H = Inches(1.0)
REF_Y = Inches(6.55)
REF_H = Inches(0.55)

# X layout constants
LM = Inches(0.4)           # left margin
FULL_W = Inches(12.5)      # full width
HALF_W = Inches(6.05)      # half width column
COL2 = Inches(6.75)        # second column X
THIRD_W = Inches(3.95)     # third width
COL3_2 = Inches(4.65)      # 2nd of 3 columns
COL3_3 = Inches(8.9)       # 3rd of 3 columns

# Row positions for 2-row layouts (no overlap guaranteed)
ROW1 = Inches(1.15)        # first content row
ROW2 = Inches(3.75)        # second content row
ROW1_H = Inches(2.4)       # row 1 box height
ROW2_H = Inches(2.5)       # row 2 box height

COLORS = {
    'deep_blue': RGBColor(0x1B, 0x3A, 0x5C),
    'teal': RGBColor(0x00, 0x7B, 0x83),
    'green': RGBColor(0x2E, 0x7D, 0x32),
    'orange': RGBColor(0xE6, 0x5C, 0x00),
    'purple': RGBColor(0x6A, 0x1B, 0x9A),
    'red': RGBColor(0xC6, 0x28, 0x28),
    'gold': RGBColor(0xF9, 0xA8, 0x25),
    'light_blue': RGBColor(0xBB, 0xDE, 0xFB),
    'light_green': RGBColor(0xC8, 0xE6, 0xC9),
    'light_purple': RGBColor(0xE1, 0xBE, 0xE7),
    'light_orange': RGBColor(0xFF, 0xE0, 0xB2),
    'light_teal': RGBColor(0xB2, 0xDF, 0xDB),
    'light_red': RGBColor(0xFF, 0xCD, 0xD2),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'cream': RGBColor(0xFF, 0xF8, 0xE1),
    'navy': RGBColor(0x0D, 0x47, 0xA1),
    'dark_green': RGBColor(0x1B, 0x5E, 0x20),
    'maroon': RGBColor(0x88, 0x00, 0x38),
}

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]



# ============================================================
# HELPER FUNCTIONS (with strict zone enforcement)
# ============================================================

def bg(slide, c1, c2=None):
    """Gradient or solid background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    if c2:
        shape.fill.gradient()
        shape.fill.gradient_stops[0].color.rgb = COLORS[c1]
        shape.fill.gradient_stops[1].color.rgb = COLORS[c2]
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS[c1]
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def title_bar(slide, text, color='deep_blue'):
    """Title bar: fixed at y=0, h=1.0."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, TITLE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[color]
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Times New Roman'

def ref_bar(slide, text):
    """Reference bar: fixed at y=6.55, h=0.55."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), REF_Y, SLIDE_W, REF_H)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p.font.name = 'Times New Roman'

def cbox(slide, x, y, w, h, lines, bgc='white', tc='dark_gray', fs=13, title=None, ttc='deep_blue', border=None):
    """Content box with text. All Y values must be between 1.15 and 6.4."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS[bgc]
    if border:
        box.line.color.rgb = COLORS[border]
        box.line.width = Pt(1.5)
    else:
        box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.05)
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS[ttc]
        p.font.name = 'Times New Roman'
        p.space_after = Pt(3)
    for i, line in enumerate(lines):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fs)
        p.font.color.rgb = COLORS[tc]
        p.font.name = 'Times New Roman'
        p.space_after = Pt(2)

def divider(title, subtitle, color='deep_blue'):
    """Section divider slide."""
    slide = prs.slides.add_slide(blank_layout)
    bg(slide, color, 'white')
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = 'Times New Roman'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLORS['cream']
    p2.font.name = 'Times New Roman'
    p2.alignment = PP_ALIGN.CENTER

def tbl(title, headers, rows, color='deep_blue', ref=''):
    """Table slide."""
    slide = prs.slides.add_slide(blank_layout)
    bg(slide, 'white', 'light_blue')
    title_bar(slide, title, color)
    nr, nc = len(rows)+1, len(headers)
    table = slide.shapes.add_table(nr, nc, LM, Inches(1.15), FULL_W, Inches(5.1)).table
    cw = int(FULL_W / nc)
    for i in range(nc):
        table.columns[i].width = cw
    for i, h in enumerate(headers):
        c = table.cell(0, i)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = COLORS[color]
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COLORS['white']
            p.font.name = 'Times New Roman'
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = table.cell(ri+1, ci)
            c.text = val
            if ri % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = COLORS['dark_gray']
                p.font.name = 'Times New Roman'
    if ref:
        ref_bar(slide, ref)

def flow(slide, y, items, colors):
    """Horizontal process flow."""
    n = len(items)
    bw = (FULL_W - Inches(0.25)*(n-1)) / n
    for i, (item, ck) in enumerate(zip(items, colors)):
        x = LM + (bw + Inches(0.25)) * i
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), y, int(bw), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS[ck]
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = 'Times New Roman'
        p.alignment = PP_ALIGN.CENTER
        if i < n-1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, int(x+bw), y+Inches(0.3), Inches(0.25), Inches(0.35))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLORS['gold']
            arr.line.fill.background()



# ============================================================
# TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
bg(slide, 'deep_blue', 'navy')
tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10), Inches(4.5))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "MOTIVATIONAL ENHANCEMENT THERAPY (MET)"; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = COLORS['white']; p.font.name = 'Times New Roman'; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "A Comprehensive Clinical Guide"; p2.font.size = Pt(22); p2.font.color.rgb = COLORS['gold']; p2.font.name = 'Times New Roman'; p2.alignment = PP_ALIGN.CENTER; p2.space_after = Pt(30)
p3 = tf.add_paragraph(); p3.text = "Based on the NIAAA Project MATCH MET Manual (Miller et al., 1992)"; p3.font.size = Pt(14); p3.font.color.rgb = COLORS['light_blue']; p3.font.name = 'Times New Roman'; p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph(); p4.text = "& NIMHANS Substance Use Disorders Manual (2016)"; p4.font.size = Pt(14); p4.font.color.rgb = COLORS['light_blue']; p4.font.name = 'Times New Roman'; p4.alignment = PP_ALIGN.CENTER
p5 = tf.add_paragraph(); p5.text = ""; p5.space_after = Pt(20)
p6 = tf.add_paragraph(); p6.text = "4-Session Brief Intervention  |  Evidence-Based  |  Client-Centered"; p6.font.size = Pt(13); p6.font.color.rgb = COLORS['gold']; p6.font.name = 'Times New Roman'; p6.alignment = PP_ALIGN.CENTER
ref_bar(slide, "Miller, W.R., Zweben, A., DiClemente, C.C., & Rychtarik, R.G. (1992). MET Manual. NIAAA Project MATCH Series Vol. 2.")

# TABLE OF CONTENTS
slide = prs.slides.add_slide(blank_layout)
bg(slide, 'white', 'light_blue')
title_bar(slide, "PRESENTATION OUTLINE", 'deep_blue')
secs = [("1. Introduction to MET","What is MET, origins, how it differs"),("2. Theoretical Foundations","TTM, Self-Efficacy, Cognitive Dissonance"),
    ("3. Principles & Spirit","FRAMES, OARS, DARES, MI Spirit, Handling Resistance"),("4. Session-by-Session Guide","4-session protocol with dialogues"),
    ("5. Clinical Techniques","Reflective listening, decisional balance, rulers, traps"),("6. Case Conceptualization","Applied case with session-wise plans"),
    ("7. Worksheets & Tools","6 printable clinical worksheets"),("8. Comparison with Other Therapies","MET vs CBT, 12-Step, BSFT; family/caregiver"),
    ("9. Comorbid Psychiatric Disorders","PD, depression, anxiety, PTSD, psychosis"),("10. Research & Effectiveness","MATCH, UKATT, meta-analyses, Indian research")]
clrs = ['light_blue','light_teal','light_green','light_purple','light_orange','light_red','light_blue','light_green','light_purple','light_orange']
tclrs = ['deep_blue','teal','green','purple','orange','maroon','deep_blue','green','purple','orange']
for i,(t,d) in enumerate(secs):
    cbox(slide, LM, Inches(1.15)+Inches(0.52)*i, FULL_W, Inches(0.48), [d], bgc=clrs[i], fs=11, title=t, ttc=tclrs[i], border=tclrs[i])
ref_bar(slide, "Structure: Miller et al. (1992). MET Manual; NIMHANS (2016). Substance Use Disorders.")



# ============================================================
# SECTION 1: INTRODUCTION
# ============================================================
divider("SECTION 1", "Introduction to Motivational Enhancement Therapy", 'deep_blue')

# What is MET
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'cream'); title_bar(slide, "What is Motivational Enhancement Therapy (MET)?", 'deep_blue')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["MET is a brief, systematic intervention designed","to produce rapid, internally motivated change.","It does NOT guide the client step-by-step through","recovery. Instead, it uses motivational strategies","to mobilize the client's OWN change resources.","","Key: 4 planned sessions over 12 weeks."], bgc='light_blue', title="Definition", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["- Developed for NIAAA Project MATCH (1992)","- Based on Motivational Interviewing principles","- Uses personal feedback + MI techniques","- Designed for alcohol/substance use disorders","- Adaptable to various clinical settings","- Therapist as collaborator, NOT expert"], bgc='light_green', title="Key Facts", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["\"MET is based on principles of motivational psychology and is designed to produce rapid,","internally motivated change. This treatment strategy does not attempt to guide and train the","client, step by step, through recovery, but instead employs motivational strategies to mobilize","the client's own change resources.\" (MET Manual, p. 1)","","Core idea: The motivation and resources for change ALREADY EXIST within the client.","The therapist's job is to create conditions that help these emerge naturally."], bgc='cream', title="From the Manual:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Miller, W.R. et al. (1992). MET Manual, p. 1. NIAAA Project MATCH Monograph Series, Vol. 2.")

# How MET Differs - Table
tbl("How MET Differs from Other Approaches", ["Feature", "Traditional Approaches", "MET Approach"],
    [["Therapist Role", "Expert / Teacher / Director", "Collaborative partner; elicits client's own motivation"],
     ["Client Role", "Passive recipient of treatment", "Active agent of own change"],
     ["View of Resistance", "Denial to be confronted", "Signal to change therapeutic strategy"],
     ["Goals", "Set by therapist or program", "Negotiated with client based on their values"],
     ["Core Techniques", "Skills training, education, confrontation", "Reflective listening, feedback, exploring ambivalence"],
     ["Duration", "Usually 12+ sessions", "4 structured sessions over 12 weeks"],
     ["Change Mechanism", "Teaching new skills/information", "Mobilizing client's own internal resources"]],
    color='green', ref="Ref: Miller et al. (1992). MET Manual, pp. 1-5; NIMHANS (2016). Ch. 8.")

# Origins
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_purple'); title_bar(slide, "Origins and Development of MET", 'purple')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["1983: Miller publishes foundational MI paper","1991: Miller & Rollnick publish first MI book","1992: MET Manual for Project MATCH","  - Largest alcohol treatment trial ever","  - 1,726 participants across 9 US sites","  - Compared MET vs CBT vs 12-Step (TSF)"], bgc='light_purple', title="Timeline", ttc='purple', border='purple', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["MET = MI + Personalized Assessment Feedback","","1. MOTIVATIONAL INTERVIEWING (MI):","   Client-centered counseling style","   Explores and resolves ambivalence","2. PERSONALIZED FEEDBACK:","   Objective data presented to client","   Creates discrepancy with goals/values"], bgc='cream', title="Core Components", ttc='orange', border='orange', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["Project MATCH Finding: MET achieved comparable outcomes to 12-session CBT and 12-session TSF","in just 4 sessions. This was revolutionary - a brief motivational intervention matching efficacy","of longer treatments. Particularly effective for clients HIGH in anger and LOW in readiness.","","Research: At 3-year follow-up, all three treatments maintained gains. MET was the most","cost-effective option, requiring only one-third of the therapist time (4 vs 12 sessions)."], bgc='light_green', title="Key Finding:", ttc='green', border='green', fs=12)
ref_bar(slide, "Ref: Project MATCH Research Group (1997). J Studies on Alcohol, 58, 7-29; Miller (1983). Behav Psychother, 11.")



# ============================================================
# SECTION 2: THEORETICAL FOUNDATIONS
# ============================================================
divider("SECTION 2", "Theoretical Foundations of MET", 'teal')

# TTM Overview
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_teal'); title_bar(slide, "The Transtheoretical Model (TTM) of Change", 'teal')
cbox(slide, LM, ROW1, FULL_W, Inches(1.8), ["The Transtheoretical Model (Prochaska & DiClemente, 1982) is the PRIMARY theoretical foundation of MET.","It describes behavior change as a PROCESS through predictable stages. People do not change abruptly.","The therapist's task is to match interventions to the client's current stage.","\"The responsibility and capability for change lie within the client.\" (MET Manual, p. 2)"], bgc='cream', title="Foundation of MET", ttc='teal', border='teal', fs=12)
flow(slide, Inches(3.15), ["Pre-\ncontemplation", "Contemplation", "Preparation", "Action", "Maintenance", "Relapse"], ['red','orange','gold','green','teal','purple'])
cbox(slide, LM, Inches(4.4), FULL_W, Inches(1.8), ["Key Principle: Change is CYCLICAL, not linear. Most people cycle through stages 3-7 times","before achieving lasting change. Relapse is a NORMAL part of the process, not failure.","Average smoker cycles 3 times before permanent quit (Prochaska et al., 1992)."], bgc='light_red', title="Important:", ttc='red', border='red', fs=12)
ref_bar(slide, "Ref: Prochaska & DiClemente (1984). The Transtheoretical Approach; MET Manual (1992), pp. 6-12.")

# Precontemplation
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_red'); title_bar(slide, "Stage 1: PRECONTEMPLATION - Not Considering Change", 'red')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["Person does NOT see substance use as a problem.","Not thinking about change at all.","","Common statements:","\"I don't have a problem\"","\"Everyone drinks like me\"","\"My family is overreacting\"","\"I'm here because the court sent me\""], bgc='light_red', title="What It Looks Like:", ttc='red', border='red', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["DO: Raise doubt gently, increase awareness","DO: Provide objective information/feedback","DO: Explore events that brought them in","DO: Plant seeds without pushing","","DON'T: Argue or confront","DON'T: Push for immediate change","DON'T: Label them as 'alcoholic'/'addict'"], bgc='cream', title="Therapist Tasks:", ttc='orange', border='orange', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["MET Technique: Present Personal Feedback Report (PFR) showing objective data.","Example: \"Your liver enzymes are at ___ level. Here's where that falls compared to norms...\"","The therapist presents FACTS without arguing, letting the DATA create discrepancy.","","\"The precontemplation stage is characterized by a lack of awareness that a problem exists.","The individual has no intention to change behavior in the foreseeable future.\" (Manual, p. 7)"], bgc='white', title="How MET Addresses This Stage:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 6-8; Prochaska & DiClemente (1982). Psychotherapy, 19, 276-288.")

# Contemplation
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "Stage 2: CONTEMPLATION - Thinking About Change", 'orange')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["Person ACKNOWLEDGES a problem but is","AMBIVALENT. Sees both pros and cons.","","Common statements:","\"I know I drink too much, but...\"","\"I want to quit, but I don't know how\"","\"Sometimes I think I should cut down\"","\"Part of me wants to change, part doesn't\""], bgc='light_orange', title="What It Looks Like:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["DO: Explore ambivalence (don't resolve it)","DO: Use Decisional Balance worksheet","DO: Elicit self-motivational statements","DO: Tip the balance toward change","DO: Develop discrepancy with values","","DON'T: Rush to action planning","DON'T: Tell them what to do"], bgc='light_green', title="Therapist Tasks:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["MET Key Strategy: EXPLORING AMBIVALENCE with Decisional Balance","Help client articulate BOTH sides - reasons to change AND reasons to stay the same.","\"Ambivalence is the contemplator's defining feature. The therapist's task is to tip","the decisional balance.\" (Manual, p. 8)","","Research: DiClemente et al. (1991) showed contemplators who explored ambivalence progressed faster."], bgc='cream', title="MET Strategy:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 8-10; DiClemente et al. (1991). J Consult Clin Psychol, 59, 295-304.")

# Preparation
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'cream'); title_bar(slide, "Stage 3: PREPARATION - Ready to Plan", 'gold')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["Person has DECIDED to change and is","planning how. Balance has tipped.","","Common statements:","\"I need to do something about this\"","\"What are my options?\"","\"I'm going to quit next Monday\"","\"I've already started cutting down\""], bgc='light_orange', title="What It Looks Like:", ttc='gold', border='gold', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["DO: Help develop a Change Plan","DO: Offer a MENU of options (not prescribe)","DO: Support self-efficacy strongly","DO: Help set realistic, achievable goals","DO: Negotiate (not impose) a plan","","DON'T: Prescribe a single path","DON'T: Miss the window of readiness"], bgc='light_green', title="Therapist Tasks:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["MET Key Strategy: CHANGE PLAN WORKSHEET","Help client complete: (1) Changes desired (2) Reasons (3) Steps planned","(4) How others can help (5) How they'll know it's working (6) Things that could interfere","","\"The window of determination is open for a period of time. If action is not taken, the person","may slip back into contemplation or precontemplation.\" (MET Manual, p. 10)"], bgc='light_blue', title="MET Strategy:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 10-11; DiClemente et al. (1991). J Consult Clin Psychol, 59, 295-304.")

# Action, Maintenance, Relapse
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "Stages 4-6: Action, Maintenance, and Relapse", 'green')
cbox(slide, LM, ROW1, THIRD_W, ROW1_H, ["ACTION: Person is actively","making changes to behavior.","","Therapist tasks:","- Affirm efforts","- Troubleshoot obstacles","- Review Change Plan","- Build self-efficacy"], bgc='light_green', title="Action", ttc='green', border='green', fs=11)
cbox(slide, COL3_2, ROW1, THIRD_W, ROW1_H, ["MAINTENANCE: Sustaining","gains over time.","","Therapist tasks:","- Prevent complacency","- Identify high-risk situations","- Plan coping strategies","- Celebrate milestones"], bgc='light_teal', title="Maintenance", ttc='teal', border='teal', fs=11)
cbox(slide, COL3_3, ROW1, THIRD_W, ROW1_H, ["RELAPSE: Return to earlier","stage - NOT a failure.","","Therapist tasks:","- Normalize (3-7 cycles avg)","- Explore without blame","- Revise Change Plan","- Re-engage motivation"], bgc='light_red', title="Relapse", ttc='red', border='red', fs=11)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["\"The Transtheoretical Model views relapse not as catastrophic failure but as a normal part of","the cycle of change. With each cycle, the person learns something new.\" (MET Manual, p. 12)","","MET Sessions 3-4 focus on these stages: reviewing progress, affirming successes, problem-solving","obstacles, and normalizing any slips as information rather than defeat.","Research: Prochaska et al. (1992) found most successful changers had 3+ prior attempts."], bgc='cream', title="From the Manual:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 11-14; Prochaska, DiClemente & Norcross (1992). Am Psychologist, 47.")

# Self-Efficacy
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_purple'); title_bar(slide, "Theoretical Foundation: Self-Efficacy (Bandura, 1977)", 'purple')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["Self-efficacy = person's BELIEF in ability to","successfully perform a behavior.","","In MET: belief that one CAN change.","","Sources (Bandura):","1. Past performance accomplishments","2. Vicarious experience (others' success)","3. Verbal persuasion (therapist affirmation)","4. Physiological/emotional states"], bgc='light_purple', title="What is Self-Efficacy?", ttc='purple', border='purple', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["How MET Builds Self-Efficacy:","","1. AFFIRMING client's strengths/past successes","   \"You stayed sober 3 months - real strength\"","2. SUPPORTING ability to change","   \"You have what it takes to do this\"","3. OFFERING menu of options (sense of control)","4. CELEBRATING small wins","5. \"Why not a lower number?\" on rulers"], bgc='light_green', title="MET Application:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["\"A person who lacks confidence that he or she can change is unlikely to try. The therapist's","task is to enhance the client's belief in the possibility of change.\" (MET Manual, p. 17)","","Research: Ilgen et al. (2005) found self-efficacy at treatment entry predicted drinking outcomes","at 1-year follow-up. Clients with higher self-efficacy had significantly better outcomes.","This is the 'S' in FRAMES: Support self-efficacy."], bgc='cream', title="Manual Quote & Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Bandura (1977). Psychol Review, 84; MET Manual (1992), pp. 16-17; Ilgen et al. (2005). Addictive Behav.")

# Cognitive Dissonance
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "Theoretical Foundation: Cognitive Dissonance (Festinger, 1957)", 'orange')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["Cognitive Dissonance = uncomfortable tension","when behavior contradicts values/beliefs.","","Substance use examples:","- \"I am a good parent\" vs \"I drink every","   night and my children suffer\"","- \"I value my health\" vs \"I am damaging","   my liver with alcohol\"","- \"I am independent\" vs \"I can't function","   without my substance\""], bgc='light_orange', title="The Theory:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["How MET Uses Cognitive Dissonance:","","1. DEVELOPING DISCREPANCY (key principle)","   Show gap between IS and WANT TO BE","2. Personal Feedback creates dissonance:","   \"You said health is your top value...\"","   \"Your liver enzymes show damage...\"","3. Client must resolve the tension","4. Resolution happens through CHANGE","   (not through therapist confrontation)"], bgc='light_blue', title="MET Application:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["\"A discrepancy between present behavior and important personal goals will motivate change.","The therapist's task is to develop and amplify such discrepancy.\" (MET Manual, p. 14)","","Research: Draycott & Dabbs (1998) confirmed discrepancy awareness predicts behavior change.","McNally et al. (2005) showed developing discrepancy is the most potent MI technique for","moving clients from contemplation to preparation stage."], bgc='cream', title="Manual Quote & Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Festinger (1957). Cognitive Dissonance; MET Manual (1992), pp. 13-15; McNally et al. (2005).")



# ============================================================
# SECTION 3: PRINCIPLES & SPIRIT
# ============================================================
divider("SECTION 3", "Principles and Spirit of MET", 'green')

# Five Principles
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "Five Principles of Motivational Interviewing in MET", 'green')
for i, (t,d,bc,tc2) in enumerate([("1. Express Empathy","Accept client where they are. Ambivalence is NORMAL. Reflective listening.",'light_blue','deep_blue'),("2. Develop Discrepancy","Help client see gap between current behavior and important goals/values.",'light_green','green'),("3. Avoid Argumentation","Arguments are counterproductive. Defending breeds defensiveness.",'light_orange','orange'),("4. Roll with Resistance","Don't fight resistance - use it. Offer new perspectives without imposing.",'light_purple','purple'),("5. Support Self-Efficacy","Client's belief in possibility of change is a key motivator.",'light_teal','teal')]):
    cbox(slide, LM, Inches(1.15)+Inches(1.05)*i, FULL_W, Inches(0.95), [d], bgc=bc, title=t, ttc=tc2, border=tc2, fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 13-17; Miller & Rollnick (1991). Motivational Interviewing, Ch. 3.")

# FRAMES
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "FRAMES: Elements of Effective Brief Interventions", 'deep_blue')
for i, (t,d,bc,tc2) in enumerate([("F - Feedback","Personal feedback about risk/impairment based on objective assessment data",'light_red','red'),("R - Responsibility","Emphasize responsibility for change lies with the CLIENT",'light_orange','orange'),("A - Advice","Clear advice to change given in a non-prescriptive manner",'cream','gold'),("M - Menu","Offer a MENU of strategies/options - client chooses their own path",'light_green','green'),("E - Empathy","Warm, reflective, empathic counseling style without judgment",'light_blue','deep_blue'),("S - Self-Efficacy","Reinforce hope and belief that change IS possible for this person",'light_purple','purple')]):
    cbox(slide, LM, Inches(1.15)+Inches(0.88)*i, FULL_W, Inches(0.8), [d], bgc=bc, title=t, ttc=tc2, border=tc2, fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 16-17; Miller & Sanchez (1994). Motivating young adults for change.")

# OARS
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_teal'); title_bar(slide, "OARS: Core Microskills of MET", 'teal')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["OPEN-ENDED QUESTIONS:","- Cannot be answered with yes/no","- \"What concerns you about your drinking?\"","- \"How has substance use affected your life?\"","","AFFIRMATIONS:","- Recognize strengths and efforts","- \"It took courage to come here today\"","- \"You've shown real resilience\""], bgc='light_blue', title="O - Open Questions & A - Affirmations", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["REFLECTIVE LISTENING (Primary skill):","- Repeat back MEANING (not just words)","- Simple: \"So you feel frustrated\"","- Complex: \"You want to quit but worry","  about losing your social circle\"","","SUMMARIES:","- Collect and link what client has said","- \"Let me see if I have this right...\""], bgc='light_green', title="R - Reflections & S - Summaries", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["PRACTICE RATIO: Aim for 2-3 reflections for every question asked.","Most common mistake: asking too many questions (feels like interrogation).","","\"The principal technique for expressing empathy is reflective listening... seeking through","your responses to understand the client's meaning and feelings.\" (MET Manual, p. 21)","","Research: Moyers et al. (2005) showed MI-consistent behaviors (reflections) predicted better outcomes."], bgc='cream', title="Clinical Tip:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 20-35; Moyers et al. (2005). J Subst Abuse Treat, 28, 19-26.")

# DARES
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "DARES: Self-Motivational Statements to Elicit", 'orange')
cbox(slide, LM, Inches(1.15), FULL_W, Inches(1.2), ["Goal: ELICIT (not provide) self-motivational statements. The more the CLIENT voices arguments","FOR change, the more likely actual change becomes. Listen for and actively evoke these:"], bgc='cream', title="Core Principle:", ttc='orange', border='gold', fs=12)
for i,(t,ex,bc) in enumerate([("D - Desire","\"I want to change\" / \"I wish I could stop\"",'light_blue'),("A - Ability","\"I think I could do it\" / \"I was able to quit before\"",'light_green'),("R - Reasons","\"My health is suffering\" / \"My family needs me sober\"",'light_purple'),("E - Emotional","\"I'm tired of feeling this way\" / \"I hate what I've become\"",'light_orange'),("S - Steps","\"I could try meetings\" / \"Maybe I'll call my doctor\"",'light_teal')]):
    cbox(slide, LM, Inches(2.55)+Inches(0.75)*i, FULL_W, Inches(0.65), [f"Example: {ex}"], bgc=bc, title=t, ttc='dark_gray', fs=11)
ref_bar(slide, "Ref: MET Manual (1992), pp. 25-30; Amrhein et al. (2003). J Consult Clin Psychol, 71, 862-878.")

# Spirit of MI
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "The Spirit of Motivational Interviewing", 'green')
cbox(slide, LM, ROW1, THIRD_W, ROW1_H, ["Working WITH the client as","equal partners. NOT","expert-to-patient.","","\"The therapeutic relationship","is more like a partnership","than an expert-recipient","one.\" (Manual, p. 13)"], bgc='light_blue', title="Collaboration", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL3_2, ROW1, THIRD_W, ROW1_H, ["Drawing out client's OWN","motivations, strengths,","and reasons for change.","","Motivation is EVOKED","from within, not installed","from outside. The client","has the answers."], bgc='light_green', title="Evocation", ttc='green', border='green', fs=12)
cbox(slide, COL3_3, ROW1, THIRD_W, ROW1_H, ["Respecting client's right","and capacity to direct","their own life.","","Client decides whether,","when, and how to change.","Even if they choose NOT","to change right now."], bgc='light_purple', title="Autonomy", ttc='purple', border='purple', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["In Practice: Therapist asks more than tells; listens more than instructs.","Client does most of the talking (aim for 70:30 client:therapist talk ratio).","No labeling, no shaming, no arguing for change on behalf of the client.","","Research: Miller et al. (1993) found therapist directiveness predicted worse outcomes.","Client resistance increased when therapists were confrontational."], bgc='cream', title="What It Looks Like:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 13-18; Miller et al. (1993). J Stud Alcohol, 54, 455-461.")

# Handling Resistance
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_red'); title_bar(slide, "Handling Client Resistance: Strategies from the Manual", 'red')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["TYPES OF RESISTANCE:","1. Arguing - challenging, hostility","2. Interrupting - cutting off, talking over","3. Denying - blaming, excusing, minimizing","4. Ignoring - inattention, sidetracking","","KEY: Resistance is a signal to CHANGE","your approach, not fight harder."], bgc='light_red', title="Recognizing Resistance:", ttc='red', border='red', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["STRATEGIES TO ROLL WITH IT:","1. Simple Reflection - acknowledge","2. Amplified Reflection - stronger form","3. Double-Sided Reflection - both sides","4. Shifting Focus - redirect topic","5. Agreement with Twist - agree then reframe","6. Reframing - new interpretation","7. Emphasize Personal Choice"], bgc='light_green', title="Rolling with Resistance:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["Client: \"I don't think I drink more than my friends.\"","Simple Reflection: \"You see yourself as a pretty normal drinker.\"","Amplified: \"So there's really nothing at all to be concerned about.\"","Double-Sided: \"Your drinking feels normal to you, AND something prompted you to come in.\"","Reframe: \"Your friends are important and you want to fit in with them.\"","","Research: Resistance in session predicts WORSE outcomes (Miller et al., 1993)."], bgc='cream', title="Clinical Examples:", ttc='orange', border='gold', fs=11)
ref_bar(slide, "Ref: MET Manual (1992), pp. 32-38; Miller et al. (1993). J Stud Alcohol, 54, 455-461.")



# ============================================================
# SECTION 4: SESSION-BY-SESSION GUIDE
# ============================================================
divider("SECTION 4", "Session-by-Session Therapy Guide\nHow to Actually DO MET", 'purple')

# Overview
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_purple'); title_bar(slide, "MET: The 4-Session Structure (Overview)", 'purple')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["SESSION 1 (Week 1):","- Build rapport and therapeutic alliance","- Present Personal Feedback Report","- Explore client's reaction to feedback","- Elicit self-motivational statements","- Gauge readiness to change","","SESSION 2 (Week 2):","- Strengthen commitment to change","- Develop Change Plan (collaborative)"], bgc='light_blue', title="Sessions 1 & 2", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["SESSION 3 (Week 6):","- Review progress since Session 2","- Renew motivation and commitment","- Problem-solve obstacles","- Revise Change Plan if needed","","SESSION 4 (Week 12):","- Review overall progress","- Consolidate gains","- Maintenance/relapse prevention planning"], bgc='light_orange', title="Sessions 3 & 4", ttc='orange', border='orange', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["KEY PRINCIPLE: Each session uses MI techniques throughout. MET is not a checklist but a clinical","style applied within a structured framework. The therapist continuously uses OARS skills, evokes","change talk, rolls with resistance, and supports self-efficacy in every interaction.","","Session spacing is intentional: Session 1-2 close together (build momentum), then space out to","allow client time to implement changes and experience results before reviewing progress."], bgc='cream', title="Important Note:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 40-80; Chapters III-V: Session protocols.")

# Session 1
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "SESSION 1: Building Motivation for Change", 'deep_blue')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["STEP 1: OPENING (10-15 min)","- Welcome warmly and genuinely","- Explain the process briefly","- Set non-judgmental tone","- Ask open question to start:","  \"What brought you here today?\"","  \"How are things going for you?\"","","Then LISTEN. Use reflections."], bgc='light_blue', title="Opening:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["STEP 2: PERSONAL FEEDBACK (30-40 min)","Present assessment data objectively:","1. Drinking/drug use patterns","2. Comparison with population norms","3. Blood test results (liver, etc.)","4. Neuropsych test results","5. Consequences reported","6. Risk factors identified","After each: \"What do you make of this?\""], bgc='light_green', title="Feedback:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["STEP 3: ELICITING CHANGE TALK (15-20 min)","- Evocative questions: \"What worries you about your use?\"","- Importance Ruler: \"How important is it to change? (0-10) Why not a lower number?\"","- Explore pros/cons: \"What do you like about using? What concerns you?\"","- Look forward: \"Where do you see yourself in 5 years if nothing changes?\"","","STEP 4: CLOSING (5-10 min) - Summarize key themes, assess readiness, assign reflection task."], bgc='cream', title="Eliciting & Closing:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 40-58; Chapter III: Session 1 - Building Motivation for Change.")

# Session 1 Feedback Detail
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "SESSION 1: Personal Feedback Report (Detail)", 'deep_blue')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["PERSONAL FEEDBACK REPORT contains:","","1. DRINKING/DRUG USE PATTERN","   - Quantity/frequency data","   - Peak blood alcohol levels","   - Where client falls on population norms","2. NEGATIVE CONSEQUENCES","   - Physical, psychological, social","3. RISK FACTORS/INDICATORS","   - Family history, tolerance, dependence"], bgc='white', title="Components:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["HOW TO PRESENT FEEDBACK:","","1. Use a structured form/booklet","2. Go through each section systematically","3. After each piece of data, PAUSE and ASK:","   - \"What do you make of this?\"","   - \"Does this surprise you?\"","4. Use reflective listening on response","5. Do NOT argue if they minimize","6. Highlight discrepancies gently"], bgc='light_green', title="Technique:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["\"The feedback session is not a lecture. It is a structured clinical interaction in which the","therapist presents information and invites the client's reaction. The style is empathic,","non-confrontational, and curious rather than didactic.\" (MET Manual, p. 44)","","Research: Agostinelli et al. (1995) found that personalized normative feedback significantly","increased motivation to change when compared to no-feedback control groups."], bgc='cream', title="Manual Guidance:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 44-55; Agostinelli et al. (1995). J Consult Clin Psychol, 63, 284-290.")

# Session 1 - Dialogue
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "SESSION 1: Clinical Dialogue Example", 'deep_blue')
cbox(slide, LM, Inches(1.15), FULL_W, Inches(5.1), ["FEEDBACK DIALOGUE:","T: \"I'd like to share results from your assessment. Here's where your drinking falls","   compared to the general adult population...\" [shows graph]","C: \"Hmm. I didn't think it was that high.\"","T: \"It's more than you expected.\" (Simple reflection)","C: \"Yeah. But a lot of people drink.\"","T: \"You're not alone in drinking, AND what surprises you is how your amount compares.\"","C: \"I mean... maybe I do drink more than I thought.\" (CHANGE TALK - problem recognition)","T: \"You're starting to see that your drinking may be on the higher end.\" (Amplify change talk)","","ELICITING CHANGE TALK:","T: \"You mentioned your wife has been concerned. What concerns her the most?\"","C: \"She says I'm different when I drink. More angry. I guess she has a point.\"","T: \"You've noticed that too - alcohol changes you in ways you don't like.\" (Reflection)","C: \"Yeah... I don't want to be that kind of person. That's not who I am.\"","T: \"Being a good person and partner is really important to you.\" (Affirming values)","","NOTE: Therapist NEVER argues. Every response is reflection, affirmation, or open question."], bgc='white', border='deep_blue', fs=11)
ref_bar(slide, "Ref: MET Manual (1992), pp. 48-55; Dialogue adapted from manual principles and examples.")

# Session 1 - Eliciting Change Talk
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "SESSION 1: Eliciting Self-Motivational Statements", 'deep_blue')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["CATEGORIES OF CHANGE TALK:","","1. Problem Recognition:","   \"I guess this is more serious than I thought\"","2. Concern:","   \"I'm worried about what this is doing to me\"","3. Intention to Change:","   \"I think I need to do something about this\"","4. Optimism:","   \"I think I can do this if I really try\""], bgc='light_blue', title="Change Talk Types:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["HOW TO ELICIT CHANGE TALK:","","1. EVOCATIVE QUESTIONS:","   \"What worries you about your drinking?\"","   \"What would be different if you changed?\"","2. IMPORTANCE RULER (0-10):","   \"Why not a lower number?\" (key question!)","3. EXPLORING PROS & CONS:","   \"What do you like about using?\"","4. LOOKING FORWARD/BACK:","   \"Where are you in 5 years if nothing changes?\""], bgc='light_green', title="Techniques to Elicit:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["SESSION 1 CLOSING: Provide comprehensive summary that includes:","- Key facts from feedback | Client's reactions | Any self-motivational statements made","- Both sides of ambivalence | Any expressions of concern or intention","","READINESS RULER: \"On a scale of 1-10, how ready are you to change right now?\"","If 1-3: Precontemplation (plant seeds) | If 4-6: Contemplation (explore ambivalence)","If 7-10: Preparation (move toward Change Plan in Session 2)"], bgc='cream', title="Closing & Readiness:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 48-58; Miller & Rollnick (1991). MI, Ch. 5: Eliciting Self-Motivational Stmts.")

# Session 2
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "SESSION 2: Strengthening Commitment & Change Plan", 'green')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["STEP 1: OPENING (10 min)","\"What's happened since last time?\"","\"Have you thought more about this?\"","Reflect on any changes already made.","","STEP 2: RECAPITULATION (5 min)","Brief summary of Session 1 themes.","Check: \"Did I capture that correctly?\"","","STEP 3: DEEPENING (15 min)","Continue exploring ambivalence if needed."], bgc='light_green', title="First Half:", ttc='green', border='green', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["STEP 4: CHANGE PLAN (25-30 min)","","\"Would you like to put together a plan?\"","","Change Plan Worksheet:","1. Changes I want to make...","2. Most important reasons...","3. Steps I plan to take...","4. How others can help...","5. I'll know it's working if...","6. Things that could interfere..."], bgc='light_blue', title="Change Plan:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["ADAPTING TO READINESS:","- NOT READY: Continue rapport, present more feedback, plant seeds gently","- UNSURE: Use Decisional Balance, explore values vs behavior, evocative questions","- READY: Negotiate Change Plan, offer menu of options, help set specific goals","","\"The Change Plan is the CLIENT's plan - not the therapist's prescription. It is negotiated","collaboratively, and the client's preferences and choices are honored.\" (MET Manual, p. 63)"], bgc='cream', title="Adapting to Client:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 59-68; Chapter IV: Session 2 - Strengthening Commitment.")

# Sessions 3-4
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "SESSIONS 3 & 4: Reviewing Progress & Maintenance", 'orange')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["SESSION 3 (Week 6):","1. \"How have things been going?\"","2. Review Change Plan progress","   - What's working? What isn't?","   - Any modifications needed?","3. Renew motivation","   - Re-explore importance/confidence","   - Affirm progress (even small steps)","4. Address slips WITHOUT judgment","5. Problem-solve new obstacles"], bgc='light_orange', title="Session 3:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["SESSION 4 (Week 12):","1. Review overall progress","   - Compare current to baseline","   - Celebrate achievements","2. Maintenance planning","   - Identify high-risk situations","   - Plan coping strategies","3. Discuss relapse as normal","4. Plan ongoing support","5. Termination with hope"], bgc='light_purple', title="Session 4:", ttc='purple', border='purple', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["IF CLIENT HAS RELAPSED:","1. Normalize: \"Slips are a common part of change. Many people experience this.\"","2. Explore without blame: \"What happened? What was going on?\"","3. Reframe: \"Now you know that [situation] is a trigger. That's useful information.\"","4. Re-engage: \"Where do you want to go from here?\"","","NEVER: Express disappointment, say \"you failed\", lecture, or label them as hopeless."], bgc='light_red', title="Handling Relapse:", ttc='red', border='red', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 69-80; Chapter V: Sessions 3-4; Marlatt & Gordon (1985).")



# ============================================================
# SECTION 5: CLINICAL TECHNIQUES
# ============================================================
divider("SECTION 5", "Clinical Techniques in Detail", 'teal')

# Reflective Listening
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_teal'); title_bar(slide, "Reflective Listening: The Foundation Skill", 'teal')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["LEVELS OF REFLECTION:","","Level 1 - REPEAT/REPHRASE:","C: \"I drink every night\" T: \"You drink every evening\"","","Level 2 - PARAPHRASE (meaning):","T: \"It's become a daily routine\"","","Level 3 - REFLECTION OF FEELING:","T: \"You sound concerned about that pattern\""], bgc='light_teal', title="Types:", ttc='teal', border='teal', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["STRATEGIC USE IN MET:","","- Reflect CHANGE TALK more than sustain talk","- AMPLIFY motivation:","  C: \"I guess I drink a bit much\"","  T: \"You've realized your drinking has become","      a real problem\" (amplified)","- DOUBLE-SIDED for ambivalence:","  \"Drinking helps you relax AND it's costing","   you your marriage.\""], bgc='light_green', title="Strategic Use:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["PRACTICE GUIDELINES:","- Aim for 2-3 reflections for every question asked","- Reflections are STATEMENTS (drop voice at end), not questions","- If your reflection is wrong, client will correct you - that's fine","- Common mistake: Too many questions, not enough reflections","","Research: Apodaca & Longabaugh (2009) meta-analysis found MI-consistent therapist behaviors","(especially reflections) were significantly linked to reduced substance use."], bgc='cream', title="Tips & Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 20-24; Apodaca & Longabaugh (2009). J Subst Abuse Treat, 37, 68-86.")

# Decisional Balance & Rulers
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "Techniques: Decisional Balance & Rulers", 'orange')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["DECISIONAL BALANCE (4-quadrant grid):","","Good things about using | Costs of using","________________________|________________","Benefits of change     | Costs of change","________________________|________________","","Start with GOOD things about using first!","(Shows understanding, builds trust)","Then: \"Where does that leave you?\""], bgc='light_orange', title="Decisional Balance:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["IMPORTANCE & CONFIDENCE RULERS:","","\"How important is it to change? (0-10)\"","\"How confident are you? (0-10)\"","","KEY QUESTION: \"Why not a lower number?\"","(Invites client to argue FOR change)","","If importance low: develop discrepancy","If confidence low: build self-efficacy","If both high: move to Change Plan"], bgc='light_blue', title="Rulers:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["WHY \"WHY NOT A LOWER NUMBER?\" WORKS:","\"Why not higher?\" = client defends why they HAVEN'T changed (sustain talk)","\"Why not lower?\" = client articulates reasons they DO want to change (change talk)","","Research: LaBrie et al. (2006) found Importance/Confidence rulers in brief MI predicted","actual behavior change at follow-up. Articulation of own reasons for change is more persuasive","than hearing reasons from others (Self-Perception Theory, Bem 1972)."], bgc='cream', title="Clinical Wisdom:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 28-37; LaBrie et al. (2006). Addictive Behaviors, 31(8), 1428-1435.")

# Therapist Traps
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_red'); title_bar(slide, "Therapist Traps to AVOID (from the Manual)", 'red')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["1. QUESTION-ANSWER TRAP","   Too many Qs = feels like interrogation","   Fix: More reflections, fewer questions","","2. CONFRONTATION-DENIAL TRAP","   Arguing = client defends drinking more","   Fix: Roll with resistance, reflect","","3. EXPERT TRAP","   \"You should...\" / \"Research shows...\"","   Fix: Elicit client's own solutions first"], bgc='light_red', title="Traps 1-3:", ttc='red', border='red', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["4. LABELING TRAP","   Insisting on \"alcoholic\" label","   Fix: Focus on behavior, not labels","","5. PREMATURE FOCUS TRAP","   Jumping to solutions too early","   Fix: Follow client's pace","","6. BLAMING TRAP","   Client feels blamed for problem","   Fix: \"What matters is what you want NOW\""], bgc='light_orange', title="Traps 4-6:", ttc='orange', border='orange', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["\"Direct argumentation and an aggressive confrontational approach tend to increase client","resistance and are associated with poorer outcomes.\" (MET Manual, pp. 18-19)","","REMEMBER: Every time YOU argue FOR change, the client argues AGAINST it.","This is the 'righting reflex' - therapist's desire to fix things actually backfires.","","Research: Miller et al. (1993) showed confrontational therapist style predicted client","drinking at 1-year follow-up: MORE confrontation = MORE drinking."], bgc='cream', title="Key Principle:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 18-20; Miller et al. (1993). J Stud Alcohol, 54, 455-461.")



# ============================================================
# SECTION 6: CASE CONCEPTUALIZATION
# ============================================================
divider("SECTION 6", "Case Conceptualization\nApplying MET to a Clinical Case", 'maroon')

# Case Intro
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_red'); title_bar(slide, "Case Study: Applying MET to a Substance Use Case", 'maroon')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["PRESENTING INFORMATION:","- Age: 25 years, Male","- Multiple substance use (polysubstance)","- History of marital dissolution (divorced)","- Referred for assessment and intervention","","ASSESSMENT FINDINGS:","- Externalizing personality organization","- Adequate psychological resources","- Poorly modulated affect (impulsive)","- Damaged, negative self-concept"], bgc='light_red', title="Client Background:", ttc='maroon', border='maroon', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["KEY PSYCHOLOGICAL FEATURES:","- Emotion-driven decision making","- Hasty, inefficient information processing","- Raw emotional discharge under stress","- Socially engaged but difficulty with","  deep intimate relationships","- Pervasive pessimism/damage-oriented self","- Elevated risk indicators (monitor closely)","","SUBSTANCE USE FUNCTION:","Substances regulate overwhelming affect"], bgc='light_orange', title="Psychological Profile:", ttc='orange', border='orange', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["WHY MET IS APPROPRIATE FOR THIS CLIENT:","1. Adequate resources but poor utilization (emotion-first style) - MET engages existing capacity","2. Confrontation would activate externalizing defenses and increase resistance","3. Damaged self-image NEEDS affirmation (the 'S' in FRAMES), not more criticism","4. Autonomy needs respect (impulsive style responds poorly to directives)","5. Substances serve affect-regulation function - MET explores this non-judgmentally"], bgc='cream', title="MET Rationale:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: Case from Rorschach assessment (Exner Comprehensive System); MET Manual (1992).")

# Case Sessions
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "Case: Session-by-Session MET Plan", 'deep_blue')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["SESSION 1 FOCUS:","- Build rapport carefully (guarded intimacy)","- Present feedback GENTLY (damaged self)","- Focus discrepancy: VALUES (being good","  partner) vs SITUATION (divorced, dependent)","- AFFIRM strengths heavily (resources exist)","- Monitor emotional escalation closely","","Caution: His affect can escalate quickly.","Keep pace slow, reflect often."], bgc='light_blue', title="Session 1:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["SESSION 2 FOCUS:","- Change Plan addresses FUNCTION of use","  (What does using DO for you? What need?)","- Small, achievable steps (failures feel","  catastrophic given damaged self-image)","- Affect regulation alternatives:","  \"What else helps when feelings get intense?\"","- Reflective delay practice:","  \"What if you took 10 min before deciding?\"","Key: Plan must match his pace/style."], bgc='light_green', title="Session 2:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["SESSIONS 3-4 FOCUS:","- Session 3: If progress - affirm heavily (counters damaged self-image). If slip - normalize,","  explore emotional triggers, revise plan to add more emotion regulation strategies.","- Session 4: Maintenance plan for high-risk situations (intense emotions, conflict, rejection).","  Referral for continued therapy (affect regulation, self-image repair, intimacy work).","","BEYOND MET: Affect regulation training, self-image repair, interpersonal skills, slower processing.","MET provides the MOTIVATIONAL FOUNDATION for longer-term therapeutic work."], bgc='cream', title="Sessions 3-4 & Beyond:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), Chs. III-V; Exner (2003). Rorschach, Vol. 1; NIMHANS (2016).")

# Case Dialogue
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "Case: Clinical Dialogue (Session 1)", 'deep_blue')
cbox(slide, LM, Inches(1.15), FULL_W, Inches(5.1), ["T: \"You mentioned that your marriage ended. What role, if any, did substance use play?\"","C: \"She said I was different when I used. Angry.\"","T: \"Being the kind of partner you want to be is important to you, and the substances were","   getting in the way of that.\" (Developing discrepancy)","C: \"Yeah... I didn't mean to be like that.\"","T: \"There's a gap between who you WANT to be and what happens when you use.\"","C: \"I guess so. I just... I don't know how to handle things without it.\"","T: \"You're wondering if you can cope differently. That's an honest question.\" (Reflection)","C: \"Nobody ever asked me that before. They just tell me to stop.\"","T: \"And being told what to do doesn't help you.\" (Rolling with resistance)","C: \"No, it makes me want to do it more.\"","T: \"You respond better when you get to decide for yourself.\" (Affirming autonomy)","","KEY TECHNIQUES USED: Developing discrepancy, reflective listening, rolling with resistance,","affirming autonomy, open questions. Notice: therapist NEVER said \"you should stop using.\""], bgc='white', border='deep_blue', fs=11)
ref_bar(slide, "Ref: MET Manual (1992); Dialogue based on MI principles applied to this case profile.")



# ============================================================
# SECTION 7: WORKSHEETS
# ============================================================
divider("SECTION 7", "Clinical Worksheets & Tools\nPrintable Resources for MET Sessions", 'deep_blue')

# Change Plan
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "WORKSHEET 1: Change Plan (from MET Manual)", 'deep_blue')
cbox(slide, LM, Inches(1.15), FULL_W, Inches(5.1), ["1. THE CHANGES I WANT TO MAKE ARE:","   _________________________________________________________________________","","2. THE MOST IMPORTANT REASONS WHY I WANT TO MAKE THESE CHANGES ARE:","   _________________________________________________________________________","","3. THE STEPS I PLAN TO TAKE IN CHANGING ARE:","   _________________________________________________________________________","","4. THE WAYS OTHER PEOPLE CAN HELP ME ARE:","   Person: _________________ How: ___________________________________________","","5. I WILL KNOW THAT MY PLAN IS WORKING IF:","   _________________________________________________________________________","","6. SOME THINGS THAT COULD INTERFERE WITH MY PLAN ARE:","   _________________________________________________________________________"], bgc='white', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), Appendix B: Change Plan Worksheet.")

# Decisional Balance
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "WORKSHEET 2: Decisional Balance Grid", 'green')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["GOOD THINGS about my substance use:","","1. ____________________________________","2. ____________________________________","3. ____________________________________","4. ____________________________________","5. ____________________________________"], bgc='light_orange', title="Benefits of Status Quo:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["COSTS of my current substance use:","","1. ____________________________________","2. ____________________________________","3. ____________________________________","4. ____________________________________","5. ____________________________________"], bgc='light_red', title="Costs of Status Quo:", ttc='red', border='red', fs=12)
cbox(slide, LM, ROW2, HALF_W, Inches(2.4), ["BENEFITS of making a change:","","1. ____________________________________","2. ____________________________________","3. ____________________________________","4. ____________________________________","5. ____________________________________"], bgc='light_green', title="Benefits of Change:", ttc='green', border='green', fs=12)
cbox(slide, COL2, ROW2, HALF_W, Inches(2.4), ["COSTS of making a change:","","1. ____________________________________","2. ____________________________________","3. ____________________________________","4. ____________________________________","5. ____________________________________"], bgc='light_purple', title="Costs of Change:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: Janis & Mann (1977). Decision Making; MET Manual (1992), p. 29.")

# Rulers
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_purple'); title_bar(slide, "WORKSHEET 3: Importance, Confidence & Readiness Rulers", 'purple')
cbox(slide, LM, Inches(1.15), FULL_W, Inches(1.5), ["IMPORTANCE: How important is it to you to make this change?","Not important  0 --- 1 --- 2 --- 3 --- 4 --- 5 --- 6 --- 7 --- 8 --- 9 --- 10  Extremely important","Why this number and not lower? _______________________________________________________"], bgc='light_blue', border='deep_blue', fs=12)
cbox(slide, LM, Inches(2.85), FULL_W, Inches(1.5), ["CONFIDENCE: How confident are you that you COULD make this change?","Not confident  0 --- 1 --- 2 --- 3 --- 4 --- 5 --- 6 --- 7 --- 8 --- 9 --- 10  Extremely confident","Why this number and not lower? _______________________________________________________"], bgc='light_green', border='green', fs=12)
cbox(slide, LM, Inches(4.55), FULL_W, Inches(1.5), ["READINESS: How ready are you to make this change RIGHT NOW?","Not ready  0 --- 1 --- 2 --- 3 --- 4 --- 5 --- 6 --- 7 --- 8 --- 9 --- 10  Completely ready","What would help you move up one point? _______________________________________________"], bgc='light_orange', border='orange', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 35-37; Rollnick, Mason & Butler (1999). Health Behavior Change.")

# Diary
tbl("WORKSHEET 4: Daily Self-Monitoring Diary", ["Day","Situation/Trigger","Feelings (0-10)","Urge (0-10)","What I Did Instead","Result"],
    [["Monday","____________","____","____","____________","______"],["Tuesday","____________","____","____","____________","______"],["Wednesday","____________","____","____","____________","______"],["Thursday","____________","____","____","____________","______"],["Friday","____________","____","____","____________","______"],["Saturday","____________","____","____","____________","______"],["Sunday","____________","____","____","____________","______"]],
    color='teal', ref="Ref: Adapted from MET Manual (1992); Marlatt & Gordon (1985). Relapse Prevention.")

# RP Plan
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "WORKSHEET 5: My Relapse Prevention Plan", 'orange')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["MY HIGH-RISK SITUATIONS:","1. ____________________________________","2. ____________________________________","3. ____________________________________","","MY EARLY WARNING SIGNS:","1. ____________________________________","2. ____________________________________","3. ____________________________________"], bgc='light_red', title="Identifying Risks:", ttc='red', border='red', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["MY COPING STRATEGIES:","1. ____________________________________","2. ____________________________________","3. ____________________________________","","SUPPORT CONTACTS:","Name: _____________ Phone: _____________","Name: _____________ Phone: _____________","Name: _____________ Phone: _____________"], bgc='light_green', title="My Resources:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["IF I HAVE A SLIP, I WILL:","1. Remember: a slip is NOT a failure - it's information about a trigger to address","2. Call: _________________________________ (my support person)","3. Do instead: _______________________________________________________________","4. Review: What triggered it? What was I feeling? What can I learn?","","MY TOP REASONS TO STAY ON TRACK: ________________________________________________"], bgc='cream', title="Emergency Plan:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), Session 4; Marlatt & Donovan (2005). Relapse Prevention, 2nd ed.")

# Values
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "WORKSHEET 6: Personal Values Exploration", 'green')
cbox(slide, LM, Inches(1.15), FULL_W, Inches(3.2), ["Rate each: IMPORTANCE (1-5) and how CURRENT BEHAVIOR aligns (1-5)","","VALUE                    IMPORTANCE    ALIGNMENT    GAP?","Being a good parent      _________     _________    ___","Physical health          _________     _________    ___","Financial security       _________     _________    ___","Being honest             _________     _________    ___","Close relationships      _________     _________    ___","Self-respect             _________     _________    ___","Career success           _________     _________    ___","Independence             _________     _________    ___"], bgc='white', title="Values-Behavior Alignment:", ttc='green', border='green', fs=12)
cbox(slide, LM, Inches(4.55), FULL_W, Inches(1.7), ["Which values have the BIGGEST gaps? How does substance use affect these values?","What would change look like for the values that matter most?","","Therapist Note: Creates cognitive dissonance by making the gap explicit. Use to develop discrepancy."], bgc='cream', title="Reflection:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Miller et al. (2001). Personal Values Card Sort; MET Manual (1992).")



# ============================================================
# SECTION 8: COMPARISON WITH OTHER THERAPIES (NEW - ~8 slides)
# ============================================================
divider("SECTION 8", "MET Compared with Other Therapies\nFor Patient, Family, and Caregivers", 'navy')

# MET vs CBT
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "MET vs Cognitive Behavioral Therapy (CBT)", 'navy')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["MET APPROACH:","- Focus: WHY to change (motivation)","- Client provides the answers","- 4 sessions over 12 weeks","- Non-directive, exploratory style","- Works best EARLY in treatment","- Addresses ambivalence and readiness","- No homework assignments given","- Therapist evokes, doesn't teach"], bgc='light_blue', title="MET:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["CBT APPROACH:","- Focus: HOW to change (skills)","- Therapist teaches coping skills","- 12-16 sessions typically","- Structured, directive style","- Works best DURING action stage","- Addresses triggers and behaviors","- Regular homework/practice assigned","- Therapist educates and trains"], bgc='light_green', title="CBT:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["EVIDENCE FOR COMBINING MET + CBT:","- Project MATCH (1997): MET alone = CBT alone in outcomes (4 vs 12 sessions)","- Cannabis Youth Treatment (Dennis et al., 2004): MET/CBT combination most effective","- COMBINE Study (Anton et al., 2006, JAMA): MET as prelude to CBT improved engagement","- Marijuana Treatment Project (Stephens et al., 2004): MET+CBT optimal for cannabis","","RECOMMENDATION: Use MET FIRST (build motivation), then CBT for skills (sequential model)."], bgc='cream', title="Research Evidence:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Project MATCH (1997); Dennis et al. (2004). J Subst Abuse Treat; COMBINE (2006). JAMA, 295.")

# MET vs 12-Step
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_purple'); title_bar(slide, "MET vs 12-Step Facilitation (TSF) / AA", 'purple')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["MET APPROACH:","- No labels required (\"alcoholic\")","- Client sets own goals (abstinence or","  moderation - client's choice)","- Brief (4 sessions)","- Individual therapy","- Therapist is collaborative partner","- No spiritual component","- \"You decide what's best for you\""], bgc='light_purple', title="MET:", ttc='purple', border='purple', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["12-STEP / TSF APPROACH:","- Acceptance of \"alcoholic\" identity","- Goal: total abstinence only","- Ongoing/lifetime involvement","- Group-based (AA/NA meetings)","- Sponsor system for support","- Spiritual foundation (Higher Power)","- \"You are powerless over alcohol\"","- Structured step-work program"], bgc='light_orange', title="12-Step:", ttc='orange', border='orange', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["EVIDENCE:","- Project MATCH (1997): Both EQUALLY effective at 1 and 3-year follow-up","- MET more cost-effective (4 sessions vs 12 for same outcomes)","- MET better for clients HIGH in anger (Karno & Longabaugh, 2005)","- TSF better for clients with high social networks supportive of drinking","","INTEGRATION: Use MET to BUILD motivation, then refer to AA/NA for ongoing community support.","MET addresses \"why change\" while 12-Step provides structure and fellowship."], bgc='cream', title="Research Comparison:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: Project MATCH (1997); Karno & Longabaugh (2005). J Stud Alcohol, 66, 488-495.")

# Comprehensive Comparison Table
tbl("Comprehensive Comparison of Evidence-Based Therapies",
    ["Therapy", "Mechanism", "Duration", "Best For", "Evidence"],
    [["MET", "Internal motivation\n(discrepancy + MI)", "4 sessions", "Low-readiness clients\nAngry clients\nBrief settings", "Level 1\n(MATCH, UKATT)"],
     ["CBT", "Skills training\n(coping, triggers)", "12-16 sessions", "Action-stage clients\nAnxiety comorbidity", "Level 1\n(Carroll, 1994)"],
     ["12-Step (TSF)", "Peer support\nSpiritual growth", "Ongoing", "Abstinence-oriented\nHigh social support", "Level 1\n(MATCH)"],
     ["Contingency Mgmt", "External rewards\nfor abstinence", "12-24 weeks", "Stimulant use\nImmediate reinforcement", "Level 1\n(Higgins, 2004)"],
     ["CRA+CRAFT", "Environmental\nrestructuring", "12-24 sessions", "Family involvement\nSocial isolation", "Level 1\n(Meyers, 1995)"],
     ["BSFT (Family)", "Family system\nrestructuring", "12-16 sessions", "Adolescents\nFamily conflict", "Level 1\n(Szapocznik, 2003)"]],
    color='navy', ref="Ref: NICE Guidelines (2011); SAMHSA TIP 35 (1999); APA Practice Guidelines (2006).")

# Why MET is More Effective
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "Why MET is More Effective Than Traditional Approaches", 'green')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["CLINICAL ADVANTAGES:","","1. COST-EFFECTIVE: 4 sessions = 12 sessions","   (70% less therapist time, same results)","","2. ENGAGES RESISTANT CLIENTS:","   Works for court-mandated, family-pressured","","3. REDUCES DROPOUT:","   Non-confrontational = clients stay longer","   (Carroll et al., 2006)","","4. MATCHES CLIENT'S STAGE:","   Doesn't force action on unready clients"], bgc='light_green', title="Clinical:", ttc='green', border='green', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["SYSTEMIC ADVANTAGES:","","5. SCALABLE:","   Delivered by trained lay counselors","   (Nadkarni et al., 2017 - Lancet)","","6. UNIVERSAL ENHANCER:","   Improves outcomes when added to ANY","   other treatment (CBT, meds, residential)","","7. CROSS-CULTURAL:","   Works across cultures (WHO, 2002)","","8. EVIDENCE FOR ANGRY CLIENTS:","   Significantly better than other Tx"], bgc='light_blue', title="Systemic:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["RESEARCH SUPPORT:","- Angry clients: MET significantly BETTER than CBT/TSF (Project MATCH matching hypothesis)","- Low-readiness clients: MET better at engaging and retaining them in treatment","- Brief settings: MET achievable where 12-session therapy is not feasible","- Stepped care: MET as first step, add intensive treatment only if MET insufficient","","Meta-analysis: Lundahl et al. (2010) - MI/MET effect size d=0.79 for substance use","compared to no treatment. This is a MEDIUM-LARGE effect for a 4-session intervention."], bgc='cream', title="Evidence:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Lundahl et al. (2010). Patient Ed & Counsel; Project MATCH (1997); Nadkarni et al. (2017). Lancet.")

# Family & Caregiver
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "Family & Caregiver Involvement in MET", 'orange')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["ROLE OF FAMILY IN MET:","","MET Manual: A \"concerned significant other\"","can join Session 2 for Change Plan:","","- Provides support for client's goals","- Helps identify triggers at home","- Can reinforce positive changes","- Learns non-confrontational approach","","KEY: Family supports, doesn't direct.","Family uses MI principles too."], bgc='light_orange', title="Family in MET:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["COMPLEMENTARY FAMILY THERAPIES:","","1. CRAFT (Community Reinforcement","   and Family Training):","   - For families of unmotivated users","   - Teaches positive communication","   - 64% engagement rate (Meyers, 2002)","","2. Behavioral Couples Therapy:","   - Recovery contracts","   - Reduces substance use AND","     relationship distress simultaneously"], bgc='light_green', title="Adjunct Family Therapies:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["CAREGIVER PSYCHOEDUCATION (use alongside MET):","- Teach CRAFT principles: reward sobriety, allow natural consequences, self-care","- Avoid enabling behaviors (covering up, making excuses, giving money)","- Understand stages of change (don't expect overnight transformation)","- Learn MI techniques for home conversations (reflective listening, not lecturing)","","Research: Smith & Meyers (2004) - CRAFT achieved 64% treatment engagement vs 30% Al-Anon","and 13% Johnson Intervention (confrontational). Non-confrontation wins for families too."], bgc='cream', title="Caregiver Education:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), p. 65; Meyers et al. (2002). JCCP; Smith & Meyers (2004). CRAFT handbook.")

# Stepped Care / Combined Treatment
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "Combined Treatment: Therapies Used Alongside MET", 'green')
cbox(slide, LM, Inches(1.15), FULL_W, Inches(1.7), ["THE STEPPED CARE MODEL (NIMHANS & international guidelines):","Step 1: MET/Brief MI (4 sessions) - for ALL clients as starting point","Step 2: If insufficient - add CBT/Skills Training (12 sessions)","Step 3: If still insufficient - add Pharmacotherapy + Intensive outpatient","Step 4: If still insufficient - Residential treatment / Therapeutic community"], bgc='light_green', title="Sequential Treatment:", ttc='green', border='green', fs=12)
cbox(slide, LM, Inches(3.05), HALF_W, Inches(3.2), ["FOR THE PATIENT:","","1. MET + Pharmacotherapy:","   - Naltrexone/Acamprosate (alcohol)","   - Buprenorphine (opioids)","   - COMBINE (2006): MI + naltrexone effective","","2. MET + CBT (sequential):","   - MET first for motivation (1-4)","   - CBT next for skills (5-16)","","3. MET + Contingency Management:","   - External + internal motivation combined"], bgc='light_blue', title="For Patient:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, Inches(3.05), HALF_W, Inches(3.2), ["FOR FAMILY & CAREGIVERS:","","1. CRAFT Training (6-12 sessions)","   Positive communication, consequences","","2. Behavioral Family Therapy:","   Communication skills, problem-solving","","3. Family Psychoeducation:","   Understanding addiction as illness","   Reducing expressed emotion (EE)","","4. Al-Anon/Nar-Anon support groups","","5. Caregiver self-care programs"], bgc='light_purple', title="For Family:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: NIMHANS (2016), Ch. 8-10; COMBINE (2006). JAMA; Meyers et al. (2002); NICE (2011).")

# MET as Pretreatment Enhancer
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_teal'); title_bar(slide, "MET as a 'Pretreatment' Enhancer for Other Therapies", 'teal')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["MET AS PRELUDE TO OTHER TREATMENTS:","","Research shows MET BEFORE other therapies:","","1. Increases treatment ENGAGEMENT","   (clients attend more sessions)","2. Reduces DROPOUT rates","3. Improves OUTCOMES of subsequent Tx","4. Increases MEDICATION adherence","5. Enhances CBT, DBT, residential Tx","","\"MI as a prelude to treatment is one of","its most promising applications\"","(Hettema et al., 2005)"], bgc='light_teal', title="MET as Enhancer:", ttc='teal', border='teal', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["RESEARCH EVIDENCE:","","- Carroll et al. (2006): MI before Tx =","  better retention in drug treatment","- Westra & Dozois (2006): MI before CBT","  enhanced outcomes for anxiety","- Martino et al. (2007): MI improved","  engagement in psychiatric treatment","- Swanson et al. (1999): MI before","  residential Tx = longer stays","- Daley et al. (1998): MI + 12-Step =","  better than 12-Step alone","","Cost: Only 1-2 extra sessions needed"], bgc='light_green', title="Studies:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["CLINICAL IMPLICATION: Even if your primary treatment modality is CBT, DBT, 12-Step, or","residential treatment, STARTING with 1-2 sessions of MET/MI will improve outcomes.","MET is not just a standalone treatment - it's a universal treatment ENHANCER.","","This is particularly important for comorbid patients who need multiple therapies.","Start with MET to build motivation, then add the disorder-specific intervention.","Research: Hettema et al. (2005) - MI as pretreatment had the largest effect sizes in their review."], bgc='cream', title="Take-Home Message:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Carroll et al. (2006). Drug Alcohol Depend; Hettema et al. (2005). Ann Rev Clin Psychol, 1, 91-111.")



# ============================================================
# SECTION 9: COMORBID PSYCHIATRIC DISORDERS (NEW - ~8 slides)
# ============================================================
divider("SECTION 9", "MET with Comorbid Psychiatric Disorders\nWhat to Do When Co-Occurring Conditions Exist", 'purple')

# Overview
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_purple'); title_bar(slide, "Dual Diagnosis: Substance Use + Psychiatric Disorders", 'purple')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["PREVALENCE OF COMORBIDITY:","","- 50-70% of substance users have at least","  one comorbid psychiatric disorder","  (Regier et al., 1990; Kessler, 2004)","- Common comorbidities:","  * Depression (30-50%)","  * Anxiety disorders (25-40%)","  * Personality disorders (40-70%)","  * PTSD (25-40%)","  * Bipolar disorder (20-30%)"], bgc='light_purple', title="The Problem:", ttc='purple', border='purple', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["CAN MET BE USED WITH COMORBIDITY?","","YES - with adaptations:","- MI/MET principles applicable regardless","  of comorbid diagnosis","- Tested in dual diagnosis populations","- Key: INTEGRATE, don't separate treatment","","\"Integrated treatment is more effective","than parallel or sequential treatment\"","(Drake et al., 2004, Schizophr Bull)"], bgc='light_green', title="MET Applicability:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["GENERAL PRINCIPLES FOR MET WITH COMORBIDITY:","1. TREAT BOTH SIMULTANEOUSLY (integrated approach superior to sequential)","2. Acknowledge substance use may be SELF-MEDICATION for psychiatric symptoms","3. Explore FUNCTION of use: \"What does drinking do for your anxiety/depression?\"","4. Develop discrepancy: \"You use to feel better, but does it make things better long-term?\"","5. Adapt pace - comorbid clients may need more sessions or slower progression","6. Coordinate with prescribing psychiatrist for pharmacotherapy decisions"], bgc='cream', title="General Approach:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Drake et al. (2004). Schizophr Bull, 30; Kessler (2004). Biol Psych; NIMHANS (2016), Ch. 12.")

# Personality Disorders
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_red'); title_bar(slide, "MET with Comorbid Personality Disorders", 'red')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["CHALLENGES:","- Interpersonal difficulties affect rapport","- Emotional dysregulation (BPD)","- Impulsivity increases relapse risk","- Identity disturbance complicates values work","- Therapist-client dynamic more complex","- Higher dropout rates","","PREVALENCE: 40-70% of substance users","meet criteria for at least one PD","(Verheul, 2001; NIMHANS data)"], bgc='light_red', title="PD + Substance Use:", ttc='red', border='red', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["MET ADAPTATIONS FOR PD:","","1. LONGER ENGAGEMENT PHASE","   May need extra sessions before feedback","2. EMPHASIS ON ROLLING WITH RESISTANCE","   PD clients show more resistance","3. VALIDATE EMOTIONS before exploring change","   (especially for Borderline PD)","4. SMALLER CHANGE GOALS (reduce overwhelm)","5. CONSISTENT boundaries + empathy","6. Monitor for therapeutic relationship ruptures"], bgc='light_green', title="Adaptations:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["RESEARCH EVIDENCE:","- Ball et al. (2007): Dual Focus Schema Therapy + MI for PD+SUD - improved outcomes","  vs standard counseling (J Nervous & Mental Disease, 195, 24-31)","- Bornovalova & Daughters (2007): MI effective for Borderline PD + SUD with distress tolerance","- Gregory et al. (2008): MI feasible with Antisocial PD (commonly seen in substance users)","- NIMHANS recommendation: MI/MET as engagement strategy for PD patients, then transition","  to DBT (Borderline) or Schema Therapy (Cluster B) for long-term personality work"], bgc='cream', title="Research:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: Ball et al. (2007). J Nerv Ment Dis; Bornovalova & Daughters (2007); Verheul (2001).")

# Depression
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "MET with Comorbid Depression", 'deep_blue')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["DEPRESSION + SUBSTANCE USE:","","- 30-50% comorbidity rate","- Alcohol is CNS depressant - worsens mood","- Substance-induced vs independent depression","- Both need addressing simultaneously","","CHALLENGES FOR MET:","- Low energy/motivation (core of depression)","- Hopelessness affects self-efficacy","- Suicidal ideation must be monitored","- Cognitive distortions complicate reasoning"], bgc='light_blue', title="The Challenge:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["MET ADAPTATIONS FOR DEPRESSION:","","1. EXTRA EMPHASIS on self-efficacy building","   (counters hopelessness of depression)","2. Explore MOOD-SUBSTANCE link:","   \"When do you drink most? When low?\"","3. Develop discrepancy with MOOD:","   \"You drink to feel better, but next day","    your mood is even worse.\"","4. Coordinate antidepressant medication","5. Safety planning if suicidal ideation","6. Smaller goals (energy is limited)"], bgc='light_green', title="Adaptations:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["RESEARCH EVIDENCE:","- Baker et al. (2010): MI + CBT for comorbid depression + SUD - significantly reduced both","  depressive symptoms AND substance use (Addiction, 105, 1560-1568)","- Satre et al. (2016): MI effective for older adults with depression + alcohol use","- Westra & Dozois (2006): MI as prelude to CBT for depression enhanced engagement","- Rao et al. (2015, NIMHANS): MI-based brief intervention effective for depression + alcohol","  in Indian primary care settings (Asian J Psychiatry, 16, 45-51)"], bgc='cream', title="Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Baker et al. (2010). Addiction; Westra & Dozois (2006). JCCP; Rao et al. (2015). Asian J Psych.")

# Anxiety & PTSD
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_teal'); title_bar(slide, "MET with Comorbid Anxiety Disorders & PTSD", 'teal')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["ANXIETY + SUBSTANCE USE:","- 25-40% comorbidity rate","- Substances \"self-medicate\" anxiety","- Alcohol reduces anxiety short-term,","  increases it long-term (rebound)","- Withdrawal mimics anxiety symptoms","","PTSD + SUBSTANCE USE:","- 25-40% of PTSD patients have SUD","- Substances numb trauma memories","- Integrated treatment essential"], bgc='light_teal', title="The Challenge:", ttc='teal', border='teal', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["MET ADAPTATIONS:","","For ANXIETY:","1. Validate anxiety as real (not excuse)","2. Explore substance-anxiety cycle","3. Include anxiety management in Change Plan","4. Confidence building is crucial","","For PTSD:","1. Safety first - stabilization before change","2. Respect avoidance (don't push trauma)","3. MET for engagement, then EMDR/PE","4. Trauma-informed MI adaptation"], bgc='light_green', title="Adaptations:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["RESEARCH EVIDENCE:","- Hien et al. (2009): Integrated MI + Seeking Safety for PTSD+SUD - effective in reducing","  both trauma symptoms AND substance use (J Consult Clin Psychol, 77, 607-619)","- Westra (2012): MI enhanced CBT outcomes for anxiety disorders","- Najavits (2002): Seeking Safety + MI as integrated approach for trauma + substance use","- Sannibale et al. (2013): MI + exposure therapy for alcohol+PTSD - feasible and effective","- Bolton et al. (2016, Lancet Psychiatry): MI effective in low-resource settings for comorbidity"], bgc='cream', title="Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Hien et al. (2009). JCCP; Westra (2012). Addictive Behaviors; Sannibale et al. (2013). Addiction.")

# Bipolar & Psychosis
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "MET with Bipolar Disorder & Psychotic Disorders", 'orange')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["BIPOLAR + SUBSTANCE USE:","- 40-60% lifetime SUD comorbidity","- Mania: impulsivity increases use","- Depression: self-medication pattern","- Non-adherence to mood stabilizers","","PSYCHOSIS + SUBSTANCE USE:","- 40-50% have comorbid SUD","- Cannabis, alcohol most common","- Complicates symptom management","- Reduced insight may limit engagement"], bgc='light_orange', title="The Challenge:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["MET ADAPTATIONS:","","For BIPOLAR:","1. Address medication adherence with MI","2. Discrepancy: substance triggers episodes","3. Change Plan includes mood monitoring","4. Extra self-efficacy during depressive phases","","For PSYCHOSIS:","1. Simplify language and concepts","2. Focus on immediate, concrete goals","3. Shorter sessions if attention limited","4. Affirm heavily (low self-esteem common)","5. Family involvement especially important"], bgc='light_green', title="Adaptations:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["RESEARCH EVIDENCE:","- Barrowclough et al. (2001): MI + CBT for schizophrenia + SUD - significant reduction","  in substance use over 18 months (BMJ, 323, 1-5). Landmark dual diagnosis trial.","- Graeber et al. (2003): MI for psychosis + SUD improved engagement and retention","- Kemp et al. (1996): Compliance Therapy (MI-based) for bipolar medication adherence","- Baker et al. (2005): MI effective for comorbid psychosis + cannabis (Addiction, 100, 1614)","- Haddock et al. (2003): Integrated MI+CBT for psychosis+SUD in community mental health"], bgc='cream', title="Research:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: Barrowclough et al. (2001). BMJ; Baker et al. (2005). Addiction; Kemp et al. (1996).")

# ADHD & Eating Disorders
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_purple'); title_bar(slide, "MET with ADHD & Other Comorbid Conditions", 'purple')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["ADHD + SUBSTANCE USE:","- 25-40% of adults with SUD have ADHD","- Impulsivity is a shared vulnerability","- Stimulant medications may be needed","- Short attention span affects session format","","MET ADAPTATIONS FOR ADHD:","- Shorter, more focused sessions","- More visual aids and written summaries","- Repeated summaries (attention issues)","- Address impulsivity as shared target","- \"How does using affect your focus?\""], bgc='light_purple', title="ADHD + SUD:", ttc='purple', border='purple', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["EATING DISORDERS + SUBSTANCE USE:","- 20-35% comorbidity rate","- Shared features: impulsivity, affect","  regulation difficulties, body image","","MET ADAPTATIONS FOR ED:","- Address both behaviors in Change Plan","- Explore function of BOTH (affect regulation)","- Validate ambivalence about ED AND substance","- Body image discrepancy can be explored","","GAMBLING DISORDER:","- MI has Level 1 evidence for gambling","  (Hodgins et al., 2001; Petry et al., 2008)"], bgc='light_green', title="Other Comorbidities:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["RESEARCH EVIDENCE:","- Zulauf et al. (2014): MI effective for ADHD+SUD when sessions adapted for attention","- Treasure & Schmidt (2001): MI for eating disorders - MI effective for treatment ambivalence","- Hodgins et al. (2001): MI reduced gambling in a brief telephone intervention","- Petry et al. (2008): MI + CBT for pathological gambling - significant improvements","- van Emmerik-van Oortmerssen et al. (2012): ADHD screening in SUD populations recommended","","PRINCIPLE: MI/MET is adaptable to ANY behavioral change target - the principles are universal."], bgc='cream', title="Research:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Zulauf et al. (2014). Curr Drug Abuse Rev; Hodgins et al. (2001). JCCP; Treasure & Schmidt (2001).")

# Summary Table
tbl("Summary: MET Adaptations for All Comorbid Disorders",
    ["Comorbid Disorder", "Key MET Adaptation", "Adjunct Therapy", "Key Research"],
    [["Personality Disorders", "Extra alliance building\nValidation before change\nSmaller goals", "DBT (Borderline)\nSchema Therapy\nMentalisation", "Ball et al. (2007)\nBornovalova (2007)"],
     ["Major Depression", "Extra self-efficacy\nMood-substance link\nSafety planning", "Antidepressants\nCBT for Depression\nBehavioral Activation", "Baker et al. (2010)\nRao et al. (2015)"],
     ["Anxiety Disorders", "Validate anxiety\nExplore anxiety cycle\nConfidence building", "SSRIs/SNRIs\nCBT for Anxiety\nRelaxation Training", "Westra (2012)\nWestra & Dozois (2006)"],
     ["PTSD", "Safety first\nTrauma-informed MI\nRespect avoidance", "EMDR / PE\nSeeking Safety\nCPT", "Hien et al. (2009)\nSannibale (2013)"],
     ["Bipolar Disorder", "Med adherence with MI\nLink use to episodes\nMood monitoring", "Mood Stabilizers\nIPSRT\nPsychoeducation", "Kemp et al. (1996)\nGraeber (2003)"],
     ["Psychotic Disorders", "Simplify language\nConcrete goals\nShorter sessions", "Antipsychotics\nIntegrated dual Dx\nFamily intervention", "Barrowclough (2001)\nBaker (2005)"]],
    color='purple', ref="Ref: All cited in individual slides. Review: Drake et al. (2004). Schizophr Bull, 30, 795-808.")



# ============================================================
# SECTION 10: RESEARCH & EFFECTIVENESS
# ============================================================
divider("SECTION 10", "Research Evidence & Effectiveness\nHow Well Does MET Work?", 'navy')

# Project MATCH
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "Project MATCH: The Landmark Trial", 'navy')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["STUDY DESIGN:","- Largest alcohol treatment trial ever","- 1,726 participants, 9 US sites","- Randomized to 3 conditions:","  1. MET (4 sessions / 12 weeks)","  2. CBT (12 sessions)","  3. TSF/12-Step (12 sessions)","- Follow-up: 1 year and 3 years","- Cost: $27 million (NIAAA funded)"], bgc='light_blue', title="The Study:", ttc='navy', border='navy', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["KEY FINDINGS:","- All three treatments significantly improved","- MET achieved COMPARABLE outcomes to","  CBT and TSF in just 4 sessions vs 12","- Percent Days Abstinent improved in ALL","- Drinks Per Drinking Day decreased in ALL","- At 3-year follow-up: gains maintained","","IMPLICATION: MET is as effective as longer","treatments = highly COST-EFFECTIVE"], bgc='light_green', title="Results:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["MATCHING FINDING: Clients HIGH in anger did significantly BETTER in MET than other treatments.","Non-confrontational approach works better for angry, reactive clients.","","\"The finding that a 4-session motivational intervention could produce outcomes comparable to","12-session treatments had profound implications for cost-effectiveness.\" (MATCH Group, 1997)","","India relevance: With limited therapist availability, a 4-session effective treatment is ideal","for resource-constrained settings (NIMHANS recommends MET as first-line, 2016)."], bgc='cream', title="Significance:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Project MATCH (1997). J Stud Alcohol, 58, 7-29; (1998). Addiction, 93, 1434-1447.")

# Meta-analyses Table
tbl("Meta-Analyses: MI/MET Effectiveness Across Studies",
    ["Study", "N Studies", "Key Finding", "Effect Size"],
    [["Burke et al. (2003)\nJ Consult Clin Psychol", "30 RCTs", "MI/MET effective for alcohol,\ndrugs, diet, adherence", "d = 0.25-0.57"],
     ["Hettema et al. (2005)\nAnn Rev Clin Psychol", "72 studies", "MI effective across substances\nvs no treatment/advice", "d = 0.77 at follow-up"],
     ["Lundahl et al. (2010)\nPatient Ed & Counsel", "119 studies", "MI moderate effect;\nstronger for substance use", "d = 0.22 (combined)\nd = 0.79 (substance)"],
     ["Vasilaki et al. (2006)\nAddiction", "15 RCTs", "Brief MI effective for\nreducing alcohol consumption", "Significant reduction"],
     ["Smedslund et al. (2011)\nCochrane Review", "59 RCTs", "MI reduces substance use\nmore than no treatment", "SMD = -0.79"]],
    color='navy', ref="Ref: All peer-reviewed meta-analyses of MI/MET interventions as cited.")

# Indian Research
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "Research Evidence: Indian Context", 'orange')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["NIMHANS EVIDENCE:","- Brief MI effective in Indian primary care","  (Murthy et al., 2009)","- NIMHANS protocols integrate MET as","  first-line psychosocial intervention","- Community MI reduced alcohol in rural","  Karnataka (Nadkarni et al., 2017)","","AIIMS, DELHI:","- Brief MI for alcohol dependence","  positive outcomes (Pal et al., 2007)"], bgc='light_orange', title="Indian Studies:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["PREMIUM TRIAL (Lancet, 2017):","- Lay counselors delivered MI in Goa","- Significant reduction in harmful drinking","- Published in The Lancet","- Proved task-shifting feasible in India","","CHAND et al. (2018, NIMHANS):","- Technology-assisted brief intervention","- Showed feasibility in Indian clinics","","OUTCOME: MI/MET effective in India"], bgc='light_green', title="Key Studies:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["CULTURAL CONSIDERATIONS FOR INDIA:","- Family involvement crucial (integrate into MET Change Plan)","- High stigma - MI's non-judgmental approach especially valuable","- Task-shifting to lay counselors feasible (PREMIUM trial conclusive)","- Adaptations: family values, community roles, spiritual beliefs in values work","- NIMHANS (2016) recommends MI/MET as evidence-based first-line intervention","- Cost-effective for India's resource-constrained mental health system"], bgc='cream', title="Indian Adaptations:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: Nadkarni et al. (2017). Lancet, 389; NIMHANS (2016); Chand et al. (2018). Indian J Psych.")

# UKATT
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "UKATT & Global Evidence Summary", 'green')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["UKATT (UK Alcohol Treatment Trial):","- Largest UK trial (742 clients)","- MET (3 sessions) vs SBNT (8 sessions)","- Result: BOTH equally effective","- MET significantly more cost-effective","- 5x less therapist time, same results","","CONCLUSION: MET achieves equivalent","outcomes at a fraction of the cost","(UKATT Research Team, 2005, BMJ)"], bgc='light_green', title="UKATT:", ttc='green', border='green', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["OTHER INTERNATIONAL EVIDENCE:","- WHO (2002): MI effective in 10 countries","- COMBINE (2006, JAMA): MI + naltrexone","- Cannabis Youth Treatment (2004):","  MET effective for adolescent cannabis","- Stephens et al. (2004):","  2 sessions MI = 6 sessions CBT","","MI/MET works across cultures, substances,","age groups, and clinical settings."], bgc='light_blue', title="Global:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["EVIDENCE SUMMARY:","- EMPIRICALLY SUPPORTED for alcohol use disorders (Level 1 evidence)","- Effective for cannabis, cocaine, opioids, polysubstance use","- Works with adolescents, adults, older adults across cultures","- Effective when delivered by trained lay counselors","- Cost-effective: same results as longer treatments","- Compatible with pharmacotherapy and enhances medication adherence","- Universal treatment enhancer - improves other therapies when added"], bgc='cream', title="Bottom Line:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: UKATT (2005). BMJ, 331; COMBINE (2006). JAMA, 295; WHO (2002); Stephens et al. (2004).")



# ============================================================
# CLOSING SLIDES
# ============================================================

# Key Takeaways
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'cream'); title_bar(slide, "Key Takeaways: MET in Clinical Practice", 'deep_blue')
for i,(t,d,bc,tc2) in enumerate([("MET is BRIEF but POWERFUL","4 sessions achieve outcomes of 12-session treatments",'light_blue','deep_blue'),("Based on SOLID THEORY","Transtheoretical Model + Self-Efficacy + Cognitive Dissonance",'light_green','green'),("The CLIENT does the work","Therapist evokes motivation, doesn't install it",'light_purple','purple'),("RESISTANCE is information","Not to fight - a signal to change approach",'light_orange','orange'),("Works with COMORBIDITY","Adaptable for PD, depression, anxiety, PTSD, psychosis",'light_teal','teal'),("EVIDENCE-BASED globally","Project MATCH, UKATT, PREMIUM trial, 100+ RCTs",'light_red','maroon')]):
    cbox(slide, LM, Inches(1.15)+Inches(0.88)*i, FULL_W, Inches(0.8), [d], bgc=bc, title=t, ttc=tc2, border=tc2, fs=12)
ref_bar(slide, "Ref: Miller et al. (1992); Project MATCH (1997); UKATT (2005); NIMHANS (2016); Nadkarni et al. (2017).")

# References
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "Comprehensive References", 'deep_blue')
cbox(slide, LM, Inches(1.15), HALF_W, Inches(5.1), ["PRIMARY SOURCES:","Miller, W.R. et al. (1992). MET Manual. NIAAA.","NIMHANS (2016). Substance Use Disorders Manual.","Miller & Rollnick (2013). MI, 3rd ed. Guilford.","Prochaska & DiClemente (1984). TTM.","Project MATCH (1997). J Stud Alcohol, 58.","UKATT (2005). BMJ, 331, 544.","","COMORBIDITY:","Barrowclough et al. (2001). BMJ, 323.","Baker et al. (2010). Addiction, 105.","Ball et al. (2007). J Nerv Ment Dis, 195.","Hien et al. (2009). JCCP, 77.","Drake et al. (2004). Schizophr Bull, 30.","Westra & Dozois (2006). JCCP, 74."], bgc='white', border='deep_blue', fs=10)
cbox(slide, COL2, Inches(1.15), HALF_W, Inches(5.1), ["INDIAN RESEARCH:","Nadkarni et al. (2017). Lancet, 389.","Chand et al. (2018). Indian J Psychiatry.","Rao et al. (2015). Asian J Psychiatry, 16.","Murthy et al. (2009). NIMHANS.","","META-ANALYSES:","Lundahl et al. (2010). Pat Ed Counsel, 80.","Hettema et al. (2005). Ann Rev Clin Psychol.","Smedslund et al. (2011). Cochrane Review.","Burke et al. (2003). JCCP, 71.","","FAMILY:","Smith & Meyers (2004). CRAFT.","Meyers et al. (2002). JCCP, 70.","Carroll et al. (2006). Drug Alcohol Depend."], bgc='white', border='green', fs=10)
ref_bar(slide, "All from peer-reviewed journals, published manuals, and institutional publications.")

# Thank You
slide = prs.slides.add_slide(blank_layout); bg(slide, 'deep_blue', 'navy')
tb = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(9), Inches(4))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "Thank You"; p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = COLORS['white']; p.font.name = 'Times New Roman'; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = ""; p2.space_after = Pt(30)
p3 = tf.add_paragraph(); p3.text = "\"People are generally better persuaded by the reasons"; p3.font.size = Pt(16); p3.font.italic = True; p3.font.color.rgb = COLORS['gold']; p3.font.name = 'Times New Roman'; p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph(); p4.text = "which they have themselves discovered"; p4.font.size = Pt(16); p4.font.italic = True; p4.font.color.rgb = COLORS['gold']; p4.font.name = 'Times New Roman'; p4.alignment = PP_ALIGN.CENTER
p5 = tf.add_paragraph(); p5.text = "than by those which have come into the minds of others.\""; p5.font.size = Pt(16); p5.font.italic = True; p5.font.color.rgb = COLORS['gold']; p5.font.name = 'Times New Roman'; p5.alignment = PP_ALIGN.CENTER
p6 = tf.add_paragraph(); p6.text = "\u2014 Blaise Pascal"; p6.font.size = Pt(13); p6.font.color.rgb = COLORS['light_blue']; p6.font.name = 'Times New Roman'; p6.alignment = PP_ALIGN.CENTER
p7 = tf.add_paragraph(); p7.text = ""; p7.space_after = Pt(20)
p8 = tf.add_paragraph(); p8.text = "Based on: MET Manual (NIAAA, 1992) & NIMHANS Substance Use Disorders (2016)"; p8.font.size = Pt(11); p8.font.color.rgb = COLORS['light_blue']; p8.font.name = 'Times New Roman'; p8.alignment = PP_ALIGN.CENTER

# (save moved to end of file)



# ============================================================
# ADDITIONAL SLIDES TO REACH ~80 (inserted before SAVE)
# ============================================================

# Move save to end - these slides will be appended

# Additional Session 2 Detail slide
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "SESSION 2: Clinical Dialogue - Change Plan Development", 'green')
cbox(slide, LM, Inches(1.15), FULL_W, Inches(5.1), ["DEVELOPING THE CHANGE PLAN - DIALOGUE:","","T: \"Last time we talked about how your drinking has been affecting your family life.","   You said you wanted things to be different. Would you like to put together a plan?\"","C: \"Yeah, I think so. But I don't know where to start.\"","T: \"That's okay. Let's think about this together. What changes would you like to make?\"","C: \"I guess... I want to cut down. Maybe not drink during the week.\"","T: \"Weekday abstinence - that's a specific, achievable goal. What's the most important","   reason for making that change?\" (Eliciting commitment language)","C: \"My kids. They deserve a better father in the evenings.\"","T: \"Being present for your children is what matters most to you.\" (Affirming value)","C: \"Yeah. And I can't be present if I'm drinking.\"","T: \"You've connected the dots between drinking and missing out on what matters.\"","   (Reflecting change talk - the client is arguing for change, not the therapist)","","T: \"What steps could you take to make weekday abstinence work?\"","C: \"Maybe not keep beer in the house... and find something else to do in the evening.\"","T: \"You already have two concrete ideas. That's great. What might get in the way?\""], bgc='white', border='green', fs=11)
ref_bar(slide, "Ref: MET Manual (1992), pp. 63-68; Dialogue based on Change Plan Worksheet protocol.")

# Additional: How to Present Normative Feedback
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "Presenting Normative Feedback: A Key MET Technique", 'deep_blue')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["WHAT IS NORMATIVE FEEDBACK?","","Showing the client WHERE they fall","compared to the general population:","","\"Of all adults in India, here's where","your drinking falls...\"","[Show graph: 95th percentile]","","\"95% of adults drink LESS than you.\"","","This creates discrepancy without arguing.","The DATA does the work."], bgc='light_blue', title="The Technique:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["HOW TO DELIVER EFFECTIVELY:","","1. Present the data NEUTRALLY","   (not with alarm or judgment)","2. Use visual aids (graphs, charts)","3. PAUSE after presenting - let it land","4. Ask: \"What do you make of this?\"","5. Reflect their reaction","6. Do NOT push if they minimize","   (they'll think about it later)","","\"Let the data speak for itself\""], bgc='light_green', title="Delivery Tips:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["AREAS FOR NORMATIVE COMPARISON:","- Drinking quantity vs population norms (percentile ranking)","- Number of alcohol-related problems vs typical drinkers","- Blood alcohol concentration estimates vs legal/medical limits","- Liver enzyme levels vs normal ranges","- Money spent on substances vs average household spending","","Research: Neighbors et al. (2004) showed personalized normative feedback alone reduced heavy","drinking in college students. When combined with MI, effects were even stronger."], bgc='cream', title="What to Compare:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 44-50; Neighbors et al. (2004). J Consult Clin Psychol, 72, 434-442.")

# Additional: Stages of Change assessment
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_purple'); title_bar(slide, "Assessing Stage of Change: Practical Guide", 'purple')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["HOW TO ASSESS READINESS:","","Use the READINESS RULER:","\"On a scale of 1-10, how ready are you","to make a change right now?\"","","1-3 = Precontemplation","  (not ready, don't push)","4-6 = Contemplation","  (thinking about it, explore ambivalence)","7-10 = Preparation/Action","  (ready, develop Change Plan)"], bgc='light_purple', title="Quick Assessment:", ttc='purple', border='purple', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["LANGUAGE CLUES BY STAGE:","","PRECONTEMPLATION:","\"I don't have a problem\" \"Everyone does it\"","","CONTEMPLATION:","\"Maybe\" \"I've been thinking\" \"I'm not sure\"","\"Part of me wants to\" \"But on the other hand\"","","PREPARATION:","\"I need to\" \"I'm going to\" \"How do I?\"","\"What are my options?\" \"I've decided\""], bgc='light_green', title="Listening for Stage:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["MATCHING YOUR RESPONSE TO THEIR STAGE:","- Precontemplation: Raise awareness gently, provide feedback, plant seeds","- Contemplation: Explore ambivalence, use decisional balance, evoke change talk","- Preparation: Develop Change Plan, offer menu of options, set goals","- Action: Affirm, troubleshoot, review plan, build confidence","- Maintenance: Identify risks, plan for slips, celebrate, consolidate","","CARDINAL RULE: Never be ahead of the client. Match YOUR energy to THEIR readiness."], bgc='cream', title="Stage-Matched Responses:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: MET Manual (1992), pp. 6-14; DiClemente & Velasquez (2002). Ch. in MI 2nd ed.")

# Additional: Evidence for MET with specific populations
tbl("MET/MI Evidence by Specific Population",
    ["Population", "Setting", "Key Study", "Finding"],
    [["Heavy-drinking college students", "University counseling", "Marlatt et al. (1998)\nBAMI Study", "Single MI session reduced\nharm at 2-year follow-up"],
     ["Pregnant women", "Prenatal care", "Handmaker et al. (1999)\nProject CHOICES", "MI reduced drinking in\npregnancy significantly"],
     ["Adolescents (cannabis)", "Outpatient", "Dennis et al. (2004)\nCYT Study", "MET/CBT5 most effective\nbrief intervention"],
     ["HIV+ substance users", "HIV clinic", "Parsons et al. (2007)\nProject PLUS", "MI reduced substance use\nand improved adherence"],
     ["Emergency department", "Hospital ED", "Bernstein et al. (2005)\nProject ASSERT", "Single MI session reduced\ncocaine use at 6 months"],
     ["Older adults (65+)", "Primary care", "Satre et al. (2016)", "MI effective for reducing\nhazardous drinking"],
     ["Inmates/prisoners", "Correctional", "Stein et al. (2006)", "MI reduced post-release\nsubstance use"]],
    color='teal', ref="Ref: All studies cited in table. Review: Hettema et al. (2005). Ann Rev Clin Psychol, 1, 91-111.")

# Additional: Clinical Skills Checklist
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_teal'); title_bar(slide, "MET Therapist Skills Checklist (Self-Assessment)", 'teal')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["MI-CONSISTENT BEHAVIORS:","(aim to maximize these)","","[ ] Open-ended questions","[ ] Simple reflections","[ ] Complex reflections","[ ] Affirmations of client strengths","[ ] Summaries that capture both sides","[ ] Emphasis on personal choice","[ ] Supporting self-efficacy","[ ] Exploring ambivalence","[ ] Developing discrepancy gently"], bgc='light_teal', title="DO MORE:", ttc='teal', border='teal', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["MI-INCONSISTENT BEHAVIORS:","(aim to minimize these)","","[ ] Advising without permission","[ ] Confronting or arguing","[ ] Directing or ordering","[ ] Warning or threatening","[ ] Moralizing or preaching","[ ] Judging or criticizing","[ ] Labeling (\"you're an alcoholic\")","[ ] Overloading with information","[ ] Closed questions (yes/no)"], bgc='light_red', title="DO LESS:", ttc='red', border='red', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["QUALITY INDICATORS (from research):","- Reflection-to-question ratio: aim for 2:1 or higher","- % Open questions: aim for >70% of all questions","- % MI-consistent behaviors: aim for >90% of all responses","- % Complex reflections: aim for >50% of all reflections","","Research: Moyers et al. (2005) showed these ratios predicted client outcomes.","The MITI (Motivational Interviewing Treatment Integrity) scale measures these in supervision."], bgc='cream', title="Quality Measures:", ttc='purple', border='purple', fs=12)
ref_bar(slide, "Ref: Moyers et al. (2005). MITI Manual; MET Manual (1992); Miller & Rollnick (2013). MI 3rd ed.")

# Additional: Training & Supervision
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_green'); title_bar(slide, "Training & Supervision in MET/MI", 'green')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["HOW TO LEARN MET/MI:","","1. Read the MET Manual (freely available)","2. Attend a 2-day MI workshop","3. Practice with role-plays and feedback","4. Record sessions and review with supervisor","5. Use MITI coding for self-assessment","6. Join MI learning community (MINT)","","KEY INSIGHT: Reading alone is insufficient.","MI is a SKILL that requires practice and","feedback, like learning a musical instrument."], bgc='light_green', title="Learning Path:", ttc='green', border='green', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["SUPERVISION RECOMMENDATIONS:","","- Regular session recording review","- Use MITI (MI Treatment Integrity) scale","- Focus on reflection-to-question ratio","- Identify MI-inconsistent behaviors","- Practice difficult scenarios (resistance)","","FIDELITY MATTERS:","Research shows MI delivered POORLY can","actually HARM outcomes. Quality matters","more than quantity of sessions.","(Miller & Rollnick, 2014)"], bgc='light_blue', title="Supervision:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["EVIDENCE ON TRAINING:","- Miller et al. (2004): Workshop alone insufficient; need ongoing feedback/coaching","- Martino et al. (2008): Supervision with feedback improved MI skill retention significantly","- NIMHANS: Has structured MI training programs for Indian mental health professionals","- Task-shifting: Nadkarni et al. (2017) showed lay counselors can learn MI effectively","","RECOMMENDATION: Minimum training = 16-hour workshop + 6 months supervised practice","with regular recording review. This applies in both Indian and international contexts."], bgc='cream', title="Research on Training:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Miller et al. (2004). Drug Alcohol Depend; Martino et al. (2008). JCCP; NIMHANS (2016).")

# Additional: When NOT to use MET
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_red'); title_bar(slide, "Limitations & When MET May Not Be Sufficient Alone", 'red')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["WHEN MET ALONE MAY BE INSUFFICIENT:","","- Severe dependence requiring detox","- Active psychosis (need stabilization first)","- Severe cognitive impairment","- Immediate safety risks (suicidal crisis)","- Client already in Action stage","  (may need skills, not motivation)","- Severe social instability (homeless)","","NOTE: Even in these cases, MI SPIRIT","should inform all clinical interactions."], bgc='light_red', title="Limitations:", ttc='red', border='red', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["MET IS NOT APPROPRIATE FOR:","","- Clients who need SKILLS training","  (already motivated, need HOW not WHY)","  -> Refer to CBT after MET","","- Emergency situations requiring","  immediate medical intervention","  -> Stabilize first, then MET","","- Mandatory abstinence settings where","  client choice is not an option","  -> MI spirit still applies, but autonomy","    support is constrained"], bgc='light_orange', title="Not Appropriate:", ttc='orange', border='orange', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["IMPORTANT CAVEAT:","Even when MET alone is insufficient, its PRINCIPLES should inform all clinical interactions.","The MI spirit (collaboration, evocation, autonomy) is beneficial in ANY therapeutic relationship.","","\"MI is not just a technique - it is a way of being with clients that enhances all treatment.\"","(Miller & Rollnick, 2013)","","Clinical recommendation: Use MET principles as the foundation, add specific interventions","(CBT, DBT, medication, residential) on top as needed based on severity and presentation."], bgc='cream', title="Important:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: Miller & Rollnick (2013). MI 3rd ed.; NICE Guidelines (2011); NIMHANS (2016), Ch. 8.")



# Additional: Motivational Enhancement for Treatment Adherence
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_blue'); title_bar(slide, "MET for Medication & Treatment Adherence", 'deep_blue')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["MI FOR MEDICATION ADHERENCE:","","MI techniques can address:","- Ambivalence about taking medication","- Side effect concerns","- Beliefs about not needing meds","- Forgetting/practical barriers","","Approach:","\"What concerns you about the medication?\"","\"What would be different if you took it?\"","Explore, don't lecture."], bgc='light_blue', title="The Problem:", ttc='deep_blue', border='deep_blue', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["EVIDENCE:","","- Kemp et al. (1996): Compliance Therapy","  (MI-based) improved medication adherence","  in psychosis patients","- Parsons et al. (2007): MI improved","  antiretroviral adherence in HIV+ users","- Zweben & Zuckoff (2002): MI for","  pharmacotherapy engagement review","","PRINCIPLE: MI explores reasons TO adhere","without pressuring. Client decides."], bgc='light_green', title="Research:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["CLINICAL APPLICATION:","- Session 4 of MET can address medication as part of maintenance planning","- For comorbid patients on antidepressants/mood stabilizers/antipsychotics","- Use same MI techniques: explore ambivalence, develop discrepancy, support autonomy","- \"What role could medication play in your plan?\" (not \"You need to take your meds\")","","Research: Medication adherence improves 40-60% when MI techniques are used vs standard","education approach (Zygmunt et al., 2002. Psychiatr Serv, 53, 1508-1515)."], bgc='cream', title="How to Apply:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Kemp et al. (1996). BMJ; Parsons et al. (2007). AIDS Behav; Zygmunt et al. (2002). Psychiatr Serv.")

# Additional: Cultural Adaptations
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_orange'); title_bar(slide, "Cultural Adaptations of MET for Indian Context", 'orange')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["INDIAN CULTURAL FACTORS:","","- Family-centric decision making","  (not just individual autonomy)","- Stigma around mental health/addiction","- Spiritual/religious values important","- Gender roles affect presentation","- Economic pressures are major triggers","- Joint family system: multiple stakeholders","- Help-seeking often delayed","- Traditional healing sought first"], bgc='light_orange', title="Cultural Context:", ttc='orange', border='orange', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["MET ADAPTATIONS FOR INDIA:","","1. Include FAMILY in values exploration","   (family honor, duty, role as provider)","2. Integrate SPIRITUAL values naturally","   (\"How does your faith view this?\")","3. Address STIGMA non-judgmentally","   (normalize help-seeking)","4. Use culturally relevant EXAMPLES","5. Consider COLLECTIVIST values in Change Plan","6. GENDER-sensitive approach","7. Economic factors in discrepancy work"], bgc='light_green', title="Adaptations:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["INDIAN RESEARCH ON CULTURAL ADAPTATION:","- Nadkarni et al. (2017): PREMIUM trial showed culturally adapted MI delivered by lay","  counselors (\"Healthy Activity Programme\") was effective in Goa (Lancet, 389)","- Patel et al. (2010): Task-shifting mental health care to community workers feasible in India","- NIMHANS (2016): Recommends family sessions as adjunct to individual MI in Indian settings","- Cultural values like family duty/honor can be POWERFUL motivators when explored with MI","","KEY INSIGHT: MI's person-centered, non-judgmental approach is especially valuable in cultures","with high stigma around substance use. The principles translate well across cultures."], bgc='cream', title="Evidence:", ttc='deep_blue', border='deep_blue', fs=12)
ref_bar(slide, "Ref: Nadkarni et al. (2017). Lancet; Patel et al. (2010). Lancet; NIMHANS (2016), Ch. 8.")

# Additional: Ethical Considerations
slide = prs.slides.add_slide(blank_layout); bg(slide, 'white', 'light_purple'); title_bar(slide, "Ethical Considerations in MET Practice", 'purple')
cbox(slide, LM, ROW1, HALF_W, ROW1_H, ["ETHICAL PRINCIPLES IN MET:","","1. AUTONOMY: Client's right to choose","   (even if they choose not to change)","2. BENEFICENCE: Acting in client's interest","   (even when they don't see the harm)","3. NON-MALEFICENCE: Do no harm","   (confrontation can harm - MI avoids it)","4. INFORMED CONSENT: Client knows the","   approach and their rights"], bgc='light_purple', title="Ethical Framework:", ttc='purple', border='purple', fs=12)
cbox(slide, COL2, ROW1, HALF_W, ROW1_H, ["ETHICAL DILEMMAS IN MET:","","- Client chooses continued use despite risks","  -> Respect autonomy but document","- Mandated clients (court-ordered)","  -> MI still works; acknowledge constraint","- Risk to self/others (suicidal, DUI)","  -> Safety overrides autonomy temporarily","- Family pressure vs client wishes","  -> Client's goals take priority","- Therapist values vs client's choice","  -> Manage countertransference"], bgc='light_green', title="Common Dilemmas:", ttc='green', border='green', fs=12)
cbox(slide, LM, ROW2, FULL_W, ROW2_H, ["\"The spirit of MI is deeply respectful of client autonomy. The therapist may feel an urge","to persuade, direct, or confront - this is the 'righting reflex' that must be managed.","Even when a client makes choices we disagree with, their autonomy is honored.\"","(Miller & Rollnick, 2013)","","Exception: When there is IMMINENT risk to self or others, standard duty-of-care protocols","override MI autonomy principles (safety planning, hospitalization if needed, reporting).","This applies to our case: elevated S-CON indicators require active risk monitoring."], bgc='cream', title="Key Principle:", ttc='orange', border='gold', fs=12)
ref_bar(slide, "Ref: Miller & Rollnick (2013). MI 3rd ed., Ch. 2; APA Ethics Code (2017); MCI Guidelines.")

# Additional: Summary Evidence Table
tbl("Evidence-Based Indications: When to Use MET",
    ["Indication", "Evidence Level", "Key Study", "Effect"],
    [["Alcohol Use Disorder", "Level 1 (Strong)", "Project MATCH (1997)\nUKATT (2005)", "Equal to 12-session Tx\nin 4 sessions"],
     ["Cannabis Use", "Level 1 (Strong)", "CYT (2004)\nStephens et al.", "Significant reduction"],
     ["Treatment Engagement", "Level 1 (Strong)", "Carroll et al. (2006)\nMartino et al. (2007)", "Increases retention"],
     ["Medication Adherence", "Level 2 (Moderate)", "Kemp (1996)\nParsons (2007)", "40-60% improvement"],
     ["Dual Diagnosis", "Level 2 (Moderate)", "Barrowclough (2001)\nBaker (2010)", "Effective integrated"],
     ["Adolescents", "Level 1 (Strong)", "CYT Study\nJensen et al. (2011)", "Developmentally apt"],
     ["Low-resource settings", "Level 1 (Strong)", "Nadkarni (2017)\nWHO (2002)", "Task-shifting works"]],
    color='navy', ref="Ref: NICE (2011); SAMHSA TIP 35; APA Guidelines (2006); NIMHANS (2016).")

# ============================================================
# FINAL SAVE
# ============================================================
output = '/projects/sandbox/Dango-kiro/MET_Comprehensive_Presentation.pptx'
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
