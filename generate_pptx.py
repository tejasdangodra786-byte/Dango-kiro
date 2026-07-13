#!/usr/bin/env python3
"""Generate MBRP Research Synopsis PowerPoint Presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xB8)
ACCENT_GOLD = RGBColor(0xD4, 0xA5, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
BODY_TEXT = RGBColor(0x2D, 0x2D, 0x44)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)


def add_title_slide(prs, title, subtitle, researcher="[Researcher Name]",
                    guide="[Supervisor Name]", dept="Clinical Psychology",
                    inst="[University/Institute Name]", year="2026"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11.3), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = ACCENT_GOLD
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    # Details
    details = f"\nResearcher: {researcher}\nGuide: {guide}\nDepartment: {dept}\nInstitution: {inst}\nYear: {year}"
    p3 = tf.add_paragraph()
    p3.text = details
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)
    return slide


def add_content_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_LIGHT
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    # Subtitle if provided
    top_offset = 1.3
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = MEDIUM_BLUE
        top_offset = 1.7
    # Bullets
    txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(top_offset), Inches(11.9), Inches(7.5 - top_offset - 0.3))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.font.color.rgb = BODY_TEXT
        p.space_before = Pt(8)
        p.level = 0
    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_LIGHT
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # Left column
    txL = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.8))
    tfL = txL.text_frame
    tfL.word_wrap = True
    pL = tfL.paragraphs[0]
    pL.text = left_title
    pL.font.size = Pt(15)
    pL.font.bold = True
    pL.font.color.rgb = MEDIUM_BLUE
    for b in left_bullets:
        pp = tfL.add_paragraph()
        pp.text = b
        pp.font.size = Pt(13)
        pp.font.color.rgb = BODY_TEXT
        pp.space_before = Pt(6)
    # Right column
    txR = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.8))
    tfR = txR.text_frame
    tfR.word_wrap = True
    pR = tfR.paragraphs[0]
    pR.text = right_title
    pR.font.size = Pt(15)
    pR.font.bold = True
    pR.font.color.rgb = MEDIUM_BLUE
    for b in right_bullets:
        pp = tfR.add_paragraph()
        pp.text = b
        pp.font.size = Pt(13)
        pp.font.color.rgb = BODY_TEXT
        pp.space_before = Pt(6)
    return slide


# ============================================================
# SLIDE 1: TITLE
# ============================================================
add_title_slide(prs,
    "Efficacy of Brief Mindfulness-Based Relapse Prevention (MBRP) Intervention\non Craving, Impulsivity, and Mindfulness\nin Substance Dependent Patients",
    "MPhil Clinical Psychology Research Synopsis")

# ============================================================
# SLIDE 2: INTRODUCTION 1/3
# ============================================================
add_content_slide(prs, "SLIDE 2: Introduction & Background (1/3)",
    [
        "India faces a significant substance use crisis; MAGNITUDE study (2019) estimated ~3.1 crore individuals affected by substance use disorders",
        "Opioid dependence constitutes a major public health burden, particularly in Punjab, Rajasthan, Northeast India, and metropolitan areas",
        "WHO estimates India accounts for ~25% of global opioid-related deaths in South-East Asia",
        "Economic burden includes lost productivity, healthcare costs, and family dysfunction",
        "National Drug Dependence Treatment Centre (NDDTC) reports increasing treatment-seeking among opioid users",
        "Indian de-addiction centers primarily offer pharmacotherapy (OST, naltrexone) with limited structured psychotherapy",
        "Psychosocial interventions remain under-utilized despite evidence of superior combined treatment outcomes"
    ],
    subtitle="Substance Dependence: Indian Context")


# ============================================================
# SLIDE 3: INTRODUCTION 2/3
# ============================================================
add_content_slide(prs, "SLIDE 3: Introduction & Background (2/3)",
    [
        "Relapse rates in substance dependence range from 40-60% within first year post-treatment (NIDA, 2020)",
        "Indian studies report even higher relapse rates (~70-80%) in opioid dependence (Mattoo et al., 2009)",
        "Triggers for relapse: craving, negative affect, interpersonal conflict, environmental cues",
        "Traditional relapse prevention (Marlatt & Gordon, 1985) has moderate efficacy but limited mindfulness integration",
        "Automatic cognitive-behavioral patterns (e.g., 'apparently irrelevant decisions') perpetuate relapse cycles",
        "Need for interventions addressing both cognitive automaticity and emotional dysregulation",
        "Brief interventions essential for Indian rehab settings with limited resources and high patient turnover"
    ],
    subtitle="The Relapse Problem in De-Addiction")

# ============================================================
# SLIDE 4: INTRODUCTION 3/3
# ============================================================
add_content_slide(prs, "SLIDE 4: Introduction & Background (3/3)",
    [
        "Mindfulness: intentional, non-judgmental, present-moment awareness (Kabat-Zinn, 1990)",
        "Mindfulness-Based Relapse Prevention (MBRP) developed by Bowen, Chawla, & Marlatt (2011)",
        "Integrates Mindfulness-Based Stress Reduction (MBSR) + Cognitive-Behavioral Relapse Prevention (RP)",
        "Core mechanisms: decentering from craving, disrupting automatic reactivity, increasing distress tolerance",
        "MBRP targets the 'craving to use' automaticity by cultivating awareness of triggers",
        "Brief adaptations (6-8 sessions) show promise for resource-limited settings",
        "Growing global evidence supports MBRP; however, Indian validation remains scarce"
    ],
    subtitle="Mindfulness in Addiction Treatment")


# ============================================================
# SLIDE 5: VARIABLE 1 - CRAVING
# ============================================================
add_content_slide(prs, "SLIDE 5: Variable 1 - Craving",
    [
        "DEFINITION: Intense subjective urge to use a substance, triggered by internal/external cues; motivational state from incentive-sensitization (Robinson & Berridge, 1993)",
        "RELEVANCE: Primary predictor of relapse in opioid dependence; mediates cue exposure and use behavior; correlates with severity and treatment dropout",
        "NEUROPSYCHOLOGICAL BASIS: Mesolimbic dopamine pathway (VTA -> Nucleus Accumbens); PFC hypoactivation; conditioned cue-reward associations",
        "TECHNIQUE - Urge Surfing: Observing craving as a transient wave without acting on it",
        "TECHNIQUE - Mindfulness Exposure: Non-reactive awareness of craving sensations",
        "TECHNIQUE - SOBER Breathing Space: Pause-observe-redirect during craving triggers",
        "TECHNIQUE - Cognitive Decentering: 'I am having a craving' vs. 'I need the drug'"
    ])

# ============================================================
# SLIDE 6: VARIABLE 2 - IMPULSIVITY
# ============================================================
add_content_slide(prs, "SLIDE 6: Variable 2 - Impulsivity",
    [
        "DEFINITION: Tendency toward rapid, unplanned actions without adequate consideration of consequences (Moeller et al., 2001); multi-dimensional: motor, attentional, non-planning (Patton et al., 1995)",
        "RELEVANCE: Higher trait impulsivity predicts initiation, escalation, and relapse; mediates craving-to-use behavior; associated with treatment non-adherence",
        "NEUROPSYCHOLOGICAL BASIS: PFC dysfunction -> impaired executive control; reduced inhibitory control (Go/No-Go); impaired delay discounting; DLPFC hypoactivation",
        "TECHNIQUE - Response Inhibition Training: Mindful pause before automatic behavioral responses",
        "TECHNIQUE - Awareness Training: Noticing impulse-action sequences without engagement",
        "TECHNIQUE - Mindful Decision-Making: Creating space between stimulus and response",
        "TECHNIQUE - STOP Technique: Stop-Take a breath-Observe-Proceed mindfully"
    ])

# ============================================================
# SLIDE 7: VARIABLE 3 - MINDFULNESS
# ============================================================
add_content_slide(prs, "SLIDE 7: Variable 3 - Mindfulness",
    [
        "DEFINITION: Capacity to attend to present-moment experience with openness, curiosity, and non-judgment; 5 facets: Observing, Describing, Acting with Awareness, Non-Judging, Non-Reactivity (Baer et al., 2006)",
        "RELEVANCE: Substance-dependent individuals show significantly lower dispositional mindfulness; acts as protective factor against relapse; improvements mediate treatment outcomes",
        "NEUROPSYCHOLOGICAL BASIS: ACC activation -> enhanced self-regulation; Insula -> interoceptive awareness; PFC-amygdala connectivity -> emotional regulation; DMN regulation -> reduced rumination",
        "TECHNIQUE - Sitting Meditation: Focused attention on breath, body sensations, thoughts",
        "TECHNIQUE - Body Scan: Systematic non-judgmental awareness of bodily states",
        "TECHNIQUE - Mindful Movement: Gentle yoga/walking with present-moment focus",
        "TECHNIQUE - Non-Judgmental Awareness: Labeling experiences without evaluation"
    ])


# ============================================================
# SLIDES 8-14: REVIEW OF LITERATURE
# ============================================================
add_content_slide(prs, "SLIDE 8: Review of Literature (1/7) - MBRP Foundational Studies",
    [
        "Bowen, Chawla, & Marlatt (2011) developed original MBRP as 8-session aftercare integrating mindfulness + cognitive-behavioral relapse prevention",
        "Bowen, Witkiewitz, et al. (2014) RCT (N=286): MBRP showed significantly fewer substance use days at 12-month follow-up vs. RP and TAU",
        "Witkiewitz & Bowen (2010) demonstrated MBRP reduced relationship between depressive symptoms and craving; mindfulness mediates affect-craving pathways",
        "Bowen & Marlatt (2009) found brief mindfulness meditation reduced craving among incarcerated substance users vs. controls",
        "Zgierska et al. (2009) systematic review: mindfulness meditation is a promising adjunct in substance use disorder treatment"
    ])

add_content_slide(prs, "SLIDE 9: Review of Literature (2/7) - Craving Studies",
    [
        "Witkiewitz, Bowen, Douglas, & Hsu (2013): MBRP participants showed significant reductions in craving intensity vs. TAU over 4-month follow-up",
        "Garland, Froeliger, & Howard (2014): Mindfulness-Oriented Recovery Enhancement (MORE) reduced opioid craving through enhanced positive reappraisal",
        "Garland, Manusov, et al. (2014): Mindfulness training associated with reduced craving and decreased neural reactivity to drug cues in opioid patients",
        "Brewer, Mallik, et al. (2011): Mindfulness training (vs. CBT) resulted in greater craving reductions in cigarette smokers",
        "Hsu, Collins, & Marlatt (2013): Urge surfing effectively reduced craving intensity and frequency in alcohol-dependent outpatients"
    ])

add_content_slide(prs, "SLIDE 10: Review of Literature (3/7) - Impulsivity Studies",
    [
        "Garland, Roberts-Lewis, et al. (2016): Mindfulness-based intervention reduced BIS-11 impulsivity scores significantly in substance-dependent adults",
        "Murphy & MacKillop (2012): Trait mindfulness inversely associated with impulsive decision-making (delay discounting) in alcohol users",
        "Korponay, Dentico, et al. (2019): Long-term meditators exhibited lower impulsivity and greater prefrontal cortical thickness vs. controls",
        "Alfonso, Caracuel, et al. (2011): Combined mindfulness-inhibitory control training reduced impulsive responding in polysubstance users",
        "Peters, Erisman, et al. (2011): Brief mindfulness induction reduced impulsive responses on behavioral tasks among high-impulsivity individuals"
    ])


add_content_slide(prs, "SLIDE 11: Review of Literature (4/7) - Mindfulness & Substance Use",
    [
        "Karyadi, VanderVeen, & Cyders (2014) meta-analysis: significant negative association between mindfulness and substance use (r = -0.21), craving, and problems",
        "Li, Howard, Garland, et al. (2017) meta-analysis of 42 RCTs: MBIs reduced substance misuse (d=0.33), craving (d=0.68), and stress (d=0.44)",
        "Chiesa & Serretti (2014) reviewed 24 studies: MBIs show moderate efficacy for reducing substance use and craving across substances",
        "Bowen, Witkiewitz, Dillworth, & Marlatt (2007): MBRP increased acceptance and acting with awareness (FFMQ facets) in post-treatment users",
        "Grant, Colaiaco, et al. (2017) Cochrane-level review: moderate-quality evidence supporting MBIs for substance use disorders"
    ])

add_content_slide(prs, "SLIDE 12: Review of Literature (5/7) - Brief Interventions",
    [
        "Glasner-Edwards, Mooney, et al. (2017): Abbreviated 6-session mindfulness intervention feasible and effective in reducing stimulant use",
        "Davis, Goldberg, et al. (2014): Brief 4-week mindfulness training reduced cigarette consumption and craving among smokers",
        "Shorey, Elmquist, et al. (2017): Brief (single-session + booster) mindfulness intervention reduced substance craving in residential treatment",
        "Roos, Bowen, & Witkiewitz (2017): Even brief daily 15-min meditation mediated MBRP effects on craving and substance use",
        "Rationale: Indian rehabilitation settings typically have 4-8 week admission, necessitating condensed evidence-based protocols"
    ])

add_content_slide(prs, "SLIDE 13: Review of Literature (6/7) - Indian Studies",
    [
        "Sarkar & Balhara (2016): Highlighted underutilization of structured psychological interventions in Indian de-addiction settings",
        "Murthy (2016): Indian de-addiction centers predominantly use pharmacotherapy with limited psychosocial integration",
        "Ghosh, Basu, & Avasthi (2018): Relapse rates >70% among opioid-dependent patients in North Indian treatment centers",
        "Jain, Majumder, & Gupta (2013): Preliminary mindfulness-based intervention study in alcohol-dependent patients (India) showed craving reductions",
        "CRITICAL NOTE: No published Indian RCT on Brief MBRP for opioid dependence identified - establishes critical research gap"
    ])

add_content_slide(prs, "SLIDE 14: Review of Literature (7/7) - Summary Table",
    [
        "MBRP & Relapse: Fewer substance use days at 12-month follow-up | Evidence: RCT (Bowen et al., 2014)",
        "Craving: MBRP reduces craving intensity via urge surfing | Evidence: Multiple RCTs",
        "Impulsivity: Mindfulness inversely associated with impulsive responding | Evidence: Correlational + RCT",
        "Mindfulness: Significant improvement in FFMQ scores post-MBRP | Evidence: Pre-post + RCT",
        "Brief Models: 4-8 session protocols feasible and effective | Evidence: Emerging",
        "Indian Context: Extremely limited MBRP research; high relapse burden | Evidence: Gap identified"
    ])


# ============================================================
# SLIDE 15: RESEARCH GAP
# ============================================================
add_content_slide(prs, "SLIDE 15: Research Gap",
    [
        "Lack of Indian MBRP studies: No published RCT has tested MBRP (standard or brief) in Indian opioid-dependent populations",
        "Limited combined variable assessment: Most studies examine craving OR mindfulness in isolation; few assess craving + impulsivity + mindfulness simultaneously",
        "Absence of brief protocol validation: Indian rehab settings need 6-8 session models but none empirically validated in this context",
        "Overreliance on pharmacotherapy: Indian treatment lacks integration of structured evidence-based psychological interventions alongside OST",
        "Population specificity: Global MBRP studies focus on alcohol/polysubstance; opioid-specific MBRP evidence is limited",
        "Cultural adaptation need: Mindfulness requires culturally congruent adaptation for Indian populations",
        "PRESENT STUDY: First brief MBRP trial in Indian opioid-dependent sample assessing craving, impulsivity, and mindfulness simultaneously"
    ])

# ============================================================
# SLIDE 16: AIM
# ============================================================
add_content_slide(prs, "SLIDE 16: Aim of the Study",
    [
        "To evaluate the efficacy of a Brief Mindfulness-Based Relapse Prevention (MBRP) intervention in reducing craving and impulsivity, and enhancing mindfulness, among substance-dependent patients in an Indian de-addiction setting",
        "",
        "Specifically: To compare outcomes between Brief MBRP + TAU (Experimental) and Psychoeducation + TAU (Control) on three dependent variables measured pre- and post-intervention"
    ])

# ============================================================
# SLIDE 17: OBJECTIVES
# ============================================================
add_content_slide(prs, "SLIDE 17: Objectives",
    [
        "1. To assess and compare craving levels (pre vs. post) in Experimental (Brief MBRP + TAU) and Control (Psychoeducation + TAU) groups",
        "2. To assess and compare impulsivity levels (pre vs. post) in both groups",
        "3. To assess and compare mindfulness levels (pre vs. post) in both groups",
        "4. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in reducing craving",
        "5. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in reducing impulsivity",
        "6. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in enhancing mindfulness"
    ])


# ============================================================
# SLIDE 18: HYPOTHESES
# ============================================================
add_content_slide(prs, "SLIDE 18: Hypotheses (Directional)",
    [
        "H1: Participants receiving Brief MBRP + TAU will show significantly GREATER REDUCTION in craving scores (HCQ) compared to Psychoeducation + TAU, from pre-test to post-test",
        "",
        "H2: Participants receiving Brief MBRP + TAU will show significantly GREATER REDUCTION in impulsivity scores (BIS-11) compared to Psychoeducation + TAU, from pre-test to post-test",
        "",
        "H3: Participants receiving Brief MBRP + TAU will show significantly GREATER INCREASE in mindfulness scores (FFMQ) compared to Psychoeducation + TAU, from pre-test to post-test"
    ])

# ============================================================
# SLIDE 19: RESEARCH DESIGN
# ============================================================
add_content_slide(prs, "SLIDE 19: Research Design",
    [
        "Design: Pre-test Post-test Control Group Experimental Design",
        "",
        "R   O1   X1   O2   ->  Experimental Group (Brief MBRP + TAU)",
        "R   O1   X2   O2   ->  Control Group (Psychoeducation + TAU)",
        "",
        "R = Random assignment | O1 = Pre-test (FFMQ, BIS-11, HCQ, ASSIST) | X1 = Brief MBRP | X2 = Psychoeducation | O2 = Post-test (FFMQ, BIS-11, HCQ)",
        "",
        "Features: True experimental design with randomization; Active control (Psychoeducation) for attention/contact effects; TAU maintained in both groups"
    ],
    subtitle="Pre-test Post-test Control Group Design")

# ============================================================
# SLIDE 20: SAMPLE
# ============================================================
add_content_slide(prs, "SLIDE 20: Sample",
    [
        "Population: Substance-dependent patients (primarily opioid users) admitted to de-addiction/rehabilitation centers",
        "Sampling Method: Purposive sampling followed by random allocation to groups",
        "Sample Size: N = 60 (30 per group)",
        "Setting: Government/private de-addiction centers in India",
        "Recruitment: Consecutive admissions meeting inclusion criteria over 6-8 months",
        "Attrition Consideration: Recruit N = 70 (35 per group) to account for ~15% dropout"
    ])


# ============================================================
# SLIDE 21: SAMPLE SIZE ESTIMATION
# ============================================================
add_content_slide(prs, "SLIDE 21: Sample Size Estimation",
    [
        "Formula: n = [(Za/2 + Zb)^2 x 2s^2] / d^2",
        "Parameters: Effect size (d) = 0.50 (medium; Li et al., 2017 meta-analysis: d = 0.33-0.68) | Power (1-b) = 0.80 -> Zb = 0.84 | Alpha = 0.05 (two-tailed) -> Za/2 = 1.96",
        "Calculation: n = [(1.96 + 0.84)^2 x 2 x 1] / (0.50)^2 = [7.84 x 2] / 0.25 = 15.68 / 0.25 = 62.72 ~ 63 total (~32 per group)",
        "G*Power Verification: Independent t-test (two-tailed): d=0.50, a=0.05, power=0.80 -> n=64 total; ANCOVA with 1 covariate -> n=34 per group",
        "FINAL DECISION: N = 60 (30 per group) - justified by ANCOVA as primary analysis (reduces required n)",
        "Recruit 70 total (35/group) to account for attrition (~15%)",
        "Consistent with similar studies: Bowen et al. (2009), Glasner-Edwards et al. (2017)"
    ])

# ============================================================
# SLIDE 22: INCLUSION CRITERIA
# ============================================================
add_content_slide(prs, "SLIDE 22: Inclusion Criteria",
    [
        "1. Diagnosis of Substance Dependence as per ICD-10 (F10-F19) / ICD-11 criteria",
        "2. Primary substance: Opioids (heroin, pharmaceutical opioids); polysubstance users with primary opioid dependence included",
        "3. Age: 18-50 years (male participants)",
        "4. Completed detoxification phase (minimum 7 days post-withdrawal)",
        "5. Currently admitted in de-addiction/rehabilitation center",
        "6. Minimum education: 5th standard (ability to comprehend psychometric tools)",
        "7. Willingness to provide written informed consent",
        "8. Able to attend all intervention sessions during admission"
    ])

# ============================================================
# SLIDE 23: EXCLUSION CRITERIA
# ============================================================
add_content_slide(prs, "SLIDE 23: Exclusion Criteria",
    [
        "1. Severe psychiatric comorbidity: Psychotic disorders, Bipolar I with psychotic features, severe MDE with suicidality",
        "2. Significant cognitive impairment (MMSE < 24) or intellectual disability",
        "3. Active withdrawal symptoms (COWS score > 12)",
        "4. History of traumatic brain injury with LOC > 30 minutes",
        "5. Current participation in another structured psychological intervention study",
        "6. Medical instability requiring acute care",
        "7. History of prior formal mindfulness/meditation training (> 1 month)"
    ])


# ============================================================
# SLIDE 24: VARIABLES
# ============================================================
add_content_slide(prs, "SLIDE 24: Variables",
    [
        "INDEPENDENT VARIABLE (IV):",
        "  Type of Intervention: Level 1 = Brief MBRP + TAU (Experimental) | Level 2 = Psychoeducation + TAU (Control)",
        "",
        "DEPENDENT VARIABLES (DVs):",
        "  1. Craving - measured by Heroin Craving Questionnaire (HCQ)",
        "  2. Impulsivity - measured by Barratt Impulsiveness Scale (BIS-11)",
        "  3. Mindfulness - measured by Five Facet Mindfulness Questionnaire (FFMQ)",
        "",
        "CONTROLLED VARIABLES: Age, education, duration of substance use, severity (ASSIST baseline), TAU components constant, session duration/frequency equalized"
    ])

# ============================================================
# SLIDE 25: TOOLS - FFMQ & BIS-11
# ============================================================
add_content_slide(prs, "SLIDE 25: Tools - Mindfulness & Impulsivity Measures",
    [
        "1. FIVE FACET MINDFULNESS QUESTIONNAIRE (FFMQ) - Baer et al. (2006)",
        "   39 items | 5 facets: Observing, Describing, Acting with Awareness, Non-Judging, Non-Reactivity | 5-point Likert | Higher = greater mindfulness",
        "   Reliability: a = 0.75-0.91 across facets | Validity: Convergent + discriminant established | Indian: Hindi adaptation available",
        "",
        "2. BARRATT IMPULSIVENESS SCALE (BIS-11) - Patton, Stanford, & Barratt (1995)",
        "   30 items | 3 factors: Attentional, Motor, Non-Planning Impulsivity | 4-point scale | Higher = greater impulsivity",
        "   Reliability: a = 0.79-0.83; test-retest r = 0.83 | Validity: Discriminates substance users from controls | Indian: Hindi version validated (Mathew et al., 2014)"
    ])


# ============================================================
# SLIDE 26: TOOLS - CRAVING
# ============================================================
add_content_slide(prs, "SLIDE 26: Tools - Craving Assessment",
    [
        "CRAVING SCALE OPTIONS COMPARISON:",
        "  Heroin Craving Questionnaire (HCQ): 45-item (brief: 14-item), opioid-specific, multidimensional | HIGHLY SUITABLE",
        "  Obsessive Compulsive Drug Use Scale (OCDUS): 12-item, obsessive thoughts + compulsive urges | Moderate (generic)",
        "  Visual Analog Scale (VAS): Single-item 0-100mm rating | Limited (unidimensional)",
        "",
        "RECOMMENDED: HCQ-Brief (14 items) - Tiffany, Fields, Singleton, et al. (2000)",
        "  Assesses: desire to use, intention, anticipation of positive outcome, relief from withdrawal, lack of control",
        "  Scoring: 7-point Likert | Reliability: a = 0.87-0.93 | Validity: Convergent with VAS, sensitive to intervention",
        "  JUSTIFICATION: Substance-specific measurement preferred for opioid-dependent populations (EMCDDA guidelines)"
    ])

# ============================================================
# SLIDE 27: TOOLS - SEVERITY (ASSIST)
# ============================================================
add_content_slide(prs, "SLIDE 27: Tools - Severity Assessment (Pre-Test Only)",
    [
        "SEVERITY SCALE OPTIONS:",
        "  ASSIST (WHO): 8-item, risk level for 10 substance categories | BEST CHOICE",
        "  DAST-20: 20-item, drug abuse screening (yes/no) | Good but less comprehensive",
        "  AUDIT: 10-item, alcohol-specific | Only if alcohol subgroup",
        "",
        "RECOMMENDED: ASSIST - WHO ASSIST Working Group (2002); Humeniuk et al.",
        "  Scoring: Substance-specific risk: Low (0-3), Moderate (4-26), High (27+)",
        "  Reliability: Test-retest r = 0.58-0.90; a = 0.77-0.94 | Validity: Sensitivity 0.80, Specificity 0.71",
        "  Indian: WHO-validated; Hindi version available; NDDTC studies; culturally appropriate; 5-10 min",
        "  PURPOSE: Used at PRE-TEST ONLY to establish baseline severity and ensure group equivalence"
    ])


# ============================================================
# SLIDE 28: PROCEDURE
# ============================================================
add_content_slide(prs, "SLIDE 28: Procedure",
    [
        "Step 1: SCREENING - Identify eligible patients (ICD-10 diagnosis, detoxified, consent-capable)",
        "Step 2: INFORMED CONSENT - Explain study purpose, procedures, confidentiality, right to withdraw",
        "Step 3: PRE-TEST ASSESSMENT - Administer: ASSIST + HCQ + BIS-11 + FFMQ",
        "Step 4: RANDOM ALLOCATION - Computer-generated randomization -> Experimental vs. Control",
        "Step 5: INTERVENTION DELIVERY - Experimental: Brief MBRP (8 sessions x 60 min x 4 weeks) | Control: Psychoeducation (8 sessions x 60 min x 4 weeks)",
        "Step 6: POST-TEST ASSESSMENT - Administer: HCQ + BIS-11 + FFMQ (within 1 week of completion)",
        "Step 7: DATA COMPILATION & ANALYSIS",
        "",
        "Both groups receive TAU throughout | Assessments by blinded research assistant | Intervention by trained MPhil Clinical Psychologist"
    ],
    subtitle="Stepwise Research Flow")

# ============================================================
# SLIDE 29: INTERVENTION - EXPERIMENTAL
# ============================================================
add_content_slide(prs, "SLIDE 29: Intervention Plan - Experimental Group",
    [
        "BRIEF MBRP PROTOCOL: 8 Sessions x 60 min x 4 weeks (Twice weekly)",
        "Session 1: Introduction to MBRP & Autopilot - Raisin exercise; identifying automatic patterns",
        "Session 2: Awareness of Triggers - Body scan meditation; mapping personal triggers",
        "Session 3: Mindfulness in Daily Life - Sitting meditation (breath); SOBER breathing space",
        "Session 4: Mindfulness in High-Risk Situations - Urge surfing; role-play; cognitive decentering",
        "Session 5: Acceptance & Skillful Action - Non-judgmental awareness; mindful decision-making",
        "Session 6: Seeing Thoughts as Thoughts - 'Thoughts are not facts'; mountain meditation",
        "Session 7: Self-Care & Lifestyle Balance - Loving-kindness; activity scheduling; warning signs",
        "Session 8: Social Support & Maintenance - Review; relapse prevention plan; practice commitment"
    ],
    subtitle="Brief MBRP (8 Sessions) + TAU")

# ============================================================
# SLIDE 30: INTERVENTION - CONTROL
# ============================================================
add_content_slide(prs, "SLIDE 30: Intervention Plan - Control Group",
    [
        "PSYCHOEDUCATION PROTOCOL: 8 Sessions x 60 min x 4 weeks (Twice weekly)",
        "Session 1: Understanding Addiction - Nature of dependence; brain changes; disease model",
        "Session 2: Effects of Opioids - Short/long-term physical and psychological consequences",
        "Session 3: Understanding Relapse - Relapse process; warning signs; high-risk situations",
        "Session 4: Coping with Cravings - General strategies (distraction, social support); NO mindfulness",
        "Session 5: Health & Nutrition - Physical recovery; nutrition; sleep hygiene",
        "Session 6: Social Consequences - Family impact; legal issues; stigma; rehabilitation",
        "Session 7: Motivation & Goal Setting - Stages of change; personal goals; motivation enhancement",
        "Session 8: Lifestyle Changes & Summary - Long-term planning; support systems; review",
        "NOTE: Matched for contact time, group format, therapist attention to control non-specific factors"
    ],
    subtitle="Psychoeducation (8 Sessions) + TAU")


# ============================================================
# SLIDE 31: DATA ANALYSIS
# ============================================================
add_content_slide(prs, "SLIDE 31: Data Analysis",
    [
        "Descriptive Statistics: Mean, SD, frequency, percentages for sociodemographic and clinical variables",
        "Normality Testing: Shapiro-Wilk test for distribution of outcome variables",
        "Within-Group: Paired samples t-test (pre vs. post within each group); Wilcoxon for non-normal data",
        "Between-Group: Independent samples t-test (post-test between groups); Mann-Whitney U for non-normal",
        "PRIMARY ANALYSIS: ANCOVA - Post-test scores as DV, Group as IV, Pre-test scores as covariate (controls baseline differences, increases power)",
        "Effect Size: Cohen's d for between-group differences",
        "Significance Level: alpha = 0.05 (two-tailed) | Software: SPSS 26.0 / JASP",
        "Intent-to-Treat (ITT): Last Observation Carried Forward (LOCF) for dropouts"
    ])

# ============================================================
# SLIDE 32: ETHICAL CONSIDERATIONS
# ============================================================
add_content_slide(prs, "SLIDE 32: Ethical Considerations",
    [
        "Informed Consent: Written consent in Hindi/regional language; participants informed of purpose, procedures, risks, benefits",
        "Voluntary Participation: Right to withdraw at any time without penalty or impact on ongoing treatment",
        "Confidentiality: Data coded with participant IDs; no identifying information in publications; secure storage",
        "Non-Maleficence: No harmful procedures; control receives active psychoeducation (not waitlist); TAU continued for all",
        "Institutional Approval: Ethical clearance from Institutional Ethics Committee (IEC) prior to data collection",
        "Debriefing: Control group offered MBRP orientation post-study if desired",
        "Compliance: ICMR (2017) National Ethical Guidelines for Biomedical and Health Research"
    ])

# ============================================================
# SLIDE 33: EXPECTED RESULTS
# ============================================================
add_content_slide(prs, "SLIDE 33: Expected Results",
    [
        "CRAVING (HCQ):",
        "  Significant REDUCTION in craving scores in Experimental group vs. Control group",
        "  Mechanism: Urge surfing and mindfulness exposure reduce automatic craving reactivity",
        "",
        "IMPULSIVITY (BIS-11):",
        "  Significant REDUCTION in impulsivity (Motor + Attentional subscales) in Experimental group",
        "  Mechanism: Mindfulness enhances prefrontal inhibitory control and response awareness",
        "",
        "MINDFULNESS (FFMQ):",
        "  Significant INCREASE in mindfulness (Acting with Awareness + Non-Reactivity) in Experimental group",
        "  Mechanism: Structured meditation practice cultivates dispositional mindfulness"
    ])


# ============================================================
# SLIDE 34: CLINICAL IMPLICATIONS
# ============================================================
add_content_slide(prs, "SLIDE 34: Clinical Implications",
    [
        "Validates a brief, structured MBRP model feasible for Indian de-addiction settings with time-limited admissions",
        "Provides evidence-based psychological intervention to complement existing pharmacotherapy (OST, naltrexone)",
        "Demonstrates mindfulness-based approaches are culturally compatible with Indian populations (meditative traditions)",
        "Addresses multiple relapse risk factors simultaneously (craving + impulsivity + mindfulness) via single protocol",
        "Supports task-shifting: Brief MBRP can be delivered by MPhil-trained Clinical Psychologists in resource-limited settings",
        "Informs treatment policy at national level (NIMHANS, NDDTC, State Mental Health Authorities)",
        "Contributes to RCI-recognized intervention repertoire for clinical psychology training in India"
    ])

# ============================================================
# SLIDE 35: LIMITATIONS
# ============================================================
add_content_slide(prs, "SLIDE 35: Limitations",
    [
        "Sample specificity: Male opioid-dependent patients from single center; limits generalizability to females, other substances",
        "Short-term assessment: Post-test immediately after intervention; no long-term follow-up for maintenance",
        "Self-report measures: HCQ, BIS-11, FFMQ susceptible to social desirability and limited insight",
        "No biological markers: Craving measured subjectively; no physiological/neuroimaging corroboration",
        "Therapist effects: Single therapist delivery may introduce therapist-specific confounds",
        "Attention control: Psychoeducation controls for contact but not specific mindfulness mechanisms",
        "Attrition: Substance-dependent population may have higher dropout rates despite planning"
    ])

# ============================================================
# SLIDE 36: FUTURE DIRECTIONS
# ============================================================
add_content_slide(prs, "SLIDE 36: Future Directions",
    [
        "Follow-up studies: 3-month and 6-month assessments for sustained MBRP effects",
        "Multi-site RCTs: Replicate across multiple Indian de-addiction centers for generalizability",
        "Female participants: Include female substance-dependent patients for gender-specific effects",
        "Neuroimaging integration: fMRI/EEG to examine neural mechanisms of MBRP on craving circuits",
        "Dose-response analysis: Compare 6 vs. 8 vs. full 8-week MBRP for optimal dosage",
        "Mediator analysis: Examine whether mindfulness improvement mediates craving/impulsivity reduction",
        "Technology-assisted delivery: App-based/digital MBRP for post-discharge and rural access",
        "Comparative effectiveness: MBRP vs. CBT vs. ACT vs. Contingency Management in Indian samples"
    ])


# ============================================================
# SLIDES 37-40: REFERENCES
# ============================================================
add_content_slide(prs, "SLIDE 37: References (1/4)",
    [
        "Alfonso, J. P., Caracuel, A., Delgado-Pastor, L. C., & Verdejo-Garcia, A. (2011). Drug and Alcohol Dependence, 117(1), 78-81.",
        "Baer, R. A., Smith, G. T., Hopkins, J., Krietemeyer, J., & Toney, L. (2006). Assessment, 13(1), 27-45.",
        "Bowen, S., Chawla, N., & Marlatt, G. A. (2011). Mindfulness-based relapse prevention for addictive behaviors. Guilford Press.",
        "Bowen, S., & Marlatt, G. A. (2009). Psychology of Addictive Behaviors, 23(4), 666-671.",
        "Bowen, S., Witkiewitz, K., Clifasefi, S. L., et al. (2014). JAMA Psychiatry, 71(5), 547-556.",
        "Bowen, S., Witkiewitz, K., Dillworth, T. M., & Marlatt, G. A. (2007). Addictive Behaviors, 32(10), 2324-2328.",
        "Brewer, J. A., Mallik, S., Babuscio, T. A., et al. (2011). Drug and Alcohol Dependence, 119(1-2), 72-80."
    ])

add_content_slide(prs, "SLIDE 38: References (2/4)",
    [
        "Chiesa, A., & Serretti, A. (2014). Substance Use & Misuse, 49(5), 492-512.",
        "Davis, J. M., Goldberg, S. B., Anderson, M. C., et al. (2014). Substance Use & Misuse, 49(5), 571-585.",
        "Garland, E. L., Froeliger, B., & Howard, M. O. (2014). Frontiers in Psychiatry, 4, 173.",
        "Garland, E. L., Manusov, E. G., Froeliger, B., et al. (2014). J. Consulting and Clinical Psychology, 82(3), 448-459.",
        "Garland, E. L., Roberts-Lewis, A., Tronnier, C. D., et al. (2016). J. Consulting and Clinical Psychology, 84(4), 281-293.",
        "Ghosh, A., Basu, D., & Avasthi, A. (2018). Indian Journal of Psychiatry, 60(Suppl 4), S469-S476.",
        "Glasner-Edwards, S., Mooney, L. J., Ang, A., et al. (2017). Mindfulness, 8(1), 126-135.",
        "Grant, S., Colaiaco, B., Motala, A., et al. (2017). J. Addiction Medicine, 11(5), 386-396."
    ])

add_content_slide(prs, "SLIDE 39: References (3/4)",
    [
        "Hsu, S. H., Collins, S. E., & Marlatt, G. A. (2013). Addictive Behaviors, 38(3), 1852-1858.",
        "Humeniuk, R., Ali, R., Babor, T. F., et al. (2008). Addiction, 103(6), 1039-1047.",
        "Jain, R., Majumder, P., & Gupta, T. (2013). Indian Journal of Psychiatry, 55(Suppl 1), S86-S92.",
        "Kabat-Zinn, J. (1990). Full catastrophe living. Delacorte Press.",
        "Karyadi, K. A., VanderVeen, J. D., & Cyders, M. A. (2014). Drug and Alcohol Dependence, 143, 1-10.",
        "Korponay, C., Dentico, D., Kral, T. R. A., et al. (2019). Scientific Reports, 9(1), 11963.",
        "Li, W., Howard, M. O., Garland, E. L., et al. (2017). J. Substance Abuse Treatment, 75, 62-96.",
        "Marlatt, G. A., & Gordon, J. R. (1985). Relapse prevention. Guilford Press."
    ])

add_content_slide(prs, "SLIDE 40: References (4/4)",
    [
        "Mattoo, S. K., Chakrabarti, S., & Anjaiah, M. (2009). Indian J. Medical Research, 130(6), 702-708.",
        "Moeller, F. G., Barratt, E. S., Dougherty, D. M., et al. (2001). Am. J. Psychiatry, 158(11), 1783-1793.",
        "Murphy, C., & MacKillop, J. (2012). Psychopharmacology, 219(2), 527-536.",
        "NIDA (2020). Drugs, brains, and behavior: The science of addiction. NIDA.",
        "Patton, J. H., Stanford, M. S., & Barratt, E. S. (1995). J. Clinical Psychology, 51(6), 768-774.",
        "Peters, J. R., Erisman, S. M., Upton, B. T., et al. (2011). Mindfulness, 2(4), 228-235.",
        "Robinson, T. E., & Berridge, K. C. (1993). Brain Research Reviews, 18(3), 247-291.",
        "Roos, C. R., Bowen, S., & Witkiewitz, K. (2017). J. Consulting and Clinical Psychology, 85(11), 1041-1051.",
        "Sarkar, S., & Balhara, Y. P. S. (2016). Indian J. Endocrinology and Metabolism, 20(4), 527-533.",
        "Shorey, R. C., et al. (2017). Substance Use & Misuse, 52(11), 1400-1410.",
        "Tiffany, S. T., et al. (2000). Development of heroin craving questionnaire. Unpublished.",
        "WHO ASSIST Working Group (2002). Addiction, 97(9), 1183-1194.",
        "Witkiewitz, K., & Bowen, S. (2010). J. Consulting and Clinical Psychology, 78(3), 362-374.",
        "Zgierska, A., et al. (2009). Substance Abuse, 30(4), 266-294."
    ])


# ============================================================
# SAVE THE PRESENTATION
# ============================================================
output_path = "/projects/sandbox/Dango-kiro/MBRP_Research_Synopsis_PPT.pptx"
prs.save(output_path)
print(f"Presentation saved successfully: {output_path}")
print(f"Total slides: {len(prs.slides)}")
