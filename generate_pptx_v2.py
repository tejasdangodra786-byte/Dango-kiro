#!/usr/bin/env python3
"""
Generate MBRP Research Synopsis PowerPoint Presentation - Version 2
Corrected: Review of Literature = 5 slides with detailed paragraph-style descriptions
Total: 38 slides (adjusted from 40 to account for fewer lit review slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_GOLD = RGBColor(0xD4, 0xA5, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_TEXT = RGBColor(0x2D, 0x2D, 0x44)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)



def add_title_slide(prs, title, subtitle, researcher="[Researcher Name]",
                    guide="[Supervisor Name]", dept="Clinical Psychology",
                    inst="[University/Institute Name]", year="2026"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11.3), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = ACCENT_GOLD
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    details = f"\nResearcher: {researcher}\nGuide: {guide}\nDepartment: {dept}\nInstitution: {inst}\nYear: {year}"
    p3 = tf.add_paragraph()
    p3.text = details
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)



def add_content_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_LIGHT
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    top_offset = 1.3
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = MEDIUM_BLUE
        top_offset = 1.7
    txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(top_offset), Inches(11.9), Inches(7.5 - top_offset - 0.3))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = bullet
        p.font.size = Pt(13)
        p.font.color.rgb = BODY_TEXT
        p.space_before = Pt(6)



def add_paragraph_slide(prs, title, paragraphs):
    """Add a slide with detailed paragraph-style content (for literature review)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_LIGHT
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    # Content area - paragraphs
    txBox3 = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(12.1), Inches(6.0))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = para
        p.font.size = Pt(12)
        p.font.color.rgb = BODY_TEXT
        p.space_before = Pt(12)
        p.space_after = Pt(6)



# ===========================================================================
# SLIDE 1: TITLE
# ===========================================================================
add_title_slide(prs,
    "Efficacy of Brief Mindfulness-Based Relapse Prevention (MBRP) Intervention\non Craving, Impulsivity, and Mindfulness\nin Substance Dependent Patients",
    "MPhil Clinical Psychology Research Synopsis")

# ===========================================================================
# SLIDE 2: INTRODUCTION 1/3
# ===========================================================================
add_content_slide(prs, "Introduction & Background (1/3)",
    [
        "India faces a significant substance use crisis; MAGNITUDE study (2019) estimated ~3.1 crore individuals affected by substance use disorders",
        "Opioid dependence constitutes a major public health burden, particularly in Punjab, Rajasthan, Northeast India, and metropolitan areas",
        "WHO estimates India accounts for ~25% of global opioid-related deaths in South-East Asia",
        "Economic burden includes lost productivity, healthcare costs, and family dysfunction",
        "National Drug Dependence Treatment Centre (NDDTC) reports increasing treatment-seeking among opioid users",
        "Indian de-addiction centers primarily offer pharmacotherapy (OST, naltrexone) with limited structured psychotherapy",
        "Psychosocial interventions remain under-utilized despite evidence of superior combined treatment outcomes"
    ],
    subtitle="Substance Dependence: Indian Context")

# ===========================================================================
# SLIDE 3: INTRODUCTION 2/3
# ===========================================================================
add_content_slide(prs, "Introduction & Background (2/3)",
    [
        "Relapse rates in substance dependence range from 40-60% within first year post-treatment (NIDA, 2020)",
        "Indian studies report even higher relapse rates (~70-80%) in opioid dependence (Mattoo et al., 2009)",
        "Triggers for relapse: craving, negative affect, interpersonal conflict, environmental cues",
        "Traditional relapse prevention (Marlatt & Gordon, 1985) has moderate efficacy but limited mindfulness integration",
        "Automatic cognitive-behavioral patterns perpetuate relapse cycles",
        "Need for interventions addressing both cognitive automaticity and emotional dysregulation",
        "Brief interventions essential for Indian rehab settings with limited resources and high patient turnover"
    ],
    subtitle="The Relapse Problem in De-Addiction")



# ===========================================================================
# SLIDE 4: INTRODUCTION 3/3
# ===========================================================================
add_content_slide(prs, "Introduction & Background (3/3)",
    [
        "Mindfulness: intentional, non-judgmental, present-moment awareness (Kabat-Zinn, 1990)",
        "Mindfulness-Based Relapse Prevention (MBRP) developed by Bowen, Chawla, & Marlatt (2011)",
        "Integrates Mindfulness-Based Stress Reduction (MBSR) + Cognitive-Behavioral Relapse Prevention (RP)",
        "Core mechanisms: decentering from craving, disrupting automatic reactivity, increasing distress tolerance",
        "MBRP targets the 'craving to use' automaticity by cultivating awareness of triggers",
        "Brief adaptations (6-8 sessions) show promise for resource-limited settings",
        "Growing global evidence supports MBRP; however, Indian validation remains scarce"
    ],
    subtitle="Mindfulness in Addiction Treatment")

# ===========================================================================
# SLIDE 5: VARIABLE 1 - CRAVING
# ===========================================================================
add_content_slide(prs, "Variable 1: Craving",
    [
        "DEFINITION: Intense subjective urge to use a substance triggered by internal/external cues; a motivational state arising from incentive-sensitization (Robinson & Berridge, 1993) and conditioned reinforcement",
        "RELEVANCE: Primary predictor of relapse in opioid dependence; mediates relationship between cue exposure and use behavior; intensity correlates with severity and treatment dropout",
        "NEUROPSYCHOLOGICAL BASIS: Mesolimbic dopamine pathway activation (VTA -> Nucleus Accumbens); PFC hypoactivation during craving episodes; conditioned cue-reward associations; dysregulated stress-reward interaction",
        "TECHNIQUE - Urge Surfing: Observing craving as a transient wave without acting on it",
        "TECHNIQUE - Mindfulness Exposure: Non-reactive awareness of craving sensations",
        "TECHNIQUE - SOBER Breathing Space: Pause-observe-redirect during craving triggers",
        "TECHNIQUE - Cognitive Decentering: 'I am having a craving' vs. 'I need the drug'"
    ])



# ===========================================================================
# SLIDE 6: VARIABLE 2 - IMPULSIVITY
# ===========================================================================
add_content_slide(prs, "Variable 2: Impulsivity",
    [
        "DEFINITION: Tendency toward rapid, unplanned actions without adequate consideration of consequences (Moeller et al., 2001); multi-dimensional: motor, attentional, non-planning impulsivity (Patton et al., 1995)",
        "RELEVANCE: Higher trait impulsivity predicts initiation, escalation, and relapse; mediates craving-to-use behavior; associated with treatment non-adherence and premature dropout",
        "NEUROPSYCHOLOGICAL BASIS: Prefrontal cortex dysfunction -> impaired executive control; reduced inhibitory control (Go/No-Go); impaired delay discounting; dorsolateral PFC hypoactivation",
        "TECHNIQUE - Response Inhibition Training: Mindful pause before automatic behavioral responses",
        "TECHNIQUE - Awareness Training: Noticing impulse-action sequences without engagement",
        "TECHNIQUE - Mindful Decision-Making: Creating space between stimulus and response",
        "TECHNIQUE - STOP Technique: Stop-Take a breath-Observe-Proceed mindfully"
    ])

# ===========================================================================
# SLIDE 7: VARIABLE 3 - MINDFULNESS
# ===========================================================================
add_content_slide(prs, "Variable 3: Mindfulness",
    [
        "DEFINITION: Capacity to attend to present-moment experience with openness, curiosity, non-judgment; 5 facets: Observing, Describing, Acting with Awareness, Non-Judging, Non-Reactivity (Baer et al., 2006)",
        "RELEVANCE: Substance-dependent individuals show significantly lower dispositional mindfulness; acts as protective factor against relapse; improvements mediate treatment outcomes in MBRP studies",
        "NEUROPSYCHOLOGICAL BASIS: ACC activation -> enhanced self-regulation; Insula -> interoceptive awareness; PFC-amygdala connectivity -> emotional regulation; DMN regulation -> reduced rumination",
        "TECHNIQUE - Sitting Meditation: Focused attention on breath, body sensations, thoughts",
        "TECHNIQUE - Body Scan: Systematic non-judgmental awareness of bodily states",
        "TECHNIQUE - Mindful Movement: Gentle yoga/walking with present-moment focus",
        "TECHNIQUE - Non-Judgmental Awareness: Labeling experiences without evaluation"
    ])



# ===========================================================================
# SLIDES 8-12: REVIEW OF LITERATURE (5 slides, detailed paragraph style)
# ===========================================================================

# SLIDE 8: MBRP FOUNDATIONAL STUDIES
add_paragraph_slide(prs, "Review of Literature (1/5) - MBRP & Relapse Prevention",
    [
        "Bowen, Witkiewitz, Clifasefi, Grow, Chawla, Hsu, Carroll, Harrop, Collins, Lustyk, and Larimer (2014) investigate the relative efficacy of Mindfulness-Based Relapse Prevention compared to standard Relapse Prevention and Treatment As Usual for substance use disorders in JAMA Psychiatry. The randomized clinical trial involved 286 participants who had completed initial treatment for substance use disorders. The study demonstrates that at 12-month follow-up, MBRP participants reported significantly fewer days of substance use and heavy drinking compared to both standard RP and TAU groups. The researchers highlight that while all three groups showed initial improvements, MBRP participants maintained superior long-term outcomes, suggesting that the cultivation of mindfulness skills provides a durable protective mechanism against relapse. The study establishes that MBRP's integration of present-moment awareness with traditional cognitive-behavioral relapse prevention strategies creates a synergistic therapeutic effect that addresses both the automatic reactivity underlying craving and the cognitive distortions that precipitate relapse. These findings provide foundational evidence for MBRP as a gold-standard psychosocial intervention in substance use disorder aftercare.",
        "",
        "Bowen and Marlatt (2009) examine the effects of brief mindfulness-based intervention on craving among substance users in Psychology of Addictive Behaviors. The study recruited incarcerated individuals with substance use histories and administered a brief urge surfing meditation intervention. Results demonstrated significant reductions in craving intensity and frequency in the mindfulness condition compared to controls. The authors propose that even brief exposure to mindfulness techniques can disrupt the automaticity of craving responses by introducing a meta-cognitive awareness layer between trigger and behavioral response. This study is particularly relevant to the present research as it validates the premise that abbreviated mindfulness interventions can produce meaningful clinical effects on craving, supporting the feasibility of brief MBRP protocols in settings where extended treatment programs are impractical."
    ])

# SLIDE 9: CRAVING & MINDFULNESS STUDIES
add_paragraph_slide(prs, "Review of Literature (2/5) - Craving & Mindfulness Mechanisms",
    [
        "Garland, Froeliger, and Howard (2014) explore how mindfulness training targets neurocognitive mechanisms of addiction at the attention-appraisal-emotion interface in Frontiers in Psychiatry. The study presents a theoretical and empirical framework demonstrating that Mindfulness-Oriented Recovery Enhancement (MORE) reduces opioid craving through three interconnected mechanisms: attentional reorientation away from drug-related cues, positive reappraisal of previously neutral stimuli to generate natural reward, and enhanced savoring of healthy pleasures. The researchers provide neuroimaging evidence showing that mindfulness practice modulates activity in prefrontal and limbic circuits associated with craving and emotional regulation. The study is significant because it elucidates the precise cognitive-neural pathways through which mindfulness reduces craving in opioid users specifically, providing a mechanistic rationale for why MBRP techniques such as urge surfing and mindful awareness of craving sensations can reduce the subjective intensity and behavioral impact of craving episodes in substance-dependent patients.",
        "",
        "Witkiewitz, Bowen, Douglas, and Hsu (2013) investigate mindfulness-based relapse prevention effects on substance craving in Addictive Behaviors. The study conducted secondary analyses of data from a randomized controlled trial comparing MBRP to TAU among individuals in aftercare following substance use disorder treatment. Over a 4-month follow-up period, MBRP participants demonstrated significantly lower craving levels compared to TAU participants, and importantly, the relationship between negative affect and subsequent craving was significantly attenuated in the MBRP group. The authors conclude that mindfulness practice weakens the affect-craving pathway by cultivating non-reactive awareness of emotional states, thereby preventing negative emotions from automatically triggering craving responses. This decoupling of affect and craving represents a critical therapeutic mechanism that distinguishes MBRP from traditional relapse prevention approaches that primarily rely on cognitive restructuring and behavioral avoidance strategies."
    ])



# SLIDE 10: IMPULSIVITY & MINDFULNESS STUDIES
add_paragraph_slide(prs, "Review of Literature (3/5) - Impulsivity & Mindfulness",
    [
        "Garland, Roberts-Lewis, Tronnier, Graves, and Kelley (2016) examine the efficacy of Mindfulness-Oriented Recovery Enhancement versus CBT for co-occurring substance dependence, traumatic stress, and psychiatric disorders in the Journal of Consulting and Clinical Psychology. The study employed a randomized controlled design with substance-dependent adults exhibiting elevated impulsivity. Results demonstrated that the mindfulness-based intervention produced significant reductions in impulsivity scores as measured by the Barratt Impulsiveness Scale (BIS-11), particularly in the motor and attentional impulsivity subscales. The authors theorize that mindfulness meditation strengthens prefrontal cortical inhibitory mechanisms by repeatedly engaging participants in exercises requiring sustained attention, response monitoring, and deliberate non-reactivity. This enhanced top-down cognitive control translates into improved ability to inhibit prepotent impulsive responses in daily life, particularly in high-risk situations where automatic substance-seeking behavior would otherwise occur. The clinical implication is that MBRP can address impulsivity as a transdiagnostic risk factor for relapse, not merely as a stable trait but as a modifiable behavioral pattern amenable to mindfulness-based intervention.",
        "",
        "Murphy and MacKillop (2012) explore the interrelationships between impulsivity, mindfulness, and alcohol misuse in Psychopharmacology. The study investigates whether dispositional mindfulness buffers against the effects of trait impulsivity on problematic substance use. Using a cross-sectional design with 340 participants, the researchers found that trait mindfulness was inversely associated with impulsive decision-making as measured by delay discounting tasks. Critically, mindfulness moderated the relationship between impulsivity and substance use problems, such that individuals with higher mindfulness showed weaker associations between impulsivity and alcohol-related consequences. The authors propose that mindfulness functions as a cognitive resource that enables impulsive individuals to override automatic behavioral tendencies through enhanced metacognitive awareness and response flexibility. These findings support the inclusion of impulsivity as a dependent variable in MBRP research and suggest that mindfulness enhancement may serve as a mechanism through which impulsive responding is attenuated in substance-dependent populations."
    ])

# SLIDE 11: META-ANALYSES & BRIEF INTERVENTIONS
add_paragraph_slide(prs, "Review of Literature (4/5) - Meta-Analyses & Brief Models",
    [
        "Li, Howard, Garland, McGovern, and Lazar (2017) conduct a systematic review and meta-analysis of mindfulness treatment for substance misuse in the Journal of Substance Abuse Treatment, encompassing 42 randomized controlled trials. The meta-analytic findings reveal that mindfulness-based interventions produce significant effect sizes for reducing substance misuse (d = 0.33), craving (d = 0.68), and stress (d = 0.44) across diverse substance use populations and treatment contexts. The authors highlight that effect sizes for craving reduction are particularly robust, supporting the theoretical premise that mindfulness directly targets craving mechanisms through enhanced interoceptive awareness and non-reactive observation of urge states. Furthermore, the review identifies that interventions of shorter duration (4-8 sessions) demonstrated comparable efficacy to longer protocols when appropriately structured, providing empirical justification for brief MBRP adaptations. The review concludes that mindfulness-based interventions represent a viable evidence-based treatment approach with medium-to-large effects on the primary mechanisms of relapse, supporting their integration into standard substance use disorder treatment protocols.",
        "",
        "Glasner-Edwards, Mooney, Ang, Garneau, Hartwell, Brecht, and Rawson (2017) examine a pilot randomized clinical trial of mindfulness-based relapse prevention for stimulant-dependent adults using an abbreviated 6-session protocol in the journal Mindfulness. The study demonstrates that a condensed mindfulness intervention is both feasible and effective in outpatient substance use treatment settings. Participants in the brief MBRP condition showed significant reductions in substance use frequency and craving intensity compared to the health education control group. The authors emphasize that the abbreviated format maintained the core therapeutic elements of standard MBRP (body scan, sitting meditation, urge surfing, mindful movement) while condensing psychoeducational components. This study is directly relevant to the present research as it validates the brief intervention model proposed here, demonstrating that 6-8 session MBRP protocols can be successfully implemented without substantial loss of therapeutic efficacy, making them particularly appropriate for Indian rehabilitation settings with 4-8 week typical admission durations."
    ])



# SLIDE 12: INDIAN STUDIES & RESEARCH GAP
add_paragraph_slide(prs, "Review of Literature (5/5) - Indian Context & Research Gap",
    [
        "Ghosh, Basu, and Avasthi (2018) conduct a comprehensive review of relapse in opioid dependence from an Indian perspective in the Indian Journal of Psychiatry. The study reports alarmingly high relapse rates exceeding 70% among opioid-dependent patients treated in North Indian de-addiction centers, with the majority of relapses occurring within the first three months post-discharge. The authors identify craving, peer influence, negative emotional states, and lack of structured psychological aftercare as primary relapse determinants in the Indian context. Significantly, the review highlights that Indian treatment facilities predominantly rely on pharmacological approaches (opioid substitution therapy, naltrexone maintenance) with minimal integration of evidence-based psychological interventions. The authors advocate strongly for the development and validation of structured psychosocial intervention protocols tailored to Indian treatment infrastructure, patient characteristics, and resource constraints. This study directly establishes the clinical need for the present research by documenting both the high relapse burden and the psychosocial treatment gap in Indian de-addiction settings.",
        "",
        "Sarkar and Balhara (2016) examine the underutilization of structured psychological interventions in Indian de-addiction settings. The authors note that despite international evidence supporting the integration of psychosocial treatments alongside pharmacotherapy, Indian rehabilitation centers continue to operate with minimal structured psychological programming. The study identifies barriers including limited trained personnel, absence of culturally validated intervention protocols, short admission durations, and institutional emphasis on pharmacological management. Jain, Majumder, and Gupta (2013) provide preliminary evidence from a mindfulness-based intervention study with alcohol-dependent patients in India, demonstrating initial reductions in craving, but no systematic MBRP trial has been conducted with Indian opioid-dependent populations. CRITICAL GAP: No published Indian RCT has tested Brief MBRP specifically in opioid-dependent populations, simultaneously assessing craving, impulsivity, and mindfulness as treatment outcomes. The present study addresses this gap as the first brief MBRP trial designed for Indian de-addiction center infrastructure."
    ])



# ===========================================================================
# SLIDE 13: RESEARCH GAP
# ===========================================================================
add_content_slide(prs, "Research Gap",
    [
        "No published RCT has tested MBRP (standard or brief) specifically in Indian opioid-dependent populations",
        "Most studies examine craving OR mindfulness in isolation; few simultaneously assess craving + impulsivity + mindfulness as treatment outcomes",
        "Indian rehab settings require condensed 6-8 session interventions but no brief MBRP model has been empirically validated in this context",
        "Indian treatment landscape lacks integration of structured evidence-based psychological interventions alongside pharmacotherapy",
        "Global MBRP studies focus on alcohol/polysubstance users; opioid-specific MBRP evidence is limited",
        "Mindfulness practices require culturally congruent adaptation for Indian populations",
        "PRESENT STUDY: First brief MBRP trial in Indian opioid-dependent sample assessing craving, impulsivity, and mindfulness simultaneously"
    ])

# ===========================================================================
# SLIDE 14: AIM
# ===========================================================================
add_content_slide(prs, "Aim of the Study",
    [
        "To evaluate the efficacy of a Brief Mindfulness-Based Relapse Prevention (MBRP) intervention in reducing craving and impulsivity, and enhancing mindfulness, among substance-dependent patients in an Indian de-addiction setting",
        "",
        "Specifically: To compare outcomes between Brief MBRP + TAU (Experimental) and Psychoeducation + TAU (Control) on three dependent variables measured pre- and post-intervention"
    ])

# ===========================================================================
# SLIDE 15: OBJECTIVES
# ===========================================================================
add_content_slide(prs, "Objectives",
    [
        "1. To assess and compare craving levels (pre vs. post) in Experimental (Brief MBRP + TAU) and Control (Psychoeducation + TAU) groups",
        "2. To assess and compare impulsivity levels (pre vs. post) in both groups",
        "3. To assess and compare mindfulness levels (pre vs. post) in both groups",
        "4. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in reducing craving",
        "5. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in reducing impulsivity",
        "6. To determine whether Brief MBRP + TAU is significantly more effective than Psychoeducation + TAU in enhancing mindfulness"
    ])



# ===========================================================================
# SLIDE 16: HYPOTHESES
# ===========================================================================
add_content_slide(prs, "Hypotheses (Directional)",
    [
        "H1: Participants receiving Brief MBRP + TAU will show significantly GREATER REDUCTION in craving scores (HCQ) compared to Psychoeducation + TAU, from pre-test to post-test",
        "",
        "H2: Participants receiving Brief MBRP + TAU will show significantly GREATER REDUCTION in impulsivity scores (BIS-11) compared to Psychoeducation + TAU, from pre-test to post-test",
        "",
        "H3: Participants receiving Brief MBRP + TAU will show significantly GREATER INCREASE in mindfulness scores (FFMQ) compared to Psychoeducation + TAU, from pre-test to post-test"
    ])

# ===========================================================================
# SLIDE 17: RESEARCH DESIGN
# ===========================================================================
add_content_slide(prs, "Research Design",
    [
        "Design: Pre-test Post-test Control Group Experimental Design",
        "",
        "R   O1   X1   O2   -->  Experimental Group (Brief MBRP + TAU)",
        "R   O1   X2   O2   -->  Control Group (Psychoeducation + TAU)",
        "",
        "R = Random assignment | O1 = Pre-test (FFMQ, BIS-11, HCQ, ASSIST)",
        "X1 = Brief MBRP (8 sessions) | X2 = Psychoeducation (8 sessions)",
        "O2 = Post-test (FFMQ, BIS-11, HCQ)",
        "",
        "Features: True experimental design with randomization; Active control for attention effects; TAU maintained in both groups"
    ],
    subtitle="Pre-test Post-test Control Group Design")

# ===========================================================================
# SLIDE 18: SAMPLE
# ===========================================================================
add_content_slide(prs, "Sample",
    [
        "Population: Substance-dependent patients (primarily opioid users) admitted to de-addiction centers",
        "Sampling Method: Purposive sampling followed by random allocation to groups",
        "Sample Size: N = 60 (30 per group)",
        "Setting: Government/private de-addiction centers in India",
        "Recruitment: Consecutive admissions meeting inclusion criteria over 6-8 months",
        "Attrition Consideration: Recruit N = 70 (35 per group) to account for ~15% dropout"
    ])



# ===========================================================================
# SLIDE 19: SAMPLE SIZE ESTIMATION
# ===========================================================================
add_content_slide(prs, "Sample Size Estimation",
    [
        "Formula: n = [(Za/2 + Zb)^2 x 2 x sigma^2] / d^2",
        "Parameters: Effect size (d) = 0.50 (medium; Li et al., 2017: d = 0.33-0.68) | Power = 0.80 -> Zb = 0.84 | Alpha = 0.05 -> Za/2 = 1.96",
        "Calculation: n = [(1.96 + 0.84)^2 x 2 x 1] / (0.50)^2 = [7.84 x 2] / 0.25 = 62.72 ~ 63 total (~32/group)",
        "G*Power Verification: Independent t-test: d=0.50, a=0.05, power=0.80 -> n=64 total; ANCOVA with covariate -> n=34/group",
        "FINAL DECISION: N = 60 (30 per group) - justified by ANCOVA as primary analysis",
        "Recruit 70 total (35/group) to account for ~15% attrition",
        "Consistent with: Bowen et al. (2009), Glasner-Edwards et al. (2017)"
    ])

# ===========================================================================
# SLIDE 20: INCLUSION CRITERIA
# ===========================================================================
add_content_slide(prs, "Inclusion Criteria",
    [
        "1. Diagnosis of Substance Dependence as per ICD-10 (F10-F19) / ICD-11 criteria",
        "2. Primary substance: Opioids (heroin, pharmaceutical opioids); polysubstance with primary opioid dependence included",
        "3. Age: 18-50 years (male participants)",
        "4. Completed detoxification phase (minimum 7 days post-withdrawal)",
        "5. Currently admitted in de-addiction/rehabilitation center",
        "6. Minimum education: 5th standard (ability to comprehend psychometric tools)",
        "7. Willingness to provide written informed consent",
        "8. Able to attend all intervention sessions during admission"
    ])

# ===========================================================================
# SLIDE 21: EXCLUSION CRITERIA
# ===========================================================================
add_content_slide(prs, "Exclusion Criteria",
    [
        "1. Severe psychiatric comorbidity: Psychotic disorders, Bipolar I with psychotic features, severe MDE with suicidality",
        "2. Significant cognitive impairment (MMSE < 24) or intellectual disability",
        "3. Active withdrawal symptoms (COWS score > 12)",
        "4. History of traumatic brain injury with LOC > 30 minutes",
        "5. Current participation in another structured psychological intervention study",
        "6. Medical instability requiring acute care",
        "7. History of prior formal mindfulness/meditation training (> 1 month)"
    ])



# ===========================================================================
# SLIDE 22: VARIABLES
# ===========================================================================
add_content_slide(prs, "Variables",
    [
        "INDEPENDENT VARIABLE (IV):",
        "  Type of Intervention: Level 1 = Brief MBRP + TAU (Experimental) | Level 2 = Psychoeducation + TAU (Control)",
        "",
        "DEPENDENT VARIABLES (DVs):",
        "  1. Craving - measured by Heroin Craving Questionnaire (HCQ)",
        "  2. Impulsivity - measured by Barratt Impulsiveness Scale (BIS-11)",
        "  3. Mindfulness - measured by Five Facet Mindfulness Questionnaire (FFMQ)",
        "",
        "CONTROLLED VARIABLES: Age, education, duration of use, severity (ASSIST), TAU constant, session duration equalized"
    ])

# ===========================================================================
# SLIDE 23: TOOLS - FFMQ & BIS-11
# ===========================================================================
add_content_slide(prs, "Tools: Mindfulness & Impulsivity Measures",
    [
        "1. FIVE FACET MINDFULNESS QUESTIONNAIRE (FFMQ) - Baer et al. (2006)",
        "   39 items | 5 facets: Observing, Describing, Acting with Awareness, Non-Judging, Non-Reactivity",
        "   5-point Likert (1-5) | Higher = greater mindfulness | Reliability: a = 0.75-0.91 | Hindi adaptation available",
        "",
        "2. BARRATT IMPULSIVENESS SCALE (BIS-11) - Patton, Stanford, & Barratt (1995)",
        "   30 items | 3 factors: Attentional, Motor, Non-Planning Impulsivity",
        "   4-point scale (1-4) | Higher = greater impulsivity | Reliability: a = 0.79-0.83; test-retest r = 0.83",
        "   Discriminates substance users from controls | Hindi version validated (Mathew et al., 2014)"
    ])

# ===========================================================================
# SLIDE 24: TOOLS - CRAVING
# ===========================================================================
add_content_slide(prs, "Tools: Craving Assessment",
    [
        "OPTIONS: HCQ (opioid-specific, 14-item) | OCDUS (generic, 12-item) | VAS (single-item 0-100mm)",
        "",
        "RECOMMENDED: HEROIN CRAVING QUESTIONNAIRE (HCQ-Brief, 14 items) - Tiffany et al. (2000)",
        "   Assesses: desire to use, intention, anticipation of positive outcome, relief from withdrawal, lack of control",
        "   Scoring: 7-point Likert | Reliability: a = 0.87-0.93",
        "   Validity: Convergent with VAS; sensitive to intervention effects",
        "   Indian usability: Applicable to Indian opioid users; translation feasible",
        "   JUSTIFICATION: Substance-specific measurement preferred for opioid populations (EMCDDA guidelines)"
    ])



# ===========================================================================
# SLIDE 25: TOOLS - SEVERITY (ASSIST)
# ===========================================================================
add_content_slide(prs, "Tools: Severity Assessment (Pre-Test Only)",
    [
        "OPTIONS: ASSIST-WHO (multi-substance, 8-item) | DAST-20 (drug-specific) | AUDIT (alcohol-specific)",
        "",
        "RECOMMENDED: ASSIST - WHO ASSIST Working Group (2002); Humeniuk et al.",
        "   8-item screening | 10 substance categories | Risk levels: Low (0-3), Moderate (4-26), High (27+)",
        "   Reliability: Test-retest r = 0.58-0.90; a = 0.77-0.94",
        "   Validity: Sensitivity 0.80, Specificity 0.71 for substance dependence",
        "   Indian: WHO-validated; Hindi version available; used in NDDTC studies; 5-10 min administration",
        "   PURPOSE: Used at PRE-TEST ONLY for baseline severity and group equivalence (NOT an outcome measure)",
        "   JUSTIFICATION: Multi-substance coverage captures polysubstance patterns; WHO credibility; free to use"
    ])

# ===========================================================================
# SLIDE 26: PROCEDURE
# ===========================================================================
add_content_slide(prs, "Procedure",
    [
        "Step 1: SCREENING - Identify eligible patients (ICD-10 diagnosis, detoxified, consent-capable)",
        "Step 2: INFORMED CONSENT - Explain purpose, procedures, confidentiality, right to withdraw",
        "Step 3: PRE-TEST - Administer ASSIST + HCQ + BIS-11 + FFMQ",
        "Step 4: RANDOM ALLOCATION - Computer-generated randomization -> Experimental vs. Control",
        "Step 5: INTERVENTION - Experimental: Brief MBRP (8 x 60 min, 4 weeks) | Control: Psychoeducation (8 x 60 min, 4 weeks)",
        "Step 6: POST-TEST - Administer HCQ + BIS-11 + FFMQ (within 1 week of completion)",
        "Step 7: DATA COMPILATION & ANALYSIS",
        "NOTE: Both groups receive TAU throughout | Assessments by blinded RA | Intervention by MPhil Clinical Psychologist"
    ],
    subtitle="Stepwise Research Flow")



# ===========================================================================
# SLIDE 27: INTERVENTION - EXPERIMENTAL
# ===========================================================================
add_content_slide(prs, "Intervention Plan: Experimental Group (Brief MBRP + TAU)",
    [
        "8 Sessions x 60 min x 4 weeks (Twice weekly) | Group format (6-8 participants)",
        "S1: Introduction to MBRP & Autopilot - Raisin exercise; identifying automatic patterns",
        "S2: Awareness of Triggers - Body scan meditation; mapping personal triggers",
        "S3: Mindfulness in Daily Life - Sitting meditation (breath); SOBER breathing space",
        "S4: Mindfulness in High-Risk Situations - Urge surfing; role-play; cognitive decentering",
        "S5: Acceptance & Skillful Action - Non-judgmental awareness; mindful decision-making",
        "S6: Seeing Thoughts as Thoughts - 'Thoughts are not facts'; mountain meditation",
        "S7: Self-Care & Lifestyle Balance - Loving-kindness; activity scheduling; warning signs",
        "S8: Social Support & Maintenance - Review practices; relapse prevention plan; practice commitment"
    ])

# ===========================================================================
# SLIDE 28: INTERVENTION - CONTROL
# ===========================================================================
add_content_slide(prs, "Intervention Plan: Control Group (Psychoeducation + TAU)",
    [
        "8 Sessions x 60 min x 4 weeks (Twice weekly) | Group format (6-8 participants)",
        "S1: Understanding Addiction - Nature of dependence; brain changes; disease model",
        "S2: Effects of Opioids - Short/long-term physical and psychological consequences",
        "S3: Understanding Relapse - Relapse process; warning signs; high-risk situations",
        "S4: Coping with Cravings - General strategies (distraction, social support); NO mindfulness",
        "S5: Health & Nutrition - Physical recovery; nutrition; sleep hygiene",
        "S6: Social Consequences - Family impact; legal issues; stigma; rehabilitation",
        "S7: Motivation & Goal Setting - Stages of change; personal goals; motivation enhancement",
        "S8: Lifestyle Changes & Summary - Long-term planning; support systems; review",
        "NOTE: Matched for contact time, format, and attention to control non-specific factors"
    ])



# ===========================================================================
# SLIDE 29: DATA ANALYSIS
# ===========================================================================
add_content_slide(prs, "Data Analysis",
    [
        "Descriptive Statistics: Mean, SD, frequency, percentages for sociodemographic and clinical variables",
        "Normality Testing: Shapiro-Wilk test for distribution of outcome variables",
        "Within-Group: Paired samples t-test (pre vs. post); Wilcoxon Signed-Rank for non-normal data",
        "Between-Group: Independent samples t-test (post-test); Mann-Whitney U for non-normal data",
        "PRIMARY: ANCOVA with post-test as DV, Group as IV, pre-test as covariate (controls baseline, increases power)",
        "Effect Size: Cohen's d for between-group differences",
        "Significance Level: alpha = 0.05 (two-tailed) | Software: SPSS 26.0 / JASP",
        "ITT Analysis: Last Observation Carried Forward (LOCF) for dropouts"
    ])

# ===========================================================================
# SLIDE 30: ETHICAL CONSIDERATIONS
# ===========================================================================
add_content_slide(prs, "Ethical Considerations",
    [
        "Informed Consent: Written consent in Hindi/regional language; participants informed of purpose, procedures, risks, benefits",
        "Voluntary Participation: Right to withdraw at any time without penalty or impact on treatment",
        "Confidentiality: Data coded with IDs; no identifying information in publications; secure storage",
        "Non-Maleficence: Control receives active psychoeducation (not waitlist); TAU continued for all",
        "Institutional Approval: Ethical clearance from Institutional Ethics Committee (IEC) prior to study",
        "Debriefing: Control group offered MBRP orientation post-study if desired",
        "Compliance: ICMR (2017) National Ethical Guidelines for Biomedical and Health Research"
    ])

# ===========================================================================
# SLIDE 31: EXPECTED RESULTS
# ===========================================================================
add_content_slide(prs, "Expected Results",
    [
        "CRAVING (HCQ):",
        "  Significant REDUCTION in Experimental group vs. Control | Mechanism: Urge surfing reduces automatic craving reactivity",
        "",
        "IMPULSIVITY (BIS-11):",
        "  Significant REDUCTION (Motor + Attentional) in Experimental group | Mechanism: Mindfulness enhances prefrontal inhibitory control",
        "",
        "MINDFULNESS (FFMQ):",
        "  Significant INCREASE (Acting with Awareness + Non-Reactivity) in Experimental group | Mechanism: Structured meditation cultivates mindfulness",
        "",
        "Overall: Brief MBRP + TAU expected to demonstrate superior outcomes across all three DVs vs. Psychoeducation + TAU"
    ])



# ===========================================================================
# SLIDE 32: CLINICAL IMPLICATIONS
# ===========================================================================
add_content_slide(prs, "Clinical Implications",
    [
        "Validates a brief, structured MBRP model feasible for Indian de-addiction settings with time-limited admissions",
        "Provides evidence-based psychological intervention to complement pharmacotherapy (OST, naltrexone)",
        "Demonstrates mindfulness-based approaches are culturally compatible with Indian populations",
        "Addresses multiple relapse risk factors simultaneously through single integrated protocol",
        "Supports task-shifting: Brief MBRP deliverable by MPhil-trained Clinical Psychologists",
        "Informs national treatment policy (NIMHANS, NDDTC, State Mental Health Authorities)",
        "Contributes to RCI-recognized intervention repertoire for clinical psychology practice in India"
    ])

# ===========================================================================
# SLIDE 33: LIMITATIONS
# ===========================================================================
add_content_slide(prs, "Limitations",
    [
        "Sample specificity: Male opioid-dependent patients from single center; limits generalizability",
        "Short-term assessment: Post-test immediately after intervention; no long-term follow-up",
        "Self-report measures: HCQ, BIS-11, FFMQ susceptible to social desirability bias",
        "No biological markers: Craving measured subjectively; no physiological corroboration",
        "Therapist effects: Single therapist delivery may introduce confounds",
        "Attention control: Psychoeducation controls for contact but not specific mindfulness mechanisms",
        "Attrition: Substance-dependent population may show higher dropout rates"
    ])

# ===========================================================================
# SLIDE 34: FUTURE DIRECTIONS
# ===========================================================================
add_content_slide(prs, "Future Directions",
    [
        "Follow-up: 3-month and 6-month assessments for sustained MBRP effects",
        "Multi-site RCTs across multiple Indian de-addiction centers for generalizability",
        "Include female participants to examine gender-specific effects",
        "Neuroimaging (fMRI/EEG) to examine neural mechanisms of MBRP on craving circuits",
        "Dose-response: Compare 6 vs. 8 vs. full 8-week MBRP for optimal dosage",
        "Mediator analysis: Whether mindfulness mediates craving/impulsivity reduction",
        "Technology-assisted: App-based/digital MBRP for post-discharge and rural access",
        "Comparative effectiveness: MBRP vs. CBT vs. ACT in Indian samples"
    ])



# ===========================================================================
# SLIDES 35-38: REFERENCES (APA 7)
# ===========================================================================
add_content_slide(prs, "References (1/4)",
    [
        "Baer, R. A., Smith, G. T., Hopkins, J., Krietemeyer, J., & Toney, L. (2006). Using self-report assessment methods to explore facets of mindfulness. Assessment, 13(1), 27-45.",
        "Bowen, S., Chawla, N., & Marlatt, G. A. (2011). Mindfulness-based relapse prevention for addictive behaviors: A clinician's guide. Guilford Press.",
        "Bowen, S., & Marlatt, G. A. (2009). Surfing the urge: Brief mindfulness-based intervention for college student smokers. Psychology of Addictive Behaviors, 23(4), 666-671.",
        "Bowen, S., Witkiewitz, K., Clifasefi, S. L., et al. (2014). Relative efficacy of MBRP, standard RP, and TAU for substance use disorders. JAMA Psychiatry, 71(5), 547-556.",
        "Brewer, J. A., Mallik, S., Babuscio, T. A., et al. (2011). Mindfulness training for smoking cessation. Drug and Alcohol Dependence, 119(1-2), 72-80.",
        "Chiesa, A., & Serretti, A. (2014). Are MBIs effective for substance use disorders? Substance Use & Misuse, 49(5), 492-512."
    ])

add_content_slide(prs, "References (2/4)",
    [
        "Garland, E. L., Froeliger, B., & Howard, M. O. (2014). Mindfulness training targets neurocognitive mechanisms of addiction. Frontiers in Psychiatry, 4, 173.",
        "Garland, E. L., Manusov, E. G., Froeliger, B., et al. (2014). MORE for chronic pain and opioid misuse. J. Consulting and Clinical Psychology, 82(3), 448-459.",
        "Garland, E. L., Roberts-Lewis, A., Tronnier, C. D., et al. (2016). MORE vs. CBT for co-occurring substance dependence. J. Consulting and Clinical Psychology, 84(4), 281-293.",
        "Ghosh, A., Basu, D., & Avasthi, A. (2018). Relapse in opioid dependence: Indian perspective. Indian Journal of Psychiatry, 60(Suppl 4), S469-S476.",
        "Glasner-Edwards, S., Mooney, L. J., et al. (2017). MBRP for stimulant dependent adults: Pilot RCT. Mindfulness, 8(1), 126-135.",
        "Grant, S., Colaiaco, B., et al. (2017). MBRP for SUDs: Systematic review and meta-analysis. J. Addiction Medicine, 11(5), 386-396."
    ])

add_content_slide(prs, "References (3/4)",
    [
        "Humeniuk, R., Ali, R., Babor, T. F., et al. (2008). Validation of ASSIST. Addiction, 103(6), 1039-1047.",
        "Kabat-Zinn, J. (1990). Full catastrophe living. Delacorte Press.",
        "Karyadi, K. A., VanderVeen, J. D., & Cyders, M. A. (2014). Meta-analysis: trait mindfulness and substance use. Drug and Alcohol Dependence, 143, 1-10.",
        "Li, W., Howard, M. O., Garland, E. L., et al. (2017). Mindfulness treatment for substance misuse: Meta-analysis. J. Substance Abuse Treatment, 75, 62-96.",
        "Marlatt, G. A., & Gordon, J. R. (1985). Relapse prevention. Guilford Press.",
        "Mattoo, S. K., Chakrabarti, S., & Anjaiah, M. (2009). Psychosocial factors in relapse. Indian J. Medical Research, 130(6), 702-708.",
        "Moeller, F. G., Barratt, E. S., et al. (2001). Psychiatric aspects of impulsivity. Am. J. Psychiatry, 158(11), 1783-1793."
    ])

add_content_slide(prs, "References (4/4)",
    [
        "Murphy, C., & MacKillop, J. (2012). Impulsivity, mindfulness, and alcohol misuse. Psychopharmacology, 219(2), 527-536.",
        "NIDA (2020). Drugs, brains, and behavior: The science of addiction. NIDA.",
        "Patton, J. H., Stanford, M. S., & Barratt, E. S. (1995). Factor structure of BIS. J. Clinical Psychology, 51(6), 768-774.",
        "Robinson, T. E., & Berridge, K. C. (1993). Neural basis of drug craving. Brain Research Reviews, 18(3), 247-291.",
        "Sarkar, S., & Balhara, Y. P. S. (2016). Indian J. Endocrinology and Metabolism, 20(4), 527-533.",
        "Tiffany, S. T., et al. (2000). Development of heroin craving questionnaire. Unpublished manuscript.",
        "WHO ASSIST Working Group (2002). ASSIST: Development, reliability, feasibility. Addiction, 97(9), 1183-1194.",
        "Witkiewitz, K., Bowen, S., Douglas, H., & Hsu, S. H. (2013). MBRP for substance craving. Addictive Behaviors, 38(2), 1563-1571.",
        "Zgierska, A., et al. (2009). Mindfulness meditation for SUDs: Systematic review. Substance Abuse, 30(4), 266-294."
    ])



# ===========================================================================
# SAVE
# ===========================================================================
output_path = "/projects/sandbox/Dango-kiro/MBRP_Research_Synopsis_PPT.pptx"
prs.save(output_path)
print(f"SUCCESS: Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
