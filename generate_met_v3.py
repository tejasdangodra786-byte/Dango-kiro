#!/usr/bin/env python3
"""
generate_met_v3.py
Generates a comprehensive 150+ slide PowerPoint presentation on
Motivational Enhancement Therapy (MET).
"""

import traceback
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# === CONFIGURATION ===
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2C, 0x5F, 0x8A)
TEAL = RGBColor(0x17, 0xA2, 0xB8)
GOLD = RGBColor(0xD4, 0xA0, 0x1C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

FONT_NAME = "Times New Roman"
OUTPUT_PATH = "/projects/sandbox/Dango-kiro/MET_Comprehensive_Presentation.pptx"

# Zone definitions (Y positions)
TITLE_Y = Inches(0)
TITLE_H = Inches(1.1)
CONTENT_Y = Inches(1.2)
CONTENT_H = Inches(4.8)
FOOTER_Y = Inches(6.1)
FOOTER_H = Inches(1.2)



def set_font(run, size=14, bold=False, color=DARK_GRAY, italic=False):
    """Set font properties for a run."""
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic


def add_title_bar(slide, title_text):
    """Add navy title bar at top of slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), TITLE_Y, SLIDE_WIDTH, TITLE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_top = Pt(12)
    tf.margin_left = Pt(24)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    set_font(run, size=26, bold=True, color=WHITE)


def add_footer(slide, takeaway, reference=""):
    """Add takeaway/reference zone at bottom."""
    # Takeaway box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), FOOTER_Y,
        Inches(12.7), FOOTER_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
    shape.line.color.rgb = TEAL
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Key Takeaway: {takeaway}"
    set_font(run, size=11, bold=True, color=BLUE)
    if reference:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = f"Reference: {reference}"
        set_font(run2, size=9, italic=True, color=DARK_GRAY)



def content_slide(prs, title, paragraphs, takeaway, reference="", notes=""):
    """Standard content slide with title bar, paragraph text, and footer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    # Content text box
    txBox = slide.shapes.add_textbox(
        Inches(0.5), CONTENT_Y, Inches(12.3), CONTENT_H
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    tf.margin_left = Pt(6)
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = para_text
        set_font(run, size=13, color=DARK_GRAY)
    add_footer(slide, takeaway, reference)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def table_slide(prs, title, headers, rows, takeaway, reference="", notes=""):
    """Slide with a formatted table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl_width = Inches(12.0)
    tbl_height = min(CONTENT_H, Inches(0.4 * num_rows))
    left = Inches(0.6)
    top = CONTENT_Y + Inches(0.1)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, tbl_width, tbl_height
    )
    table = table_shape.table
    # Set headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                set_font(run, size=11, bold=True, color=WHITE)
    # Set data rows
    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(cell_text)
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_font(run, size=10, color=DARK_GRAY)
    add_footer(slide, takeaway, reference)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide



def two_col_slide(prs, title, left_title, left_items, right_title, right_items,
                  takeaway, reference="", notes=""):
    """Two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    # Left column header
    left_hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), CONTENT_Y,
        Inches(6.0), Inches(0.5)
    )
    left_hdr.fill.solid()
    left_hdr.fill.fore_color.rgb = BLUE
    left_hdr.line.fill.background()
    tf = left_hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = left_title
    set_font(run, size=13, bold=True, color=WHITE)
    # Left column content
    left_box = slide.shapes.add_textbox(
        Inches(0.5), CONTENT_Y + Inches(0.6), Inches(5.8), Inches(4.0)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"• {item}"
        set_font(run, size=11, color=DARK_GRAY)
    # Right column header
    right_hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), CONTENT_Y,
        Inches(6.0), Inches(0.5)
    )
    right_hdr.fill.solid()
    right_hdr.fill.fore_color.rgb = TEAL
    right_hdr.line.fill.background()
    tf = right_hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = right_title
    set_font(run, size=13, bold=True, color=WHITE)
    # Right column content
    right_box = slide.shapes.add_textbox(
        Inches(6.9), CONTENT_Y + Inches(0.6), Inches(5.8), Inches(4.0)
    )
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"• {item}"
        set_font(run, size=11, color=DARK_GRAY)
    add_footer(slide, takeaway, reference)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide



def process_slide(prs, title, steps, takeaway, reference="", notes=""):
    """Process flow slide with connected step boxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    num_steps = len(steps)
    box_w = Inches(min(2.5, 11.0 / num_steps))
    box_h = Inches(1.2)
    spacing = (Inches(12.3) - box_w * num_steps) / max(num_steps - 1, 1)
    start_x = Inches(0.5)
    y_pos = CONTENT_Y + Inches(1.5)
    for i, step_text in enumerate(steps):
        x = start_x + i * (box_w + spacing)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y_pos),
            int(box_w), int(box_h)
        )
        colors = [NAVY, BLUE, TEAL, GOLD, NAVY, BLUE, TEAL, GOLD]
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step_text
        set_font(run, size=10, bold=True, color=WHITE)
        # Add arrow between boxes
        if i < num_steps - 1:
            arrow_x = int(x) + int(box_w) + int(spacing * 0.2)
            arrow_shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x,
                int(y_pos + box_h * 0.35),
                int(spacing * 0.6), int(box_h * 0.3)
            )
            arrow_shape.fill.solid()
            arrow_shape.fill.fore_color.rgb = GOLD
            arrow_shape.line.fill.background()
    add_footer(slide, takeaway, reference)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def section_divider(prs, section_num, section_title, subtitle=""):
    """Full-page navy section divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Full navy background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    # Section number
    num_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.0), Inches(11), Inches(1.5)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"SECTION {section_num}"
    set_font(run, size=18, bold=False, color=GOLD)
    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.2), Inches(11), Inches(2.0)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_title
    set_font(run, size=36, bold=True, color=WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        set_font(run2, size=16, color=TEAL)
    return slide



def build_presentation():
    """Build the complete MET presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # =====================================================================
    # TITLE SLIDE
    # =====================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "MOTIVATIONAL ENHANCEMENT THERAPY (MET)"
    set_font(run, size=36, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "A Comprehensive Clinical Training Program"
    set_font(run2, size=22, color=GOLD)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Evidence-Based Approach to Behavior Change"
    set_font(run3, size=16, color=TEAL)
    sub_box = slide.shapes.add_textbox(Inches(2), Inches(5.0), Inches(9), Inches(1.5))
    tf2 = sub_box.text_frame
    p4 = tf2.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = "Based on Miller & Rollnick (2023) MI 4th Edition, Project MATCH (1997), and Current Research"
    set_font(run4, size=12, color=WHITE)



    # =====================================================================
    # SECTION 1: Introduction and Overview of MET
    # =====================================================================
    section_divider(prs, 1, "Introduction and Overview of MET",
                    "Understanding the Foundations of Motivational Enhancement")

    content_slide(prs,
        "What is Motivational Enhancement Therapy?",
        [
            "Motivational Enhancement Therapy (MET) is a systematic, client-centered therapeutic approach designed to produce rapid, internally motivated change in addictive behaviors and other problematic patterns. Unlike approaches that guide the client through sequential recovery steps, MET employs motivational strategies to mobilize the client's own change resources.",
            "MET was originally developed as a brief intervention (typically 4 sessions) within the Project MATCH Research Group (1997), which was the largest psychotherapy trial ever conducted for alcohol use disorders with 1,726 participants across multiple sites in the United States.",
            "The therapy integrates principles from motivational interviewing (Miller & Rollnick, 1991, 2002, 2013, 2023) with structured assessment feedback to create a powerful catalyst for behavioral change. It is distinct from MI in that it includes specific structured feedback components based on comprehensive assessment.",
            "MET operates on the fundamental premise that motivation for change is not something that a therapist instills in a passive client, but rather is elicited from the client's own values, goals, and perceptions of discrepancy between current behavior and desired life outcomes.",
        ],
        "MET is a brief, evidence-based therapy that mobilizes internal motivation for change through structured feedback and MI principles.",
        "Miller, W.R., Zweben, A., DiClemente, C.C., & Rychtarik, R.G. (1995). MET Manual, Project MATCH. NIAAA.",
        notes="Emphasize that MET is NOT simply motivational interviewing. It is a structured protocol that combines MI spirit and techniques with personalized assessment feedback. The distinction matters because MET has specific session structures and feedback components that differentiate it from MI as a communication style. Historical context: Developed specifically for Project MATCH to provide a brief alternative to 12-session CBT and 12-step facilitation."
    )

    content_slide(prs,
        "Historical Development of MET",
        [
            "The development of MET traces back to William R. Miller's seminal 1983 paper in Behavioural Psychotherapy, where he first articulated the principles of motivational interviewing as an alternative to confrontational approaches that dominated addiction treatment in the 1970s and 1980s.",
            "The formalization of MET occurred through Project MATCH (Matching Alcoholism Treatments to Client Heterogeneity), funded by the National Institute on Alcohol Abuse and Alcoholism (NIAAA) in 1989. The research team needed a brief, principled intervention that could serve as a comparison condition against longer treatments.",
            "Miller, Zweben, DiClemente, and Rychtarik (1995) published the definitive MET treatment manual, outlining a 4-session protocol: Session 1 focused on structured feedback from assessment, Session 2 on strengthening commitment, and Sessions 3-4 on reinforcing progress and managing setbacks.",
            "Since Project MATCH's landmark findings (1997), MET has been adapted globally, including significant adaptations for low- and middle-income countries (LMICs), where Medrxiv (2023) systematic review found 7 of 11 studies showed improved outcomes when MI/MET was adapted for local contexts.",
        ],
        "MET evolved from Miller's 1983 MI paper through Project MATCH into a globally adapted evidence-based treatment.",
        "Project MATCH Research Group (1997). JOSA, 58, 7-29. N=1726 participants across 9 sites.",
        notes="Key historical points: 1) Miller's 1983 paper was revolutionary in its rejection of confrontation; 2) Project MATCH was groundbreaking in scale and rigor; 3) The 4-session format was intentionally brief to test whether outcomes could match 12-session treatments. Discuss how the addiction field was dominated by confrontational approaches and the Minnesota Model before MI/MET emerged."
    )



    content_slide(prs,
        "MET vs. Motivational Interviewing: Key Distinctions",
        [
            "While MET incorporates the spirit and techniques of Motivational Interviewing (MI), it is a distinct therapeutic protocol with structured elements that go beyond MI as a communication approach. Understanding this distinction is critical for clinical practice and research integrity.",
            "MI, as defined by Miller and Rollnick (2023, 4th edition), is a collaborative conversation style for strengthening a person's own motivation and commitment to change. It is organized around four tasks: engaging, focusing, evoking, and planning. MI can be applied across virtually any context where behavior change is relevant.",
            "MET, in contrast, is a specific treatment protocol that includes: (1) comprehensive baseline assessment using standardized instruments, (2) structured personalized feedback delivered in MI style, (3) a fixed number of sessions (typically 2-4), and (4) specific session agendas including check-in procedures, feedback delivery, and commitment strengthening.",
            "The clinical implication is that a therapist can practice MI without doing MET, but cannot properly deliver MET without proficiency in MI skills. MET requires both the relational foundation of MI and the structural elements of assessment-based feedback delivery.",
        ],
        "MET = MI spirit + structured assessment feedback protocol. MI is the communication style; MET is the complete treatment package.",
        "Miller & Rollnick (2023). MI 4th Ed: Helping People Change and Grow. Guilford Press.",
        notes="This distinction trips up many clinicians and researchers. MI is like a language you speak; MET is a specific conversation you have using that language with particular content (assessment feedback). Use analogies: MI is like English, MET is like a specific type of job interview conducted in English. Both use the same communication principles but MET has specific structural requirements."
    )

    two_col_slide(prs,
        "MET vs. Traditional Addiction Counseling",
        "Traditional Confrontational Approach",
        [
            "Therapist is the expert who diagnoses",
            "Accepts label (alcoholic) as necessary",
            "Confrontation used to break denial",
            "Resistance seen as pathological trait",
            "Focus on powerlessness and disease",
            "One pathway to recovery (abstinence)",
            "Patterson & Forgatch (1985): confrontation increases resistance",
        ],
        "MET Approach",
        [
            "Client is expert on own life and values",
            "Labels de-emphasized; focus on behavior",
            "Ambivalence explored without judgment",
            "Resistance signals need to adjust approach",
            "Focus on personal choice and self-efficacy",
            "Multiple pathways; client chooses goals",
            "Valle (1981): empathy predicts 2-year outcomes",
        ],
        "MET replaces confrontation with collaboration, viewing resistance as an interpersonal signal rather than a client deficit.",
        "Patterson & Forgatch (1985). JCCP, 53, 846; Valle (1981). J Studies Alcohol, 42, 783.",
        notes="This slide highlights the paradigm shift. Patterson & Forgatch (1985) experimentally demonstrated that when therapists were instructed to confront, client resistance increased dramatically. Valle (1981) showed that counselor empathy (not confrontation skill) predicted client drinking outcomes at 2-year follow-up. These two studies together provide the empirical rationale for the MI/MET approach: confrontation harms, empathy heals."
    )

    content_slide(prs,
        "Evidence Base: Project MATCH Landmark Findings",
        [
            "Project MATCH (1997) remains the largest and most rigorous psychotherapy trial in addiction treatment history. It randomized 1,726 alcohol-dependent outpatients and 774 aftercare clients to one of three treatments: 4-session MET, 12-session Cognitive-Behavioral Therapy (CBT), or 12-session Twelve-Step Facilitation (TSF).",
            "The primary finding was remarkable: all three treatments produced substantial and sustained improvements in drinking outcomes, with few significant differences between them. This meant that 4-session MET achieved comparable outcomes to treatments requiring three times as many sessions, representing enormous efficiency gains.",
            "At 1-year follow-up, clients in all three conditions showed approximately 50% reduction in drinking days and drinks per drinking day. At 3-year follow-up, gains were maintained across conditions. The few matching hypotheses that emerged (e.g., angry clients did better in MET than in other treatments) provided modest support for treatment matching.",
            "The cost-effectiveness implications were profound: MET achieved equivalent outcomes at approximately one-third the clinical contact time, making it particularly suitable for resource-limited settings - a finding that later motivated adaptation efforts in India and other LMICs.",
        ],
        "4-session MET matched 12-session therapies in outcomes, demonstrating remarkable efficiency for addiction treatment.",
        "Project MATCH Research Group (1997). Matching treatments to client heterogeneity. JOSA, 58, 7-29.",
        notes="Project MATCH cost approximately $27 million and involved 9 clinical sites. Key methodological strengths: random assignment, manualized treatments, extensive therapist training and supervision, standardized assessment battery, long-term follow-up. The matching hypothesis (that certain client characteristics would predict differential response) received limited support, but the equivalence finding was itself groundbreaking."
    )



    # =====================================================================
    # SECTION 2: Theoretical Foundations
    # =====================================================================
    section_divider(prs, 2, "Theoretical Foundations of MET",
                    "Psychological Theories Underpinning Motivational Enhancement")

    content_slide(prs,
        "Self-Determination Theory and MET",
        [
            "Self-Determination Theory (SDT), developed by Deci and Ryan (1985, 2000), provides perhaps the most robust theoretical foundation for understanding why MET works. SDT posits that humans have three basic psychological needs: autonomy (feeling volitional), competence (feeling effective), and relatedness (feeling connected to others).",
            "MET directly supports all three basic needs: Autonomy is supported through the emphasis on personal choice, avoidance of prescriptive advice-giving, and recognition of the client as the ultimate decision-maker. Competence is supported through affirmation of client strengths, building self-efficacy, and recognizing past successes. Relatedness is supported through the empathic therapeutic relationship.",
            "Research by Markland et al. (2005) demonstrated that MI-consistent approaches activate autonomous motivation (identified and integrated regulation) rather than controlled motivation (external pressure or introjected guilt). Autonomous motivation predicts sustained behavior change, while controlled motivation predicts only short-term compliance.",
            "In the Indian context, autonomy support may need cultural adaptation. NIMHANS (2008) emphasizes that in collectivist cultures, autonomy can be framed within family and community values rather than purely individual choice, while still preserving the client's sense of agency.",
        ],
        "SDT explains MET's effectiveness: supporting autonomy, competence, and relatedness activates lasting internal motivation.",
        "Deci, E.L. & Ryan, R.M. (1985, 2000). Self-Determination Theory. Rochester, NY.",
        notes="SDT is arguably the most important theoretical lens for MET. When we tell clients what to do, we undermine autonomy and trigger reactance. When we affirm their strengths, we build competence. When we express empathy, we meet relatedness needs. All three pathways increase intrinsic/autonomous motivation. The cultural adaptation point is crucial for Indian settings."
    )

    content_slide(prs,
        "Cognitive Dissonance Theory",
        [
            "Leon Festinger's (1957) Cognitive Dissonance Theory explains a core mechanism of MET: the therapeutic development of discrepancy. When individuals simultaneously hold two contradictory cognitions (e.g., 'I value being a good parent' and 'My drinking is harming my children'), they experience psychological discomfort (dissonance) that motivates change.",
            "In MET, the therapist strategically amplifies this discrepancy - not through confrontation, but by helping clients articulate both their values/goals and their current behavior patterns. The personalized feedback session is specifically designed to make discrepancies salient and personally meaningful.",
            "Critical distinction from confrontation: In confrontational approaches, the therapist points out the discrepancy and argues for change. In MET, the client discovers and articulates the discrepancy themselves. This self-generated insight produces stronger, more durable motivation because it emerges from the client's own reasoning process.",
            "DiClemente et al. (2017) in their analysis of MI mechanisms confirmed that perceived discrepancy between current behavior and personal values is one of the primary mechanisms through which MI/MET produces behavior change, operating through increased change talk and decreased sustain talk.",
        ],
        "MET leverages cognitive dissonance by helping clients see the gap between their values and behaviors - without confrontation.",
        "Festinger, L. (1957). A Theory of Cognitive Dissonance. Stanford University Press.",
        notes="The key pedagogical point: discrepancy development is NOT confrontation. Confrontation = therapist tells client there's a problem. Discrepancy development = therapist helps client discover it themselves. The difference in psychological process is enormous. Self-generated discrepancy activates autonomous motivation; externally imposed discrepancy activates reactance."
    )

    content_slide(prs,
        "Self-Perception Theory and Change Talk",
        [
            "Daryl Bem's (1972) Self-Perception Theory offers a complementary explanation for why eliciting change talk is therapeutic. The theory proposes that individuals infer their own attitudes and beliefs by observing their own behavior - including their verbal behavior. When a client hears themselves arguing for change, they infer that they must want to change.",
            "This mechanism explains the critical finding by Amrhein et al. (2003) that commitment language (a form of change talk) spoken by the client during MI sessions significantly predicted actual behavior change outcomes. Specifically, the strength of commitment statements in the final third of sessions predicted drug use outcomes at follow-up.",
            "Moyers et al. (2007) extended this understanding by demonstrating the causal chain: therapist MI-consistent behaviors (open questions, reflections, affirmations) lead to increased client change talk, which in turn leads to improved behavioral outcomes. This confirmed that the therapist's role is to evoke the client's own arguments for change.",
            "The practical implication for MET is clear: the therapist's primary task is not to argue for change but to create conditions where the client argues for change. Every open question, strategic reflection, and affirmation should be designed to evoke and reinforce the client's own motivational statements.",
        ],
        "When clients hear themselves argue for change, they strengthen their own commitment - making evocation the therapist's primary task.",
        "Bem, D.J. (1972). Self-Perception Theory. Advances in Experimental Social Psychology, 6, 1-62.",
        notes="This is one of the most powerful insights for new MI/MET practitioners. Many therapists instinctively want to give the 'change speech' - to convince, persuade, argue. But self-perception theory tells us that the person who gives the change speech is the person who changes. Our job is to get the CLIENT talking about change, not to talk about change AT them."
    )

    content_slide(prs,
        "The Transtheoretical Model and Stages of Change",
        [
            "Prochaska and DiClemente's (1982) Transtheoretical Model (TTM) provides the stage-based framework that originally informed MET's design. The model identifies five stages: Precontemplation (no awareness of need for change), Contemplation (ambivalent), Preparation (intending change), Action (actively changing), and Maintenance (sustaining change).",
            "MET was specifically designed to address precontemplation and contemplation stages, where traditional action-oriented treatments often fail because they assume a readiness that has not yet developed. The structured feedback component of MET is designed to move precontemplators toward contemplation by making personally relevant information salient.",
            "However, modern MI/MET practice has moved beyond rigid stage-matching. Miller and Rollnick (2023, 4th edition) note that motivation fluctuates within and across sessions, and that the four processes of MI (engaging, focusing, evoking, planning) are not strictly sequential but recursive. A client may need re-engagement after a setback.",
            "The clinical relevance remains: meeting clients where they are, not where we want them to be, is a foundational principle. Assessment of readiness (using instruments like the URICA or rulers) helps clinicians adjust their approach, spending more time on evoking with ambivalent clients and more on planning with those ready for action.",
        ],
        "Stages of change inform MET by emphasizing readiness assessment and meeting clients at their current motivational stage.",
        "Prochaska, J.O. & DiClemente, C.C. (1982). Transtheoretical therapy: Toward a more integrative model of change.",
        notes="Important nuance: While stages of change informed MET's development, contemporary practice is less rigid about stage-matching. Miller has noted that motivation is more fluid than discrete stages suggest. However, the core insight - that trying to plan action with an ambivalent client is premature and counterproductive - remains clinically vital."
    )

    content_slide(prs,
        "Bandura's Self-Efficacy Theory",
        [
            "Albert Bandura's (1977) Self-Efficacy Theory identifies perceived self-efficacy - one's belief in their ability to successfully perform a specific behavior - as a critical determinant of behavior change. Self-efficacy influences choice of activities, effort expenditure, persistence in the face of obstacles, and ultimately, success.",
            "In MET, supporting self-efficacy is one of the four guiding principles of the MI spirit (along with expressing empathy, developing discrepancy, and rolling with resistance). The therapist actively seeks opportunities to affirm the client's past successes, current strengths, and capacity for change.",
            "Bandura identified four sources of self-efficacy: mastery experiences (past successes), vicarious experiences (seeing similar others succeed), verbal persuasion (encouragement from credible sources), and physiological states (feeling capable vs. anxious). MET leverages all four through affirmation, normalizing change, strategic encouragement, and creating a calm therapeutic environment.",
            "Research specific to addiction shows that self-efficacy for abstinence or moderation is one of the strongest predictors of treatment outcome across substances (Maisto et al., 2000). MET's emphasis on building confidence thus has direct empirical support as a mechanism of therapeutic change.",
        ],
        "Self-efficacy - the client's belief in their ability to change - is both a target and mechanism of MET intervention.",
        "Bandura, A. (1977). Self-efficacy: Toward a unifying theory of behavioral change. Psychological Review, 84, 191-215.",
        notes="Practical applications: Never argue with a client's pessimism about change - instead, ask about times they DID succeed at difficult things. Use scaling questions: 'On a scale of 0-10, how confident are you that you could reduce your drinking if you decided to?' Then explore what gives them even that level of confidence. Build from existing strengths rather than dwelling on failures."
    )



    # =====================================================================
    # SECTION 3: The Spirit of MI in MET
    # =====================================================================
    section_divider(prs, 3, "The Spirit of MI in MET",
                    "Partnership, Acceptance, Compassion, and Evocation")

    content_slide(prs,
        "The Four Components of MI Spirit (2023 Update)",
        [
            "Miller and Rollnick (2023, 4th edition) define the spirit of MI through four interconnected components: Partnership, Acceptance, Compassion, and Evocation. These form the relational foundation without which specific MI techniques become manipulative rather than therapeutic.",
            "Partnership means the therapeutic relationship is collaborative rather than authoritarian. The therapist contributes expertise about change processes; the client contributes expertise about their own life, values, and experiences. Neither party holds all the knowledge needed for successful change.",
            "Acceptance encompasses four sub-components: absolute worth (unconditional positive regard, following Rogers, 1957), accurate empathy (understanding the client's perspective), autonomy support (honoring the client's right to choose), and affirmation (recognizing strengths and efforts). Acceptance does not mean approval of harmful behaviors.",
            "Compassion means the therapist actively prioritizes the client's welfare and interests. Evocation means drawing out the client's own motivations, values, and ideas rather than installing external ones. Together, these four elements create a therapeutic atmosphere where honest self-exploration becomes possible.",
        ],
        "MI spirit (partnership, acceptance, compassion, evocation) is the essential foundation - without it, techniques become manipulation.",
        "Miller & Rollnick (2023). Motivational Interviewing, 4th Ed. Guilford Press.",
        notes="The spirit is MORE important than techniques. Research consistently shows that MI delivered with proper spirit but imperfect technique outperforms technically proficient MI delivered without spirit. This is why the 4th edition places even more emphasis on the relational foundation. Discuss Rogers' influence - his 1957 paper on necessary and sufficient conditions for therapeutic change profoundly shaped MI."
    )

    content_slide(prs,
        "Rogers' Influence: Empathy as the Foundation",
        [
            "Carl Rogers' (1957) landmark paper identified empathy, unconditional positive regard, and genuineness as necessary and sufficient conditions for therapeutic personality change. Miller has consistently credited Rogers as the primary intellectual ancestor of MI/MET, noting that accurate empathy remains the single most important therapist skill.",
            "Valle (1981) provided dramatic empirical evidence for Rogers' claims in addiction treatment: counselor empathy scores predicted client drinking outcomes at 6-month and 2-year follow-up more strongly than any other variable studied. Clients of high-empathy counselors showed significantly less relapse.",
            "In MET practice, empathy is expressed primarily through reflective listening - the skill of offering back to the client your understanding of what they have said, with slightly deeper meaning or implication. This creates a cycle of feeling understood, exploring further, and deepening self-awareness.",
            "Modern neuroscience research supports Rogers' intuition: feeling understood activates the same brain regions involved in reward processing and reduces activation in threat-detection circuits. Empathy literally creates the neurological conditions for open exploration rather than defensive self-protection.",
        ],
        "Empathy is not just 'nice' - it is the active ingredient that creates neurological safety for honest self-exploration.",
        "Rogers, C.R. (1957). The necessary and sufficient conditions of therapeutic personality change. J Consulting Psychology, 21, 95-103.",
        notes="Valle's 1981 study is extraordinarily important for training purposes. It demonstrates that therapist empathy (not insight, not confrontation skill, not clinical experience) predicts addiction outcomes. This study should be a touchstone for trainees who wonder whether 'soft skills' really matter. They are not soft - they are the hardest and most important skills to develop."
    )

    process_slide(prs,
        "The Four Processes of MI in MET Sessions",
        ["ENGAGING\n(Build rapport)", "FOCUSING\n(Set agenda)", "EVOKING\n(Elicit change talk)", "PLANNING\n(Commit to action)"],
        "The four MI processes are recursive, not linear - therapists may need to return to earlier processes throughout treatment.",
        "Miller & Rollnick (2023). MI 4th Ed. The four processes replace the earlier two-phase model.",
        notes="The 4th edition reorganized MI around four processes rather than two phases. This is more clinically realistic - in practice, you may need to re-engage a client who becomes defensive, or refocus when the conversation drifts. The processes are like floors of a building: you need the foundation (engaging) to support the upper floors, but you can move between them."
    )

    two_col_slide(prs,
        "MI Spirit in Practice: What It Looks Like",
        "MI-Consistent (Spirit-Aligned)",
        [
            "What brings you here today?",
            "Help me understand your perspective",
            "You're the expert on your own life",
            "What matters most to you?",
            "You've shown real strength in...",
            "Where would you like to go from here?",
            "It's ultimately your decision",
        ],
        "MI-Inconsistent (Spirit Violations)",
        [
            "You need to accept you're an alcoholic",
            "Let me tell you what's wrong",
            "I know what's best for you",
            "You should want to be sober",
            "You're in denial about your problem",
            "Here's what you need to do",
            "You have to follow this plan",
        ],
        "The difference between MI-consistent and MI-inconsistent practice lies in who holds the power and expertise.",
        "Miller & Rollnick (2023). MI 4th Edition. Guilford Press.",
        notes="Use these contrasting examples in role-plays. Have trainees identify what makes each statement consistent or inconsistent with MI spirit. The key pattern: MI-consistent statements position the client as the agent; MI-inconsistent statements position the therapist as the authority. This is not about being 'nice' - it's about who is doing the psychological work of change."
    )

    content_slide(prs,
        "Cultural Adaptation of MI Spirit in India",
        [
            "Applying MI spirit in the Indian context requires thoughtful cultural adaptation while preserving core principles. NIMHANS (2008) Psychosocial Interventions Manual emphasizes that India has approximately 62.5 million alcohol users, and treatment approaches must be culturally congruent to be effective.",
            "Family involvement presents both an opportunity and a challenge. Indian culture's collectivist orientation means that family members often attend sessions and expect to participate. MET can be adapted to include family perspectives while maintaining the individual client's autonomy - for example, by exploring how family values align with the client's own desire for change.",
            "Nadkarni et al. (2023) in the AMBIT trial successfully adapted MI for delivery via mobile phone by community health workers (Lay counselors) in Goa, India. This demonstrated that MI spirit can be maintained even in brief, technology-mediated interventions when counselors receive appropriate training and supervision.",
            "Patel et al. (2024) further demonstrated cultural adaptation by combining MI with Behavioral Couples Therapy for men with alcohol use and intimate partner violence in south India, showing that collaborative approaches are culturally acceptable and effective in Indian settings.",
        ],
        "MI spirit can be successfully adapted to Indian cultural contexts while preserving core principles of respect and collaboration.",
        "Nadkarni et al. (2023). AMBIT Trial - mobile MI for alcohol in Goa. BMC Psychiatry.",
        notes="Cultural adaptation is not about abandoning MI principles - it's about expressing them through culturally meaningful channels. In India, respect for elders, family harmony, duty (dharma), and community standing can all be explored as personal values that motivate change. The therapist honors these cultural values while still supporting individual autonomy."
    )



    # =====================================================================
    # SECTION 4: Core Principles and Techniques
    # =====================================================================
    section_divider(prs, 4, "Core Principles and Techniques of MET",
                    "OARS Skills and Strategic Interventions")

    content_slide(prs,
        "OARS: The Foundation Skills of MET Practice",
        [
            "OARS represents the four foundational microskills of MI/MET practice: Open questions, Affirmations, Reflections, and Summaries. These skills form the building blocks from which all MI/MET interventions are constructed. Mastery of OARS is essential before attempting more complex strategies.",
            "Open Questions invite elaboration rather than yes/no answers. They begin with 'what,' 'how,' 'tell me about,' or 'in what ways.' Their purpose is to evoke the client's own perspective, particularly around values, concerns, and ideas about change. Strategic open questions can be directed toward eliciting change talk.",
            "Affirmations are genuine statements recognizing client strengths, efforts, values, and positive qualities. They differ from praise (which implies evaluation) by focusing on the client's inherent qualities rather than the therapist's judgment. Effective affirmations build self-efficacy and strengthen the therapeutic relationship.",
            "Reflections are statements (not questions) that offer back the therapist's understanding of what the client has communicated. They range from simple (parroting content) to complex (reflecting deeper meaning, feeling, or implication). Skilled reflections form the majority of MI-consistent therapist responses.",
        ],
        "OARS skills are the essential microskills - particularly reflections, which should comprise 50%+ of all therapist responses in MI/MET.",
        "Miller & Rollnick (2023). MI 4th Edition. Guilford Press.",
        notes="In fidelity-coded MI sessions, the reflection-to-question ratio should be at least 2:1. Many beginning practitioners ask too many questions and reflect too little. Practice tip: for every question you want to ask, first try converting it to a reflection. 'Did that make you angry?' becomes 'That made you angry.' The latter invites deeper exploration."
    )

    content_slide(prs,
        "The Art of Reflective Listening",
        [
            "Reflective listening is simultaneously the most important and most difficult skill in MET practice. It involves offering back to the client a statement of what you understand them to be saying or feeling, slightly extending the meaning to facilitate deeper exploration. Rogers called this 'accurate empathy.'",
            "Simple reflections repeat or slightly rephrase what the client said: Client: 'I drink most nights.' Therapist: 'Drinking has become a nightly pattern for you.' Complex reflections add meaning, feeling, or implication: 'And you're noticing that it's become more automatic than intentional.'",
            "Strategic reflections in MET are used to selectively attend to change talk while not amplifying sustain talk. When a client says 'I know I should cut back, but I enjoy relaxing with a beer,' the MI-consistent response reflects the change talk: 'Part of you recognizes that cutting back would be beneficial.' This is not dishonest - it accurately reflects one side of their ambivalence.",
            "Double-sided reflections capture both sides of ambivalence: 'On one hand, drinking helps you unwind; on the other hand, you're concerned about the impact on your health.' The conjunction used matters - 'and' or 'on the other hand' (linking both as valid) is preferred over 'but' (which negates the first part).",
        ],
        "Reflective listening is the core MI/MET skill - strategic use of reflections guides the conversation toward change.",
        "Moyers et al. (2007). JCCP, 75, 790. Therapist MI behaviors → client change talk → outcomes.",
        notes="Training exercise: Have participants practice the 'reflection formation' exercise - client makes a statement, therapist must respond with a reflection (not a question, not advice, not agreement). This is surprisingly difficult for most clinicians trained in assessment-heavy approaches. The shift from question-asking to statement-making is transformative."
    )

    table_slide(prs,
        "Types of Reflections in MET Practice",
        ["Type", "Definition", "Example", "Purpose"],
        [
            ["Simple", "Repeats/rephrases content", "'You've been drinking more lately'", "Shows listening, clarifies"],
            ["Complex - Feeling", "Identifies underlying emotion", "'That sounds frustrating'", "Deepens emotional awareness"],
            ["Complex - Meaning", "Adds deeper significance", "'It's become something you can't control'", "Builds discrepancy"],
            ["Amplified", "Slightly overstates position", "'So nothing about drinking concerns you at all'", "Evokes pushback/change talk"],
            ["Double-sided", "Captures both sides", "'You enjoy it AND worry about effects'", "Acknowledges ambivalence"],
            ["Continuing paragraph", "Extends story forward", "'And you're wondering where this leads'", "Invites further exploration"],
        ],
        "Master multiple reflection types to match the clinical moment and guide conversations toward change talk.",
        "Miller & Rollnick (2023). MI 4th Edition. Chapter on reflective listening.",
        notes="Amplified reflections should be used judiciously - if overdone, they feel sarcastic. The tone must be genuinely curious, not mocking. Double-sided reflections are particularly useful with highly ambivalent clients. Continuing-the-paragraph reflections show deep understanding and invite the client to go further."
    )

    content_slide(prs,
        "Developing Discrepancy Without Confrontation",
        [
            "Developing discrepancy is one of the four original MI principles and remains central to MET practice. It involves helping clients perceive a gap between their current behavior and their deeper values, goals, or self-image. This perceived gap creates internal motivation for change without requiring external pressure.",
            "In MET, discrepancy development is facilitated primarily through the personalized feedback session, where objective assessment data (drinking patterns, health markers, risk levels) is presented alongside the client's stated values and goals. The therapist does not argue - the data speaks for itself.",
            "Key discrepancy-developing questions include: 'How does your current drinking fit with being the kind of parent you want to be?' 'What would your life look like in 5 years if nothing changes?' 'How does this match with what you told me matters most to you?' These questions invite the client to identify their own discrepancy.",
            "DiClemente et al. (2017) identified discrepancy perception as one of the key mechanisms through which MI produces change. When clients perceive discrepancy, they experience cognitive dissonance that motivates resolution - ideally through behavior change rather than value modification.",
        ],
        "Discrepancy development is strategic - the therapist creates conditions for the client to discover their own reasons for change.",
        "DiClemente et al. (2017). Mechanisms of change in MI. Addiction, 112(S2), 92-100.",
        notes="Common mistake: telling the client about the discrepancy rather than helping them see it. 'Your drinking is inconsistent with your family values' is confrontation. 'How does your drinking fit with what you told me about wanting to be present for your children?' is discrepancy development. Same content, completely different process."
    )

    content_slide(prs,
        "Rolling with Resistance: Updated Understanding",
        [
            "The concept of 'rolling with resistance' has evolved significantly in MI/MET. Miller and Rollnick (2013, 2023) moved away from the term 'resistance' (which implies the client is the problem) toward 'sustain talk' and 'discord' (which acknowledge the interpersonal nature of the interaction).",
            "Sustain talk refers to client statements favoring the status quo - arguments against change, minimizing problems, pessimism about ability to change. These are normal expressions of ambivalence, not pathology. The therapeutic response is to acknowledge without amplifying: 'You're not sure change is possible right now.'",
            "Discord refers to disruption in the therapeutic relationship - the client arguing against the therapist, becoming defensive, or disengaging. Discord signals that the therapist needs to adjust their approach, typically by backing off, increasing reflection, or returning to engagement.",
            "Practical strategies for responding to sustain talk and discord include: simple reflection (acknowledging without arguing), amplified reflection (gently overstating to evoke the other side), shifting focus (moving to a different topic), reframing (offering a new perspective), and emphasizing personal choice ('It's completely up to you what you do with this information').",
        ],
        "Resistance is relational, not characterological - when discord appears, the therapist adjusts rather than pushing harder.",
        "Miller & Rollnick (2013, 2023). MI 3rd & 4th Editions. Guilford Press.",
        notes="The shift from 'resistance' to 'sustain talk and discord' reflects a profound philosophical change. Resistance locates the problem in the client; sustain talk and discord locate it in the interaction. This is not just semantic - it changes what the therapist does. Instead of 'how do I break through their resistance?' the question becomes 'what am I doing that's creating discord?'"
    )



    # =====================================================================
    # SECTION 5: Change Talk and Sustain Talk
    # =====================================================================
    section_divider(prs, 5, "Change Talk and Sustain Talk",
                    "The Language of Ambivalence and Commitment")

    content_slide(prs,
        "Understanding Change Talk: The DARN-CAT Framework",
        [
            "Change talk is any client speech that favors movement toward change. It is the primary therapeutic target in MET because research consistently demonstrates that client change talk predicts actual behavior change outcomes. The DARN-CAT framework categorizes types of change talk by strength.",
            "Preparatory change talk (DARN) reflects movement toward readiness but not yet commitment: Desire ('I want to stop'), Ability ('I could cut back if I tried'), Reasons ('My health would improve'), and Need ('I have to do something'). These statements indicate growing motivation but have not yet crossed into commitment.",
            "Mobilizing change talk (CAT) reflects stronger movement toward action: Commitment ('I will stop drinking this week'), Activation ('I'm ready to try'), and Taking steps ('I poured the bottles out yesterday'). Amrhein et al. (2003) demonstrated that commitment language strength in session predicts actual outcomes.",
            "The clinical implication is clear: therapists should track the proportion and strength of change talk versus sustain talk throughout sessions, adjusting their approach to evoke more change talk while gently acknowledging sustain talk without amplifying it.",
        ],
        "DARN-CAT framework: preparatory change talk (DARN) builds toward mobilizing change talk (CAT) which predicts actual outcomes.",
        "Amrhein et al. (2003). JCCP, 71, 862. Commitment language predicts drug use outcomes.",
        notes="Training exercise: Play recorded sessions and have trainees code each client statement as change talk (specify DARN-CAT type), sustain talk, or neutral. This builds the critical skill of hearing change talk in real-time. Initially, trainees miss much of the change talk because they're focused on content rather than direction."
    )

    content_slide(prs,
        "Evoking Change Talk: Strategic Questions and Reflections",
        [
            "Evoking change talk is the core clinical skill that distinguishes competent MI/MET practice from general supportive counseling. Rather than waiting passively for change talk to appear, the skilled clinician actively creates conversational conditions that make change talk more likely.",
            "Evocative questions target specific DARN categories: Desire ('What would you like to be different?'), Ability ('When have you successfully made changes before?'), Reasons ('What would be the benefits of changing?'), Need ('What concerns you most about continuing as you are?'). Each question type targets a different facet of motivation.",
            "The 'looking forward' technique invites clients to envision life after change: 'If you were to make this change, how would your life be different in a year?' 'What's the best outcome you can imagine?' The 'looking back' technique contrasts with pre-problem times: 'What were things like before drinking became a concern?'",
            "Importance and confidence rulers ('On a scale of 0-10, how important is it to you to make this change?') followed by the strategic follow-up ('Why are you at a 6 and not a 3?') reliably evoke change talk by directing attention to existing motivation rather than deficits.",
        ],
        "Active evocation of change talk through strategic questions is what makes MI/MET more effective than passive supportive listening.",
        "Miller & Rollnick (2023). MI 4th Ed. Chapter on Evoking.",
        notes="The ruler follow-up question is brilliantly designed: 'Why a 6 and not a 3?' forces the client to argue FOR change (explaining why it's not lower). If you asked 'Why not a 10?' they'd argue AGAINST change (explaining barriers). This simple linguistic move reliably produces change talk. Practice it extensively."
    )

    content_slide(prs,
        "Responding to Sustain Talk: Clinical Strategies",
        [
            "Sustain talk - client arguments for maintaining current behavior - is normal and expected, especially early in treatment. It represents one side of ambivalence and should be acknowledged without being amplified or argued against. The goal is not to eliminate sustain talk but to shift the balance toward change talk.",
            "Simple reflection of sustain talk acknowledges without reinforcing: Client says 'I really enjoy drinking with my friends' - Therapist reflects 'The social aspect is important to you' and then redirects: 'And what concerns, if any, have you noticed?' This acknowledges without dwelling.",
            "Coming alongside (sometimes called 'agreeing with a twist') paradoxically reduces sustain talk: 'You're right, nobody can make you change, and maybe the benefits of drinking outweigh any concerns right now. Only you can weigh that up.' This removes the oppositional dynamic and often evokes the client's own concerns.",
            "Research by Magill et al. (2018) in their updated meta-analysis confirmed that the balance of change talk to sustain talk (not just the presence of change talk) is what matters for outcomes. Sessions where sustain talk dominates over change talk predict poorer outcomes, highlighting the importance of skillful responding.",
        ],
        "Sustain talk is normal ambivalence, not pathology. Acknowledge it briefly, then redirect attention toward change talk.",
        "Magill et al. (2018). Updated MI meta-analysis: Technical, relational, and conditional process models.",
        notes="Common trap: engaging with sustain talk as if it were a debate to be won. When a client says 'But I don't think I have a problem,' the natural response is to present evidence of the problem. This creates an argument. Instead: 'You're not seeing much reason for concern at this point. What made you come in today?' This redirects without arguing."
    )

    two_col_slide(prs,
        "Change Talk vs. Sustain Talk: Clinical Examples",
        "Change Talk (DARN-CAT)",
        [
            "D: 'I wish I could just stop'",
            "A: 'I quit smoking, so I know I can do hard things'",
            "R: 'My wife would be so happy'",
            "N: 'I need to do this for my kids'",
            "C: 'I'm going to my first AA meeting Tuesday'",
            "A: 'I'm ready to give this a real try'",
            "T: 'I already called the treatment center'",
        ],
        "Sustain Talk Examples",
        [
            "'I don't think it's that bad'",
            "'I've tried before and failed'",
            "'Everyone in my family drinks'",
            "'I'd lose all my friends'",
            "'It's the only thing that helps my anxiety'",
            "'I'm not ready for this'",
            "'Maybe later, not now'",
        ],
        "Track the ratio of change talk to sustain talk in sessions - this ratio predicts outcomes more than any single technique.",
        "Amrhein et al. (2003). JCCP, 71, 862; Moyers et al. (2007). JCCP, 75, 790.",
        notes="Have trainees practice coding these in real-time. In actual sessions, statements are often mixed or ambiguous. 'I want to stop but I don't think I can' contains both desire (change talk) and low ability (sustain talk). The skilled clinician reflects the change talk component: 'You want things to be different.'"
    )

    content_slide(prs,
        "The Mechanism: How Change Talk Produces Outcomes",
        [
            "Moyers et al. (2007) demonstrated the critical causal chain in MI/MET: Therapist MI-consistent behaviors (OARS, MI spirit) → increased client change talk → improved drinking outcomes. This established that change talk is not merely a correlate of motivation but a mediating mechanism.",
            "The mechanism operates through multiple pathways: Self-perception (Bem, 1972) - hearing oneself argue for change strengthens change identity; Cognitive dissonance (Festinger, 1957) - publicly articulating change intentions creates pressure for consistency; Social commitment - verbal commitments made to another person carry more weight than private decisions.",
            "Amrhein et al. (2003) added important nuance: it is specifically COMMITMENT language (as opposed to desire, ability, or reason statements) that most strongly predicts outcomes. Furthermore, the trajectory of commitment across a session matters - sessions ending with stronger commitment than they began predict better outcomes.",
            "DiClemente et al. (2017) integrated these findings into a comprehensive mechanism model: MI works by activating the client's own change processes through strategic evocation, leading to self-generated arguments for change that are more persuasive and durable than externally imposed reasons.",
        ],
        "The causal chain is established: MI-consistent therapy → change talk → commitment → behavior change outcomes.",
        "Moyers et al. (2007). JCCP, 75, 790; DiClemente et al. (2017). Addiction, 112(S2), 92-100.",
        notes="This is one of the best-supported mechanism models in psychotherapy research. It answers the 'why does MI work?' question empirically. The clinical implication: your job is to increase the frequency, strength, and especially commitment level of change talk. Every technique choice should be evaluated against this criterion."
    )



    # =====================================================================
    # SECTION 6: MET Session Structure
    # =====================================================================
    section_divider(prs, 6, "MET Session Structure",
                    "The Four-Session Protocol from Project MATCH")

    content_slide(prs,
        "MET Session 1: Assessment Feedback Session",
        [
            "The first MET session is the cornerstone of the protocol, consisting of structured personalized feedback from a comprehensive baseline assessment. This session typically lasts 60-90 minutes and follows a specific sequence designed to build motivation through personally relevant data presentation.",
            "The session begins with relationship building (5-10 minutes) using open questions about the client's experience with assessment and their understanding of why they are here. This engaging phase establishes rapport and helps the client feel comfortable before potentially challenging information is presented.",
            "The core of Session 1 is the Personalized Feedback Report (PFR), which typically covers: drinking pattern summary compared to normative data, consequences endorsed, risk factors identified, blood test results (GGT, MCV if available), family history, and the client's own goals and values from the assessment.",
            "Throughout feedback delivery, the therapist uses MI style: asking permission before sharing ('Would it be okay if I shared some of what we found?'), eliciting the client's reaction ('What do you make of this?'), reflecting their responses, and explicitly not telling them what to do. The session ends with a summary and invitation to return.",
        ],
        "Session 1 combines structured feedback with MI style - the data provides the content, MI provides the relational framework.",
        "Miller et al. (1995). MET Manual, Project MATCH. NIAAA. Session 1 protocol.",
        notes="Key training point: therapists often rush through feedback delivery, treating it like a medical report. Each piece of feedback should be delivered, paused, and the client's reaction elicited. 'Your weekly drinking is above 90% of adults your age... What's your reaction to seeing that?' The client's response determines what comes next."
    )

    content_slide(prs,
        "MET Session 2: Strengthening Commitment",
        [
            "Session 2 (typically 1-2 weeks after Session 1) focuses on deepening engagement, exploring ambivalence further, and beginning to strengthen commitment to specific change plans. It opens with an invitation to share reflections since the last session.",
            "The therapist reviews and builds upon themes from Session 1: 'Last time we talked, you mentioned being surprised by where your drinking falls compared to others your age. I'm curious what you've been thinking about since then.' This capitalizes on between-session processing.",
            "If the client has moved toward contemplation or preparation, the therapist can introduce a decisional balance exploration, helping the client articulate both benefits of drinking/status quo and costs/concerns. The therapist strategically reflects and summarizes, with slightly more emphasis on the change side of ambivalence.",
            "Session 2 may also introduce the Change Plan Worksheet if the client demonstrates readiness. This structured tool asks: What changes do I want to make? What are my most important reasons? What steps do I plan to take? How will others help me? What might interfere? What will tell me my plan is working? The therapist guides completion collaboratively.",
        ],
        "Session 2 builds on feedback impact by deepening exploration and, when ready, beginning specific change planning.",
        "Miller et al. (1995). MET Manual. Session 2: Strengthening commitment to change.",
        notes="Critical judgment call: is the client ready for a change plan, or do they need more evocation? Premature planning (pushing a change plan when the client is still ambivalent) creates discord. Signs of readiness include: increased change talk, decreased sustain talk, asking questions about how to change, envisioning future change. If in doubt, stay with evocation."
    )

    process_slide(prs,
        "MET Four-Session Protocol Overview",
        ["Session 1\nFeedback", "Session 2\nCommitment", "Session 3\nReview Progress", "Session 4\nConsolidate"],
        "Four sessions delivered over 12 weeks: Session 1 (Week 1), Session 2 (Week 2), Session 3 (Week 6), Session 4 (Week 12).",
        "Miller et al. (1995). MET Manual, Project MATCH. NIAAA.",
        notes="The spacing is intentional: sessions 1-2 are close together (1-2 weeks) to capitalize on feedback impact. Sessions 3-4 are spaced further apart (weeks 6 and 12) to allow time for change attempts and to address setbacks. The therapist acts as a check-in resource while the client does the work of change between sessions."
    )

    content_slide(prs,
        "MET Sessions 3-4: Progress Review and Consolidation",
        [
            "Sessions 3 and 4 serve to reinforce progress, address obstacles, and consolidate gains. They are typically scheduled at weeks 6 and 12 following the initial sessions. The extended spacing between sessions is deliberate - it communicates confidence in the client's ability to make progress independently.",
            "Session 3 opens with an exploration of what has happened since Session 2: 'Tell me about how things have been going.' The therapist affirms any positive changes (however small), explores what worked, and gently inquires about challenges without catastrophizing setbacks.",
            "If the client has made progress, Sessions 3-4 focus on identifying what is working, strengthening commitment, anticipating future challenges, and developing specific coping strategies. The therapist helps the client attribute success to their own efforts and qualities (supporting self-efficacy).",
            "If the client has not made changes or has experienced setbacks, the therapist avoids criticism or disappointment, instead exploring what happened with curiosity: 'What got in the way?' Setbacks are reframed as learning opportunities, and the change plan may be revised collaboratively. The therapist maintains optimism while respecting the client's autonomy.",
        ],
        "Later sessions reinforce progress, normalize setbacks, and help clients attribute change to their own capabilities.",
        "Miller et al. (1995). MET Manual. Sessions 3-4: Review and consolidation.",
        notes="Key principle for sessions 3-4: the therapist's reaction to setbacks is critically important. If the therapist shows disappointment or frustration, it communicates that change was being done FOR the therapist rather than for the client. Maintain genuine curiosity and equanimity. 'Interesting - what do you think happened?' is always better than 'What went wrong?'"
    )

    table_slide(prs,
        "MET Session Components and Timing",
        ["Session", "Timing", "Primary Focus", "Key Activities", "Duration"],
        [
            ["1", "Week 1", "Structured Feedback", "PFR delivery, reaction elicitation, MI style", "60-90 min"],
            ["2", "Week 2", "Commitment Building", "Review reflections, decisional balance, change plan", "45-60 min"],
            ["3", "Week 6", "Progress Review", "Affirm progress, address obstacles, revise plan", "45-60 min"],
            ["4", "Week 12", "Consolidation", "Reinforce gains, anticipate challenges, closure", "45-60 min"],
        ],
        "MET is remarkably brief (4 sessions over 12 weeks) yet produces outcomes comparable to longer treatments.",
        "Project MATCH (1997). JOSA, 58, 7-29. 4-session MET matched 12-session treatments.",
        notes="Total clinical contact time in MET: approximately 3-4 hours over 12 weeks. Compare to CBT (12 hours) or TSF (12 hours). This efficiency makes MET particularly suitable for resource-limited settings, waiting list management, and as a first-step intervention before more intensive treatment if needed."
    )



    # =====================================================================
    # SECTION 7: Personalized Feedback
    # =====================================================================
    section_divider(prs, 7, "The Personalized Feedback Report",
                    "Assessment-Based Feedback as a Catalyst for Change")

    content_slide(prs,
        "Components of the Personalized Feedback Report (PFR)",
        [
            "The Personalized Feedback Report is the signature element that distinguishes MET from general MI practice. It transforms abstract clinical assessment data into a concrete, personally meaningful document that clients can see, react to, and take home. The visual and comparative nature of feedback is designed to create discrepancy awareness.",
            "Standard PFR components include: (1) Summary of drinking patterns (quantity, frequency, peak BAC estimates), (2) Normative comparison (where the client's consumption falls relative to general population), (3) Negative consequences endorsed on standardized measures, (4) Risk factors identified (family history, early onset, physiological indicators).",
            "Additional components may include: (5) Biological markers (GGT, MCV, liver enzymes if available), (6) Neuropsychological test results if administered, (7) Expenditure calculations (money spent on alcohol annually), (8) Summary of the client's own stated values and goals from the assessment, (9) Readiness to change ruler scores.",
            "The PFR is NOT a diagnostic report delivered to a passive recipient. It is a conversational tool - each section is presented, the client's reaction is elicited, and the therapist reflects and explores before moving to the next section. The client's perspective guides which sections receive more or less attention.",
        ],
        "The PFR transforms clinical data into a personal mirror - its power lies in self-comparison, not therapist judgment.",
        "Miller et al. (1995). MET Manual. Chapter on Personalized Feedback delivery.",
        notes="Practical tips: Print the PFR on quality paper. Use graphs and visuals where possible (pie charts showing normative position, bar graphs of consequences). Give the client a copy to take home. The physical document continues working between sessions as the client reviews it. Some clients report it sitting on their nightstand for weeks, prompting ongoing reflection."
    )

    content_slide(prs,
        "Delivering Feedback: The MI-Consistent Approach",
        [
            "The manner of feedback delivery is as important as the content itself. Miller et al. (1995) emphasized that the same information can motivate change or provoke resistance depending entirely on how it is communicated. MI-consistent feedback delivery follows specific principles.",
            "Always ask permission before sharing: 'I have some information from your assessment that might be interesting. Would it be okay if we looked at this together?' This simple step respects autonomy and sets a collaborative tone. Even mandated clients can be asked: 'You're here because the court required it. Would you like to see what the assessment showed?'",
            "After each feedback element, elicit the client's reaction rather than interpreting for them: 'What do you make of that?' 'Does that surprise you at all?' 'How does that fit with what you expected?' The client's own interpretation is more powerful than the therapist's explanation.",
            "Avoid the 'expert trap': presenting feedback as evidence of a diagnosis the client must accept. Instead, offer information neutrally: 'Your drinking places you above 95% of adults your age. That's what the data shows. I'm curious what that means to you.' This allows the client to draw their own conclusions.",
        ],
        "Feedback is OFFERED for consideration, not IMPOSED as a verdict. The client's interpretation drives the therapeutic process.",
        "Miller et al. (1995). MET Manual. Feedback delivery guidelines.",
        notes="Common mistake: therapists who prepare elaborate feedback reports but then deliver them like a medical consultation - running through all sections rapidly without pausing for client reaction. The REACTION is the therapy, not the information. A session where you cover only 2 of 8 feedback sections because the client deeply engaged with those 2 is more successful than rushing through all 8."
    )

    content_slide(prs,
        "Normative Feedback: The Power of Social Comparison",
        [
            "Normative feedback - showing clients where their consumption falls relative to the general population - is one of the most reliably powerful components of the PFR. Research on personalized normative feedback (PNF) consistently shows that correcting overestimations of drinking norms produces behavior change even as a standalone intervention.",
            "The mechanism operates through social comparison: most heavy drinkers significantly overestimate how much 'everyone else' drinks, believing their behavior is normal. When shown that they drink more than 90-95% of adults their age and gender, this perception is powerfully disrupted. 'I had no idea I was that far out' is a common reaction.",
            "In the Indian context, normative data from NIMHANS (2008) and Benegal (2005) can be used. India has approximately 62.5 million alcohol users. While prevalence is lower than Western countries, consumption among drinkers is often heavier, and 'hazardous drinking' rates among male drinkers are substantial.",
            "D'Costa et al. (2019) provided prevalence data from Goa, India, showing significant rates of alcohol use disorders that can serve as local normative comparison data. Using country-specific or regional data enhances the personal relevance and credibility of normative feedback.",
        ],
        "Normative comparison is powerful because most heavy drinkers vastly overestimate population drinking - correction creates discrepancy.",
        "Benegal, V. (2005). India: alcohol and public health. Addiction, 100, 1051-1056.",
        notes="Adaptation consideration: In some Indian communities, ANY alcohol use may be stigmatized (e.g., in certain religious communities). Normative feedback needs cultural calibration - comparing to 'Indian males in your age group' rather than global norms. Also consider that in some social circles (e.g., corporate culture), heavy drinking IS the norm - here, health consequence feedback may be more impactful than normative comparison."
    )

    two_col_slide(prs,
        "Feedback Elements: What to Include vs. What to Avoid",
        "Effective Feedback Elements",
        [
            "Objective data (quantities, frequencies)",
            "Visual comparisons (graphs, charts)",
            "Normative positioning (percentile)",
            "Personal consequences endorsed",
            "Client's own words about values/goals",
            "Money spent calculations",
            "Health markers (GGT, MCV) if available",
            "Family history context (non-judgmental)",
        ],
        "Pitfalls to Avoid in Feedback",
        [
            "Diagnostic labels ('You are an alcoholic')",
            "Prescriptive conclusions ('You must stop')",
            "Catastrophizing ('This will kill you')",
            "Comparing to extreme cases",
            "Overwhelming with too much data",
            "Ignoring client's emotional reactions",
            "Rushing through without pausing",
            "Arguing with minimization responses",
        ],
        "Feedback should inform and evoke, never diagnose or prescribe. The client determines what the data means for them.",
        "Miller et al. (1995). MET Manual. Feedback delivery principles.",
        notes="Role-play exercise: Give trainees a sample PFR and have them practice delivery in pairs, with 'clients' coached to react in various ways (surprise, minimization, anger, tears). Observe whether the therapist can maintain MI style under pressure. Common failure point: when the client minimizes ('That doesn't seem that bad'), the therapist feels compelled to argue for severity."
    )



    # =====================================================================
    # SECTION 8: Evidence Base - International Research
    # =====================================================================
    section_divider(prs, 8, "Evidence Base: International Research",
                    "Meta-Analyses, Cochrane Reviews, and RCTs")

    content_slide(prs,
        "Schwenker et al. (2023): Cochrane Review of MI for Substance Use",
        [
            "The most comprehensive and recent meta-analytic evidence for MI/MET comes from Schwenker et al. (2023), a Cochrane Systematic Review (CD008063.pub3) examining MI for substance use reduction. Cochrane Reviews represent the gold standard in evidence synthesis due to their rigorous methodology.",
            "The review found a small-to-medium effect size (SMD = 0.48) for MI compared to no-treatment controls in reducing substance use. This effect, while modest in magnitude, is clinically meaningful given the brief nature of MI interventions (typically 1-4 sessions) compared to the chronic, relapsing nature of substance use disorders.",
            "Important moderators identified: MI was more effective when delivered by trained therapists with fidelity monitoring than when delivered by practitioners without supervision. The effect was robust across substances (alcohol, cannabis, tobacco, other drugs) and settings (primary care, specialty treatment, emergency departments).",
            "The review also noted that MI showed advantages over other active treatments in some comparisons, suggesting the effect is not merely attributable to therapeutic contact time. MI's specific ingredients (evocation of change talk, empathic reflection, autonomy support) appear to contribute unique therapeutic value beyond common factors.",
        ],
        "Cochrane Review (2023): MI shows SMD 0.48 for substance reduction - modest but clinically meaningful for a brief intervention.",
        "Schwenker et al. (2023). Cochrane Review: MI for substance use. CD008063.pub3. SMD 0.48.",
        notes="Context for interpreting SMD 0.48: This is a small-medium effect by Cohen's standards. But consider: this comes from typically 1-4 sessions of MI versus extensive control conditions. The effect per session hour is actually quite large. Also, this is averaged across diverse populations, settings, and substances - in specific well-delivered applications, effects may be larger."
    )

    content_slide(prs,
        "UKATT Trial: UK Evidence for Brief Intervention",
        [
            "The United Kingdom Alcohol Treatment Trial (UKATT, 2005) was designed as a pragmatic effectiveness trial comparing 3-session MET (Motivational Enhancement Therapy adapted to UK context) against 8-session Social Behaviour and Network Therapy (SBNT) for alcohol-dependent adults. N=742 participants were randomized across multiple UK treatment sites.",
            "Results paralleled Project MATCH: both treatments produced substantial reductions in drinking (approximately 40% reduction in alcohol consumption at 12-month follow-up) with no significant difference between them. Again, the briefer motivational approach matched a longer treatment in outcomes.",
            "UKATT additionally conducted cost-effectiveness analyses showing that MET was significantly more cost-effective than SBNT due to equivalent outcomes with fewer sessions. At 12-month follow-up, MET produced savings of approximately 1000 British Pounds per client while achieving equivalent clinical outcomes.",
            "The UKATT findings were particularly important because they demonstrated that Project MATCH results were not specific to the US healthcare context. The robustness of MI/MET across cultural settings (USA and UK) encouraged subsequent adaptation efforts in LMICs including India.",
        ],
        "UKATT (2005): 3-session MET matched 8-session therapy in the UK, confirming cross-cultural robustness and cost-effectiveness.",
        "UKATT Research Team (2005). BMJ, 331, 541. N=742 across UK sites.",
        notes="UKATT's practical importance: it moved MI/MET from research efficacy into real-world effectiveness. Unlike Project MATCH (which used highly selected therapists with intensive supervision), UKATT used routine NHS therapists with standard training. That the effects held confirms MI/MET is transportable to routine clinical practice."
    )

    table_slide(prs,
        "Summary of Major MET/MI Clinical Trials",
        ["Trial", "Year", "N", "Comparison", "Key Finding"],
        [
            ["Project MATCH", "1997", "1726", "MET vs CBT vs TSF", "4-session MET = 12-session treatments"],
            ["UKATT", "2005", "742", "MET vs SBNT", "3 sessions = 8 sessions; MET more cost-effective"],
            ["COMBINE", "2006", "1383", "MI + naltrexone", "MI enhanced medication adherence"],
            ["AMBIT (India)", "2023", "280", "Phone MI vs EUC", "Mobile MI feasible with lay counselors"],
            ["Cochrane Review", "2023", "Multiple", "MI vs controls", "SMD 0.48 for substance reduction"],
            ["Patel (India)", "2024", "N/A", "MI+BCT vs control", "Combined MI approach effective for IPV+alcohol"],
        ],
        "Consistent evidence across decades and continents: MET/MI produces meaningful change efficiently.",
        "Multiple sources - see individual citations for each trial.",
        notes="This table provides an evidence overview for clinicians and administrators who need to justify implementing MI/MET. The consistency across trials, settings, populations, and decades is remarkable. No other brief intervention for substance use has this level of evidence support."
    )

    content_slide(prs,
        "MI in Low- and Middle-Income Countries: Systematic Review",
        [
            "A systematic review published in Medrxiv (2023) specifically examined the effectiveness of MI when adapted for low- and middle-income countries (LMICs), where the majority of the global disease burden from substance use occurs but resources for treatment are most limited.",
            "Of 11 studies meeting inclusion criteria, 7 (64%) demonstrated improved outcomes when MI was adapted for LMIC contexts. This success rate is notable given the substantial cultural, linguistic, and resource differences between LMICs and the high-income settings where MI was developed.",
            "Successful adaptations shared key features: use of local languages, training of non-specialist health workers to deliver MI, integration with existing healthcare systems, cultural adaptation of values exploration, and abbreviated session formats to accommodate clients with limited time for treatment.",
            "India figured prominently in this literature, with the AMBIT trial (Nadkarni et al., 2023) demonstrating successful mobile-phone delivery by lay counselors in Goa, and multiple NIMHANS-related projects showing feasibility in Indian clinical settings. D'Costa et al. (2019) provided the epidemiological foundation for such work in Goa.",
        ],
        "MI works in LMICs: 7/11 studies showed improved outcomes, with successful task-shifting to non-specialist workers.",
        "Medrxiv (2023). Systematic Review: MI in LMICs - 7/11 studies showed improved outcomes.",
        notes="This is crucial evidence for global health: MI/MET doesn't require doctoral-level psychologists to deliver. With proper training and supervision, community health workers, nurses, and lay counselors can deliver MI effectively. This makes it one of the most scalable evidence-based treatments available for substance use in resource-limited settings."
    )

    content_slide(prs,
        "Recent Expansion: MET for PTSD, Tobacco, and Opioids",
        [
            "Recent research has expanded MET's evidence base beyond alcohol to diverse clinical applications. Hussey et al. (2023) demonstrated MET's effectiveness for improving treatment engagement among military veterans with PTSD, where motivation for trauma-focused treatment is often ambivalent.",
            "PMC (2024) published a comprehensive review of MET for tobacco cessation in schizophrenia, where the challenge of motivation is compounded by cognitive symptoms, negative symptoms, and the rewarding properties of nicotine for dopamine dysregulation. MET's non-confrontational approach proved particularly suitable for this population.",
            "For opioid use disorders, Olgac et al. (2024, Case Western Reserve) reviewed MI and MET applications, finding evidence for MET as both a standalone brief intervention and as an adjunct to medication-assisted treatment (MAT) with methadone or buprenorphine. PMC (2024) further described a Group MET protocol for opioid dependence on maintenance treatment.",
            "These expansions demonstrate MET's transdiagnostic potential: wherever motivation for behavior change is uncertain or ambivalent, MET's structured feedback approach combined with MI communication style can enhance engagement and outcomes. The principle transfers across conditions.",
        ],
        "MET's principles transfer across conditions: PTSD engagement, tobacco cessation, opioid treatment - wherever motivation matters.",
        "Olgac et al. (2024). MI/MET for alcohol and opioid use disorders. Case Western Reserve University.",
        notes="The transdiagnostic principle: MET works not because it treats a specific disorder but because it addresses a universal human challenge - the ambivalence about changing established behavior patterns. This means any clinician working with behavior change can benefit from MET training, regardless of their specialty."
    )



    # =====================================================================
    # SECTION 9: Indian Context Research
    # =====================================================================
    section_divider(prs, 9, "MET/MI in the Indian Context",
                    "Adaptations, Trials, and Cultural Considerations")

    content_slide(prs,
        "The AMBIT Trial: Mobile MI in Goa, India",
        [
            "Nadkarni et al. (2023) published the AMBIT (Alcohol use and Motivation: Brief Intervention using Technology) trial in BMC Psychiatry, representing one of the most rigorous evaluations of MI in an Indian setting. The trial tested whether MI delivered via mobile phone by community health workers could reduce harmful alcohol use.",
            "The innovation was twofold: (1) using mobile phones to deliver MI, overcoming geographic barriers in a population where treatment-seeking is low, and (2) using trained lay counselors (community health workers) rather than mental health professionals, addressing the severe workforce shortage in India.",
            "Participants were adult male hazardous drinkers identified through community screening in Goa. The intervention consisted of brief MI-based counseling sessions delivered via mobile phone, compared to enhanced usual care. The counselors received structured MI training with ongoing supervision.",
            "The trial demonstrated feasibility and acceptability of this approach, with important lessons for scalability: mobile delivery overcame stigma barriers (clients didn't need to visit a 'psychiatric' clinic), scheduling was flexible around work hours, and lay counselors maintained MI fidelity with supervision. This model has potential for massive scale-up across India.",
        ],
        "AMBIT showed mobile-phone MI by lay counselors is feasible and acceptable in India - a scalable model for 62.5 million users.",
        "Nadkarni et al. (2023). AMBIT Trial - mobile phone MI for alcohol use in Goa, India. BMC Psychiatry.",
        notes="Context: India has approximately 62.5 million alcohol users (NIMHANS, 2008) but fewer than 5,000 psychiatrists and virtually no addiction specialists outside major cities. Task-shifting MI to lay counselors via mobile phones could theoretically reach millions who currently have no access to evidence-based treatment. The AMBIT model is a prototype for this scaling."
    )

    content_slide(prs,
        "Combined MI + BCT for Alcohol and IPV in South India",
        [
            "Patel et al. (2024) published in PLOS ONE a study combining Motivational Interviewing with Behavioral Couples Therapy (BCT) for men who had both alcohol use problems and perpetrated intimate partner violence (IPV) in south India. This addresses a critical intersection of public health challenges.",
            "The rationale for combination was that alcohol use and IPV are strongly associated in Indian settings, with alcohol involved in 25-50% of IPV incidents. Treating either alone without addressing the other leads to limited sustained improvement. MI provided the motivational foundation while BCT addressed relationship skills and communication.",
            "The MI component focused on building motivation for both reducing drinking and reducing violent behavior, exploring how both behaviors conflicted with participants' values as partners and fathers. Cultural values around family respect and provider role were incorporated into discrepancy development.",
            "This study is significant because it demonstrates that MI can be combined with other evidence-based approaches in Indian settings, and that cultural values (family honor, respect for spouse, role as provider) can be effectively incorporated into MI-style exploration without becoming confrontational or prescriptive.",
        ],
        "Combined MI+BCT addresses the alcohol-IPV intersection in India, using cultural values as motivational resources.",
        "Patel et al. (2024). Combined MI + BCT for IPV and alcohol in south India. PLOS ONE.",
        notes="Clinical relevance: In Indian clinical practice, alcohol and domestic violence frequently co-occur. This study provides a model for addressing both simultaneously using MI principles. The key insight is that cultural values (family respect, dharma as a husband/father) can be leveraged FOR change rather than serving as barriers to engagement."
    )

    content_slide(prs,
        "NIMHANS Psychosocial Interventions Framework",
        [
            "The National Institute of Mental Health and Neurosciences (NIMHANS) in Bangalore published the Psychosocial Interventions Manual (Murthy, P., 2008) which established the framework for evidence-based psychosocial treatments in India, including motivational approaches. This manual is the reference standard for Indian clinical practice.",
            "Key statistics from NIMHANS (2008): India has approximately 62.5 million alcohol users, with significant variation across states (highest in northeastern states, lowest in Gujarat/Rajasthan due to prohibition/religious factors). Among drinkers, a substantial proportion engage in hazardous patterns.",
            "NIMHANS recommended motivational interviewing as a first-line psychosocial intervention for alcohol use disorders, noting its cultural compatibility with Indian values of respect, non-confrontation, and collaborative relationship styles. The manual provides Indian-adapted protocols for clinical use.",
            "Subsequent NIMHANS research has focused on training non-specialist health workers in MI techniques, integrating MI into primary healthcare settings, and developing culturally appropriate assessment tools that can feed into MET-style personalized feedback. These efforts address India's massive treatment gap.",
        ],
        "NIMHANS established MI as a recommended first-line psychosocial intervention for alcohol use in India - culturally compatible and scalable.",
        "NIMHANS (2008). Psychosocial Interventions Manual. Murthy, P. India: 62.5M alcohol users.",
        notes="NIMHANS's endorsement is crucial for Indian clinical practice because it provides institutional authority for MI/MET approaches. Clinicians working in government hospitals and medical colleges can cite NIMHANS guidelines when proposing MI training or service development. The cultural compatibility argument is important: MI's non-confrontational approach aligns with Indian therapeutic traditions that emphasize respect and collaboration."
    )

    content_slide(prs,
        "Alcohol Epidemiology in India: D'Costa et al. (2019) and Benegal (2005)",
        [
            "Understanding the epidemiology of alcohol use in India is essential for contextualizing MET delivery. D'Costa et al. (2019) published data on alcohol use disorders in Goa, India in Drug and Alcohol Review, providing prevalence estimates that can inform personalized feedback and normative comparisons.",
            "Benegal (2005) in Addiction journal provided a comprehensive overview of alcohol and public health in India, noting patterns distinct from Western countries: lower overall prevalence but higher per-drinking-day consumption, primarily male drinkers (female drinking stigmatized in most communities), and strong associations between alcohol and domestic violence.",
            "These epidemiological patterns have implications for MET adaptation: (1) Normative feedback must use Indian/regional comparison data, not Western norms; (2) Gender-specific approaches are needed given that male and female drinking patterns differ dramatically; (3) Family consequences (especially IPV and economic hardship) may be more motivationally salient than health consequences in Indian settings.",
            "Both studies emphasize the treatment gap: despite significant alcohol-related harm, the vast majority of affected individuals never receive any formal treatment. Brief interventions like MET that can be delivered in primary care or by non-specialists represent the most realistic pathway to narrowing this gap.",
        ],
        "Indian alcohol patterns differ from Western norms - MET feedback and motivational exploration must be culturally calibrated.",
        "D'Costa et al. (2019). Alcohol use disorders in Goa, India. Drug & Alcohol Review; Benegal (2005). Addiction, 100, 1051-1056.",
        notes="Practical adaptation: When creating PFRs for Indian clients, use Indian normative data where available. Consider that motivators may differ: in Western settings, health and personal freedom are common change motivators; in Indian settings, family responsibility, economic impact, social respect, and religious/spiritual values may be more salient. Always explore the individual client's values rather than assuming."
    )



    # =====================================================================
    # SECTION 10: Mechanisms of Change
    # =====================================================================
    section_divider(prs, 10, "Mechanisms of Change in MET",
                    "How and Why MET Produces Behavior Change")

    content_slide(prs,
        "DiClemente et al. (2017): Understanding MI Mechanisms",
        [
            "DiClemente et al. (2017) published a comprehensive analysis of mechanisms of change in MI in Addiction journal (Supplement 2, pages 92-100). This paper synthesized decades of process research to identify how MI produces its effects, moving beyond the question of whether it works to understanding why.",
            "Four primary mechanisms were identified: (1) Therapist MI-consistent behaviors that create the relational conditions for change exploration; (2) Client change talk that reflects growing motivation and commitment; (3) Perceived discrepancy between current behavior and personal values; (4) Self-efficacy enhancement that builds confidence for change.",
            "The mechanism model is sequential: MI-consistent therapist behavior → therapeutic relationship (engagement) → change talk evocation → discrepancy awareness → self-efficacy building → commitment → behavior change. Each step depends on the preceding one, explaining why technique without spirit fails.",
            "This mechanism understanding has practical implications for training and supervision: rather than teaching isolated techniques, training should focus on the entire chain from spirit/relationship through to behavior change. Supervisors can identify where in the chain a trainee is struggling and target intervention accordingly.",
        ],
        "MI works through a causal chain: therapist behavior → relationship → change talk → discrepancy → self-efficacy → commitment → change.",
        "DiClemente et al. (2017). Mechanisms of change in MI. Addiction, 112(S2), 92-100.",
        notes="This paper is essential reading for advanced MI/MET practitioners. Understanding mechanisms allows you to troubleshoot when therapy isn't working: Is the relationship established? Is change talk being evoked? Is discrepancy salient? Is self-efficacy being built? Each 'no' points to a different clinical intervention."
    )

    content_slide(prs,
        "The Causal Chain: Therapist → Client → Outcomes",
        [
            "Moyers et al. (2007) provided the first empirical demonstration of the MI causal chain using sequential analysis of therapy sessions. They showed that MI-consistent therapist behaviors (open questions, complex reflections, affirmations) were followed by increased client change talk within sessions, and that client change talk predicted subsequent drinking outcomes.",
            "The reverse was also demonstrated: MI-inconsistent therapist behaviors (confrontation, warning, directing without permission) were followed by increased sustain talk and decreased change talk, which predicted poorer outcomes. This confirmed that therapist behavior directly influences the change talk mechanism.",
            "This finding resolved a debate in the field: was change talk merely a marker of pre-existing motivation (which would reduce the therapist's role to assessment), or was it genuinely influenced by therapist behavior (which would make it a therapeutic target)? The answer was clearly the latter.",
            "Clinical implication: every therapist utterance either moves the client toward or away from change talk. There is no 'neutral' interaction. Questions, reflections, and statements all have directional consequences. Skilled MI/MET practice involves conscious moment-to-moment choices about which direction to guide the conversation.",
        ],
        "Every therapist utterance has directional consequences - there is no neutral interaction in MI/MET.",
        "Moyers et al. (2007). JCCP, 75, 790. Sequential analysis demonstrating the MI causal chain.",
        notes="Supervision application: record sessions, code therapist behaviors as MI-consistent or MI-inconsistent, then examine what follows each behavior. Trainees are often surprised to see that their well-intentioned advice-giving is consistently followed by client sustain talk, while their reflective listening is followed by deeper change talk."
    )

    process_slide(prs,
        "The MI Mechanism Pathway",
        ["MI Spirit &\nSkills", "Therapeutic\nRelationship", "Change Talk\nEvocation", "Discrepancy\nAwareness", "Behavior\nChange"],
        "The mechanism pathway shows MI works through relationship and evocation, not through information or persuasion alone.",
        "DiClemente et al. (2017). Addiction, 112(S2), 92-100.",
        notes="This visual helps clinicians understand that jumping straight to 'giving information about risks' (common in medical settings) misses the relational and evocative mechanisms that actually produce change. Information alone rarely changes behavior; information delivered within an MI relationship and evocative conversation can be transformative."
    )

    content_slide(prs,
        "Commitment Language: The Strongest Predictor",
        [
            "Amrhein et al. (2003) published a landmark study in the Journal of Consulting and Clinical Psychology (71, 862) demonstrating that the strength of commitment language within MI sessions was the strongest predictor of post-treatment drug use outcomes. This was more predictive than other types of change talk.",
            "Specifically, they found that: (1) Strength of commitment statements increased over the course of sessions; (2) End-of-session commitment strength predicted outcomes better than beginning-of-session; (3) The trajectory of commitment (increasing vs. flat) mattered more than absolute level; (4) Desire, ability, and reason statements alone did not predict outcomes as strongly.",
            "This finding has profound clinical implications: the therapist's goal is not merely to evoke any change talk, but specifically to guide the conversation toward increasingly strong commitment statements. Preparatory change talk (DARN) creates conditions for commitment, but commitment itself is the proximal predictor of behavior change.",
            "In MET specifically, the Change Plan Worksheet at the end of Session 2 is designed to crystallize commitment into specific written commitments: 'I will... by [date]... with support from...' This transforms diffuse motivation into concrete commitment that has been publicly stated to another person.",
        ],
        "Commitment language is the strongest in-session predictor of behavior change - stronger than desire, ability, or reason talk.",
        "Amrhein et al. (2003). JCCP, 71, 862. Commitment language predicts drug outcomes.",
        notes="Practical application: Track the strength of commitment in your sessions. Early in treatment, you may hear 'I might try' or 'I'll think about it' (weak commitment). Your goal is to help this evolve toward 'I will' and 'I'm going to' (strong commitment). Don't push - let it emerge naturally through the MI process. Premature demands for commitment backfire."
    )

    content_slide(prs,
        "What Does NOT Work: The Anti-Mechanism",
        [
            "Understanding what does NOT produce change is as important as understanding what does. Patterson and Forgatch (1985) experimentally demonstrated that when therapists were instructed to use confrontation ('teach' and 'confront' behaviors), client resistance increased significantly compared to when they used supportive behaviors.",
            "This finding has been replicated across multiple studies: therapist behaviors that are MI-inconsistent (advising without permission, confronting, warning, labeling, moralizing) consistently produce increases in sustain talk, client defensiveness, and poorer outcomes. The mechanism works in reverse: anti-MI behavior produces anti-change talk.",
            "Valle (1981) showed the flip side: counselor empathy (rated by observers, not self-reported) was the strongest predictor of client drinking outcomes at 2-year follow-up. Counselors in the top third of empathy had clients with significantly better outcomes than counselors in the bottom third. Empathy is not a luxury - it is a therapeutic necessity.",
            "These findings together establish that the therapist's relational stance is not merely a precondition for technique delivery but is itself a primary mechanism. MI/MET works not because of specific techniques deployed within an empathic relationship, but because the empathic relationship IS the treatment, with techniques serving as vehicles for its expression.",
        ],
        "Confrontation increases resistance; empathy predicts outcomes. The relationship IS the treatment, not just its context.",
        "Patterson & Forgatch (1985). JCCP, 53, 846; Valle (1981). J Studies Alcohol, 42, 783.",
        notes="This is perhaps the single most important finding for clinical training: what you DON'T do matters as much as what you DO. Eliminating MI-inconsistent behaviors (confrontation, unsolicited advice, lecturing) may be more impactful than adding MI-consistent behaviors. Many trainees focus on learning new skills while maintaining old confrontational habits."
    )



    # =====================================================================
    # SECTION 11: Assessment in MET
    # =====================================================================
    section_divider(prs, 11, "Assessment in MET",
                    "Comprehensive Evaluation for Personalized Feedback")

    content_slide(prs,
        "The Role of Assessment in MET: More Than Diagnosis",
        [
            "In MET, assessment serves a dual purpose that distinguishes it from standard clinical evaluation. First, it provides the data necessary for the Personalized Feedback Report. Second, and often underappreciated, the assessment process itself can be therapeutic - it communicates that the therapist is taking the client seriously and gathering comprehensive information.",
            "Standard MET assessment batteries include: quantity/frequency measures (Timeline Followback, Drinking Days questionnaire), consequences inventories (Drinker Inventory of Consequences), dependence severity (AUDIT, SADQ), readiness to change (URICA, readiness rulers), self-efficacy (AASE), values exploration, and biological markers where available.",
            "The assessment should be administered BEFORE Session 1, typically in a separate assessment session or through self-report questionnaires. This allows the therapist time to prepare the Personalized Feedback Report and organize the feedback in a coherent, impactful sequence.",
            "Cultural adaptation of assessment is critical in Indian settings. NIMHANS (2008) recommends using instruments validated in Indian populations, translating measures into local languages, and including culturally relevant consequence domains (family dishonor, workplace consequences specific to Indian employment contexts, financial impact on family).",
        ],
        "Assessment in MET serves both clinical evaluation AND therapeutic purposes - it generates the raw material for personalized feedback.",
        "Miller et al. (1995). MET Manual. Assessment and feedback preparation guidelines.",
        notes="Practical tip: Never waste good assessment data. Every piece of information gathered should potentially appear in the PFR. When administering assessment, note the client's reactions - these can inform how to sequence feedback delivery. A client who tears up during the family impact questions has signaled that family consequences are motivationally salient."
    )

    table_slide(prs,
        "Key Assessment Instruments for MET",
        ["Instrument", "What It Measures", "Feedback Use", "Time"],
        [
            ["AUDIT", "Alcohol use severity", "Normative comparison, risk level", "5 min"],
            ["Timeline Followback", "Daily drinking pattern", "Frequency/quantity graphs", "20 min"],
            ["DrInC", "Consequences (50 items)", "Consequence domains affected", "15 min"],
            ["URICA/Rulers", "Readiness to change", "Motivational staging", "5 min"],
            ["AASE", "Self-efficacy", "Confidence in specific situations", "10 min"],
            ["Values Sort", "Personal values hierarchy", "Discrepancy identification", "15 min"],
            ["GGT/MCV", "Biological markers", "Physical health impact", "Lab"],
        ],
        "A comprehensive assessment battery takes 60-90 minutes but generates the material for a powerful personalized feedback session.",
        "Miller et al. (1995). MET Manual. Assessment battery recommendations.",
        notes="Not all instruments are essential for every client. At minimum: AUDIT (severity), some drinking pattern measure, consequences inventory, and values exploration. Biological markers add impact when available but are not always feasible. In primary care settings, even AUDIT + brief drinking diary + values discussion can generate meaningful personalized feedback."
    )

    content_slide(prs,
        "Values Exploration: The Heart of Discrepancy Development",
        [
            "The values exploration component of MET assessment is uniquely important because it provides the 'other side' of the discrepancy equation. Assessment of drinking patterns and consequences shows where the client IS; values exploration shows where they WANT TO BE. The gap between them generates motivation.",
            "Values exploration can be conducted through card sorts (having clients rank pre-printed value cards), open-ended questions ('What matters most to you in life?'), or structured questionnaires. The MET manual recommends identifying the client's top 5-10 values and explicitly connecting these to the feedback discussion.",
            "In Indian cultural contexts, commonly endorsed values include: family harmony and duty, respect from community, being a good provider, religious/spiritual practice, education of children, personal dignity (izzat), and health for longevity with family. These culturally specific values are powerful motivational resources.",
            "The therapist weaves value statements throughout the feedback session: 'You mentioned that being a respected community member is one of your highest values. Let me show you where your drinking falls relative to the community...' This creates personally meaningful discrepancy without the therapist ever saying 'you should change.'",
        ],
        "Values exploration creates the motivational target - discrepancy emerges naturally when values are compared with behavior patterns.",
        "Miller et al. (1995). MET Manual. Values assessment and feedback integration.",
        notes="Training exercise: Have trainees practice the values exploration conversation. Many clinicians have never been trained to discuss values - they focus on symptoms, problems, and diagnoses. The values conversation often feels different - warmer, more personal, more connecting. It builds the relationship while gathering essential motivational data."
    )



    # =====================================================================
    # SECTION 12: Decisional Balance
    # =====================================================================
    section_divider(prs, 12, "Decisional Balance in MET",
                    "Exploring Ambivalence Systematically")

    content_slide(prs,
        "The Decisional Balance: Theory and Practice",
        [
            "The decisional balance is a structured exploration of the perceived benefits and costs of both the current behavior (drinking) and the proposed change (reducing/stopping). Rooted in Janis and Mann's (1977) decision theory, it acknowledges that problematic behaviors serve functions - they provide something valued.",
            "The four cells of the decisional balance matrix are: (1) Benefits of drinking (what does it give you?), (2) Costs of drinking (what are the downsides?), (3) Benefits of changing (what would improve?), and (4) Costs of changing (what would you lose?). All four must be explored to honor the client's full experience.",
            "A common clinical error is exploring only the costs of drinking and benefits of changing, essentially building a one-sided case for change. This violates MI spirit by positioning the therapist as an advocate for change rather than a collaborative explorer of ambivalence. Paradoxically, exploring the BENEFITS of drinking often evokes change talk: 'Well, it relaxes me, but honestly there are better ways to relax.'",
            "The timing of decisional balance in MET matters: it is most useful in Session 2 for clients who remain ambivalent after feedback. For clients who received feedback and are already moving toward commitment, jumping to a change plan may be more appropriate. Assessment of readiness guides this clinical decision.",
        ],
        "Explore ALL four cells of the decisional balance - including benefits of drinking - to honor ambivalence and evoke autonomous change talk.",
        "Janis, I.L. & Mann, L. (1977). Decision Making. Free Press.",
        notes="Counterintuitive but important: spending time on 'what do you like about drinking?' can actually increase motivation to change. When clients feel their positive experiences are acknowledged and respected, they become more willing to examine the negative side. It also builds the therapeutic relationship by demonstrating genuine curiosity rather than a hidden agenda."
    )

    two_col_slide(prs,
        "Decisional Balance: Four Quadrant Exploration",
        "CURRENT BEHAVIOR (Drinking)",
        [
            "BENEFITS: Relaxation, social connection",
            "BENEFITS: Stress relief, sleep aid",
            "BENEFITS: Confidence in social situations",
            "BENEFITS: Reward after hard work",
            "---",
            "COSTS: Hangovers, health problems",
            "COSTS: Arguments with spouse",
            "COSTS: Money spent, missed work",
            "COSTS: Guilt, self-disappointment",
        ],
        "CHANGE (Reducing/Stopping)",
        [
            "BENEFITS: Better health, more energy",
            "BENEFITS: Improved family relationships",
            "BENEFITS: More money available",
            "BENEFITS: Self-respect, sense of control",
            "---",
            "COSTS: Loss of social activity",
            "COSTS: Need to find new coping",
            "COSTS: Boredom, missing the ritual",
            "COSTS: Social pressure from peers",
        ],
        "All four quadrants are valid. The client weighs them - the therapist facilitates exploration, not decision-making.",
        "Miller et al. (1995). MET Manual. Decisional balance exercise.",
        notes="In practice, fill this in collaboratively with the client. Start with benefits of drinking (less threatening), then costs of drinking, then benefits of changing, then costs of changing. After all four cells are completed, summarize the whole picture and ask: 'Looking at all of this together, what stands out to you?' Let them draw conclusions."
    )

    # =====================================================================
    # SECTION 13: The Change Plan
    # =====================================================================
    section_divider(prs, 13, "The Change Plan Worksheet",
                    "Translating Motivation into Specific Action")

    content_slide(prs,
        "The Change Plan Worksheet: Structure and Purpose",
        [
            "The Change Plan Worksheet is a structured document used in MET Session 2 (or later) to help clients translate their developing motivation into specific, concrete plans for action. It serves as both a planning tool and a written commitment that can be reviewed in subsequent sessions.",
            "The standard Change Plan Worksheet includes these sections: (1) The changes I want to make are... (2) The most important reasons why I want to make these changes are... (3) The steps I plan to take in changing are... (4) The ways other people can help me are... (5) I will know that my plan is working if... (6) Some things that could interfere with my plan are...",
            "The worksheet is completed collaboratively - the therapist guides and reflects, but the client writes (or dictates) using their own words. This is critical because self-generated commitment is stronger than therapist-generated prescriptions. The therapist may ask clarifying questions but should not edit the client's language.",
            "Timing is critical: the Change Plan should only be introduced when the client shows clear signs of readiness (predominant change talk, asking 'how' questions, envisioning change). Introducing it prematurely to an ambivalent client creates pressure that undermines the MI process and may generate discord.",
        ],
        "The Change Plan crystallizes motivation into written commitment - but only when the client is ready. Timing is everything.",
        "Miller et al. (1995). MET Manual. Change Plan Worksheet protocol.",
        notes="Signs that a client is ready for the Change Plan: increased change talk, decreased sustain talk, questions about how to change, imagining life after change, taking small steps between sessions, decreased resistance. If in doubt, ask a readiness ruler: 'On a scale of 0-10, how ready are you to make a specific plan today?' Below 7, stay with evocation."
    )

    content_slide(prs,
        "Collaborative Goal Setting in MET",
        [
            "Goal setting in MET respects client autonomy: the therapist does not prescribe abstinence or any specific goal, but helps the client clarify their own desired outcome. This is one of the most philosophically challenging aspects for clinicians trained in disease model approaches.",
            "Research supports goal choice: clients who choose their own goals (including moderation goals) show better outcomes than those assigned goals by clinicians. Even when moderation may not be sustainable long-term, allowing the client to discover this through their own experience preserves the therapeutic relationship and internal motivation.",
            "The therapist's role is to ensure the goal is specific, measurable, and realistic while reflecting any ambivalence about goal choice: 'So your goal is to cut back to weekends only, with no more than 3 drinks. How confident are you in that plan? What might make it difficult?' This explores without prescribing.",
            "For some clients, particularly those with severe dependence, the therapist may share clinical information about risk (with permission) while maintaining autonomy: 'From a medical standpoint, with your level of physical dependence, stopping abruptly could be dangerous. Can I share some information about safer options?' This informs without demanding.",
        ],
        "MET supports client-chosen goals - even imperfect ones. Autonomous goal pursuit predicts better outcomes than therapist-imposed goals.",
        "Miller & Rollnick (2023). MI 4th Ed. Goal negotiation and planning.",
        notes="This creates ethical tension for many clinicians: 'What if they choose a goal I think is wrong?' The MI response: trust the process. A client who chooses moderation and discovers it doesn't work has learned something they couldn't learn from being told. If you respect their autonomy now, they'll return to you when they're ready for a different goal. If you override them, you lose them entirely."
    )



    # =====================================================================
    # SECTION 14: Self-Efficacy Enhancement
    # =====================================================================
    section_divider(prs, 14, "Self-Efficacy Enhancement in MET",
                    "Building Confidence for Change")

    content_slide(prs,
        "Strategies for Building Self-Efficacy in MET",
        [
            "Self-efficacy enhancement - building the client's confidence that they CAN change - is a core principle of MI/MET practice. Bandura (1977) demonstrated that self-efficacy beliefs are often more predictive of behavior than actual ability, making confidence-building a therapeutic priority.",
            "Key strategies include: (1) Affirming past successes: 'You quit smoking 5 years ago - that shows you can make and maintain difficult changes.' (2) Normalizing the difficulty: 'Most people take several attempts before finding what works. Each attempt teaches you something.' (3) Highlighting existing strengths: 'You've shown a lot of courage just coming here today.'",
            "The confidence ruler is a particularly powerful tool: 'On a scale of 0-10, how confident are you that you could reduce your drinking if you decided to?' Regardless of the number, the follow-up explores existing confidence: 'You said a 4. What gives you that much confidence? What would it take to move from a 4 to a 5 or 6?'",
            "In Indian settings, self-efficacy can be connected to cultural strengths: family support systems, religious faith, community belonging, past demonstrations of discipline (fasting, religious observances), and professional accomplishments. Nadkarni et al. (2023) found that connecting confidence to existing cultural practices enhanced engagement.",
        ],
        "Self-efficacy is built through affirming past successes, normalizing difficulty, and exploring existing confidence rather than deficits.",
        "Bandura, A. (1977). Self-efficacy theory; Nadkarni et al. (2023). AMBIT Trial, Goa.",
        notes="Common trap: trying to build self-efficacy by minimizing the difficulty of change ('It's not that hard, you can do it!'). This actually undermines confidence because the client knows it IS hard. Better approach: acknowledge difficulty while affirming capability: 'This is one of the hardest things you could undertake, and you have qualities that suggest you can succeed.'"
    )

    content_slide(prs,
        "Responding to Low Confidence and Hopelessness",
        [
            "Some clients present with extremely low self-efficacy: 'I've tried everything, nothing works, I can't do this.' This challenges even experienced MET clinicians, who may feel tempted to argue against the client's pessimism. However, arguing FOR confidence creates the same dynamic as arguing FOR change - it triggers resistance.",
            "Instead of arguing, explore with curiosity: 'You've been through a lot of attempts. Tell me about what you've tried.' Within their history of 'failures,' there are often periods of success (however brief) that can be highlighted: 'So you were alcohol-free for 3 months after the last hospitalization. What was different during that time?'",
            "Reframing 'failures' as learning: 'It sounds like you've learned a lot about what doesn't work for you. That's actually valuable information. What does your experience tell you about what might work differently this time?' This honors their experience while gently redirecting toward possibility.",
            "Sometimes the most honest response is compassionate acknowledgment: 'You're feeling pretty hopeless right now. That must be exhausting.' Simply being heard in their despair can paradoxically free clients to consider possibility. The empathy mechanism (Rogers, 1957) operates even in moments of low confidence.",
        ],
        "Never argue against hopelessness - explore it with curiosity, find hidden successes within 'failures,' and provide empathic presence.",
        "Rogers, C.R. (1957). Necessary conditions for therapeutic change; Bandura (1977). Self-efficacy.",
        notes="Advanced technique: 'Suppose someone very much like you, who had been through similar experiences, came to you for advice. What would you tell them?' This externalizes the perspective and often allows clients to access wisdom and compassion they cannot direct toward themselves."
    )

    # =====================================================================
    # SECTION 15: Relapse Prevention Integration
    # =====================================================================
    section_divider(prs, 15, "Relapse Prevention Integration",
                    "Combining MET with RP Strategies")

    content_slide(prs,
        "Integrating Relapse Prevention with MET",
        [
            "Marlatt and Gordon's (1985) Relapse Prevention (RP) model provides practical strategies that complement MET's motivational approach. While MET focuses on building motivation to initiate change, RP focuses on maintaining change by identifying and managing high-risk situations, coping with cravings, and recovering from lapses.",
            "In MET Sessions 3-4, RP concepts are introduced within the MI spirit: rather than prescribing coping strategies, the therapist explores with the client: 'What situations do you think will be most challenging?' 'What has helped you in the past when you felt tempted?' 'What would you like to have in place before those situations arise?'",
            "The MET approach to lapse/relapse differs from traditional RP in tone: traditional RP can become prescriptive ('You must identify triggers and use these coping skills'). MET-informed RP explores collaboratively: 'It sounds like weekends with your drinking friends are the hardest times. What ideas do you have about how to handle those situations?'",
            "Marlatt's concept of the 'abstinence violation effect' (AVE) - where a single lapse triggers catastrophic thinking and return to full relapse - can be addressed in MET style: 'If you do have a drink after deciding to stop, what would you want to tell yourself?' This prepares without prescribing.",
        ],
        "MET + Relapse Prevention: motivation to START change combined with skills to MAINTAIN it - delivered collaboratively.",
        "Marlatt, G.A. & Gordon, J.R. (1985). Relapse Prevention. Guilford Press.",
        notes="The integration point: MET builds motivation (Sessions 1-2), then RP concepts are woven in during consolidation (Sessions 3-4) once the client has begun making changes. Never introduce RP strategies before motivation is established - premature skill training for unmotivated clients is ineffective and can feel coercive."
    )

    content_slide(prs,
        "High-Risk Situations: Collaborative Identification",
        [
            "Identifying high-risk situations for relapse is a collaborative process in MET. The therapist guides exploration using the MI principle of evocation: 'Based on your experience, what situations do you think pose the greatest risk for returning to heavy drinking?' The client's self-knowledge is honored as the primary data source.",
            "Common high-risk categories include: emotional states (stress, anger, loneliness, boredom, celebration), social situations (drinking friends, parties, workplace events), environmental cues (walking past a bar, time of day, specific locations), and interpersonal conflicts (arguments, criticism, rejection).",
            "In Indian contexts, specific high-risk situations may include: festival celebrations where drinking is socially expected, salary day patterns, drinking with colleagues after work, family conflicts triggering escape drinking, and peer pressure in social gatherings. Understanding culturally specific risk situations enhances planning.",
            "The MET therapist helps clients develop personalized coping plans for their specific high-risk situations, again using evocative questions: 'What could you do differently when your colleagues invite you for drinks after work?' 'What has worked for you in the past in similar situations?' 'Who could support you in those moments?'",
        ],
        "High-risk situations are personally identified, not prescribed from a list. Cultural context shapes which situations pose greatest risk.",
        "Marlatt & Gordon (1985). Relapse Prevention; NIMHANS (2008). Indian adaptation considerations.",
        notes="Exercise: Have clients keep a 'risk diary' for one week, noting when they felt most tempted and what was happening. Review this collaboratively in session to identify patterns. The diary itself can be therapeutic - awareness of patterns is the first step toward managing them."
    )



    # =====================================================================
    # SECTION 16: Special Populations
    # =====================================================================
    section_divider(prs, 16, "MET with Special Populations",
                    "Adapting MET for Diverse Clinical Needs")

    content_slide(prs,
        "MET for Co-occurring PTSD and Substance Use",
        [
            "Hussey et al. (2023) demonstrated the effectiveness of MET for improving treatment engagement among military veterans with co-occurring PTSD and substance use disorders. These clients often avoid trauma-focused therapy due to fear of emotional flooding, making motivational enhancement particularly relevant.",
            "The adaptation involves using MET principles to address ambivalence about entering trauma-focused treatment rather than (or in addition to) ambivalence about substance use. The personalized feedback component includes information about how PTSD and substance use interact, maintaining awareness of the trauma-substance connection.",
            "MI-consistent exploration of avoidance: 'You mentioned that you started drinking more after your deployment. I notice that you've been offered trauma therapy before but decided not to go. Tell me about that decision.' This opens the door without pressure.",
            "The integration of MET with PTSD treatment demonstrates MET's versatility as a pre-treatment intervention: preparing clients for intensive therapy by building motivation and addressing ambivalence before asking them to engage in emotionally demanding treatment protocols.",
        ],
        "MET can serve as a motivational bridge to other treatments - building readiness for interventions clients initially avoid.",
        "Hussey et al. (2023). MET for PTSD treatment engagement in military populations.",
        notes="Clinical pearl: Many clients with co-occurring PTSD and substance use have been through multiple failed treatment attempts, often because programs demanded immediate abstinence before addressing trauma. MET's non-prescriptive approach allows exploration of the bidirectional relationship between trauma symptoms and substance use without demanding a particular treatment sequence."
    )

    content_slide(prs,
        "MET for Tobacco Cessation in Schizophrenia",
        [
            "PMC (2024) published a comprehensive review of MET for tobacco cessation in individuals with schizophrenia, addressing a population where smoking rates remain 2-3 times higher than the general population and conventional cessation approaches show limited effectiveness.",
            "Unique challenges in this population include: cognitive symptoms that make complex behavioral planning difficult, negative symptoms that reduce motivation generally, positive symptoms (delusions, hallucinations) that may complicate therapeutic communication, medication side effects that nicotine partially ameliorates, and social isolation that makes smoking a primary social activity.",
            "MET adaptations for schizophrenia include: shorter sessions, simpler feedback presentation, more repetition of key points, involving support persons in planning, addressing the specific function of smoking (managing medication side effects, providing structure, social connection), and being realistic about goal-setting.",
            "The MI spirit is particularly important here: individuals with schizophrenia have often experienced coercive treatment and paternalistic decision-making. An approach that genuinely respects their autonomy and explores their perspective can be powerfully therapeutic regardless of smoking outcomes.",
        ],
        "MET adapted for serious mental illness respects autonomy in populations accustomed to paternalistic care - itself therapeutic.",
        "PMC (2024). MET for tobacco cessation in schizophrenia. Comprehensive review.",
        notes="Important principle: meeting clients where they are. A person with schizophrenia who reduces from 30 to 20 cigarettes daily has made a significant health improvement, even if 'cessation' hasn't been achieved. MET's non-judgmental approach allows celebration of ANY positive movement rather than demanding abstinence."
    )

    content_slide(prs,
        "Group MET for Opioid Dependence",
        [
            "PMC (2024) published an RCT protocol for Group MET for opioid-dependent individuals on maintenance treatment (methadone or buprenorphine). This adaptation addresses the need for efficient, scalable interventions in opioid treatment programs where individual sessions are resource-intensive.",
            "The group format adapts MET principles: personalized feedback is delivered individually (each participant receives their own PFR) but processed in a group context. Group members provide vicarious experiences (seeing peers at different stages), social support, and normalized sharing of ambivalence.",
            "MI principles in group format require specific adaptations: the facilitator must maintain MI spirit with each individual while managing group dynamics, ensure that group members do not confront or pressure each other, and use group reflection techniques ('What did others notice about what Maria just shared?').",
            "Olgac et al. (2024) at Case Western Reserve reviewed MI/MET for opioid use disorders more broadly, finding evidence for MET as both standalone intervention and adjunct to MAT. The group format represents a middle path between individual MET and standardized group therapy.",
        ],
        "Group MET combines individual feedback with group support - scalable for opioid programs with high caseloads.",
        "PMC (2024). Group MET for opioid dependence on maintenance treatment - RCT protocol.",
        notes="Group MI/MET is challenging to deliver with fidelity because group dynamics can undermine MI spirit. Key rules: no cross-talk that involves confrontation or advice-giving between members, therapist models MI-consistent responses, use of structured exercises that allow individual reflection within group context."
    )

    # =====================================================================
    # SECTION 17: Training and Supervision
    # =====================================================================
    section_divider(prs, 17, "Training and Supervision in MET",
                    "Developing and Maintaining Clinical Competence")

    content_slide(prs,
        "Components of Effective MI/MET Training",
        [
            "Research on MI/MET training consistently shows that workshop-only training (even intensive workshops) produces knowledge gains but limited sustained behavior change in clinical practice. Effective training requires ongoing feedback, supervision, and practice opportunities to develop and maintain competence.",
            "The recommended training sequence includes: (1) Initial workshop (2-3 days) covering theory, demonstration, and initial practice; (2) Supervised practice with actual clients over 3-6 months; (3) Ongoing coaching with feedback on recorded sessions; (4) Fidelity monitoring using validated coding instruments.",
            "Key training content areas: understanding the spirit of MI, OARS skills practice, recognizing and responding to change talk and sustain talk, avoiding MI-inconsistent traps (expert trap, premature focus, etc.), delivering personalized feedback in MI style, and managing complex clinical situations (discord, ambivalence about treatment itself).",
            "In India, Nadkarni et al. (2023) demonstrated that lay counselors in Goa could be trained to deliver MI with fidelity through a structured program of initial training, ongoing supervision, and fidelity monitoring. This suggests that the training model can be adapted for non-specialist health workers in LMIC settings.",
        ],
        "Effective MI/MET training requires more than workshops - ongoing supervised practice with feedback is essential for competence.",
        "Nadkarni et al. (2023). AMBIT Trial - training model for lay counselors in Goa, India.",
        notes="Training myth: 'I went to a 2-day MI workshop, so I know MI.' Research shows that workshop knowledge without supervised practice leads to minimal skill change. The actual skill development happens in the months AFTER the workshop when trainees practice with real clients and receive feedback. Budget for this ongoing component."
    )

    content_slide(prs,
        "Fidelity Assessment: The MITI and Other Tools",
        [
            "The Motivational Interviewing Treatment Integrity (MITI) coding system is the gold standard for assessing MI/MET fidelity. It evaluates both global qualities (partnership, empathy, evoking, cultivating change talk) and specific behaviors (reflections, questions, MI-consistent and MI-inconsistent behaviors).",
            "MITI behavioral counts track: reflections (simple and complex), questions (open and closed), affirmations, MI-consistent behaviors (asking permission, emphasizing autonomy, seeking collaboration), and MI-inconsistent behaviors (confronting, advising without permission, directing). These yield key ratios (reflection:question, complex:simple reflection, %MI-consistent).",
            "Competency benchmarks: Beginning proficiency requires reflection:question ratio of 1:1, with 40% complex reflections and 90% MI-consistent behaviors. Expert proficiency requires 2:1 reflection:question ratio, 50% complex reflections, and 95% MI-consistent behaviors. These benchmarks guide training progression.",
            "In clinical supervision, recorded sessions (with client permission) are coded using MITI and reviewed with the clinician, identifying strengths and specific areas for skill development. This targeted feedback produces more improvement than general supervision discussions about cases.",
        ],
        "MITI coding provides objective feedback on MI/MET practice - essential for skill development beyond self-assessment.",
        "Moyers et al. (2007). MITI coding system; Miller & Rollnick (2023). Fidelity standards.",
        notes="Self-assessment of MI competence is notoriously unreliable - most practitioners rate themselves higher than objective coding suggests. Regular MITI coding of sessions (even periodic sampling) provides reality-checking and specific behavioral targets. It also demonstrates growth over time, which motivates continued skill development."
    )



    # =====================================================================
    # SECTION 18: MET in Primary Care and Brief Settings
    # =====================================================================
    section_divider(prs, 18, "MET in Primary Care Settings",
                    "Brief Interventions and Screening-Based Approaches")

    content_slide(prs,
        "Screening and Brief Intervention: MET Principles in Primary Care",
        [
            "Screening, Brief Intervention, and Referral to Treatment (SBIRT) is a public health approach that applies MET/MI principles in primary care settings. After screening identifies at-risk individuals (typically using AUDIT), a brief motivational intervention of 5-15 minutes can produce significant reductions in drinking.",
            "The brief intervention follows MET structure in miniature: feedback on screening results ('Your AUDIT score places you in the hazardous drinking range - here's what that means'), exploration of the client's reaction, brief evocation of motivation, and suggestion of resources if the client is interested.",
            "This represents the minimal effective dose of MET principles: even a single brief motivational conversation during a medical visit can produce meaningful behavior change. The Cochrane Review (Schwenker et al., 2023) includes brief interventions in its evidence base, confirming effectiveness of even minimal-dose MI.",
            "In India, integrating brief MI into primary healthcare presents enormous potential given the treatment gap. With 62.5 million alcohol users (NIMHANS, 2008) and limited specialist resources, training primary care physicians and nurses in brief motivational interventions could reach millions who will never access specialty treatment.",
        ],
        "Brief MET-informed interventions in primary care can reach millions who would never access specialty addiction treatment.",
        "Schwenker et al. (2023). Cochrane Review; NIMHANS (2008). Indian healthcare integration.",
        notes="The key message for primary care: you don't need to be an MI expert to deliver a brief motivational intervention. The essential elements are: provide feedback non-judgmentally, ask about the client's perspective, explore importance and confidence, and leave the door open. Even imperfect MI in primary care is better than no intervention."
    )

    content_slide(prs,
        "Emergency Department Applications of MET",
        [
            "Emergency departments (EDs) represent a unique opportunity for MET-informed intervention because alcohol-related presentations are common, patients may be particularly receptive during a crisis moment, and the medical context provides natural authority for raising health concerns.",
            "Brief motivational interventions in EDs (typically 20-30 minutes) have shown effectiveness in multiple RCTs for reducing subsequent drinking, injury recurrence, and hospital readmission. The intervention typically occurs after medical stabilization, either during the ED visit or in follow-up.",
            "The 'teachable moment' concept suggests that the immediate aftermath of an alcohol-related event (injury, DUI, family crisis) may create temporary openness to change. MET capitalizes on this by providing feedback ('The reason you're here today is directly related to alcohol') within an empathic framework ('I'm not here to lecture you, but I am concerned').",
            "Adaptation for Indian emergency settings requires sensitivity to: lack of privacy in crowded EDs, family members often present (which can be leveraged or may inhibit honest discussion), potential police involvement in injury cases, and stigma around admitting alcohol problems in a medical setting.",
        ],
        "Emergency departments offer teachable moments - brief MET during crisis can catalyze change when motivation is temporarily heightened.",
        "Schwenker et al. (2023). Cochrane Review includes ED-based MI interventions.",
        notes="Practical challenge: ED physicians and nurses have limited time and competing priorities. The solution is training dedicated 'health advisors' or 'brief intervention specialists' who can be called to consult when screening identifies at-risk patients. The AMBIT model (Nadkarni et al., 2023) of lay counselor delivery could be adapted for Indian EDs."
    )

    # =====================================================================
    # SECTION 19: Ethical Considerations
    # =====================================================================
    section_divider(prs, 19, "Ethical Considerations in MET",
                    "Autonomy, Honesty, and Clinical Boundaries")

    content_slide(prs,
        "Autonomy vs. Beneficence: The Central Ethical Tension",
        [
            "MET's emphasis on client autonomy creates genuine ethical tension with the principle of beneficence (acting in the client's best interest). What do we do when a client's autonomous choice appears to be self-destructive? How do we respect autonomy while fulfilling our duty of care?",
            "The MET resolution of this tension: autonomy and beneficence are not truly opposed because coerced change is rarely sustained. Respecting autonomy IS acting in the client's best interest because only autonomous change endures. Imposed change collapses when external pressure is removed.",
            "However, boundaries exist: when a client's choices pose immediate danger to self or others (acute suicidality, impaired driving, child endangerment), the clinician's duty to protect overrides autonomy just as it would in any therapeutic modality. MET does not require passivity in the face of acute risk.",
            "Informed consent in MET should explicitly address the approach: 'In our work together, I won't be telling you what to do or pressuring you to make changes you're not ready for. Instead, I'll help you explore your own thoughts and feelings about drinking and what, if anything, you'd like to be different.' This sets appropriate expectations.",
        ],
        "Respecting autonomy IS beneficence because only self-determined change endures - but acute safety always takes precedence.",
        "Miller & Rollnick (2023). MI 4th Ed. Ethics of MI practice.",
        notes="Discussion prompt: A client says they plan to drive home from session after disclosing heavy drinking today. How do you maintain MI spirit while addressing immediate safety? Answer: Safety trumps style. You can be direct about safety while being compassionate: 'I care about you getting home safely. I can't let you drive right now - let's figure out alternatives together.'"
    )

    content_slide(prs,
        "Manipulation vs. Influence: Drawing the Line",
        [
            "A sophisticated critique of MI/MET raises the question: is strategic use of reflections and questions to evoke change talk a form of manipulation? Where is the line between therapeutic influence (which all therapies exert) and unethical manipulation?",
            "Miller and Rollnick (2023) address this directly: the distinction lies in (1) transparency about the approach, (2) the client's interests being genuinely prioritized over the therapist's, (3) the client's autonomy being consistently respected including the right to not change, and (4) the therapist genuinely holding the outcome as the client's choice.",
            "Manipulation occurs when: the therapist has a hidden agenda (e.g., protecting institutional liability rather than serving the client), techniques are used to produce outcomes the client doesn't want, or the therapist's needs (appearing successful, meeting program targets) override the client's autonomous decision-making.",
            "The safeguard is genuine compassion: a therapist who truly cares about the client's welfare and truly respects their autonomy cannot manipulate, even while using strategic communication skills. Spirit precedes and constrains technique. Without spirit, the same techniques become manipulative.",
        ],
        "The line between therapeutic influence and manipulation lies in genuinely prioritizing client welfare and respecting their right to choose.",
        "Miller & Rollnick (2023). MI 4th Ed. Ethical practice and the spirit of MI.",
        notes="Self-reflection exercise for practitioners: 'Whose interests am I primarily serving right now? Am I frustrated because this client isn't changing, and if so, whose need does their change serve?' If the answer is anything other than 'the client's freely chosen interests,' examine your clinical stance."
    )



    # =====================================================================
    # SECTION 20: MET with Mandated Clients
    # =====================================================================
    section_divider(prs, 20, "MET with Mandated Clients",
                    "Preserving Autonomy Within Coerced Treatment")

    content_slide(prs,
        "Applying MI Spirit with Court-Mandated Clients",
        [
            "Mandated clients (those required to attend treatment by courts, employers, or family pressure) present a specific challenge for MET: how do you support autonomy when someone has no choice about attending? Paradoxically, research suggests MI may be MOST effective with mandated clients because of its respect for autonomy.",
            "The approach begins with honest acknowledgment: 'You're here because the court required it, not because you chose to come. I want to acknowledge that and ask - given that you have to be here, how would you like to use this time?' This respects the reality while opening space for genuine engagement.",
            "Project MATCH (1997) found that angry clients (who may overlap substantially with mandated clients) showed better outcomes in MET than in other treatments. The interpretation: MI's non-confrontational approach de-escalates defensiveness that other treatments exacerbate through demands for compliance and admission of problems.",
            "Strategies for mandated clients include: separating attendance requirement (non-negotiable) from change expectation (client's choice), exploring what brought them to the court's attention from their perspective, finding areas of genuine concern or desire for change even within a hostile presentation, and being transparent about reporting requirements.",
        ],
        "MI may be MOST effective with mandated/angry clients because it de-escalates rather than intensifies the confrontational dynamic.",
        "Project MATCH (1997). Angry clients showed differential benefit from MET vs. other treatments.",
        notes="Key insight: mandatory attendance does NOT mean mandatory change. You can satisfy the court requirement while still respecting the client's autonomy about what they do with the information. Many mandated clients, once they feel genuinely heard and un-judged, discover authentic motivation they didn't know they had. The mandate forced them through the door; MI keeps them engaged once inside."
    )

    # =====================================================================
    # SECTION 21: MET and Pharmacotherapy
    # =====================================================================
    section_divider(prs, 21, "MET and Pharmacotherapy",
                    "Enhancing Medication Adherence and Combined Treatment")

    content_slide(prs,
        "MET as Adjunct to Medication-Assisted Treatment",
        [
            "MET can serve as an effective adjunct to pharmacotherapy for substance use disorders, addressing a critical challenge: medication adherence. Even effective medications like naltrexone, acamprosate, or buprenorphine require consistent adherence to produce benefits, and many patients discontinue prematurely.",
            "The COMBINE study (2006) demonstrated that MI-based counseling enhanced outcomes when combined with naltrexone for alcohol dependence. The motivational component addressed ambivalence about medication use, explored side effects collaboratively, and built commitment to continued adherence.",
            "MET principles apply to medication discussions: rather than lecturing about medication benefits, explore the client's perspective: 'What do you know about naltrexone?' 'What concerns, if any, do you have about taking medication?' 'How does medication fit with how you think about recovery?' This evokes the client's own reasons for adherence.",
            "For opioid maintenance treatment, Olgac et al. (2024) and the PMC (2024) Group MET protocol specifically address motivation for continued MAT engagement, where clients may face ambivalence about being 'on medication' indefinitely versus the stability it provides.",
        ],
        "MET enhances medication outcomes by addressing adherence ambivalence - clients are more likely to stay on medications they've chosen collaboratively.",
        "COMBINE Study (2006); Olgac et al. (2024). MI/MET for opioid use disorders. Case Western Reserve.",
        notes="Practical example: A client on naltrexone says 'I don't think I need the medication anymore - I feel fine.' Instead of arguing for continued adherence, explore: 'You're feeling good, which is great. What do you think is contributing to feeling good?' Often the client will identify the medication itself as one factor, generating their own argument for continuation."
    )

    # =====================================================================
    # SECTION 22: MET Across the Lifespan
    # =====================================================================
    section_divider(prs, 22, "MET Across the Lifespan",
                    "Adaptations for Adolescents, Adults, and Older Adults")

    content_slide(prs,
        "MET with Adolescents and Young Adults",
        [
            "Adapting MET for adolescents and young adults requires sensitivity to developmental factors: identity formation (substance use may serve identity exploration purposes), peer influence (often stronger than family influence), developmental trajectory (many adolescent substance use resolves naturally), and cognitive development (abstract reasoning about future consequences is still developing).",
            "Normative feedback is particularly powerful with young drinkers who typically overestimate peer consumption even more than adults do. College students shown that they drink more than 85% of peers often show immediate reduction in consumption - the discrepancy between perceived and actual norms drives change.",
            "Autonomy support is both more challenging (parents/schools may want the therapist to 'fix' the adolescent) and more important (adolescents are developmentally primed for reactance against authority). The therapist must maintain MI spirit even when external pressures push toward directive approaches.",
            "Key adaptations include: shorter sessions, more visual/interactive feedback delivery, exploring valued identity goals ('What kind of person do you want to be known as?'), acknowledging the social functions of use honestly, and separating the therapeutic relationship from parental/institutional authority.",
        ],
        "MET with adolescents leverages developmental sensitivity to norms and identity while respecting heightened need for autonomy.",
        "Miller & Rollnick (2023). MI 4th Ed. Adaptations across populations.",
        notes="With adolescents, the 'values exploration' often centers on identity: 'Who do you want to become? What kind of reputation do you want? What do you want people to say about you?' These identity-based values can be powerful motivators when connected to substance use patterns."
    )

    content_slide(prs,
        "MET with Older Adults",
        [
            "Older adults present unique considerations for MET adaptation: late-onset drinking (which may respond differently to intervention than early-onset), physiological changes that increase alcohol sensitivity, polypharmacy interactions, social isolation that alcohol may address, grief and loss as precipitants, and potential cognitive decline affecting engagement.",
            "The respectful, non-confrontational nature of MET is particularly well-suited for older adults who may feel patronized by directive approaches. Older adults often respond well to discussion of values like independence, health maintenance, and relationship with family - all potentially affected by alcohol use.",
            "Feedback components may emphasize: age-related changes in alcohol metabolism (same amount has greater effect with age), interaction with medications, fall risk, cognitive impact, and loneliness/isolation patterns. Normative comparison can note that drinking decreases with age in the general population.",
            "Cultural considerations in India: older adults may have drinking patterns established over decades, may feel entitled to drink as an elder, or may be drinking in response to changed family dynamics (children leaving, retirement, widowhood). Respect for elders must be maintained while gently exploring the impact of drinking.",
        ],
        "MET's respectful approach suits older adults well - focus on values of independence, health, and family connection.",
        "Miller & Rollnick (2023). MI 4th Ed. Age-related adaptations.",
        notes="Clinical pearl: Many older adults have never been asked about their drinking in a non-judgmental way. The simple act of respectful inquiry ('Tell me about your drinking - how it fits into your daily life') can open profound conversations about loneliness, loss, and meaning in later life. Sometimes the alcohol issue leads to deeper therapeutic work."
    )



    # =====================================================================
    # SECTION 23: MET for Polysubstance Use
    # =====================================================================
    section_divider(prs, 23, "MET for Polysubstance Use",
                    "Addressing Multiple Substances Within the MI Framework")

    content_slide(prs,
        "Prioritizing Change Targets in Polysubstance Use",
        [
            "Many clients presenting for MET use multiple substances (alcohol, tobacco, cannabis, stimulants, opioids). The MI principle of focusing becomes critical: which substance(s) to address, in what order, and who decides? The MET approach honors client autonomy in prioritization while providing relevant clinical information.",
            "The therapist explores: 'You've mentioned drinking, smoking, and occasional cocaine use. If you were to think about making changes, which of these concerns you most?' This evocative question lets the client prioritize based on their own assessment of importance and readiness.",
            "Personalized feedback can be provided for each substance, allowing the client to see the full picture while choosing their starting point. Some clients are surprised to discover that a substance they considered 'minor' (like tobacco) actually has greater long-term health risk than their primary concern.",
            "The MET principle applies across substances: never argue about which substance the client 'should' address first. A client who is motivated to quit smoking but ambivalent about alcohol will make more progress starting with tobacco. Success with one change builds self-efficacy for tackling other substances later.",
        ],
        "Client-directed prioritization of change targets produces better outcomes than clinician-imposed substance hierarchies.",
        "Miller & Rollnick (2023). MI 4th Ed. Focusing with multiple change targets.",
        notes="Exception: when immediate medical danger exists (e.g., severe alcohol or benzodiazepine dependence with withdrawal risk), the clinician should share this information with permission. 'I need to let you know that stopping alcohol suddenly at your level of use could be medically dangerous. Can we discuss safe options?'"
    )

    content_slide(prs,
        "Cannabis and MET: Unique Considerations",
        [
            "Cannabis presents unique challenges for MET because perceived harm is lower than for other substances, social acceptance is increasing in many contexts, and many users genuinely experience benefits (relaxation, sleep, pain management) that make the decisional balance more evenly weighted.",
            "The MET approach respects this complexity: rather than arguing that cannabis is harmful (which clients may accurately perceive as overstated), the therapist explores the client's specific experience: 'What's working about cannabis use for you? What, if anything, concerns you about it?'",
            "Feedback components specific to cannabis may include: driving performance data, cognitive impact studies, impact on motivation and goal pursuit, financial calculation, legal consequences where relevant, and comparison of actual vs. perceived use frequency (many daily users underestimate their consumption).",
            "In the Indian context, cannabis (bhang, ganja, charas) has complex cultural and religious associations that differ from Western contexts. MET delivered in India must be sensitive to these cultural meanings while still exploring whether use is aligned with the individual's broader life goals and values.",
        ],
        "MET for cannabis use requires balanced exploration - acknowledging genuine benefits while exploring costs and alignment with values.",
        "Miller & Rollnick (2023). MI with ambivalence about 'low-harm' substances.",
        notes="With cannabis users, the 'importance ruler' often reveals that the client themselves doesn't see much reason to change (importance = 2-3 out of 10). The MI response is to explore whether their use is genuinely aligned with their values and goals rather than arguing for higher importance. Sometimes the answer is that their use is genuinely not problematic."
    )

    # =====================================================================
    # SECTION 24: Technology-Enhanced MET
    # =====================================================================
    section_divider(prs, 24, "Technology-Enhanced MET",
                    "Digital, Mobile, and Telehealth Delivery")

    content_slide(prs,
        "Mobile Phone Delivery: Lessons from AMBIT",
        [
            "The AMBIT trial (Nadkarni et al., 2023) demonstrated that MI can be delivered effectively via mobile phone by community health workers in Goa, India. This technology-enhanced approach addresses multiple barriers: geographic access, stigma of visiting mental health facilities, and scheduling constraints.",
            "Key adaptations for mobile delivery: sessions are typically shorter (20-30 minutes vs. 45-60 for in-person), more structured to maintain focus without visual cues, and scheduled around client availability. The loss of non-verbal communication is partially offset by increased attention to verbal tone and pacing.",
            "The AMBIT counselors received initial training in MI principles and techniques, followed by ongoing supervision that included review of recorded phone sessions. Fidelity was maintained through regular supervision and quality monitoring, demonstrating that MI quality can be sustained in remote delivery.",
            "Scalability implications: India has over 1 billion mobile phone subscriptions. If MI can be delivered effectively by trained lay counselors via phone, the potential reach is enormous - potentially addressing the needs of millions of alcohol users who currently have no access to evidence-based treatment.",
        ],
        "Mobile MI by trained lay counselors could potentially reach millions in India - the AMBIT model provides proof of concept.",
        "Nadkarni et al. (2023). AMBIT Trial. BMC Psychiatry. Mobile phone MI in Goa, India.",
        notes="The combination of task-shifting (lay counselors) and technology (mobile phones) represents a quantum leap in potential reach. Traditional models requiring trained psychologists meeting clients in clinics can serve thousands; the AMBIT model could potentially serve millions. This is the future of addiction treatment in LMICs."
    )

    content_slide(prs,
        "Telehealth MET: Post-Pandemic Adaptations",
        [
            "The COVID-19 pandemic accelerated telehealth adoption for all psychotherapies, including MI/MET. Post-pandemic research suggests that telehealth MI can be as effective as in-person delivery for many clients, while offering significant advantages in terms of access, convenience, and reach.",
            "Telehealth MI requires specific adaptations: ensuring visual connection (video preferred over phone-only for maintaining empathic connection), managing technical issues without losing therapeutic momentum, creating a sense of privacy and safety despite the remote context, and adapting reflective listening to work through a screen.",
            "Advantages of telehealth MET include: reduced stigma (no need to enter an addiction treatment facility), improved attendance rates (reduced travel and time barriers), ability to see the client's home environment (providing ecological validity), and flexibility in scheduling that can reduce no-show rates.",
            "Challenges include: digital divide (not all clients have reliable internet access), privacy concerns (others may overhear in shared living spaces), reduced ability to read non-verbal cues, and potential for distraction. In India, internet connectivity varies enormously between urban and rural areas.",
        ],
        "Telehealth MI expands access while maintaining effectiveness - but the digital divide must be addressed for equitable implementation.",
        "Post-pandemic telehealth literature; Nadkarni et al. (2023). Technology-assisted MI delivery.",
        notes="The AMBIT model (phone-based) may be more equitable than video-based telehealth in India, where smartphones with reliable internet are less universally available than basic mobile phones. Meeting the client where they are technologically is just as important as meeting them where they are motivationally."
    )

    # =====================================================================
    # SECTION 25: Cultural Competence in MET
    # =====================================================================
    section_divider(prs, 25, "Cultural Competence in MET",
                    "Delivering MET Across Diverse Cultural Contexts")

    content_slide(prs,
        "Cultural Adaptation Principles for MET",
        [
            "Cultural adaptation of MET goes beyond translation of materials - it requires understanding how culture shapes the experience of substance use, the meaning of change, and the therapeutic relationship. The systematic review of MI in LMICs (Medrxiv, 2023) identified key principles for successful cultural adaptation.",
            "Values exploration must be culturally informed: in collectivist cultures (common across Asia, including India), personal values may include family harmony, community reputation, filial duty, and collective wellbeing alongside individual goals. MET can explore discrepancy between drinking and these collective values: 'How does your drinking affect your family's standing?'",
            "Communication style adaptations: in some cultures, direct discussion of personal problems with a stranger is unusual. Building relationship and trust may require more time. The 'engaging' process may need to be extended, with more attention to establishing credibility, showing respect for cultural norms, and demonstrating understanding of the client's cultural context.",
            "Decision-making processes differ: in some cultures, major life decisions are made collectively (with family, elders, religious leaders). MET can accommodate this by exploring: 'Who else would be involved in this decision? What would they think? How important are their views to you?' while still centering the client's own agency.",
        ],
        "Cultural adaptation requires understanding how culture shapes values, communication, and decision-making - not just translating materials.",
        "Medrxiv (2023). Systematic Review: MI in LMICs; NIMHANS (2008). Cultural considerations.",
        notes="The MI principle of autonomy doesn't mean individualism. In collectivist cultures, autonomous choice might mean choosing to honor family wishes - this is still autonomous if it reflects the individual's genuine values rather than pure coercion. The therapist explores: 'Is this what YOU want, or what others expect?' Either answer can be valid."
    )

    content_slide(prs,
        "Gender Considerations in MET Delivery",
        [
            "Gender profoundly shapes the experience of substance use and the delivery of MET. In India particularly, male and female substance use patterns, social consequences, treatment access, and therapeutic needs differ dramatically due to cultural norms around gender, substance use, and helpseeking.",
            "For men in India: drinking may be socially normative in certain contexts, associated with masculinity, and facilitated by peer groups. MET can explore how drinking-related consequences conflict with valued masculine roles: provider, protector, respected community member. Patel et al. (2024) specifically addressed the alcohol-IPV connection in south Indian men.",
            "For women in India: substance use carries severe stigma, treatment-seeking is extremely rare, family consequences may include abandonment or violence, and women often drink in secret. MET adaptations must address safety concerns, stigma-reduction, and practical barriers (childcare, family permission to attend).",
            "Therapist gender matching may be important: some clients (particularly women in conservative settings) may feel more comfortable disclosing substance use to a same-gender therapist. However, the quality of the therapeutic relationship (empathy, non-judgment) may matter more than gender matching per se.",
        ],
        "Gender shapes substance use experience profoundly - MET must address gendered values, roles, barriers, and safety concerns.",
        "Patel et al. (2024). MI+BCT for men with alcohol and IPV in south India. PLOS ONE.",
        notes="In India, female alcohol use remains highly stigmatized in most communities. Women who drink face social ostracism, marital violence, and loss of custody. MET with these women must prioritize safety: 'Who knows about your drinking? Is it safe for you if others find out?' before addressing motivation for change."
    )

    # =====================================================================
    # SECTION 26: Advanced MI Techniques in MET
    # =====================================================================
    section_divider(prs, 26, "Advanced MI Techniques in MET",
                    "Complex Clinical Strategies for Experienced Practitioners")

    content_slide(prs,
        "Eliciting and Strengthening Commitment",
        [
            "Building on Amrhein et al.'s (2003) finding that commitment language predicts outcomes, advanced MET practice involves deliberate strategies for eliciting and strengthening commitment once preparatory change talk (DARN) has been sufficiently developed.",
            "The 'key question' is a powerful commitment-eliciting tool used when change talk is predominant: 'So where does this leave you?' or 'Given everything you've said, what do you think you'll do?' These questions invite the client to synthesize their own change talk into a commitment statement.",
            "Strengthening techniques include: asking for elaboration on commitment ('Tell me more about how you'll do that'), asking about next steps ('What would be the first thing you'd do?'), exploring the commitment's significance ('How important is this decision to you?'), and inviting the client to voice their commitment to significant others.",
            "The Change Plan Worksheet (used in MET Session 2) provides a structured vehicle for crystallizing commitment: writing it down, specifying when and how, identifying support, and anticipating barriers. The written commitment serves as both a planning tool and a public record of intention.",
        ],
        "Advanced MET moves from evoking change talk toward crystallizing commitment - the strongest predictor of actual behavior change.",
        "Amrhein et al. (2003). JCCP, 71, 862. Commitment language predicts outcomes.",
        notes="Timing is everything: premature commitment-eliciting (before DARN talk is dominant) produces weak or retracted commitments. Signs of readiness: decreased sustain talk, increased tempo and energy in change talk, asking 'how' questions rather than 'whether' questions, imagining specifics of changed behavior."
    )

    content_slide(prs,
        "Working with Ambivalence: Advanced Strategies",
        [
            "Deep ambivalence - where the client genuinely sees compelling arguments both for and against change - is perhaps the most clinically challenging presentation in MET. The temptation is to resolve the ambivalence by taking the 'change' side, but this simply pushes the client to the 'status quo' side.",
            "Advanced strategies for ambivalent clients include: amplifying the discrepancy through values-behavior contrast without taking sides, sitting with ambivalence (tolerating the discomfort rather than rushing to resolve it), exploring the ambivalence itself as the problem ('What's it like to be pulled in two directions?'), and using the decisional balance to make both sides explicit.",
            "The 'coming alongside' technique can paradoxically break ambivalence: 'Maybe this isn't the right time for you to make changes. You have a lot going on, and drinking is helping you cope. Perhaps this is something you'll want to revisit later when things are more stable.' This removes pressure and often evokes the change side.",
            "Columbo approach (named after the TV detective): playing intentionally naive - 'Help me understand something. You've told me that drinking is causing all these problems, and also that you want to keep doing it. I'm a bit confused about how both things can be true. Can you help me understand?' This invites the client to resolve the paradox themselves.",
        ],
        "Deep ambivalence is not resolved by argument but by creating conditions where the client resolves it themselves.",
        "Miller & Rollnick (2023). MI 4th Ed. Advanced strategies for persistent ambivalence.",
        notes="The 'Columbo approach' works because it avoids telling the client they have a contradiction (confrontation) while making the contradiction salient through genuine-seeming puzzlement. The client then feels compelled to explain - and in explaining, often resolves their own ambivalence. It respects intelligence while creating discrepancy."
    )

    content_slide(prs,
        "Responding to Sustain Talk Without Reinforcing It",
        [
            "One of the most nuanced clinical skills in advanced MET practice is responding to sustain talk in ways that acknowledge without reinforcing it. The therapist must validate the client's experience without amplifying the status quo arguments.",
            "The 'reflect and redirect' strategy: briefly reflect the sustain talk, then use a transitional phrase to redirect toward change talk. Example: Client says 'I just don't think I have a problem.' Therapist: 'You're not seeing drinking as a concern right now [reflect]. Tell me, what brought you here today? [redirect]'",
            "Strategic use of silence: when a client makes a strong sustain talk statement, sometimes the most effective response is a brief silence followed by a simple 'Hmm' or 'I see' before redirecting. This neither reinforces nor confronts. Dwelling on sustain talk (asking multiple follow-up questions about it) inadvertently reinforces it.",
            "Reframing offers a new perspective without arguing: 'I've never been able to quit' can be reframed as 'You've had the courage to try multiple times, even though it's been difficult' - same facts, different emphasis. The reframe highlights persistence (a strength) rather than failure.",
        ],
        "Acknowledge sustain talk briefly, then redirect - dwelling on it or arguing against it both make it stronger.",
        "Miller & Rollnick (2023). MI 4th Ed; Moyers et al. (2007). Therapist behaviors and outcomes.",
        notes="Training exercise: present a sustain talk statement and have trainees generate 3 different responses: one that reinforces (asking more about it), one that confronts (arguing against it), and one that acknowledges and redirects (MI-consistent). Most trainees instinctively choose one of the first two options - training shifts them toward the third."
    )



    # =====================================================================
    # SECTION 27: MET Outcomes and Cost-Effectiveness
    # =====================================================================
    section_divider(prs, 27, "MET Outcomes and Cost-Effectiveness",
                    "Clinical and Economic Evidence for Implementation")

    content_slide(prs,
        "Cost-Effectiveness of MET: The Economic Argument",
        [
            "The cost-effectiveness of MET is one of its strongest arguments for implementation, particularly in resource-limited settings. UKATT (2005) demonstrated that 3-session MET was significantly more cost-effective than 8-session SBNT, producing equivalent clinical outcomes at approximately one-third the clinical contact time.",
            "Project MATCH (1997) similarly showed that 4-session MET produced outcomes equivalent to 12-session CBT and TSF. When calculated per session hour, MET's therapeutic yield is approximately 3 times greater than longer treatments - each MET session produces proportionally more benefit.",
            "For health systems in LMICs, this efficiency is not merely advantageous but necessary. With limited trained workforce (NIMHANS, 2008 notes India's psychiatrist shortage) and massive unmet need (62.5 million alcohol users), only highly efficient interventions can begin to close the treatment gap.",
            "Beyond direct treatment costs, MET reduces downstream healthcare utilization: reduced emergency department visits, fewer hospitalizations, decreased work absenteeism, and reduced social costs (domestic violence, accidents, family breakdown). The economic case for MET implementation is compelling across multiple cost domains.",
        ],
        "MET is 3x more efficient per session than longer treatments - essential for closing treatment gaps in resource-limited settings.",
        "UKATT (2005). BMJ, 331, 541; Project MATCH (1997). JOSA, 58, 7-29.",
        notes="When making the case for MI/MET implementation to administrators: emphasize not just clinical effectiveness but economic value. A 4-session treatment that matches 12-session treatments allows you to serve 3 times as many clients with the same clinical workforce. In public health terms, this is transformative for treatment access."
    )

    table_slide(prs,
        "Comparative Outcomes: MET vs. Extended Treatments",
        ["Outcome Measure", "MET (4 sessions)", "CBT (12 sessions)", "TSF (12 sessions)", "Significant Difference"],
        [
            ["% Days Abstinent (1yr)", "~65-70%", "~65-70%", "~70-75%", "TSF slightly higher"],
            ["Drinks per Drinking Day", "Reduced ~50%", "Reduced ~50%", "Reduced ~50%", "No difference"],
            ["Cost per Client", "$500-800", "$1500-2400", "$1500-2400", "MET significantly lower"],
            ["Therapist Hours", "3-4 hours", "12 hours", "12 hours", "MET 3x more efficient"],
            ["Client Attendance", "High (brief)", "Moderate", "Moderate", "MET better retention"],
        ],
        "MET achieves equivalent clinical outcomes at 25-30% of the cost of longer treatments - a compelling economic argument.",
        "Project MATCH (1997). JOSA, 58, 7-29. Comparative outcomes across treatments.",
        notes="Note: these are approximate figures from Project MATCH. Individual studies vary. The key message is consistent across trials: MET's brevity does not compromise its effectiveness. Some have argued that MET SHOULD be the default first intervention, with longer treatments reserved for those who don't respond to brief approaches."
    )

    content_slide(prs,
        "Long-Term Outcomes: Durability of MET Effects",
        [
            "A critical question for any brief intervention is whether effects are sustained over time or merely produce short-term improvements that fade. Project MATCH (1997) addressed this with 3-year follow-up data, showing that improvements were maintained across all conditions including MET.",
            "The 3-year follow-up data showed that the initial substantial reduction in drinking (approximately 50% from baseline) was maintained with minimal deterioration. This suggests that MET's brief intervention catalyzes a change process that continues beyond therapy termination - consistent with the self-determination theory prediction.",
            "Magill et al. (2018) in their updated meta-analysis found that MI effects were generally maintained at follow-up timepoints, though some decay was observed at very long follow-up periods. The pattern suggests that MI/MET initiates change that is largely self-sustaining but may benefit from booster sessions.",
            "The mechanism of sustained effects likely involves the internalization of change motivation (Deci & Ryan, 2000): because MET builds autonomous rather than controlled motivation, clients continue to be motivated by their own values and goals long after therapy ends, unlike externally motivated change which collapses when the external pressure is removed.",
        ],
        "MET effects are durable at 3-year follow-up - autonomous motivation sustains itself because it comes from within the client.",
        "Project MATCH (1997). 3-year follow-up data; Deci & Ryan (2000). Self-Determination Theory.",
        notes="This addresses the skeptic's concern: 'Surely 4 sessions can't produce lasting change.' The answer from research is clear: it can, and the mechanism is the quality of motivation built (autonomous vs. controlled) rather than the quantity of therapist contact. Four sessions of genuine autonomous motivation-building outperforms 12 sessions of externally imposed treatment."
    )

    # =====================================================================
    # SECTION 28: Implementation Science
    # =====================================================================
    section_divider(prs, 28, "Implementation Science and MET",
                    "Scaling MET in Health Systems")

    content_slide(prs,
        "Barriers to MET Implementation in Health Systems",
        [
            "Despite strong evidence, MET implementation faces significant barriers in many health systems. These include: workforce training needs, organizational resistance to brief interventions (perceived as 'not enough treatment'), incompatibility with fee-for-service models that incentivize longer treatment, and misconceptions about MI being 'just talking nicely.'",
            "In the Indian context, specific barriers include: severe shortage of trained mental health professionals, limited supervision infrastructure for maintaining MI fidelity, stigma around seeking treatment for alcohol use, cultural expectations that treatment should be directive, and lack of standardized assessment tools validated in Indian languages.",
            "The AMBIT trial (Nadkarni et al., 2023) addressed several of these barriers simultaneously: task-shifting to lay counselors addressed workforce shortage, mobile delivery addressed access and stigma, structured protocols addressed fidelity concerns, and embedded supervision addressed quality maintenance.",
            "Implementation strategies from the systematic review (Medrxiv, 2023) that predicted success include: building local champions within organizations, integrating MI into existing care pathways rather than creating parallel systems, providing ongoing supervision rather than one-time training, and demonstrating cost-effectiveness to administrators.",
        ],
        "Successful MET implementation requires addressing workforce, organizational, and systemic barriers - not just clinical training.",
        "Nadkarni et al. (2023). AMBIT Trial implementation; Medrxiv (2023). MI in LMICs.",
        notes="Implementation is a science in itself. Having an effective treatment is necessary but not sufficient for improving population health. The gap between evidence and practice is enormous in mental health globally. MI/MET's relative simplicity and efficiency make it more implementable than many treatments, but barriers remain."
    )

    content_slide(prs,
        "Task-Shifting MET to Non-Specialist Providers",
        [
            "Task-shifting - training non-specialist health workers to deliver interventions traditionally provided by specialists - is the most promising strategy for scaling MET in resource-limited settings. The AMBIT trial (Nadkarni et al., 2023) provided proof-of-concept for MI delivery by community health workers in India.",
            "Key requirements for successful task-shifting include: structured training programs (typically 5-10 days initial training), ongoing supervision (weekly individual or group supervision), fidelity monitoring (regular coding of recorded sessions), simplified treatment manuals adapted for non-specialist literacy levels, and clear referral pathways for complex cases.",
            "WHO's Mental Health Gap Action Programme (mhGAP) explicitly recommends brief motivational interventions for alcohol use disorders deliverable by non-specialists in primary care. This global endorsement supports the case for task-shifting MI/MET as a strategy for closing the treatment gap.",
            "In India, potential non-specialist providers include: Accredited Social Health Activists (ASHAs), Auxiliary Nurse Midwives (ANMs), community health workers, trained peer counselors (individuals in recovery), primary care physicians with brief training, and teachers/counselors in educational settings.",
        ],
        "Task-shifting MI to community health workers is feasible and essential for reaching the millions who need but cannot access treatment.",
        "Nadkarni et al. (2023). AMBIT; WHO mhGAP recommendations for substance use.",
        notes="The ASHA (Accredited Social Health Activist) network in India comprises approximately 1 million community health workers. If even a fraction were trained in brief MI, the reach would be extraordinary. The challenge is supervision infrastructure - MI without supervision degrades over time. Technology-assisted supervision (phone-based, recorded sessions) may help."
    )

    # =====================================================================
    # SECTION 29: Future Directions
    # =====================================================================
    section_divider(prs, 29, "Future Directions for MET",
                    "Emerging Research and Innovation")

    content_slide(prs,
        "Emerging Technologies: AI-Enhanced MET Delivery",
        [
            "Artificial intelligence and natural language processing technologies are being explored as potential tools for enhancing MET delivery. Applications include: automated coding of MI fidelity (reducing supervision burden), AI-generated personalized feedback from electronic health records, chatbot-based brief motivational interventions, and real-time coaching for clinicians during sessions.",
            "The potential for AI-assisted training is particularly promising: systems that can analyze trainee recordings and provide immediate, objective feedback on MI-consistent and MI-inconsistent behaviors could dramatically reduce the supervision bottleneck that limits training scale-up.",
            "However, critical questions remain about whether AI can genuinely replicate the empathic human connection that research identifies as MI's primary mechanism (Valle, 1981; Rogers, 1957). Technology may be best positioned as a supplement to human-delivered MI rather than a replacement.",
            "In India, where smartphone penetration is increasing rapidly, AI-enhanced MI applications could leverage existing digital infrastructure. The AMBIT model's phone-based approach could be augmented with AI tools that support lay counselors with real-time suggestions and quality monitoring.",
        ],
        "AI may enhance MI delivery, training, and fidelity monitoring - but human empathic connection likely remains irreplaceable as the core mechanism.",
        "Emerging technology literature; Valle (1981). Empathy predicts outcomes.",
        notes="The key insight: technology should serve the relationship, not replace it. AI-coded fidelity feedback helps therapists improve their human skills. AI-generated normative data makes human-delivered feedback sessions more powerful. The moment we try to replace human empathy with algorithmic empathy, we lose the mechanism."
    )

    content_slide(prs,
        "Personalized and Precision Approaches to MET",
        [
            "The future of MET may involve greater personalization based on individual client characteristics, biological markers, genetics, and real-time data. Rather than a one-size-fits-all protocol, 'precision MET' would tailor feedback content, timing, and intensity to individual needs.",
            "Project MATCH (1997) attempted but largely failed to find robust treatment-matching variables (which clients benefit most from which treatments). However, modern data science approaches with larger datasets and more variables may yet identify meaningful matching criteria.",
            "Biological personalization: as understanding of addiction neurobiology advances, MET feedback could incorporate genetic risk information, brain imaging data, or biomarker profiles. This raises ethical questions about how to present biological information without undermining perceived self-efficacy or agency.",
            "Ecological Momentary Assessment (EMA) combined with Just-In-Time Adaptive Interventions (JITAIs) could deliver MI-consistent microdoses of motivation at high-risk moments identified through smartphone sensors. This represents a radical expansion of MET beyond the clinic into daily life.",
        ],
        "Precision MET may combine biological data, AI, and real-time monitoring for individually tailored motivational interventions.",
        "Project MATCH (1997); Emerging precision medicine literature.",
        notes="The philosophical question for precision MET: does personalizing the 'content' of feedback change the fundamental MI mechanism? If the mechanism is relationship + evocation (not information delivery), then precision content matters less than we might think. The relationship still does the heavy lifting regardless of how precisely targeted the feedback is."
    )

    content_slide(prs,
        "Integrating MET with Behavioral Economics",
        [
            "Behavioral economics offers insights that can enhance MET practice. Concepts like present bias (overweighting immediate rewards), loss aversion (losses feel bigger than equivalent gains), and choice architecture (how options are presented affects decisions) can inform how feedback is delivered and change plans are structured.",
            "Present bias explains why alcohol's immediate rewards (relaxation, social ease) outweigh its delayed consequences (health problems, relationship damage). MET can address this by making delayed consequences more vivid and personally salient through feedback, and by identifying immediate benefits of change (better sleep tonight, more energy tomorrow).",
            "Loss aversion suggests that framing change in terms of what the client loses by continuing (losing family, losing health, losing money) may be more motivating than framing it in terms of what they gain by changing. MET's discrepancy development naturally leverages this through showing what drinking costs.",
            "Choice architecture informs the Change Plan Worksheet: offering a limited menu of specific, concrete options (rather than open-ended 'what will you do?') can reduce decision fatigue and increase commitment. This doesn't violate autonomy if the client chooses among options they've generated.",
        ],
        "Behavioral economics insights can enhance MET by addressing cognitive biases that maintain problematic behavior patterns.",
        "Behavioral economics literature; Miller & Rollnick (2023). MI 4th Ed.",
        notes="The integration of behavioral economics with MI is an emerging frontier. Key insight: humans don't make rational decisions about health behaviors - they're influenced by framing, defaults, salience, and social norms. MI already implicitly addresses some of these (normative feedback leverages social norms). Explicit integration could enhance effectiveness."
    )

    # =====================================================================
    # SECTION 30: Summary and Integration
    # =====================================================================
    section_divider(prs, 30, "Summary and Integration",
                    "Bringing It All Together for Clinical Practice")

    content_slide(prs,
        "Core Messages: What Every MET Practitioner Must Remember",
        [
            "After 30 sections of detailed content, certain core messages emerge as foundational to effective MET practice. First and foremost: the spirit of MI (partnership, acceptance, compassion, evocation) is more important than any specific technique. Without spirit, techniques become manipulation.",
            "The mechanism is clear and well-supported: MI-consistent therapist behavior → therapeutic relationship → client change talk → commitment → behavior change (Moyers et al., 2007; DiClemente et al., 2017). Every clinical decision should serve this causal chain.",
            "Evidence is robust and consistent across populations, cultures, settings, and substances: from Project MATCH (1997) through Cochrane Review (Schwenker et al., 2023), MET produces clinically meaningful outcomes with remarkable efficiency. The evidence base now includes Indian-specific data (Nadkarni et al., 2023; Patel et al., 2024).",
            "Cultural adaptation is essential but should never abandon core principles. In India and other collectivist cultures, MI's emphasis on collaboration, respect, and non-confrontation is already culturally aligned. The adaptation lies in how values are explored and how autonomy is expressed within cultural norms.",
        ],
        "Spirit over technique, relationship over information, evocation over persuasion, autonomy over coercion - these are MET's core truths.",
        "Miller & Rollnick (2023). MI 4th Ed; DiClemente et al. (2017). Addiction, 112(S2), 92-100.",
        notes="Final synthesis for trainees: if you remember nothing else from this training, remember these four principles: (1) Listen more than you talk. (2) Ask open questions that evoke the client's own reasons for change. (3) Reflect what you hear with slightly deeper meaning. (4) Resist the urge to fix, advise, or convince. If you master these, you are practicing MI."
    )



    content_slide(prs,
        "The Practitioner's Self-Assessment Checklist",
        [
            "Effective MET practice requires ongoing self-monitoring and skill development. The following self-assessment questions help practitioners maintain fidelity to MI principles throughout their clinical work and identify areas for continued growth.",
            "Relationship quality: Am I genuinely curious about this client's experience? Do they feel understood? Am I maintaining positive regard even when I disagree with their choices? Would an observer rate my empathy as high based on my verbal responses? Am I experiencing frustration or judgment that might leak into my communication?",
            "Technical proficiency: Is my reflection-to-question ratio at least 2:1? Am I using complex reflections (not just parroting)? Am I strategically attending to change talk? Am I avoiding MI-inconsistent behaviors (unsolicited advice, confrontation, labeling)? Am I summarizing effectively?",
            "Directional skill: Am I evoking change talk effectively? Can I hear the difference between change talk and sustain talk in real-time? Am I responding to sustain talk without reinforcing it? Is the client's commitment language growing stronger across the session? Am I gauging readiness before moving to planning?",
        ],
        "Ongoing self-assessment and supervision are essential - MI/MET competence requires continuous development, not one-time training.",
        "Miller & Rollnick (2023). MI 4th Ed. Self-assessment and skill development.",
        notes="Encourage practitioners to record sessions (with client consent) and self-assess before supervision. Identifying your own MI-inconsistent moments before a supervisor points them out accelerates learning and reduces defensiveness in supervision."
    )

    process_slide(prs,
        "The MET Clinical Pathway: From Assessment to Outcomes",
        ["Assessment\n& Feedback", "Ambivalence\nExploration", "Commitment\nBuilding", "Change\nPlanning", "Long-term\nMaintenance"],
        "The MET pathway is flexible and client-paced - not all clients need all stages, and movement between stages is expected.",
        "Miller et al. (1995). MET Manual. Clinical pathway overview.",
        notes="This final process slide captures the entire MET journey. Remember: the pathway is not strictly linear. Some clients loop between ambivalence exploration and commitment building multiple times. Others move rapidly from assessment feedback directly to change planning. Meet the client where they are."
    )

    two_col_slide(prs,
        "What MET IS vs. What MET Is NOT",
        "MET IS:",
        [
            "A structured 4-session protocol",
            "Based on comprehensive assessment feedback",
            "Delivered in MI spirit (partnership, evocation)",
            "Client-centered but strategically guided",
            "Evidence-based (Project MATCH, UKATT, Cochrane)",
            "Culturally adaptable across settings",
            "Cost-effective and efficient",
            "Scalable through task-shifting",
        ],
        "MET IS NOT:",
        [
            "General supportive counseling",
            "Confrontational feedback delivery",
            "Just 'being nice' to clients",
            "Passive non-directional therapy",
            "Limited to alcohol use alone",
            "Only for Western contexts",
            "A replacement for medical treatment",
            "Appropriate without proper training",
        ],
        "MET is a specific, evidence-based protocol requiring training, fidelity monitoring, and cultural adaptation for effective delivery.",
        "Miller et al. (1995). MET Manual; Miller & Rollnick (2023). MI 4th Edition.",
        notes="This summary slide is useful for communicating what MET is to administrators, referral sources, and colleagues. Many people have vague notions of MI/MET as 'just being empathetic' or 'not confronting people.' The precision of the actual protocol - with its assessment, feedback, session structure, and training requirements - is often underappreciated."
    )

    content_slide(prs,
        "Recommended Reading and Resources",
        [
            "Essential texts: (1) Miller, W.R. & Rollnick, S. (2023). Motivational Interviewing: Helping People Change and Grow (4th ed.). Guilford Press. (2) Miller, W.R., Zweben, A., DiClemente, C.C., & Rychtarik, R.G. (1995). Motivational Enhancement Therapy Manual. NIAAA.",
            "Key research: Project MATCH Research Group (1997). JOSA, 58, 7-29. UKATT (2005). BMJ, 331, 541. Schwenker et al. (2023). Cochrane Review CD008063.pub3. DiClemente et al. (2017). Addiction, 112(S2), 92-100. Moyers et al. (2007). JCCP, 75, 790.",
            "Indian context: Nadkarni et al. (2023). AMBIT Trial, BMC Psychiatry. Patel et al. (2024). PLOS ONE. NIMHANS (2008). Psychosocial Interventions Manual. D'Costa et al. (2019). Drug & Alcohol Review. Benegal (2005). Addiction, 100, 1051-1056.",
            "Training resources: Motivational Interviewing Network of Trainers (MINT) - www.motivationalinterviewing.org. MITI 4.2 coding manual for fidelity assessment. MI Assessment: Supervisory Tools for Enhancing Proficiency (MIA-STEP) for training programs.",
        ],
        "Continued learning through reading, training, supervision, and practice is essential for maintaining and advancing MI/MET competence.",
        "Multiple sources - comprehensive reference list for continued professional development.",
        notes="Encourage practitioners to join MINT (Motivational Interviewing Network of Trainers) for ongoing learning community, attend annual MI conferences, seek certified MI training, and commit to regular supervision with MI-proficient supervisors. Skill development in MI is a career-long journey, not a single training event."
    )

    # Final closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    set_font(run, size=40, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Motivational Enhancement Therapy: Evidence-Based, Client-Centered, Culturally Adaptable"
    set_font(run2, size=18, color=GOLD)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = '"People are generally better persuaded by the reasons which they have themselves discovered'
    set_font(run3, size=14, italic=True, color=TEAL)
    p4 = tf.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = 'than by those which have come into the mind of others." — Blaise Pascal'
    set_font(run4, size=14, italic=True, color=TEAL)

    return prs




def add_supplementary_slides(prs):
    """Add additional detailed slides to ensure 150+ total."""

    # =====================================================================
    # Additional slides for Section 1 - Extended Introduction
    # =====================================================================
    content_slide(prs,
        "Key Principles Underlying MET Practice",
        [
            "MET rests upon several key principles that distinguish it from other therapeutic approaches. The first principle is that motivation is not a fixed trait but a dynamic state that fluctuates and can be influenced through interpersonal interaction. This means that apparent 'lack of motivation' is not a client deficit but a clinical opportunity.",
            "The second principle is that the therapeutic relationship itself is a potent catalyst for change. Research consistently shows that the quality of the working alliance predicts treatment outcomes across all psychotherapy modalities, but in MI/MET this relationship is not merely a vehicle for technique delivery - it IS the primary mechanism of change.",
            "The third principle is that ambivalence is normal, not pathological. Most people considering behavior change hold simultaneously competing motivations - reasons to change and reasons to maintain the status quo. Acknowledging and exploring both sides without judgment is more therapeutically productive than attempting to argue one side.",
            "The fourth principle is that direct confrontation and argumentation are counterproductive for promoting behavior change. As Patterson and Forgatch (1985) demonstrated experimentally, confrontation increases client resistance rather than compliance. The more you push, the more they push back - a phenomenon called psychological reactance.",
        ],
        "Motivation is dynamic, relationships heal, ambivalence is normal, and confrontation backfires - these are MET's axioms.",
        "Patterson & Forgatch (1985). JCCP, 53, 846. Confrontation effects on resistance.",
        notes="These four principles should be internalized before learning any specific MI technique. They represent a paradigm shift from the dominant treatment culture of confrontation and compliance. Many trainees intellectually accept these principles but struggle to implement them when faced with resistant or ambivalent clients in practice."
    )

    content_slide(prs,
        "Who Benefits Most from MET? Client Characteristics",
        [
            "Research has identified several client characteristics that may predict differential benefit from MET compared to other approaches. Project MATCH (1997) found that clients high in anger showed significantly better outcomes in MET than in CBT or TSF, suggesting that angry clients respond particularly well to the non-confrontational approach.",
            "Clients low in readiness to change (precontemplation and early contemplation) may benefit more from MET than from action-oriented treatments that assume readiness. Attempting cognitive-behavioral skill training with an unmotivated client is like teaching someone to swim who doesn't want to get in the water.",
            "Clients with high psychological reactance (tendency to resist perceived threats to freedom) are particularly well-served by MET's emphasis on personal choice and autonomy support. For these clients, any approach that tells them what to do or prescribes specific behaviors will be met with opposition.",
            "Conversely, clients who are already highly motivated and seeking specific guidance may prefer and benefit more from directive, skills-based approaches. MET is not universally superior - it is particularly suited for ambivalent, pre-contemplative, angry, or reactant clients who are not yet ready for action-oriented treatment.",
        ],
        "MET may be particularly effective for angry, ambivalent, reactant, and pre-contemplative clients who resist directive approaches.",
        "Project MATCH (1997). JOSA, 58, 7-29. Matching findings for angry clients.",
        notes="Clinical application: screen for readiness, anger, and reactance. Clients scoring high on these dimensions should be offered MET as a first-line intervention. Clients already in preparation or action stages may benefit from adding skills-based components (CBT, RP) to the motivational foundation."
    )

    # =====================================================================
    # Additional slides for Theoretical Foundations
    # =====================================================================
    content_slide(prs,
        "Psychological Reactance Theory and MET",
        [
            "Psychological reactance theory (Brehm, 1966) explains why direct persuasion attempts often backfire in clinical settings. When individuals perceive that their freedom of choice is being threatened, they experience a motivational state (reactance) that drives them to restore their freedom by doing the opposite of what is being demanded.",
            "This mechanism explains the common clinical observation that the harder a therapist pushes for change, the more resistant the client becomes. It is not that the client is 'in denial' or 'unmotivated' - they are experiencing a natural psychological response to perceived coercion. The confrontation-resistance cycle is maintained by the therapist's behavior.",
            "MET resolves this by explicitly supporting autonomy: 'It's entirely your choice what you do with this information.' 'No one can make you change if you don't want to.' These statements paradoxically INCREASE motivation because they remove the threat to freedom that was triggering reactance.",
            "Miller and Rollnick have consistently emphasized that MI was developed specifically as an alternative to confrontational approaches that trigger reactance. When a client says 'Don't tell me what to do,' the MI response ('I wouldn't dream of it - only you can decide') de-activates reactance and opens space for genuine exploration.",
        ],
        "Reactance theory explains why pushing harder backfires - supporting autonomy paradoxically increases motivation by removing perceived threats to freedom.",
        "Brehm, J.W. (1966). Psychological Reactance Theory; Miller & Rollnick (2023). MI 4th Ed.",
        notes="Clinical pearl: when a client becomes more resistant, check whether you've been pushing too hard. The moment you find yourself arguing with a client about whether they have a problem or need to change, you have left MI territory and entered the confrontation-resistance cycle. Back off, reflect, emphasize autonomy."
    )

    content_slide(prs,
        "Social Learning Theory and Modeling in MET",
        [
            "Bandura's (1977) Social Learning Theory extends beyond self-efficacy to encompass the role of observational learning and social modeling in behavior change. In MET, the therapist models a way of thinking about behavior change that the client can internalize - curious, non-judgmental, balanced, and optimistic.",
            "The therapist also models self-regulation skills implicitly: noticing discrepancies between values and behavior, weighing pros and cons thoughtfully, making decisions based on personal values, and planning specific action steps. By observing this process with the therapist, clients learn to conduct similar self-examination independently.",
            "Vicarious self-efficacy can be enhanced through normalization: 'Many people I've worked with who were in similar situations have found that...' or 'It's common for people to feel stuck at this point. Others have found that...' These statements provide hope through social comparison without minimizing the individual's struggle.",
            "In group MET formats (PMC, 2024 protocol for opioid dependence), the social learning function is amplified: group members at different stages of change provide natural models of the change process, and hearing peers articulate change talk can be more powerful than hearing it from a therapist.",
        ],
        "MET therapists model a way of thinking about change - clients internalize this process for independent self-regulation.",
        "Bandura, A. (1977). Social Learning Theory; PMC (2024). Group MET protocol.",
        notes="The therapist's modeling function is often unconscious but powerful. When you demonstrate curious, non-judgmental self-reflection in session ('Hmm, let me think about that differently'), you are teaching the client that it's possible to examine oneself without self-attack. This implicit curriculum may be as important as explicit MI techniques."
    )

    # =====================================================================
    # Additional detailed clinical slides
    # =====================================================================
    content_slide(prs,
        "The Assessment-Feedback-Action Loop in MET",
        [
            "MET follows a distinctive clinical logic: comprehensive assessment generates data, data is transformed into personalized feedback, feedback creates discrepancy and motivation, motivation leads to action planning, and action outcomes are assessed in subsequent sessions. This creates an iterative loop of assessment-feedback-action.",
            "The assessment phase is not merely diagnostic - it is designed to gather specific data points that will have motivational impact when fed back. Timeline Follow-Back data becomes a vivid drinking diary. AUDIT scores become normative positioning. Consequence inventories become a personal cost inventory. Each instrument has a feedback function.",
            "The feedback phase transforms raw data into personally meaningful information. A GGT level of 85 U/L means nothing in isolation; when presented as 'Your liver enzyme is nearly double the normal range, which suggests your liver is under significant stress from alcohol processing,' it becomes motivationally relevant.",
            "The action phase follows naturally from internalized motivation - not from therapist prescription. When the client has seen their data, reacted to it, explored its implications for their values, and developed their own conclusions, action planning feels like a natural next step rather than an imposed task.",
        ],
        "MET's assess-feedback-action loop creates a natural progression from data to awareness to motivation to change.",
        "Miller et al. (1995). MET Manual. The clinical logic of structured feedback.",
        notes="This loop can be repeated within a single session (mini-feedback moments) or across the full MET protocol (assessment session → feedback session → action planning → progress review). Understanding the loop helps therapists maintain momentum and direction without being directive."
    )

    content_slide(prs,
        "Handling the Expert Trap in MET",
        [
            "The 'expert trap' occurs when the therapist positions themselves (or is positioned by the client) as the authority who will diagnose the problem and prescribe the solution. This is perhaps the most common violation of MI spirit, particularly for clinicians trained in medical or cognitive-behavioral models where expertise is central to the therapeutic role.",
            "Signs of falling into the expert trap: giving unsolicited advice, speaking more than the client, asking closed diagnostic questions, offering interpretations without invitation, planning treatment unilaterally, and responding to 'What should I do?' with direct prescription rather than exploration.",
            "MI-consistent alternatives: when asked 'What do you think I should do?' respond with 'I'm more interested in what you think might work for you' or 'I have some ideas, but first I'd like to hear your thoughts.' When the urge to advise arises, convert it to a question: instead of 'You should try AA,' ask 'What kind of support, if any, do you think would be helpful?'",
            "When sharing information IS appropriate (and it sometimes is in MET, particularly during feedback), use the Elicit-Provide-Elicit framework: first ask what the client already knows, then share information with permission, then ask for their reaction. This preserves the collaborative stance while allowing information exchange.",
        ],
        "The expert trap - positioning yourself as the one with the answers - is the most common MI spirit violation. Evocation prevents it.",
        "Miller & Rollnick (2023). MI 4th Ed. Avoiding the expert trap.",
        notes="Paradox of expertise in MI: the more you know about addiction, the harder it can be to practice MI because your expertise creates pressure to share it. The MI-skilled expert holds their knowledge lightly, offering it only when invited and always as information for the client to weigh rather than as prescription to follow."
    )



    content_slide(prs,
        "The Elicit-Provide-Elicit Framework",
        [
            "The Elicit-Provide-Elicit (E-P-E) framework is the MI-consistent method for sharing information, advice, or feedback when it is clinically appropriate. It preserves the collaborative stance while allowing the therapist to contribute their expertise without falling into the expert trap.",
            "ELICIT (first): Before sharing information, ask what the client already knows or what they are curious about. 'What do you know about how alcohol affects sleep?' 'What have you heard about medications for alcohol dependence?' This respects existing knowledge and identifies gaps.",
            "PROVIDE: Share information with permission, neutrally, and in digestible amounts. 'Would it be okay if I shared what research shows about that?' Present information as data, not prescription: 'Research suggests that...' rather than 'You need to know that...' Keep it brief and relevant to the client's stated concerns.",
            "ELICIT (second): After providing information, immediately return focus to the client: 'What do you make of that?' 'How does that fit with your experience?' 'What questions does that raise for you?' This ensures that information is processed through the client's own values and perspective rather than being passively received.",
        ],
        "E-P-E allows information sharing within MI spirit: ask what they know, share with permission, ask for their reaction.",
        "Miller & Rollnick (2023). MI 4th Ed. Information exchange in MI.",
        notes="E-P-E is particularly important during the feedback session in MET. Every piece of PFR data should be delivered using this framework: 'What did you expect your drinking level to be compared to others your age? [Elicit] Here's where you actually fall - above 92% of adults. [Provide] What's your reaction to seeing that? [Elicit]'"
    )

    content_slide(prs,
        "Summaries in MET: Types and Strategic Use",
        [
            "Summaries are the 'S' in OARS and serve crucial functions in MET sessions: they demonstrate that the therapist has been listening, they organize complex material, they can strategically highlight change talk over sustain talk, and they create transition points between topics or session segments.",
            "Collecting summaries gather material that has been offered over a period of time: 'Let me see if I've got this right. You started drinking in college, it escalated during your divorce, you've noticed it affecting your health and your relationship with your children, and part of you wants things to be different.' This validates and organizes.",
            "Linking summaries connect current statements to previous ones: 'That reminds me of what you said earlier about wanting to be a better father. It sounds like that value keeps coming up as important to you.' This reinforces change talk themes across the session.",
            "Transitional summaries signal movement to a new topic or phase: 'So we've talked about what brought you here and what your drinking looks like. If it's okay with you, I'd like to show you some of the assessment results.' This creates smooth shifts while maintaining the client's sense of control over the process.",
        ],
        "Strategic summaries organize information, demonstrate empathy, reinforce change talk themes, and create smooth transitions.",
        "Miller & Rollnick (2023). MI 4th Ed. The art of summarizing.",
        notes="Advanced summary skill: the 'bouquet' summary that collects all change talk from a session and presents it back as a coherent narrative of motivation. 'Let me pull together what you've said today: you want to be healthier, you're concerned about your marriage, you know you can make changes because you've done hard things before, and you're thinking about cutting back. That's a lot of reasons and confidence.' This is powerfully motivating."
    )

    table_slide(prs,
        "MI-Consistent vs. MI-Inconsistent Therapist Behaviors",
        ["MI-Consistent (DO)", "MI-Inconsistent (AVOID)", "Impact on Change Talk"],
        [
            ["Open questions", "Closed/leading questions", "Open → + change talk"],
            ["Complex reflections", "Simple parroting only", "Complex → deeper exploration"],
            ["Affirmations", "Praise (evaluative)", "Affirm → + self-efficacy"],
            ["Asking permission", "Giving unsolicited advice", "Permission → autonomy preserved"],
            ["Emphasizing autonomy", "Directing/demanding", "Autonomy → - reactance"],
            ["Reflecting change talk", "Arguing with sustain talk", "Reflect CT → + CT"],
            ["Evoking client views", "Providing expert opinion", "Evoke → internal motivation"],
        ],
        "Every therapist behavior either promotes or inhibits change talk - awareness of this guides moment-to-moment clinical choices.",
        "Moyers et al. (2007). JCCP, 75, 790. MITI behavioral codes and outcomes.",
        notes="Use this table as a quick reference during supervision. When reviewing a session, categorize each therapist response as MI-consistent or MI-inconsistent, then examine what the client said immediately after. The pattern becomes obvious: consistent behaviors pull for change talk, inconsistent behaviors push toward sustain talk."
    )

    content_slide(prs,
        "The Importance-Confidence Matrix in MET",
        [
            "The importance-confidence matrix is a simple but powerful assessment tool that distinguishes between two dimensions of motivation: (1) How important is change to the client? and (2) How confident are they that they can change? These require different therapeutic responses.",
            "High importance + High confidence = Client is ready for action. Move to planning. Ask: 'What's your next step?' 'When will you start?' Do not over-process with someone who is ready to go. High importance + Low confidence = Client wants to change but doesn't believe they can. Focus on self-efficacy building: past successes, strengths, small achievable steps.",
            "Low importance + High confidence = Client could change but doesn't see why they should. Focus on discrepancy development: values exploration, personalized feedback, exploring concerns. This is the classic precontemplator who says 'I could stop anytime, I just don't want to.'",
            "Low importance + Low confidence = Client neither wants to nor believes they can change. Most challenging. Start with engagement, find any area of concern (however small), and build from there. Consider whether the therapeutic relationship itself might gradually shift both dimensions over time.",
        ],
        "Importance and confidence require different interventions - always assess both before choosing your therapeutic strategy.",
        "Miller & Rollnick (2023). MI 4th Ed. Importance and confidence as separate dimensions.",
        notes="Quick assessment in session: 'On a scale of 0-10, how important is it to you to make this change? ... And on that same scale, how confident are you that you could do it if you decided to?' These two numbers instantly tell you where to focus. Most common error: trying to build confidence in someone who doesn't yet see change as important."
    )

    content_slide(prs,
        "Working with Significant Others in MET",
        [
            "Significant others (spouses, family members, friends) play an important role in MET, both as sources of information for the PFR and as potential supports for change. The MET manual (Miller et al., 1995) includes provision for a concerned significant other (CSO) session to gather additional perspective.",
            "The CSO session follows MI principles: the significant other is treated with the same empathy and respect as the client. Their concerns are heard, their perspective is validated, and they are invited to share observations that might contribute to the PFR. They are NOT recruited as allies in confronting the client.",
            "In Indian family-centered culture, involving family members may be expected and appropriate. Patel et al. (2024) demonstrated that Behavioral Couples Therapy combined with MI can address relationship dynamics alongside substance use. Family involvement should be guided by the client's preference and comfort.",
            "Boundaries with significant others: information from CSO sessions is used in the PFR only with the client's knowledge and consent. The CSO is not given clinical information about the client without permission. The primary therapeutic relationship remains with the client, and their autonomy is never undermined by collusion with family members.",
        ],
        "Significant others can enhance MET through additional feedback data and change support - but the client's autonomy always takes priority.",
        "Miller et al. (1995). MET Manual. CSO involvement protocol; Patel et al. (2024). MI+BCT.",
        notes="Cultural consideration: In India, family members may attend sessions uninvited or may pressure the therapist to 'make him stop drinking.' The MI-consistent response: 'I can see how much you care about him. In our approach, we help people find their own motivation for change rather than telling them what to do. Would you be willing to share your observations that I can use in our feedback session?'"
    )

    content_slide(prs,
        "Exploring Values: The Card Sort Technique",
        [
            "The values card sort is a structured technique for identifying client values that will inform discrepancy development throughout MET. The client is given a set of cards, each naming a value (family, health, success, independence, spirituality, etc.) and asked to sort them into categories: very important, somewhat important, and not important.",
            "Once sorted, the therapist explores the top 5-10 values: 'Tell me about why family is so important to you.' 'What does being a good provider mean to you?' 'How does health connect to your other goals?' This exploration deepens understanding and creates rich material for feedback delivery.",
            "The strategic connection: after values are explored, the therapist can reference them during feedback. 'You mentioned that your children's respect is one of your highest values. When you look at these consequences of your drinking that you endorsed - missing your daughter's concert, your son seeing you intoxicated - how does that fit with that value?'",
            "For Indian clients, values cards should include culturally relevant options: dharma/duty, izzat/family honor, service to elders, community standing, religious devotion, economic stability for extended family, children's education, and maintaining cultural traditions. These may resonate more than Western-centric options.",
        ],
        "Values card sorts provide the motivational target for discrepancy development - culturally adapted values enhance personal relevance.",
        "Miller et al. (1995). MET Manual. Values assessment methods.",
        notes="Practical tip: Create your own values card set appropriate for your cultural context. Include blank cards so clients can write their own values. The process of sorting and discussing values is often experienced as affirming and respectful - quite different from the usual clinical focus on problems and deficits."
    )

    content_slide(prs,
        "Managing Discord in the Therapeutic Relationship",
        [
            "Discord - a disruption in the therapeutic alliance signaled by argumentativeness, disengagement, or hostility - is a critical clinical signal in MET. Unlike sustain talk (which is about ambivalence regarding change), discord is about the relationship between therapist and client. It indicates that the therapist needs to adjust.",
            "Common therapist behaviors that trigger discord: pushing too hard for commitment before the client is ready, reflecting inaccurately (the client feels misunderstood), providing feedback without permission, making assumptions about the client's experience, and violating cultural norms of respect or communication.",
            "Repair strategies include: backing off from the current direction ('I may have gotten ahead of where you are'), apologizing for missteps ('I apologize - I shouldn't have assumed'), shifting to listening mode (increase reflections, decrease questions), emphasizing autonomy ('You're the one who gets to decide about all of this'), and simply pausing to ask how the client is feeling about the conversation.",
            "The skilled MET therapist uses discord as information rather than experiencing it as failure. Discord tells you something important about what the client needs right now - usually more empathy, more space, more autonomy support, or a different topic entirely. It is a compass, not a catastrophe.",
        ],
        "Discord is a relationship signal, not a client problem. It tells the therapist to adjust approach, not push harder.",
        "Miller & Rollnick (2023). MI 4th Ed. Responding to discord; Rogers (1957). Therapeutic conditions.",
        notes="The experienced MI practitioner actually welcomes mild discord because it provides information about what the client needs. The response to discord ('I think I may have gotten off track - let me ask what would be most helpful for you right now') can actually STRENGTHEN the alliance by demonstrating responsiveness and humility."
    )



    content_slide(prs,
        "Biological Markers in MET Feedback",
        [
            "Biological markers provide objective data for the MET Personalized Feedback Report that can have powerful motivational impact. Unlike self-report measures, biological markers cannot be minimized or denied - they provide irrefutable evidence of physical impact that can create significant discrepancy with health values.",
            "Common markers used in MET feedback include: Gamma-Glutamyl Transferase (GGT) - elevated levels indicate liver stress from alcohol processing; Mean Corpuscular Volume (MCV) - elevated levels suggest chronic heavy drinking affecting red blood cell formation; Liver enzymes (AST, ALT) - indicators of liver damage; and Carbohydrate-Deficient Transferrin (CDT) - a specific marker for recent heavy drinking.",
            "Delivery of biological results follows MI principles: 'Your GGT level came back at 85 - the normal range is under 50. This suggests your liver is working harder than it should to process alcohol. What's your reaction to hearing that?' The therapist provides context without catastrophizing or diagnosing.",
            "In Indian primary care settings, basic liver function tests may be more readily available than specialized markers like CDT. Even a simple GGT and liver enzyme panel can provide meaningful feedback data. Physical examination findings (hepatomegaly, tremor, peripheral neuropathy) can also be incorporated into the PFR when relevant.",
        ],
        "Biological markers provide undeniable physical evidence that can create powerful discrepancy - delivered with empathy, not as punishment.",
        "Miller et al. (1995). MET Manual. Biological feedback component.",
        notes="Important: never use biological results to scare or confront. 'Your liver is going to fail' is confrontation. 'Your liver is showing signs of stress, and the good news is these markers typically improve when drinking is reduced' is MI-consistent. Always pair concerning information with hope for reversibility and the client's capacity for change."
    )

    content_slide(prs,
        "MET in Criminal Justice Settings",
        [
            "Criminal justice settings present unique challenges for MET delivery: clients are mandated and often resentful, confidentiality has limits (reports to probation/parole), and the power differential between therapist and client is amplified by the institutional context. Yet research suggests MI/MET is well-suited for this population.",
            "Project MATCH (1997) finding that angry clients responded particularly well to MET is directly relevant: criminal justice clients often present with anger, defensiveness, and perceived coercion. MET's non-confrontational approach can de-escalate these emotional states and create unexpected engagement.",
            "Practical adaptations: be completely transparent about reporting requirements and confidentiality limits BEFORE beginning clinical work. Acknowledge the coercive context honestly: 'You're here because you have to be, and I respect your honesty about that. Within that constraint, how would you like to use our time together?'",
            "Find areas of genuine intrinsic motivation: even within a mandated context, most clients have some personal goals (maintaining custody of children, keeping a job, reducing legal involvement) that align with behavior change. These autonomous motivators are more powerful than the mandate itself for sustaining change.",
        ],
        "Mandated clients respond well to MET because it respects their autonomy within an inherently coercive context.",
        "Project MATCH (1997). Angry clients benefited differentially from MET.",
        notes="Common mistake in criminal justice MI: pretending the coercion doesn't exist. Authenticity requires acknowledging reality: 'Neither of us chose for you to be here. The court made that decision. But since you are here, we have some options about how to use this time.' This honest acknowledgment paradoxically increases engagement because the client doesn't have to maintain a facade."
    )

    content_slide(prs,
        "MET and Harm Reduction: Compatible Philosophies",
        [
            "Harm reduction - the philosophical position that any reduction in harm is worthwhile, even without complete abstinence - is highly compatible with MET's non-judgmental, autonomy-supporting approach. MET does not require abstinence as the only acceptable goal and can support clients in pursuing whatever change they choose.",
            "This compatibility extends to practical integration: MET can address ambivalence about naloxone carrying (for opioid users), safer injection practices, moderation rather than abstinence goals, substitution (switching from spirits to beer), reducing use frequency, or any other harm-reducing behavior the client values.",
            "The philosophical alignment: both harm reduction and MI/MET reject the idea that a single endpoint (complete abstinence) must be accepted before any help is available. Both meet people where they are. Both prioritize engagement over requirements. Both trust that people are capable of making informed decisions about their own lives.",
            "In India, harm reduction approaches to alcohol may be particularly relevant given that many clients will not accept abstinence as a goal. Cultural drinking patterns, social expectations, and the absence of strong 'recovery culture' mean that moderation goals may be more culturally acceptable and therefore more likely to be pursued.",
        ],
        "MET and harm reduction share core philosophy: meet people where they are, respect their choices, and celebrate any positive movement.",
        "Marlatt & Gordon (1985). Harm reduction principles; Miller & Rollnick (2023). Goal flexibility in MI.",
        notes="This can be controversial in settings that require abstinence. The MI response: research shows that clients who choose their own goals show better outcomes than those with imposed goals. If a client eventually needs abstinence, they may discover this through experience with moderation attempts. Respect the learning process."
    )

    content_slide(prs,
        "Supervision Models for MET Practice",
        [
            "Effective supervision is essential for developing and maintaining MI/MET competence. Research indicates that without ongoing supervision, MI skills degrade over time even after excellent initial training. Several supervision models have demonstrated effectiveness.",
            "Individual supervision with recorded session review is the gold standard. The supervisor listens to or watches a session recording, provides MITI-based behavioral feedback (specific ratios, examples of MI-consistent and inconsistent moments), and helps the supervisee identify patterns and improvement targets.",
            "Group supervision offers efficiency and peer learning: members present recorded segments, receive feedback from peers and supervisor, practice alternative responses to difficult moments, and share strategies for common challenges. The group format normalizes the struggle of learning MI.",
            "Technology-assisted supervision (phone-based or video-conferenced) expands access, particularly important in India where supervision expertise is geographically concentrated in major cities. The AMBIT trial (Nadkarni et al., 2023) used regular supervision calls to maintain lay counselor fidelity in Goa.",
        ],
        "MI competence requires ongoing supervision with recorded session feedback - workshop training alone is insufficient for lasting skill development.",
        "Nadkarni et al. (2023). AMBIT supervision model; Moyers et al. (2007). MITI for supervision.",
        notes="Budget recommendation for organizations implementing MI/MET: allocate at least 30% of training resources to ongoing supervision rather than putting everything into initial workshops. The workshop creates awareness and enthusiasm; supervision creates competence. Without the latter, the former fades within 3-6 months."
    )

    content_slide(prs,
        "Common Traps and How to Avoid Them",
        [
            "Thomas Gordon originally identified 'communication roadblocks' that MI has adopted as common traps to avoid. These are natural counselor responses that inadvertently undermine the MI process by taking the client's work away from them or triggering reactance.",
            "The Assessment Trap: Asking so many questions that the interaction feels like an interrogation rather than a conversation. Solution: balance questions with reflections (2:1 ratio minimum). Convert many questions into reflections: 'How did that make you feel?' becomes 'That sounds like it was difficult for you.'",
            "The Premature Focus Trap: Jumping directly to the substance use issue before establishing rapport or understanding the client's priorities. Solution: ask what the client wants to discuss first. 'What would be most helpful to talk about today?' Their priority may differ from yours.",
            "The Labeling Trap: Insisting on diagnostic labels ('You're an alcoholic') that trigger defensiveness. Solution: focus on behaviors and consequences rather than labels. 'You're drinking more than you're comfortable with' is MI-consistent; 'You have alcoholism' is not. Labels are for charts, not conversations.",
        ],
        "Common traps (assessment, premature focus, labeling, expert) derail MI process - awareness and practice prevent them.",
        "Miller & Rollnick (2023). MI 4th Ed. Communication traps and avoiding them.",
        notes="Self-monitoring exercise: after each session, ask yourself: Did I fall into any traps today? Which one? What triggered it? What could I do differently next time? This honest self-reflection, perhaps shared in supervision, accelerates learning. Nobody practices MI perfectly - the goal is progressive improvement and self-awareness."
    )

    two_col_slide(prs,
        "MET Session 1: Structured Feedback Sequence",
        "First Half (30 minutes)",
        [
            "Welcome and agenda setting (5 min)",
            "Check-in about assessment experience",
            "Ask permission to share feedback",
            "Present drinking summary data",
            "Show normative comparison (graph)",
            "Elicit reaction to each data point",
            "Reflect responses empathically",
        ],
        "Second Half (30 minutes)",
        [
            "Present consequence inventory results",
            "Connect to stated values",
            "Share biological markers if available",
            "Explore risk factors (family history)",
            "Provide summary of key themes",
            "Ask: 'Where does this leave you?'",
            "Schedule Session 2, provide PFR copy",
        ],
        "Session 1 is structured but flexible - the client's reactions determine depth and pacing of feedback delivery.",
        "Miller et al. (1995). MET Manual. Session 1 detailed protocol.",
        notes="Key judgment calls in Session 1: How much feedback to present (some clients need all sections; others are deeply moved by one or two). How long to spend on each section (follow the client's emotional energy). When to stop and explore vs. when to move on. Always prioritize depth over breadth."
    )

    content_slide(prs,
        "The Readiness Ruler in Clinical Practice",
        [
            "Readiness rulers are simple visual or verbal scales (typically 0-10) used to assess motivation, importance, confidence, and readiness for change. Despite their simplicity, they are among the most clinically useful tools in MET practice because of the strategic follow-up questions they enable.",
            "Standard questions: 'On a scale of 0-10, where 0 is not at all important and 10 is the most important thing in your life, how important is it to you to make a change in your drinking?' Followed by: 'And on that same scale, how confident are you that you could make that change if you decided to?'",
            "The STRATEGIC follow-up: After the client gives a number (say, 6 for importance), ask: 'Why are you at a 6 and not a 3?' This forces the client to argue FOR importance (explaining why it's not lower) and reliably produces change talk. NEVER ask 'Why not a 10?' - this forces them to argue AGAINST importance.",
            "Tracking rulers across sessions shows motivational trajectory: are importance and confidence increasing over time? Rulers provide both clinical data and therapeutic conversation starters. They can be used at the beginning of each session as a check-in: 'Last time you were at a 6 for importance and 4 for confidence. Where are you today?'",
        ],
        "Readiness rulers are deceptively simple tools - their power lies in the strategic follow-up that reliably elicits change talk.",
        "Miller & Rollnick (2023). MI 4th Ed. Ruler techniques and strategic follow-up.",
        notes="This is one of the easiest MI techniques to teach and implement. Even clinicians with minimal MI training can use rulers effectively. The follow-up question ('why not lower?') is the key innovation - it reverses the natural tendency to ask about barriers ('why not higher?') and instead directs attention to existing motivation."
    )

    content_slide(prs,
        "Affirmations: Beyond Simple Praise",
        [
            "Affirmations in MI/MET are distinct from praise. Praise is evaluative ('Good job!' 'I'm proud of you!') and positions the therapist as the judge of the client's worth. Affirmation is descriptive and recognizes inherent qualities: 'You showed real courage in coming here today' or 'That took a lot of honesty to share.'",
            "Effective affirmations target character strengths and values rather than specific behaviors: 'You clearly care deeply about your family' (character) vs. 'Good job staying sober this week' (behavior/praise). Character affirmations build identity-based motivation: 'I'm the kind of person who cares about family, therefore I will act accordingly.'",
            "Timing matters: affirmations are most powerful when unexpected and when the client has just expressed vulnerability, effort, or change talk. They should feel earned and genuine rather than formulaic. Indiscriminate affirmation ('Everything you say is wonderful!') is patronizing and therapeutically empty.",
            "In Indian cultural contexts, affirmations can reference culturally valued qualities: dedication to family, hard work, respect for elders, spiritual strength, community contribution, perseverance through adversity. These culturally resonant affirmations may carry more weight than generic compliments.",
        ],
        "True affirmations recognize character strengths and values - they build identity-based motivation, not dependent behavior.",
        "Miller & Rollnick (2023). MI 4th Ed. The art of affirmation.",
        notes="Practice exercise: write 10 affirmations that could apply to a client who is ambivalent about changing their drinking. Make sure each one is descriptive (not evaluative), targets character (not behavior), and would feel genuine to the client. This is harder than it sounds - we're trained to evaluate rather than describe."
    )

    content_slide(prs,
        "MET and Mindfulness Integration",
        [
            "Emerging research explores integrating mindfulness-based practices with MET, leveraging the shared emphasis on non-judgmental awareness, present-moment focus, and acceptance. Mindfulness can enhance MET by helping clients observe their cravings and urges without automatically acting on them.",
            "In the context of MET Sessions 3-4, mindfulness techniques can be introduced as coping strategies within the MI framework: 'Some people find that practicing awareness of their urges - noticing them without acting on them - helps them make more deliberate choices. Is that something you'd be interested in exploring?'",
            "The compatibility is philosophical: both MI/MET and mindfulness traditions respect the individual's experience, avoid judgmental labeling, and trust in the person's inherent capacity for wisdom and growth. Both approaches are fundamentally about awareness - MI about awareness of discrepancy, mindfulness about awareness of present-moment experience.",
            "In Indian contexts, mindfulness has deep cultural roots through meditation traditions (dhyana, vipassana). Connecting MET with familiar contemplative practices can enhance cultural resonance and provide clients with additional tools for managing cravings and maintaining awareness of their values during high-risk situations.",
        ],
        "Mindfulness and MET share philosophical foundations of non-judgment and awareness - integration enhances both approaches.",
        "Miller & Rollnick (2023). MI integration with complementary approaches; Indian meditation traditions.",
        notes="This integration works best when introduced organically in response to the client's expressed interest rather than as a prescribed homework assignment. 'You mentioned wanting to be more aware of when you're reaching for a drink automatically. Would you like to explore some techniques for building that awareness?'"
    )

    content_slide(prs,
        "Outcome Monitoring in MET: Tracking Progress",
        [
            "Systematic outcome monitoring is essential for effective MET practice, providing both clinician and client with data on whether the intervention is producing the desired effects. Outcome monitoring can be done through repeated assessment using the same measures administered at baseline.",
            "Brief weekly or session-to-session measures include: number of standard drinks per week (self-monitoring diary), AUDIT-C (3-item screen for current drinking), readiness rulers (importance and confidence scores), and goal attainment scaling (progress toward self-defined goals from the Change Plan).",
            "Outcome data can be incorporated into subsequent MET sessions as additional personalized feedback: 'At your first session, you were drinking 35 standard drinks per week. Over the past month, your diary shows an average of 18. That's nearly a 50% reduction. What do you make of that?' This reinforces progress and builds self-efficacy.",
            "If outcomes show deterioration or no progress, this too is explored non-judgmentally: 'Your drinking has stayed about the same over the past month. Tell me about what's been happening.' The therapist neither expresses disappointment nor minimizes - curiosity guides the exploration of what's interfering with the client's stated goals.",
        ],
        "Outcome monitoring provides ongoing personalized feedback and allows responsive adjustment of the MET approach across sessions.",
        "Miller et al. (1995). MET Manual. Progress monitoring and session adjustment.",
        notes="Practical tip: create a simple graph that shows the client's drinking trajectory across sessions. Visual representation of progress (even modest progress) is motivating and makes the abstract concept of 'improvement' concrete. Keep the graph in the client's file and show it at the beginning of each session."
    )

    content_slide(prs,
        "Termination and Follow-Up in MET",
        [
            "MET's brief structure means that termination planning is built into the protocol from the beginning. Session 4 explicitly addresses consolidation and closure, helping clients prepare for continued self-management without ongoing therapeutic support. The therapist communicates confidence in the client's ability to maintain gains independently.",
            "Key termination tasks include: reviewing the entire MET journey (from initial assessment through current progress), affirming the client's growth and efforts, updating the Change Plan based on lessons learned, identifying ongoing support resources (mutual aid groups, social support, community resources), and establishing conditions for return if needed.",
            "The MI-consistent termination message: 'You've done significant work here - this change came from you, your values, and your determination. I'm confident in your ability to continue this path. If challenges arise in the future, my door is always open - but I don't think you'll need it as much as you might fear.' This affirms self-efficacy while maintaining the safety net.",
            "Follow-up options in MET: some protocols include brief phone check-ins at 3 and 6 months post-treatment to reinforce gains and catch early signs of relapse. These brief contacts (5-10 minutes) maintain the therapeutic connection without creating dependence on ongoing treatment.",
        ],
        "MET termination affirms client capability and independence while leaving the door open - the goal is autonomy, not dependency.",
        "Miller et al. (1995). MET Manual. Session 4: Consolidation and termination.",
        notes="The spirit of termination in MET: celebrate the client's autonomy and self-direction. They are not 'graduating' from your program - they are continuing a journey that was always theirs. Your role was temporary catalyst, not ongoing support system. This framing prevents dependency and reinforces internal attribution of change."
    )

    content_slide(prs,
        "The Ambivalence Model: Normal, Not Pathological",
        [
            "A paradigm-shifting insight of MI/MET is that ambivalence about change is NORMAL, not pathological. Virtually everyone considering significant behavior change simultaneously holds reasons for and against change. This is not 'denial,' 'resistance,' or 'lack of motivation' - it is the ordinary human experience of being of two minds.",
            "The clinical implication: if ambivalence is normal, then the therapist's task is not to eliminate it (impossible) or overwhelm it with arguments (counterproductive) but to explore it fully and allow the client to resolve it in their own way and time. This requires patience that many clinicians find challenging.",
            "Ambivalence can be visualized as a seesaw: on one side, reasons to maintain current behavior; on the other, reasons to change. The therapist's role is to help the client see both sides clearly, not to add weight to one side. Paradoxically, when both sides are fully acknowledged, clients typically find their own way forward.",
            "Research supports this approach: Magill et al. (2018) confirmed that the ratio of change talk to sustain talk (not the elimination of sustain talk) predicts outcomes. Some sustain talk is expected and normal even in clients who ultimately make significant changes. It is the BALANCE that matters.",
        ],
        "Ambivalence is the normal human response to change - acknowledging both sides fully allows natural resolution.",
        "Magill et al. (2018). Updated MI meta-analysis; Miller & Rollnick (2023). MI 4th Ed.",
        notes="This is perhaps the most liberating insight for new MI practitioners: you don't have to 'fix' the client's ambivalence or make them certain about change. You only need to create conditions where they can fully explore both sides. Human beings naturally move toward resolution when given space and empathy. Trust the process."
    )

    return prs




def main():
    """Main entry point."""
    try:
        print("Building MET Comprehensive Presentation...")
        print("Creating slides...")
        prs = build_presentation()
        print(f"  Main presentation built: {len(prs.slides)} slides")

        print("Adding supplementary slides...")
        prs = add_supplementary_slides(prs)
        print(f"  Total slides after supplement: {len(prs.slides)} slides")

        print(f"Saving to {OUTPUT_PATH}...")
        prs.save(OUTPUT_PATH)
        print(f"SUCCESS: Presentation saved with {len(prs.slides)} slides.")
        print(f"File: {OUTPUT_PATH}")
    except Exception as e:
        print(f"ERROR: {e}")
        traceback.print_exc()
        raise


if __name__ == "__main__":
    main()
