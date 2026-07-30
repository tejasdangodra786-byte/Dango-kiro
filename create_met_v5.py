#!/usr/bin/env python3
"""
MET Comprehensive 80-Slide Presentation Generator
Based on: MET Manual (Miller et al., 1992) & NIMHANS SUD Manual (2016)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# === CONSTANTS ===
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "Times New Roman"
TITLE_FONT_SIZE = Pt(36)
CONTENT_FONT_SIZE = Pt(14)
REF_FONT_SIZE = Pt(9)

# Layout zones
TITLE_Y = Inches(0)
TITLE_HEIGHT = Inches(1.0)
CONTENT_Y = Inches(1.2)
CONTENT_BOTTOM = Inches(6.3)
REF_Y = Inches(6.5)
REF_HEIGHT = Inches(0.5)
LEFT_MARGIN = Inches(0.5)
FULL_WIDTH = Inches(12.333)

# Two-column layout
COL1_X = Inches(0.5)
COL1_WIDTH = Inches(5.9)
COL2_X = Inches(6.8)
COL2_WIDTH = Inches(5.9)

# Two-row layout
ROW1_Y = Inches(1.2)
ROW1_HEIGHT = Inches(2.3)
ROW2_Y = Inches(3.7)
ROW2_HEIGHT = Inches(2.5)


# === COLOR SCHEMES ===
COLORS = {
    'deep_blue': RGBColor(0x1B, 0x3A, 0x5C),
    'light_blue': RGBColor(0xD6, 0xEA, 0xF8),
    'teal': RGBColor(0x00, 0x7B, 0x83),
    'light_teal': RGBColor(0xD0, 0xF0, 0xF2),
    'green': RGBColor(0x2E, 0x7D, 0x32),
    'light_green': RGBColor(0xE8, 0xF5, 0xE9),
    'purple': RGBColor(0x6A, 0x1B, 0x9A),
    'light_purple': RGBColor(0xF3, 0xE5, 0xF5),
    'maroon': RGBColor(0x88, 0x00, 0x38),
    'light_red': RGBColor(0xFC, 0xE4, 0xEC),
    'navy': RGBColor(0x0D, 0x47, 0xA1),
    'orange': RGBColor(0xE6, 0x5C, 0x00),
    'light_orange': RGBColor(0xFF, 0xF3, 0xE0),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'black': RGBColor(0x00, 0x00, 0x00),
    'gold': RGBColor(0xFF, 0xD7, 0x00),
    'cream': RGBColor(0xFF, 0xFD, 0xF0),
    'light_gray': RGBColor(0xF0, 0xF0, 0xF0),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'medium_gray': RGBColor(0x66, 0x66, 0x66),
}

SECTION_COLORS = {
    1: ('deep_blue', 'light_blue'),
    2: ('teal', 'light_teal'),
    3: ('green', 'light_green'),
    4: ('purple', 'light_purple'),
    5: ('deep_blue', 'light_blue'),
    6: ('maroon', 'light_red'),
    7: ('navy', 'light_blue'),
    8: ('orange', 'light_orange'),
    9: ('purple', 'light_purple'),
    10: ('navy', 'light_blue'),
}


# === HELPER FUNCTIONS ===
def set_shape_fill(shape, color):
    """Set solid fill color on a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, section_num):
    """Add colored title bar at top of slide."""
    color_name = SECTION_COLORS[section_num][0]
    color = COLORS[color_name]
    
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), TITLE_Y, SLIDE_WIDTH, TITLE_HEIGHT
    )
    set_shape_fill(title_shape, color)
    title_shape.line.fill.background()
    
    tf = title_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "   " + title_text
    run.font.name = FONT_NAME
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

def add_reference_bar(slide, ref_text):
    """Add reference bar at bottom."""
    ref_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), REF_Y, SLIDE_WIDTH, REF_HEIGHT
    )
    set_shape_fill(ref_shape, COLORS['light_gray'])
    ref_shape.line.fill.background()
    
    tf = ref_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "   " + ref_text
    run.font.name = FONT_NAME
    run.font.size = REF_FONT_SIZE
    run.font.italic = True
    run.font.color.rgb = COLORS['medium_gray']


def add_content_box(slide, x, y, width, height, text_lines, border_color, fill_color, 
                    font_size=CONTENT_FONT_SIZE, bold_first=False, alignment=PP_ALIGN.LEFT):
    """Add a content box with colored border and light fill."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
    )
    set_shape_fill(shape, fill_color)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.name = FONT_NAME
        run.font.size = font_size
        run.font.color.rgb = COLORS['dark_gray']
        if bold_first and i == 0:
            run.font.bold = True
            run.font.color.rgb = border_color

def add_section_divider(slide, section_title, section_num):
    """Add a section divider slide with gradient-like appearance."""
    color_name = SECTION_COLORS[section_num][0]
    color = COLORS[color_name]
    light_name = SECTION_COLORS[section_num][1]
    light_color = COLORS[light_name]
    
    # Large colored background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(bg_shape, color)
    bg_shape.line.fill.background()
    
    # Accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.2), Inches(9.333), Inches(0.08)
    )
    set_shape_fill(accent, COLORS['gold'])
    accent.line.fill.background()
    
    # Title text
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_title
    run.font.name = FONT_NAME
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = " "
    run2.font.size = Pt(14)


def make_standard_slide(prs, title, content_boxes, section_num, reference):
    """Create a standard slide with title bar, content boxes, and reference."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_title_bar(slide, title, section_num)
    
    color_name = SECTION_COLORS[section_num][0]
    light_name = SECTION_COLORS[section_num][1]
    border_color = COLORS[color_name]
    fill_color = COLORS[light_name]
    
    for box in content_boxes:
        x = box.get('x', LEFT_MARGIN)
        y = box.get('y', CONTENT_Y)
        w = box.get('w', FULL_WIDTH)
        h = box.get('h', Inches(4.9))
        lines = box.get('lines', [])
        bc = box.get('border_color', border_color)
        fc = box.get('fill_color', fill_color)
        bf = box.get('bold_first', False)
        fs = box.get('font_size', CONTENT_FONT_SIZE)
        align = box.get('alignment', PP_ALIGN.LEFT)
        add_content_box(slide, x, y, w, h, lines, bc, fc, fs, bf, align)
    
    add_reference_bar(slide, reference)
    return slide


# === PRESENTATION CREATION ===
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ============================================================
# SECTION 1: TITLE & INTRODUCTION (Slides 1-8)
# ============================================================

# --- Slide 1: Title Slide ---
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
set_shape_fill(bg1, COLORS['deep_blue'])
bg1.line.fill.background()

# Gold accent line
accent1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.8), Inches(9.333), Inches(0.06))
set_shape_fill(accent1, COLORS['gold'])
accent1.line.fill.background()

# Title
tb1 = slide1.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.333), Inches(2.0))
tf1 = tb1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run()
r1.text = "Motivational Enhancement Therapy (MET)"
r1.font.name = FONT_NAME
r1.font.size = Pt(42)
r1.font.bold = True
r1.font.color.rgb = COLORS['gold']

p1b = tf1.add_paragraph()
p1b.alignment = PP_ALIGN.CENTER
r1b = p1b.add_run()
r1b.text = "A Comprehensive Clinical Guide"
r1b.font.name = FONT_NAME
r1b.font.size = Pt(28)
r1b.font.color.rgb = COLORS['white']

# Subtitle info
tb1s = slide1.shapes.add_textbox(Inches(2), Inches(3.2), Inches(9.333), Inches(3.0))
tf1s = tb1s.text_frame
tf1s.word_wrap = True
lines_s1 = [
    "Based on the MET Manual (Miller, Zweben, DiClemente & Rychtarik, 1992)",
    "NIAAA Project MATCH Monograph Series, Volume 2",
    "& NIMHANS Substance Use Disorders Manual (2016)",
    "",
    "80-Slide Educational Presentation",
    "For Clinical Training & Academic Use"
]
for i, line in enumerate(lines_s1):
    if i == 0:
        p = tf1s.paragraphs[0]
    else:
        p = tf1s.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)
    r = p.add_run()
    r.text = line
    r.font.name = FONT_NAME
    r.font.size = Pt(16)
    r.font.color.rgb = COLORS['light_blue']

# Bottom bar
bot1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), SLIDE_WIDTH, Inches(0.7))
set_shape_fill(bot1, RGBColor(0x0F, 0x27, 0x3F))
bot1.line.fill.background()


# --- Slide 2: Table of Contents ---
toc_lines_left = [
    "TABLE OF CONTENTS",
    "",
    "Section 1: Introduction (Slides 1-8)",
    "Section 2: Theoretical Foundations (Slides 9-18)",
    "Section 3: Principles & Spirit of MI (Slides 19-26)",
    "Section 4: Session-by-Session Guide (Slides 27-38)",
    "Section 5: Clinical Techniques (Slides 39-44)",
]
toc_lines_right = [
    "",
    "",
    "Section 6: Case Conceptualization (Slides 45-50)",
    "Section 7: Worksheets & Tools (Slides 51-56)",
    "Section 8: Comparison with Other Therapies (57-64)",
    "Section 9: Comorbid Disorders (Slides 65-72)",
    "Section 10: Research & Effectiveness (Slides 73-80)",
]
make_standard_slide(prs, "Table of Contents", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': toc_lines_left, 'bold_first': True},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': toc_lines_right, 'bold_first': True},
], 1, "Ref: Miller et al. (1992). MET Manual, NIAAA Project MATCH Monograph Vol. 2")

# --- Slide 3: What is MET? ---
met_def_lines = [
    "WHAT IS MOTIVATIONAL ENHANCEMENT THERAPY?",
    "",
    "MET is a brief, systematic intervention designed to produce rapid,",
    "internally motivated change in persons with substance use problems.",
    "",
    "Key Features:",
    "  * Uses 4 carefully structured sessions (not 12 or more)",
    "  * Does NOT teach skills or give advice directly",
    "  * Instead, it mobilizes the client's OWN internal motivation",
    "  * Uses empathic, non-confrontational strategies",
    "  * Relies on personalized assessment feedback",
    "  * Based on Motivational Interviewing principles (Miller & Rollnick)",
    "",
    "Think of it this way: MET does not push you to change.",
    "It helps YOU find your own reasons to change.",
]
make_standard_slide(prs, "What is MET? - Definition", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': met_def_lines, 'bold_first': True},
], 1, "Ref: Miller et al. (1992). MET Manual, p.1; NIMHANS SUD Manual (2016), Ch. 5")


# --- Slide 4: How MET Differs ---
diff_left = [
    "TRADITIONAL APPROACHES",
    "",
    "* Therapist is the expert",
    "* Confrontation of denial",
    "* Prescribes specific steps",
    "* Labels client (e.g., 'alcoholic')",
    "* Teaches coping skills directly",
    "* Long-term (12+ sessions typical)",
    "* Resistance = client's fault",
    "* External motivation used",
]
diff_right = [
    "MET APPROACH",
    "",
    "* Client is the expert on self",
    "* Empathy and reflection",
    "* Elicits client's own solutions",
    "* Avoids labels entirely",
    "* Mobilizes internal resources",
    "* Brief (4 sessions total)",
    "* Resistance = therapist signal",
    "* Internal motivation activated",
]
make_standard_slide(prs, "MET vs Traditional Approaches", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': diff_left, 'bold_first': True,
     'fill_color': COLORS['light_blue'], 'border_color': COLORS['deep_blue']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': diff_right, 'bold_first': True,
     'fill_color': COLORS['light_green'], 'border_color': COLORS['green']},
], 1, "Ref: Miller et al. (1992). MET Manual, pp. 1-5; Miller & Rollnick (2002)")

# --- Slide 5: Origins ---
origins_lines = [
    "ORIGINS & HISTORY OF MET",
    "",
    "1983: William R. Miller publishes foundational paper on Motivational Interviewing",
    "1991: Miller & Rollnick publish first MI book - clinical framework",
    "1989-1997: Project MATCH - largest psychotherapy trial for alcohol (1,726 clients)",
    "1992: MET Manual developed for Project MATCH (4-session adaptation of MI)",
    "",
    "Project MATCH compared three treatments:",
    "  1. Motivational Enhancement Therapy (MET) - 4 sessions",
    "  2. Cognitive-Behavioral Therapy (CBT) - 12 sessions",
    "  3. Twelve-Step Facilitation (TSF) - 12 sessions",
    "",
    "Key Finding: MET achieved COMPARABLE outcomes in only 4 sessions!",
    "This proved brief motivational approaches can match intensive treatments.",
]
make_standard_slide(prs, "Origins: Project MATCH & MI History", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': origins_lines, 'bold_first': True},
], 1, "Ref: Miller et al. (1992); Project MATCH Research Group (1997); Miller & Rollnick (1991)")


# --- Slide 6: Who is MET For? ---
target_lines = [
    "WHO IS MET DESIGNED FOR?",
    "",
    "Primary Target Populations:",
    "  * Adults with alcohol use disorders (originally developed for)",
    "  * Adults with other substance use disorders (adapted)",
    "  * Dual diagnosis clients (substance + mental health)",
    "  * Individuals ambivalent about change (contemplation stage)",
    "",
    "Settings Where MET is Used:",
    "  * Outpatient addiction clinics",
    "  * Primary healthcare (brief interventions)",
    "  * Emergency departments (screening & brief MET)",
    "  * Community mental health centers",
    "  * Correctional / forensic settings",
    "",
    "MET is especially effective for those NOT yet ready to change.",
]
make_standard_slide(prs, "Who is MET For? - Target Populations", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': target_lines, 'bold_first': True},
], 1, "Ref: Miller et al. (1992), pp. 5-8; NIMHANS SUD Manual (2016), Ch. 5")

# --- Slide 7: The MET Therapist's Role ---
role_lines = [
    "THE MET THERAPIST'S ROLE",
    "",
    "The therapist in MET is NOT a teacher, advisor, or expert.",
    "Instead, the therapist is a skilled GUIDE who:",
    "",
    "  * Expresses genuine empathy (understands without judging)",
    "  * Creates a safe, non-threatening atmosphere",
    "  * Avoids arguments and confrontation at all costs",
    "  * Elicits the client's own concerns and motivations",
    "  * Provides objective personal feedback (from assessments)",
    "  * Highlights discrepancies between goals and current behavior",
    "  * Supports self-efficacy (belief in ability to change)",
    "  * Rolls with resistance rather than opposing it",
    "",
    "The therapist's tone is warm, curious, and collaborative.",
    "Think: 'Dancing with the client' rather than 'wrestling.'",
]
make_standard_slide(prs, "The MET Therapist's Role", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': role_lines, 'bold_first': True},
], 1, "Ref: Miller et al. (1992), pp. 12-18; Miller & Rollnick (2002), Ch. 3")


# --- Slide 8: Overview of 4-Session Structure ---
struct_left = [
    "SESSION 1: Building Motivation",
    "",
    "* Review assessment results",
    "* Provide Personal Feedback Report",
    "* Elicit client's reactions",
    "* Develop discrepancy",
    "* Assess readiness to change",
    "",
    "SESSION 2: Strengthening Commitment",
    "",
    "* Review progress since Session 1",
    "* Reinforce change talk",
    "* Develop written Change Plan",
    "* Consolidate commitment",
]
struct_right = [
    "SESSION 3: Review & Renewal",
    "",
    "* Review progress on Change Plan",
    "* Affirm positive changes",
    "* Address barriers encountered",
    "* Renew or revise commitment",
    "",
    "SESSION 4: Final Review",
    "",
    "* Review overall progress",
    "* Anticipate future challenges",
    "* Handle any relapse/slip",
    "* Plan for continued self-change",
    "* Termination & referrals if needed",
]
make_standard_slide(prs, "Overview of the 4-Session MET Structure", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': struct_left, 'bold_first': True},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': struct_right, 'bold_first': True},
], 1, "Ref: Miller et al. (1992). MET Manual, pp. 29-70 (Session-by-session guide)")


# ============================================================
# SECTION 2: THEORETICAL FOUNDATIONS (Slides 9-18)
# ============================================================

# --- Slide 9: Section Divider ---
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide9, "SECTION 2\nTheoretical Foundations of MET", 2)

# --- Slide 10: Transtheoretical Model ---
ttm_lines = [
    "THE TRANSTHEORETICAL MODEL (Prochaska & DiClemente, 1983)",
    "",
    "This model describes HOW people change. Change is not an event - it is a PROCESS.",
    "People move through predictable stages:",
    "",
    "  1. PRECONTEMPLATION - Not thinking about change (unaware/unwilling)",
    "  2. CONTEMPLATION - Aware of problem, weighing pros & cons",
    "  3. PREPARATION - Decided to change, making plans",
    "  4. ACTION - Actively modifying behavior",
    "  5. MAINTENANCE - Sustaining new behavior over time",
    "  6. RELAPSE - Return to old patterns (normal part of cycle)",
    "",
    "Key Insight for MET: The therapist must MATCH the intervention",
    "to the client's current stage. Pushing action on a precontemplator",
    "creates resistance. MET meets clients WHERE THEY ARE.",
]
make_standard_slide(prs, "Transtheoretical Model - Stages of Change", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': ttm_lines, 'bold_first': True},
], 2, "Ref: Prochaska & DiClemente (1983); Miller et al. (1992), pp. 8-12")

# --- Slide 11: Precontemplation ---
precon_lines = [
    "PRECONTEMPLATION STAGE - 'I Don't Have a Problem'",
    "",
    "Characteristics of the Precontemplator:",
    "  * Does not perceive their substance use as problematic",
    "  * Often enters therapy due to external pressure (family, court, job)",
    "  * May be defensive, minimizing, or rationalizing",
    "  * Unaware of consequences or in active denial",
    "",
    "MET Therapist's Task at This Stage:",
    "  * DO NOT argue or confront - this increases resistance",
    "  * Raise doubt gently by exploring consequences",
    "  * Use personalized feedback to create awareness",
    "  * Plant seeds of concern without demanding change",
    "  * Express empathy and build rapport first",
    "",
    "Goal: Move from 'no concern' to 'some awareness' (raise doubt).",
]
make_standard_slide(prs, "Precontemplation Stage (Detailed)", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': precon_lines, 'bold_first': True},
], 2, "Ref: DiClemente & Prochaska (1998); Miller et al. (1992), p. 10")


# --- Slide 12: Contemplation ---
contemp_lines = [
    "CONTEMPLATION STAGE - 'I Know It's a Problem, But...'",
    "",
    "Characteristics of the Contemplator:",
    "  * Acknowledges the problem exists",
    "  * Weighing pros and cons of change (ambivalence)",
    "  * May have been thinking about change for months/years",
    "  * 'Stuck' - aware but not committed to action",
    "  * Classic statement: 'I want to quit, but I'm not sure I can'",
    "",
    "MET Therapist's Task at This Stage:",
    "  * Tip the decisional balance (help pros of change outweigh cons)",
    "  * Develop discrepancy between values and current behavior",
    "  * Explore the 'good things' and 'less good things' about use",
    "  * Use reflective listening to amplify change talk",
    "  * Normalize ambivalence - it's okay to be unsure",
    "",
    "Goal: Resolve ambivalence in favor of change.",
]
make_standard_slide(prs, "Contemplation Stage (Detailed)", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': contemp_lines, 'bold_first': True},
], 2, "Ref: Miller et al. (1992), pp. 10-11; Miller & Rollnick (2002), Ch. 2")

# --- Slide 13: Preparation ---
prep_lines = [
    "PREPARATION STAGE - 'I'm Going to Do Something About This'",
    "",
    "Characteristics of Someone in Preparation:",
    "  * Has decided to change (tipping point reached)",
    "  * May have already taken small steps (cut down, told someone)",
    "  * Seeks information about HOW to change",
    "  * May set a quit date or ask for help",
    "  * Window of opportunity - relatively brief stage",
    "",
    "MET Therapist's Task at This Stage:",
    "  * Help develop a concrete, written Change Plan",
    "  * Explore options (not prescribe) for change methods",
    "  * Affirm the decision and build confidence",
    "  * Negotiate goals that are realistic and client-chosen",
    "  * Address practical barriers proactively",
    "",
    "Goal: Assist in creating a workable change plan.",
]
make_standard_slide(prs, "Preparation Stage (Detailed)", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': prep_lines, 'bold_first': True},
], 2, "Ref: Miller et al. (1992), pp. 11-12; DiClemente (2003)")


# --- Slide 14: Action & Maintenance ---
action_left = [
    "ACTION STAGE",
    "",
    "* Client is actively changing behavior",
    "* May have stopped or reduced use",
    "* Implementing the Change Plan",
    "* Visible behavioral modifications",
    "",
    "Therapist Tasks:",
    "* Affirm steps taken",
    "* Help identify and address barriers",
    "* Support self-efficacy",
    "* Review and adjust Change Plan",
    "* Connect to additional resources",
]
action_right = [
    "MAINTENANCE STAGE",
    "",
    "* Sustaining changes over time (6+ months)",
    "* Developing new lifestyle patterns",
    "* Building relapse prevention skills",
    "* Integrating new identity",
    "",
    "Therapist Tasks:",
    "* Reinforce long-term commitment",
    "* Anticipate high-risk situations",
    "* Normalize ongoing challenges",
    "* Plan for future without therapy",
    "* Celebrate progress and growth",
]
make_standard_slide(prs, "Action & Maintenance Stages", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': action_left, 'bold_first': True},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': action_right, 'bold_first': True},
], 2, "Ref: Prochaska & DiClemente (1983); Miller et al. (1992), pp. 60-70")

# --- Slide 15: Relapse ---
relapse_lines = [
    "RELAPSE AS PART OF THE CHANGE CYCLE",
    "",
    "Critical Concept: Relapse is NOT failure - it is a NORMAL part of change.",
    "",
    "Facts About Relapse:",
    "  * Most people cycle through stages 3-7 times before sustained change",
    "  * Relapse provides learning opportunities about triggers",
    "  * MET reframes relapse as a 'slip' not a catastrophe",
    "  * Each attempt builds skills and self-knowledge",
    "",
    "MET Response to Relapse:",
    "  * Avoid shame, blame, or labeling (never say 'you failed')",
    "  * Explore what led to the slip (functional analysis)",
    "  * Reconnect with original motivations for change",
    "  * Renew and revise the Change Plan",
    "  * Strengthen self-efficacy: 'You did it before, you can do it again'",
    "  * Re-engage in the change cycle at the appropriate stage",
]
make_standard_slide(prs, "Relapse as Part of the Change Cycle", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': relapse_lines, 'bold_first': True},
], 2, "Ref: Miller et al. (1992), pp. 60-65; Marlatt & Gordon (1985)")


# --- Slide 16: Self-Efficacy Theory ---
se_lines = [
    "SELF-EFFICACY THEORY (Bandura, 1977) & MET",
    "",
    "Self-efficacy = A person's BELIEF in their ability to succeed at a task.",
    "",
    "Why It Matters in MET:",
    "  * People who believe they CAN change are more likely to actually change",
    "  * Low self-efficacy predicts dropout and continued substance use",
    "  * MET actively builds self-efficacy throughout treatment",
    "",
    "How MET Builds Self-Efficacy:",
    "  * Affirming past successes (however small)",
    "  * Highlighting personal strengths and resources",
    "  * Using confidence rulers to track improvement",
    "  * Providing examples of others who have succeeded",
    "  * Avoiding criticism that would undermine confidence",
    "  * Supporting autonomy: 'Only YOU can make this change'",
    "",
    "MET Principle: Support Self-Efficacy is one of the 5 core MI principles.",
]
make_standard_slide(prs, "Self-Efficacy Theory (Bandura) & MET", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': se_lines, 'bold_first': True},
], 2, "Ref: Bandura (1977); Miller et al. (1992), pp. 18-20; Miller & Rollnick (2002)")

# --- Slide 17: Cognitive Dissonance ---
cd_lines = [
    "COGNITIVE DISSONANCE THEORY (Festinger, 1957) & MET",
    "",
    "Core Idea: When a person holds two contradictory beliefs simultaneously,",
    "they experience psychological discomfort (dissonance) that motivates change.",
    "",
    "Application in MET - 'Developing Discrepancy':",
    "  * Client says: 'I value being a good parent'",
    "  * Reality: 'My drinking has caused me to miss my child's events'",
    "  * This GAP between values and behavior = DISCREPANCY",
    "  * The discomfort from this gap drives motivation to change",
    "",
    "How MET Creates Therapeutic Discrepancy:",
    "  * Personal feedback shows objective consequences",
    "  * Exploring 'what matters most to you' vs current behavior",
    "  * Reflecting back contradictions gently (not confrontationally)",
    "  * Letting the CLIENT notice the gap (not pointing it out directly)",
    "",
    "Key: The motivation comes FROM WITHIN when the person sees the gap.",
]
make_standard_slide(prs, "Cognitive Dissonance Theory & MET", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': cd_lines, 'bold_first': True},
], 2, "Ref: Festinger (1957); Miller et al. (1992), pp. 14-16; Miller & Rollnick (2002)")


# --- Slide 18: Decision-Making Theory ---
dm_lines = [
    "DECISION-MAKING THEORY (Janis & Mann, 1977) & MET",
    "",
    "Core Idea: People make decisions by weighing utilities - gains and losses.",
    "Decisional conflict occurs when all options have significant pros AND cons.",
    "",
    "Janis & Mann's Conflict Model:",
    "  * Unconflicted adherence - ignore the problem (precontemplation)",
    "  * Unconflicted change - impulsive action without planning",
    "  * Defensive avoidance - procrastination, buck-passing",
    "  * Hypervigilance - panic-driven decisions (poor quality)",
    "  * Vigilance - careful weighing of all options (BEST outcomes)",
    "",
    "How MET Uses This Theory:",
    "  * Decisional Balance exercise helps structure the weighing process",
    "  * Therapist creates conditions for 'vigilant' decision-making",
    "  * Avoids creating hypervigilance (no pressure, no deadlines)",
    "  * Supports thorough consideration of all options",
    "",
    "MET aims to help clients make HIGH-QUALITY decisions about change.",
]
make_standard_slide(prs, "Decision-Making Theory (Janis & Mann) & MET", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': dm_lines, 'bold_first': True},
], 2, "Ref: Janis & Mann (1977); Miller et al. (1992), pp. 16-18")


# ============================================================
# SECTION 3: PRINCIPLES & SPIRIT (Slides 19-26)
# ============================================================

# --- Slide 19: Section Divider ---
slide19 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide19, "SECTION 3\nPrinciples & Spirit of\nMotivational Interviewing", 3)

# --- Slide 20: Five Principles of MI ---
princ_left = [
    "THE 5 PRINCIPLES OF MI",
    "",
    "1. EXPRESS EMPATHY",
    "   Understand the client's perspective",
    "   without judging. Acceptance facilitates",
    "   change. Ambivalence is normal.",
    "",
    "2. DEVELOP DISCREPANCY",
    "   Help client see gap between current",
    "   behavior and important personal goals/values.",
    "",
    "3. AVOID ARGUMENTATION",
    "   Arguments are counterproductive.",
    "   Defending creates defensiveness.",
]
princ_right = [
    "",
    "",
    "4. ROLL WITH RESISTANCE",
    "   Don't fight resistance - use it.",
    "   New perspectives are invited, not imposed.",
    "   The client is the primary resource.",
    "",
    "5. SUPPORT SELF-EFFICACY",
    "   Belief in possibility of change is a",
    "   powerful motivator. The client IS capable.",
    "   Hope and optimism are therapeutic.",
    "",
    "These 5 principles guide EVERY",
    "interaction in MET sessions.",
]
make_standard_slide(prs, "Five Principles of Motivational Interviewing", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': princ_left, 'bold_first': True,
     'fill_color': COLORS['light_green'], 'border_color': COLORS['green']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': princ_right, 'bold_first': True,
     'fill_color': COLORS['cream'], 'border_color': COLORS['green']},
], 3, "Ref: Miller & Rollnick (1991, 2002); Miller et al. (1992), pp. 12-20")


# --- Slide 21: FRAMES ---
frames_lines = [
    "FRAMES - Brief Intervention Elements (Miller & Sanchez, 1994)",
    "",
    "F = FEEDBACK",
    "     Provide personal, objective feedback about the client's substance use",
    "     (e.g., how their drinking compares to norms, health risks identified)",
    "",
    "R = RESPONSIBILITY",
    "     Emphasize that change is the CLIENT's choice and responsibility",
    "     (e.g., 'Only you can decide what to do about this')",
    "",
    "A = ADVICE",
    "     Offer clear advice to change (when appropriate, with permission)",
    "",
    "M = MENU OF OPTIONS",
    "     Provide multiple options for change (not just one 'right way')",
    "",
    "E = EMPATHY",
    "     Use warm, reflective, understanding counseling style",
    "",
    "S = SELF-EFFICACY",
    "     Express confidence in client's ability to change",
]
make_standard_slide(prs, "FRAMES - Brief Intervention Elements", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': frames_lines, 'bold_first': True},
], 3, "Ref: Miller & Sanchez (1994); Miller et al. (1992), pp. 21-28; NIMHANS (2016)")

# --- Slide 22: OARS ---
oars_lines = [
    "OARS - Core Microskills of MI",
    "",
    "O = OPEN-ENDED QUESTIONS",
    "     Questions that cannot be answered with just 'yes' or 'no'",
    "     Example: 'Tell me about your drinking' vs 'Do you drink?'",
    "",
    "A = AFFIRMATIONS",
    "     Genuine statements recognizing client's strengths and efforts",
    "     Example: 'It took courage to come here today'",
    "",
    "R = REFLECTIVE LISTENING",
    "     Restating/rephrasing what client says to show understanding",
    "     Example: Client: 'I guess my wife worries' / Therapist: 'Your drinking",
    "     is affecting your marriage in ways that concern you'",
    "",
    "S = SUMMARIES",
    "     Collecting and organizing what client has said",
    "     Links themes, reinforces change talk, transitions conversation",
    "",
    "These 4 skills are the TOOLS the MET therapist uses in every moment.",
]
make_standard_slide(prs, "OARS - Core Microskills", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': oars_lines, 'bold_first': True},
], 3, "Ref: Miller & Rollnick (2002); Miller et al. (1992), pp. 21-25; NIMHANS (2016)")


# --- Slide 23: DARES - Change Talk ---
dares_lines = [
    "DARES - Recognizing Change Talk",
    "",
    "Change talk = Any statement by the client that favors change.",
    "The therapist's job is to ELICIT and REINFORCE change talk.",
    "",
    "D = DESIRE ('I want to quit', 'I wish things were different')",
    "",
    "A = ABILITY ('I could probably cut down', 'I've done it before')",
    "",
    "R = REASONS ('My health is suffering', 'I'm losing my family')",
    "",
    "E = NEED ('I've got to do something', 'I need to change')",
    "",
    "S = COMMITMENT ('I will stop', 'I'm going to do this')",
    "",
    "Opposite = SUSTAIN TALK (arguments for NOT changing)",
    "Example: 'Drinking helps me relax' or 'I can handle it'",
    "",
    "MET Principle: Reinforce DARES, reflect sustain talk without arguing.",
]
make_standard_slide(prs, "DARES - Change Talk Categories", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': dares_lines, 'bold_first': True},
], 3, "Ref: Miller & Rollnick (2002); Amrhein et al. (2003); NIMHANS (2016)")

# --- Slide 24: Spirit of MI ---
spirit_left = [
    "THE SPIRIT OF MI",
    "(The 'Way of Being' with Clients)",
    "",
    "1. COLLABORATION",
    "   * Partnership, not expert-patient",
    "   * Working together side-by-side",
    "   * Honoring client's expertise on self",
    "   * Power is shared equally",
    "",
    "2. EVOCATION",
    "   * Drawing out what's already there",
    "   * Client has the answers within",
    "   * Therapist's role: ask, not tell",
    "   * Motivation is elicited, not installed",
]
spirit_right = [
    "",
    "",
    "",
    "3. AUTONOMY",
    "   * Client's right to choose respected",
    "   * No coercion, no pressure",
    "   * 'It's up to you' is genuine",
    "   * Freedom to choose NOT to change",
    "",
    "WHY THIS MATTERS:",
    "When people feel controlled, they resist.",
    "When people feel free, they choose wisely.",
    "MI Spirit creates conditions where change",
    "becomes the client's own natural choice.",
]
make_standard_slide(prs, "Spirit of MI: Collaboration, Evocation, Autonomy", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': spirit_left, 'bold_first': True,
     'fill_color': COLORS['light_green'], 'border_color': COLORS['green']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': spirit_right, 'bold_first': True,
     'fill_color': COLORS['cream'], 'border_color': COLORS['green']},
], 3, "Ref: Miller & Rollnick (2002), Ch. 2; Miller et al. (1992), pp. 12-14")


# --- Slide 25: Handling Resistance ---
resist_lines = [
    "HANDLING RESISTANCE - Types & Strategies",
    "",
    "Types of Resistance (Miller et al., 1992):",
    "  * Arguing - Challenging, discounting, being hostile",
    "  * Interrupting - Talking over, cutting off therapist",
    "  * Denying - Blaming, disagreeing, excusing, minimizing",
    "  * Ignoring - Inattention, sidetracking, no response",
    "",
    "MET Strategies for Resistance:",
    "  * Simple Reflection - Acknowledge without judgment",
    "  * Amplified Reflection - Slightly exaggerate to prompt correction",
    "  * Double-Sided Reflection - Reflect both sides of ambivalence",
    "  * Shifting Focus - Move away from contentious topic",
    "  * Reframing - Offer new interpretation of information",
    "  * Agreement with a Twist - Agree but add new direction",
    "  * Emphasizing Personal Choice - 'It's entirely up to you'",
    "",
    "Golden Rule: If resistance increases, CHANGE WHAT YOU'RE DOING.",
]
make_standard_slide(prs, "Handling Resistance - Types & Strategies", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': resist_lines, 'bold_first': True},
], 3, "Ref: Miller et al. (1992), pp. 22-28; Miller & Rollnick (2002), Ch. 8")

# --- Slide 26: Clinical Dialogue Examples ---
dialog_left = [
    "RESISTANCE DIALOGUE EXAMPLE",
    "",
    "Client: 'I don't have a problem. My wife",
    "is just overreacting as usual.'",
    "",
    "BAD Response (confrontational):",
    "'Well, your blood tests say otherwise.",
    "You clearly have a drinking problem.'",
    "(This creates MORE resistance)",
    "",
    "GOOD Response (rolling with resistance):",
    "'So from your perspective, your wife's",
    "concerns seem exaggerated. Tell me",
    "more about how you see things.'",
    "(This opens dialogue)",
]
dialog_right = [
    "CHANGE TALK DIALOGUE EXAMPLE",
    "",
    "Client: 'I suppose my drinking has caused",
    "some problems at work lately.'",
    "",
    "BAD Response (missing change talk):",
    "'Okay, what else is going on?'",
    "(Missed opportunity to reinforce)",
    "",
    "GOOD Response (reinforcing change talk):",
    "'You're noticing that alcohol is starting",
    "to affect important areas of your life,",
    "like your career. That takes awareness.",
    "Tell me more about what you've noticed.'",
    "(Amplifies and reinforces change talk)",
]
make_standard_slide(prs, "Clinical Dialogue: Resistance & Change Talk", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': dialog_left, 'bold_first': True,
     'fill_color': COLORS['light_red'], 'border_color': COLORS['maroon']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': dialog_right, 'bold_first': True,
     'fill_color': COLORS['light_green'], 'border_color': COLORS['green']},
], 3, "Ref: Miller et al. (1992), pp. 22-28; Miller & Rollnick (2002), Clinical examples")


# ============================================================
# SECTION 4: SESSION-BY-SESSION GUIDE (Slides 27-38)
# ============================================================

# --- Slide 27: Section Divider ---
slide27 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide27, "SECTION 4\nSession-by-Session Guide", 4)

# --- Slide 28: Session 1 Overview ---
s1_overview = [
    "SESSION 1 OVERVIEW: BUILDING MOTIVATION FOR CHANGE",
    "",
    "Duration: 60-90 minutes (the longest session)",
    "Timing: After assessment battery is completed",
    "",
    "Primary Goals:",
    "  * Establish therapeutic rapport and trust",
    "  * Provide structured Personal Feedback Report (PFR)",
    "  * Elicit client's reactions to feedback",
    "  * Begin developing discrepancy",
    "  * Assess current readiness to change",
    "",
    "Key Materials Needed:",
    "  * Completed assessment results (drinking profile, consequences)",
    "  * Personal Feedback Report prepared in advance",
    "  * Readiness to Change Ruler",
    "  * Change Plan Worksheet (in case client is ready)",
    "",
    "Therapist Mindset: Be curious, not corrective.",
]
make_standard_slide(prs, "Session 1 Overview: Building Motivation", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': s1_overview, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 29-42 (Session 1 protocol)")

# --- Slide 29: Session 1 Opening ---
s1_open = [
    "SESSION 1: OPENING & BUILDING RAPPORT",
    "",
    "Opening the Session (First 10-15 minutes):",
    "",
    "  * Start with an open-ended question about the client's experience",
    "    Example: 'What brought you here today?'",
    "  * Acknowledge any coercion or external pressure without judgment",
    "    Example: 'So your boss suggested you come. How do you feel about that?'",
    "  * Establish a collaborative tone from the first moment",
    "  * Explain the session structure briefly:",
    "    'Today I'd like to share some results from your assessment",
    "     and hear your thoughts about them.'",
    "",
    "Rapport-Building Tips:",
    "  * Use client's name naturally",
    "  * Match their language and pace",
    "  * Affirm them for showing up",
    "  * Be genuinely curious about their world",
    "  * Avoid premature focus on 'the problem'",
]
make_standard_slide(prs, "Session 1: Opening & Rapport Building", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': s1_open, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 30-33; Miller & Rollnick (2002)")


# --- Slide 30: Personal Feedback Report ---
pfr_lines = [
    "SESSION 1: THE PERSONAL FEEDBACK REPORT (PFR)",
    "",
    "The PFR is the CORE TOOL of Session 1. It contains:",
    "",
    "  1. DRINKING/USE PROFILE: Quantity, frequency, patterns",
    "  2. NORMATIVE COMPARISON: How client compares to general population",
    "     (e.g., 'Your drinking places you in the top 5% of adults')",
    "  3. RISK FACTORS: Blood alcohol levels, tolerance indicators",
    "  4. CONSEQUENCES: Physical, social, legal, emotional problems reported",
    "  5. DEPENDENCE INDICATORS: Signs of physiological dependence",
    "  6. NEUROPSYCHOLOGICAL FUNCTIONING: Any cognitive impacts detected",
    "  7. FAMILY RISK FACTORS: Genetic/family history of substance problems",
    "",
    "How to Present the PFR:",
    "  * Give the report TO the client (their copy)",
    "  * Go through each section, asking for reactions",
    "  * 'What do you make of this?' / 'Does this surprise you?'",
    "  * Let the client process - don't rush or interpret for them",
]
make_standard_slide(prs, "Session 1: Personal Feedback Report", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': pfr_lines, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 33-38 (Feedback protocol)")

# --- Slide 31: Eliciting Change Talk in Session 1 ---
ct_s1 = [
    "SESSION 1: ELICITING CHANGE TALK",
    "",
    "After presenting feedback, actively draw out the client's motivation:",
    "",
    "Key Questions to Elicit Change Talk:",
    "  * 'What concerns you most about what you've heard?'",
    "  * 'In what ways does this worry you?'",
    "  * 'What would happen if you didn't change anything?'",
    "  * 'What are the best reasons for making a change?'",
    "  * 'What would be different in your life if you changed?'",
    "  * 'On a scale of 0-10, how important is it for you to change?'",
    "",
    "When You Hear Change Talk:",
    "  * Reflect it back (amplify it)",
    "  * Ask for more: 'Tell me more about that'",
    "  * Affirm: 'That's an important realization'",
    "  * Summarize change talk you've heard so far",
    "",
    "When You Hear Sustain Talk: Reflect it, then ask the other side.",
]
make_standard_slide(prs, "Session 1: Eliciting Change Talk", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': ct_s1, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 38-40; Miller & Rollnick (2002), Ch. 7")


# --- Slide 32: Session 1 Closing ---
s1_close = [
    "SESSION 1: CLOSING & READINESS ASSESSMENT",
    "",
    "Closing the First Session (Last 10-15 minutes):",
    "",
    "1. Provide a SUMMARY of what you've discussed:",
    "   'Let me pull together what you've shared today...'",
    "   Include both sides of ambivalence, ending with change talk.",
    "",
    "2. Assess READINESS TO CHANGE:",
    "   Use the Readiness Ruler (0-10):",
    "   'On a scale of 0 to 10, how ready are you to make a change?'",
    "   Follow up: 'Why a [number] and not a lower number?'",
    "",
    "3. Based on Readiness Level:",
    "   * Low (1-3): 'Thank you for being open. Let's talk again next time.'",
    "   * Medium (4-6): 'You're considering change. Let's explore that more.'",
    "   * High (7-10): 'You seem quite ready. Would you like to start",
    "     making a plan?' (Introduce Change Plan Worksheet)",
    "",
    "4. Schedule Session 2 and give client the PFR to take home.",
]
make_standard_slide(prs, "Session 1: Closing & Readiness Assessment", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': s1_close, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 40-42; Rollnick et al. (1999)")

# --- Slide 33: Session 2 Overview ---
s2_overview = [
    "SESSION 2 OVERVIEW: STRENGTHENING COMMITMENT",
    "",
    "Duration: 45-60 minutes",
    "Timing: 1-2 weeks after Session 1",
    "",
    "Primary Goals:",
    "  * Review what happened since Session 1",
    "  * Reinforce and amplify any change talk or steps taken",
    "  * Address ongoing ambivalence if present",
    "  * Develop a formal, written CHANGE PLAN",
    "  * Consolidate commitment to the plan",
    "",
    "Key Principle: The Session 2 approach depends on the client's",
    "current readiness level:",
    "  * If still ambivalent: Continue building motivation (don't force)",
    "  * If ready: Move to Change Plan development",
    "  * If already acting: Affirm, support, help structure efforts",
    "",
    "Session 2 is where the 'bridge' from motivation to action is built.",
]
make_standard_slide(prs, "Session 2 Overview: Strengthening Commitment", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': s2_overview, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 43-54 (Session 2 protocol)")


# --- Slide 34: Change Plan Development ---
cp_lines = [
    "SESSION 2: THE CHANGE PLAN",
    "",
    "The Change Plan Worksheet has these components:",
    "",
    "1. THE CHANGES I WANT TO MAKE ARE:",
    "   (Specific, client-chosen goals)",
    "",
    "2. THE MOST IMPORTANT REASONS TO CHANGE ARE:",
    "   (Client's own reasons, in their own words)",
    "",
    "3. THE STEPS I PLAN TO TAKE ARE:",
    "   (Concrete, achievable action steps)",
    "",
    "4. THE WAYS OTHER PEOPLE CAN HELP ME ARE:",
    "   (Social support identified)",
    "",
    "5. I WILL KNOW MY PLAN IS WORKING IF:",
    "   (Observable, measurable indicators)",
    "",
    "6. THINGS THAT COULD INTERFERE AND HOW I'LL HANDLE THEM:",
    "   (Anticipating obstacles and creating coping strategies)",
]
make_standard_slide(prs, "Session 2: Change Plan Development", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': cp_lines, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 46-50 (Change Plan Worksheet)")

# --- Slide 35: Adapting to Readiness ---
adapt_left = [
    "IF CLIENT IS NOT YET READY",
    "(Readiness = 1-4)",
    "",
    "Do NOT push the Change Plan.",
    "Instead:",
    "* Continue exploring ambivalence",
    "* Ask evocative questions",
    "* Use decisional balance",
    "* Provide more information",
    "* Respect their pace",
    "* Review personal feedback again",
    "* End with door open for next time",
    "",
    "'There's no rush. I'm here when",
    "you're ready to talk more.'",
]
adapt_right = [
    "IF CLIENT IS READY",
    "(Readiness = 7-10)",
    "",
    "Move to Change Plan directly.",
    "Steps:",
    "* Set clear, specific goals",
    "* Brainstorm strategies together",
    "* Identify support people",
    "* Write out the plan formally",
    "* Have client sign/commit",
    "* Give them a copy",
    "* Plan for obstacles",
    "",
    "'Let's put this into a plan that",
    "works for you.'",
]
make_standard_slide(prs, "Session 2: Adapting to Readiness Level", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': adapt_left, 'bold_first': True,
     'fill_color': COLORS['light_orange'], 'border_color': COLORS['orange']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': adapt_right, 'bold_first': True,
     'fill_color': COLORS['light_green'], 'border_color': COLORS['green']},
], 4, "Ref: Miller et al. (1992), pp. 43-54; Rollnick et al. (1999)")


# --- Slide 36: Sessions 3-4 ---
s34_lines = [
    "SESSIONS 3 AND 4: REVIEW & RENEWAL",
    "",
    "Timing: Typically at weeks 6 and 12 (spaced out intentionally)",
    "",
    "Session 3 Goals:",
    "  * Review progress on the Change Plan",
    "  * Affirm positive changes observed",
    "  * Address any barriers or difficulties encountered",
    "  * Renew or revise commitment as needed",
    "  * Reinforce self-efficacy with evidence of progress",
    "",
    "Session 4 Goals:",
    "  * Final review of treatment goals and progress",
    "  * Anticipate future challenges and high-risk situations",
    "  * Plan for long-term maintenance without therapy",
    "  * Discuss options for continued care if needed",
    "  * Celebrate achievements and mark the ending",
    "",
    "Key Principle: These sessions match the client's current stage.",
]
make_standard_slide(prs, "Sessions 3-4: Review & Renewal", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': s34_lines, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 55-70 (Sessions 3-4 protocol)")

# --- Slide 37: Handling Relapse in Sessions 3-4 ---
rel_s34 = [
    "HANDLING RELAPSE IN SESSIONS 3-4",
    "",
    "If the client has relapsed or slipped since Session 2:",
    "",
    "Step 1: NORMALIZE",
    "  'Setbacks are a normal part of change. Most people experience them.'",
    "",
    "Step 2: EXPLORE (without blame)",
    "  'What happened? Walk me through the situation.'",
    "  'What were you feeling before the slip?'",
    "",
    "Step 3: LEARN",
    "  'What did you learn from this about your triggers?'",
    "  'What might you do differently next time?'",
    "",
    "Step 4: RECONNECT WITH MOTIVATION",
    "  'What were your main reasons for wanting to change?'",
    "  'Are those reasons still important to you?'",
    "",
    "Step 5: REVISE THE PLAN",
    "  'Let's look at your Change Plan and see what needs adjusting.'",
]
make_standard_slide(prs, "Handling Relapse in Sessions 3-4", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': rel_s34, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 60-65; Marlatt & Gordon (1985)")


# --- Slide 38: Full Session Dialogue Example ---
full_dialog = [
    "CLINICAL DIALOGUE: SESSION 1 EXCERPT (Feedback Phase)",
    "",
    "Therapist: 'I'd like to share some results from your assessment. Your",
    "  drinking pattern places you at the 95th percentile - meaning you drink",
    "  more than 95% of adults your age. What do you make of that?'",
    "",
    "Client: 'I guess I didn't realize it was that much compared to others.'",
    "",
    "Therapist: 'It's a bit surprising to you. [Reflection] You also reported",
    "  some memory lapses and missed days at work. How do these fit together",
    "  for you?'",
    "",
    "Client: 'I mean... I suppose the drinking is causing some of that. But",
    "  I still think I can handle it.'",
    "",
    "Therapist: 'So on one hand, you're seeing some connections between your",
    "  drinking and these consequences. On the other hand, part of you feels",
    "  you have it under control. [Double-sided reflection]'",
    "",
    "Client: 'Yeah... but honestly, the memory thing worries me.'",
    "",
    "Therapist: 'That worry tells me something matters to you here.' [Affirm]",
]
make_standard_slide(prs, "Clinical Dialogue: Full Session Example", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': full_dialog, 'bold_first': True},
], 4, "Ref: Miller et al. (1992), pp. 33-40 (Session 1 examples); adapted")


# ============================================================
# SECTION 5: CLINICAL TECHNIQUES (Slides 39-44)
# ============================================================

# --- Slide 39: Section Divider ---
slide39 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide39, "SECTION 5\nClinical Techniques in MET", 5)

# --- Slide 40: Reflective Listening ---
refl_lines = [
    "REFLECTIVE LISTENING - Levels & Examples",
    "",
    "Level 1: SIMPLE REFLECTION (Repeat/Rephrase)",
    "  Client: 'I drink every evening after work.'",
    "  Therapist: 'You drink every evening.' (Minimal change)",
    "",
    "Level 2: PARAPHRASE (Same meaning, different words)",
    "  Client: 'My wife keeps nagging me about my drinking.'",
    "  Therapist: 'Your wife's concern about your drinking frustrates you.'",
    "",
    "Level 3: REFLECTION OF FEELING (Identify unstated emotion)",
    "  Client: 'I just don't know what to do anymore.'",
    "  Therapist: 'You're feeling overwhelmed and stuck.'",
    "",
    "Level 4: REFLECTION OF MEANING (Connect to values/identity)",
    "  Client: 'I missed my daughter's recital because I was hungover.'",
    "  Therapist: 'Being a present father matters deeply to you, and",
    "  alcohol got in the way of that.'",
    "",
    "MET uses primarily Level 2-4 reflections. Aim for 2:1 reflection-to-question ratio.",
]
make_standard_slide(prs, "Reflective Listening - Levels & Examples", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': refl_lines, 'bold_first': True},
], 5, "Ref: Miller et al. (1992), pp. 21-23; Miller & Rollnick (2002), Ch. 5")

# --- Slide 41: Decisional Balance ---
db_left = [
    "DECISIONAL BALANCE TECHNIQUE",
    "",
    "A structured way to explore ambivalence.",
    "Creates a 4-cell matrix:",
    "",
    "BENEFITS OF STAYING THE SAME:",
    "* Relaxation, social connection",
    "* Avoids withdrawal discomfort",
    "* Familiar routine, pleasure",
    "",
    "COSTS OF STAYING THE SAME:",
    "* Health problems, money lost",
    "* Relationship damage",
    "* Work problems, legal issues",
]
db_right = [
    "",
    "",
    "",
    "",
    "",
    "BENEFITS OF CHANGING:",
    "* Better health, save money",
    "* Improved relationships",
    "* More energy, clarity, pride",
    "",
    "COSTS OF CHANGING:",
    "* Withdrawal, loss of coping tool",
    "* Social pressure, boredom",
    "* Hard work, uncertainty",
    "",
    "The therapist explores ALL four quadrants",
    "without favoring one side.",
]
make_standard_slide(prs, "Decisional Balance Technique", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': db_left, 'bold_first': True,
     'fill_color': COLORS['light_blue'], 'border_color': COLORS['deep_blue']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': db_right, 'bold_first': True,
     'fill_color': COLORS['cream'], 'border_color': COLORS['deep_blue']},
], 5, "Ref: Miller et al. (1992), pp. 16-18; Janis & Mann (1977)")


# --- Slide 42: Importance & Confidence Rulers ---
ruler_lines = [
    "IMPORTANCE & CONFIDENCE RULERS",
    "",
    "Simple, powerful tools for assessing motivation:",
    "",
    "IMPORTANCE RULER:",
    "'On a scale of 0 to 10, how IMPORTANT is it for you to change?'",
    "  0 = Not at all important .............. 10 = Most important thing",
    "",
    "  Follow-up: 'Why did you say [N] and not a lower number?'",
    "  (This elicits change talk - reasons for change)",
    "",
    "  Follow-up: 'What would it take to move from [N] to [N+2]?'",
    "  (This identifies barriers and needed conditions)",
    "",
    "CONFIDENCE RULER:",
    "'On a scale of 0 to 10, how CONFIDENT are you that you could change?'",
    "  0 = Not at all confident .............. 10 = Completely confident",
    "",
    "  Follow-up: 'Why not a zero?' (Elicits self-efficacy)",
    "  Follow-up: 'What would increase your confidence?' (Identifies support needs)",
    "",
    "These rulers make abstract feelings concrete and measurable.",
]
make_standard_slide(prs, "Importance & Confidence Rulers", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': ruler_lines, 'bold_first': True},
], 5, "Ref: Rollnick et al. (1999); Miller et al. (1992), pp. 40-41")

# --- Slide 43: Normative Feedback ---
norm_lines = [
    "NORMATIVE FEEDBACK - Showing How the Client Compares",
    "",
    "What is Normative Feedback?",
    "  Presenting the client's data alongside population norms.",
    "  This creates discrepancy through objective comparison.",
    "",
    "Example Presentation:",
    "  'The average adult in your age group drinks about 3 drinks per week.",
    "   Your assessment shows you're averaging about 28 drinks per week.",
    "   That puts you in the top 3% of drinkers nationally.'",
    "",
    "Types of Normative Data Used in MET:",
    "  * Quantity/frequency compared to age/gender norms",
    "  * Blood alcohol content estimates",
    "  * Liver function test results",
    "  * Number of consequences experienced vs. typical",
    "  * Dependence severity compared to clinical populations",
    "",
    "Key: Present OBJECTIVELY. Let the numbers speak.",
    "Don't add judgment. Ask: 'What do you make of this?'",
]
make_standard_slide(prs, "Normative Feedback Presentation", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': norm_lines, 'bold_first': True},
], 5, "Ref: Miller et al. (1992), pp. 33-38; Agostinelli et al. (1995)")


# --- Slide 44: Therapist Traps ---
traps_lines = [
    "THERAPIST TRAPS TO AVOID IN MET",
    "",
    "1. THE QUESTION-ANSWER TRAP",
    "   Asking too many closed questions in a row (feels like interrogation)",
    "   Fix: Use more reflections, fewer questions (2:1 ratio)",
    "",
    "2. THE CONFRONTATION-DENIAL TRAP",
    "   Arguing with the client about their problem",
    "   Fix: Roll with resistance, avoid 'you have a problem' statements",
    "",
    "3. THE EXPERT TRAP",
    "   Telling the client what to do (advice-giving without permission)",
    "   Fix: Elicit the client's own ideas first",
    "",
    "4. THE LABELING TRAP",
    "   Insisting on diagnostic labels ('You're an alcoholic')",
    "   Fix: Focus on behavior and consequences, not labels",
    "",
    "5. THE PREMATURE FOCUS TRAP",
    "   Jumping to 'the problem' before the client is ready",
    "   Fix: Follow the client's pace; build rapport first",
    "",
    "6. THE BLAMING TRAP",
    "   Focusing on whose 'fault' the problem is",
    "   Fix: Focus forward on solutions, not backward on blame",
]
make_standard_slide(prs, "Therapist Traps to Avoid", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': traps_lines, 'bold_first': True},
], 5, "Ref: Miller et al. (1992), pp. 23-28; Miller & Rollnick (2002), Ch. 6")


# ============================================================
# SECTION 6: CASE CONCEPTUALIZATION (Slides 45-50)
# ============================================================

# --- Slide 45: Section Divider ---
slide45 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide45, "SECTION 6\nCase Conceptualization\n(MET Application)", 6)

# --- Slide 46: Case Background ---
case_bg = [
    "CASE BACKGROUND (Anonymized)",
    "",
    "Demographics: 25-year-old male, recently divorced",
    "",
    "Presenting Problem:",
    "  * Polysubstance use (alcohol + other substances)",
    "  * Relationship breakdown (divorce)",
    "  * Occupational difficulties",
    "  * Emotional dysregulation",
    "",
    "History:",
    "  * Substance use escalated over past 3 years",
    "  * Multiple failed attempts at controlled use",
    "  * Family history of substance use problems",
    "  * Divorce triggered increase in consumption",
    "  * Currently referred by family for treatment",
    "",
    "Stage of Change at Intake: Precontemplation/early Contemplation",
    "(Acknowledges some problems but minimizes their connection to substance use)",
]
make_standard_slide(prs, "Case Background - Anonymized Client", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': case_bg, 'bold_first': True},
], 6, "Ref: Case formulation based on MET Manual framework (Miller et al., 1992)")

# --- Slide 47: Psychological Profile ---
psych_left = [
    "PSYCHOLOGICAL ASSESSMENT FINDINGS",
    "",
    "Cognitive-Emotional Profile:",
    "* Externalizing personality organization",
    "* Pervasively extratensive style",
    "  (EB = 2:10, emotion-driven coping)",
    "* Poorly modulated affect",
    "  (FC:CF+C = 3:7 - emotions overflow)",
    "* Adequate resources (EA=12) but",
    "  poorly utilized under stress",
    "",
    "Information Processing:",
    "* Extreme underincorporative style",
    "  (Zd = -14, hasty scanning)",
    "* Rushes through decisions",
    "* Misses important information",
]
psych_right = [
    "",
    "",
    "Self-Image & Relationships:",
    "* Damaged self-concept (MOR=5)",
    "* Low self-esteem (Egocentricity low)",
    "* Socially engaged but avoids",
    "  emotional intimacy (Sum T = 0)",
    "* Guarded attachment style",
    "",
    "Risk & Clinical Concerns:",
    "* Elevated risk indicators (S-CON=7)",
    "* Poor impulse control",
    "* Limited reflective capacity",
    "",
    "Therapeutic Needs:",
    "* Affect regulation training",
    "* Reflective delay development",
    "* Self-image repair",
    "* Building distress tolerance",
]
make_standard_slide(prs, "Psychological Profile & Assessment", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': psych_left, 'bold_first': True,
     'fill_color': COLORS['light_red'], 'border_color': COLORS['maroon']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': psych_right, 'bold_first': True,
     'fill_color': COLORS['cream'], 'border_color': COLORS['maroon']},
], 6, "Ref: Assessment interpreted per Exner (2003); Rorschach CS norms")


# --- Slide 48: Session 1 Plan ---
s1_plan = [
    "MET SESSION 1 PLAN FOR THIS CLIENT",
    "",
    "Given Profile: Precontemplator, externalizing, emotion-driven, guarded",
    "",
    "Opening (10 min): Build rapport without pressure",
    "  T: 'I appreciate you coming today. What's brought you here?'",
    "  C: 'My family dragged me. I don't think I need this.'",
    "  T: 'So this wasn't your idea. Tell me about your situation in your words.'",
    "",
    "Feedback Phase (30 min): Present gently, ask for reactions",
    "  T: 'Your assessment shows some interesting patterns. For example, your",
    "  substance use is in a range that puts you at elevated health risk.",
    "  What do you make of that?'",
    "  Adaptation: Slow pace (matches his hasty processing - model reflective delay)",
    "",
    "Eliciting Change Talk (15 min):",
    "  T: 'What matters most to you in life right now?'",
    "  T: 'How does your substance use fit with those things that matter?'",
    "  Adaptation: Connect to damaged self-image - 'who do you want to be?'",
    "",
    "Closing: Readiness ruler, affirm courage, schedule Session 2.",
]
make_standard_slide(prs, "Session 1 Plan: This Client", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': s1_plan, 'bold_first': True},
], 6, "Ref: Miller et al. (1992), Session 1 protocol; adapted for case profile")

# --- Slide 49: Session 2 Change Plan ---
s2_plan = [
    "SESSION 2: CHANGE PLAN FOR THIS CLIENT",
    "",
    "Assuming client moved to early Contemplation after Session 1:",
    "",
    "THE CHANGES I WANT TO MAKE:",
    "  'Reduce drinking to weekends only; stop using other substances'",
    "",
    "MY MOST IMPORTANT REASONS:",
    "  'I want to feel better about myself; I want my family to respect me'",
    "  (Leverages damaged self-image as motivator)",
    "",
    "STEPS I PLAN TO TAKE:",
    "  * Avoid triggers (certain friends, places) - addresses hasty processing",
    "  * Pause before acting on urges (builds reflective delay)",
    "  * Tell one person about my goals (addresses intimacy avoidance)",
    "  * Track use daily (self-monitoring)",
    "",
    "WHAT COULD GO WRONG:",
    "  'Strong emotions, loneliness, old friends' - (matches extratensive style)",
    "  Plan: Call support person; leave situation; use breathing exercises",
    "",
    "PROGRESS INDICATORS: Fewer use days, better mood, family engagement",
]
make_standard_slide(prs, "Session 2: Change Plan for This Client", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': s2_plan, 'bold_first': True},
], 6, "Ref: Miller et al. (1992), pp. 46-50; adapted for case formulation")


# --- Slide 50: Sessions 3-4 & Long-term ---
long_plan = [
    "SESSIONS 3-4 & LONG-TERM PLAN",
    "",
    "SESSION 3 (Week 6) - Expected Focus:",
    "  * Review Change Plan adherence (likely partial compliance)",
    "  * Address relapses with empathy and normalization",
    "  * Reconnect with motivations (self-image, family respect)",
    "  * Modify plan based on what worked/didn't work",
    "  * Affirm: 'You came back. That shows commitment.'",
    "",
    "SESSION 4 (Week 12) - Expected Focus:",
    "  * Review overall progress across 12 weeks",
    "  * Celebrate gains, no matter how small",
    "  * Plan for ongoing maintenance without MET",
    "  * Discuss referral options (group therapy, individual therapy)",
    "",
    "LONG-TERM RECOMMENDATIONS:",
    "  * Follow-up with CBT for affect regulation skills",
    "  * Schema therapy for self-image repair",
    "  * Possible 12-step or SMART Recovery for social support",
    "  * Periodic MET 'booster' sessions as needed",
    "  * Family involvement through CRAFT approach",
]
make_standard_slide(prs, "Sessions 3-4 & Long-Term Plan", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': long_plan, 'bold_first': True},
], 6, "Ref: Miller et al. (1992), pp. 55-70; NIMHANS SUD Manual (2016)")


# ============================================================
# SECTION 7: WORKSHEETS (Slides 51-56)
# ============================================================

# --- Slide 51: Section Divider ---
slide51 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide51, "SECTION 7\nWorksheets & Clinical Tools", 7)

# --- Slide 52: Change Plan Worksheet ---
cpw_lines = [
    "CHANGE PLAN WORKSHEET (Miller et al., 1992)",
    "",
    "-------------------------------------------------------------------",
    "The changes I want to make are:",
    "_______________________________________________________________",
    "",
    "The most important reasons why I want to make these changes are:",
    "_______________________________________________________________",
    "",
    "The steps I plan to take in changing are:",
    "_______________________________________________________________",
    "",
    "The ways other people can help me are:",
    "Person: _________________ How: _______________________________",
    "",
    "I will know that my plan is working if:",
    "_______________________________________________________________",
    "",
    "Some things that could interfere with my plan are:",
    "_______________________________________________________________",
    "-------------------------------------------------------------------",
    "Client Signature: _________________ Date: _____________",
]
make_standard_slide(prs, "Change Plan Worksheet", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': cpw_lines, 'bold_first': True},
], 7, "Ref: Miller et al. (1992), pp. 46-50 (Change Plan Worksheet - Appendix)")

# --- Slide 53: Decisional Balance Grid ---
db_grid_left = [
    "DECISIONAL BALANCE GRID",
    "",
    "BENEFITS OF CURRENT BEHAVIOR",
    "(What I like about my substance use)",
    "",
    "1. _________________________________",
    "2. _________________________________",
    "3. _________________________________",
    "4. _________________________________",
    "",
    "COSTS OF CHANGING",
    "(What I would lose/miss if I changed)",
    "",
    "1. _________________________________",
    "2. _________________________________",
    "3. _________________________________",
]
db_grid_right = [
    "",
    "",
    "COSTS OF CURRENT BEHAVIOR",
    "(What I don't like about my use)",
    "",
    "1. _________________________________",
    "2. _________________________________",
    "3. _________________________________",
    "4. _________________________________",
    "",
    "BENEFITS OF CHANGING",
    "(What I would gain if I changed)",
    "",
    "1. _________________________________",
    "2. _________________________________",
    "3. _________________________________",
]
make_standard_slide(prs, "Decisional Balance Grid", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': db_grid_left, 'bold_first': True,
     'fill_color': COLORS['light_blue'], 'border_color': COLORS['navy']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': db_grid_right, 'bold_first': True,
     'fill_color': COLORS['cream'], 'border_color': COLORS['navy']},
], 7, "Ref: Janis & Mann (1977); Miller et al. (1992), pp. 16-18")


# --- Slide 54: Importance/Confidence/Readiness Rulers ---
ruler_ws = [
    "IMPORTANCE / CONFIDENCE / READINESS RULERS",
    "",
    "IMPORTANCE: How important is it for you to make this change?",
    "",
    "  0----1----2----3----4----5----6----7----8----9----10",
    "  Not at all                                    Extremely",
    "  important                                     important",
    "",
    "CONFIDENCE: How confident are you that you COULD change?",
    "",
    "  0----1----2----3----4----5----6----7----8----9----10",
    "  Not at all                                    Completely",
    "  confident                                     confident",
    "",
    "READINESS: How ready are you to make a change RIGHT NOW?",
    "",
    "  0----1----2----3----4----5----6----7----8----9----10",
    "  Not ready                                     Very ready",
    "  at all                                        right now",
    "",
    "Follow-up Questions: 'Why not a lower number?' / 'What would move you up?'",
]
make_standard_slide(prs, "Importance / Confidence / Readiness Rulers", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': ruler_ws, 'bold_first': True},
], 7, "Ref: Rollnick et al. (1999); Miller et al. (1992), pp. 40-41")

# --- Slide 55: Daily Self-Monitoring Diary ---
diary_lines = [
    "DAILY SELF-MONITORING DIARY",
    "",
    "Instructions: Record each day's substance use, triggers, and mood.",
    "",
    "| Day    | Substance | Amount | Trigger/Situation | Mood (0-10) | Urge (0-10) |",
    "|--------|-----------|--------|-------------------|-------------|-------------|",
    "| Mon    | _________ | ______ | _________________ | ___________ | ___________ |",
    "| Tue    | _________ | ______ | _________________ | ___________ | ___________ |",
    "| Wed    | _________ | ______ | _________________ | ___________ | ___________ |",
    "| Thu    | _________ | ______ | _________________ | ___________ | ___________ |",
    "| Fri    | _________ | ______ | _________________ | ___________ | ___________ |",
    "| Sat    | _________ | ______ | _________________ | ___________ | ___________ |",
    "| Sun    | _________ | ______ | _________________ | ___________ | ___________ |",
    "",
    "Weekly Summary:",
    "  Total days used: ___  Total amount: ___  Highest urge day: ___",
    "  Most common trigger: _______________________________________________",
    "  Coping strategy used: ______________________________________________",
    "",
    "Bring this diary to your next MET session for review.",
]
make_standard_slide(prs, "Daily Self-Monitoring Diary", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': diary_lines, 'bold_first': True},
], 7, "Ref: NIMHANS SUD Manual (2016); Sobell & Sobell (1992) - Timeline Follow-back")


# --- Slide 56: Relapse Prevention Plan + Values ---
rp_left = [
    "RELAPSE PREVENTION PLAN",
    "",
    "My High-Risk Situations:",
    "1. _________________________________",
    "2. _________________________________",
    "3. _________________________________",
    "",
    "My Warning Signs:",
    "1. _________________________________",
    "2. _________________________________",
    "",
    "My Coping Strategies:",
    "1. _________________________________",
    "2. _________________________________",
    "3. _________________________________",
    "",
    "My Emergency Contacts:",
    "1. _________________________________",
    "2. _________________________________",
]
rp_right = [
    "VALUES CARD SORT (Abbreviated)",
    "",
    "My Top 5 Values:",
    "1. _________________________________",
    "2. _________________________________",
    "3. _________________________________",
    "4. _________________________________",
    "5. _________________________________",
    "",
    "How my substance use CONFLICTS",
    "with these values:",
    "_________________________________",
    "_________________________________",
    "",
    "How CHANGING would serve",
    "these values:",
    "_________________________________",
    "_________________________________",
]
make_standard_slide(prs, "Relapse Prevention Plan & Values Worksheet", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': rp_left, 'bold_first': True,
     'fill_color': COLORS['light_blue'], 'border_color': COLORS['navy']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': rp_right, 'bold_first': True,
     'fill_color': COLORS['cream'], 'border_color': COLORS['navy']},
], 7, "Ref: Marlatt & Gordon (1985); Miller et al. (1992); W. Miller Values Card Sort")


# ============================================================
# SECTION 8: COMPARISON WITH OTHER THERAPIES (Slides 57-64)
# ============================================================

# --- Slide 57: Section Divider ---
slide57 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide57, "SECTION 8\nMET Compared with\nOther Therapies", 8)

# --- Slide 58: MET vs CBT ---
cbt_left = [
    "MET (Motivational Enhancement Therapy)",
    "",
    "Focus: Building internal motivation",
    "Sessions: 4 (brief)",
    "Approach: Non-directive, client-led",
    "Technique: MI principles, feedback",
    "Goal: Client finds OWN reasons",
    "Therapist Role: Guide, mirror",
    "Homework: Minimal (self-monitoring)",
    "Best For: Ambivalent clients,",
    "  early stages of change",
    "Theory: Stages of Change, MI",
    "Resistance: Rolled with",
    "Skills: Not taught directly",
]
cbt_right = [
    "CBT (Cognitive-Behavioral Therapy)",
    "",
    "Focus: Changing thoughts & behaviors",
    "Sessions: 12-16 (intensive)",
    "Approach: Structured, therapist-led",
    "Technique: Skills training, practice",
    "Goal: Therapist teaches new skills",
    "Therapist Role: Teacher, coach",
    "Homework: Extensive (skill practice)",
    "Best For: Action/maintenance stage,",
    "  skill deficits identified",
    "Theory: Learning theory, cognition",
    "Resistance: Addressed as cognition",
    "Skills: Directly taught & practiced",
]
make_standard_slide(prs, "MET vs CBT - Detailed Comparison", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': cbt_left, 'bold_first': True,
     'fill_color': COLORS['light_orange'], 'border_color': COLORS['orange']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': cbt_right, 'bold_first': True,
     'fill_color': COLORS['light_blue'], 'border_color': COLORS['deep_blue']},
], 8, "Ref: Project MATCH Research Group (1997); Miller et al. (1992); Carroll (1998)")

# --- Slide 59: MET vs 12-Step ---
tsf_left = [
    "MET",
    "",
    "Philosophy: Client-driven change",
    "Higher Power: Not involved",
    "Labels: Avoided entirely",
    "Disease Model: Not assumed",
    "Goal: Flexible (abstinence or reduction)",
    "Groups: Not required",
    "Sponsor: Not required",
    "Duration: 4 sessions",
    "Steps/Stages: Stages of Change model",
    "Confrontation: Never used",
    "Evidence: Strong RCT evidence",
]
tsf_right = [
    "12-STEP FACILITATION (TSF)",
    "",
    "Philosophy: Surrender to higher power",
    "Higher Power: Central to recovery",
    "Labels: 'I am an alcoholic' is expected",
    "Disease Model: Core assumption",
    "Goal: Total abstinence only",
    "Groups: AA/NA attendance required",
    "Sponsor: Encouraged strongly",
    "Duration: 12 sessions + lifelong meetings",
    "Steps/Stages: 12 sequential steps",
    "Confrontation: Common in groups",
    "Evidence: Moderate evidence (MATCH, Cochrane)",
]
make_standard_slide(prs, "MET vs 12-Step Facilitation", [
    {'x': COL1_X, 'y': CONTENT_Y, 'w': COL1_WIDTH, 'h': Inches(4.9),
     'lines': tsf_left, 'bold_first': True,
     'fill_color': COLORS['light_orange'], 'border_color': COLORS['orange']},
    {'x': COL2_X, 'y': CONTENT_Y, 'w': COL2_WIDTH, 'h': Inches(4.9),
     'lines': tsf_right, 'bold_first': True,
     'fill_color': COLORS['light_purple'], 'border_color': COLORS['purple']},
], 8, "Ref: Project MATCH (1997); Nowinski et al. (1992); Ferri et al. (2006, Cochrane)")


# --- Slide 60: Comprehensive Comparison Table ---
comp_table = [
    "COMPREHENSIVE COMPARISON TABLE: 6 THERAPY APPROACHES FOR SUD",
    "",
    "Therapy      | Sessions | Focus           | Best For          | Evidence",
    "-------------|----------|-----------------|-------------------|----------",
    "MET          | 4        | Motivation      | Ambivalent clients| Strong",
    "CBT          | 12-16    | Skills/Thoughts | Skill deficits    | Strong",
    "TSF(12-Step) | 12+life  | Spiritual/Group | Community-oriented| Moderate",
    "CRA          | 12-24    | Environmental   | Social isolation  | Strong",
    "CRAFT        | 6-12     | Family system   | Concerned others  | Strong",
    "Contingency  | Ongoing  | Reinforcement   | Concrete thinkers | Strong",
    "Mgmt (CM)    |          |                 |                   |",
    "",
    "Key Insight from Research:",
    "  * No single therapy works for everyone",
    "  * MET's advantage: cost-effective, brief, works across stages",
    "  * Combined approaches often work best (MET + CBT, MET + CM)",
    "  * Client preference and matching matter",
    "",
    "MET is uniquely cost-effective: achieves comparable results in 1/3 the time.",
]
make_standard_slide(prs, "6 Therapies Compared - Comprehensive Table", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': comp_table, 'bold_first': True},
], 8, "Ref: Project MATCH (1997); UKATT (2005); NIMHANS (2016); Magill & Ray (2009)")

# --- Slide 61: Why MET is More Effective ---
why_met = [
    "WHY MET IS MORE EFFECTIVE IN SPECIFIC CONTEXTS",
    "",
    "MET outperforms other therapies when:",
    "",
    "1. CLIENTS ARE AMBIVALENT OR NOT READY TO CHANGE",
    "   * CBT and TSF assume readiness - MET doesn't",
    "   * MET builds motivation FIRST, then supports action",
    "",
    "2. CLIENTS ARE ANGRY OR RESISTANT TO TREATMENT",
    "   * Project MATCH: Angry clients did BETTER with MET",
    "   * Confrontation makes resistance worse; empathy reduces it",
    "",
    "3. RESOURCES ARE LIMITED (Time, Money, Staff)",
    "   * 4 sessions vs 12-16: massive cost savings",
    "   * Equal outcomes with fewer resources",
    "",
    "4. CLIENTS HAVE HIGH SOCIAL NETWORKS",
    "   * MATCH: High social support + MET = best outcomes",
    "",
    "5. AS A PRETREATMENT ENHANCER",
    "   * MET before other therapies increases retention & outcomes",
]
make_standard_slide(prs, "Why MET is More Effective in Specific Contexts", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': why_met, 'bold_first': True},
], 8, "Ref: Project MATCH (1997); Bien et al. (1993); Hettema et al. (2005)")


# --- Slide 62: Family & Caregiver Involvement ---
family_lines = [
    "FAMILY & CAREGIVER INVOLVEMENT",
    "",
    "CRAFT (Community Reinforcement & Family Training):",
    "  * Teaches family members MI skills to support change",
    "  * Family learns to reinforce sobriety, not enable use",
    "  * 64% of resistant clients entered treatment via CRAFT",
    "  * Works WITHOUT the user needing to attend initially",
    "",
    "Stepped Care Model with MET:",
    "  Step 1: Brief MET (4 sessions) for all clients",
    "  Step 2: Add CBT skills if insufficient progress",
    "  Step 3: Add intensive outpatient or residential if needed",
    "  Step 4: Long-term aftercare (meetings, booster sessions)",
    "",
    "Family Involvement in MET Sessions:",
    "  * Significant other can attend feedback session (with consent)",
    "  * Family member can provide collateral information",
    "  * Support person can be part of the Change Plan",
    "  * Miller (1999): 'The social network is a powerful motivator'",
]
make_standard_slide(prs, "Family & Caregiver Involvement", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': family_lines, 'bold_first': True},
], 8, "Ref: Meyers et al. (2002) CRAFT; NIMHANS (2016); Miller et al. (1999)")

# --- Slide 63: Combined Treatments ---
combined_lines = [
    "THERAPIES USED ALONGSIDE MET (Combined Treatment)",
    "",
    "MET + CBT: Most common and most researched combination",
    "  * MET first (builds motivation) then CBT (teaches skills)",
    "  * Used in COMBINE study for alcohol dependence",
    "  * Particularly effective for clients with skill deficits",
    "",
    "MET + Contingency Management (CM):",
    "  * MET provides motivation; CM provides immediate reinforcement",
    "  * Effective for stimulant and opioid use disorders",
    "",
    "MET + Pharmacotherapy:",
    "  * MET improves medication adherence",
    "  * Naltrexone + MET: better outcomes for alcohol (COMBINE, 2006)",
    "  * Nicotine Replacement + MI: improved quit rates",
    "",
    "MET + Group Therapy:",
    "  * Individual MET sessions prepare clients for group work",
    "  * Reduces dropout from group treatment",
    "",
    "Principle: MET works as a 'motivational booster' for ANY treatment.",
]
make_standard_slide(prs, "Therapies Used Alongside MET", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': combined_lines, 'bold_first': True},
], 8, "Ref: COMBINE Study (Anton et al., 2006); Carroll et al. (2006); NIMHANS (2016)")


# --- Slide 64: MET as Pretreatment ---
pretreat_lines = [
    "MET AS A PRETREATMENT ENHANCER",
    "",
    "What Does 'Pretreatment' Mean?",
    "  Using MET sessions BEFORE another therapy begins to:",
    "  * Increase engagement and retention in primary treatment",
    "  * Reduce dropout rates",
    "  * Build therapeutic alliance faster",
    "  * Resolve ambivalence before skills training starts",
    "",
    "Research Evidence:",
    "  * Bien et al. (1993): Brief MI before treatment increased completion by 50%",
    "  * Connors et al. (2002): MET pretreatment improved 12-month outcomes",
    "  * Carroll et al. (2006): MI + CBT > CBT alone for cannabis",
    "  * Hettema et al. (2005) meta-analysis: pretreatment MI effect size d=0.77",
    "",
    "Practical Application:",
    "  * 1-2 MET sessions before ANY therapy program",
    "  * Especially valuable for mandated/coerced clients",
    "  * Can reduce no-shows to first treatment session",
    "  * Cost-effective: adds only 1-2 hours of therapist time",
]
make_standard_slide(prs, "MET as Pretreatment Enhancer", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': pretreat_lines, 'bold_first': True},
], 8, "Ref: Bien et al. (1993); Hettema et al. (2005); Carroll et al. (2006)")


# ============================================================
# SECTION 9: COMORBID PSYCHIATRIC DISORDERS (Slides 65-72)
# ============================================================

# --- Slide 65: Section Divider ---
slide65 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide65, "SECTION 9\nMET for Comorbid\nPsychiatric Disorders", 9)

# --- Slide 66: Dual Diagnosis Overview ---
dd_lines = [
    "DUAL DIAGNOSIS: OVERVIEW & PREVALENCE",
    "",
    "Definition: Co-occurring substance use disorder AND another mental disorder.",
    "",
    "Prevalence (Key Statistics):",
    "  * 50-60% of people with SUD have a comorbid mental disorder",
    "  * 20-30% of people with mental disorders have comorbid SUD",
    "  * Depression + SUD: ~40% co-occurrence",
    "  * Anxiety + SUD: ~30% co-occurrence",
    "  * Personality Disorders + SUD: ~50-70% (ASPD highest)",
    "  * PTSD + SUD: ~25-40% co-occurrence",
    "",
    "Why This Matters for MET:",
    "  * Comorbidity complicates motivation and readiness",
    "  * Mental health symptoms may maintain substance use",
    "  * Treatment must address both conditions",
    "  * MET can be adapted for dual diagnosis populations",
    "  * Standard MET protocols may need modification",
    "",
    "NIMHANS Data: Comorbidity is the RULE, not the exception, in Indian settings.",
]
make_standard_slide(prs, "Dual Diagnosis Overview - Prevalence", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': dd_lines, 'bold_first': True},
], 9, "Ref: NIMHANS SUD Manual (2016); Regier et al. (1990); Kessler et al. (1994)")


# --- Slide 67: Personality Disorders + MET ---
pd_lines = [
    "PERSONALITY DISORDERS + MET: ADAPTATIONS & RESEARCH",
    "",
    "Challenges:",
    "  * Interpersonal difficulties affect therapeutic alliance",
    "  * Borderline PD: Emotional instability, splitting, impulsivity",
    "  * Antisocial PD: Low empathy, manipulation, treatment resistance",
    "  * Narcissistic PD: Defensiveness, difficulty with feedback",
    "",
    "MET Adaptations:",
    "  * Longer rapport-building phase (patience is critical)",
    "  * Avoid power struggles at all costs (especially ASPD)",
    "  * Use more affirmations (builds damaged self-image)",
    "  * Focus on client-defined values, not therapist agenda",
    "  * Set clear boundaries without being punitive",
    "  * Expect slower progress - adjust expectations",
    "  * Address relationship patterns that maintain use",
    "",
    "Research: Ball et al. (2007): Dual-focus Schema Therapy + MI",
    "improved outcomes in PD + SUD. MET helps engage PD clients who",
    "typically drop out of confrontational programs.",
]
make_standard_slide(prs, "Personality Disorders + MET Adaptations", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': pd_lines, 'bold_first': True},
], 9, "Ref: Ball et al. (2007); Martino et al. (2002); NIMHANS (2016)")

# --- Slide 68: Depression + MET ---
dep_lines = [
    "DEPRESSION + MET: ADAPTATIONS & RESEARCH",
    "",
    "Challenges:",
    "  * Low energy, hopelessness reduce motivation for change",
    "  * Anhedonia: Substance use may be only source of pleasure",
    "  * Negative thinking: 'I can't change' / 'What's the point?'",
    "  * Suicidal ideation requires careful monitoring",
    "  * Withdrawal from social support",
    "",
    "MET Adaptations:",
    "  * Extra emphasis on self-efficacy (counter hopelessness)",
    "  * Smaller, more achievable goals in Change Plan",
    "  * Acknowledge depression openly as part of the picture",
    "  * Connect substance reduction to mood improvement",
    "  * Provide psychoeducation: alcohol is a depressant",
    "  * Safety planning if suicidal ideation present",
    "  * More frequent check-ins between sessions",
    "",
    "Research: Baker et al. (2010): MI + CBT for depression + SUD",
    "reduced both drinking and depression symptoms at 12 months.",
]
make_standard_slide(prs, "Depression + MET Adaptations", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': dep_lines, 'bold_first': True},
], 9, "Ref: Baker et al. (2010); Hides et al. (2010); NIMHANS (2016)")


# --- Slide 69: Anxiety & PTSD + MET ---
anx_lines = [
    "ANXIETY & PTSD + MET: ADAPTATIONS & RESEARCH",
    "",
    "Challenges (Anxiety):",
    "  * Substance use as self-medication for anxiety",
    "  * Fear of withdrawal symptoms (anxiety about anxiety)",
    "  * Social anxiety may prevent group treatments",
    "  * Avoidance patterns complicate engagement",
    "",
    "Challenges (PTSD):",
    "  * Substances numb traumatic memories",
    "  * Trauma disclosure requires extreme safety",
    "  * Hypervigilance makes trust-building slower",
    "  * Re-traumatization risk with confrontational approaches",
    "",
    "MET Adaptations:",
    "  * Create exceptionally safe therapeutic environment",
    "  * Psychoeducation: 'using to cope with feelings makes sense'",
    "  * Develop alternative coping BEFORE reducing substance",
    "  * Don't force trauma processing - respect avoidance initially",
    "  * For PTSD: Integrated MI + trauma therapy (COPE model)",
    "",
    "Research: Sannibale et al. (2013): MI + CBT for comorbid PTSD/AUD effective.",
]
make_standard_slide(prs, "Anxiety & PTSD + MET Adaptations", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': anx_lines, 'bold_first': True},
], 9, "Ref: Sannibale et al. (2013); Back et al. (2014) COPE; NIMHANS (2016)")

# --- Slide 70: Bipolar & Psychotic Disorders ---
bip_lines = [
    "BIPOLAR & PSYCHOTIC DISORDERS + MET ADAPTATIONS",
    "",
    "BIPOLAR DISORDER + SUD:",
    "  * Substance use often escalates during manic episodes",
    "  * Medication non-adherence common",
    "  * MET can address medication motivation alongside SUD",
    "  * Adaptations: Simpler language during episodes, shorter sessions",
    "  * Graeber et al. (2003): MI improved lithium adherence",
    "",
    "PSYCHOTIC DISORDERS + SUD:",
    "  * 40-50% of people with schizophrenia have comorbid SUD",
    "  * Cognitive impairment requires simpler, slower approach",
    "  * Paranoia may affect trust (extra rapport-building needed)",
    "  * MET Adaptations for psychosis:",
    "    - Use concrete, simple language",
    "    - Shorter sessions (20-30 min instead of 60)",
    "    - More repetition and visual aids",
    "    - Focus on immediate, tangible benefits",
    "    - Integrate with psychiatric medication management",
    "",
    "Research: Barrowclough et al. (2010): MI + CBT for psychosis + SUD - effective.",
]
make_standard_slide(prs, "Bipolar & Psychotic Disorders + MET", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': bip_lines, 'bold_first': True},
], 9, "Ref: Graeber et al. (2003); Barrowclough et al. (2010); NIMHANS (2016)")


# --- Slide 71: Summary Table - All Comorbid Adaptations ---
summary_comorbid = [
    "SUMMARY TABLE: MET ADAPTATIONS FOR COMORBID DISORDERS",
    "",
    "Disorder         | Key Challenge           | MET Adaptation                  ",
    "-----------------|-------------------------|----------------------------------",
    "Depression       | Hopelessness, low energy| Extra self-efficacy, small goals ",
    "Anxiety          | Self-medication, fear   | Safety first, alternative coping ",
    "PTSD             | Numbing, trust issues   | Slow pace, no forced disclosure  ",
    "Borderline PD    | Instability, splitting  | Consistency, clear boundaries    ",
    "Antisocial PD    | Power struggles         | Avoid confrontation, use autonomy",
    "Bipolar          | Episode-driven use      | Medication adherence focus       ",
    "Psychosis        | Cognitive impairment    | Simple language, shorter sessions",
    "ADHD             | Impulsivity, boredom    | Structure, novelty, brief tasks  ",
    "",
    "Universal Principles Across All Comorbidities:",
    "  * Maintain MI spirit (empathy, collaboration, autonomy)",
    "  * Adapt pace and complexity to cognitive capacity",
    "  * Address the RELATIONSHIP between conditions",
    "  * Integrated treatment > sequential treatment",
    "  * Safety always comes first",
]
make_standard_slide(prs, "Summary: Comorbid Adaptations Table", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': summary_comorbid, 'bold_first': True},
], 9, "Ref: NIMHANS SUD Manual (2016); Martino et al. (2002); Mueser et al. (2003)")

# --- Slide 72: General Principles Integrated Treatment ---
integ_lines = [
    "GENERAL PRINCIPLES FOR INTEGRATED TREATMENT",
    "",
    "1. TREAT BOTH CONDITIONS SIMULTANEOUSLY",
    "   * Don't wait for one to resolve before addressing the other",
    "   * Substance use and mental health interact bidirectionally",
    "",
    "2. ONE TREATMENT TEAM / COORDINATED CARE",
    "   * Avoid split treatment where providers don't communicate",
    "   * MET therapist should know the psychiatric treatment plan",
    "",
    "3. STAGE-MATCHED INTERVENTIONS",
    "   * Match both SUD readiness AND mental health readiness",
    "   * Client may be in action for depression but precontemplation for SUD",
    "",
    "4. HARM REDUCTION ORIENTATION",
    "   * Any positive change is worthwhile",
    "   * Abstinence is ideal but not the only acceptable goal",
    "",
    "5. LONG-TERM PERSPECTIVE",
    "   * Both conditions are chronic/relapsing",
    "   * Recovery is a process, not an event",
    "   * Ongoing monitoring and booster sessions recommended",
]
make_standard_slide(prs, "General Principles for Integrated Treatment", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': integ_lines, 'bold_first': True},
], 9, "Ref: NIMHANS SUD Manual (2016); Mueser et al. (2003); Drake et al. (2001)")


# ============================================================
# SECTION 10: RESEARCH & EFFECTIVENESS (Slides 73-80)
# ============================================================

# --- Slide 73: Section Divider ---
slide73 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide73, "SECTION 10\nResearch Evidence &\nEffectiveness of MET", 10)

# --- Slide 74: Project MATCH Results ---
match_lines = [
    "PROJECT MATCH RESULTS - The Landmark Study",
    "",
    "Study Design (1989-1997):",
    "  * 1,726 clients with alcohol use disorders",
    "  * 9 sites across USA, randomized controlled trial",
    "  * Three treatments: MET (4 sessions), CBT (12), TSF (12)",
    "  * Follow-up: 1 year, then 3 years",
    "",
    "Key Findings:",
    "  * ALL THREE treatments produced significant improvement",
    "  * MET achieved COMPARABLE outcomes to CBT and TSF",
    "  * MET required only 33% of the sessions!",
    "  * At 3-year follow-up: no significant differences between groups",
    "",
    "Matching Findings (Client-Treatment):",
    "  * Angry clients: BETTER with MET (less confrontation)",
    "  * Low readiness clients: MET outperformed others",
    "  * High social support: MET worked especially well",
    "  * Network supportive of drinking: TSF worked better",
    "",
    "Conclusion: MET is a cost-effective, time-efficient treatment.",
]
make_standard_slide(prs, "Project MATCH Results", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': match_lines, 'bold_first': True},
], 10, "Ref: Project MATCH Research Group (1997, 1998); Matching hypotheses report")

# --- Slide 75: UKATT & International Evidence ---
ukatt_lines = [
    "UKATT & INTERNATIONAL EVIDENCE",
    "",
    "UKATT Study (UK Alcohol Treatment Trial, 2005):",
    "  * 742 clients across 7 UK sites",
    "  * Compared: MET (3 sessions) vs Social Behaviour & Network Therapy (8)",
    "  * Result: MET equally effective with FEWER sessions",
    "  * Cost analysis: MET significantly more cost-effective",
    "",
    "WHO Brief Intervention Study (International):",
    "  * 10 countries, primary care settings",
    "  * Brief motivational intervention (single session) effective",
    "  * 20% reduction in alcohol consumption at follow-up",
    "",
    "Australian SHADE Project (2010):",
    "  * MI + CBT for comorbid depression + alcohol",
    "  * Superior to standard treatment on both outcomes",
    "",
    "European Studies:",
    "  * Netherlands: MI effective for cannabis use (De Jonge et al., 2016)",
    "  * Sweden: MI in emergency departments reduced drinking (Soderstrom et al.)",
    "  * Multi-site reviews confirm cross-cultural effectiveness",
]
make_standard_slide(prs, "UKATT & International Evidence", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': ukatt_lines, 'bold_first': True},
], 10, "Ref: UKATT Research Team (2005); WHO Brief Intervention Study; Baker et al. (2010)")


# --- Slide 76: Meta-Analyses Summary ---
meta_lines = [
    "META-ANALYSES SUMMARY TABLE",
    "",
    "Study                  | Year | Finding",
    "-----------------------|------|--------------------------------------------------",
    "Bien et al.            | 1993 | Brief interventions effective (review of 32 studies)",
    "Noonan & Moyers        | 1997 | MI effective across substances",
    "Burke et al.           | 2003 | 30 RCTs: MI effect size d=0.25-0.57",
    "Hettema et al.         | 2005 | 72 studies: MI works for alcohol, drugs, diet",
    "                       |      | Effect size: d=0.77 as pretreatment adjunct",
    "Lundahl et al.         | 2010 | 119 studies: MI > no treatment, = active treatments",
    "                       |      | With fewer sessions needed",
    "Smedslund et al.       | 2011 | Cochrane: MI reduces substance use vs no treatment",
    "DiClemente et al.      | 2017 | MI retains effectiveness in diverse populations",
    "",
    "Consistent Findings Across Meta-Analyses:",
    "  * MI/MET is effective (better than no treatment)",
    "  * MI/MET is as effective as other active treatments",
    "  * MI/MET achieves results in fewer sessions (cost-effective)",
    "  * Effects are durable (maintained at long-term follow-up)",
    "  * Works across cultures, substances, and settings",
]
make_standard_slide(prs, "Meta-Analyses Summary Table", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': meta_lines, 'bold_first': True},
], 10, "Ref: Hettema et al. (2005); Lundahl et al. (2010); Smedslund et al. (2011, Cochrane)")

# --- Slide 77: Indian Research ---
indian_lines = [
    "INDIAN RESEARCH: NIMHANS, PREMIUM, AIIMS",
    "",
    "NIMHANS Contributions (Bangalore):",
    "  * NIMHANS SUD Manual (2016): Comprehensive treatment guidelines",
    "  * MI/MET integrated into standard SUD treatment protocols",
    "  * Training programs for counselors across India",
    "  * Brief interventions adapted for Indian cultural context",
    "",
    "PREMIUM Project (Programme for Effective Mental health Interventions):",
    "  * Patel et al. (2014): Task-shifting MI to lay counselors in Goa",
    "  * PREMIUM Alcohol: Brief MI by non-specialists - EFFECTIVE",
    "  * Proved MI works even when delivered by trained community workers",
    "  * Lancet publication: Global mental health milestone",
    "",
    "AIIMS (Delhi) & Other Indian Centers:",
    "  * Adapted MI protocols for tobacco cessation (Varghese et al.)",
    "  * MI for opioid use in Punjab (community health workers)",
    "  * Integration of MI in ICTC (HIV prevention) settings",
    "",
    "Key Insight: MI/MET is feasible and effective in Indian public health settings.",
]
make_standard_slide(prs, "Indian Research: NIMHANS, PREMIUM, AIIMS", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': indian_lines, 'bold_first': True},
], 10, "Ref: NIMHANS SUD Manual (2016); Patel et al. (2014, Lancet); Varghese et al.")


# --- Slide 78: Evidence by Population ---
pop_lines = [
    "EVIDENCE BY POPULATION - WHO BENEFITS FROM MET?",
    "",
    "Population              | Evidence Level | Key Findings",
    "------------------------|---------------|-----------------------------------",
    "Adults with AUD         | Strong (A)    | Project MATCH, UKATT - comparable",
    "Adults with drug use    | Strong (A)    | Cannabis, cocaine, opioids",
    "Adolescents             | Strong (A)    | Brief MI effective (Jensen et al.)",
    "Dual diagnosis          | Moderate (B)  | Adapted MI effective (Baker et al.)",
    "Mandated/coerced clients| Strong (A)    | MI reduces resistance significantly",
    "Emergency Dept. patients| Strong (A)    | Brief MI reduces repeat visits",
    "Pregnant women          | Moderate (B)  | Reduced substance use in pregnancy",
    "College students        | Strong (A)    | BASICS program (Dimeff et al.)",
    "HIV+ individuals        | Moderate (B)  | Improved ART adherence + reduced use",
    "Smokers/Tobacco         | Strong (A)    | MI + NRT improves quit rates",
    "Gamblers                | Moderate (B)  | Brief MI reduces gambling severity",
    "Older adults            | Emerging      | Few studies, promising results",
    "",
    "Key Pattern: MI/MET works across populations, substances, and settings.",
    "It is one of the most versatile evidence-based approaches available.",
]
make_standard_slide(prs, "Evidence by Population Table", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': pop_lines, 'bold_first': True},
], 10, "Ref: Hettema et al. (2005); Lundahl et al. (2010); NIMHANS (2016)")

# --- Slide 79: Key Takeaways ---
takeaway_lines = [
    "KEY TAKEAWAYS - 6 Essential Points",
    "",
    "1. MET IS BRIEF BUT POWERFUL",
    "   4 sessions achieve what others need 12+ to accomplish.",
    "   Cost-effective, time-efficient, resource-friendly.",
    "",
    "2. MOTIVATION IS ELICITED, NOT INSTALLED",
    "   People already have reasons to change. MET helps them find and voice these.",
    "",
    "3. THE THERAPIST'S STYLE MATTERS MORE THAN TECHNIQUE",
    "   Empathy, warmth, and collaboration predict outcomes more than any tool.",
    "",
    "4. RESISTANCE IS A SIGNAL, NOT A CHARACTER FLAW",
    "   When resistance rises, change YOUR approach, not the client.",
    "",
    "5. CHANGE IS A PROCESS, NOT AN EVENT",
    "   Stages are cyclical. Relapse is learning. Patience is essential.",
    "",
    "6. MET WORKS ACROSS POPULATIONS AND SETTINGS",
    "   From rural India to urban ERs, from adolescents to elderly,",
    "   from alcohol to gambling - the principles are universal.",
]
make_standard_slide(prs, "Key Takeaways - 6 Essential Points", [
    {'x': LEFT_MARGIN, 'y': CONTENT_Y, 'w': FULL_WIDTH, 'h': Inches(4.9),
     'lines': takeaway_lines, 'bold_first': True},
], 10, "Ref: Miller et al. (1992); Miller & Rollnick (2002); NIMHANS (2016)")


# --- Slide 80: Thank You + References ---
slide80 = prs.slides.add_slide(prs.slide_layouts[6])
# Dark blue background
bg80 = slide80.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
set_shape_fill(bg80, COLORS['navy'])
bg80.line.fill.background()

# Gold accent
acc80 = slide80.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.6), Inches(9.333), Inches(0.06))
set_shape_fill(acc80, COLORS['gold'])
acc80.line.fill.background()

# Thank you text
tb80 = slide80.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10.333), Inches(2.0))
tf80 = tb80.text_frame
tf80.word_wrap = True
p80 = tf80.paragraphs[0]
p80.alignment = PP_ALIGN.CENTER
r80 = p80.add_run()
r80.text = "Thank You"
r80.font.name = FONT_NAME
r80.font.size = Pt(48)
r80.font.bold = True
r80.font.color.rgb = COLORS['gold']

p80b = tf80.add_paragraph()
p80b.alignment = PP_ALIGN.CENTER
r80b = p80b.add_run()
r80b.text = "Motivational Enhancement Therapy - A Comprehensive Guide"
r80b.font.name = FONT_NAME
r80b.font.size = Pt(20)
r80b.font.color.rgb = COLORS['white']

# References box
ref_box = slide80.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.0), Inches(11.733), Inches(3.8)
)
set_shape_fill(ref_box, RGBColor(0x0A, 0x2A, 0x5E))
ref_box.line.color.rgb = COLORS['gold']
ref_box.line.width = Pt(2)

tf_ref = ref_box.text_frame
tf_ref.word_wrap = True
tf_ref.margin_left = Inches(0.2)
tf_ref.margin_right = Inches(0.2)
tf_ref.margin_top = Inches(0.15)

ref_texts = [
    "PRIMARY REFERENCES:",
    "",
    "Miller, W.R., Zweben, A., DiClemente, C.C., & Rychtarik, R.G. (1992). Motivational Enhancement",
    "  Therapy Manual. NIAAA Project MATCH Monograph Series, Vol. 2.",
    "NIMHANS (2016). Substance Use Disorders: Manual for Physicians. Bangalore: NIMHANS.",
    "Miller, W.R. & Rollnick, S. (2002). Motivational Interviewing (2nd ed.). Guilford Press.",
    "Project MATCH Research Group (1997). Matching treatments to client heterogeneity. JSAD, 58.",
    "UKATT Research Team (2005). UK Alcohol Treatment Trial. BMJ, 331, 541.",
    "Prochaska, J.O. & DiClemente, C.C. (1983). Stages of change. J Consulting & Clin Psych, 51.",
    "Hettema, J., Steele, J., & Miller, W.R. (2005). MI meta-analysis. Annual Rev Clin Psych, 1.",
    "Lundahl, B.W. et al. (2010). MI meta-analysis. Clinical Psychology Review, 30(1).",
    "Patel, V. et al. (2014). PREMIUM: Lancet, 389(10065).",
]
for i, line in enumerate(ref_texts):
    if i == 0:
        p = tf_ref.paragraphs[0]
    else:
        p = tf_ref.add_paragraph()
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = line
    run.font.name = FONT_NAME
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS['light_blue']
    if i == 0:
        run.font.bold = True
        run.font.color.rgb = COLORS['gold']


# ============================================================
# SAVE PRESENTATION
# ============================================================

output_path = "/projects/sandbox/Dango-kiro/MET_Comprehensive_Presentation.pptx"
prs.save(output_path)

total_slides = len(prs.slides)
print(f"\n{'='*60}")
print(f"  MET Comprehensive Presentation Generated Successfully!")
print(f"{'='*60}")
print(f"  Total Slides: {total_slides}")
print(f"  Output File: {output_path}")
print(f"  Font: Times New Roman throughout")
print(f"  Layout: Widescreen 13.333 x 7.5 inches")
print(f"  Sections: 10 complete sections")
print(f"{'='*60}")

if total_slides != 80:
    print(f"\n  WARNING: Expected 80 slides, got {total_slides}!")
    print(f"  Difference: {80 - total_slides} slides")
else:
    print(f"\n  CONFIRMED: Exactly 80 slides generated!")
